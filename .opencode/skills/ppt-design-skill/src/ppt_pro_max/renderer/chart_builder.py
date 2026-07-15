"""Chart Builder — python-pptx chart generation with theme colors."""

from __future__ import annotations

from typing import Any

try:
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


_LEGACY_CHART_TYPE_MAP: dict[str, Any] = {}
_ENTERPRISE_CHART_TYPE_MAP: dict[str, Any] = {}
if _PPTX_AVAILABLE:
    _LEGACY_CHART_TYPE_MAP = {
        "Bar Chart Vertical": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "Bar Chart Horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
        "Line Chart": XL_CHART_TYPE.LINE,
        "Area Chart": XL_CHART_TYPE.AREA,
        "Pie Chart": XL_CHART_TYPE.PIE,
        "Doughnut Chart": XL_CHART_TYPE.DOUGHNUT,
        "Scatter Plot": XL_CHART_TYPE.XY_SCATTER,
        "Radar Chart": XL_CHART_TYPE.RADAR_FILLED,
    }
    _ENTERPRISE_CHART_TYPE_MAP = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "bar_100": XL_CHART_TYPE.COLUMN_STACKED_100,
        "bar_3d": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "pie_3d": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
    }


class ChartBuilder:
    def build(
        self,
        slide,
        chart_type_or_config,
        data=None,
        style=None,
        position=None,
        brand_colors=None,
    ) -> Any:
        if not _PPTX_AVAILABLE:
            return None

        if isinstance(chart_type_or_config, dict):
            return self._build_from_config(slide, chart_type_or_config, position, brand_colors)

        return self._build_legacy(slide, chart_type_or_config, data, style, position)

    def _build_legacy(
        self,
        slide,
        chart_type: str,
        data: dict[str, Any],
        style: dict[str, Any],
        position: dict[str, float],
    ) -> Any:
        pptx_type = _LEGACY_CHART_TYPE_MAP.get(chart_type)
        if pptx_type is None:
            return None

        chart_data = CategoryChartData()
        chart_data.categories = data.get("labels", ["Q1", "Q2", "Q3", "Q4"])

        values = data.get("values", [10, 25, 45, 80])
        if isinstance(values, list) and isinstance(values[0], list):
            for i, series in enumerate(values):
                chart_data.add_series(f"Series {i + 1}", series)
        else:
            chart_data.add_series("Data", values)

        try:
            chart_frame = slide.shapes.add_chart(
                pptx_type,
                Inches(position.get("x", 1.5)),
                Inches(position.get("y", 1.5)),
                Inches(position.get("width", 10.333)),
                Inches(position.get("height", 4.5)),
                chart_data,
            )

            chart = chart_frame.chart
            chart.has_legend = False

            plot = chart.plots[0]
            colors = style.get("colors", {})
            primary_color = colors.get("primary", "#2563EB")

            try:
                series = plot.series[0]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor.from_string(primary_color.lstrip("#"))
            except Exception:
                pass

            return chart_frame
        except Exception:
            return None

    def _build_from_config(
        self,
        slide,
        chart_config: dict[str, Any],
        position: dict[str, float] | None = None,
        brand_colors: dict[str, str] | None = None,
    ) -> Any:
        chart_type_str = chart_config.get("type", "bar")
        pptx_type = _ENTERPRISE_CHART_TYPE_MAP.get(chart_type_str)
        if pptx_type is None:
            pptx_type = _ENTERPRISE_CHART_TYPE_MAP.get("bar")

        is_xy_chart = chart_type_str == "scatter"

        try:
            if is_xy_chart:
                chart_data = XyChartData()
                series_list = chart_config.get("series", [])
                if series_list:
                    for s_idx, s in enumerate(series_list):
                        xy_series = chart_data.add_series(s.get("name", f"Series {s_idx + 1}"))
                        for cat_idx, val in enumerate(s.get("values", [])):
                            xy_series.add_data_point(cat_idx + 1, val)
                else:
                    values = chart_config.get("values", [10, 25, 45, 80])
                    xy_series = chart_data.add_series("Data")
                    for i, v in enumerate(values):
                        xy_series.add_data_point(i + 1, v)
            else:
                chart_data = CategoryChartData()
                chart_data.categories = chart_config.get("categories", ["Q1", "Q2", "Q3", "Q4"])

                series_list = chart_config.get("series", [])
                if series_list:
                    for s in series_list:
                        chart_data.add_series(s.get("name", "Data"), s.get("values", []))
                else:
                    values = chart_config.get("values", [10, 25, 45, 80])
                    if isinstance(values, list) and len(values) > 0 and isinstance(values[0], list):
                        for i, series in enumerate(values):
                            chart_data.add_series(f"Series {i + 1}", series)
                    else:
                        chart_data.add_series("Data", values)

            if position is None:
                position = {}

            chart_frame = slide.shapes.add_chart(
                pptx_type,
                Inches(position.get("x", 1.5)),
                Inches(position.get("y", 1.5)),
                Inches(position.get("width", 10.333)),
                Inches(position.get("height", 4.5)),
                chart_data,
            )

            chart = chart_frame.chart

            style = chart_config.get("style", {})
            show_legend = style.get("show_legend", True)
            show_labels = style.get("show_labels", False)
            legend_position = style.get("legend_position", "bottom")
            color_scheme = style.get("color_scheme", "brand")

            chart.has_legend = show_legend
            if show_legend:
                pos_map = {
                    "bottom": XL_LEGEND_POSITION.BOTTOM,
                    "top": XL_LEGEND_POSITION.TOP,
                    "left": XL_LEGEND_POSITION.LEFT,
                    "right": XL_LEGEND_POSITION.RIGHT,
                }
                chart.legend.position = pos_map.get(legend_position, XL_LEGEND_POSITION.BOTTOM)

            if show_labels:
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.font.size = Pt(9)

            num_series = len(chart_config.get("series", []))
            if num_series == 0:
                num_series = 1
            chart_colors = self._resolve_colors(color_scheme, brand_colors, num_series)
            plot = chart.plots[0]
            for i, series in enumerate(plot.series):
                try:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor.from_string(
                        chart_colors[i % len(chart_colors)].lstrip("#")
                    )
                except Exception:
                    pass

            chart_title = chart_config.get("title")
            if chart_title:
                chart.has_title = True
                chart.chart_title.text_frame.paragraphs[0].text = chart_title

            return chart_frame
        except Exception:
            return None

    def _resolve_colors(self, color_scheme, brand_colors: dict[str, str] | None, num_series: int) -> list[str]:
        if color_scheme == "auto":
            return ["#4472C4", "#ED7D31", "#A5A5A5", "#FFC000", "#5B9BD5", "#70AD47"]

        if color_scheme == "brand" and brand_colors:
            primary = brand_colors.get("primary", "#2563EB")
            secondary = brand_colors.get("secondary", "#64748B")
            accent = brand_colors.get("accent", "#F97316")
            base = [primary, secondary, accent]
            while len(base) < num_series:
                base.append(self._rotate_hue(base[0], len(base) * 60))
            return base[:num_series]

        if isinstance(color_scheme, list):
            return color_scheme

        return ["#2563EB", "#F97316", "#10B981", "#8B5CF6", "#EF4444", "#06B6D4"]

    @staticmethod
    def _rotate_hue(hex_color: str, degrees: int) -> str:
        import colorsys

        r = int(hex_color[1:3], 16) / 255.0
        g = int(hex_color[3:5], 16) / 255.0
        b = int(hex_color[5:7], 16) / 255.0
        h, lum, s = colorsys.rgb_to_hls(r, g, b)
        h = (h + degrees / 360.0) % 1.0
        r2, g2, b2 = colorsys.hls_to_rgb(h, lum, s)
        return f"#{int(r2 * 255):02X}{int(g2 * 255):02X}{int(b2 * 255):02X}"
