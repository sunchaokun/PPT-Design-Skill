"""SlideExtractor — extract content + layout from existing PPT.

Used by the beautify mode to extract content from a client's PPT,
then re-render with new visual style using PrecisionRenderer.

Extraction flow per slide:
1. Iterate shapes → identify title, subtitle, bullets, images, tables, groups, SmartArt
2. Record shape positions → layout_hint
3. Infer goal from content semantics
4. Extract speaker notes
5. Return pages[] compatible with PrecisionRenderer.render_slide()
"""

from __future__ import annotations

import os
import tempfile
import zipfile
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches, Pt

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

_SKIP_PT_TYPES = {"doc", "parTrans", "sibTrans", "pres"}

_ALG_TYPE_MAP = {
    "lin": "process",
    "snake": "process",
    "cycle": "cycle",
    "hier": "hierarchy",
    "pyra": "pyramid",
    "matrix": "matrix",
    "seg": "relationship",
    "pic": "picture",
}

_CTA_KEYWORDS = {"thank", "thanks", "谢谢", "联系", "contact", "questions", "next step"}
_PROBLEM_KEYWORDS = {"problem", "challenge", "痛点", "问题", "挑战", "pain point"}
_SOLUTION_KEYWORDS = {"solution", "solve", "方案", "解决", "how we"}
_FEATURES_KEYWORDS = {"feature", "capability", "功能", "特性", "能力"}
_DATA_KEYWORDS = {"data", "metric", "数据", "指标", "kpi", "growth"}
_CODE_KEYWORDS = {"code", "demo", "代码", "示例", "quick start", "getting started"}
_MARKET_KEYWORDS = {"market", "tam", "市场", "规模"}


class SlideExtractor:

    def __init__(self, temp_dir: str | None = None):
        self._temp_dir = temp_dir or os.path.join(tempfile.gettempdir(), "ppt-extracted-images")
        os.makedirs(self._temp_dir, exist_ok=True)

    def extract(self, pptx_path: str) -> list[dict[str, Any]]:
        if not os.path.isfile(pptx_path):
            return []

        prs: Presentation | None = None
        try:
            prs = Presentation(pptx_path)
        except Exception:
            return []

        total_slides = len(prs.slides)
        pages: list[dict[str, Any]] = []

        smartart_map = self._extract_smartart_map(pptx_path)

        for idx, slide in enumerate(prs.slides):
            page = self._extract_slide(slide, idx, total_slides, smartart_map, pptx_path)
            pages.append(page)

        return pages

    def _extract_slide(
        self,
        slide,
        slide_idx: int,
        total_slides: int,
        smartart_map: dict,
        pptx_path: str,
    ) -> dict[str, Any]:
        title = ""
        subtitle = ""
        bullets: list[str] = []
        image_path: str | None = None
        chart_data: dict | None = None
        diagram_type: str | None = None
        diagram_data: dict | None = None
        notes: str | None = None
        complex_elements: list[dict] = []
        title_pos: tuple | None = None
        body_pos: tuple | None = None
        image_pos: tuple | None = None
        has_image = False
        is_full_bleed = False
        has_table = False

        for shape in slide.shapes:
            shape_type = shape.shape_type
            pos = self._emu_to_inches_tuple(shape)

            if shape.has_text_frame:
                text = shape.text.strip()
                if not text:
                    continue

                if not title and self._is_title_shape(shape, slide_idx):
                    title = text
                    title_pos = pos
                elif not subtitle and title and not bullets:
                    if self._is_likely_subtitle(shape):
                        subtitle = text
                    else:
                        for para in shape.text_frame.paragraphs:
                            ptext = para.text.strip()
                            if ptext:
                                bullets.append(ptext)
                        body_pos = pos
                else:
                    for para in shape.text_frame.paragraphs:
                        ptext = para.text.strip()
                        if ptext:
                            bullets.append(ptext)
                    if body_pos is None:
                        body_pos = pos

            elif shape_type == MSO_SHAPE_TYPE.PICTURE or (hasattr(shape, 'image') and shape.shape_type == 13):
                try:
                    img_blob = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    img_filename = f"slide{slide_idx + 1}_img_{shape.shape_id}.{ext}"
                    img_path = os.path.join(self._temp_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(img_blob)
                    image_path = img_path
                    has_image = True
                    image_pos = pos
                except Exception:
                    pass

            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                has_table = True
                table = shape.table
                rows_data = []
                for row in table.rows:
                    row_cells = []
                    for cell in row.cells:
                        row_cells.append(cell.text.strip())
                    rows_data.append(row_cells)
                diagram_type = "table"
                diagram_data = {"headers": rows_data[0] if rows_data else [], "rows": rows_data[1:] if len(rows_data) > 1 else []}

            elif shape_type == MSO_SHAPE_TYPE.CHART:
                try:
                    chart = shape.chart
                    chart_data = self._extract_chart_data(chart)
                except Exception:
                    pass

            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                group_info = self._extract_group(shape)
                if group_info:
                    group_info["bounds"] = pos
                    complex_elements.append(group_info)

            elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        try:
                            ph_type = shape.placeholder_format.type
                        except (ValueError, AttributeError):
                            ph_type = None
                        if not title and (ph_type == 1 or ph_type == 3):
                            title = text
                            title_pos = pos
                        elif not subtitle and ph_type == 4:
                            subtitle = text
                        else:
                            for para in shape.text_frame.paragraphs:
                                ptext = para.text.strip()
                                if ptext:
                                    bullets.append(ptext)
                            if body_pos is None:
                                body_pos = pos

        graphic_frames = self._find_smartart_graphic_frames(slide)
        for gf_shape, rel_ids in graphic_frames:
            dm_rid = rel_ids.get("dm")
            if dm_rid and dm_rid in smartart_map:
                sa_info = smartart_map[dm_rid]
                ce = {
                    "type": "smartart",
                    "category": sa_info.get("category", "process"),
                    "variant": sa_info.get("variant", ""),
                    "nodes": sa_info.get("nodes", []),
                    "xml_parts": sa_info.get("xml_parts", {}),
                    "bounds": self._emu_to_inches_tuple(gf_shape),
                }
                complex_elements.append(ce)

        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes = notes_text
        except Exception:
            pass

        goal = self._infer_goal(
            slide_idx, total_slides, title, bullets,
            has_image, has_table, chart_data, complex_elements,
        )

        if slide_idx == 0 and has_image and not bullets and not has_table:
            is_full_bleed = True

        layout_hint = {
            "title_pos": title_pos,
            "body_pos": body_pos,
            "image_pos": image_pos,
            "has_image": has_image,
            "bullet_count": len(bullets),
            "is_full_bleed": is_full_bleed,
        }

        return {
            "goal": goal,
            "title": title,
            "subtitle": subtitle or None,
            "bullets": bullets if bullets else None,
            "image": image_path,
            "cards": None,
            "diagram_type": diagram_type,
            "diagram_data": diagram_data,
            "code": None,
            "exercise": None,
            "chart": chart_data,
            "notes": notes,
            "links": [],
            "layout_hint": layout_hint,
            "complex_elements": complex_elements,
        }

    def _is_title_shape(self, shape, slide_idx: int) -> bool:
        try:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (1, 3):
                    return True
        except (ValueError, AttributeError):
            pass
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size >= Emu(Pt(28)):
                        return True
        if slide_idx == 0 and shape.top < Emu(Inches(2)):
            return True
        return False

    def _is_likely_subtitle(self, shape) -> bool:
        try:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == 4:
                    return True
        except (ValueError, AttributeError):
            pass
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and Pt(16) <= run.font.size < Pt(28):
                        return True
        return False

    def _emu_to_inches_tuple(self, shape) -> tuple[float, float, float, float]:
        try:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if left is None or top is None or width is None or height is None:
                return (0.0, 0.0, 0.0, 0.0)
            x = left / 914400
            y = top / 914400
            w = width / 914400
            h = height / 914400
            return (round(x, 2), round(y, 2), round(w, 2), round(h, 2))
        except Exception:
            return (0.0, 0.0, 0.0, 0.0)

    def _extract_group(self, group_shape) -> dict[str, Any] | None:
        texts: list[str] = []
        images: list[str] = []
        shapes_info: list[dict] = []

        try:
            for child in group_shape.shapes:
                child_type = child.shape_type
                if child.has_text_frame and child.text.strip():
                    texts.append(child.text.strip())
                if child_type == MSO_SHAPE_TYPE.PICTURE or (hasattr(child, 'image') and child.shape_type == 13):
                    try:
                        img_blob = child.image.blob
                        ext = child.image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        img_filename = f"group_img_{child.shape_id}.{ext}"
                        img_path = os.path.join(self._temp_dir, img_filename)
                        with open(img_path, "wb") as f:
                            f.write(img_blob)
                        images.append(img_path)
                    except Exception:
                        pass
                if hasattr(child, 'left') and hasattr(child, 'width'):
                    shapes_info.append({
                        "type": self._shape_type_name(child),
                        "bounds": self._emu_to_inches_tuple(child),
                    })
        except Exception:
            all_t = group_shape._element.findall(f".//{{{_NS['a']}}}t")
            texts = [t.text for t in all_t if t.text]

        category = self._infer_group_category(shapes_info, texts)

        return {
            "type": "group",
            "texts": texts,
            "images": images,
            "shapes": shapes_info,
            "inferred_category": category,
        }

    def _shape_type_name(self, shape) -> str:
        stype = shape.shape_type
        mapping = {
            MSO_SHAPE_TYPE.AUTO_SHAPE: "rectangle",
            MSO_SHAPE_TYPE.PICTURE: "picture",
            MSO_SHAPE_TYPE.GROUP: "group",
            MSO_SHAPE_TYPE.TABLE: "table",
            MSO_SHAPE_TYPE.CHART: "chart",
        }
        return mapping.get(stype, f"type_{stype}")

    def _infer_group_category(self, shapes_info: list[dict], texts: list[str]) -> str:
        if len(shapes_info) >= 3:
            xs = [s["bounds"][0] for s in shapes_info if s.get("bounds") and len(s["bounds"]) >= 2]
            ys = [s["bounds"][1] for s in shapes_info if s.get("bounds") and len(s["bounds"]) >= 2]
            if len(xs) >= 2:
                x_range = max(xs) - min(xs)
                y_range = max(ys) - min(ys) if ys else 0
                if x_range > y_range * 2:
                    return "process"
                if y_range > x_range * 2:
                    return "hierarchy"
        return "infographic"

    def _extract_chart_data(self, chart) -> dict:
        data = {"type": str(chart.chart_type), "series": []}
        try:
            for series in chart.series:
                values = []
                for val in series.values:
                    values.append(val)
                data["series"].append({"name": str(series.format), "values": values})
        except Exception:
            pass
        return data

    def _find_smartart_graphic_frames(self, slide) -> list[tuple[Any, dict]]:
        frames = []
        dgm_ns = _NS["dgm"]
        r_ns = _NS["r"]

        for shape in slide.shapes:
            try:
                elem = shape._element
                rel_ids_elems = elem.findall(f".//{{{dgm_ns}}}relIds")
                if not rel_ids_elems:
                    continue
                for ri in rel_ids_elems:
                    ids = {}
                    for attr in ("dm", "lo", "qs", "cs"):
                        val = ri.get(f"{{{r_ns}}}{attr}")
                        if val:
                            ids[attr] = val
                    if ids:
                        frames.append((shape, ids))
            except Exception:
                pass
        return frames

    def _extract_smartart_map(self, pptx_path: str) -> dict[str, dict]:
        result: dict[str, dict] = {}
        try:
            with zipfile.ZipFile(pptx_path) as z:
                names = z.namelist()
                slide_rels = sorted([
                    n for n in names
                    if n.startswith("ppt/slides/_rels/") and n.endswith(".rels")
                ])

                for sr in slide_rels:
                    content = z.read(sr)
                    root = etree.fromstring(content)
                    rel_ns = _NS["rel"]

                    dgm_refs: dict[str, str] = {}
                    for rel in root.findall(f"{{{rel_ns}}}Relationship"):
                        target = rel.get("Target", "")
                        rid = rel.get("Id", "")
                        if "diagram" in target.lower():
                            dgm_refs[rid] = target

                    if not dgm_refs:
                        continue

                    for rid, target in dgm_refs.items():
                        parts_dir = os.path.dirname(target)
                        abs_dir = os.path.normpath(f"ppt/slides/{parts_dir}").replace("\\", "/")
                        fname = os.path.basename(target)

                        if not fname.startswith("data"):
                            continue

                        abs_path = f"{abs_dir}/{fname}"
                        if abs_path not in names:
                            continue

                        data_xml = z.read(abs_path)
                        nodes = self._parse_data_xml(data_xml)

                        layout_fname = fname.replace("data", "layout")
                        layout_path = f"{abs_dir}/{layout_fname}"
                        category = "process"
                        variant = ""
                        if layout_path in names:
                            layout_xml = z.read(layout_path)
                            cat_var = self._parse_layout_xml(layout_xml)
                            category = cat_var.get("category", "process")
                            variant = cat_var.get("variant", "")

                        xml_parts: dict[str, bytes] = {}
                        for part_name in ("data", "layout", "colors", "quickStyle"):
                            part_fname = fname.replace("data", part_name)
                            part_path = f"{abs_dir}/{part_fname}"
                            if part_path in names:
                                xml_parts[part_name] = z.read(part_path)

                        result[rid] = {
                            "category": category,
                            "variant": variant,
                            "nodes": nodes,
                            "xml_parts": xml_parts,
                        }
        except Exception:
            pass

        return result

    def _parse_data_xml(self, data_xml: bytes) -> list[dict]:
        root = etree.fromstring(data_xml)
        a_ns = _NS["a"]
        dgm_ns = _NS["dgm"]
        nodes = []

        for pt in root.iter(f"{{{dgm_ns}}}pt"):
            pt_type = pt.get("type", "")
            if pt_type in _SKIP_PT_TYPES:
                continue
            t_elem = pt.find(f".//{{{a_ns}}}t")
            if t_elem is not None and t_elem.text:
                level = 0
                prlo = pt.find(f"{{{dgm_ns}}}prLo")
                if prlo is not None:
                    try:
                        level = int(prlo.get("val", "0"))
                    except ValueError:
                        pass
                nodes.append({"id": len(nodes), "text": t_elem.text, "level": level})

        return nodes

    def _parse_layout_xml(self, layout_xml: bytes) -> dict:
        root = etree.fromstring(layout_xml)
        dgm_ns = _NS["dgm"]

        category = "process"
        variant = ""

        alg_elems = root.findall(f".//{{{dgm_ns}}}alg")
        for alg in alg_elems:
            alg_type = alg.get("type", "")
            category = _ALG_TYPE_MAP.get(alg_type, category)
            break

        cat_elems = root.findall(f".//{{{dgm_ns}}}cat")
        for cat in cat_elems:
            cat_type = cat.get("type", "")
            if cat_type:
                variant = cat_type.split("/")[-1] if "/" in cat_type else cat_type
                break

        return {"category": category, "variant": variant}

    def _infer_goal(
        self,
        slide_idx: int,
        total_slides: int,
        title: str,
        bullets: list[str],
        has_image: bool,
        has_table: bool,
        chart_data: dict | None,
        complex_elements: list[dict],
    ) -> str:
        title_lower = title.lower()
        all_text = title_lower + " " + " ".join(bullets).lower()

        if slide_idx == 0 and total_slides > 1:
            return "hook"

        for kw in _CTA_KEYWORDS:
            if kw in title_lower or kw in all_text:
                if slide_idx == total_slides - 1:
                    return "cta"

        for kw in _PROBLEM_KEYWORDS:
            if kw in title_lower or kw in all_text:
                return "problem"

        for kw in _SOLUTION_KEYWORDS:
            if kw in title_lower or kw in all_text:
                return "solution"

        for kw in _FEATURES_KEYWORDS:
            if kw in title_lower or kw in all_text:
                return "features"

        for kw in _DATA_KEYWORDS:
            if kw in title_lower or kw in all_text:
                return "data"

        for kw in _CODE_KEYWORDS:
            if kw in title_lower or kw in all_text:
                return "code"

        for kw in _MARKET_KEYWORDS:
            if kw in title_lower or kw in all_text:
                return "market"

        if slide_idx == 0 and not bullets and not has_table and not chart_data:
            return "hook"

        if slide_idx == total_slides - 1:
            if not bullets and not has_table and not chart_data:
                return "cta"

        if chart_data:
            return "data"

        if has_table:
            return "data"

        if any(e["type"] == "smartart" for e in complex_elements):
            sa = next(e for e in complex_elements if e["type"] == "smartart")
            sa_cat = sa.get("category", "")
            if sa_cat in ("process", "cycle", "hierarchy", "pyramid", "matrix"):
                return "data"

        return "content"
