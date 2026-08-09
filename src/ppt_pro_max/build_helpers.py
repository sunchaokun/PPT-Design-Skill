"""build_helpers — LLM build script toolbox.

Usage in build scripts:
    from ppt_pro_max.build_helpers import *
    C = {'primary': '#2E6504', 'accent': '#7DA92F', ...}
    t = TYPOGRAPHY['mckinsey']
    sp = SPACING['mckinsey']
    prs = Presentation(template_path)
    s = add_slide(prs)
    page_header(s, 'Title', 'Subtitle', C, typo=t, spacing=sp)
    kpi_card(s, 1.0, 1.5, 3.0, 1.35, '12.8亿', '年度产值', '+8.3%', C=C, typo=t)
    prs.save('output.pptx')
"""
from __future__ import annotations

import copy

from lxml import etree
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation as _Presentation

from ppt_pro_max.renderer.text_effects import (
    apply_text_gradient, apply_text_gradient_preset, set_vertical_text,
    set_text_rotation, apply_text_outline,
)
from ppt_pro_max.renderer.blip_fill import (
    add_circle_image as _add_circle_image,
    add_image_in_shape as _add_image_in_shape,
    apply_blip_duotone, apply_blip_artistic,
)
from ppt_pro_max.renderer.visual_effects import apply_soft_edge
from ppt_pro_max.renderer.visual_effects import (
    apply_3d, apply_bevel, apply_pattern_fill, apply_frosted_glass,
)
from ppt_pro_max.renderer.decoration_library import (
    add_brush_divider as _add_brush_divider,
    add_neon_border as _add_neon_border,
    add_glass_panel as _add_glass_panel,
    add_grid_background as _add_grid_background,
    add_ink_splash as _add_ink_splash,
)
from ppt_pro_max.renderer.animation import (
    add_slide_transition, add_entrance_animation,
    add_exit_animation, add_emphasis_animation,
)


def _resolve_color(val, C):
    if val is None:
        return '#000000'
    if val.startswith('#'):
        return val
    return (C or {}).get(val, '#000000')


def _rgb(hex_str):
    return RGBColor.from_string(hex_str.lstrip('#'))


def _set_cjk_font(run, font_name):
    if not font_name:
        return
    rPr = run._r.find(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr'
    )
    if rPr is None:
        return
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ea = rPr.find(f'{{{ns}}}ea')
    if ea is None:
        ea = etree.SubElement(rPr, f'{{{ns}}}ea')
    ea.set('typeface', font_name)
    cs = rPr.find(f'{{{ns}}}cs')
    if cs is None:
        cs = etree.SubElement(rPr, f'{{{ns}}}cs')
    cs.set('typeface', font_name)


def _strip_style(shape):
    sp = shape._element
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    style_el = sp.find(f'{{{ns}}}style')
    if style_el is not None:
        sp.remove(style_el)


def _add_shape(shapes, mso_type, left, top, width, height):
    sh = shapes.add_shape(mso_type, left, top, width, height)
    _strip_style(sh)
    return sh


def _set_run(paragraph, txt, font_size=12, color='text_body', bold=False,
             font_name=None, C=None):
    run = paragraph.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(_resolve_color(color, C))
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    cjk_font = (C or {}).get('font_cjk') or (C or {}).get('font_body')
    if cjk_font:
        _set_cjk_font(run, cjk_font)
    return run


def _lighten(hex_color, amount=30):
    h = hex_color.lstrip('#')
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = min(255, r + amount)
    g = min(255, g + amount)
    b = min(255, b + amount)
    return f'{r:02X}{g:02X}{b:02X}'


class Typography:
    """Font size scale for a design style.

    Each style defines its own scale so that headings, body, and captions
    maintain consistent visual rhythm across all slides.
    """

    def __init__(self, hero=44, h1=28, h2=20, h3=16, body=12, caption=10, micro=8):
        self.hero = hero
        self.h1 = h1
        self.h2 = h2
        self.h3 = h3
        self.body = body
        self.caption = caption
        self.micro = micro

    def scale(self, level):
        return getattr(self, level, self.body)


class Spacing:
    """Spacing system for a design style.

    Controls margins, padding, gaps, and rhythm so that every component
    has consistent breathing room.
    """

    def __init__(self, page_margin=0.65, section_gap=0.5, card_gap=0.35,
                 card_padding=0.2, line_height=1.4, bar_gap=0.2):
        self.page_margin = page_margin
        self.section_gap = section_gap
        self.card_gap = card_gap
        self.card_padding = card_padding
        self.line_height = line_height
        self.bar_gap = bar_gap


TYPOGRAPHY = {
    'mckinsey': Typography(hero=44, h1=28, h2=20, h3=16, body=12, caption=10, micro=8),
    'cyberpunk': Typography(hero=48, h1=28, h2=18, h3=14, body=11, caption=9, micro=7),
    'creative': Typography(hero=44, h1=28, h2=22, h3=18, body=13, caption=11, micro=9),
    'professional': Typography(hero=44, h1=28, h2=20, h3=16, body=12, caption=10, micro=8),
    'minimal': Typography(hero=40, h1=24, h2=18, h3=14, body=11, caption=9, micro=7),
    'cjk_mckinsey': Typography(hero=44, h1=30, h2=22, h3=18, body=14, caption=12, micro=11),
    'cjk_professional': Typography(hero=44, h1=30, h2=22, h3=18, body=14, caption=12, micro=11),
    'cjk_creative': Typography(hero=44, h1=30, h2=24, h3=20, body=15, caption=13, micro=11),
}

SPACING = {
    'mckinsey': Spacing(page_margin=0.65, section_gap=0.5, card_gap=0.35,
                        card_padding=0.2, line_height=1.4, bar_gap=0.2),
    'cyberpunk': Spacing(page_margin=0.8, section_gap=0.6, card_gap=0.4,
                         card_padding=0.25, line_height=1.3, bar_gap=0.25),
    'creative': Spacing(page_margin=0.8, section_gap=0.6, card_gap=0.4,
                        card_padding=0.25, line_height=1.5, bar_gap=0.25),
    'professional': Spacing(page_margin=0.65, section_gap=0.5, card_gap=0.35,
                            card_padding=0.2, line_height=1.4, bar_gap=0.2),
    'minimal': Spacing(page_margin=1.0, section_gap=0.6, card_gap=0.5,
                       card_padding=0.3, line_height=1.5, bar_gap=0.3),
}


def set_widescreen(prs):
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sldSz = prs._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz'
    )
    if sldSz is not None and 'type' in sldSz.attrib:
        del sldSz.attrib['type']


def set_dark_theme(prs, C=None):
    C = C or {}
    from lxml import etree as _et
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    bg = C.get('background', '#0B1020')
    fg = C.get('text_dark', '#E2E8F0')
    theme_part = None
    for rel in prs.part.rels.values():
        if 'theme' in rel.reltype:
            theme_part = rel.target_part
            break
    if theme_part is None:
        return
    theme_el = _et.fromstring(theme_part.blob)
    clrScheme = theme_el.find(f'{{{ns}}}themeElements/{{{ns}}}clrScheme')
    if clrScheme is None:
        return
    for tag, val in [('dk1', fg), ('lt1', bg)]:
        el = clrScheme.find(f'{{{ns}}}{tag}')
        if el is not None:
            for child in list(el):
                el.remove(child)
            srgb = _et.SubElement(el, f'{{{ns}}}srgbClr')
            srgb.set('val', val.lstrip('#'))
    theme_part._blob = _et.tostring(theme_el, xml_declaration=True, encoding='UTF-8', standalone=True)


def clean_save(prs, path):
    import os, zipfile, shutil
    from lxml import etree as _et
    os.makedirs(os.path.dirname(path) if os.path.dirname(path) else '.', exist_ok=True)
    prs.save(path)
    tmp = path + '.tmp'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_ct = 'http://schemas.openxmlformats.org/package/2006/content-types'
    with zipfile.ZipFile(path, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            items = {}
            for item in zin.namelist():
                if 'printerSettings' in item:
                    continue
                items[item] = zin.read(item)
    rid_map = None
    for name, data in items.items():
        if name == 'ppt/_rels/presentation.xml.rels':
            root = _et.fromstring(data)
            to_remove = [r for r in root
                         if 'printerSettings' in r.get('Target', '')]
            if to_remove:
                for r in to_remove:
                    root.remove(r)
                rid_map = {}
                for r in root:
                    old = r.get('Id', '')
                    if old.startswith('rId'):
                        rid_map[old] = 'rId%d' % (len(rid_map) + 1)
                for r in root:
                    old = r.get('Id', '')
                    if old in rid_map:
                        r.set('Id', rid_map[old])
            data = _et.tostring(root, xml_declaration=True,
                                encoding='UTF-8', standalone=True)
            items[name] = data
    if rid_map:
        for name in list(items.keys()):
            if name == 'ppt/presentation.xml':
                root = _et.fromstring(items[name])
                for el in root.iter():
                    rid = el.get(f'{{{ns_r}}}id', '')
                    if rid in rid_map:
                        el.set(f'{{{ns_r}}}id', rid_map[rid])
                items[name] = _et.tostring(root, xml_declaration=True,
                                           encoding='UTF-8', standalone=True)
    ct_data = items.get('[Content_Types].xml')
    if ct_data is not None:
        ct_root = _et.fromstring(ct_data)
        to_remove = []
        for el in ct_root:
            if el.tag == f'{{{ns_ct}}}Default':
                ct_val = el.get('ContentType', '')
                if 'printerSettings' in ct_val:
                    to_remove.append(el)
        for el in to_remove:
            ct_root.remove(el)
        if to_remove:
            items['[Content_Types].xml'] = _et.tostring(
                ct_root, xml_declaration=True, encoding='UTF-8', standalone=True)
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in items.items():
            if name.endswith('.xml') and name.startswith('ppt/slides/slide'):
                try:
                    root = _et.fromstring(data)
                    changed = False
                    for ln in root.iter(f'{{{ns_a}}}ln'):
                        if len(ln) == 0:
                            _et.SubElement(ln, f'{{{ns_a}}}noFill')
                            changed = True
                    if changed:
                        data = _et.tostring(root, xml_declaration=True,
                                            encoding='UTF-8', standalone=True)
                except Exception:
                    pass
            zout.writestr(name, data)
    shutil.move(tmp, path)


def add_slide(prs, layout_index=None):
    if layout_index is not None:
        return prs.slides.add_slide(prs.slide_layouts[layout_index])
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[-1])


def rect(slide, left, top, width, height, fill, line=None, C=None):
    shape = _add_shape(slide.shapes, MSO_SHAPE.RECTANGLE,
                       Inches(left), Inches(top),
                       Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        shape.line.color.rgb = _rgb(_resolve_color(line, C))
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def rrect(slide, left, top, width, height, fill, line=None, C=None):
    shape = _add_shape(slide.shapes, MSO_SHAPE.ROUNDED_RECTANGLE,
                       Inches(left), Inches(top),
                       Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        shape.line.color.rgb = _rgb(_resolve_color(line, C))
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def oval(slide, left, top, width, height, fill, line=None, C=None):
    shape = _add_shape(slide.shapes, MSO_SHAPE.OVAL,
                       Inches(left), Inches(top),
                       Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        shape.line.color.rgb = _rgb(_resolve_color(line, C))
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def shape(slide, shape_type, left, top, width, height, fill, line=None, C=None):
    _type = shape_type
    if isinstance(_type, str):
        _type = getattr(MSO_SHAPE, _type.upper(), MSO_SHAPE.RECTANGLE)
    sh = _add_shape(slide.shapes, _type,
                    Inches(left), Inches(top),
                    Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        sh.line.color.rgb = _rgb(_resolve_color(line, C))
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def _centered_shape(slide, mso_type, cx, cy, width, height, fill, line=None, C=None):
    sh = _add_shape(slide.shapes, mso_type,
                    Inches(cx - width / 2), Inches(cy - height / 2),
                    Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(_resolve_color(fill, C))
    if line:
        sh.line.color.rgb = _rgb(_resolve_color(line, C))
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def hexagon(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.HEXAGON, cx, cy, size, size * 0.87,
                           fill, line, C)


def pentagon(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.REGULAR_PENTAGON, cx, cy, size, size,
                           fill, line, C)


def octagon(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.OCTAGON, cx, cy, size, size,
                           fill, line, C)


def diamond(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.DIAMOND, cx, cy, size, size,
                           fill, line, C)


def triangle(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height,
                 fill, line, C)


def right_triangle(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.RIGHT_TRIANGLE, left, top, width, height,
                 fill, line, C)


def parallelogram(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.PARALLELOGRAM, left, top, width, height,
                 fill, line, C)


def trapezoid(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.TRAPEZOID, left, top, width, height,
                 fill, line, C)


def star5(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_5_POINT, cx, cy, size, size,
                           fill, line, C)


def star6(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_6_POINT, cx, cy, size, size,
                           fill, line, C)


def star8(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_8_POINT, cx, cy, size, size,
                           fill, line, C)


def star10(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_10_POINT, cx, cy, size, size,
                           fill, line, C)


def star12(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.STAR_12_POINT, cx, cy, size, size,
                           fill, line, C)


def donut(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.DONUT, cx, cy, size, size,
                           fill, line, C)


def heart(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.HEART, cx, cy, size, size,
                           fill, line, C)


def cross(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.CROSS, cx, cy, size, size,
                           fill, line, C)


def arrow(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, height,
                 fill, line, C)


def chevron(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.CHEVRON, left, top, width, height,
                 fill, line, C)


def cloud(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.CLOUD, left, top, width, height,
                 fill, line, C)


def lightning(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.LIGHTNING_BOLT, left, top, width, height,
                 fill, line, C)


def gear(slide, cx, cy, size, fill, line=None, C=None, teeth=6):
    mso = MSO_SHAPE.GEAR_9 if teeth >= 9 else MSO_SHAPE.GEAR_6
    return _centered_shape(slide, mso, cx, cy, size, size, fill, line, C)


def funnel(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FUNNEL, left, top, width, height,
                 fill, line, C)


def moon(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.MOON, cx, cy, size, size,
                           fill, line, C)


def sun(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.SUN, cx, cy, size, size,
                           fill, line, C)


def wave(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.WAVE, left, top, width, height,
                 fill, line, C)


def block_arc(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.BLOCK_ARC, cx, cy, size, size,
                           fill, line, C)


def callout(slide, left, top, width, height, fill, line=None, C=None, style='rect'):
    _MAP = {
        'rect': MSO_SHAPE.RECTANGULAR_CALLOUT,
        'round': MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        'oval': MSO_SHAPE.OVAL_CALLOUT,
        'cloud': MSO_SHAPE.CLOUD_CALLOUT,
    }
    mso = _MAP.get(style, MSO_SHAPE.RECTANGULAR_CALLOUT)
    return shape(slide, mso, left, top, width, height, fill, line, C)


def flow_process(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FLOWCHART_PROCESS, left, top, width, height,
                 fill, line, C)


def flow_decision(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.FLOWCHART_DECISION, cx, cy, size, size,
                           fill, line, C)


def flow_data(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FLOWCHART_DATA, left, top, width, height,
                 fill, line, C)


def flow_document(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FLOWCHART_DOCUMENT, left, top, width, height,
                 fill, line, C)


def flow_connector(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.FLOWCHART_CONNECTOR, cx, cy, size, size,
                           fill, line, C)


def no_symbol(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.NO_SYMBOL, cx, cy, size, size,
                           fill, line, C)


def plaque(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.PLAQUE, left, top, width, height,
                 fill, line, C)


def frame(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FRAME, left, top, width, height,
                 fill, line, C)


def cube(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.CUBE, left, top, width, height,
                 fill, line, C)


def bevel(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.BEVEL, left, top, width, height,
                 fill, line, C)


def folded_corner(slide, left, top, width, height, fill, line=None, C=None):
    return shape(slide, MSO_SHAPE.FOLDED_CORNER, left, top, width, height,
                 fill, line, C)


def tear(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.TEAR, cx, cy, size, size,
                           fill, line, C)


def math_plus(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.MATH_PLUS, cx, cy, size, size,
                           fill, line, C)


def math_multiply(slide, cx, cy, size, fill, line=None, C=None):
    return _centered_shape(slide, MSO_SHAPE.MATH_MULTIPLY, cx, cy, size, size,
                           fill, line, C)


def text(slide, left, top, width, height, txt, font_size=12,
         color='text_body', bold=False, align='left',
         font_name=None, C=None, anchor='top'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor == 'middle':
        try:
            txBox.text_frame._txBody.bodyPr.set('anchor', 'ctr')
        except Exception:
            pass
    p = tf.paragraphs[0]
    _set_run(p, txt, font_size=font_size, color=color, bold=bold,
             font_name=font_name, C=C)
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    return txBox


def multiline(slide, left, top, width, height, lines, font_size=12,
              color='text_body', bold=False, align='left',
              font_name=None, C=None, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_run(p, line, font_size=font_size, color=color, bold=bold,
                 font_name=font_name, C=C)
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                       'right': PP_ALIGN.RIGHT}[align]
        if line_spacing:
            p.space_before = Pt(line_spacing)
            p.space_after = Pt(line_spacing)
        else:
            p.space_before = Pt(2)
            p.space_after = Pt(2)
    return txBox


def copy_decorations(slide, template_slide, skip_long_text=True, skip_image=True):
    for shape in template_slide.shapes:
        if skip_image and shape.shape_type == 13:
            continue
        if skip_long_text and shape.has_text_frame:
            if len(shape.text_frame.text) > 50:
                continue
        el = copy.deepcopy(shape._element)
        slide.shapes._spTree.append(el)


def copy_logo(slide, template_slide, color_hints=None):
    for shape in template_slide.shapes:
        if shape.shape_type != 6:
            continue
        if color_hints:
            sp = shape._element
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            for hint in color_hints:
                if sp.find(f'.//{{{ns}}}srgbClr[@val="{hint.lstrip("#")}"]') is not None:
                    el = copy.deepcopy(sp)
                    slide.shapes._spTree.append(el)
                    return
        else:
            el = copy.deepcopy(shape._element)
            slide.shapes._spTree.append(el)
            return


def top_bar(slide, color, width=13.333, height=0.08, C=None):
    color_val = _resolve_color(color, C)
    return rect(slide, 0, 0, width, height, color_val)


def page_header(slide, title, subtitle='', C=None, left=0.65, width=None,
                typo=None, spacing=None):
    C = C or {}
    cw = width or (13.333 - 2 * left)
    t = typo or TYPOGRAPHY.get('mckinsey')
    sp = spacing or SPACING.get('mckinsey')
    text(slide, left, 0.45, cw, 0.5, title,
         font_size=t.h1, color=C.get('text_dark', '#000000'), bold=True,
         font_name=C.get('font_heading'), C=C)
    if subtitle:
        text(slide, left, 0.95, cw, 0.25, subtitle,
             font_size=t.caption, color=C.get('text_muted', '#666666'),
             font_name=C.get('font_body'), C=C)
    rect(slide, left, 1.25, cw, 0.004, C.get('divider', '#CCCCCC'))


def kpi_card(slide, left, top, width, height, number, label,
             trend='', trend_up=True, C=None, typo=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')
    pad = 0.2

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        bg = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(left), Inches(top),
                        Inches(width), Inches(height))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(C.get('card_bg', C.get('white', '#FFFFFF')))
        bg.line.color.rgb = _rgb(C.get('card_line', C.get('light', '#DDDDDD')))
        bg.line.width = Pt(1)

        accent_bar = _add_shape(gs, MSO_SHAPE.RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(0.06))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = _rgb(C.get('accent', '#4CAF50'))
        accent_bar.line.fill.background()

        num_box = gs.add_textbox(Inches(left + pad), Inches(top + 0.25),
                                 Inches(width - 2 * pad), Inches(0.5))
        p = num_box.text_frame.paragraphs[0]
        _set_run(p, number, font_size=t.h1, color=C.get('primary', '#1B5E20'),
                 bold=True, font_name=C.get('font_heading'), C=C)

        lbl_box = gs.add_textbox(Inches(left + pad), Inches(top + 0.75),
                                 Inches(width - 2 * pad), Inches(0.3))
        p2 = lbl_box.text_frame.paragraphs[0]
        _set_run(p2, label, font_size=t.caption, color=C.get('text_muted', '#666666'),
                 font_name=C.get('font_body'), C=C)

        if trend:
            tc = C.get('primary', '#1B5E20') if trend_up else '#C53030'
            trend_box = gs.add_textbox(Inches(left + pad), Inches(top + 1.05),
                                       Inches(width - 2 * pad), Inches(0.25))
            p3 = trend_box.text_frame.paragraphs[0]
            _set_run(p3, trend, font_size=t.micro, color=tc, bold=True,
                     font_name=C.get('font_body'), C=C)

        return group
    else:
        rrect(slide, left, top, width, height,
              C.get('card_bg', C.get('white', '#FFFFFF')),
              line=C.get('card_line', C.get('light', '#DDDDDD')), C=C)
        rect(slide, left, top, width, 0.06, C.get('accent', '#4CAF50'), C=C)
        text(slide, left + pad, top + 0.25, width - 2 * pad, 0.5, number,
             font_size=t.h1, color=C.get('primary', '#1B5E20'), bold=True,
             font_name=C.get('font_heading'), C=C)
        text(slide, left + pad, top + 0.75, width - 2 * pad, 0.3, label,
             font_size=t.caption, color=C.get('text_muted', '#666666'),
             font_name=C.get('font_body'), C=C)
        if trend:
            tc = C.get('primary', '#1B5E20') if trend_up else '#C53030'
            text(slide, left + pad, top + 1.05, width - 2 * pad, 0.25, trend,
                 font_size=t.micro, color=tc, bold=True,
                 font_name=C.get('font_body'), C=C)


def bar_chart(slide, left, top, data, max_width=5.0, bar_height=0.3, C=None,
              typo=None, spacing=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')
    sp = spacing or SPACING.get('mckinsey')
    bar_colors = [C.get('primary', '#1B5E20'), C.get('accent', '#4CAF50'),
                  C.get('muted', '#81C784'), C.get('light', '#C8E6C9')]
    gap = sp.bar_gap

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        for i, (label, pct, val) in enumerate(data):
            y = top + i * (bar_height + gap)

            bg = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(left), Inches(y),
                            Inches(max_width), Inches(bar_height))
            bg.fill.solid()
            bg.fill.fore_color.rgb = _rgb(C.get('bg_tint', '#F5F5F5'))
            bg.line.fill.background()

            bar_w = max_width * pct
            if bar_w > 0:
                bar = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(y),
                                 Inches(bar_w), Inches(bar_height))
                bar.fill.solid()
                bar.fill.fore_color.rgb = _rgb(bar_colors[i % len(bar_colors)])
                bar.line.fill.background()

            lbl_box = gs.add_textbox(Inches(left - 0.9), Inches(y - 0.03),
                                     Inches(0.85), Inches(bar_height))
            p = lbl_box.text_frame.paragraphs[0]
            _set_run(p, label, font_size=t.caption, color=C.get('text_body', '#333333'),
                     font_name=C.get('font_body'), C=C)
            p.alignment = PP_ALIGN.RIGHT

            val_box = gs.add_textbox(Inches(left + max_width + 0.08), Inches(y - 0.03),
                                     Inches(0.6), Inches(bar_height))
            p2 = val_box.text_frame.paragraphs[0]
            _set_run(p2, val, font_size=t.caption, color=C.get('text_dark', '#000000'),
                     bold=True, font_name=C.get('font_body'), C=C)

        return group
    else:
        for i, (label, pct, val) in enumerate(data):
            y = top + i * (bar_height + gap)
            rrect(slide, left, y, max_width, bar_height, C.get('bg_tint', '#F5F5F5'), C=C)
            bar_w = max_width * pct
            if bar_w > 0:
                rrect(slide, left, y, bar_w, bar_height,
                      bar_colors[i % len(bar_colors)], C=C)
            text(slide, left - 0.9, y - 0.03, 0.85, bar_height, label,
                 font_size=t.caption, color=C.get('text_body', '#333333'), align='right',
                 font_name=C.get('font_body'), C=C)
            text(slide, left + max_width + 0.08, y - 0.03, 0.6, bar_height, val,
                 font_size=t.caption, color=C.get('text_dark', '#000000'), bold=True,
                 font_name=C.get('font_body'), C=C)


def comparison_bars(slide, left, top, metrics, max_width=4.0, C=None,
                    typo=None, spacing=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')
    sp = spacing or SPACING.get('mckinsey')
    row_h = 0.55

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        for idx, (label, v_old, v_new, pct_old, pct_new) in enumerate(metrics):
            y = top + idx * row_h

            lbl_box = gs.add_textbox(Inches(left - 1.1), Inches(y - 0.02),
                                     Inches(1.0), Inches(0.2))
            p = lbl_box.text_frame.paragraphs[0]
            _set_run(p, label, font_size=t.caption, color=C.get('text_body', '#333333'),
                     bold=True, font_name=C.get('font_body'), C=C)
            p.alignment = PP_ALIGN.RIGHT

            bg1 = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(left), Inches(y),
                             Inches(max_width), Inches(0.18))
            bg1.fill.solid()
            bg1.fill.fore_color.rgb = _rgb(C.get('bg_tint', '#F5F5F5'))
            bg1.line.fill.background()

            bar_old = max_width * pct_old
            if bar_old > 0:
                b1 = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(y),
                                Inches(bar_old), Inches(0.18))
                b1.fill.solid()
                b1.fill.fore_color.rgb = _rgb(C.get('muted', '#81C784'))
                b1.line.fill.background()

            bg2 = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(left), Inches(y + 0.22),
                             Inches(max_width), Inches(0.18))
            bg2.fill.solid()
            bg2.fill.fore_color.rgb = _rgb(C.get('bg_tint', '#F5F5F5'))
            bg2.line.fill.background()

            bar_new = max_width * pct_new
            if bar_new > 0:
                b2 = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(y + 0.22),
                                Inches(bar_new), Inches(0.18))
                b2.fill.solid()
                b2.fill.fore_color.rgb = _rgb(C.get('primary', '#1B5E20'))
                b2.line.fill.background()

            old_box = gs.add_textbox(Inches(left + max_width + 0.1), Inches(y - 0.03),
                                     Inches(0.8), Inches(0.2))
            p2 = old_box.text_frame.paragraphs[0]
            _set_run(p2, v_old, font_size=t.micro, color=C.get('text_muted', '#666666'),
                     font_name=C.get('font_body'), C=C)

            new_box = gs.add_textbox(Inches(left + max_width + 0.1), Inches(y + 0.19),
                                     Inches(0.8), Inches(0.2))
            p3 = new_box.text_frame.paragraphs[0]
            _set_run(p3, v_new, font_size=t.micro, color=C.get('text_dark', '#000000'),
                     bold=True, font_name=C.get('font_body'), C=C)

        return group
    else:
        for label, v_old, v_new, pct_old, pct_new in metrics:
            text(slide, left - 1.1, top - 0.02, 1.0, 0.2, label,
                 font_size=t.caption, color=C.get('text_body', '#333333'), bold=True,
                 align='right', font_name=C.get('font_body'), C=C)
            rrect(slide, left, top, max_width, 0.18, C.get('bg_tint', '#F5F5F5'), C=C)
            bar_old = max_width * pct_old
            if bar_old > 0:
                rrect(slide, left, top, bar_old, 0.18, C.get('muted', '#81C784'), C=C)
            rrect(slide, left, top + 0.22, max_width, 0.18, C.get('bg_tint', '#F5F5F5'), C=C)
            bar_new = max_width * pct_new
            if bar_new > 0:
                rrect(slide, left, top + 0.22, bar_new, 0.18, C.get('primary', '#1B5E20'), C=C)
            text(slide, left + max_width + 0.1, top - 0.03, 0.8, 0.2, v_old,
                 font_size=t.micro, color=C.get('text_muted', '#666666'),
                 font_name=C.get('font_body'), C=C)
            text(slide, left + max_width + 0.1, top + 0.19, 0.8, 0.2, v_new,
                 font_size=t.micro, color=C.get('text_dark', '#000000'), bold=True,
                 font_name=C.get('font_body'), C=C)
            top += row_h
        return top


def donut_chart(slide, cx, cy, radius, inner_radius, sectors, C=None,
                typo=None, grouped=True, native=True):
    """Donut/pie chart. When native=True and sectors>1, uses PowerPoint native
    doughnut chart for accurate sector angles. When native=False or sectors==1,
    falls back to Shape-based rendering for maximum visual customization."""
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    if native and len(sectors) > 1:
        chart_w = radius * 2 + 2.5
        chart_h = radius * 2 + 0.6
        chart_left = cx - radius - 0.2
        chart_top = cy - radius - 0.3

        categories = [s[0] for s in sectors]
        pct_values = []
        for s in sectors:
            pct_str = s[1].replace('%', '').strip()
            try:
                pct_values.append(float(pct_str))
            except (ValueError, AttributeError):
                pct_values.append(0)
        sector_colors = [s[2] for s in sectors]

        series = [{"name": "Share", "values": pct_values}]
        chart_style = {
            "show_legend": True,
            "legend_position": "right",
            "show_labels": True,
            "show_percentage": True,
            "show_value": False,
            "label_position": "best_fit",
            "color_scheme": sector_colors,
        }

        result = native_chart(slide, chart_left, chart_top, chart_w, chart_h,
                              "doughnut", categories=categories, series=series,
                              style=chart_style, C=C)
        if result is not None:
            return result

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        for name, pct_str, clr in sectors:
            outer = _add_shape(gs, MSO_SHAPE.OVAL,
                               Inches(cx - radius), Inches(cy - radius),
                               Inches(radius * 2), Inches(radius * 2))
            outer.fill.solid()
            outer.fill.fore_color.rgb = _rgb(clr)
            outer.line.fill.background()

        inner = _add_shape(gs, MSO_SHAPE.OVAL,
                           Inches(cx - inner_radius), Inches(cy - inner_radius),
                           Inches(inner_radius * 2), Inches(inner_radius * 2))
        inner.fill.solid()
        inner.fill.fore_color.rgb = _rgb(C.get('background', '#FFFFFF'))
        inner.line.fill.background()

        center_box = gs.add_textbox(
            Inches(cx - 0.5), Inches(cy - 0.2),
            Inches(1.0), Inches(0.4))
        p = center_box.text_frame.paragraphs[0]
        _set_run(p, '100%', font_size=t.h2, color=C.get('primary', '#1B5E20'),
                 bold=True, font_name=C.get('font_heading'), C=C)
        p.alignment = PP_ALIGN.CENTER

        ly = cy - radius
        lx = cx + radius + 0.5
        for name, pct_str, clr in sectors:
            dot = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(lx), Inches(ly),
                             Inches(0.2), Inches(0.2))
            dot.fill.solid()
            dot.fill.fore_color.rgb = _rgb(clr)
            dot.line.fill.background()

            lbl = gs.add_textbox(Inches(lx + 0.3), Inches(ly - 0.02),
                                 Inches(1.5), Inches(0.25))
            p2 = lbl.text_frame.paragraphs[0]
            _set_run(p2, f'{name}  {pct_str}', font_size=t.caption,
                     color=C.get('text_body', '#333333'),
                     font_name=C.get('font_body'), C=C)
            ly += 0.35

        return group
    else:
        for name, pct_str, clr in sectors:
            oval(slide, cx - radius, cy - radius, radius * 2, radius * 2, clr, C=C)
        oval(slide, cx - inner_radius, cy - inner_radius,
             inner_radius * 2, inner_radius * 2,
             C.get('background', '#FFFFFF'), C=C)
        text(slide, cx - 0.5, cy - 0.2, 1.0, 0.4, '100%',
             font_size=t.h2, color=C.get('primary', '#1B5E20'), bold=True,
             align='center', font_name=C.get('font_heading'), C=C)
        ly = cy - radius
        lx = cx + radius + 0.5
        for name, pct_str, clr in sectors:
            rrect(slide, lx, ly, 0.2, 0.2, clr, C=C)
            text(slide, lx + 0.3, ly - 0.02, 1.5, 0.25, f'{name}  {pct_str}',
                 font_size=t.caption, color=C.get('text_body', '#333333'),
                 font_name=C.get('font_body'), C=C)
            ly += 0.35


def native_chart(slide, left, top, width, height, chart_type,
                 categories=None, series=None, style=None, C=None):
    """Native PowerPoint chart — editable data, axes, gridlines, legend.

    chart_type: 'bar'|'bar_stacked'|'bar_100'|'bar_3d'|
                'bar_horizontal'|'bar_horizontal_stacked'|'bar_horizontal_100'|
                'line'|'line_markers'|'line_stacked'|'line_stacked_100'|
                'pie'|'pie_3d'|'pie_exploded'|
                'doughnut'|'doughnut_exploded'|
                'area'|'area_stacked'|'area_stacked_100'|
                'scatter'|'scatter_lines'|'scatter_smooth'|
                'radar'|'radar_markers'|'bubble'
    categories: ['Q1','Q2','Q3','Q4']  (not used for scatter/bubble)
    series: [{'name':'Revenue','values':[30,45,60,75]}, ...]
            For scatter: values = [[x1,y1],[x2,y2],...]
            For bubble:  values = [[x1,y1,size1],[x2,y2,size2],...]
    style: {
        'show_legend': True,
        'legend_position': 'bottom',  # bottom|top|left|right
        'show_labels': False,
        'show_value': True,
        'show_percentage': False,
        'show_category_name': False,
        'label_font_size': 9,
        'label_position': 'outside_end',  # center|inside_end|outside_end|best_fit
        'number_format': '#,##0',
        'color_scheme': 'brand',  # 'brand'|'auto'|['#hex',...]
        'title': 'Chart Title',
        'value_axis_title': 'Revenue ($M)',
        'category_axis_title': 'Quarter',
        'gridlines': 'major_y',  # 'none'|'major_y'|'major_x'|'major_xy'
        'tick_number_format': '#,##0',
        'chart_style': 2,  # 1-48 built-in PowerPoint chart style
    }
    C: color dictionary (used for 'brand' color_scheme)
    """
    from ppt_pro_max.renderer.chart_builder import ChartBuilder

    if categories is None:
        categories = ["Q1", "Q2", "Q3", "Q4"]
    if series is None:
        series = [{"name": "Data", "values": [10, 25, 45, 80]}]
    if style is None:
        style = {}

    chart_config = {
        "type": chart_type,
        "categories": categories,
        "series": series,
        "style": style,
    }

    brand_colors = None
    if C:
        brand_colors = {
            "primary": C.get("primary", "#2563EB"),
            "secondary": C.get("secondary", "#64748B"),
            "accent": C.get("accent", "#F97316"),
        }

    position = {"x": left, "y": top, "width": width, "height": height}
    builder = ChartBuilder()
    return builder.build(slide, chart_config, position=position, brand_colors=brand_colors)


def _draw_connector(slide, x1, y1, x2, y2, width_pt=1.0, color='#A0A0A0', dash=False):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.width = Pt(width_pt)
    c.line.color.rgb = _rgb(color)
    if dash:
        c.line.dash_style = 2
    return c


def _transparent_textbox(slide, left, top, width, height, txt,
                         font_size=160, font_name='SimSun', color='#000000',
                         bold=False, anchor='ctr'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.color.rgb = _rgb(color)
    run.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    tb.fill.background()
    tb.line.fill.background()
    bodyPr = tf._txBody.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if anchor:
        bodyPr.set('anchor', anchor)
    bodyPr.set('marL', '0')
    bodyPr.set('marR', '0')
    bodyPr.set('marT', '0')
    bodyPr.set('marB', '0')
    return tb


def mizi_grid(slide, left, top, size, char=None,
              border_color='#4CAF50', guide_color='#A0A0A0',
              border_pt=2.5, guide_pt=1.0, diag_pt=0.75,
              font_size=160, font_name='SimSun', font_color='#000000'):
    """米字格 — cross + diagonal guide lines with optional character overlay.

    Layout:  border (solid, green), cross (dashed, gray), diagonals (dashed, gray)
    Usage:   mizi_grid(s, 1.0, 1.5, 2.5, char='永')
    """
    x1, y1, x2, y2 = left, top, left + size, top + size

    _draw_connector(slide, x1, y1, x2, y1, border_pt, border_color)
    _draw_connector(slide, x1, y2, x2, y2, border_pt, border_color)
    _draw_connector(slide, x1, y1, x1, y2, border_pt, border_color)
    _draw_connector(slide, x2, y1, x2, y2, border_pt, border_color)

    cx, cy = left + size / 2, top + size / 2
    _draw_connector(slide, x1, cy, x2, cy, guide_pt, guide_color, dash=True)
    _draw_connector(slide, cx, y1, cx, y2, guide_pt, guide_color, dash=True)

    _draw_connector(slide, x1, y1, x2, y2, diag_pt, guide_color, dash=True)
    _draw_connector(slide, x2, y1, x1, y2, diag_pt, guide_color, dash=True)

    if char:
        _transparent_textbox(slide, left, top, size, size, char,
                             font_size=font_size, font_name=font_name, color=font_color)
    return slide


def tian_grid(slide, left, top, size, char=None,
              border_color='#4CAF50', guide_color='#A0A0A0',
              border_pt=2.5, guide_pt=1.0,
              font_size=160, font_name='SimSun', font_color='#000000'):
    """田字格 — cross guide lines only, no diagonals.

    Layout:  border (solid, green), cross (dashed, gray)
    Usage:   tian_grid(s, 1.0, 1.5, 2.5, char='永')
    """
    x1, y1, x2, y2 = left, top, left + size, top + size

    _draw_connector(slide, x1, y1, x2, y1, border_pt, border_color)
    _draw_connector(slide, x1, y2, x2, y2, border_pt, border_color)
    _draw_connector(slide, x1, y1, x1, y2, border_pt, border_color)
    _draw_connector(slide, x2, y1, x2, y2, border_pt, border_color)

    cx, cy = left + size / 2, top + size / 2
    _draw_connector(slide, x1, cy, x2, cy, guide_pt, guide_color, dash=True)
    _draw_connector(slide, cx, y1, cx, y2, guide_pt, guide_color, dash=True)

    if char:
        _transparent_textbox(slide, left, top, size, size, char,
                             font_size=font_size, font_name=font_name, color=font_color)
    return slide


def pinyin_grid(slide, left, top, width, pinyin=None,
                baseline_y=None, line_spacing=0.3,
                light_color='#A0A0A0', dark_color='#424242',
                light_pt=0.75, dark_pt=1.5,
                font_size=36, font_name='SimSun', font_color='#000000'):
    """四线格（拼音格）— 4 horizontal lines for pinyin writing.

    Layout:  line1 (light), line2 (dark), line3/baseline (dark), line4 (light)
    If pinyin given, a transparent textbox is placed with baseline aligned to line3.
    Usage:   pinyin_grid(s, 1.0, 2.0, 3.0, pinyin='yǒng')
    """
    if baseline_y is None:
        baseline_y = top + line_spacing * 2

    x1, x2 = left, left + width

    _draw_connector(slide, x1, baseline_y - line_spacing * 2, x2, baseline_y - line_spacing * 2, light_pt, light_color)
    _draw_connector(slide, x1, baseline_y - line_spacing, x2, baseline_y - line_spacing, dark_pt, dark_color)
    _draw_connector(slide, x1, baseline_y, x2, baseline_y, dark_pt, dark_color)
    _draw_connector(slide, x1, baseline_y + line_spacing, x2, baseline_y + line_spacing, light_pt, light_color)

    if pinyin:
        grid_h = line_spacing * 4
        tb_top = baseline_y - line_spacing * 2
        _transparent_textbox(slide, left, tb_top, width, grid_h, pinyin,
                             font_size=font_size, font_name=font_name, color=font_color)
    return slide


def hanzi_row(slide, left, top, size, chars, grid_type='mizi',
              gap=0.3, border_color='#4CAF50', guide_color='#A0A0A0',
              border_pt=2.5, guide_pt=1.0, diag_pt=0.75,
              font_size=160, font_name='SimSun', font_color='#000000'):
    """Draw a row of character grids — convenience function for multiple grids.

    grid_type: 'mizi' | 'tian'
    chars: list of characters; None entries draw empty grids
    Usage:   hanzi_row(s, 1.0, 1.5, 2.0, ['永', None, '和'], grid_type='mizi')
    """
    grid_fn = mizi_grid if grid_type == 'mizi' else tian_grid
    for i, ch in enumerate(chars):
        x = left + i * (size + gap)
        if grid_type == 'mizi':
            grid_fn(slide, x, top, size, char=ch,
                    border_color=border_color, guide_color=guide_color,
                    border_pt=border_pt, guide_pt=guide_pt,
                    diag_pt=diag_pt, font_size=font_size,
                    font_name=font_name, font_color=font_color)
        else:
            grid_fn(slide, x, top, size, char=ch,
                    border_color=border_color, guide_color=guide_color,
                    border_pt=border_pt, guide_pt=guide_pt,
                    font_size=font_size,
                    font_name=font_name, font_color=font_color)
    return slide


def pinyin_hanzi_block(slide, left, top, size, items, gap=0.3,
                       grid_type='mizi', pinyin_line_spacing=0.3,
                       border_color='#4CAF50', guide_color='#A0A0A0',
                       border_pt=2.5, guide_pt=1.0, diag_pt=0.75,
                       pinyin_light_color='#A0A0A0', pinyin_dark_color='#424242',
                       pinyin_light_pt=0.75, pinyin_dark_pt=1.5,
                       char_font_size=160, char_font_name='SimSun', char_font_color='#000000',
                       pinyin_font_size=36, pinyin_font_name='SimSun', pinyin_font_color='#000000'):
    """Draw pinyin grid + character grid as a paired block for each item.

    grid_type: 'mizi' | 'tian'
    items: list of (pinyin, char) tuples; char can be None for empty grid
    Usage:  pinyin_hanzi_block(s, 1.0, 0.5, 2.0, [('yǒng','永'), ('hé','和'), (None, None)])
    """
    pinyin_height = pinyin_line_spacing * 4
    pinyin_gap = 0.15
    grid_fn = mizi_grid if grid_type == 'mizi' else tian_grid

    for i, item in enumerate(items):
        py, ch = item if item else (None, None)
        x = left + i * (size + gap)

        if py:
            baseline_y = top + pinyin_line_spacing * 2
            pinyin_grid(slide, x, top, size, pinyin=py,
                        baseline_y=baseline_y, line_spacing=pinyin_line_spacing,
                        light_color=pinyin_light_color, dark_color=pinyin_dark_color,
                        light_pt=pinyin_light_pt, dark_pt=pinyin_dark_pt,
                        font_size=pinyin_font_size, font_name=pinyin_font_name,
                        font_color=pinyin_font_color)
        else:
            baseline_y = top + pinyin_line_spacing * 2
            pinyin_grid(slide, x, top, size,
                        baseline_y=baseline_y, line_spacing=pinyin_line_spacing,
                        light_color=pinyin_light_color, dark_color=pinyin_dark_color,
                        light_pt=pinyin_light_pt, dark_pt=pinyin_dark_pt)

        char_top = top + pinyin_height + pinyin_gap
        if grid_type == 'mizi':
            grid_fn(slide, x, char_top, size, char=ch,
                    border_color=border_color, guide_color=guide_color,
                    border_pt=border_pt, guide_pt=guide_pt, diag_pt=diag_pt,
                    font_size=char_font_size, font_name=char_font_name,
                    font_color=char_font_color)
        else:
            grid_fn(slide, x, char_top, size, char=ch,
                    border_color=border_color, guide_color=guide_color,
                    border_pt=border_pt, guide_pt=guide_pt,
                    font_size=char_font_size, font_name=char_font_name,
                    font_color=char_font_color)
    return slide


def highlight_cards(slide, left, top, cards, total_width=12.0, C=None,
                    typo=None, spacing=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')
    sp = spacing or SPACING.get('mckinsey')
    n = len(cards)
    gap = sp.card_gap
    card_w = (total_width - gap * (n - 1)) / n
    card_h = 1.4
    pad = sp.card_padding

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        for i, (title, desc, accent) in enumerate(cards):
            x = left + i * (card_w + gap)

            bg = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(top),
                            Inches(card_w), Inches(card_h))
            bg.fill.solid()
            bg.fill.fore_color.rgb = _rgb(C.get('card_bg', '#F9F9F9'))
            bg.line.color.rgb = _rgb(C.get('light', '#DDDDDD'))
            bg.line.width = Pt(1)

            accent_bar = _add_shape(gs, MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(top),
                                    Inches(card_w), Inches(0.06))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = _rgb(accent)
            accent_bar.line.fill.background()

            title_box = gs.add_textbox(Inches(x + pad), Inches(top + 0.18),
                                       Inches(card_w - 2 * pad), Inches(0.3))
            p = title_box.text_frame.paragraphs[0]
            _set_run(p, title, font_size=t.h3, color=C.get('text_dark', '#000000'),
                     bold=True, font_name=C.get('font_heading'), C=C)

            desc_box = gs.add_textbox(Inches(x + pad), Inches(top + 0.52),
                                      Inches(card_w - 2 * pad), Inches(0.7))
            p2 = desc_box.text_frame.paragraphs[0]
            _set_run(p2, desc, font_size=t.caption, color=C.get('text_muted', '#666666'),
                     font_name=C.get('font_body'), C=C)

        return group
    else:
        for i, (title, desc, accent) in enumerate(cards):
            x = left + i * (card_w + gap)
            rrect(slide, x, top, card_w, 1.2, C.get('card_bg', '#F9F9F9'),
                  line=C.get('light', '#DDDDDD'), C=C)
            rect(slide, x, top, card_w, 0.06, accent, C=C)
            text(slide, x + pad, top + 0.18, card_w - 2 * pad, 0.3, title,
                 font_size=t.h3, color=C.get('text_dark', '#000000'), bold=True,
                 font_name=C.get('font_heading'), C=C)
            text(slide, x + pad, top + 0.52, card_w - 2 * pad, 0.5, desc,
                 font_size=t.caption, color=C.get('text_muted', '#666666'),
                 font_name=C.get('font_body'), C=C)


def code_block(slide, left, top, width, height, lines, language='python',
               C=None, typo=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        bg = _add_shape(gs, MSO_SHAPE.RECTANGLE,
                        Inches(left), Inches(top),
                        Inches(width), Inches(height))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb('#1E1E1E')
        bg.line.fill.background()

        badge_w = len(language) * 0.12 + 0.3
        badge = _add_shape(gs, MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(left + 0.1), Inches(top - 0.28),
                           Inches(badge_w), Inches(0.25))
        badge.fill.solid()
        badge.fill.fore_color.rgb = _rgb(C.get('accent', '#4CAF50'))
        badge.line.fill.background()

        badge_txt = gs.add_textbox(Inches(left + 0.15), Inches(top - 0.28),
                                   Inches(badge_w - 0.1), Inches(0.25))
        p = badge_txt.text_frame.paragraphs[0]
        _set_run(p, language, font_size=t.micro, color=C.get('white', '#FFFFFF'),
                 bold=True, font_name=C.get('font_body'), C=C)

        code_box = gs.add_textbox(Inches(left + 0.2), Inches(top + 0.15),
                                  Inches(width - 0.4), Inches(height - 0.3))
        tf = code_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p2 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p2.add_run()
            run.text = line
            run.font.size = Pt(t.body)
            run.font.color.rgb = _rgb('#D4D4D4')
            run.font.name = 'Consolas'
            p2.space_before = Pt(3)
            p2.space_after = Pt(3)

        return group
    else:
        rect(slide, left, top, width, height, '#1E1E1E')
        rrect(slide, left, top - 0.28, len(language) * 0.12 + 0.3, 0.25,
              C.get('accent', '#4CAF50'), C=C)
        text(slide, left + 0.05, top - 0.28, len(language) * 0.12 + 0.2, 0.25,
             language, font_size=t.micro, color=C.get('white', '#FFFFFF'),
             bold=True, font_name='Consolas', C=C)
        multiline(slide, left + 0.2, top + 0.15, width - 0.4, height - 0.3,
                  lines, font_size=t.body, color='#D4D4D4',
                  font_name='Consolas', line_spacing=3)


def section_divider(slide, number, title, C=None, typo=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        bg = _add_shape(gs, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(C.get('primary', '#1B5E20'))
        bg.line.fill.background()

        num_box = gs.add_textbox(Inches(1.2), Inches(1.5),
                                 Inches(3.0), Inches(2.0))
        p = num_box.text_frame.paragraphs[0]
        _set_run(p, str(number).zfill(2), font_size=72, color=C.get('light', '#C8E6C9'),
                 bold=True, font_name=C.get('font_heading'), C=C)

        line = _add_shape(gs, MSO_SHAPE.RECTANGLE,
                          Inches(1.2), Inches(3.6),
                          Inches(2.0), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = _rgb(C.get('accent', '#4CAF50'))
        line.line.fill.background()

        title_box = gs.add_textbox(Inches(1.2), Inches(3.9),
                                   Inches(10.0), Inches(1.5))
        p2 = title_box.text_frame.paragraphs[0]
        _set_run(p2, title, font_size=t.hero, color=C.get('white', '#FFFFFF'),
                 bold=True, font_name=C.get('font_heading'), C=C)

        return group
    else:
        rect(slide, 0, 0, 13.333, 7.5, C.get('primary', '#1B5E20'))
        text(slide, 1.2, 1.5, 3.0, 2.0, str(number).zfill(2),
             font_size=72, color=C.get('light', '#C8E6C9'), bold=True,
             font_name=C.get('font_heading'), C=C)
        rect(slide, 1.2, 3.6, 2.0, 0.04, C.get('accent', '#4CAF50'))
        text(slide, 1.2, 3.9, 10.0, 1.5, title,
             font_size=t.hero, color=C.get('white', '#FFFFFF'), bold=True,
             font_name=C.get('font_heading'), C=C)


def hero_slide(slide, title, subtitle='', C=None, typo=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        bg = _add_shape(gs, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(C.get('primary', '#1B5E20'))
        bg.line.fill.background()

        title_box = gs.add_textbox(Inches(1.2), Inches(2.0),
                                   Inches(10.0), Inches(1.5))
        p = title_box.text_frame.paragraphs[0]
        _set_run(p, title, font_size=t.hero, color=C.get('white', '#FFFFFF'),
                 bold=True, font_name=C.get('font_heading'), C=C)

        if subtitle:
            sub_box = gs.add_textbox(Inches(1.2), Inches(3.6),
                                     Inches(10.0), Inches(0.5))
            p2 = sub_box.text_frame.paragraphs[0]
            _set_run(p2, subtitle, font_size=t.h2, color=C.get('light', '#C8E6C9'),
                     font_name=C.get('font_body'), C=C)

        return group
    else:
        rect(slide, 0, 0, 13.333, 7.5, C.get('primary'))
        text(slide, 1.2, 2.0, 10.0, 1.5, title,
             font_size=t.hero, color=C.get('white', '#FFFFFF'), bold=True,
             font_name=C.get('font_heading'), C=C)
        if subtitle:
            text(slide, 1.2, 3.6, 10.0, 0.5, subtitle,
                 font_size=t.h2, color=C.get('light', '#C8E6C9'),
                 font_name=C.get('font_body'), C=C)


def cta_slide(slide, title, subtitle='', C=None, typo=None, grouped=True):
    C = C or {}
    t = typo or TYPOGRAPHY.get('mckinsey')

    if grouped:
        group = slide.shapes.add_group_shape()
        gs = group.shapes

        bg = _add_shape(gs, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(C.get('primary', '#1B5E20'))
        bg.line.fill.background()

        title_box = gs.add_textbox(Inches(1.2), Inches(2.5),
                                   Inches(10.0), Inches(1.5))
        p = title_box.text_frame.paragraphs[0]
        _set_run(p, title, font_size=t.h1 + 12, color=C.get('white', '#FFFFFF'),
                 bold=True, font_name=C.get('font_heading'), C=C)

        if subtitle:
            sub_box = gs.add_textbox(Inches(1.2), Inches(4.0),
                                     Inches(10.0), Inches(0.5))
            p2 = sub_box.text_frame.paragraphs[0]
            _set_run(p2, subtitle, font_size=t.h3, color=C.get('light', '#C8E6C9'),
                     font_name=C.get('font_body'), C=C)

        return group
    else:
        rect(slide, 0, 0, 13.333, 7.5, C.get('primary'))
        text(slide, 1.2, 2.5, 10.0, 1.5, title,
             font_size=t.h1 + 12, color=C.get('white', '#FFFFFF'), bold=True,
             font_name=C.get('font_heading'), C=C)
        if subtitle:
         text(slide, 1.2, 4.0, 10.0, 0.5, subtitle,
              font_size=t.h3, color=C.get('light', '#C8E6C9'),
              font_name=C.get('font_body'), C=C)


def gradient_text(slide, left, top, width, height, txt, preset='gold-shine',
                  stops=None, font_size=44, bold=False, font_name=None,
                  cjk_font=None, align='left'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    if cjk_font:
        _set_cjk_font(run, cjk_font)
    if stops:
        apply_text_gradient(run, stops)
    else:
        apply_text_gradient_preset(run, preset)
    return txBox


def vertical_text(slide, left, top, width, height, txt, direction='ea',
                  font_name='STKaiti', font_size=24, color='#000000',
                  bold=False, align='center'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_vertical_text(tf, direction)
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(_resolve_color(color, None))
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    return txBox


def seal_stamp(slide, left, top, size, txt, fill_hex='#C41E3A',
               font_name='STZhongsong', rotation=-15, style='zhu',
               border_width_pt=4.0):
    from ppt_pro_max.renderer.decoration_library import add_seal_stamp as _add_seal_stamp
    _add_seal_stamp(slide, left, top, size, txt,
                    fill_hex=fill_hex, font_name=font_name,
                    rotation=rotation, style=style,
                    border_width_pt=border_width_pt)
    sh = slide.shapes[-1]
    return sh


def check_contrast(color1, color2, min_ratio=3.0):
    """Check WCAG contrast ratio between two hex colors.

    Returns (ratio, ok). For body text: min_ratio=4.5 (AA),
    for large text: min_ratio=3.0 (AA). Returns (0, False) on parse error.
    """
    def _lum(h):
        h = h.lstrip('#')
        if len(h) != 6:
            return 0
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        def _lin(c):
            c = c / 255
            return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
        return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)
    l1, l2 = _lum(color1), _lum(color2)
    if l1 == 0 and l2 == 0:
        return (0, False)
    lighter, darker = max(l1, l2), min(l1, l2)
    ratio = (lighter + 0.05) / (darker + 0.05)
    return (round(ratio, 1), ratio >= min_ratio)


def contrast_text(bg_color, min_ratio=4.5):
    """Return '#FFFFFF' or '#1A1A1A' depending on which has better contrast against bg_color."""
    w_ratio, w_ok = check_contrast(bg_color, '#FFFFFF', min_ratio)
    d_ratio, d_ok = check_contrast(bg_color, '#1A1A1A', min_ratio)
    if w_ratio >= d_ratio:
        return '#FFFFFF'
    return '#1A1A1A'


def cover_image(slide, left, top, width, height, image_path):
    """Add image with cover-fit (crop to fill, no stretch).

    Uses Pillow to pre-crop the image to the exact aspect ratio,
    then places it at the specified position. This is the correct
    way to add images to PPT — never use add_picture with stretch.
    """
    import os as _os
    import tempfile
    import hashlib
    if not _os.path.isfile(image_path):
        return None
    from PIL import Image as PILImage
    img = PILImage.open(image_path)
    img_w, img_h = img.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        crop_w = int(img_h * box_ratio)
        crop_h = img_h
        cleft = (img_w - crop_w) // 2
        ctop = 0
    else:
        crop_w = img_w
        crop_h = int(img_w / box_ratio)
        cleft = 0
        ctop = (img_h - crop_h) // 2
    cropped = img.crop((cleft, ctop, cleft + crop_w, ctop + crop_h))
    cache_dir = _os.path.join(tempfile.gettempdir(), "ppt-cropped")
    _os.makedirs(cache_dir, exist_ok=True)
    crop_key = f"crop:{image_path}:{width}x{height}"
    crop_hash = hashlib.md5(crop_key.encode()).hexdigest()
    cropped_path = _os.path.join(cache_dir, f"{crop_hash}.png")
    if not _os.path.exists(cropped_path):
        cropped.save(cropped_path, "PNG")
    return slide.shapes.add_picture(
        cropped_path, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )


def circle_image(slide, cx, cy, radius, image_path, border_color=None):
    return _add_circle_image(slide, cx, cy, radius, image_path,
                             border_hex=border_color)


def hex_image(slide, cx, cy, size, image_path, border_color=None):
    x = cx - size / 2
    y = cy - size * 0.87 / 2
    return _add_image_in_shape(slide, MSO_SHAPE.HEXAGON, x, y, size, size * 0.87,
                               image_path, border_hex=border_color)


def star_image(slide, cx, cy, size, image_path, points=5, border_color=None):
    _STAR_MAP = {
        5: MSO_SHAPE.STAR_5_POINT,
        6: MSO_SHAPE.STAR_6_POINT,
        8: MSO_SHAPE.STAR_8_POINT,
        10: MSO_SHAPE.STAR_10_POINT,
        12: MSO_SHAPE.STAR_12_POINT,
    }
    mso = _STAR_MAP.get(points, MSO_SHAPE.STAR_5_POINT)
    x = cx - size / 2
    y = cy - size / 2
    return _add_image_in_shape(slide, mso, x, y, size, size,
                               image_path, border_hex=border_color)


def diamond_image(slide, cx, cy, size, image_path, border_color=None):
    x = cx - size / 2
    y = cy - size / 2
    return _add_image_in_shape(slide, MSO_SHAPE.DIAMOND, x, y, size, size,
                               image_path, border_hex=border_color)


def heart_image(slide, cx, cy, size, image_path, border_color=None):
    x = cx - size / 2
    y = cy - size / 2
    return _add_image_in_shape(slide, MSO_SHAPE.HEART, x, y, size, size,
                               image_path, border_hex=border_color)


def shape_image(slide, shape_type, left, top, width, height, image_path,
                border_color=None):
    _type = shape_type
    if isinstance(_type, str):
        _type = getattr(MSO_SHAPE, _type.upper(), MSO_SHAPE.OVAL)
    return _add_image_in_shape(slide, _type, left, top, width, height,
                               image_path, border_hex=border_color)


def soft_edge_image(slide, left, top, width, height, image_path,
                    soft_radius=10):
    import os as _os
    if not _os.path.isfile(image_path):
        return None
    shape = slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    apply_soft_edge(shape, radius_pt=soft_radius)
    return shape


def duotone_image(slide, left, top, width, height, image_path,
                  color1='#0000FF', color2='#FF0000'):
    shape = _add_image_in_shape(slide, MSO_SHAPE.RECTANGLE,
                                left, top, width, height, image_path)
    apply_blip_duotone(shape, color1, color2)
    return shape


def artistic_image(slide, left, top, width, height, image_path,
                   effect='watercolor_sponge', params=None):
    shape = _add_image_in_shape(slide, MSO_SHAPE.RECTANGLE,
                                left, top, width, height, image_path)
    apply_blip_artistic(shape, effect, params)
    return shape


def shape_3d(slide, left, top, width, height, depth=10.0, material='powder',
             extrusion_color='#000000', shape_type=MSO_SHAPE.RECTANGLE):
    sh = _add_shape(slide.shapes, shape_type, Inches(left), Inches(top),
                    Inches(width), Inches(height))
    apply_3d(sh, depth_pt=depth, material=material,
             extrusion_color=extrusion_color)
    return sh


def bevel_shape(slide, left, top, width, height, top_w=4.0, top_h=2.0,
                material='powder', shape_type=MSO_SHAPE.RECTANGLE):
    sh = _add_shape(slide.shapes, shape_type, Inches(left), Inches(top),
                    Inches(width), Inches(height))
    apply_bevel(sh, top_w=top_w, top_h=top_h, material=material)
    return sh


def pattern_fill(slide, left, top, width, height, pattern_type, fg_color,
                 bg_color, fg_alpha=None,
                 shape_type=MSO_SHAPE.RECTANGLE):
    sh = _add_shape(slide.shapes, shape_type, Inches(left), Inches(top),
                    Inches(width), Inches(height))
    apply_pattern_fill(sh, pattern_type, fg_color, bg_color,
                       fg_alpha=fg_alpha)
    return sh


def frosted_panel(slide, left, top, width, height, tint='#FFFFFF',
                  alpha=50, soft_edge=8):
    sh = _add_shape(slide.shapes, MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                    Inches(width), Inches(height))
    apply_frosted_glass(sh, tint_color=tint, tint_alpha=alpha,
                        soft_edge=soft_edge)
    return sh


def brush_divider(slide, left, top, width, color='#2C2C2C', thickness=0.08):
    return _add_brush_divider(slide, left, top, width, color=color,
                              thickness=thickness)


def neon_border(slide, left, top, width, height, color='#8B5CF6', radius=0.1):
    return _add_neon_border(slide, left, top, width, height, color=color,
                            radius=radius)


def glass_panel(slide, left, top, width, height, tint='#FFFFFF', alpha=50,
                soft_edge=8):
    return _add_glass_panel(slide, left, top, width, height, tint=tint,
                            alpha=alpha, soft_edge=soft_edge)


def grid_background(slide, spacing=1.0, color='#E0E0E0', alpha=15):
    return _add_grid_background(slide, spacing=spacing, color=color,
                                alpha=alpha)


def ink_splash(slide, left, top, size, color='#2C2C2C', alpha=100):
    return _add_ink_splash(slide, left, top, size, color=color, alpha=alpha)


# ── Boolean shape functions (require shapely) ──


def _bool_import():
    from ppt_pro_max.renderer.boolean_shapes import HAS_SHAPELY
    return HAS_SHAPELY


def spotlight(slide, cx, cy, radius, alpha=70, color='#000000'):
    if not _bool_import():
        rrect(slide, 0, 0, 13.333, 7.5, color)
        return None
    from ppt_pro_max.renderer.boolean_shapes import (
        poly_rect, poly_circle, bool_subtract, bool_shape,
    )
    overlay = poly_rect(0, 0, 13.333, 7.5)
    cutout = poly_circle(cx, cy, radius)
    geom = bool_subtract(overlay, cutout)
    return bool_shape(geom, slide, 0, 0, 13.333, 7.5,
                      fill=color, alpha=alpha)


def bool_donut(slide, cx, cy, outer_r, inner_r, fill='#1D78FA',
               line=None, C=None):
    if not _bool_import():
        return donut(slide, cx, cy, outer_r * 2, fill, line, C)
    from ppt_pro_max.renderer.boolean_shapes import (
        poly_circle, bool_subtract, bool_shape,
    )
    outer = poly_circle(cx, cy, outer_r)
    inner = poly_circle(cx, cy, inner_r)
    geom = bool_subtract(outer, inner)
    size = outer_r * 2
    return bool_shape(geom, slide, cx - outer_r, cy - outer_r,
                      size, size, fill=fill, line=line, C=C)


def bool_frame(slide, x, y, w, h, border, fill=None, line=None, C=None):
    if not _bool_import():
        rrect(slide, x, y, w, h, fill or '#1D78FA', line, C)
        return None
    from ppt_pro_max.renderer.boolean_shapes import (
        poly_rect, bool_subtract, bool_shape,
    )
    outer = poly_rect(x, y, w, h)
    inner = poly_rect(x + border, y + border, w - 2 * border, h - 2 * border)
    geom = bool_subtract(outer, inner)
    return bool_shape(geom, slide, x, y, w, h, fill=fill or '#1D78FA',
                      line=line, C=C)


def bool_clipped_card(slide, x, y, w, h, clip_corners, clip_size=0.3,
                      fill=None, line=None, C=None):
    if not _bool_import():
        return rrect(slide, x, y, w, h, fill or '#1D78FA', line, C)
    from ppt_pro_max.renderer.boolean_shapes import (
        poly_rect, bool_subtract, bool_union, bool_shape, Polygon,
    )
    base = poly_rect(x, y, w, h)
    clips = []
    if 'tl' in clip_corners:
        clips.append(Polygon([(x, y), (x + clip_size, y), (x, y + clip_size)]))
    if 'tr' in clip_corners:
        clips.append(Polygon([(x + w - clip_size, y), (x + w, y), (x + w, y + clip_size)]))
    if 'bl' in clip_corners:
        clips.append(Polygon([(x, y + h - clip_size), (x + clip_size, y + h), (x, y + h)]))
    if 'br' in clip_corners:
        clips.append(Polygon([(x + w, y + h - clip_size), (x + w, y + h), (x + w - clip_size, y + h)]))
    if not clips:
        return rrect(slide, x, y, w, h, fill or '#1D78FA', line, C)
    from ppt_pro_max.renderer.boolean_shapes import Polygon as _Poly, bool_union
    all_clips = bool_union(*clips)
    geom = bool_subtract(base, all_clips)
    return bool_shape(geom, slide, x, y, w, h, fill=fill or '#1D78FA',
                      line=line, C=C)


def bool_neon_tube(slide, x, y, w, h, wall=0.06, fill=None, C=None):
    if not _bool_import():
        return neon_border(slide, x, y, w, h, color=fill or '#8B5CF6')
    from ppt_pro_max.renderer.boolean_shapes import (
        poly_rounded_rect, bool_subtract, bool_shape,
    )
    outer = poly_rounded_rect(x, y, w, h, wall * 2)
    inner = poly_rounded_rect(x + wall, y + wall, w - 2 * wall, h - 2 * wall, wall)
    geom = bool_subtract(outer, inner)
    return bool_shape(geom, slide, x, y, w, h, fill=fill or '#8B5CF6', C=C)


def bool_star(slide, cx, cy, r, points=5, inner_ratio=0.4, fill=None,
              line=None, C=None):
    if not _bool_import():
        mso_map = {5: MSO_SHAPE.STAR_5_POINT, 6: MSO_SHAPE.STAR_6_POINT,
                   8: MSO_SHAPE.STAR_8_POINT}
        mso = mso_map.get(points, MSO_SHAPE.STAR_5_POINT)
        return _centered_shape(slide, mso, cx, cy, r * 2, r * 2, fill or '#1D78FA', line, C)
    from ppt_pro_max.renderer.boolean_shapes import poly_star, bool_shape
    geom = poly_star(cx, cy, r, points=points, inner_ratio=inner_ratio)
    return bool_shape(geom, slide, cx - r, cy - r, r * 2, r * 2,
                      fill=fill or '#1D78FA', line=line, C=C)


def bool_cross(slide, cx, cy, w, h, bar_ratio=0.33, fill=None,
               line=None, C=None):
    if not _bool_import():
        return cross(slide, cx, cy, max(w, h), fill or '#1D78FA', line, C)
    from ppt_pro_max.renderer.boolean_shapes import poly_rect, bool_union, bool_shape
    bar_w = w * bar_ratio
    bar_h = h * bar_ratio
    v_bar = poly_rect(cx - bar_w / 2, cy - h / 2, bar_w, h)
    h_bar = poly_rect(cx - w / 2, cy - bar_h / 2, w, bar_h)
    geom = bool_union(v_bar, h_bar)
    return bool_shape(geom, slide, cx - w / 2, cy - h / 2, w, h,
                      fill=fill or '#1D78FA', line=line, C=C)


def slide_transition(slide, transition_type='fade', speed='medium',
                     advance_on_click=True, advance_after_ms=None):
    add_slide_transition(slide, transition_type=transition_type, speed=speed,
                         advance_on_click=advance_on_click,
                         advance_after_ms=advance_after_ms)


def entrance_animation(slide, shape_id, effect='fade_in', delay_ms=0,
                       duration_ms=500, click_triggered=True):
    add_entrance_animation(slide, shape_id, effect=effect, delay_ms=delay_ms,
                           duration_ms=duration_ms,
                           click_triggered=click_triggered)


def exit_animation(slide, shape_id, effect='fade_out', delay_ms=0,
                   duration_ms=500, click_triggered=True):
    add_exit_animation(slide, shape_id, effect=effect, delay_ms=delay_ms,
                       duration_ms=duration_ms,
                       click_triggered=click_triggered)


def emphasis_animation(slide, shape_id, effect='pulse', delay_ms=0,
                       duration_ms=500, click_triggered=True):
    add_emphasis_animation(slide, shape_id, effect=effect, delay_ms=delay_ms,
                           duration_ms=duration_ms,
                           click_triggered=click_triggered)


def text_outline(slide, left, top, width, height, txt,
                 color='#FFFFFF', width_pt=1.5,
                 font_size=44, bold=False, font_name=None, C=None,
                 align='left'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    cjk_font = (C or {}).get('font_cjk') or (C or {}).get('font_body')
    if cjk_font:
        _set_cjk_font(run, cjk_font)
    apply_text_outline(run, color=color, width_pt=width_pt)
    return txBox


def text_shadow(slide, left, top, width, height, txt,
                blur_pt=8, distance_pt=3, direction_deg=90,
                color='#000000', alpha_pct=25,
                font_size=44, bold=False, font_name=None, C=None,
                align='left'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    cjk_font = (C or {}).get('font_cjk') or (C or {}).get('font_body')
    if cjk_font:
        _set_cjk_font(run, cjk_font)
    from ppt_pro_max.renderer.visual_effects import apply_shadow as _apply_shadow
    _apply_shadow(txBox, blur_pt=blur_pt, distance_pt=distance_pt,
                  direction_deg=direction_deg, color=color, alpha_pct=alpha_pct)
    return txBox


def text_glow(slide, left, top, width, height, txt,
              color='#00FFFF', size_pt=8, alpha_pct=40,
              font_size=44, bold=False, font_name=None, C=None,
              align='left'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    cjk_font = (C or {}).get('font_cjk') or (C or {}).get('font_body')
    if cjk_font:
        _set_cjk_font(run, cjk_font)
    from ppt_pro_max.renderer.visual_effects import apply_glow as _apply_glow
    _apply_glow(txBox, radius_pt=size_pt, color=color, alpha_pct=alpha_pct)
    return txBox


def add_shadow(shape, blur_pt=8, distance_pt=3, direction_deg=90,
               color='#000000', alpha_pct=25):
    from ppt_pro_max.renderer.visual_effects import apply_shadow as _apply_shadow
    _apply_shadow(shape, blur_pt=blur_pt, distance_pt=distance_pt,
                  direction_deg=direction_deg, color=color, alpha_pct=alpha_pct)


def add_glow(shape, color='#00FFFF', size_pt=8, alpha_pct=40):
    from ppt_pro_max.renderer.visual_effects import apply_glow as _apply_glow
    _apply_glow(shape, radius_pt=size_pt, color=color, alpha_pct=alpha_pct)


def adjust_image(shape, brightness=0, contrast=0, saturation=100):
    from ppt_pro_max.renderer.blip_fill import (
        apply_blip_brightness_contrast, apply_blip_saturation,
    )
    if brightness != 0 or contrast != 0:
        apply_blip_brightness_contrast(shape, bright_pct=brightness,
                                       contrast_pct=contrast)
    if saturation != 100:
        apply_blip_saturation(shape, saturation_pct=saturation)


def analyze_pptx(pptx_path):
    from ppt_pro_max import extract_design_dna
    return extract_design_dna(pptx_path)


def set_theme_colors(prs, C=None):
    """Write C dict colors into the PowerPoint theme clrScheme.

    Maps: primary→accent1, secondary→accent2, tertiary→accent3,
    muted→accent4, light→accent5, text_dark→dk2, text_body→lt2.
    Makes theme colors recognizable by PowerPoint (fills default to C palette).
    """
    C = C or {}
    _ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    theme_part = None
    for rel in prs.part.rels.values():
        if 'theme' in rel.reltype:
            theme_part = rel.target_part
            break
    if theme_part is None:
        return
    theme_el = etree.fromstring(theme_part.blob)
    clrScheme = theme_el.find(f'{{{_ns}}}themeElements/{{{_ns}}}clrScheme')
    if clrScheme is None:
        return

    mapping = {
        'dk1': C.get('text_dark'),
        'lt1': C.get('background'),
        'dk2': C.get('text_dark'),
        'lt2': C.get('card_bg'),
        'accent1': C.get('primary'),
        'accent2': C.get('secondary', C.get('text_body')),
        'accent3': C.get('tertiary', C.get('accent')),
        'accent4': C.get('muted'),
        'accent5': C.get('light'),
        'accent6': C.get('divider'),
    }
    for tag, val in mapping.items():
        if not val:
            continue
        el = clrScheme.find(f'{{{_ns}}}{tag}')
        if el is None:
            continue
        for child in list(el):
            el.remove(child)
        srgb = etree.SubElement(el, f'{{{_ns}}}srgbClr')
        srgb.set('val', val.lstrip('#'))

    theme_part._blob = etree.tostring(theme_el, xml_declaration=True,
                                      encoding='UTF-8', standalone=True)


Presentation = _Presentation
