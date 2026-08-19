"""BuildQA — quality checks for Build-mode PPTs.

Inherits best checks from DeliveryGate, adds Build-mode-specific checks,
and introduces a three-tier severity system (fatal / warning / review).

Inherited from DeliveryGate:
  - residual_placeholder (fatal)
  - blank_page (fatal)
  - broken_image_ref (fatal)
  - font_too_small (warning, threshold 11pt)
  - text_overflow (warning)
  - title_duplicate (warning)
  - page_count_mismatch (fatal)

Changed from DeliveryGate:
  - color_break → review (Build mode has more color freedom)
  - font_mismatch → review
  - toc_mismatch → warning (unchanged)

Removed (template-clone / component-library specific):
  - background_missing
  - decoration_missing
  - component_placeholder_residual

New Build-mode checks:
  - element_out_of_bounds (fatal for content, review for decorative bleed)
  - image_stretched (warning)

Auto-fixable: residual_placeholder, font_too_small
"""

from __future__ import annotations

import os
import re
from dataclasses import dataclass, field

from pptx import Presentation

from ppt_pro_max.build_helpers import check_contrast

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

_PLACEHOLDER_PATTERNS = [
    re.compile(r"单击此处", re.IGNORECASE),
    re.compile(r"请输入", re.IGNORECASE),
    re.compile(r"click\s+to", re.IGNORECASE),
    re.compile(r"your\s+logo", re.IGNORECASE),
    re.compile(r"^xxx$", re.IGNORECASE),
    re.compile(r"汇报[：:]\s*x+", re.IGNORECASE),
    re.compile(r"placeholder", re.IGNORECASE),
    re.compile(r"请添加", re.IGNORECASE),
    re.compile(r"此处添加", re.IGNORECASE),
    re.compile(r"输入你的正文", re.IGNORECASE),
    re.compile(r"enter\s+text\s+here", re.IGNORECASE),
    re.compile(r"insert\s+title", re.IGNORECASE),
    re.compile(r"your\s+title\s+here", re.IGNORECASE),
    re.compile(r"add\s+your\s+content", re.IGNORECASE),
    re.compile(r"type\s+something", re.IGNORECASE),
    re.compile(r"double\s+click\s+to\s+edit", re.IGNORECASE),
]

_MIN_FONT_PT = 11
_MIN_FONT_HUNDREDTHS = 1100


def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _is_visible_text_color(rPr, a_ns: str) -> bool:
    """Skip runs whose color is transparent or inherited from theme."""
    for el in rPr.iter(f"{{{a_ns}}}solidFill"):
        if el.find(f"{{{a_ns}}}srgbClr") is not None:
            return True
        if el.find(f"{{{a_ns}}}schemeClr") is not None:
            srgb_el = el.find(f"{{{a_ns}}}schemeClr")
            if srgb_el.find(f"{{{a_ns}}}alpha") is not None:
                alpha_el = srgb_el.find(f"{{{a_ns}}}alpha")
                if int(alpha_el.get("val", "100000")) == 0:
                    return False
            return True
    return False


def _get_rpr_color_hex(rPr, a_ns: str) -> str | None:
    for el in rPr.iter(f"{{{a_ns}}}solidFill"):
        srgb = el.find(f"{{{a_ns}}}srgbClr")
        if srgb is not None:
            return f"#{srgb.get('val', '').upper()}"
    return None


@dataclass
class CheckItem:
    category: str
    check_id: str
    severity: str
    slide_index: int
    message: str
    detail: str = ""
    auto_fixable: bool = False


@dataclass
class QAReport:
    total_slides: int
    total_checks: int
    passed: int
    fatals: list[CheckItem] = field(default_factory=list)
    warnings: list[CheckItem] = field(default_factory=list)
    reviews: list[CheckItem] = field(default_factory=list)

    @property
    def is_passable(self) -> bool:
        return len(self.fatals) == 0


class BuildQA:

    _ALL_CHECK_IDS = {
        "residual_placeholder",
        "blank_page",
        "broken_image_ref",
        "font_too_small",
        "text_overflow",
        "title_duplicate",
        "page_count_mismatch",
        "color_break",
        "font_mismatch",
        "toc_mismatch",
        "element_out_of_bounds",
        "image_stretched",
        "low_contrast_text",
        "spacing_tight",
    }

    def check(
        self,
        pptx_path: str,
        plans: list[dict] | None = None,
        mode: str = "business",
    ) -> QAReport:
        """Run QA checks.

        mode='business' (default): min font 11pt.
        mode='scientific': min font 8pt — figure captions, sequence fonts, and
            superscript citations in journal decks are intentionally small
            (Nature/Cell convention). All other checks unchanged.
        """
        if not os.path.isfile(pptx_path):
            return QAReport(
                total_slides=0, total_checks=0, passed=0,
                fatals=[CheckItem("content", "file_missing", "fatal", -1, "Output file not found")],
            )

        min_font_pt = 8 if mode == "scientific" else _MIN_FONT_PT

        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        all_items: list[CheckItem] = []

        for idx, slide in enumerate(prs.slides):
            slide_elem = slide._element
            all_text = self._extract_all_text(slide_elem)

            all_items.extend(self._check_residual_placeholders(idx, all_text))
            all_items.extend(self._check_blank_page(idx, all_text))
            all_items.extend(self._check_broken_image_ref(idx, slide_elem, prs))
            all_items.extend(self._check_font_too_small(idx, slide_elem, min_font_pt))
            all_items.extend(self._check_text_overflow(idx, slide))
            all_items.extend(self._check_title_duplicate(idx, slide_elem))
            all_items.extend(self._check_color_break(idx, slide_elem))
            all_items.extend(self._check_font_mismatch(idx, slide_elem))
            all_items.extend(self._check_element_out_of_bounds(idx, slide, slide_width, slide_height))
            all_items.extend(self._check_image_stretched(idx, slide))
            all_items.extend(self._check_low_contrast_text(idx, slide_elem, prs=prs))
            all_items.extend(self._check_spacing(idx, slide))

        planned_count = len(plans) if plans else 0
        all_items.extend(self._check_page_count(total_slides, planned_count))
        all_items.extend(self._check_toc_mismatch(prs, plans))

        fatals = [i for i in all_items if i.severity == "fatal"]
        warnings = [i for i in all_items if i.severity == "warning"]
        reviews = [i for i in all_items if i.severity == "review"]
        passed = len(all_items) - len(fatals) - len(warnings) - len(reviews)

        return QAReport(
            total_slides=total_slides,
            total_checks=len(all_items),
            passed=max(0, passed),
            fatals=fatals,
            warnings=warnings,
            reviews=reviews,
        )

    def auto_fix(self, pptx_path: str, report: QAReport) -> None:
        if not os.path.isfile(pptx_path):
            return

        prs = Presentation(pptx_path)

        for item in report.fatals + report.warnings:
            if not item.auto_fixable:
                continue

            if item.check_id == "residual_placeholder":
                self._fix_residual_placeholders(prs, item.slide_index)

            elif item.check_id == "font_too_small":
                self._fix_font_too_small(prs, item.slide_index)

        prs.save(pptx_path)

    def format_report(self, report: QAReport) -> str:
        lines = [
            "=" * 60,
            "  Build-Mode QA Report",
            "=" * 60,
            f"Total slides: {report.total_slides} | Checks: {report.total_checks} | "
            f"Passed: {report.passed} | Fatals: {len(report.fatals)} | "
            f"Warnings: {len(report.warnings)} | Reviews: {len(report.reviews)}",
            "",
        ]

        categories: dict[str, list[CheckItem]] = {}
        for item in report.fatals + report.warnings + report.reviews:
            categories.setdefault(item.category, []).append(item)

        for cat, items in categories.items():
            has_fatal = any(i.severity == "fatal" for i in items)
            has_warning = any(i.severity == "warning" for i in items)
            has_review = any(i.severity == "review" for i in items)
            if has_fatal:
                lines.append(f"  {cat}: FATAL")
            elif has_warning:
                lines.append(f"  {cat}: WARNING")
            elif has_review:
                lines.append(f"  {cat}: REVIEW")
            for item in items:
                prefix = item.severity.upper()
                lines.append(f"    - Slide {item.slide_index + 1}: [{prefix}] {item.message}")
                if item.detail:
                    lines.append(f"      Detail: {item.detail}")

        lines.append("")
        if report.is_passable:
            lines.append("Result: PASS - ready for delivery")
        else:
            lines.append("Result: FAIL - fatal issues remain")

        return "\n".join(lines)

    # ── Check implementations ──

    def _extract_all_text(self, slide_elem) -> list[str]:
        a_ns = _NS["a"]
        texts = []
        for t in slide_elem.iter(f"{{{a_ns}}}t"):
            if t.text:
                texts.append(t.text.strip())
        return texts

    def _check_residual_placeholders(self, slide_idx: int, all_text: list[str]) -> list[CheckItem]:
        items = []
        for text in all_text:
            for pattern in _PLACEHOLDER_PATTERNS:
                if pattern.search(text):
                    items.append(CheckItem(
                        "content", "residual_placeholder", "fatal", slide_idx,
                        "Residual placeholder text found",
                        detail=text, auto_fixable=True,
                    ))
                    break
        return items

    def _check_blank_page(self, slide_idx: int, all_text: list[str]) -> list[CheckItem]:
        non_empty = [t for t in all_text if t and len(t) > 1]
        if not non_empty:
            return [CheckItem(
                "content", "blank_page", "fatal", slide_idx,
                "Blank page — no text content",
                auto_fixable=False,
            )]
        return []

    def _check_broken_image_ref(self, slide_idx: int, slide_elem, prs) -> list[CheckItem]:
        items = []
        a_ns = _NS["a"]
        r_ns = _NS["r"]

        for blip in slide_elem.iter(f"{{{a_ns}}}blip"):
            embed = blip.get(f"{{{r_ns}}}embed")
            if embed:
                try:
                    slide_part = prs.slides[slide_idx].part
                    rel = slide_part.rels.get(embed)
                    if rel is None:
                        items.append(CheckItem(
                            "material", "broken_image_ref", "fatal", slide_idx,
                            "Broken image reference",
                            detail=embed, auto_fixable=False,
                        ))
                except Exception:
                    pass
        return items

    def _check_font_too_small(self, slide_idx: int, slide_elem, min_font_pt: int = _MIN_FONT_PT) -> list[CheckItem]:
        a_ns = _NS["a"]
        items = []
        for rPr in slide_elem.iter(f"{{{a_ns}}}rPr"):
            sz = rPr.get("sz")
            if sz:
                try:
                    pt = int(sz) / 100
                    if pt < min_font_pt:
                        items.append(CheckItem(
                            "layout", "font_too_small", "warning", slide_idx,
                            f"Font size too small ({pt:.0f}pt)",
                            auto_fixable=True,
                        ))
                except ValueError:
                    pass
        return items[:3]

    def _check_text_overflow(self, slide_idx: int, slide) -> list[CheckItem]:
        items = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if not tf.word_wrap:
                continue
            text = tf.text
            if not text or not text.strip():
                continue
            try:
                width_inches = shape.width / 914400
                height_inches = shape.height / 914400
                total_est_lines = 0
                for p in tf.paragraphs:
                    if not p.text or not p.text.strip():
                        continue
                    font_size_pt = 12
                    if p.font.size:
                        font_size_pt = p.font.size / 12700
                    else:
                        for r in p.runs:
                            if r.font.size:
                                font_size_pt = r.font.size / 12700
                                break
                    chars_per_line = max(1, int(width_inches * 72 / font_size_pt * 0.65))
                    # CJK-aware: full-width chars weigh ~1.5x Latin (est. width ratio)
                    p_text = p.text
                    cjk = sum(1 for ch in p_text if ord(ch) > 0x2E80)
                    latin = len(p_text) - cjk
                    weighted = cjk * 1.5 + latin
                    total_est_lines += max(1, weighted / chars_per_line)
                if total_est_lines == 0:
                    continue
                avg_font_pt = 12
                for p in tf.paragraphs:
                    if p.font.size:
                        avg_font_pt = p.font.size / 12700
                        break
                    for r in p.runs:
                        if r.font.size:
                            avg_font_pt = r.font.size / 12700
                            break
                    if avg_font_pt != 12:
                        break
                line_height_in = avg_font_pt * 1.4 / 72
                max_lines = height_inches / line_height_in if line_height_in > 0 else 999
                if total_est_lines > max_lines * 2.0:
                    items.append(CheckItem(
                        "layout", "text_overflow", "warning", slide_idx,
                        "Text may overflow textbox",
                        detail=f"est_lines={total_est_lines:.0f} vs max={max_lines:.0f}",
                        auto_fixable=False,
                    ))
            except Exception:
                pass
        return items[:3]

    def _check_title_duplicate(self, slide_idx: int, slide_elem) -> list[CheckItem]:
        items = []
        a_ns = _NS["a"]
        title_texts: list[str] = []
        for t in slide_elem.iter(f"{{{a_ns}}}t"):
            if t.text and t.text.strip():
                parent_run = t.getparent()
                if parent_run is not None:
                    rPr = parent_run.find(f"{{{a_ns}}}rPr")
                    if rPr is not None:
                        sz = rPr.get("sz")
                        if sz and int(sz) >= 2800:
                            title_texts.append(t.text.strip())

        if len(title_texts) >= 2:
            from collections import Counter
            counts = Counter(title_texts)
            for text, count in counts.items():
                if count >= 2:
                    items.append(CheckItem(
                        "content", "title_duplicate", "warning", slide_idx,
                        f"Duplicate title text ({count}x)",
                        detail=text, auto_fixable=False,
                    ))
        return items

    def _check_color_break(self, slide_idx: int, slide_elem) -> list[CheckItem]:
        a_ns = _NS["a"]
        items = []
        seen_colors: set[str] = set()
        for srgb in slide_elem.iter(f"{{{a_ns}}}srgbClr"):
            val = srgb.get("val", "").upper()
            if val and len(val) == 6 and val not in seen_colors:
                seen_colors.add(val)
                items.append(CheckItem(
                    "design", "color_break", "review", slide_idx,
                    "Non-palette color found (review for intent)",
                    detail=f"#{val}", auto_fixable=False,
                ))
        return items[:3]

    def _find_slide_bg_hex(self, prs, slide_elem) -> str | None:
        """Detect the dominant slide background color.

        Checks, in order of priority:
        1. Full-bleed solid-fill rectangles (incl. inside groups) sized >= 50% of slide
        2. The slide's <p:bg> element (solidFill)
        3. Theme background color
        Returns hex like '#RRGGBB' or None (caller falls back to white).
        """
        a_ns = _NS["a"]
        p_ns = _NS["p"]
        sw = prs.slide_width or 12192000
        sh = prs.slide_height or 6858000

        def _rect_is_full_bleed(sp):
            try:
                if sp.get("x") is None or sp.get("y") is None:
                    return None
                x = int(sp.get("x")); y = int(sp.get("y"))
                w = int(sp.get("cx", 0)); h = int(sp.get("cy", 0))
            except (TypeError, ValueError):
                return None
            if x <= 0 and y <= 0 and w >= sw * 0.5 and h >= sh * 0.5:
                return True
            return False

        # Walk shape tree (incl. groups) for full-bleed solidFill rects
        for sp in slide_elem.iter(f"{{{p_ns}}}sp"):
            xfrm = sp.find(f"{{{p_ns}}}spPr")
            if xfrm is None:
                continue
            xfrm = xfrm.find(f"{{{a_ns}}}xfrm")
            if xfrm is None:
                continue
            off = xfrm.find(f"{{{a_ns}}}off")
            ext = xfrm.find(f"{{{a_ns}}}ext")
            if off is None or ext is None:
                continue
            try:
                x = int(off.get("x")); y = int(off.get("y"))
                w = int(ext.get("cx")); h = int(ext.get("cy"))
            except (TypeError, ValueError):
                continue
            if x <= 0 and y <= 0 and w >= sw * 0.5 and h >= sh * 0.5:
                spPr = sp.find(f"{{{p_ns}}}spPr")
                for el in spPr.iter(f"{{{a_ns}}}solidFill"):
                    srgb = el.find(f"{{{a_ns}}}srgbClr")
                    if srgb is not None:
                        return f"#{srgb.get('val', '').upper()}"
        # Slide background
        bg = slide_elem.find(f"{{{p_ns}}}cSld")
        if bg is not None:
            bg_el = bg.find(f"{{{p_ns}}}bg")
            if bg_el is not None:
                for el in bg_el.iter(f"{{{a_ns}}}srgbClr"):
                    return f"#{el.get('val', '').upper()}"
        return None

    def _sp_box(self, sp) -> tuple | None:
        """Return (x, y, cx, cy) of a <p:sp> in its parent coordinate space, or None."""
        a_ns = _NS["a"]
        p_ns = _NS["p"]
        spPr = sp.find(f"{{{p_ns}}}spPr")
        if spPr is None:
            spPr = sp.find(f"{{{a_ns}}}spPr")
        if spPr is None:
            return None
        xfrm = spPr.find(f"{{{a_ns}}}xfrm")
        if xfrm is None:
            return None
        off = xfrm.find(f"{{{a_ns}}}off")
        ext = xfrm.find(f"{{{a_ns}}}ext")
        if off is None or ext is None:
            return None
        try:
            return (int(off.get("x")), int(off.get("y")),
                    int(ext.get("cx")), int(ext.get("cy")))
        except (TypeError, ValueError):
            return None

    def _sp_fill_hex(self, sp) -> str | None:
        """Direct spPr solidFill hex (not recursive into text runs), or None if not solid."""
        a_ns = _NS["a"]
        p_ns = _NS["p"]
        spPr = sp.find(f"{{{p_ns}}}spPr")
        if spPr is None:
            spPr = sp.find(f"{{{a_ns}}}spPr")
        if spPr is None:
            return None
        solid = spPr.find(f"{{{a_ns}}}solidFill")
        if solid is None:
            return None
        srgb = solid.find(f"{{{a_ns}}}srgbClr")
        if srgb is None:
            return None
        return f"#{srgb.get('val', '').upper()}"

    def _rpr_parent_sp(self, rPr):
        """Walk up from an a:rPr to its enclosing <p:sp>, or None."""
        p_ns = _NS["p"]
        el = rPr.getparent()
        while el is not None:
            if el.tag == f"{{{p_ns}}}sp":
                return el
            el = el.getparent()
        return None

    def _find_text_bg_hex(self, text_sp, container) -> str | None:
        """Find the filled sibling rect in `container` backing this text shape.

        Uses point-in-shape on the text shape's center point (coordinates are
        relative to `container`, so group nesting needs no slide-level math) and
        picks the smallest containing filled rect (most specific = the direct
        backing, e.g. the code-block dark panel behind inner text).
        Returns hex or None.
        """
        tb = self._sp_box(text_sp)
        if tb is None:
            return None
        tcx, tcy = tb[0] + tb[2] / 2.0, tb[1] + tb[3] / 2.0
        best = None
        best_area = None
        for sp in container:
            if sp is text_sp:
                continue
            fill = self._sp_fill_hex(sp)
            if fill is None:
                continue
            b = self._sp_box(sp)
            if b is None:
                continue
            sx, sy, sw, sh = b
            if sx <= tcx <= sx + sw and sy <= tcy <= sy + sh:
                area = sw * sh
                if best_area is None or area < best_area:
                    best, best_area = fill, area
        return best

    def _check_low_contrast_text(self, slide_idx: int, slide_elem, slide_bg_hex: str = "#FFFFFF",
                                 prs=None) -> list[CheckItem]:
        """Check text-to-background contrast (WCAG AA).

        Background resolution per text run, in priority order:
        1. A filled sibling rect in the same container that backs the text shape
           (handles group-relative coords; e.g. code-block dark panel).
        2. The slide-level background (fallback for un-backed text like subtitles).
        """
        a_ns = _NS["a"]
        items = []
        if prs is not None:
            detected = self._find_slide_bg_hex(prs, slide_elem)
            if detected:
                slide_bg_hex = detected
        for rPr in slide_elem.iter(f"{{{a_ns}}}rPr"):
            if not _is_visible_text_color(rPr, a_ns):
                continue
            color = _get_rpr_color_hex(rPr, a_ns)
            if not color:
                continue
            bg_hex = slide_bg_hex
            text_sp = self._rpr_parent_sp(rPr)
            if text_sp is not None:
                container = text_sp.getparent()
                if container is not None:
                    bg = self._find_text_bg_hex(text_sp, container)
                    if bg:
                        bg_hex = bg
            sz = rPr.get("sz")
            pt = float(int(sz)) / 100.0 if sz else 12.0
            try:
                ratio, _ok = check_contrast(color, bg_hex)
            except Exception:
                continue
            if pt >= 24:  # large text — AA threshold 3:1
                if ratio < 3.0:
                    items.append(CheckItem(
                        "design", "low_contrast_text", "fatal", slide_idx,
                        "Large text contrast below WCAG AA (3:1)",
                        detail=f"{color} on {bg_hex} = {ratio:.2f}:1",
                        auto_fixable=False,
                    ))
            else:  # body text — AA threshold 4.5:1
                if ratio < 4.5:
                    items.append(CheckItem(
                        "design", "low_contrast_text", "fatal", slide_idx,
                        "Text contrast below WCAG AA (4.5:1)",
                        detail=f"{color} on {bg_hex} = {ratio:.2f}:1",
                        auto_fixable=False,
                    ))
        return items[:5]

    def _check_font_mismatch(self, slide_idx: int, slide_elem) -> list[CheckItem]:
        a_ns = _NS["a"]
        items = []
        seen_fonts: set[str] = set()
        for latin in slide_elem.iter(f"{{{a_ns}}}latin"):
            typeface = latin.get("typeface", "")
            if typeface and typeface not in seen_fonts and not typeface.startswith("+"):
                seen_fonts.add(typeface)
                items.append(CheckItem(
                    "design", "font_mismatch", "review", slide_idx,
                    "Non-standard font found (review for intent)",
                    detail=typeface, auto_fixable=False,
                ))
        return items[:3]

    _BLEED_MAX_FRACTION = 0.5
    _MINOR_BLEED_FRACTION = 0.05
    _DECO_ALPHA_THRESHOLD = 0.20

    def _is_decorative_bleed(self, shape, slide_width: int, slide_height: int) -> bool:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return False

        alpha = self._get_fill_alpha(shape)
        if alpha is not None and alpha <= self._DECO_ALPHA_THRESHOLD:
            return True

        return False

    def _get_fill_alpha(self, shape) -> float | None:
        try:
            spPr = shape._element.find(f"{{{_NS['p']}}}spPr")
            if spPr is None:
                spPr = shape._element.find(f"{{{_NS['a']}}}spPr")
            if spPr is None:
                return None
            for solidFill in spPr.iter(f"{{{_NS['a']}}}solidFill"):
                for srgb in solidFill.iter(f"{{{_NS['a']}}}srgbClr"):
                    alpha_el = srgb.find(f"{{{_NS['a']}}}alpha")
                    if alpha_el is not None:
                        return int(alpha_el.get("val", "100000")) / 100000.0
                for scheme in solidFill.iter(f"{{{_NS['a']}}}schemeClr"):
                    alpha_el = scheme.find(f"{{{_NS['a']}}}alpha")
                    if alpha_el is not None:
                        return int(alpha_el.get("val", "100000")) / 100000.0
        except Exception:
            pass
        return None

    def _check_element_out_of_bounds(
        self,
        slide_idx: int,
        slide,
        slide_width: int,
        slide_height: int,
    ) -> list[CheckItem]:
        items = []
        for shape in slide.shapes:
            x = shape.left
            y = shape.top
            w = shape.width
            h = shape.height

            if x < 0 or y < 0 or (x + w) > slide_width or (y + h) > slide_height:
                bleed_x = max(-x, 0) + max((x + w) - slide_width, 0)
                bleed_y = max(-y, 0) + max((y + h) - slide_height, 0)
                bleed_total = bleed_x + bleed_y
                slide_max = slide_width + slide_height

                is_deco = self._is_decorative_bleed(shape, slide_width, slide_height)
                bleed_fraction = bleed_total / slide_max if slide_max > 0 else 1.0

                if is_deco and bleed_fraction <= self._BLEED_MAX_FRACTION:
                    items.append(CheckItem(
                        "design", "element_out_of_bounds", "review", slide_idx,
                        "Decorative shape extends beyond slide (intentional bleed)",
                        detail=(
                            f"shape='{shape.name}' pos=({x},{y}) size=({w},{h}) "
                            f"slide=({slide_width},{slide_height})"
                        ),
                        auto_fixable=False,
                    ))
                elif bleed_fraction <= self._MINOR_BLEED_FRACTION:
                    items.append(CheckItem(
                        "layout", "element_out_of_bounds", "warning", slide_idx,
                        "Shape slightly extends beyond slide boundaries",
                        detail=(
                            f"shape='{shape.name}' pos=({x},{y}) size=({w},{h}) "
                            f"slide=({slide_width},{slide_height}) bleed={bleed_fraction:.1%}"
                        ),
                        auto_fixable=False,
                    ))
                else:
                    items.append(CheckItem(
                        "layout", "element_out_of_bounds", "fatal", slide_idx,
                        "Shape extends beyond slide boundaries",
                        detail=(
                            f"shape='{shape.name}' pos=({x},{y}) size=({w},{h}) "
                            f"slide=({slide_width},{slide_height}) bleed={bleed_fraction:.1%}"
                        ),
                        auto_fixable=False,
                    ))
        return items

    def _check_image_stretched(self, slide_idx: int, slide) -> list[CheckItem]:
        items = []
        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue
            try:
                img = shape.image
                img_width_px = img.size[0]
                img_height_px = img.size[1]
                if img_width_px <= 0 or img_height_px <= 0:
                    continue

                img_ratio = img_width_px / img_height_px
                shape_ratio = shape.width / shape.height

                if abs(img_ratio - shape_ratio) > 0.15:
                    items.append(CheckItem(
                        "layout", "image_stretched", "warning", slide_idx,
                        "Image aspect ratio differs from shape",
                        detail=(
                            f"shape='{shape.name}' img_ratio={img_ratio:.2f} "
                            f"shape_ratio={shape_ratio:.2f}"
                        ),
                        auto_fixable=False,
                    ))
            except Exception:
                pass
        return items

    def _check_page_count(self, total_slides: int, planned_count: int) -> list[CheckItem]:
        if planned_count > 0 and total_slides != planned_count:
            return [CheckItem(
                "count", "page_count_mismatch", "fatal", -1,
                f"Page count mismatch: {total_slides} actual vs {planned_count} planned",
                auto_fixable=False,
            )]
        return []

    def _check_toc_mismatch(self, prs, plans: list[dict] | None) -> list[CheckItem]:
        items = []
        if not plans:
            return items

        toc_slide = None
        toc_idx = -1
        for idx, slide in enumerate(prs.slides):
            a_ns = _NS["a"]
            for t in slide._element.iter(f"{{{a_ns}}}t"):
                if t.text and any(kw in t.text.lower() for kw in ("目录", "contents", "agenda")):
                    toc_slide = slide
                    toc_idx = idx
                    break
            if toc_slide is not None:
                break

        if toc_slide is None or toc_idx == -1:
            return items

        toc_text: list[str] = []
        a_ns = _NS["a"]
        for t in toc_slide._element.iter(f"{{{a_ns}}}t"):
            if t.text and t.text.strip():
                toc_text.append(t.text.strip())

        content_sections = sum(
            1 for p in plans
            if isinstance(p, dict) and p.get("page_type") == "transition"
        )

        if content_sections > 0 and len(toc_text) > 1 and content_sections != len(toc_text) - 1:
            items.append(CheckItem(
                "content", "toc_mismatch", "warning", toc_idx,
                f"TOC lists {len(toc_text) - 1} items but content has {content_sections} sections",
                auto_fixable=False,
            ))
        return items

    def _check_spacing(self, slide_idx: int, slide) -> list[CheckItem]:
        """Check minimum spacing between shapes on a slide."""
        items = []
        shapes = list(slide.shapes)

        if len(shapes) < 2:
            return items

        # Minimum gap: 30px = 30/96 inches = 0.3125 inches
        min_gap_inches = 0.3125

        for i, s1 in enumerate(shapes):
            for s2 in shapes[i + 1:]:
                gap = self._calc_gap(s1, s2)
                if 0 < gap < min_gap_inches:
                    items.append(CheckItem(
                        "layout", "spacing_tight", "review", slide_idx,
                        f"Shapes too close: {gap:.2f}in gap (min: {min_gap_inches:.2f}in)",
                        auto_fixable=False,
                    ))
                    break  # Only report once per shape pair
        return items

    def _calc_gap(self, s1, s2) -> float:
        """Calculate the minimum gap between two shapes in inches.

        Returns positive value if shapes are separated, 0 if overlapping,
        negative if they overlap.
        """
        try:
            # Get bounding boxes in EMU
            x1, y1 = s1.left, s1.top
            w1, h1 = s1.width, s1.height
            x2, y2 = s2.left, s2.top
            w2, h2 = s2.width, s2.height

            # Calculate gap
            # Horizontal gap
            if x1 + w1 < x2:  # s1 is left of s2
                h_gap = x2 - (x1 + w1)
            elif x2 + w2 < x1:  # s2 is left of s1
                h_gap = x1 - (x2 + w2)
            else:  # Overlapping horizontally
                h_gap = 0

            # Vertical gap
            if y1 + h1 < y2:  # s1 is above s2
                v_gap = y2 - (y1 + h1)
            elif y2 + h2 < y1:  # s2 is above s1
                v_gap = y1 - (y2 + h2)
            else:  # Overlapping vertically
                v_gap = 0

            # Return the minimum gap
            if h_gap > 0 and v_gap > 0:
                return min(h_gap, v_gap) / 914400  # EMU to inches
            elif h_gap > 0:
                return h_gap / 914400
            elif v_gap > 0:
                return v_gap / 914400
            else:
                return 0  # Overlapping

        except Exception:
            return 0  # Skip if calculation fails

    # ── Auto-fix implementations ──

    def _fix_residual_placeholders(self, prs, slide_idx: int) -> bool:
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            return False

        slide_elem = prs.slides[slide_idx]._element
        a_ns = _NS["a"]
        fixed = False

        for t in list(slide_elem.iter(f"{{{a_ns}}}t")):
            if not t.text:
                continue
            text = t.text.strip()
            for pattern in _PLACEHOLDER_PATTERNS:
                if pattern.search(text):
                    t.text = ""
                    for child in list(t):
                        t.remove(child)
                    fixed = True
                    break
        return fixed

    def _fix_font_too_small(self, prs, slide_idx: int) -> bool:
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            return False

        slide_elem = prs.slides[slide_idx]._element
        a_ns = _NS["a"]
        fixed = False

        for tag_name in (f"{{{a_ns}}}rPr", f"{{{a_ns}}}defRPr"):
            for rPr in slide_elem.iter(tag_name):
                sz = rPr.get("sz")
                if sz:
                    try:
                        pt = int(sz) / 100
                        if pt < _MIN_FONT_PT:
                            rPr.set("sz", str(_MIN_FONT_HUNDREDTHS))
                            fixed = True
                    except ValueError:
                        pass
        return fixed
