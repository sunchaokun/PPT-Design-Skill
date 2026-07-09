"""Phase 4: PPT Renderer — professional python-pptx generation."""

from __future__ import annotations

import hashlib
import os
from datetime import datetime
from pathlib import Path
from typing import Any

from ppt_pro_max.decider.design_decider import PageDesign
from ppt_pro_max.content.content_generator import PageContent
from ppt_pro_max.renderer.theme_mapper import ThemeMapper
from ppt_pro_max.renderer.layout_registry import LayoutRegistry, SLIDE_WIDTH, SLIDE_HEIGHT
from ppt_pro_max.renderer.chart_builder import ChartBuilder
from ppt_pro_max.renderer.effects import Effects
from ppt_pro_max.renderer.image_fetcher import ImageFetcher
from ppt_pro_max.qa.qa_gates import QAGates

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree

    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


class PPTRenderer:
    def __init__(self, image_mode: str = "placeholder", image_config: dict[str, Any] | None = None):
        if not _PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is required. Install with: pip install python-pptx")
        self.theme_mapper = ThemeMapper()
        self.layout_registry = LayoutRegistry()
        self.chart_builder = ChartBuilder()
        self.effects = Effects()
        self.qa = QAGates()
        self.image_fetcher = ImageFetcher(mode=image_mode, **(image_config or {}))
    def render(
        self,
        page_designs: list[PageDesign],
        page_contents: list[PageContent],
        output_path: str | None = None,
        fetch_images: bool = False,
        theme_name: str | None = None,
        design_system: dict[str, Any] | None = None,
        composed_theme: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        if not page_designs:
            raise ValueError("No page designs to render")

        ds = design_system or self._default_design_system()
        ppt_theme = self.theme_mapper.map(ds, theme_name=theme_name, composed_theme=composed_theme)

        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)

        self.theme_mapper.apply_theme(prs, ppt_theme)

        for design, content in zip(page_designs, page_contents):
            self._render_slide(prs, design, content, ppt_theme)

        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"presentation_{timestamp}.pptx"

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        prs.save(output_path)

        qa_result = self.qa.check(output_path, page_designs, page_contents)

        return {
            "output_path": os.path.abspath(output_path),
            "page_count": len(page_designs),
            "strategy": "generated",
            "theme": ppt_theme.get("name", "default"),
            "qa": qa_result,
        }

    def _default_design_system(self) -> dict[str, Any]:
        return {
            "colors": {
                "primary": "#2563EB",
                "on-primary": "#FFFFFF",
                "secondary": "#64748B",
                "accent": "#F97316",
                "background": "#FFFFFF",
                "foreground": "#1E293B",
                "muted": "#F1F5F9",
                "muted-foreground": "#94A3B8",
                "border": "#E2E8F0",
                "destructive": "#EF4444",
            },
            "typography": {
                "heading": "Inter",
                "body": "Inter",
            },
        }

    def _render_slide(self, prs: Presentation, design: PageDesign, content: PageContent, theme: dict) -> None:
        layout_def = self.layout_registry.get_layout(design.layout)
        if layout_def is None:
            layout_def = self.layout_registry.get_layout("content-with-title")

        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        is_hero_slide = layout_def.get("name") in ("Title Slide", "CTA Closing")

        if is_hero_slide:
            self._apply_background(slide, design, theme)
            hero_image_rendered = False
            for ph_name, ph_def in layout_def.get("placeholders", {}).items():
                if ph_def.get("type") == "image":
                    self._render_image_placeholder(slide, ph_def, content, theme)
                    hero_image_rendered = True
            if hero_image_rendered:
                self._apply_dark_overlay(slide)
            for ph_name, ph_def in layout_def.get("placeholders", {}).items():
                if ph_def.get("type") != "image":
                    self._render_placeholder(slide, ph_name, ph_def, content, theme, design)
        else:
            self._apply_background(slide, design, theme)
            if design.full_bleed:
                self._apply_full_bleed_bg(slide, content, theme)
            if design.break_pattern:
                self._apply_pattern_break(slide, theme)
            for ph_name, ph_def in layout_def.get("placeholders", {}).items():
                self._render_placeholder(slide, ph_name, ph_def, content, theme, design)

        self._apply_transition(slide, design.transition)

    def _apply_dark_overlay(self, slide) -> None:
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT),
        )
        fill = overlay.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
        overlay.line.fill.background()
        overlay.rotation = 0

        try:
            sp_pr = overlay._element.find(qn("p:spPr"))
            solid_fill = sp_pr.find(qn("a:solidFill"))
            if solid_fill is not None:
                srgb = solid_fill.find(qn("a:srgbClr"))
                if srgb is not None:
                    alpha = etree.SubElement(srgb, qn("a:alpha"))
                    alpha.set("val", "55000")
        except Exception:
            pass

    def _apply_background(self, slide, design: PageDesign, theme: dict) -> None:
        bg_color = theme.get("colors", {}).get("background", "#FFFFFF")
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip("#"))

    def _apply_full_bleed_bg(self, slide, content: PageContent, theme: dict) -> None:
        accent = theme.get("colors", {}).get("accent", "#F97316")
        muted = theme.get("colors", {}).get("muted", "#F1F5F9")
        self.effects.add_gradient_background(slide, muted, accent)

    def _apply_pattern_break(self, slide, theme: dict) -> None:
        accent = theme.get("colors", {}).get("accent", "#F97316")
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(SLIDE_WIDTH), Inches(0.06),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(accent.lstrip("#"))
        bar.line.fill.background()

    def _render_placeholder(
        self, slide, ph_name: str, ph_def: dict, content: PageContent, theme: dict, design: PageDesign
    ) -> None:
        ph_type = ph_def.get("type", "text")

        if ph_type == "decoration":
            self._render_decoration(slide, ph_name, ph_def, theme)
            return

        if ph_type == "image":
            self._render_image_placeholder(slide, ph_def, content, theme)
            return

        if ph_type == "chart":
            self._render_chart_placeholder(slide, ph_def, content, theme, design)
            return

        if ph_name == "cta_button":
            self._render_cta_button(slide, ph_def, content, theme)
            return

        if ph_name.startswith("card"):
            self._render_card(slide, ph_name, ph_def, content, theme)
            return

        if ph_name.startswith("metric"):
            self._render_metric(slide, ph_name, ph_def, content, theme)
            return

        if ph_name == "big_number":
            self._render_big_number(slide, ph_def, content, theme)
            return

        if ph_name == "label":
            self._render_label(slide, ph_def, content, theme)
            return

        if ph_name == "quote_mark":
            self._render_quote_mark(slide, ph_def, theme)
            return

        text = self._get_placeholder_text(ph_name, content)
        if text is None:
            return

        is_body = ph_name == "body"
        self._add_textbox(slide, ph_def, text, theme, is_body=is_body)

    def _render_decoration(self, slide, ph_name: str, ph_def: dict, theme: dict) -> None:
        deco_type = ph_def.get("decoration_type", "")
        accent = theme.get("colors", {}).get("accent", "#F97316")
        primary = theme.get("colors", {}).get("primary", "#2563EB")

        if deco_type == "accent_bar":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(ph_def["x"]), Inches(ph_def["y"]),
                Inches(ph_def["width"]), Inches(ph_def["height"]),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor.from_string(accent.lstrip("#"))
            bar.line.fill.background()

        elif deco_type == "left_accent":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(ph_def["x"]), Inches(ph_def["y"]),
                Inches(ph_def["width"]), Inches(ph_def["height"]),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor.from_string(primary.lstrip("#"))
            bar.line.fill.background()

        elif deco_type == "title_underline":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(ph_def["x"]), Inches(ph_def["y"]),
                Inches(ph_def["width"]), Inches(ph_def["height"]),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor.from_string(accent.lstrip("#"))
            bar.line.fill.background()

    def _add_textbox(self, slide, ph_def: dict, text: str, theme: dict, extra_style: dict | None = None, is_body: bool = False) -> None:
        left = Inches(ph_def["x"])
        top = Inches(ph_def["y"])
        width = Inches(ph_def["width"])
        height = Inches(ph_def["height"])

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        if is_body and "\n" in text:
            lines = text.split("\n")
            for idx, line in enumerate(lines):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                clean = line.strip()
                if clean.startswith("• ") or clean.startswith("- "):
                    clean = clean[2:]

                p.text = clean
                p.level = 1

                font_size = ph_def.get("font_size", 16)
                p.font.size = Pt(font_size)
                p.font.bold = False

                alignment = ph_def.get("alignment", "left")
                align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
                p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)

                color_role = ph_def.get("color_role", "foreground")
                color_hex = theme.get("colors", {}).get(color_role, "#1E293B")
                p.font.color.rgb = RGBColor.from_string(color_hex.lstrip("#"))

                font_name = ph_def.get("font_role", "body")
                font_family = theme.get("typography", {}).get(font_name, "Inter")
                p.font.name = font_family

                p.space_after = Pt(8)
                p.space_before = Pt(4)
        else:
            p = tf.paragraphs[0]
            p.text = text

            font_size = ph_def.get("font_size", 16)
            p.font.size = Pt(font_size)
            p.font.bold = ph_def.get("font_weight") in ("bold", "semibold", "600", "700")

            alignment = ph_def.get("alignment", "left")
            align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
            p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)

            color_role = ph_def.get("color_role", "foreground")
            color_hex = theme.get("colors", {}).get(color_role, "#1E293B")
            p.font.color.rgb = RGBColor.from_string(color_hex.lstrip("#"))

            font_name = ph_def.get("font_role", "body")
            font_family = theme.get("typography", {}).get(font_name, "Inter")
            p.font.name = font_family

            if extra_style and extra_style.get("vertical_center"):
                p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)

    def _render_card(self, slide, ph_name: str, ph_def: dict, content: PageContent, theme: dict) -> None:
        card_index = int(ph_name.replace("card", "")) - 1
        bullets = content.bullets
        card_text = bullets[card_index] if bullets and card_index < len(bullets) else f"Point {card_index + 1}"

        left = Inches(ph_def["x"])
        top = Inches(ph_def["y"])
        width = Inches(ph_def["width"])
        height = Inches(ph_def["height"])

        border_color = theme.get("colors", {}).get("border", "#E2E8F0")
        bg_color = theme.get("colors", {}).get("muted", "#F1F5F9")
        fg_color = theme.get("colors", {}).get("foreground", "#1E293B")
        accent_color = theme.get("colors", {}).get("accent", "#F97316")

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip("#"))
        card.line.color.rgb = RGBColor.from_string(border_color.lstrip("#"))
        card.line.width = Pt(1)

        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(ph_def["x"] + 0.15), Inches(ph_def["y"] + 0.2),
            Inches(0.5), Inches(0.06),
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor.from_string(accent_color.lstrip("#"))
        accent_bar.line.fill.background()

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.5)

        p = tf.paragraphs[0]
        p.text = card_text
        p.font.size = Pt(ph_def.get("font_size", 14))
        p.font.color.rgb = RGBColor.from_string(fg_color.lstrip("#"))
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)

    def _render_metric(self, slide, ph_name: str, ph_def: dict, content: PageContent, theme: dict) -> None:
        metric_index = int(ph_name.replace("metric", "")) - 1
        metrics = content.metrics or []

        if metric_index < len(metrics):
            value = metrics[metric_index].get("value", "—")
            label = metrics[metric_index].get("label", "")
        else:
            value = "—"
            label = ""

        left = Inches(ph_def["x"])
        top = Inches(ph_def["y"])
        width = Inches(ph_def["width"])
        height = Inches(ph_def["height"])

        muted_color = theme.get("colors", {}).get("muted", "#F1F5F9")
        border_color = theme.get("colors", {}).get("border", "#E2E8F0")
        accent_color = theme.get("colors", {}).get("accent", "#F97316")
        fg_color = theme.get("colors", {}).get("muted-foreground", "#94A3B8")

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(muted_color.lstrip("#"))
        card.line.color.rgb = RGBColor.from_string(border_color.lstrip("#"))
        card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.3)

        p_value = tf.paragraphs[0]
        p_value.text = value
        p_value.font.size = Pt(ph_def.get("font_size", 40))
        p_value.font.bold = True
        p_value.font.color.rgb = RGBColor.from_string(accent_color.lstrip("#"))
        p_value.alignment = PP_ALIGN.CENTER

        p_label = tf.add_paragraph()
        p_label.text = label
        p_label.font.size = Pt(12)
        p_label.font.color.rgb = RGBColor.from_string(fg_color.lstrip("#"))
        p_label.alignment = PP_ALIGN.CENTER
        p_label.space_before = Pt(6)

    def _render_big_number(self, slide, ph_def: dict, content: PageContent, theme: dict) -> None:
        big_text = content.title if content.title else ""
        self._add_textbox(slide, ph_def, big_text, theme)

    def _render_label(self, slide, ph_def: dict, content: PageContent, theme: dict) -> None:
        label_text = content.subtitle if content.subtitle else ""
        if not label_text and content.bullets:
            label_text = content.bullets[0]
        if not label_text:
            return
        self._add_textbox(slide, ph_def, label_text, theme)

    def _render_quote_mark(self, slide, ph_def: dict, theme: dict) -> None:
        accent = theme.get("colors", {}).get("accent", "#F97316")
        left = Inches(ph_def["x"])
        top = Inches(ph_def["y"])
        width = Inches(ph_def["width"])
        height = Inches(ph_def["height"])

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "\u201C"
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(accent.lstrip("#"))
        p.alignment = PP_ALIGN.LEFT

    def _render_chart_placeholder(
        self, slide, ph_def: dict, content: PageContent, theme: dict, design: PageDesign
    ) -> None:
        if content.chart_data is None:
            muted_color = theme.get("colors", {}).get("muted", "#F1F5F9")
            border_color = theme.get("colors", {}).get("border", "#E2E8F0")
            hint_color = theme.get("colors", {}).get("muted-foreground", "#94A3B8")
            left = Inches(ph_def["x"])
            top = Inches(ph_def["y"])
            width = Inches(ph_def["width"])
            height = Inches(ph_def["height"])
            placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor.from_string(muted_color.lstrip("#"))
            placeholder.line.color.rgb = RGBColor.from_string(border_color.lstrip("#"))
            placeholder.line.width = Pt(1)
            tf = placeholder.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = design.chart_type or "Chart"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor.from_string(hint_color.lstrip("#"))
            p.alignment = PP_ALIGN.CENTER
            return

        chart_type = content.chart_data.get("type", design.chart_type or "Line Chart")
        data = content.chart_data.get("data", {})
        position = {
            "x": ph_def["x"],
            "y": ph_def["y"],
            "width": ph_def["width"],
            "height": ph_def["height"],
        }
        style = {"colors": theme.get("colors", {})}
        self.chart_builder.build(slide, chart_type, data, style, position)

    def _render_cta_button(self, slide, ph_def: dict, content: PageContent, theme: dict) -> None:
        button_text = content.title if content.title else "Get Started"

        left = Inches(ph_def["x"])
        top = Inches(ph_def["y"])
        width = Inches(ph_def["width"])
        height = Inches(ph_def["height"])

        bg_role = ph_def.get("bg_color_role", "primary")
        bg_color = theme.get("colors", {}).get(bg_role, "#2563EB")

        button = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        button.fill.solid()
        button.fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip("#"))
        button.line.fill.background()

        tf = button.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = button_text
        p.font.size = Pt(ph_def.get("font_size", 18))
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        self.effects.add_shadow(button, blur_pt=8, offset_pt=3, color="#000000", alpha=20)

    def _render_image_placeholder(self, slide, ph_def: dict, content: PageContent, theme: dict) -> None:
        left = Inches(ph_def["x"])
        top = Inches(ph_def["y"])
        width = Inches(ph_def["width"])
        height = Inches(ph_def["height"])

        image_path = self.image_fetcher.fetch(
            keywords=content.image_keywords or content.goal,
            emotion=content.goal,
            goal=content.goal,
            width=int(ph_def["width"] * 96),
            height=int(ph_def["height"] * 96),
        )

        if image_path and os.path.exists(image_path):
            try:
                self._add_picture_cover(slide, image_path, ph_def)
                return
            except Exception:
                pass

        generated = self._generate_placeholder_image(
            content.image_keywords or content.goal,
            ph_def["width"], ph_def["height"], theme,
        )
        if generated:
            try:
                self._add_picture_cover(slide, generated, ph_def)
                return
            except Exception:
                pass

        muted_color = theme.get("colors", {}).get("muted", "#F1F5F9")
        border_color = theme.get("colors", {}).get("border", "#E2E8F0")
        hint_color = theme.get("colors", {}).get("muted-foreground", "#94A3B8")

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(muted_color.lstrip("#"))
        shape.line.color.rgb = RGBColor.from_string(border_color.lstrip("#"))
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        hint = content.image_keywords if content.image_keywords else "Image"
        p.text = hint
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor.from_string(hint_color.lstrip("#"))
        p.alignment = PP_ALIGN.CENTER

    def _add_picture_cover(self, slide, image_path: str, ph_def: dict) -> None:
        box_w = ph_def["width"]
        box_h = ph_def["height"]
        box_ratio = box_w / box_h

        from PIL import Image as PILImage
        img = PILImage.open(image_path)
        img_w, img_h = img.size
        img_ratio = img_w / img_h

        if img_ratio > box_ratio:
            crop_w = int(img_h * box_ratio)
            crop_h = img_h
            left = (img_w - crop_w) // 2
            top = 0
        else:
            crop_w = img_w
            crop_h = int(img_w / box_ratio)
            left = 0
            top = (img_h - crop_h) // 2

        cropped = img.crop((left, top, left + crop_w, top + crop_h))

        import tempfile
        cache_dir = os.path.join(tempfile.gettempdir(), "ppt-cropped")
        os.makedirs(cache_dir, exist_ok=True)
        crop_key = f"crop:{image_path}:{box_w}x{box_h}"
        crop_hash = hashlib.md5(crop_key.encode()).hexdigest()
        cropped_path = os.path.join(cache_dir, f"{crop_hash}.png")
        if not os.path.exists(cropped_path):
            cropped.save(cropped_path, "PNG")

        slide.shapes.add_picture(
            cropped_path,
            Inches(ph_def["x"]),
            Inches(ph_def["y"]),
            Inches(box_w),
            Inches(box_h),
        )

    def _generate_placeholder_image(self, keywords: str, w_inches: float, h_inches: float, theme: dict) -> str | None:
        try:
            from PIL import Image, ImageDraw, ImageFont
            import tempfile, hashlib

            primary = theme.get("colors", {}).get("primary", "#2563EB")
            accent = theme.get("colors", {}).get("accent", "#F97316")
            muted = theme.get("colors", {}).get("muted", "#F1F5F9")
            foreground = theme.get("colors", {}).get("foreground", "#1E293B")

            px_w = max(int(w_inches * 96), 200)
            px_h = max(int(h_inches * 96), 150)

            img = Image.new("RGB", (px_w, px_h))
            draw = ImageDraw.Draw(img)

            r1, g1, b1 = int(primary[1:3], 16), int(primary[3:5], 16), int(primary[5:7], 16)
            r2, g2, b2 = int(accent[1:3], 16), int(accent[3:5], 16), int(accent[5:7], 16)

            for y in range(px_h):
                t = y / px_h
                r = int(r1 + (r2 - r1) * t)
                g = int(g1 + (g2 - g1) * t)
                b = int(b1 + (b2 - b1) * t)
                draw.line([(0, y), (px_w, y)], fill=(r, g, b))

            margin = min(60, px_w // 8, px_h // 8)
            inner_w = px_w - 2 * margin
            inner_h = px_h - 2 * margin
            if inner_w > 100 and inner_h > 80:
                draw.rounded_rectangle(
                    [margin, margin, px_w - margin, px_h - margin],
                    radius=20,
                    fill=(255, 255, 255, 180),
                    outline=(255, 255, 255),
                    width=2,
                )

            cx = px_w // 2
            cy = px_h // 2 - 20

            icon_size = min(60, px_w // 6, px_h // 6)
            if icon_size > 20:
                draw.ellipse(
                    [cx - icon_size // 2, cy - icon_size // 2, cx + icon_size // 2, cy + icon_size // 2],
                    fill=(255, 255, 255),
                    outline=(200, 200, 200),
                    width=2,
                )
                tri_size = icon_size // 4
                draw.polygon(
                    [
                        (cx - tri_size // 2, cy - tri_size),
                        (cx + tri_size, cy),
                        (cx - tri_size // 2, cy + tri_size),
                    ],
                    fill=(r1, g1, b1),
                )

            label_y = cy + icon_size // 2 + 15
            try:
                font = ImageFont.truetype("arial.ttf", 18)
            except Exception:
                font = ImageFont.load_default()
            draw.text((cx, label_y), keywords[:30], fill=(255, 255, 255), font=font, anchor="mt")

            cache_key = hashlib.md5(f"gen:{keywords}:{px_w}x{px_h}".encode()).hexdigest()
            cache_dir = os.path.join(tempfile.gettempdir(), "ppt-gen-images")
            os.makedirs(cache_dir, exist_ok=True)
            img_path = os.path.join(cache_dir, f"{cache_key}.png")
            img.save(img_path, "PNG")
            return img_path
        except Exception:
            return None

    def _get_placeholder_text(self, ph_name: str, content: PageContent) -> str | None:
        if ph_name == "title":
            return content.title
        if ph_name in ("subtitle", "tagline"):
            return content.subtitle
        if ph_name == "body":
            return "\n".join(f"\u2022 {b}" for b in content.bullets) if content.bullets else content.subtitle
        if ph_name == "section_number":
            return str(content.position)
        if ph_name == "section_title":
            return content.title
        if ph_name == "quote_text":
            return content.quote.get("text", "") if content.quote else None
        if ph_name == "quote_author":
            return content.quote.get("author", "") if content.quote else None
        if ph_name == "insight":
            return "\n".join(f"\u2022 {b}" for b in content.bullets) if content.bullets else None
        if ph_name in ("left_col", "right_col"):
            return "\n".join(f"\u2022 {b}" for b in content.bullets) if content.bullets else None
        return None

    def _apply_transition(self, slide, transition_name: str) -> None:
        if not _PPTX_AVAILABLE:
            return
        try:
            transition_xml_map = {
                "fade": "<p:fade/>",
                "fade-slow": "<p:fade/>",
                "push": '<p:push dir="l"/>',
                "push-left": '<p:push dir="l"/>',
                "push-right": '<p:push dir="r"/>',
                "wipe": '<p:wipe dir="d"/>',
                "wipe-down": '<p:wipe dir="d"/>',
                "cut": "<p:cut/>",
            }
            child_xml = transition_xml_map.get(transition_name, "<p:fade/>")
            transition = etree.SubElement(slide._element, qn("p:transition"))
            if transition_name == "fade-slow":
                transition.set("spd", "slow")
            else:
                transition.set("spd", "med")
            child = etree.fromstring(child_xml)
            transition.append(child)
        except Exception:
            pass
