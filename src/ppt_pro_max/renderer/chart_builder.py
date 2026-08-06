"""Chart Builder — python-pptx chart generation with theme colors."""

from __future__ import annotations

from typing import Any

try:
    from lxml import etree
    from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

_LEGACY_CHART_TYPE_MAP: dict[str, Any] = {}
_ENTERPRISE_CHART_TYPE_MAP: dict[str, Any] = {}
if _PPTX_AVAILABLE:
    _LEGACY_CHART_TYPE_MAP = {
        "Bar Chart Vertical": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "Bar Chart Horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
        "Line Chart": XL_CHART_TYPE.LINE,
        "Area Chart": XL_CHART_TYPE.AREA,
        "Pie Chart": XL_CHART_TYPE.PIE,
        "Doughnut Chart": XL_CHART_TYPE.DOUGHNUT,
        "Scatter Plot": XL_CHART_TYPE.XY_SCATTER,
        "Radar Chart": XL_CHART_TYPE.RADAR_FILLED,
    }
    _ENTERPRISE_CHART_TYPE_MAP = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "bar_100": XL_CHART_TYPE.COLUMN_STACKED_100,
        "bar_3d": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_horizontal_stacked": XL_CHART_TYPE.BAR_STACKED,
        "bar_horizontal_100": XL_CHART_TYPE.BAR_STACKED_100,
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "line_stacked": XL_CHART_TYPE.LINE_STACKED,
        "line_stacked_100": XL_CHART_TYPE.LINE_STACKED_100,
        "pie": XL_CHART_TYPE.PIE,
        "pie_3d": XL_CHART_TYPE.PIE,
        "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "doughnut_exploded": XL_CHART_TYPE.DOUGHNUT_EXPLODED,
        "area": XL_CHART_TYPE.AREA,
        "area_stacked": XL_CHART_TYPE.AREA_STACKED,
        "area_stacked_100": XL_CHART_TYPE.AREA_STACKED_100,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "scatter_lines": XL_CHART_TYPE.XY_SCATTER_LINES,
        "scatter_smooth": XL_CHART_TYPE.XY_SCATTER_SMOOTH,
        "radar": XL_CHART_TYPE.RADAR_FILLED,
        "radar_markers": XL_CHART_TYPE.RADAR_MARKERS,
        "bubble": XL_CHART_TYPE.BUBBLE,
        "stock_hlc": XL_CHART_TYPE.STOCK_HLC,
        "stock_ohlc": XL_CHART_TYPE.STOCK_OHLC,
    }

_PIE_TYPES = {"pie", "pie_3d", "pie_exploded", "doughnut", "doughnut_exploded"}
_XY_TYPES = {"scatter", "scatter_lines", "scatter_smooth"}
_BUBBLE_TYPES = {"bubble"}


class ChartBuilder:
    def build(
        self,
        slide,
        chart_type_or_config,
        data=None,
        style=None,
        position=None,
        brand_colors=None,
        brand_fonts=None,
    ) -> Any:
        if not _PPTX_AVAILABLE:
            return None

        if isinstance(chart_type_or_config, dict):
            return self._build_from_config(slide, chart_type_or_config, position,
                                           brand_colors, brand_fonts)

        return self._build_legacy(slide, chart_type_or_config, data, style, position)

    def _build_legacy(
        self,
        slide,
        chart_type: str,
        data: dict[str, Any],
        style: dict[str, Any],
        position: dict[str, float],
    ) -> Any:
        pptx_type = _LEGACY_CHART_TYPE_MAP.get(chart_type)
        if pptx_type is None:
            return None

        chart_data = CategoryChartData()
        chart_data.categories = data.get("labels", ["Q1", "Q2", "Q3", "Q4"])

        values = data.get("values", [10, 25, 45, 80])
        if isinstance(values, list) and len(values) > 0 and isinstance(values[0], list):
            for i, series in enumerate(values):
                chart_data.add_series(f"Series {i + 1}", series)
        else:
            chart_data.add_series("Data", values)

        try:
            chart_frame = slide.shapes.add_chart(
                pptx_type,
                Inches(position.get("x", 1.5)),
                Inches(position.get("y", 1.5)),
                Inches(position.get("width", 10.333)),
                Inches(position.get("height", 4.5)),
                chart_data,
            )

            chart = chart_frame.chart
            chart.has_legend = False

            plot = chart.plots[0]
            colors = style.get("colors", {})
            primary_color = colors.get("primary", "#2563EB")

            try:
                series = plot.series[0]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor.from_string(primary_color.lstrip("#"))
            except Exception:
                pass

            return chart_frame
        except Exception:
            return None

    def _build_from_config(
        self,
        slide,
        chart_config: dict[str, Any],
        position: dict[str, float] | None = None,
        brand_colors: dict[str, str] | None = None,
        brand_fonts: dict[str, str] | None = None,
    ) -> Any:
        chart_type_str = chart_config.get("type", "bar")
        pptx_type = _ENTERPRISE_CHART_TYPE_MAP.get(chart_type_str)
        if pptx_type is None:
            pptx_type = _ENTERPRISE_CHART_TYPE_MAP.get("bar")

        is_xy = chart_type_str in _XY_TYPES
        is_bubble = chart_type_str in _BUBBLE_TYPES
        is_pie = chart_type_str in _PIE_TYPES

        try:
            chart_data = self._build_chart_data(chart_config, is_xy, is_bubble)

            if position is None:
                position = {}

            chart_frame = slide.shapes.add_chart(
                pptx_type,
                Inches(position.get("x", 1.5)),
                Inches(position.get("y", 1.5)),
                Inches(position.get("width", 10.333)),
                Inches(position.get("height", 4.5)),
                chart_data,
            )

            chart = chart_frame.chart
            style = chart_config.get("style", {})

            self._apply_legend(chart, style)
            self._apply_data_labels(chart, style, is_pie)
            self._apply_colors(chart, chart_config, style, brand_colors, is_pie)
            self._apply_chart_title(chart, chart_config)
            self._apply_axes(chart, style, is_pie)
            self._apply_chart_style(chart, style)
            self._apply_text_theme(chart, brand_colors, brand_fonts)

            return chart_frame
        except Exception:
            return None

    def _build_chart_data(self, chart_config: dict, is_xy: bool, is_bubble: bool) -> Any:
        if is_bubble:
            return self._build_bubble_data(chart_config)
        if is_xy:
            return self._build_xy_data(chart_config)
        return self._build_category_data(chart_config)

    def _build_category_data(self, chart_config: dict) -> Any:
        chart_data = CategoryChartData()
        chart_data.categories = chart_config.get("categories", ["Q1", "Q2", "Q3", "Q4"])

        series_list = chart_config.get("series", [])
        if series_list:
            for s in series_list:
                chart_data.add_series(s.get("name", "Data"), s.get("values", []))
        else:
            values = chart_config.get("values", [10, 25, 45, 80])
            if isinstance(values, list) and len(values) > 0 and isinstance(values[0], list):
                for i, series in enumerate(values):
                    chart_data.add_series(f"Series {i + 1}", series)
            else:
                chart_data.add_series("Data", values)
        return chart_data

    def _build_xy_data(self, chart_config: dict) -> Any:
        chart_data = XyChartData()
        series_list = chart_config.get("series", [])
        if series_list:
            for s_idx, s in enumerate(series_list):
                xy_series = chart_data.add_series(s.get("name", f"Series {s_idx + 1}"))
                for point in s.get("values", []):
                    if isinstance(point, (list, tuple)) and len(point) >= 2:
                        xy_series.add_data_point(point[0], point[1])
                    else:
                        xy_series.add_data_point(s_idx + 1, point if not isinstance(point, (list, tuple)) else 0)
        else:
            values = chart_config.get("values", [10, 25, 45, 80])
            xy_series = chart_data.add_series("Data")
            for i, v in enumerate(values):
                if isinstance(v, (list, tuple)) and len(v) >= 2:
                    xy_series.add_data_point(v[0], v[1])
                else:
                    xy_series.add_data_point(i + 1, v)
        return chart_data

    def _build_bubble_data(self, chart_config: dict) -> Any:
        chart_data = BubbleChartData()
        series_list = chart_config.get("series", [])
        if series_list:
            for s_idx, s in enumerate(series_list):
                b_series = chart_data.add_series(s.get("name", f"Series {s_idx + 1}"))
                for point in s.get("values", []):
                    if isinstance(point, (list, tuple)) and len(point) >= 3:
                        b_series.add_data_point(point[0], point[1], point[2])
                    elif isinstance(point, (list, tuple)) and len(point) >= 2:
                        b_series.add_data_point(point[0], point[1], 1)
        else:
            values = chart_config.get("values", [])
            if values:
                b_series = chart_data.add_series("Data")
                for i, v in enumerate(values):
                    if isinstance(v, (list, tuple)) and len(v) >= 3:
                        b_series.add_data_point(v[0], v[1], v[2])
        return chart_data

    def _apply_legend(self, chart: Any, style: dict) -> None:
        show_legend = style.get("show_legend", True)
        chart.has_legend = show_legend
        if show_legend:
            pos_map = {
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "top": XL_LEGEND_POSITION.TOP,
                "left": XL_LEGEND_POSITION.LEFT,
                "right": XL_LEGEND_POSITION.RIGHT,
            }
            legend_position = style.get("legend_position", "bottom")
            chart.legend.position = pos_map.get(legend_position, XL_LEGEND_POSITION.BOTTOM)
            chart.legend.include_in_layout = False

    def _apply_data_labels(self, chart: Any, style: dict, is_pie: bool) -> None:
        show_labels = style.get("show_labels", False)
        if not show_labels:
            return

        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(style.get("label_font_size", 9))

        show_value = style.get("show_value", True)
        show_percentage = style.get("show_percentage", False)
        show_category = style.get("show_category_name", False)

        if is_pie:
            show_percentage = style.get("show_percentage", True)
            show_value = style.get("show_value", False)

        try:
            data_labels.show_value = show_value
        except Exception:
            pass
        try:
            data_labels.show_percentage = show_percentage
        except Exception:
            pass
        try:
            data_labels.show_category_name = show_category
        except Exception:
            pass

        number_format = style.get("number_format")
        if number_format:
            try:
                data_labels.number_format = number_format
            except Exception:
                pass

        label_position = style.get("label_position")
        if label_position and _PPTX_AVAILABLE:
            pos_map = {
                "center": XL_LABEL_POSITION.CENTER,
                "inside_end": XL_LABEL_POSITION.INSIDE_END,
                "outside_end": XL_LABEL_POSITION.OUTSIDE_END,
                "best_fit": XL_LABEL_POSITION.BEST_FIT,
            }
            try:
                data_labels.label_position = pos_map.get(label_position, XL_LABEL_POSITION.OUTSIDE_END)
            except Exception:
                pass

    def _apply_colors(self, chart: Any, chart_config: dict, style: dict,
                      brand_colors: dict[str, str] | None, is_pie: bool) -> None:
        color_scheme = style.get("color_scheme", "brand")
        series_list = chart_config.get("series", [])
        num_series = len(series_list) if series_list else 1
        if is_pie and num_series <= 1:
            categories = chart_config.get("categories", [])
            num_points = len(categories) if categories else 4
            chart_colors = self._resolve_colors(color_scheme, brand_colors, num_points)
        else:
            chart_colors = self._resolve_colors(color_scheme, brand_colors, num_series)
        plot = chart.plots[0]

        if is_pie and num_series <= 1:
            try:
                series = plot.series[0]
                for i, point in enumerate(series.points):
                    try:
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = RGBColor.from_string(
                            chart_colors[i % len(chart_colors)].lstrip("#")
                        )
                    except Exception:
                        break
            except Exception:
                pass
        else:
            for i, series in enumerate(plot.series):
                try:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor.from_string(
                        chart_colors[i % len(chart_colors)].lstrip("#")
                    )
                except Exception:
                    pass

    def _apply_chart_title(self, chart: Any, chart_config: dict) -> None:
        chart_title = chart_config.get("title")
        if chart_title is None:
            chart_title = chart_config.get("style", {}).get("title")
        if chart_title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = chart_title

    def _apply_axes(self, chart: Any, style: dict, is_pie: bool) -> None:
        if is_pie:
            return

        value_axis_title = style.get("value_axis_title")
        category_axis_title = style.get("category_axis_title")
        gridlines = style.get("gridlines", "major_y")
        tick_format = style.get("tick_number_format")

        try:
            value_axis = chart.value_axis
            if value_axis_title:
                value_axis.has_title = True
                value_axis.axis_title.text_frame.paragraphs[0].text = value_axis_title
            if tick_format:
                value_axis.tick_labels.number_format = tick_format

            if "major_y" in gridlines:
                value_axis.has_major_gridlines = True
                try:
                    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                except Exception:
                    pass
            else:
                value_axis.has_major_gridlines = False

            if "major_x" in gridlines:
                try:
                    cat_axis = chart.category_axis
                    cat_axis.has_major_gridlines = True
                except Exception:
                    pass
        except Exception:
            pass

        try:
            category_axis = chart.category_axis
            if category_axis_title:
                category_axis.has_title = True
                category_axis.axis_title.text_frame.paragraphs[0].text = category_axis_title
        except Exception:
            pass

    def _apply_chart_style(self, chart: Any, style: dict) -> None:
        chart_style = style.get("chart_style")
        if chart_style is not None:
            try:
                chart.style = int(chart_style)
            except Exception:
                pass

    def _apply_text_theme(self, chart: Any, brand_colors: dict[str, str] | None,
                          brand_fonts: dict[str, str] | None = None) -> None:
        """Theme all chart text (title, legend, axis labels, data labels) to the brand."""
        if not _PPTX_AVAILABLE or not brand_colors:
            return

        text_color = (brand_colors.get("foreground") or brand_colors.get("text")
                      or brand_colors.get("on-primary") or "#1F2937")
        dark = False
        try:
            bg = (brand_colors.get("background") or "#FFFFFF").lstrip("#")
            r, g, b = int(bg[0:2], 16), int(bg[2:4], 16), int(bg[4:6], 16)
            dark = (0.299 * r + 0.587 * g + 0.114 * b) / 255 < 0.5
        except Exception:
            pass
        if dark and not (brand_colors.get("foreground") or brand_colors.get("text")):
            text_color = "#E2E8F0"

        body_font = (brand_fonts or {}).get("body") or (brand_fonts or {}).get("heading")
        heading_font = (brand_fonts or {}).get("heading") or body_font
        rgb = RGBColor.from_string(text_color.lstrip("#"))

        try:
            if chart.has_title:
                for p in chart.chart_title.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = rgb
                        if heading_font:
                            r.font.name = heading_font
        except Exception:
            pass
        try:
            if chart.has_legend:
                chart.legend.font.color.rgb = rgb
                chart.legend.font.size = Pt(11)
                if body_font:
                    chart.legend.font.name = body_font
        except Exception:
            pass
        try:
            plot = chart.plots[0]
            if plot.has_data_labels:
                plot.data_labels.font.color.rgb = rgb
                if body_font:
                    plot.data_labels.font.name = body_font
        except Exception:
            pass

        # Pie/doughnut charts have no axes — guard the access.
        try:
            cat_axis = chart.category_axis
        except Exception:
            cat_axis = None
        try:
            val_axis = chart.value_axis
        except Exception:
            val_axis = None
        for axis in (cat_axis, val_axis):
            if axis is None:
                continue
            try:
                axis.tick_labels.font.color.rgb = rgb
                axis.tick_labels.font.size = Pt(11)
            except Exception:
                pass
            try:
                if axis.has_title:
                    for p in axis.axis_title.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.color.rgb = rgb
                            if body_font:
                                r.font.name = body_font
            except Exception:
                pass

        # Global chartSpace txPr default — themes any text not overridden above.
        try:
            self._set_chartspace_text_color(chart, text_color, body_font)
        except Exception:
            pass

        # Theme-aware gridlines.
        if val_axis is not None:
            try:
                grid = brand_colors.get("border") or ("#334155" if dark else "#E2E8F0")
                if val_axis.has_major_gridlines:
                    val_axis.major_gridlines.format.line.color.rgb = RGBColor.from_string(grid.lstrip("#"))
            except Exception:
                pass

    def _set_chartspace_text_color(self, chart: Any, text_color: str,
                                   body_font: str | None = None) -> None:
        cs = chart._chartSpace
        txPr = cs.find(qn("c:txPr"))
        if txPr is None:
            return
        p = txPr.find(qn("a:p"))
        if p is None:
            return
        pPr = p.find(qn("a:pPr"))
        if pPr is None:
            return
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            return
        latin = defRPr.find(qn("a:latin"))
        if latin is not None:
            defRPr.remove(latin)
        for sf in defRPr.findall(qn("a:solidFill")):
            defRPr.remove(sf)
        solidFill = etree.SubElement(defRPr, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", text_color.lstrip("#"))
        if body_font:
            latin_el = etree.SubElement(defRPr, qn("a:latin"))
            latin_el.set("typeface", body_font)

    def _resolve_colors(self, color_scheme, brand_colors: dict[str, str] | None, num_series: int) -> list[str]:
        if color_scheme == "auto":
            return ["#4472C4", "#ED7D31", "#A5A5A5", "#FFC000", "#5B9BD5", "#70AD47"]

        if color_scheme == "brand" and brand_colors:
            primary = brand_colors.get("primary", "#2563EB")
            secondary = brand_colors.get("secondary", "#64748B")
            accent = brand_colors.get("accent", "#F97316")
            base = [primary, secondary, accent]
            while len(base) < num_series:
                base.append(self._rotate_hue(base[0], len(base) * 60))
            return base[:num_series]

        if isinstance(color_scheme, list):
            return color_scheme

        return ["#2563EB", "#F97316", "#10B981", "#8B5CF6", "#EF4444", "#06B6D4"]

    @staticmethod
    def _rotate_hue(hex_color: str, degrees: int) -> str:
        import colorsys

        r = int(hex_color[1:3], 16) / 255.0
        g = int(hex_color[3:5], 16) / 255.0
        b = int(hex_color[5:7], 16) / 255.0
        h, lum, s = colorsys.rgb_to_hls(r, g, b)
        h = (h + degrees / 360.0) % 1.0
        r2, g2, b2 = colorsys.hls_to_rgb(h, lum, s)
        return f"#{int(r2 * 255):02X}{int(g2 * 255):02X}{int(b2 * 255):02X}"
