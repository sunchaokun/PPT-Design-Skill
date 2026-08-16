"""SVGCompiler — compile SVG subset to native editable PPTX shapes.

Public API::

    from ppt_pro_max.renderer.svg_compiler import SVGCompiler, SVGCompileError, SVGResult
    result = SVGCompiler(C=context).compile(svg_text, slide, rect)
"""
from __future__ import annotations

import math
import re
import time
from dataclasses import dataclass, field

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from ppt_pro_max.renderer.boolean_shapes import bool_shape
from ppt_pro_max.renderer.freeform_builder import FreeformBuilder
from ppt_pro_max.renderer.visual_effects import set_solid_fill_with_alpha

from ._affine import Affine, parse_transform
from ._errors import SVGCompileError
from ._paint import GradientDef
from ._paint import _parse_percent_or_float as _parse_pct
from ._paint import apply_gradient as _apply_gradient
from ._paint import resolve_paint as _resolve_paint
from ._path import parse_path, to_beziers
from ._sanitizer import sanitize
from ._text import render_svg_text as _render_svg_text

SVG_NS = "http://www.w3.org/2000/svg"
SVG = f"{{{SVG_NS}}}"

# ─────────────────────────── color helpers ───────────────────────────

_NAMED_COLORS: dict[str, str] = {
    "aliceblue": "F0F8FF", "antiquewhite": "FAEBD7", "aqua": "00FFFF",
    "aquamarine": "7FFFD4", "azure": "F0FFFF", "beige": "F5F5DC",
    "bisque": "FFE4C4", "black": "000000", "blanchedalmond": "FFEBCD",
    "blue": "0000FF", "blueviolet": "8A2BE2", "brown": "A52A2A",
    "burlywood": "DEB887", "cadetblue": "5F9EA0", "chartreuse": "7FFF00",
    "chocolate": "D2691E", "coral": "FF7F50", "cornflowerblue": "6495ED",
    "cornsilk": "FFF8DC", "crimson": "DC143C", "cyan": "00FFFF",
    "darkblue": "00008B", "darkcyan": "008B8B", "darkgoldenrod": "B8860B",
    "darkgray": "A9A9A9", "darkgreen": "006400", "darkgrey": "A9A9A9",
    "darkkhaki": "BDB76B", "darkmagenta": "8B008B", "darkolivegreen": "556B2F",
    "darkorange": "FF8C00", "darkorchid": "9932CC", "darkred": "8B0000",
    "darksalmon": "E9967A", "darkseagreen": "8FBC8F", "darkslateblue": "483D8B",
    "darkslategray": "2F4F4F", "darkslategrey": "2F4F4F",
    "darkturquoise": "00CED1", "darkviolet": "9400D3", "deeppink": "FF1493",
    "deepskyblue": "00BFFF", "dimgray": "696969", "dimgrey": "696969",
    "dodgerblue": "1E90FF", "firebrick": "B22222", "floralwhite": "FFFAF0",
    "forestgreen": "228B22", "fuchsia": "FF00FF", "gainsboro": "DCDCDC",
    "ghostwhite": "F8F8FF", "gold": "FFD700", "goldenrod": "DAA520",
    "gray": "808080", "green": "008000", "greenyellow": "ADFF2F",
    "grey": "808080", "honeydew": "F0FFF0", "hotpink": "FF69B4",
    "indianred": "CD5C5C", "indigo": "4B0082", "ivory": "FFFFF0",
    "khaki": "F0E68C", "lavender": "E6E6FA", "lavenderblush": "FFF0F5",
    "lawngreen": "7CFC00", "lemonchiffon": "FFFACD", "lightblue": "ADD8E6",
    "lightcoral": "F08080", "lightcyan": "E0FFFF", "lightgoldenrodyellow": "FAFAD2",
    "lightgray": "D3D3D3", "lightgreen": "90EE90", "lightgrey": "D3D3D3",
    "lightpink": "FFB6C1", "lightsalmon": "FFA07A", "lightseagreen": "20B2AA",
    "lightskyblue": "87CEFA", "lightslategray": "778899",
    "lightslategrey": "778899", "lightsteelblue": "B0C4DE",
    "lightyellow": "FFFFE0", "lime": "00FF00", "limegreen": "32CD32",
    "linen": "FAF0E6", "magenta": "FF00FF", "maroon": "800000",
    "mediumaquamarine": "66CDAA", "mediumblue": "0000CD",
    "mediumorchid": "BA55D3", "mediumpurple": "9370DB",
    "mediumseagreen": "3CB371", "mediumslateblue": "7B68EE",
    "mediumspringgreen": "00FA9A", "mediumturquoise": "48D1CC",
    "mediumvioletred": "C71585", "midnightblue": "191970",
    "mintcream": "F5FFFA", "mistyrose": "FFE4E1", "moccasin": "FFE4B5",
    "navajowhite": "FFDEAD", "navy": "000080", "oldlace": "FDF5E6",
    "olive": "808000", "olivedrab": "6B8E23", "orange": "FFA500",
    "orangered": "FF4500", "orchid": "DA70D6", "palegoldenrod": "EEE8AA",
    "palegreen": "98FB98", "paleturquoise": "AFEEEE",
    "palevioletred": "DB7093", "papayawhip": "FFEFD5", "peachpuff": "FFDAB9",
    "peru": "CD853F", "pink": "FFC0CB", "plum": "DDA0DD",
    "powderblue": "B0E0E6", "purple": "800080", "rebeccapurple": "663399",
    "red": "FF0000", "rosybrown": "BC8F8F", "royalblue": "4169E1",
    "saddlebrown": "8B4513", "salmon": "FA8072", "sandybrown": "F4A460",
    "seagreen": "2E8B57", "seashell": "FFF5EE", "sienna": "A0522D",
    "silver": "C0C0C0", "skyblue": "87CEEB", "slateblue": "6A5ACD",
    "slategray": "708090", "slategrey": "708090", "snow": "FFFAFA",
    "springgreen": "00FF7F", "steelblue": "4682B4", "tan": "D2B48C",
    "teal": "008080", "thistle": "D8BFD8", "tomato": "FF6347",
    "turquoise": "40E0D0", "violet": "EE82EE", "wheat": "F5DEB3",
    "white": "FFFFFF", "whitesmoke": "F5F5F5", "yellow": "FFFF00",
    "yellowgreen": "9ACD32",
}


def _resolve_svg_color(raw: str | None, C: dict | None, fallback: str) -> str:
    """Resolve an SVG fill/stroke color value against a C (context) dict.

    Priority:
      1. var(--name) → look up C[name], then C[f"palette_{name}"]
      2. known C key
      3. hex
      4. rgb()/rgba()
      5. hsl()/hsla()
      6. named-color
      7. currentColor → C.get("text_dark", "#000000")
    """
    if not raw or raw == "none":
        return None
    v = raw.strip()

    if v.startswith("var("):
        name = v[4:-1].strip().lstrip("-")
        if name in (C or {}):
            return C[name]
        pk = f"palette_{name}"
        if pk in (C or {}):
            return C[pk]
        raise SVGCompileError(f"unresolved color token: var({name})")

    if v in (C or {}):
        return C[v]

    if v == "currentColor":
        return C.get("text_dark", "#000000") if C else "#000000"

    if v.startswith("#"):
        h = v.lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return "#" + h

    lower = v.lower()
    if lower in _NAMED_COLORS:
        return "#" + _NAMED_COLORS[lower]

    # rgb(…) / rgba(…)
    m = _RGB_RE.match(v)
    if m:
        r, g, b, a = int(m.group(1)), int(m.group(2)), int(m.group(3)), m.group(4)
        if a is not None:
            return f"#{r:02X}{g:02X}{b:02X}{int(a*255):02X}"
        return f"#{r:02X}{g:02X}{b:02X}"

    # hsl(…) / hsla(…)
    m = _HSL_RE.match(v)
    if m:
        import colorsys
        h_val, s_val, l_val = float(m.group(1)) / 360.0, float(m.group(2)) / 100.0, float(m.group(3)) / 100.0
        r, g, b = colorsys.hls_to_rgb(h_val, l_val, s_val)
        a = float(m.group(4)) if m.group(4) else 1.0
        if a < 1.0:
            return f"#{int(r*255):02X}{int(g*255):02X}{int(b*255):02X}{int(a*255):02X}"
        return f"#{int(r*255):02X}{int(g*255):02X}{int(b*255):02X}"

    raise SVGCompileError(f"unsupported color value: {v!r}")


_RGB_RE = re.compile(
    r"rgba?\(\s*(\d{1,3})\s*,\s*(\d{1,3})\s*,\s*(\d{1,3})\s*(?:,\s*([\d.]+)\s*)?\)"
)
_HSL_RE = re.compile(
    r"hsla?\(\s*(\d+(?:\.\d+)?)\s*,\s*(\d+(?:\.\d+)?)%\s*,\s*(\d+(?:\.\d+)?)%\s*(?:,\s*([\d.]+)\s*)?\)"
)


# ─────────────────────────── data classes ───────────────────────────

@dataclass
class SVGResult:
    shapes: list = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    features: set = field(default_factory=set)
    compile_ms: float = 0.0
    shape_count: int = 0


# ─────────────��───────────── compiler ───────────────────────────

class SVGCompiler:
    """Compile a subset of SVG to native editable PPTX shapes."""

    def __init__(self, C: dict | None = None) -> None:
        self.C = C or {}

    def compile(
        self,
        svg_text: str,
        slide,
        rect: tuple[float, float, float, float],
        vb: tuple[float, float, float, float] | None = None,
    ) -> SVGResult:
        result = SVGResult()
        t0 = time.perf_counter()

        root = sanitize(svg_text)

        if vb is None:
            vb_el = root.get("viewBox")
            if vb_el:
                vals = [float(v) for v in vb_el.replace(",", " ").split()]
                if len(vals) >= 4:
                    vb = (vals[0], vals[1], vals[2], vals[3])
            if vb is None:
                vb = (0.0, 0.0, 400.0, 300.0)

        self._vb = vb
        self._rect = rect
        self._slide = slide
        self._grads: dict[str, GradientDef] = {}
        self._clips: dict[str, list[list[tuple[float, float]]]] = {}
        self._shape_count = 0
        self._features: set[str] = set()
        self._warnings: list[str] = []
        self._text_rects: list[tuple[float, float, float, float]] = []

        pre_shape_count = len(slide.shapes)

        self._collect_defs(root)
        self._walk(root, Affine(), [])

        result.features = self._features
        result.warnings = self._warnings
        result.shape_count = self._shape_count
        result.compile_ms = (time.perf_counter() - t0) * 1000

        self._detect_text_overlaps(slide, pre_shape_count, result)

        return result

    # ── coordinate transforms ────────────────────────────────

    def _to_inches(self, x: float, y: float) -> tuple[float, float]:
        lx, ly, w, h = self._rect
        vx, vy, vw, vh = self._vb
        s = min(w / vw, h / vh)
        ox = lx + (w - vw * s) / 2.0
        oy = ly + (h - vh * s) / 2.0
        return ox + (x - vx) * s, oy + (y - vy) * s

    # ── defs collection ──────────────────────────────────────

    def _collect_defs(self, root: etree._Element) -> None:
        for g in root.iter(f"{SVG}linearGradient"):
            stops = []
            for s in g.iter(f"{SVG}stop"):
                off = s.get("offset", "0")
                if off.endswith("%"):
                    pos = float(off.rstrip("%")) / 100.0
                else:
                    pos = float(off)
                col = s.get("stop-color", "#000000")
                col = _resolve_svg_color(col, self.C, "#000000")
                op = float(s.get("stop-opacity", "1"))
                stops.append((pos, col, op))
            x1 = _parse_pct(g.get("x1", "0"))
            y1 = _parse_pct(g.get("y1", "0"))
            x2 = _parse_pct(g.get("x2", "1"))
            y2 = _parse_pct(g.get("y2", "0"))
            self._grads[g.get("id")] = GradientDef(stops=stops, x1=x1, y1=y1, x2=x2, y2=y2)

        for g in root.iter(f"{SVG}radialGradient"):
            stops = []
            for s in g.iter(f"{SVG}stop"):
                off = s.get("offset", "0")
                pos = float(off.rstrip("%")) / 100.0 if off.endswith("%") else float(off)
                col = _resolve_svg_color(s.get("stop-color", "#000000"), self.C, "#000000")
                op = float(s.get("stop-opacity", "1"))
                stops.append((pos, col, op))
            cx = _parse_pct(g.get("cx", "50%"))
            cy = _parse_pct(g.get("cy", "50%"))
            r_ = _parse_pct(g.get("r", "50%"))
            self._grads[g.get("id")] = GradientDef(
                stops=stops, cx=cx, cy=cy, r=r_, gradient_type="radial"
            )

        for c in root.iter(f"{SVG}clipPath"):
            polys = []
            for child in c:
                tf = Affine()
                sub = self._svg_polygon(child, tf)
                if sub:
                    polys.extend(sub)
            self._clips[c.get("id")] = polys

    # ── SVG element → polygon points (SVG-space) ─────────────

    def _svg_polygon(
        self, el: etree._Element, tf: Affine
    ) -> list[list[tuple[float, float]]]:
        if not isinstance(el.tag, str):
            return []
        tag = el.tag.split("}")[-1]
        if tag == "rect":
            x = float(el.get("x", 0))
            y = float(el.get("y", 0))
            w = float(el.get("width", 0))
            h = float(el.get("height", 0))
            pts = [(x, y), (x + w, y), (x + w, y + h), (x, y + h)]
        elif tag == "circle":
            cx = float(el.get("cx", 0))
            cy = float(el.get("cy", 0))
            r = float(el.get("r", 0))
            # Use 4 cubic Bezier arcs (each π/2) for circles — edit-point friendly
            pts = self._circle_cubics(cx, cy, r)
        elif tag == "ellipse":
            cx = float(el.get("cx", 0))
            cy = float(el.get("cy", 0))
            rx = float(el.get("rx", 0))
            ry = float(el.get("ry", 0))
            pts = self._ellipse_cubics(cx, cy, rx, ry)
        elif tag in ("polygon", "polyline"):
            flat = [float(v) for v in el.get("points", "").replace(",", " ").split()]
            pts = list(zip(flat[0::2], flat[1::2]))
        elif tag == "line":
            pts = [
                (float(el.get("x1")), float(el.get("y1"))),
                (float(el.get("x2")), float(el.get("y2"))),
            ]
        elif tag == "path":
            cmds, _ = parse_path(el.get("d", ""))
            subs = to_beziers(cmds)
            out: list[list[tuple[float, float]]] = []
            for sub in subs:
                flattened: list[tuple[float, float]] = []
                for seg in sub:
                    flattened.extend(self._flatten_cubic(seg, n=6))
                if flattened:
                    out.append([tf.apply(px, py) for px, py in flattened])
            return out
        else:
            return []
        return [[tf.apply(px, py) for px, py in pts]]

    @staticmethod
    def _circle_cubics(cx: float, cy: float, r: float) -> list[tuple[float, float]]:
        """Return 4 cubic Bezier control points for a circle (16 points total)."""
        pts: list[tuple[float, float]] = []
        for i in range(4):
            a0 = i * math.pi / 2
            a1 = (i + 1) * math.pi / 2
            p0x = cx + r * math.cos(a0)
            p0y = cy + r * math.sin(a0)
            p3x = cx + r * math.cos(a1)
            p3y = cy + r * math.sin(a1)
            alpha = (4 / 3) * math.tan((a1 - a0) / 4) * r
            # handle sign
            if (a1 - a0) > math.pi:
                alpha = -alpha
            c1x = cx + r * (math.cos(a0) - alpha * math.sin(a0))
            c1y = cy + r * (math.sin(a0) + alpha * math.cos(a0))
            c2x = cx + r * (math.cos(a1) + alpha * math.sin(a1))
            c2y = cy + r * (math.sin(a1) - alpha * math.cos(a1))
            pts.extend([(p0x, p0y), (c1x, c1y), (c2x, c2y), (p3x, p3y)])
        return pts

    @staticmethod
    def _ellipse_cubics(cx: float, cy: float, rx: float, ry: float) -> list[tuple[float, float]]:
        pts: list[tuple[float, float]] = []
        for i in range(4):
            a0 = i * math.pi / 2
            a1 = (i + 1) * math.pi / 2
            p0x = cx + rx * math.cos(a0)
            p0y = cy + ry * math.sin(a0)
            p3x = cx + rx * math.cos(a1)
            p3y = cy + ry * math.sin(a1)
            k = (4 / 3) * (math.sqrt(2) - 1)
            c1x = p0x + k * rx * math.sin(a0)
            c1y = p0y - k * ry * math.cos(a0)
            c2x = p3x - k * rx * math.sin(a1)
            c2y = p3y + k * ry * math.cos(a1)
            pts.extend([(p0x, p0y), (c1x, c1y), (c2x, c2y), (p3x, p3y)])
        return pts

    @staticmethod
    def _flatten_cubic(
        seg: tuple[tuple[float, float], tuple[float, float], tuple[float, float], tuple[float, float]],
        n: int = 6,
    ) -> list[tuple[float, float]]:
        (x0, y0), (x1, y1), (x2, y2), (x3, y3) = seg
        out: list[tuple[float, float]] = []
        for i in range(n + 1):
            t = i / n
            mt = 1.0 - t
            px = (
                mt ** 3 * x0
                + 3 * mt * mt * t * x1
                + 3 * mt * t * t * x2
                + t ** 3 * x3
            )
            py = (
                mt ** 3 * y0
                + 3 * mt * mt * t * y1
                + 3 * mt * t * t * y2
                + t ** 3 * y3
            )
            out.append((px, py))
        return out

    # ── paint resolution ─────────────────────────────────────

    def _paint(
        self, el: etree._Element, which: str
    ) -> tuple[str, object | None, int]:
        return _resolve_paint(el, which, self._grads, self.C, _resolve_svg_color, self._features)

    # ── rendering ────────────────────────────────────────────

    def _walk(self, el: etree._Element, tf: Affine, clip_stack: list) -> None:
        # Skip non-element nodes (e.g., lxml Comment, ProcessingInstruction)
        if not isinstance(el.tag, str):
            return
        tag = el.tag.split("}")[-1] if el.tag else ""
        tf = tf.compose(parse_transform(el.get("transform")))

        c = el.get("clip-path")
        if c and c.startswith("url(#"):
            gid = c[c.index("#") + 1 : -1]
            polys = self._clips.get(gid, [])
            clip_stack = clip_stack + polys

        if tag == "g" or tag == "svg":
            for child in el:
                self._walk(child, tf, clip_stack)
        elif tag == "defs":
            pass  # already collected
        elif tag in ("image", "filter", "mask"):
            self._features.add(tag)
            self._warnings.append(f"unsupported SVG feature: <{tag}> element (skipped)")
        elif tag == "text":
            self._render_text(el, tf)
        elif tag in (
            "rect", "circle", "ellipse", "polygon", "polyline", "line", "path",
        ):
            self._render_shape(el, tag, tf, clip_stack)

    def _render_shape(
        self, el: etree._Element, tag: str, tf: Affine, clip_stack: list
    ) -> None:
        self._features.add(tag)

        fkind, fval, fa = self._paint(el, "fill")
        skind, sval, _ = self._paint(el, "stroke")
        sw = float(el.get("stroke-width", "1"))

        subpaths = self._svg_polygon(el, tf)
        if not subpaths:
            return

        # collect all points for bbox
        all_pts: list[tuple[float, float]] = []
        for sub in subpaths:
            all_pts.extend(sub)
        if not all_pts:
            return

        minx = min(p[0] for p in all_pts)
        maxx = max(p[0] for p in all_pts)
        miny = min(p[1] for p in all_pts)
        maxy = max(p[1] for p in all_pts)

        # map bounding box to slide inches
        corners_in = [self._to_inches(minx, miny), self._to_inches(maxx, maxy)]
        ix0, iy0 = corners_in[0]
        ix1, iy1 = corners_in[1]
        iw = ix1 - ix0
        ih = iy1 - iy0

        # rebuild subpaths in slide inches, local coords
        local: list[list[tuple[float, float]]] = []
        for sub in subpaths:
            in_pts = [self._to_inches(px, py) for px, py in sub]
            local.append([(px - ix0, py - iy0) for px, py in in_pts])

        # clip polygons in slide inches, local coords (relative to shape bbox)
        clip_local: list[list[tuple[float, float]]] = []
        for clip_subpath in clip_stack:
            cpts = [self._to_inches(px, py) for px, py in clip_subpath]
            clip_local.append([(px - ix0, py - iy0) for px, py in cpts])

        needs_bool = bool(clip_stack) or el.get("fill-rule") == "evenodd"

        if needs_bool:
            self._render_bool(
                local, clip_local, ix0, iy0, iw, ih,
                fkind, fval, fa, skind, sval,
            )
            return

        # Fast path: rect with solid fill and no stroke → native shape
        if (
            tag == "rect"
            and not el.get("rx")
            and fkind == "solid"
            and skind == "none"
            and not clip_stack
        ):
            self._add_native(ix0, iy0, iw, ih, fval, fa)
            return

        # Freeform path
        fill_hex = _resolve_svg_color(fval, self.C, "") if fkind == "solid" else None
        stroke_hex = _resolve_svg_color(sval, self.C, "") if skind == "solid" else None
        elem = self._add_freeform(
            local, ix0, iy0, iw, ih, fill_hex or "#FFFFFF", fa, stroke_hex, sw
        )
        if fkind == "grad":
            self._apply_gradient_to_elem(elem, fval)

    def _render_bool(
        self,
        local_subpaths: list[list[tuple[float, float]]],
        clip_local: list[list[tuple[float, float]]],
        ix0: float, iy0: float, iw: float, ih: float,
        fkind: str, fval, fa: int, skind: str, sval,
    ) -> None:
        """Boolean-compute shapes and render, intersecting with clip regions."""
        try:
            from shapely.errors import GEOSException, TopologicalError
            from shapely.geometry import Polygon as ShapelyPoly
            from shapely.validation import make_valid
        except ImportError:
            raise SVGCompileError("shapely is required for boolean operations")

        poly = None
        for sub in local_subpaths:
            pts = [(p[0] + ix0, p[1] + iy0) for p in sub]
            if len(pts) < 3:
                continue
            try:
                p = ShapelyPoly(pts).buffer(0)
                p = make_valid(p)
            except (TopologicalError, ValueError, GEOSException):
                p = None
            if p is not None and not p.is_empty:
                poly = p if poly is None else poly.union(p)

        if poly is None or poly.is_empty:
            return

        # Apply clip regions: intersect shape with clip in slide inches
        for clip_poly in clip_local:
            cpts = [(p[0] + ix0, p[1] + iy0) for p in clip_poly]
            if len(cpts) < 3:
                continue
            try:
                cpoly = ShapelyPoly(cpts).buffer(0)
                cpoly = make_valid(cpoly)
            except (TopologicalError, ValueError, GEOSException):
                cpoly = None
            if cpoly is not None and not cpoly.is_empty:
                poly = poly.intersection(cpoly)

        if poly.is_empty:
            return

        stroke_hex = _resolve_svg_color(sval, self.C, "") if skind == "solid" else None

        if fkind == "grad":
            elem = bool_shape(
                poly, self._slide, ix0, iy0, iw, ih, fill="#FFFFFF", line=stroke_hex, alpha=fa
            )
            if elem is not None:
                self._apply_gradient_to_elem(elem, fval)
        else:
            fill_hex = _resolve_svg_color(fval, self.C, "") if fkind == "solid" else "#FFFFFF"
            elem = bool_shape(
                poly, self._slide, ix0, iy0, iw, ih, fill=fill_hex, line=stroke_hex, alpha=fa
            )
        if elem is not None:
            self._shape_count += 1

    def _add_native(
        self, x: float, y: float, w: float, h: float, fill: str | None, alpha: int
    ) -> None:
        sh = self._slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor.from_string(fill.lstrip("#"))
            if alpha < 100:
                set_solid_fill_with_alpha(sh, fill, alpha)
        else:
            sh.fill.background()
        sh.line.fill.background()
        self._shape_count += 1

    def _add_freeform(
        self,
        local_subs: list[list[tuple[float, float]]],
        ix0: float, iy0: float, iw: float, ih: float,
        fill: str | None,
        alpha: int,
        stroke: str | None,
        sw: float,
    ) -> object:
        fb = FreeformBuilder()
        for sub in local_subs:
            if sub:
                fb.move_to(sub[0][0], sub[0][1])
                for pt in sub[1:]:
                    fb.line_to(pt[0], pt[1])
                fb.close()
        elem = fb.build(
            self._slide, ix0, iy0, iw, ih,
            fill_color=fill, line_color=stroke,
            line_width_pt=max(sw, 0.5),
            no_fill=(fill is None),
        )
        if alpha < 100 and fill:
            self._apply_alpha_to_elem(elem, fill, alpha)
        self._shape_count += 1
        return elem

    def _apply_gradient_to_elem(
        self, elem: object, grad: GradientDef
    ) -> None:
        _apply_gradient(elem, grad, self._wrap_elem)

    def _apply_alpha_to_elem(self, elem: object, fill: str, alpha: int) -> None:
        if alpha >= 100:
            return
        wrapper = self._wrap_elem(elem)
        set_solid_fill_with_alpha(wrapper, fill, alpha)

    @staticmethod
    def _wrap_elem(elem):
        """Wrap a raw lxml element so visual_effects can access shape._element."""
        if hasattr(elem, "_element"):
            return elem
        class _ShapeProxy:
            def __init__(self, el):
                self._element = el
        return _ShapeProxy(elem)

    # ── text ─────────────────────────────────────────────────

    def _render_text(self, el: etree._Element, tf: Affine) -> None:
        # scale for font-size: svg region size / viewBox size
        lx, ly, rw, rh = self._rect
        svg_w, svg_h = self._vb[2], self._vb[3]
        _render_svg_text(
            el=el,
            tf=tf,
            to_inches_fn=self._to_inches,
            slide=self._slide,
            C=self.C,
            resolve_color_fn=_resolve_svg_color,
            features=self._features,
            svg_w=svg_w,
            svg_h=svg_h,
            slide_w=rw,
            slide_h=rh,
            rect_x=lx,
            rect_y=ly,
        )
        self._shape_count += 1

    # ── collision detection ──────────────────────────────────

    @staticmethod
    def _detect_text_overlaps(slide, pre_count: int, result: SVGResult) -> None:
        """Detect overlapping text boxes among shapes added during this compile.

        Walks all text boxes added since pre_count, collects their bounding
        rects in slide inches, and emits a warning for each pair that overlaps
        by more than a small tolerance (avoiding false positives from adjacent
        text boxes that touch at edges).
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        text_rects: list[tuple[float, float, float, float]] = []
        for shape in slide.shapes:
            if slide.shapes.index(shape) < pre_count:
                continue
            if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            try:
                x = shape.left / 914400.0
                y = shape.top / 914400.0
                w = shape.width / 914400.0
                h = shape.height / 914400.0
            except (AttributeError, TypeError):
                continue
            text_rects.append((x, y, x + w, y + h))

        n = len(text_rects)
        if n < 2:
            return

        tol = 0.05
        for i in range(n):
            for j in range(i + 1, n):
                ax0, ay0, ax1, ay1 = text_rects[i]
                bx0, by0, bx1, by1 = text_rects[j]
                dx = min(ax1, bx1) - max(ax0, bx0)
                dy = min(ay1, by1) - max(ay0, by0)
                if dx > tol and dy > tol:
                    result.warnings.append(
                        f"text box overlap detected: "
                        f"rect1=({ax0:.2f},{ay0:.2f},{ax1:.2f},{ay1:.2f}) "
                        f"rect2=({bx0:.2f},{by0:.2f},{bx1:.2f},{by1:.2f}) "
                        f"overlap=({dx:.2f}x{dy:.2f})"
                    )
