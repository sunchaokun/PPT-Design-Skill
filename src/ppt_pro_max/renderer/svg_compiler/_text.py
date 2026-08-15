"""SVG text rendering with baseline-aware positioning and <tspan> support.

Extracts text rendering from _compiler.py and adds:
- Pillow-based text measurement (with graceful fallback to estimate_text_size)
- SVG baseline (dominant-baseline, alignment-baseline) → PPT vertical anchor mapping
- Multi-line <tspan> support — each tspan with x/y becomes a new paragraph;
  tspans without x/y are inline runs within the same paragraph
- font-family → PPT font mapping
"""
from __future__ import annotations

import re
from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ppt_pro_max.renderer.text_measurer import estimate_text_size

from ._affine import Affine

_FONT_MAP: dict[str, str] = {
    "arial": "Arial",
    "helvetica": "Arial",
    "sans-serif": "Arial",
    "times new roman": "Times New Roman",
    "times": "Times New Roman",
    "serif": "Times New Roman",
    "georgia": "Georgia",
    "courier": "Courier New",
    "courier new": "Courier New",
    "monospace": "Courier New",
    "consolas": "Consolas",
    "verdana": "Verdana",
    "tahoma": "Tahoma",
    "calibri": "Calibri",
    "impact": "Impact",
    "comic sans ms": "Comic Sans MS",
}

_BASELINE_MAP: dict[str, MSO_ANCHOR] = {
    "auto": MSO_ANCHOR.MIDDLE,
    "alphabetic": MSO_ANCHOR.BOTTOM,
    "text-after-edge": MSO_ANCHOR.BOTTOM,
    "text-before-edge": MSO_ANCHOR.TOP,
    "central": MSO_ANCHOR.MIDDLE,
    "middle": MSO_ANCHOR.MIDDLE,
    "hanging": MSO_ANCHOR.TOP,
    "mathematical": MSO_ANCHOR.MIDDLE,
}

_ANCHOR_MAP: dict[str, PP_ALIGN] = {
    "middle": PP_ALIGN.CENTER,
    "end": PP_ALIGN.RIGHT,
    "start": PP_ALIGN.LEFT,
}


@dataclass
class TextMetrics:
    width_inches: float
    height_inches: float
    ascent_ratio: float = 0.8
    descent_ratio: float = 0.2


@dataclass
class _SpanSpec:
    text: str = ""
    x: float | None = None
    y: float | None = None
    dx: float = 0.0
    dy: float = 0.0
    font_size: float | None = None
    font_family: str | None = None
    fill: str | None = None
    bold: bool = False
    italic: bool = False
    is_new_line: bool = False


def _resolve_font_family(raw: str | None) -> str:
    if not raw:
        return "Calibri"
    families = re.split(r"[,\s]+", raw.strip("'\""))
    for f in families:
        key = f.strip().lower()
        if key in _FONT_MAP:
            return _FONT_MAP[key]
    return families[0].strip().strip("'\"")


def _resolve_baseline(el) -> MSO_ANCHOR:
    db = el.get("dominant-baseline", el.get("alignment-baseline", "auto"))
    return _BASELINE_MAP.get(db.lower(), MSO_ANCHOR.MIDDLE)


def _measure_text(
    content: str, font_size_pt: float, font_family: str, max_width_inches: float
) -> TextMetrics:
    try:
        from PIL import ImageFont

        try:
            font = ImageFont.truetype(font_family + ".ttf", int(font_size_pt))
        except OSError:
            font = ImageFont.load_default()

        bbox = font.getbbox(content)
        w_px = bbox[2] - bbox[0]
        h_px = bbox[3] - bbox[1]
        ascent = font.getmetrics()[0]
        total_h = font.getmetrics()[0] + font.getmetrics()[1]

        px_per_inch = 96.0
        w_in = w_px / px_per_inch
        h_in = h_px / px_per_inch

        return TextMetrics(
            width_inches=max(w_in, 0.5),
            height_inches=max(h_in, 0.3),
            ascent_ratio=ascent / total_h if total_h > 0 else 0.8,
            descent_ratio=1.0 - (ascent / total_h if total_h > 0 else 0.8),
        )
    except ImportError:
        _, h_est = estimate_text_size(content, max(8, int(font_size_pt)), max_width_inches)
        return TextMetrics(
            width_inches=max_width_inches,
            height_inches=h_est,
            ascent_ratio=0.8,
            descent_ratio=0.2,
        )


def _resolve_fill(raw: str | None, C: dict, resolve_color_fn) -> str:
    if not raw or raw == "none":
        return "#000000"
    if raw.startswith("url(#"):
        return "#000000"
    resolved = resolve_color_fn(raw, C, "")
    return resolved if resolved else "#000000"


def _parse_font_size(raw: str | None, parent_fs: float) -> float:
    if not raw:
        return parent_fs
    clean = re.sub(r"[^\d.]", "", raw)
    return float(clean) if clean else parent_fs


def _parse_font_weight(raw: str | None) -> bool:
    if not raw:
        return False
    return raw not in ("normal", "100", "200", "300")


def _collect_spans(el, parent_fs: float, parent_ff: str, parent_fill: str,
                   C: dict, resolve_color_fn) -> list[_SpanSpec]:
    spans: list[_SpanSpec] = []

    direct_text = el.text
    if direct_text and direct_text.strip():
        spans.append(_SpanSpec(
            text=direct_text.strip(),
            font_size=parent_fs,
            font_family=parent_ff,
            fill=parent_fill,
        ))

    for child in el:
        tag = child.tag.split("}")[-1] if child.tag else ""
        if tag != "tspan":
            tail = child.tail
            if tail and tail.strip():
                spans.append(_SpanSpec(
                    text=tail.strip(),
                    font_size=parent_fs,
                    font_family=parent_ff,
                    fill=parent_fill,
                ))
            continue

        fs = _parse_font_size(child.get("font-size"), parent_fs)
        ff = _resolve_font_family(child.get("font-family")) if child.get("font-family") else parent_ff
        fill = _resolve_fill(child.get("fill", parent_fill), C, resolve_color_fn)
        bold = _parse_font_weight(child.get("font-weight"))
        italic = child.get("font-style") == "italic"

        has_x = child.get("x") is not None
        has_y = child.get("y") is not None

        span = _SpanSpec(
            text=(child.text or "").strip(),
            x=float(child.get("x")) if has_x else None,
            y=float(child.get("y")) if has_y else None,
            dx=float(child.get("dx", "0")),
            dy=float(child.get("dy", "0")),
            font_size=fs,
            font_family=ff,
            fill=fill,
            bold=bold,
            italic=italic,
            is_new_line=has_x or has_y,
        )
        spans.append(span)

        tail = child.tail
        if tail and tail.strip():
            spans.append(_SpanSpec(
                text=tail.strip(),
                font_size=parent_fs,
                font_family=parent_ff,
                fill=parent_fill,
            ))

    return spans


def _group_spans_into_lines(spans: list[_SpanSpec]) -> list[list[_SpanSpec]]:
    lines: list[list[_SpanSpec]] = []
    current: list[_SpanSpec] = []
    for sp in spans:
        if sp.is_new_line and current:
            lines.append(current)
            current = []
        current.append(sp)
    if current:
        lines.append(current)
    return lines if lines else [[]]


def render_svg_text(
    el,
    tf: Affine,
    to_inches_fn,
    slide,
    C: dict,
    resolve_color_fn,
    features: set,
) -> None:
    features.add("text")

    x = float(el.get("x", 0))
    y = float(el.get("y", 0))
    ix, iy = to_inches_fn(*tf.apply(x, y))

    parent_fs = _parse_font_size(el.get("font-size"), 14.0)
    parent_ff = _resolve_font_family(el.get("font-family"))
    parent_fill = _resolve_fill(el.get("fill"), C, resolve_color_fn)
    anchor = el.get("text-anchor", "start")
    v_anchor = _resolve_baseline(el)

    spans = _collect_spans(el, parent_fs, parent_ff, parent_fill, C, resolve_color_fn)

    has_tspan_children = any(
        c.tag.split("}")[-1] == "tspan" for c in el
    )

    if not has_tspan_children:
        content = "".join(el.itertext()).strip()
        if not content:
            return
        _render_simple_text(
            content, ix, iy, parent_fs, parent_ff, parent_fill,
            anchor, v_anchor, el, slide, C, resolve_color_fn,
        )
        return

    if not any(s.text for s in spans):
        return

    _render_tspan_text(
        spans, ix, iy, parent_fs, parent_ff, parent_fill,
        anchor, v_anchor, el, slide, to_inches_fn, tf,
        C, resolve_color_fn,
    )


def _render_simple_text(
    content: str, ix: float, iy: float,
    fs: float, ff: str, fill: str,
    anchor: str, v_anchor, el, slide,
    C: dict, resolve_color_fn,
) -> None:
    metrics = _measure_text(content, fs, ff, 8.0)

    if anchor == "middle":
        left = ix - metrics.width_inches / 2
    elif anchor == "end":
        left = ix - metrics.width_inches
    else:
        left = ix - 0.1

    if v_anchor == MSO_ANCHOR.TOP:
        top = iy
    elif v_anchor == MSO_ANCHOR.BOTTOM:
        top = iy - metrics.height_inches
    else:
        top = iy - metrics.height_inches * metrics.ascent_ratio

    width = metrics.width_inches + 0.2
    height = metrics.height_inches * 1.5

    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf_el = tb.text_frame
    tf_el.word_wrap = False
    tf_el.vertical_anchor = v_anchor
    p = tf_el.paragraphs[0]
    p.alignment = _ANCHOR_MAP.get(anchor, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = content
    run.font.size = Pt(fs)
    run.font.color.rgb = RGBColor.from_string(fill.lstrip("#"))
    run.font.name = ff

    bold = el.get("font-weight")
    if bold and bold not in ("normal", "100", "200", "300"):
        run.font.bold = True

    italic = el.get("font-style")
    if italic == "italic":
        run.font.italic = True


def _render_tspan_text(
    spans: list[_SpanSpec], ix: float, iy: float,
    parent_fs: float, parent_ff: str, parent_fill: str,
    anchor: str, v_anchor, el, slide,
    to_inches_fn, tf: Affine,
    C: dict, resolve_color_fn,
) -> None:
    lines = _group_spans_into_lines(spans)

    all_text = " ".join(s.text for line in lines for s in line if s.text)
    metrics = _measure_text(all_text, parent_fs, parent_ff, 8.0)
    line_h = metrics.height_inches * 1.3

    if anchor == "middle":
        left = ix - metrics.width_inches / 2
    elif anchor == "end":
        left = ix - metrics.width_inches
    else:
        left = ix - 0.1

    if v_anchor == MSO_ANCHOR.TOP:
        top = iy
    elif v_anchor == MSO_ANCHOR.BOTTOM:
        top = iy - line_h * len(lines)
    else:
        top = iy - metrics.height_inches * metrics.ascent_ratio

    width = metrics.width_inches + 0.4
    height = line_h * len(lines) + 0.2

    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf_el = tb.text_frame
    tf_el.word_wrap = True
    tf_el.vertical_anchor = v_anchor

    first_para = True
    for line in lines:
        if first_para:
            p = tf_el.paragraphs[0]
            first_para = False
        else:
            p = tf_el.add_paragraph()

        p.alignment = _ANCHOR_MAP.get(anchor, PP_ALIGN.LEFT)
        p.space_after = Pt(0)

        for sp in line:
            if not sp.text:
                continue
            run = p.add_run()
            run.text = sp.text
            run.font.size = Pt(sp.font_size or parent_fs)
            run.font.color.rgb = RGBColor.from_string((sp.fill or parent_fill).lstrip("#"))
            run.font.name = sp.font_family or parent_ff
            if sp.bold:
                run.font.bold = True
            if sp.italic:
                run.font.italic = True
