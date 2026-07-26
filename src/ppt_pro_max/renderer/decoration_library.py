"""Decoration Library — pre-built decorative elements for various styles.

Provides high-level helpers for common decorative patterns:
  - Brush stroke dividers (ink-wash style)
  - Seal stamps (Chinese traditional)
  - Scroll frames (ink-wash / zen)
  - Neon borders (cyberpunk)
  - Grid backgrounds (scientific)
  - Glass panels (semi-transparent overlay)
  - Ink splashes (ink-wash decoration)
"""
from __future__ import annotations

import math

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from ppt_pro_max.renderer.freeform_builder import FreeformBuilder
from ppt_pro_max.renderer.visual_effects import apply_glow, apply_frosted_glass


def add_brush_divider(
    slide, x: float, y: float, w: float,
    color: str = "#2C2C2C", thickness: float = 0.08,
) -> object:
    builder = FreeformBuilder()
    ht = thickness / 2
    builder.move_to(0, ht * 0.3)
    builder.cubic_bezier_to(
        w * 0.03, ht * 0.1,
        w * 0.06, ht * 0.8,
        w * 0.1, ht * 0.5,
    )
    builder.cubic_bezier_to(
        w * 0.15, ht * 0.2,
        w * 0.18, ht * 1.6,
        w * 0.25, ht,
    )
    builder.cubic_bezier_to(
        w * 0.32, ht * 0.4,
        w * 0.38, ht * 1.5,
        w * 0.45, ht * 0.8,
    )
    builder.cubic_bezier_to(
        w * 0.52, ht * 0.1,
        w * 0.58, ht * 1.4,
        w * 0.65, ht,
    )
    builder.cubic_bezier_to(
        w * 0.72, ht * 0.6,
        w * 0.78, ht * 1.3,
        w * 0.85, ht * 0.7,
    )
    builder.cubic_bezier_to(
        w * 0.9, ht * 0.3,
        w * 0.95, ht * 1.2,
        w, ht * 0.5,
    )
    builder.line_to(w * 0.98, ht + thickness * 0.4)
    builder.cubic_bezier_to(
        w * 0.93, ht + thickness * 0.6,
        w * 0.88, ht + thickness * 0.1,
        w * 0.82, ht + thickness * 0.3,
    )
    builder.cubic_bezier_to(
        w * 0.75, ht + thickness * 0.5,
        w * 0.68, ht + thickness * 0.1,
        w * 0.6, ht + thickness * 0.3,
    )
    builder.cubic_bezier_to(
        w * 0.52, ht + thickness * 0.5,
        w * 0.45, ht + thickness * 0.1,
        w * 0.38, ht + thickness * 0.3,
    )
    builder.cubic_bezier_to(
        w * 0.3, ht + thickness * 0.5,
        w * 0.22, ht + thickness * 0.1,
        w * 0.15, ht + thickness * 0.3,
    )
    builder.cubic_bezier_to(
        w * 0.1, ht + thickness * 0.5,
        w * 0.05, ht + thickness * 0.2,
        0, ht + thickness * 0.5,
    )
    builder.close()
    return builder.build(
        slide, x, y - thickness, w, thickness * 2,
        fill_color=color, no_fill=False,
        line_color=None, line_width_pt=0,
    )


def add_seal_stamp(
    slide, x: float, y: float, size: float, text: str,
    fill_hex: str = "#C41E3A", font_name: str = "STZhongsong",
    rotation: float = -15, style: str = "zhu",
    border_width_pt: float = 4.0,
) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(size), Inches(size),
    )
    if style == "zhu":
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string(fill_hex.lstrip("#"))
        sh.line.color.rgb = RGBColor.from_string(fill_hex.lstrip("#"))
        sh.line.width = Pt(border_width_pt)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(max(14, int(size * 22)))
        run.font.color.rgb = RGBColor.from_string("FFFFFF")
        run.font.bold = True
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string("F5F0E8")
        sh.line.color.rgb = RGBColor.from_string(fill_hex.lstrip("#"))
        sh.line.width = Pt(border_width_pt)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(max(14, int(size * 22)))
        run.font.color.rgb = RGBColor.from_string(fill_hex.lstrip("#"))
        run.font.bold = True
    if rotation != 0:
        spPr = sh._element.find(qn("p:spPr"))
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is not None:
            xfrm.set("rot", str(int(rotation * 60000)))
    return sh._element


def add_scroll_frame(
    slide, x: float, y: float, w: float, h: float,
    style: str = "xuan",
) -> object:
    style_colors = {
        "xuan": {"border": "#D4C5A0", "fill": "#F5F0E8"},
        "silk": {"border": "#8B7355", "fill": "#FAF5EF"},
    }
    colors = style_colors.get(style, style_colors["xuan"])
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(colors["fill"].lstrip("#"))
    sh.line.color.rgb = RGBColor.from_string(colors["border"].lstrip("#"))
    sh.line.width = Pt(1.5)
    return sh._element


def add_neon_border(
    slide, x: float, y: float, w: float, h: float,
    color: str = "#8B5CF6", radius: float = 0.1,
) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    spPr = sh._element.find(qn("p:spPr"))
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    etree.SubElement(spPr, qn("a:noFill"))
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(2.0 * 12700)))
    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr")).set("val", color.lstrip("#"))
    apply_glow(sh, radius_pt=6, color=color, alpha_pct=60)
    return sh._element


def add_grid_background(
    slide, spacing: float = 1.0, color: str = "#E0E0E0", alpha: int = 15,
) -> list:
    sw = 13.333
    sh = 7.5
    elems = []
    x = spacing
    while x < sw:
        line_sh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(0), Inches(0.005), Inches(sh),
        )
        spPr = line_sh._element.find(qn("p:spPr"))
        for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
            el = spPr.find(qn(tag))
            if el is not None:
                spPr.remove(el)
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", color.lstrip("#"))
        alpha_el = etree.SubElement(srgb, qn("a:alpha"))
        alpha_el.set("val", str(alpha * 1000))
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            spPr.remove(ln)
        etree.SubElement(spPr, qn("a:ln"))
        elems.append(line_sh._element)
        x += spacing

    y = spacing
    while y < sh:
        line_sh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(y), Inches(sw), Inches(0.005),
        )
        spPr = line_sh._element.find(qn("p:spPr"))
        for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
            el = spPr.find(qn(tag))
            if el is not None:
                spPr.remove(el)
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", color.lstrip("#"))
        alpha_el = etree.SubElement(srgb, qn("a:alpha"))
        alpha_el.set("val", str(alpha * 1000))
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            spPr.remove(ln)
        etree.SubElement(spPr, qn("a:ln"))
        elems.append(line_sh._element)
        y += spacing

    return elems


def add_glass_panel(
    slide, x: float, y: float, w: float, h: float,
    tint: str = "#FFFFFF", alpha: int = 15, soft_edge: float = 8,
) -> object:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    apply_frosted_glass(sh, tint_color=tint, tint_alpha=alpha,
                        soft_edge=soft_edge)
    return sh._element


def add_ink_splash(
    slide, x: float, y: float, size: float,
    color: str = "#2C2C2C", alpha: int = 100,
) -> object:
    import random
    rng = random.Random(42)
    builder = FreeformBuilder()
    r = size / 2
    points = 16
    for i in range(points):
        angle = 2 * math.pi * i / points
        variance = 0.5 + 0.8 * rng.random()
        if i % 4 == 0:
            variance *= 1.4
        px = r + r * variance * math.cos(angle)
        py = r + r * variance * math.sin(angle)
        if i == 0:
            builder.move_to(px, py)
        else:
            prev_angle = 2 * math.pi * (i - 0.5) / points
            cp_var = 0.6 + 0.5 * rng.random()
            cp1x = r + r * cp_var * math.cos(prev_angle)
            cp1y = r + r * cp_var * math.sin(prev_angle)
            cp2x = px + (px - cp1x) * 0.15
            cp2y = py + (py - cp1y) * 0.15
            builder.cubic_bezier_to(cp1x, cp1y, cp2x, cp2y, px, py)
    builder.close()

    builder.new_path()
    for _ in range(3):
        dx = r * (0.3 + 0.7 * rng.random())
        dy = r * (0.3 + 0.7 * rng.random())
        dr = r * (0.05 + 0.1 * rng.random())
        builder.move_to(dx + dr, dy)
        for j in range(6):
            a = 2 * math.pi * j / 6
            sx = dx + dr * (0.7 + 0.3 * rng.random()) * math.cos(a)
            sy = dy + dr * (0.7 + 0.3 * rng.random()) * math.sin(a)
            builder.line_to(sx, sy)
        builder.close()

    elem = builder.build(
        slide, x, y, size, size,
        fill_color=color, no_fill=False,
        line_color=None, line_width_pt=0,
    )

    if alpha < 100:
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            srgb = solidFill.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha_el = etree.SubElement(srgb, qn("a:alpha"))
                alpha_el.set("val", str(alpha * 1000))
    return elem


def _lighten(hex_color: str, amount: int) -> str:
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    r = min(255, r + amount)
    g = min(255, g + amount)
    b = min(255, b + amount)
    return f"{r:02X}{g:02X}{b:02X}"
