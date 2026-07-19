"""PrecisionRenderer — brand-compliant + pixel-perfect rendering.

Combines Pipeline's brand constraint system (BrandSpec, template layouts,
logo/footer/watermark) with Build Script's rendering precision (run-level
fonts, Pillow cropping, per-element coordinates).

This is the third rendering path:
  - EnterpriseRenderer: template-driven, paragraph-level fonts → low quality
  - PPTRenderer (freestyle fallback): layout-registry driven → medium quality
  - PrecisionRenderer: brand-aware + run-level precision → delivery quality
"""

from __future__ import annotations

import hashlib
import os
import tempfile
from typing import Any

from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.renderer.layout_registry import SLIDE_WIDTH, SLIDE_HEIGHT
from ppt_pro_max.renderer.visual_effects import (
    apply_gradient, apply_shadow, apply_glow, GradientFill, GradientStop,
)
from ppt_pro_max.renderer.shape_factory import ShapeFactory

CORNER_RADIUS_SCALE = {
    "sm": 4,
    "md": 8,
    "lg": 12,
    "pill": 50,
}


class PrecisionRenderer:

    def __init__(self, brand_spec: BrandSpec | None = None, template_path: str | None = None):
        self._brand = brand_spec
        self._template_path = template_path
        self._has_template = bool(template_path and os.path.exists(template_path))
        self._crop_cache_dir = os.path.join(tempfile.gettempdir(), "ppt-precision-crops")
        os.makedirs(self._crop_cache_dir, exist_ok=True)
        self._shape_factory: ShapeFactory | None = None
        from ppt_pro_max.renderer.layout_engine import LayoutEngine
        from ppt_pro_max.renderer.decoration_renderer import DecorationRenderer
        self._layout_engine = LayoutEngine()
        self._decoration_renderer = DecorationRenderer()

    @property
    def brand(self) -> BrandSpec:
        return self._brand or BrandSpec()

    @property
    def shape_factory(self) -> ShapeFactory:
        if self._shape_factory is None:
            brand_colors = {}
            if self._brand and self._brand.colors:
                brand_colors = dict(self._brand.colors)
            self._shape_factory = ShapeFactory(brand_colors=brand_colors)
        return self._shape_factory

    def _c(self, role: str, fallback: str = "#000000") -> str:
        if self._brand and self._brand.colors:
            return self._brand.colors.get(role, fallback)
        return fallback

    def _font_h(self) -> str:
        if self._brand and self._brand.fonts:
            return self._brand.fonts.get("heading", "Inter")
        return "Inter"

    def _font_b(self) -> str:
        if self._brand and self._brand.fonts:
            return self._brand.fonts.get("body", "Inter")
        return "Inter"

    def _is_dark(self) -> bool:
        if self._brand and self._brand.dark_mode:
            return True
        bg_hex = self._c("background", "#FFFFFF").lstrip("#")
        r, g, b = int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16)
        luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
        return luminance < 0.5

    @staticmethod
    def _rgb(h: str) -> RGBColor:
        return RGBColor.from_string(h.lstrip("#"))

    def create_presentation(self) -> Presentation:
        if self._has_template:
            try:
                from ppt_pro_max.enterprise.slide_utils import remove_slide
                prs = Presentation(self._template_path)
                while len(prs.slides) > 0:
                    remove_slide(prs, 0)
                return prs
            except Exception:
                pass
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        return prs

    def add_slide(self, prs: Presentation, layout_name: str | None = None):
        if self._has_template:
            if layout_name:
                for layout in prs.slide_layouts:
                    if layout.name == layout_name:
                        return prs.slides.add_slide(layout)
            for layout in prs.slide_layouts:
                if "blank" in layout.name.lower():
                    return prs.slides.add_slide(layout)
            return prs.slides.add_slide(prs.slide_layouts[0])
        blank = None
        for layout in prs.slide_layouts:
            if "blank" in layout.name.lower():
                blank = layout
                break
        return prs.slides.add_slide(blank or prs.slide_layouts[-1])

    # ── Text helpers (run-level fonts) ──

    def add_text(self, slide, text: str, x: float, y: float, w: float, h: float,
                 font: str | None = None, size: int = 20, color_role: str = "foreground",
                 color_hex: str | None = None,
                 bold: bool = False, align: str = "left") -> object:
        font = font or self._font_h()
        color = color_hex or self._c(color_role)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = self._rgb(color)
        run.font.bold = bold
        from ppt_pro_max.renderer.theme_mapper import get_cjk_companion
        cjk = get_cjk_companion(font, "heading" if size >= 20 else "body")
        self._set_font_with_cjk(run, font, cjk)
        return tb

    def add_multiline(self, slide, lines: list[str], x: float, y: float, w: float, h: float,
                      font: str | None = None, size: int = 14, color_role: str = "foreground",
                      color_hex: str | None = None,
                      bold: bool = False, align: str = "left", spacing: int = 6) -> object:
        font = font or self._font_b()
        color = color_hex or self._c(color_role)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        from ppt_pro_max.renderer.theme_mapper import get_cjk_companion
        cjk = get_cjk_companion(font, "heading" if size >= 20 else "body")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
            run = p.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = self._rgb(color)
            run.font.bold = bold
            self._set_font_with_cjk(run, font, cjk)
            p.space_after = Pt(spacing)
        return tb

    # ── Image helpers (Pillow pre-crop) ──

    def add_image(self, slide, path: str, x: float, y: float, w: float, h: float) -> None:
        if not os.path.isfile(path):
            return
        with PILImage.open(path) as img:
            iw, ih = img.size
            box_ratio = w / h
            img_ratio = iw / ih
            if img_ratio > box_ratio:
                cw, ch = int(ih * box_ratio), ih
                cl, ct = (iw - cw) // 2, 0
            else:
                cw, ch = iw, int(iw / box_ratio)
                cl, ct = 0, (ih - ch) // 2
            cropped = img.crop((cl, ct, cl + cw, ct + ch))
            cp = os.path.join(self._crop_cache_dir,
                              hashlib.md5(f"{path}:{w}x{h}".encode()).hexdigest() + ".png")
            if not os.path.isfile(cp):
                cropped.save(cp, "PNG")
        slide.shapes.add_picture(cp, Inches(x), Inches(y), Inches(w), Inches(h))

    # ── Shape helpers ──

    def add_rect(self, slide, x: float, y: float, w: float, h: float,
                 fill_role: str | None = None, fill_hex: str | None = None,
                 border_role: str | None = None, border_hex: str | None = None,
                 gradient: bool = False, shadow: bool = False) -> object:
        fill = fill_hex or self._c(fill_role or "muted")
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        if gradient:
            apply_gradient(sh, self._lighten(fill), fill, gradient_type="linear", angle=5400000)
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=4, distance_pt=2, alpha_pct=15)
        border = border_hex or (self._c(border_role) if border_role else None)
        if border:
            sh.line.color.rgb = self._rgb(border)
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    def add_rounded_rect(self, slide, x: float, y: float, w: float, h: float,
                         fill_role: str | None = None, fill_hex: str | None = None,
                         border_role: str | None = None, border_hex: str | None = None,
                         gradient: bool = False, shadow: bool = False,
                         corner_radius: str | int = "md") -> object:
        fill = fill_hex or self._c(fill_role or "muted")
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        if isinstance(corner_radius, str):
            radius_pt = CORNER_RADIUS_SCALE.get(corner_radius, 8)
        else:
            radius_pt = corner_radius
        min_dim = min(Inches(w), Inches(h))
        adj_val = int((radius_pt * 12700) / (min_dim / 2) * 100000) if min_dim > 0 else 16667
        adj_val = min(100000, max(0, adj_val))
        spPr = sh._element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn("a:avLst"))
            for existing_gd in avLst.findall(qn("a:gd")):
                if existing_gd.get("name") == "adj":
                    existing_gd.set("fmla", f"val {adj_val}")
                    break
            else:
                gd = etree.SubElement(avLst, qn("a:gd"))
                gd.set("name", "adj")
                gd.set("fmla", f"val {adj_val}")
        if gradient:
            apply_gradient(sh, self._lighten(fill), fill, gradient_type="linear", angle=2700000)
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=4, distance_pt=2, alpha_pct=15)
        border = border_hex or (self._c(border_role) if border_role else None)
        if border:
            sh.line.color.rgb = self._rgb(border)
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    def add_oval(self, slide, x: float, y: float, w: float, h: float,
                 fill_role: str | None = None, fill_hex: str | None = None,
                 gradient: bool = True, shadow: bool = True,
                 label: str = "", font_size: int = 16,
                 font_color: str | None = None) -> object:
        if fill_role is None and fill_hex is None:
            fill_role = "muted" if self._is_dark() else "primary"
        fill = fill_hex or self._c(fill_role)
        effective_font_color = font_color or (self._c("foreground", "#FFFFFF") if self._is_dark() else self._c("on-primary", "#FFFFFF"))
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        if gradient:
            apply_gradient(sh, fill, self._darken(fill), gradient_type="path")
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=4, distance_pt=2, alpha_pct=25)
        sh.line.fill.background()
        if label:
            tf = sh.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = self._rgb(effective_font_color)
            run.font.bold = True
        return sh

    def add_donut(self, slide, x: float, y: float, size: float,
                  fill_role: str | None = None, fill_hex: str | None = None,
                  gradient: bool = True, shadow: bool = True,
                  label: str = "", font_size: int = 18) -> object:
        if fill_role is None and fill_hex is None:
            fill_role = "muted" if self._is_dark() else "primary"
        fill = fill_hex or self._c(fill_role)
        font_color = self._c("foreground", "#FFFFFF") if self._is_dark() else self._c("on-primary", "#FFFFFF")
        sh = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(x), Inches(y), Inches(size), Inches(size))
        if gradient:
            apply_gradient(sh, fill, self._darken(fill), gradient_type="path")
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=6, distance_pt=3, alpha_pct=20)
        sh.line.fill.background()
        if label:
            tf = sh.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = self._rgb(font_color)
            run.font.bold = True
        return sh

    def add_hexagon(self, slide, x: float, y: float, size: float,
                    fill_role: str | None = None, fill_hex: str | None = None,
                    gradient: bool = True, shadow: bool = True,
                    label: str = "", font_size: int = 16) -> object:
        if fill_role is None and fill_hex is None:
            fill_role = "muted" if self._is_dark() else "primary"
        fill = fill_hex or self._c(fill_role)
        font_color = self._c("on-primary", "#FFFFFF") if not self._is_dark() else self._c("foreground", "#FFFFFF")
        sh = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(size), Inches(size * 0.87))
        if gradient:
            apply_gradient(sh, fill, self._darken(fill), gradient_type="linear", angle=5400000)
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._rgb(fill)
        if shadow:
            apply_shadow(sh, blur_pt=4, distance_pt=2, alpha_pct=25)
        sh.line.fill.background()
        if label:
            tf = sh.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = self._rgb(font_color)
            run.font.bold = True
        return sh

    @staticmethod
    def _lighten(hex_color: str, amount: int = 30) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = min(255, r + amount)
        g = min(255, g + amount)
        b = min(255, b + amount)
        return f"{r:02X}{g:02X}{b:02X}"

    @staticmethod
    def _darken(hex_color: str, amount: int = 30) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = max(0, r - amount)
        g = max(0, g - amount)
        b = max(0, b - amount)
        return f"{r:02X}{g:02X}{b:02X}"

    # ── Overlay helpers ──

    def add_dark_overlay(self, slide, opacity: float = 0.65) -> None:
        bg_hex = self._c("background", "#060B18" if self._is_dark() else "#000000")
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT))
        ov.fill.solid()
        ov.fill.fore_color.rgb = self._rgb(bg_hex)
        ov.line.fill.background()
        el = ov._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if el is not None:
            a = etree.SubElement(el, qn("a:alpha"))
            a.set("val", str(int(opacity * 100000)))

    def add_gradient_overlay(self, slide, opacity_bottom: float = 0.72,
                              opacity_top: float = 0.0,
                              color_role: str = "background") -> None:
        bg_hex = self._c(color_role, "#000000")
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT))
        ov.line.fill.background()
        top_alpha = int(opacity_top * 100000)
        mid_alpha = int(opacity_bottom * 0.4 * 100000)
        bot_alpha = int(opacity_bottom * 100000)
        grad = GradientFill(
            stops=[
                GradientStop(color=bg_hex, position=0, alpha=top_alpha),
                GradientStop(color=bg_hex, position=40000, alpha=mid_alpha),
                GradientStop(color=bg_hex, position=100000, alpha=bot_alpha),
            ],
            angle=5400000,
        )
        grad.apply(ov)

    def add_overlay(self, slide, x: float, y: float, w: float, h: float,
                    color_hex: str = "#000000", opacity: float = 0.65) -> None:
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        ov.fill.solid()
        ov.fill.fore_color.rgb = self._rgb(color_hex)
        ov.line.fill.background()
        el = ov._element.find(qn('p:spPr')).find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if el is not None:
            a = etree.SubElement(el, qn('a:alpha'))
            a.set('val', str(int(opacity * 100000)))

    # ── Brand visual design ──

    def render_slide(self, prs: Presentation, page: dict[str, Any], 
                     component_lib=None, layout_variant: dict | None = None,
                     page_index: int = 0, total_pages: int = 0) -> object:
        elements = page.get("elements")
        if elements:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            self._render_elements(slide, elements)
            return slide

        goal = page.get("goal", "content")
        title = page.get("title", "")
        subtitle = page.get("subtitle")
        bullets = page.get("bullets") or []
        image_path = page.get("image")
        cards = page.get("cards") or []
        diagram_type = page.get("diagram_type")
        diagram_data = page.get("diagram_data")
        code = page.get("code")
        exercise = page.get("exercise")
        chart = page.get("chart")
        notes = page.get("notes")
        _image_grid = page.get("image_grid")
        _icons = page.get("icons")
        component_type = page.get("component_type")
        component_category = page.get("component_category")
        explicit_layout = page.get("layout")

        if component_lib is not None and not component_type:
            if bullets or cards:
                p_type, p_cat = self._proactive_component_match(
                    bullets, cards, explicit_layout or "auto", component_lib,
                    existing_component_type=component_type,
                    mood=page.get("mood"),
                )
                if p_type:
                    component_type = p_type
                    component_category = p_cat

        if explicit_layout:
            goal = self._remap_layout_to_goal(explicit_layout, goal)

        is_hero = goal in ("hook", "cta")
        is_section = goal == "section"

        variant = layout_variant or {}
        margin_left = variant.get("content_margin_left", 0.9)
        title_align = variant.get("title_alignment", "left")
        decoration_style = variant.get("decoration_style", "accent-bar")
        if isinstance(page, dict):
            deco = page.get("decoration")
            if isinstance(deco, dict):
                decoration_style = deco.get("style", decoration_style)

        layout_name = page.get("template_layout_name")
        slide = self.add_slide(prs, layout_name=layout_name)

        if is_section:
            section_num = page.get("section_number", page_index + 1)
            section_sub = subtitle or ""
            self.render_section_divider(slide, section_num, title, section_sub)
            return slide

        if is_hero:
            blocks = page.get("blocks")
            has_blocks = bool(blocks)
            self.apply_hero_overlay(slide, prs, image_path=image_path if image_path and os.path.isfile(image_path) else None)
            title_y = 1.0 if has_blocks else 2.0
            subtitle_y = title_y + 1.4
            if title:
                self.add_text(slide, title, 1.2, title_y, 8, 1.5,
                              size=52 if not has_blocks else 44, color_role="on-primary" if not (image_path and os.path.isfile(image_path)) else "foreground", bold=True)
            if subtitle:
                self.add_text(slide, subtitle, 1.2, subtitle_y, 8, 0.8,
                              font=self._font_b(), size=28, color_role="foreground")
            if bullets and not has_blocks:
                bullet_text = "  \u2022  ".join(bullets[:5])
                self.add_text(slide, bullet_text, 1.2, subtitle_y + 0.8, 8, 0.5,
                              font=self._font_b(), size=14, color_role="muted-foreground")
            if blocks:
                self._render_blocks(slide, blocks, is_hero=True)
        else:
            self.apply_brand_background(slide, prs, goal=goal, 
                                        page_index=page_index, total_pages=total_pages)

            if title:
                align = "center" if title_align == "center" else "left"
                self.add_text(slide, title, margin_left, 0.5, 11, 0.8,
                              size=36, color_role="foreground", bold=True, align=align)
                accent_hex = self._c("accent", self._c("primary", "#2563EB"))
                self.add_gradient_line(slide, margin_left, 1.2, 2.5, 0.04, accent_hex)

                colors = {
                    "primary": self._c("primary", "#2563EB"),
                    "accent": self._c("accent", "#F97316"),
                    "foreground": self._c("foreground", "#1E293B"),
                    "muted": self._c("muted", "#F1F5F9"),
                    "border": self._c("border", "#E2E8F0"),
                }
                self._decoration_renderer.apply_title_decoration(
                    slide, margin_left, 0.5, 11, decoration_style, colors,
                    add_rect_fn=lambda sl, x, y, w, h, **kw: self.add_rect(sl, x, y, w, h, **kw),
                    add_oval_fn=lambda sl, x, y, w, h, **kw: self.add_oval(sl, x, y, w, h, **kw),
                    add_text_fn=lambda sl, t, x, y, w, h, **kw: self.add_text(sl, t, x, y, w, h, **kw),
                    apply_glow_fn=apply_glow,
                )

            if subtitle:
                self.add_text(slide, subtitle, margin_left, 1.5, 10, 0.5,
                              font=self._font_b(), size=14, color_role="muted-foreground")

            blocks = page.get("blocks")
            if blocks:
                self._render_blocks(slide, blocks)
            elif cards:
                n = len(cards)
                card_w = min(3.6, (11.5 - 0.4 * (n - 1)) / n)
                for i, card in enumerate(cards):
                    xx = margin_left + i * (card_w + 0.4)
                    card_title = card.get("title", "")
                    card_body = card.get("text", card.get("body", ""))
                    self.add_card(slide, xx, 1.6, card_w, 4.5, card_title, card_body,
                                  featured=(i == 0))

            elif component_type and component_lib is not None:
                enriched_page = dict(page)
                enriched_page["component_type"] = component_type
                if component_category:
                    enriched_page["component_category"] = component_category
                self._render_component_on_slide(slide, enriched_page, component_lib)

            elif diagram_type and diagram_data:
                self._render_diagram_on_slide(slide, diagram_type, diagram_data)

            elif code:
                self._render_code_on_slide(slide, code)

            elif exercise:
                self._render_exercise_on_slide(slide, exercise)

            elif bullets:
                self._render_bullets_on_slide(slide, bullets, margin_left)

            if image_path and os.path.isfile(image_path) and not is_hero:
                self.add_masked_image(slide, image_path, 8.3, 1.2, 4.2, 5.3)

            if chart:
                self._render_chart_on_slide(slide, chart)

            if total_pages > 0 and not is_hero:
                self.add_progress_bar(slide, page_index + 1, total_pages)

        if notes:
            self._render_notes_on_slide(slide, notes)

        return slide

    def _render_blocks(self, slide, blocks: list[dict], is_hero: bool = False) -> None:
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(self)
        br.render(slide, blocks, is_hero=is_hero)

    def _render_elements(self, slide, elements: list[dict]) -> None:
        for el in elements:
            self._dispatch_element(slide, el)

    def _dispatch_element(self, slide, el: dict) -> None:
        kind = el.get("type")
        x = el.get("x", 0.0)
        y = el.get("y", 0.0)
        w = el.get("w", 4.0)
        h = el.get("h", 1.0)

        if kind == "text":
            color = el.get("color")
            color_role = el.get("color_role")
            effective_color = color or self._c(color_role or "foreground", "#1A1A1A")
            self.add_text(slide, el.get("text", ""), x, y, w, h,
                          font=el.get("font", self._font_h()),
                          size=el.get("size", 18),
                          color_hex=effective_color,
                          bold=el.get("bold", False),
                          align=el.get("align", "left"))

        elif kind == "multiline":
            lines = el.get("lines", [])
            color = el.get("color")
            color_role = el.get("color_role")
            effective_color = color or self._c(color_role or "foreground", "#1A1A1A")
            self.add_multiline(slide, lines, x, y, w, h,
                               font=el.get("font", self._font_b()),
                               size=el.get("size", 14),
                               color_hex=effective_color,
                               bold=el.get("bold", False),
                               align=el.get("align", "left"),
                               spacing=el.get("spacing", 6))

        elif kind == "rect":
            fill = el.get("fill")
            fill_role = el.get("fill_role")
            border = el.get("border")
            border_role = el.get("border_role")
            self.add_rect(slide, x, y, w, h,
                          fill_hex=fill, fill_role=fill_role,
                          border_hex=border, border_role=border_role,
                          gradient=el.get("gradient", False),
                          shadow=el.get("shadow", False))

        elif kind == "rounded_rect":
            fill = el.get("fill")
            fill_role = el.get("fill_role")
            border = el.get("border")
            border_role = el.get("border_role")
            self.add_rounded_rect(slide, x, y, w, h,
                                  fill_hex=fill, fill_role=fill_role,
                                  border_hex=border, border_role=border_role,
                                  gradient=el.get("gradient", False),
                                  shadow=el.get("shadow", False),
                                  corner_radius=el.get("radius", "md"))

        elif kind == "image":
            path = el.get("path", "")
            if os.path.isfile(path):
                self.add_image(slide, path, x, y, w, h)

        elif kind == "overlay":
            color = el.get("color") or self._c("background", "#000000")
            opacity = el.get("opacity", 0.65)
            self.add_overlay(slide, x, y, w, h, color, opacity)

        elif kind == "gradient_line":
            color = el.get("color") or self._c("accent", self._c("primary", "#2563EB"))
            self.add_gradient_line(slide, x, y, w, h, color)

        elif kind == "hexagon":
            size = el.get("size", min(w, h))
            fill = el.get("fill")
            fill_role = el.get("fill_role", "primary")
            self.add_hexagon(slide, x, y, size,
                             fill_hex=fill, fill_role=fill_role,
                             gradient=el.get("gradient", True),
                             shadow=el.get("shadow", True),
                             label=el.get("label", ""),
                             font_size=el.get("font_size", 16))

        elif kind == "oval":
            fill = el.get("fill")
            fill_role = el.get("fill_role", "primary")
            self.add_oval(slide, x, y, w, h,
                          fill_hex=fill, fill_role=fill_role,
                          gradient=el.get("gradient", True),
                          shadow=el.get("shadow", True),
                          label=el.get("label", ""),
                          font_size=el.get("font_size", 16))

        elif kind == "donut":
            size = el.get("size", min(w, h))
            fill = el.get("fill")
            fill_role = el.get("fill_role", "primary")
            self.add_donut(slide, x, y, size,
                           fill_hex=fill, fill_role=fill_role,
                           gradient=el.get("gradient", True),
                           shadow=el.get("shadow", True),
                           label=el.get("label", ""),
                           font_size=el.get("font_size", 18))

    def _render_bullets_on_slide(self, slide, bullets: list, margin_left: float = 0.9) -> None:
        if len(bullets) >= 6:
            col_w = (11.5 - 0.3) / 2
            mid = (len(bullets) + 1) // 2
            left_lines = [f"\u2022  {b}" for b in bullets[:mid]]
            self.add_multiline(slide, left_lines, margin_left, 1.6, col_w, 4.5,
                               size=14, color_role="foreground", spacing=8)
            right_lines = [f"\u2022  {b}" for b in bullets[mid:]]
            self.add_multiline(slide, right_lines, margin_left + col_w + 0.3, 1.6, col_w, 4.5,
                               size=14, color_role="foreground", spacing=8)
            self.add_rect(slide, margin_left + col_w + 0.15, 1.7, 0.01, 4.0,
                          fill_hex=self._c("border", "#E2E8F0"))
        else:
            bullet_lines = [f"\u2022  {b}" for b in bullets[:8]]
            self.add_multiline(slide, bullet_lines, margin_left, 1.6, 7, 4.5,
                               size=14, color_role="foreground", spacing=6)

    def _render_component_on_slide(self, slide, page: dict[str, Any], component_lib) -> None:
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        component_type = page.get("component_type", "")
        component_category = page.get("component_category", "")
        component_variant = page.get("component_variant", "")
        component_fit = page.get("component_fit", "contain")
        bullets = page.get("bullets") or []

        component_bounds = page.get("component_bounds")
        if component_bounds and len(component_bounds) == 4:
            content_area = tuple(float(v) for v in component_bounds)
        else:
            content_area = self._compute_content_area(page)

        element = {
            "type": component_type,
            "category": component_category,
            "variant": component_variant,
            "texts": bullets,
            "nodes": [{"text": b} for b in bullets],
            "node_count": len(bullets) if bullets else 0,
            "bounds": content_area,
            "component_fit": component_fit,
        }

        renderer = ComponentRenderer()
        success = renderer.render(slide, element, self._brand, component_lib)

        if not success:
            if page.get("diagram_type") and page.get("diagram_data"):
                self._render_diagram_on_slide(slide, page["diagram_type"], page["diagram_data"])
            elif bullets:
                bullet_lines = [f"\u2022  {b}" for b in bullets[:8]]
                self.add_multiline(slide, bullet_lines, 0.9, 1.6, 7, 4.5,
                                   size=14, color_role="foreground", spacing=6)

    def _compute_content_area(self, page: dict[str, Any]) -> tuple[float, float, float, float]:
        goal = page.get("goal", "content")
        has_image = bool(page.get("image"))
        has_title = bool(page.get("title"))
        has_subtitle = bool(page.get("subtitle"))

        slide_w = SLIDE_WIDTH
        slide_h = SLIDE_HEIGHT

        margin_left = 0.9
        margin_right = 0.4
        margin_bottom = 0.5

        top = 0.5
        if has_title:
            top += 0.9
        if has_subtitle:
            top += 0.5

        right_edge = slide_w - margin_right
        if has_image and goal not in ("hook", "cta"):
            right_edge = 8.0

        content_w = right_edge - margin_left
        content_h = slide_h - top - margin_bottom

        return (margin_left, top, content_w, content_h)

    def _remap_layout_to_goal(self, layout_name: str, fallback_goal: str) -> str:
        layout_to_goal = {
            "title-slide": "hook",
            "cta-closing": "cta",
            "content-with-title": "content",
            "two-column": "content",
            "three-column-cards": "features",
            "four-metrics": "traction",
            "big-number": "agitation",
            "quote": "testimonials",
            "chart-focus": "data",
            "image-plus-text": "content",
            "sidebar-left": "overview",
            "exercise-layout": "exercise",
            "code-block": "code",
            "funnel": "funnel",
            "grid-2x2-cards": "features",
            "dense-bullets": "content",
            "swot-matrix": "content",
            "cycle-diagram": "content",
            "timeline-horizontal": "content",
            "table-layout": "content",
            "section-header": "content",
            "blank": "content",
        }
        return layout_to_goal.get(layout_name, fallback_goal)

    def _render_diagram_on_slide(self, slide, diagram_type: str, diagram_data: dict) -> None:
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        style = DiagramStyle.from_brand_spec(self.brand) if self._brand else DiagramStyle()
        region = Region(left=0.9, top=1.5, width=7.0, height=5.0)
        engine = DiagramEngine()
        try:
            engine.render(slide, diagram_type, diagram_data, style, region)
        except Exception:
            pass

    def _render_code_on_slide(self, slide, code_data) -> None:
        code_text = code_data if isinstance(code_data, str) else code_data.get("source", code_data.get("code", ""))
        language = code_data.get("language", "") if isinstance(code_data, dict) else ""
        code_bg = "#1E293B"
        self.add_rounded_rect(slide, 0.9, 1.5, 11.533, 5.0, fill_hex=code_bg, corner_radius="lg")
        if language:
            badge_text = language.upper()
            badge_w = len(badge_text) * 0.1 + 0.3
            self.add_rounded_rect(slide, 1.0, 1.2, badge_w, 0.3,
                                  fill_role="primary", corner_radius="sm")
            self.add_text(slide, badge_text, 1.1, 1.22, badge_w - 0.2, 0.26,
                          size=11, color_role="on-primary", bold=True)
        lines = code_text.split("\n")
        all_lines = [f"  {line}" for line in lines[:30]]
        self.add_multiline(slide, all_lines, 1.2, 1.7, 11, 4.5,
                           font="Consolas", size=11, color_role="muted-foreground", spacing=4)

    def _render_exercise_on_slide(self, slide, exercise_data) -> None:
        instructions = exercise_data.get("instructions", "") if isinstance(exercise_data, dict) else str(exercise_data)
        duration = exercise_data.get("duration", "") if isinstance(exercise_data, dict) else ""
        steps = exercise_data.get("steps", []) if isinstance(exercise_data, dict) else []
        badge_text = f"Exercise {duration}" if duration else "Exercise"
        self.add_badge(slide, badge_text, 0.9, 1.2, variant="solid")
        if instructions:
            self.add_text(slide, instructions, 0.9, 1.8, 11.533, 1.2,
                          font=self._font_b(), size=13, color_role="muted-foreground")
        if steps:
            step_lines = [f"{i+1}. {s}" for i, s in enumerate(steps)]
            self.add_multiline(slide, step_lines, 0.9, 3.2, 11.533, 3.5,
                               size=13, color_role="foreground", spacing=6)

    def _render_chart_on_slide(self, slide, chart_config: dict) -> None:
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        brand_colors = None
        if self._brand and self._brand.colors:
            brand_colors = self._brand.colors
        position = {"x": 1.5, "y": 1.5, "width": 10.333, "height": 4.5}
        builder = ChartBuilder()
        try:
            builder.build(slide, chart_config, position=position, brand_colors=brand_colors)
        except Exception:
            pass

    def _render_notes_on_slide(self, slide, notes_text: str) -> None:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    def apply_brand_background(self, slide, prs: Presentation, goal: str = "content",
                                page_index: int = 0, total_pages: int = 0) -> None:
        bg_hex = self._c("background", "#FFFFFF")
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self._rgb(bg_hex)
        except Exception:
            pass

        strip_style = "auto"
        if self._brand and self._brand.spacing:
            strip_style = self._brand.spacing.get("strip_style", "auto")

        if strip_style == "none":
            pass
        elif strip_style == "auto":
            variant = page_index % 3
            if variant == 0 and page_index > 0:
                accent_hex = self._c("accent", self._c("primary", "#2563EB"))
                self.add_rect(slide, 0, 0, 0.06, SLIDE_HEIGHT, fill_hex=accent_hex, gradient=True)
            elif variant == 1:
                accent_hex = self._c("accent", self._c("primary", "#2563EB"))
                self.add_rect(slide, 0, SLIDE_HEIGHT - 0.03, SLIDE_WIDTH, 0.03, fill_hex=accent_hex)
        elif strip_style == "left":
            accent_hex = self._c("accent", self._c("primary", "#2563EB"))
            self.add_rect(slide, 0, 0, 0.06, SLIDE_HEIGHT, fill_hex=accent_hex, gradient=True)

        muted_hex = self._c("muted", "#F1F5F9")
        show_footer = self._brand and self._brand.footer and self._brand.footer.get("show_footer_text")
        has_progress = total_pages > 0
        if not show_footer and not has_progress:
            self.add_rect(slide, 0, SLIDE_HEIGHT - 0.15, SLIDE_WIDTH, 0.15, fill_hex=muted_hex, gradient=True)

    def apply_hero_overlay(self, slide, prs: Presentation, image_path: str | None = None) -> None:
        if image_path and os.path.isfile(image_path):
            self.add_image(slide, image_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
            self.add_gradient_overlay(slide, opacity_bottom=0.72, opacity_top=0.0)
        else:
            primary_hex = self._c("primary", "#2563EB")
            self.add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_hex=primary_hex, gradient=True)

    # ── Card component ──

    def add_card(self, slide, x: float, y: float, w: float, h: float,
                 title: str, body: str, accent_role: str = "primary",
                 featured: bool = False) -> None:
        accent_hex = self._c(accent_role)
        card_bg = self._c("muted", "#0D152A" if self._is_dark() else "#F8FAFC")
        card_bd = self._c("border", "#1A2A4A" if self._is_dark() else "#E2E8F0")
        title_size = 24 if featured else 20
        self.add_rounded_rect(slide, x, y, w, h, fill_hex=card_bg, border_hex=card_bd, gradient=True, shadow=True)
        if featured:
            self.add_rect(slide, x, y, w, 0.15, fill_hex=accent_hex, gradient=True)
            self.add_text(slide, title, x + 0.2, y + 0.25, w - 0.4, 0.5,
                          self._font_h(), title_size, accent_role, bold=True)
            self.add_multiline(slide, body.split("\n"), x + 0.2, y + 0.8, w - 0.4, h - 1.1,
                               self._font_b(), 15, "muted-foreground")
        else:
            self.add_rect(slide, x + 0.2, y + 0.2, w - 0.4, 0.04, fill_hex=accent_hex, gradient=True)
            self.add_text(slide, title, x + 0.2, y + 0.4, w - 0.4, 0.5,
                          self._font_h(), title_size, accent_role, bold=True)
            self.add_multiline(slide, body.split("\n"), x + 0.2, y + 0.9, w - 0.4, h - 1.2,
                               self._font_b(), 15, "muted-foreground")

    # ── Brand compliance: logo, footer, watermark ──

    def apply_logo(self, slide, logo_path: str, prs: Presentation,
                   current_goal: str | None = None) -> None:
        if not self._brand or not self._brand.logo:
            return
        logo_spec = self._brand.logo
        skip_slides = logo_spec.get("skip_slides", [])
        if current_goal and current_goal in skip_slides:
            return
        if not os.path.isfile(logo_path):
            return

        position = logo_spec.get("position", "top_right")
        width_inches = logo_spec.get("width_inches", 1.0)
        width_emu = Inches(width_inches)

        try:
            with PILImage.open(logo_path) as img:
                aspect = img.size[1] / img.size[0] if img.size[0] > 0 else 0.5
        except Exception:
            aspect = 0.5
        height_emu = int(width_emu * aspect)

        slide_w, slide_h = prs.slide_width, prs.slide_height
        if position == "top_right":
            left, top = slide_w - width_emu - Inches(0.3), Inches(0.3)
        elif position == "top_left":
            left, top = Inches(0.3), Inches(0.3)
        elif position == "bottom_right":
            left, top = slide_w - width_emu - Inches(0.3), slide_h - height_emu - Inches(0.3)
        else:
            left, top = slide_w - width_emu - Inches(0.3), Inches(0.3)

        slide.shapes.add_picture(logo_path, left, top, width=width_emu, height=height_emu)

    def apply_footer(self, prs: Presentation) -> None:
        if not self._brand or not self._brand.footer:
            return
        config = self._brand.footer
        total = len(prs.slides)
        show_page_number = config.get("show_page_number", False)
        page_number_format = config.get("page_number_format", "{n}")
        page_number_position = config.get("page_number_position", "bottom_right")
        show_footer_text = config.get("show_footer_text", False)
        footer_text = config.get("footer_text", "")
        footer_position = config.get("footer_position", "bottom_center")
        font_size_pt = config.get("font_size_pt", 10)
        skip_pages = config.get("skip_pages", [])

        if not show_page_number and not show_footer_text:
            return

        muted_color = self._rgb(self._c("muted-foreground", "#999999"))
        position_map = {
            "bottom_left": Inches(0.9),
            "bottom_center": Inches(5.833),
            "bottom_right": Inches(11.433),
        }

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            is_cover = idx == 0
            should_skip = is_cover or slide_num in skip_pages

            if show_page_number and not should_skip:
                page_text = page_number_format.replace("{n}", str(slide_num)).replace("{total}", str(total))
                left = position_map.get(page_number_position, Inches(11.433))
                tb = slide.shapes.add_textbox(left, Inches(7.0), Inches(2.0), Inches(0.3))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT
                run = p.add_run()
                run.text = page_text
                run.font.size = Pt(font_size_pt)
                run.font.color.rgb = muted_color
                run.font.name = self._font_b()

            if show_footer_text and footer_text and not should_skip:
                left = position_map.get(footer_position, Inches(5.833))
                tb = slide.shapes.add_textbox(left, Inches(7.0), Inches(4.0), Inches(0.3))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = footer_text
                run.font.size = Pt(font_size_pt)
                run.font.color.rgb = muted_color
                run.font.name = self._font_b()

    def apply_watermark(self, prs: Presentation) -> None:
        if not self._brand or not self._brand.watermark:
            return
        config = self._brand.watermark
        text = config.get("text", "CONFIDENTIAL")
        opacity = config.get("opacity", 0.15)
        rotation = config.get("rotation", -45)
        font_size_pt = config.get("font_size_pt", 72)
        skip_pages = config.get("skip_pages", [1])

        muted_hex = self._c("muted-foreground", "#999999")

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            if slide_num in skip_pages:
                continue

            tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = self._rgb(muted_hex)
            run.font.bold = True
            run.font.name = self._font_h()

            sp = tb._element
            srgbClr_el = sp.find(".//" + qn("a:srgbClr"))
            if srgbClr_el is not None:
                alpha_elem = etree.SubElement(srgbClr_el, qn("a:alpha"))
                alpha_elem.set("val", str(int(opacity * 100000)))

            xfrm_el = sp.find(".//" + qn("a:xfrm"))
            if xfrm_el is not None:
                xfrm_el.set("rot", str(int(rotation * 60000)))

    # ── Save ──

    def save(self, prs: Presentation, output_path: str) -> None:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        prs.save(output_path)

    # ── §2.1 CJK font pairing ──

    def _set_font_with_cjk(self, run, latin_font: str, cjk_font: str | None = None) -> None:
        run.font.name = latin_font
        if cjk_font:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", cjk_font)
            cs = rPr.find(qn("a:cs"))
            if cs is None:
                cs = etree.SubElement(rPr, qn("a:cs"))
            cs.set("typeface", cjk_font)

    # ── §2.3 Badge/Tag styling ──

    def add_badge(self, slide, text: str, x: float, y: float,
                  variant: str = "default", size: str = "sm") -> None:
        font_size = {"sm": 11, "md": 11, "lg": 12}.get(size, 11)
        badge_text = text.upper()
        cjk_count = sum(1 for ch in badge_text if '\u4e00' <= ch <= '\u9fff' or '\u3000' <= ch <= '\u303f')
        latin_count = len(badge_text) - cjk_count
        char_width = font_size * 0.009
        cjk_width = font_size * 0.016
        badge_w = latin_count * char_width + cjk_count * cjk_width + 0.3
        badge_h = 0.35
        if variant == "default":
            tint_hex = self._lighten(self._c("primary", "#2563EB"))
            self.add_rounded_rect(slide, x, y, badge_w, badge_h,
                                  fill_hex=tint_hex, corner_radius="sm")
            self.add_text(slide, badge_text, x + 0.1, y + 0.02, badge_w - 0.2, badge_h - 0.04,
                          size=font_size, color_role="primary", bold=True)
        elif variant == "solid":
            self.add_rounded_rect(slide, x, y, badge_w, badge_h,
                                  fill_role="primary", corner_radius="sm")
            self.add_text(slide, badge_text, x + 0.1, y + 0.02, badge_w - 0.2, badge_h - 0.04,
                          size=font_size, color_role="on-primary", bold=True)
        elif variant == "outline":
            self.add_rounded_rect(slide, x, y, badge_w, badge_h,
                                  fill_hex=self._c("background", "#FFFFFF"),
                                  border_role="primary", corner_radius="sm")
            self.add_text(slide, badge_text, x + 0.1, y + 0.02, badge_w - 0.2, badge_h - 0.04,
                          size=font_size, color_role="primary", bold=True)

    # ── §2.4 Section divider ──

    def render_section_divider(self, slide_or_prs, section_number: int,
                                section_title: str, section_subtitle: str = "") -> object:
        if hasattr(slide_or_prs, 'background'):
            slide = slide_or_prs
        else:
            slide = self.add_slide(slide_or_prs)
        primary_hex = self._c("primary", "#2563EB")
        if self._is_dark():
            bg = self._lighten(primary_hex, 120)
        else:
            bg = self._lighten(primary_hex, 80)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(bg)
        self.add_text(slide, f"{section_number:02d}",
                      2.0, 1.5, 9.333, 2.0,
                      size=72, color_role="primary", bold=True)
        self.add_text(slide, section_title,
                      2.0, 3.5, 9.333, 1.0,
                      size=40, color_role="foreground", bold=True)
        self.add_rect(slide, 2.0, 4.6, 3.0, 0.03, fill_role="accent", gradient=True)
        if section_subtitle:
            self.add_text(slide, section_subtitle,
                          2.0, 4.8, 9.333, 0.5,
                          font=self._font_b(), size=18, color_role="muted-foreground")
        return slide

    # ── §3.2 Progress bar ──

    def add_progress_bar(self, slide, current: int, total: int) -> None:
        bar_y = SLIDE_HEIGHT - 0.04
        bar_h = 0.03
        self.add_rect(slide, 0, bar_y, SLIDE_WIDTH, bar_h,
                      fill_hex=self._c("border", "#E2E8F0"))
        fill_w = SLIDE_WIDTH * (current / total)
        self.add_rect(slide, 0, bar_y, fill_w, bar_h,
                      fill_hex=self._c("primary", "#2563EB"))

    # ── §3.4 Gradient line ──

    def add_gradient_line(self, slide, x: float, y: float, width: float, height: float,
                           from_color: str, to_color: str = "transparent") -> object:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(y), Inches(width), Inches(height))
        rect.line.fill.background()
        grad = GradientFill(
            stops=[
                GradientStop(color=from_color, position=0, alpha=100000),
                GradientStop(color=from_color if to_color == "transparent" else to_color,
                             position=100000, alpha=0),
            ],
            angle=0,
        )
        grad.apply(rect)
        return rect

    # ── §3.5 Image masking ──

    def add_masked_image(self, slide, image_path: str, x: float, y: float,
                          w: float, h: float, padding: float = 0.15,
                          corner_radius: str = "md") -> None:
        frame_bg = "#FFFFFF" if not self._is_dark() else "#1E293B"
        self.add_rounded_rect(slide, x, y, w, h, fill_hex=frame_bg,
                              shadow=True, corner_radius=corner_radius)
        img_x = x + padding
        img_y = y + padding
        img_w = w - 2 * padding
        img_h = h - 2 * padding
        self.add_image(slide, image_path, img_x, img_y, img_w, img_h)

    # ── §3.7 Hero pattern selection ──

    def _select_hero_pattern(self, page: dict, mood: str = "professional") -> str:
        subtitle = page.get("subtitle", "")
        has_image = bool(page.get("image"))
        if not has_image:
            return "gradient"
        if len(subtitle) > 60:
            return "split-left"
        if mood in ("creative", "bold", "vibrant", "startup"):
            return "asymmetric"
        return "bottom-fade"

    _MOOD_CATEGORY_MAP: dict[str, str] = {
        "mckinsey": "hierarchy",
        "consulting": "process",
        "tech": "cycle",
        "creative": "radial",
        "dark": "chart",
        "fintech": "funnel",
        "education": "pyramid",
        "health": "cycle",
        "government": "hierarchy",
        "industrial": "process",
        "startup": "infographic",
        "luxury": "pyramid",
        "nature": "cycle",
        "minimal": "infographic",
        "bold": "comparison",
        "international": "matrix",
        "cream": "infographic",
        "frosted": "matrix",
        "pastel": "infographic",
        "retro": "timeline",
        "legal": "hierarchy",
        "pharma": "funnel",
        "realestate": "pyramid",
        "automotive": "process",
        "aviation": "process",
        "energy": "funnel",
        "telecom": "cycle",
        "logistics": "process",
    }

    def _mood_to_preferred_category(self, mood: str | None) -> str | None:
        if not mood:
            return None
        return self._MOOD_CATEGORY_MAP.get(mood)

    def _proactive_component_match(
        self,
        bullets: list[str],
        cards: list[dict],
        layout: str,
        component_lib,
        existing_component_type: str | None = None,
        mood: str | None = None,
    ) -> tuple[str | None, str | None]:
        if existing_component_type:
            return (None, None)
        if component_lib is None:
            return (None, None)
        if not bullets and not cards:
            return (None, None)

        from ppt_pro_max.enterprise.content_parser import infer_component_category

        items = bullets if bullets else [c.get("title", "") for c in cards]
        n = len(items)
        if n < 2:
            return (None, None)

        _, category = infer_component_category(items)

        if layout == "cards" and category in ("process", "pyramid", "matrix"):
            category = "infographic"

        matched = self._search_and_validate(component_lib, category, n)
        if matched:
            return ("group", matched)

        mood_cat = self._mood_to_preferred_category(mood)
        if mood_cat and mood_cat != category:
            matched = self._search_and_validate(component_lib, mood_cat, n)
            if matched:
                return ("group", matched)

        return (None, None)

    def _search_and_validate(self, component_lib, category: str, n: int) -> str | None:
        candidates = component_lib.search(type="group", category=category, node_count=n, limit=5)
        if not candidates:
            candidates = component_lib.search(type="group", category=category, limit=5)
        if not candidates:
            return None
        for cand in candidates:
            xml_parts = component_lib.load_xml(cand["id"])
            if xml_parts and "group" in xml_parts:
                return category
        return None
