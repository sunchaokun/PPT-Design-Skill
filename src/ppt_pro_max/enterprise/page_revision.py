"""PageRevisionEngine — page-level CRUD operations for Enterprise Pipeline."""

from __future__ import annotations

import os
from dataclasses import dataclass
from typing import Any

from pptx import Presentation

from ppt_pro_max.enterprise.slide_utils import remove_slide


@dataclass
class PageOp:
    page: int
    action: str
    target: int | None = None
    insert_after: bool = True


def parse_pages(pages_str: str, num_slides: int = 0) -> list[PageOp]:
    ops: list[PageOp] = []
    for token in pages_str.split(","):
        token = token.strip()
        if not token:
            continue

        try:
            if "<>" in token:
                parts = token.split("<>")
                if len(parts) != 2:
                    raise ValueError(f"无效的交换格式: {token}")
                a, b = int(parts[0]), int(parts[1])
                ops.append(PageOp(page=a, action="swap", target=b))
            elif ">" in token and not token.startswith(">"):
                parts = token.split(">")
                if len(parts) != 2:
                    raise ValueError(f"无效的移动格式: {token}")
                a, b = int(parts[0]), int(parts[1])
                ops.append(PageOp(page=a, action="move", target=b))
            elif token.startswith("+"):
                page = int(token[1:])
                ops.append(PageOp(page=page, action="insert", insert_after=True))
            elif token.startswith("-") and len(token) > 1 and token[1:].isdigit():
                page = int(token[1:])
                ops.append(PageOp(page=page, action="delete"))
            elif "-" in token and not token.startswith("-"):
                parts = token.split("-", 1)
                if len(parts) != 2:
                    raise ValueError(f"无效的范围格式: {token}")
                start, end = int(parts[0]), int(parts[1])
                for p in range(start, end + 1):
                    ops.append(PageOp(page=p, action="modify"))
            elif token.isdigit():
                ops.append(PageOp(page=int(token), action="modify"))
            else:
                raise ValueError(f"无法识别的页码操作: {token}")
        except ValueError as e:
            if "无法识别" in str(e) or "无效" in str(e) or "超出范围" in str(e):
                raise
            raise ValueError(f"无效的页码操作格式: {token}") from e

    if num_slides > 0:
        for op in ops:
            if op.action in ("modify", "delete", "move", "swap"):
                if op.page < 1 or op.page > num_slides:
                    raise ValueError(f"页码 {op.page} 超出范围（1-{num_slides}）")
            if op.action in ("move", "swap") and op.target is not None:
                if op.target < 1 or op.target > num_slides:
                    raise ValueError(f"目标页码 {op.target} 超出范围（1-{num_slides}）")
            if op.action == "insert":
                if op.page < 1 or op.page > num_slides:
                    raise ValueError(f"插入位置 {op.page} 超出范围（1-{num_slides}）")

    return ops


def compute_target_sequence(
    num_slides: int,
    ops: list[PageOp],
) -> tuple[list[int], dict[int, list[PageOp]]]:
    sequence = list(range(num_slides))
    new_slides: dict[int, list[PageOp]] = {}

    swap_ops = [o for o in ops if o.action == "swap"]
    for op in swap_ops:
        idx_a = op.page - 1
        idx_b = op.target - 1
        if idx_a in sequence and idx_b in sequence:
            pos_a = sequence.index(idx_a)
            pos_b = sequence.index(idx_b)
            sequence[pos_a], sequence[pos_b] = sequence[pos_b], sequence[pos_a]

    move_ops = [o for o in ops if o.action == "move"]
    for op in move_ops:
        idx_source = op.page - 1
        idx_target = op.target - 1
        if idx_source == idx_target:
            continue
        if idx_source in sequence and idx_target in sequence:
            pos_source = sequence.index(idx_source)
            sequence.pop(pos_source)
            pos_target = sequence.index(idx_target)
            sequence.insert(pos_target + 1, idx_source)

    delete_ops = sorted([o for o in ops if o.action == "delete"], key=lambda o: o.page)
    for op in delete_ops:
        idx = op.page - 1
        if idx in sequence:
            sequence.remove(idx)

    insert_ops = sorted([o for o in ops if o.action == "insert"], key=lambda o: o.page)
    for op in insert_ops:
        idx_original = op.page - 1
        if idx_original in sequence:
            pos = sequence.index(idx_original)
            insert_pos = pos + 1
        else:
            insert_pos = len(sequence)
        new_slides.setdefault(insert_pos, []).append(op)

    return sequence, new_slides


class PageRevisionEngine:

    def __init__(self, template_path: str):
        self._template_path = template_path

    def revise(
        self,
        page_ops: list[str],
        output_path: str | None = None,
    ) -> dict[str, Any]:
        prs_src = Presentation(self._template_path)
        num_slides = len(prs_src.slides)

        if not page_ops:
            out = output_path or self._default_output()
            self._save_as(prs_src, out)
            return {"num_slides": num_slides, "output_path": out, "ops_applied": []}

        all_ops: list[PageOp] = []
        for op_str in page_ops:
            all_ops.extend(parse_pages(op_str, num_slides))

        sequence, new_slides = compute_target_sequence(num_slides, all_ops)

        prs_new = Presentation(self._template_path)
        self._rebuild(prs_new, prs_src, sequence, new_slides)

        out = output_path or self._default_output()
        self._save_as(prs_new, out)

        return {
            "num_slides": len(prs_new.slides),
            "output_path": out,
            "ops_applied": [f"{o.action}({o.page})" for o in all_ops],
        }

    def _default_output(self) -> str:
        base = os.path.splitext(self._template_path)[0]
        return f"{base}_revised.pptx"

    def _rebuild(
        self,
        prs_new: Presentation,
        prs_src: Presentation,
        sequence: list[int],
        new_slides: dict[int, list[PageOp]],
    ) -> None:
        slides_src = list(prs_src.slides)
        layout_map = {layout.name: layout for layout in prs_new.slide_layouts}

        while len(prs_new.slides) > 0:
            remove_slide(prs_new, 0)

        existing_slides: list[tuple[int, PageOp | None]] = []
        for i, src_idx in enumerate(sequence):
            existing_slides.append((i, None))

        insert_entries: list[tuple[int, list[PageOp]]] = sorted(new_slides.items())
        for insert_pos, op_list in insert_entries:
            for offset, op in enumerate(op_list):
                existing_slides.insert(insert_pos + offset, (insert_pos, op))

        for entry in existing_slides:
            pos, op = entry
            if op is None:
                if pos < len(sequence):
                    src_idx = sequence[pos]
                    src_slide = slides_src[src_idx]
                    layout = self._match_layout(src_slide, prs_new, layout_map)
                    new_slide = prs_new.slides.add_slide(layout)
                    self._copy_content(src_slide, new_slide)
                else:
                    layout = prs_new.slide_layouts[0]
                    new_slide = prs_new.slides.add_slide(layout)
            else:
                insert_layout = prs_new.slide_layouts[-1]
                for layout in prs_new.slide_layouts:
                    if "title" in layout.name.lower() and "blank" not in layout.name.lower():
                        insert_layout = layout
                        break
                new_slide = prs_new.slides.add_slide(insert_layout)
                if new_slide.shapes.title:
                    new_slide.shapes.title.text = f"New Page (after page {op.page})"

    def _match_layout(self, src_slide, prs_new, layout_map):
        layout_name = src_slide.slide_layout.name
        if layout_name in layout_map:
            return layout_map[layout_name]
        for layout in prs_new.slide_layouts:
            if "title" in layout.name.lower() and "blank" not in layout.name.lower():
                return layout
        return prs_new.slide_layouts[0]

    def _copy_content(self, src_slide, dst_slide) -> None:
        from lxml import etree
        from pptx.oxml.ns import qn

        src_cSld = src_slide._element.find(qn("p:cSld"))
        dst_cSld = dst_slide._element.find(qn("p:cSld"))

        if src_cSld is not None and dst_cSld is not None:
            parent = dst_cSld.getparent()
            parent.remove(dst_cSld)
            new_cSld = etree.fromstring(etree.tostring(src_cSld))
            parent.append(new_cSld)

            dst_slide._element.remove(dst_slide._element.find(qn("p:cSld")))
            insert_pos = 0
            children = list(dst_slide._element)
            for ci, child in enumerate(children):
                if child.tag == qn("p:clrMapOvr") or child.tag == qn("p:transition"):
                    insert_pos = ci
                    break
            dst_slide._element.insert(insert_pos, new_cSld)

        src_timing = src_slide._element.find(qn("p:timing"))
        if src_timing is not None:
            dst_timing = dst_slide._element.find(qn("p:timing"))
            if dst_timing is not None:
                dst_slide._element.remove(dst_timing)
            new_timing = etree.fromstring(etree.tostring(src_timing))
            cSld_elem = dst_slide._element.find(qn("p:cSld"))
            if cSld_elem is not None:
                cSld_idx = list(dst_slide._element).index(cSld_elem)
                dst_slide._element.insert(cSld_idx + 1, new_timing)

        src_transition = src_slide._element.find(qn("p:transition"))
        if src_transition is not None:
            dst_transition = dst_slide._element.find(qn("p:transition"))
            if dst_transition is not None:
                dst_slide._element.remove(dst_transition)
            new_transition = etree.fromstring(etree.tostring(src_transition))
            cSld_elem = dst_slide._element.find(qn("p:cSld"))
            if cSld_elem is not None:
                cSld_idx = list(dst_slide._element).index(cSld_elem)
                dst_slide._element.insert(cSld_idx + 1, new_transition)

    def _recover_dropped_content(self, dst_slide, copied_texts: list[str]) -> None:
        existing = set()
        for shape in dst_slide.placeholders:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        existing.add(para.text.strip())
        if dst_slide.shapes.title and dst_slide.shapes.title.text.strip():
            existing.add(dst_slide.shapes.title.text.strip())

        dropped = [t for t in copied_texts if t not in existing]
        if not dropped:
            return

        from pptx.util import Inches, Pt
        left = Inches(0.9)
        top = Inches(4.0)
        width = Inches(7.0)
        height = Inches(2.0)

        txBox = dst_slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, text in enumerate(dropped):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(11)

    def _save_as(self, prs: Presentation, path: str) -> None:
        out_dir = os.path.dirname(path)
        if out_dir:
            os.makedirs(out_dir, exist_ok=True)
        prs.save(path)
