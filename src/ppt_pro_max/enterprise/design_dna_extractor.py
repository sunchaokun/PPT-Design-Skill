"""DesignDNAExtractor — zero-loss design DNA extraction from .pptx.

Two extraction modes:
1. **Template mode** (has template): Clone XML + locate text zones → clone_and_patch
   - Cover: patch title/subtitle only
   - TOC: patch directory text only
   - Content: patch title + body, optionally inject GroupShape/SmartArt
   - Transition: patch title only
   - Back cover: patch author/details only
   → 100% visual fidelity, zero design loss

2. **VI mode** (no template): Extract brand identity elements → BrandSpec
   - Full color palette (theme + actual usage)
   - Font scheme (heading/body/CJK + sizes)
   - Logo image
   - Decorative elements (GroupShape XML)
   - Page structure patterns
   → Feed into Build mode for fine-grained design

Key data structures:
- TextZone: a single editable text region in a slide (XML-level located)
- SlideDNA: one slide's complete design DNA (XML clone + text zones + page type)
- DesignDNA: entire PPT's design DNA (all slides + VI + theme)
"""

from __future__ import annotations

import copy
import os
import re
import tempfile
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from ppt_pro_max.enterprise.brand_spec import BrandSpec

_GOAL_TO_PAGE_TYPE = {
    "hook": "cover",
    "overview": "toc",
    "section": "transition",
    "problem": "content",
    "solution": "content",
    "features": "content",
    "data": "content",
    "code": "content",
    "exercise": "content",
    "content": "content",
    "cta": "back_cover",
}

_LAYOUT_AUTO_RULES = {
    "cards": {"min_bullets": 2, "max_bullets": 4, "needs_cards": True},
    "bullets": {"min_bullets": 3, "max_bullets": 8},
}


@dataclass
class PagePlan:
    page_type: str
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    image_path: str | None = None
    layout: str = "auto"
    cards: list[dict] = field(default_factory=list)
    section_number: int = 0
    notes: str = ""
    component_type: str | None = None
    component_category: str | None = None
    mood: str | None = None


@dataclass
class ExpandResult:
    output_path: str
    template_mapping: list[dict] = field(default_factory=list)
    total_slides: int = 0
    cloned_count: int = 0
    deleted_originals: list[int] = field(default_factory=list)

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

_SYS_CLR_FALLBACK = {
    "windowText": "#000000",
    "WindowText": "#000000",
    "window": "#FFFFFF",
    "Window": "#FFFFFF",
}

_COVER_KEYWORDS = {"汇报", "报告", "方案", "规划", "战略", "年度", "总结", "展望", "report", "plan", "strategy", "annual"}
_TOC_KEYWORDS = {"目录", "contents", "agenda", "议程", "概览", "overview", "壹", "贰", "叁", "肆", "伍"}
_CTA_KEYWORDS = {"谢谢", "感谢", "thank", "thanks", "联系", "contact", "观看", "聆听", "谢", "谢观", "谢观看"}
_TRANSITION_KEYWORDS = {"part", "章节", "第.*章", "section"}

_BACK_COVER_PATTERNS = {
    "large_title_count_ge_3": True,
    "body_is_decorative": True,
}


@dataclass
class TextZone:
    zone_id: str
    role: str  # title | subtitle | body | badge | decoration | logo_text | page_number
    text: str
    font_name: str | None = None
    font_size_pt: float | None = None
    bold: bool = False
    color_hex: str | None = None
    xpath: str = ""
    run_index: int = 0
    para_index: int = 0
    shape_name: str = ""
    bounds: tuple[float, float, float, float] = (0.0, 0.0, 0.0, 0.0)


@dataclass
class SlideDNA:
    slide_index: int
    page_type: str  # cover | toc | content | transition | back_cover | unknown
    slide_xml: bytes  # full <p:sld> XML — zero loss
    rels_xml: bytes  # slide relationships XML
    rels_map: dict[str, str] = field(default_factory=dict)  # rId → target
    text_zones: list[TextZone] = field(default_factory=list)
    image_refs: list[dict[str, Any]] = field(default_factory=list)
    group_xmls: list[bytes] = field(default_factory=list)
    smartart_refs: list[dict[str, Any]] = field(default_factory=list)
    layout_name: str = ""
    layout_index: int = 0
    has_background_image: bool = False
    background_colors: list[str] = field(default_factory=list)
    notes_text: str = ""
    brand_spec: BrandSpec | None = None


@dataclass
class DesignDNA:
    source_path: str
    slide_width_emu: int = 0
    slide_height_emu: int = 0
    slides: list[SlideDNA] = field(default_factory=list)
    theme_xml: bytes = b""
    color_palette: dict[str, str] = field(default_factory=dict)
    font_scheme: dict[str, str] = field(default_factory=dict)
    cjk_font_scheme: dict[str, str] = field(default_factory=dict)
    actual_colors: dict[str, int] = field(default_factory=dict)
    actual_fonts: dict[str, int] = field(default_factory=dict)
    actual_font_sizes: dict[int, int] = field(default_factory=dict)
    actual_bg_colors: dict[str, int] = field(default_factory=dict)
    logo_blob: bytes | None = None
    logo_content_type: str = ""
    decorative_groups: list[dict[str, Any]] = field(default_factory=list)
    brand_spec: BrandSpec | None = None
    image_blobs: dict[str, bytes] = field(default_factory=dict)


class DesignDNAExtractor:

    def __init__(self, temp_dir: str | None = None):
        self._temp_dir = temp_dir or os.path.join(tempfile.gettempdir(), "ppt-dna-images")
        os.makedirs(self._temp_dir, exist_ok=True)

    def extract(self, pptx_path: str) -> DesignDNA:
        dna = DesignDNA(source_path=pptx_path)

        prs = Presentation(pptx_path)
        dna.slide_width_emu = prs.slide_width
        dna.slide_height_emu = prs.slide_height

        dna.theme_xml = self._extract_theme_xml(prs)
        dna.color_palette, dna.font_scheme, dna.cjk_font_scheme = self._extract_theme_scheme(prs)

        with zipfile.ZipFile(pptx_path) as z:
            self._extract_slides_from_zip(z, prs, dna)
            self._extract_all_images(z, dna)
            self._extract_logo(z, prs, dna)

        self._extract_actual_usage(prs, dna)
        self._classify_page_types(dna)
        self._extract_decorative_groups(prs, dna)

        dna.brand_spec = self._build_brand_spec(prs, dna)

        if dna.brand_spec:
            for slide_dna in dna.slides:
                slide_dna.brand_spec = dna.brand_spec

        self._last_dna = dna
        return dna

    def clone_and_patch(
        self,
        dna: DesignDNA,
        patches: list[dict[str, Any]],
        output_path: str,
    ) -> str:
        prs = Presentation(dna.source_path)

        total_original = len(prs.slides)
        total_target = len(patches)

        for idx in range(min(total_original, total_target)):
            slide = prs.slides[idx]
            patch = patches[idx]
            self._patch_slide(slide, dna.slides[idx], patch)

        if total_target > total_original:
            last_slide_dna = dna.slides[-1] if dna.slides else None
            content_dna = None
            for s in dna.slides:
                if s.page_type == "content":
                    content_dna = s
                    break

            for idx in range(total_original, total_target):
                template_dna = content_dna or last_slide_dna
                if template_dna:
                    self._inject_new_slide(prs, template_dna, patches[idx], dna)

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)
        return output_path

    def plan_pages(
        self,
        story_plan: dict[str, Any],
        page_contents: list[dict[str, Any]] | None = None,
    ) -> list[PagePlan]:
        pages = story_plan.get("pages", [])
        plans: list[PagePlan] = []
        section_counter = 0

        for i, page in enumerate(pages):
            goal = page.get("goal", "content")
            page_type = _GOAL_TO_PAGE_TYPE.get(goal, "content")

            if page_type == "transition":
                section_counter += 1

            layout = page.get("layout", "auto")
            if layout == "auto" and page_type == "content":
                layout = self._auto_layout(page)

            plan = PagePlan(
                page_type=page_type,
                title=page.get("title", ""),
                subtitle=page.get("subtitle", ""),
                bullets=page.get("bullets", []),
                image_path=page.get("image"),
                layout=layout,
                cards=page.get("cards", []),
                section_number=section_counter,
                notes=page.get("notes", ""),
            )

            if page_contents and i < len(page_contents):
                pc = page_contents[i]
                if pc.get("bullets"):
                    plan.bullets = pc["bullets"]
                if pc.get("cards"):
                    plan.cards = pc["cards"]
                pc_layout = pc.get("layout")
                if pc_layout and pc_layout != "auto":
                    plan.layout = pc_layout
                if pc.get("component_type"):
                    plan.component_type = pc["component_type"]
                if pc.get("component_category"):
                    plan.component_category = pc["component_category"]

            if plan.component_type == "group" and plan.component_category:
                _category_to_layout = {
                    "process": "process", "hierarchy": "hierarchy",
                    "swot": "swot", "timeline": "timeline",
                    "cycle": "process", "pyramid": "hierarchy",
                    "matrix": "swot", "infographic": "numbered_row",
                    "comparison": "bullets", "funnel": "process",
                    "radial": "bullets", "chart": "bullets",
                }
                derived_layout = _category_to_layout.get(plan.component_category)
                if derived_layout:
                    plan.layout = derived_layout

            plans.append(plan)

        return plans

    def expand_and_patch(
        self,
        dna: DesignDNA,
        plans: list[PagePlan],
        output_path: str,
        component_lib=None,
        brand_spec: BrandSpec | None = None,
    ) -> ExpandResult:
        prs = Presentation(dna.source_path)

        template_index: dict[str, list[int]] = {}
        for slide_dna in dna.slides:
            pt = slide_dna.page_type
            template_index.setdefault(pt, []).append(slide_dna.slide_index)

        used_originals: set[int] = set()
        use_counts: dict[int, int] = {}
        template_mapping: list[dict] = []
        cloned_count = 0
        plan_to_physical: list[int] = []
        card_page_idx = 0

        for plan_idx, plan in enumerate(plans):
            page_type = plan.page_type
            candidates = template_index.get(page_type, [])

            if page_type == "content" and plan.cards:
                card_page_idx += 1

            if not candidates:
                for fallback_type in ("content", "transition", "cover"):
                    candidates = template_index.get(fallback_type, [])
                    if candidates:
                        break

            if not candidates:
                template_mapping.append({
                    "plan_index": plan_idx,
                    "page_type": page_type,
                    "template_slide_index": -1,
                    "is_clone": True,
                    "status": "skipped_no_template",
                })
                continue

            template_slide_idx = self._select_best_template(candidates, plan, dna, used_originals, use_counts)
            use_counts[template_slide_idx] = use_counts.get(template_slide_idx, 0) + 1

            if template_slide_idx not in used_originals:
                used_originals.add(template_slide_idx)
                slide = prs.slides[template_slide_idx]
                slide_dna = dna.slides[template_slide_idx]
                self._patch_slide(slide, slide_dna, self._plan_to_patch(plan))
                if self.needs_content_rerender(slide_dna, plan):
                    self.rerender_content_zone(slide, slide_dna, plan, dna.brand_spec, card_page_idx, component_lib, brand_spec_override=brand_spec)
                template_mapping.append({
                    "plan_index": plan_idx,
                    "page_type": page_type,
                    "template_slide_index": template_slide_idx,
                    "is_clone": False,
                })
                plan_to_physical.append(template_slide_idx)
            else:
                template_dna = dna.slides[template_slide_idx]
                self._inject_new_slide(prs, template_dna, self._plan_to_patch(plan), dna)
                new_slide_idx = len(prs.slides) - 1
                new_slide = prs.slides[new_slide_idx]
                if self.needs_content_rerender(template_dna, plan):
                    self.rerender_content_zone(new_slide, template_dna, plan, dna.brand_spec, card_page_idx, component_lib, brand_spec_override=brand_spec)
                cloned_count += 1
                template_mapping.append({
                    "plan_index": plan_idx,
                    "page_type": page_type,
                    "template_slide_index": template_slide_idx,
                    "is_clone": True,
                })
                plan_to_physical.append(new_slide_idx)

        unused_originals = sorted(
            set(range(len(dna.slides))) - used_originals,
            reverse=True,
        )

        for idx in unused_originals:
            from ppt_pro_max.enterprise.slide_utils import remove_slide
            remove_slide(prs, idx)
            plan_to_physical = [
                (i - 1 if i > idx else i) if i != idx else -1
                for i in plan_to_physical
            ]

        plan_to_physical = [i for i in plan_to_physical if i >= 0]

        self._reorder_slides(prs, plan_to_physical)
        plan_to_physical = list(range(len(plan_to_physical)))

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)

        self._post_save_inject_groups(output_path, prs, plans, plan_to_physical, dna, component_lib, brand_spec=brand_spec)

        if brand_spec and brand_spec.fonts:
            self._replace_theme_fonts_in_pptx(output_path, brand_spec)

        try:
            from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
            gate = DeliveryGate()
            report = gate.check(output_path, dna, plans)
            if report.fatals:
                for _round in range(2):
                    gate.auto_fix(output_path, dna, plans, report)
                    report = gate.check(output_path, dna, plans)
                    if not report.fatals:
                        break
        except Exception:
            pass

        unused_originals_result = sorted(
            set(range(len(dna.slides))) - used_originals,
        )

        return ExpandResult(
            output_path=output_path,
            template_mapping=template_mapping,
            total_slides=len(plan_to_physical),
            cloned_count=cloned_count,
            deleted_originals=unused_originals_result,
        )

    def _select_best_template(
        self,
        candidates: list[int],
        plan: PagePlan,
        dna: DesignDNA,
        used_originals: set[int],
        use_counts: dict[int, int],
    ) -> int:
        if len(candidates) == 1:
            return candidates[0]

        scored = []
        for idx in candidates:
            slide_dna = dna.slides[idx]
            score = 0
            zone_count = len([z for z in slide_dna.text_zones if z.role in ("body", "subtitle")])
            target_count = len(plan.bullets or []) + len(plan.cards or [])
            score -= abs(zone_count - target_count) * 10
            if idx not in used_originals:
                score += 50
            deco_count = len([z for z in slide_dna.text_zones if z.role == "decoration"])
            score += deco_count * 5
            score -= use_counts.get(idx, 0) * 20
            scored.append((score, idx))
        scored.sort(reverse=True)
        return scored[0][1] if scored else candidates[0]

    def _post_save_inject_groups(self, output_path, prs, plans, plan_to_physical, dna, component_lib=None, brand_spec: BrandSpec | None = None) -> None:
        p_ns = _NS["p"]
        a_ns = _NS["a"]
        needs_fix = []
        for i, plan in enumerate(plans):
            if plan.page_type != "content":
                continue
            if i >= len(plan_to_physical):
                continue
            phys_idx = plan_to_physical[i]
            if phys_idx < 0 or phys_idx >= len(prs.slides):
                continue
            slide = prs.slides[phys_idx]
            sp_tree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
            grps = sp_tree.findall(f"{{{p_ns}}}grpSp")
            has_content_group = False
            for grp in grps:
                grp_texts = grp.findall(f".//{{{a_ns}}}t")
                if any(t.text and t.text.strip() for t in grp_texts):
                    has_content_group = True
                    break
            if not has_content_group:
                needs_fix.append((phys_idx, plan, slide))

            if not needs_fix:
                return

            import re
            import zipfile
            import shutil
            from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
            from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
            from ppt_pro_max.enterprise.brand_spec import BrandSpec

            lib = component_lib
            if lib is None:
                db_path = find_db_path()
                if db_path is None:
                    return
                lib = ComponentLibrary(db_path)
            renderer = ComponentRenderer()

            body_color = self._extract_body_color_from_dna(dna)
            accent_color = self._extract_accent_color_from_dna(dna, body_color)

            slide_patches = {}
            for phys_idx, plan, slide in needs_fix:
                layout = plan.layout if plan.layout != "auto" else self._auto_layout_from_plan(plan)
                category_map = {"chevron": "process", "bullets": "infographic", "cards": "infographic",
                                "process": "process", "hierarchy": "hierarchy", "two_col": "infographic",
                                "numbered_row": "process", "swot": "swot", "timeline": "process",
                                "table": "infographic", "code": "infographic", "chart": "infographic"}
                category = plan.component_category or category_map.get(layout, "infographic")

                content_bounds = self._compute_default_content_bounds(plan)
                bullets = plan.bullets or []
                cards = plan.cards or []
                texts = bullets if bullets else [c.get("title", "") + " " + c.get("body", c.get("text", "")) for c in cards]
                if not texts:
                    texts = [plan.title or ""]

                candidates = lib.search(type="group", category=category, node_count=len(texts), limit=50)
                if not candidates:
                    candidates = lib.search(type="group", category=category, limit=50)

                component = None
                xml_parts = None
                best_shape_count = 0
                for cand in candidates:
                    xp = lib.load_xml(cand["id"])
                    if not xp or "group" not in xp:
                        continue
                    gx = xp["group"]
                    if isinstance(gx, bytes):
                        gx = gx.decode("utf-8")
                    if re.search(r"c:chart|graphicFrame", gx):
                        continue
                    try:
                        grp_root = etree.fromstring(gx)
                        shape_count = len(grp_root.findall(f"{{{p_ns}}}sp")) + len(grp_root.findall(f"{{{p_ns}}}cxnSp"))
                    except Exception:
                        shape_count = 0
                    if shape_count > best_shape_count:
                        best_shape_count = shape_count
                        component = cand
                        xml_parts = xp

                if component is None or xml_parts is None:
                    fallback_body = body_color or "#333333"
                    fallback_accent = accent_color or "#6096E6"
                    if layout == "process" and plan.bullets:
                        self._render_chevron_group(slide, plan.bullets, *content_bounds, fallback_body, fallback_accent)
                    elif layout == "cards" and plan.cards:
                        self._render_numbered_card_group(slide, plan.cards, *content_bounds, fallback_body, fallback_accent)
                    elif plan.bullets:
                        self._render_numbered_row_group(slide, plan.bullets, *content_bounds, fallback_body, fallback_accent)
                    continue

                effective_brand = brand_spec
                if effective_brand is None:
                    accent_color_val = accent_color or body_color or "#6096E6"
                    body_color_val = body_color or "#333333"
                    effective_brand = BrandSpec(colors={
                        "primary": accent_color_val,
                        "accent": accent_color_val,
                        "muted": self._lighten_color(accent_color_val, 0.85),
                        "foreground": body_color_val,
                        "on-primary": "#FFFFFF",
                        "background": "#FFFFFF",
                    })

                filled = renderer._fill_group_data(xml_parts, texts)
                styled = renderer._apply_brand_colors(filled, effective_brand)
                if effective_brand and effective_brand.fonts:
                    styled["group"] = ComponentRenderer._replace_group_fonts(styled["group"], effective_brand)

                grp_xml_bytes = styled["group"]
                if isinstance(grp_xml_bytes, str):
                    grp_xml_bytes = grp_xml_bytes.encode("utf-8")

                grp_elem = etree.fromstring(grp_xml_bytes)

                next_id = 2
                for cNv in slide._element.iter(f"{{{p_ns}}}cNvPr"):
                    try:
                        next_id = max(next_id, int(cNv.get("id", "1")) + 1)
                    except (ValueError, TypeError):
                        pass
                for cNv in grp_elem.iter(f"{{{p_ns}}}cNvPr"):
                    cNv.set("id", str(next_id))
                    cNv.set("name", "Component %d" % next_id)
                    next_id += 1

                renderer._denormalize_coordinates(grp_elem, *content_bounds)

                grp_sp_xml = etree.tostring(grp_elem, encoding="unicode", pretty_print=True)

                partname = str(slide.part.partname).lstrip("/")

                with zipfile.ZipFile(output_path, "r") as zin:
                    slide_xml_bytes = zin.read(partname)

                slide_xml_str = slide_xml_bytes.decode("utf-8")
                insert_point = "</p:spTree>"
                if insert_point in slide_xml_str:
                    slide_xml_str = slide_xml_str.replace(insert_point, grp_sp_xml + insert_point, 1)
                    slide_patches[partname] = slide_xml_str.encode("utf-8")

            if slide_patches:
                tmp_path = output_path + ".tmp"
                with zipfile.ZipFile(output_path, "r") as zin:
                    with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
                        for item in zin.infolist():
                            data = zin.read(item.filename)
                            if item.filename in slide_patches:
                                data = slide_patches[item.filename]
                            zout.writestr(item, data)
                shutil.move(tmp_path, output_path)

    def _replace_theme_fonts_in_pptx(self, pptx_path: str, brand_spec: BrandSpec) -> None:
        import re
        import zipfile
        import shutil

        fonts = brand_spec.fonts
        if not fonts:
            return

        heading = fonts.get("heading", "")
        body = fonts.get("body", "")
        cjk_heading = fonts.get("cjk_heading", heading)
        cjk_body = fonts.get("cjk_body", body)

        if not any([heading, body, cjk_heading, cjk_body]):
            return

        replacements = {}
        if heading:
            replacements["+mj-lt"] = heading
        if body:
            replacements["+mn-lt"] = body
        if cjk_heading:
            replacements["+mj-ea"] = cjk_heading
        if cjk_body:
            replacements["+mn-ea"] = cjk_body

        try:
            tmp_path = pptx_path + ".tmp"
            changed = False
            with zipfile.ZipFile(pptx_path, "r") as zin:
                with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
                    for item in zin.infolist():
                        data = zin.read(item.filename)
                        if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                            text = data.decode("utf-8")
                            for old, new in replacements.items():
                                if old in text:
                                    text = text.replace(f'typeface="{old}"', f'typeface="{new}"')
                                    changed = True
                            if body:
                                for non_template in ["Lato Light", "Nevis", "Calibri", "Segoe UI", "Roboto", "Open Sans", "Montserrat", "Raleway", "Poppins", "Inter"]:
                                    if non_template in text:
                                        text = text.replace(f'typeface="{non_template}"', f'typeface="{body}"')
                                        changed = True
                            data = text.encode("utf-8")
                        zout.writestr(item, data)
            if changed:
                shutil.move(tmp_path, pptx_path)
            else:
                import os
                os.remove(tmp_path)
        except Exception:
            pass

    @staticmethod
    def _reorder_slides(prs, plan_to_physical: list[int]) -> None:
        n = len(plan_to_physical)
        if n == 0 or plan_to_physical == list(range(n)):
            return
        sldIdLst = prs.slides._sldIdLst
        elements = list(sldIdLst)
        if max(plan_to_physical) >= len(elements):
            return
        reordered = [elements[i] for i in plan_to_physical]
        for child in list(sldIdLst):
            sldIdLst.remove(child)
        for elem in reordered:
            sldIdLst.append(elem)

    def _plan_to_patch(self, plan: PagePlan) -> dict[str, Any]:
        patch: dict[str, Any] = {}
        if plan.title:
            patch["title"] = plan.title
        if plan.subtitle:
            patch["subtitle"] = plan.subtitle
        if plan.bullets:
            patch["bullets"] = plan.bullets
        if plan.image_path:
            patch["image_path"] = plan.image_path
        return patch

    def _auto_layout(self, page: dict[str, Any]) -> str:
        cards = page.get("cards", [])
        if cards:
            return "cards"

        bullets = page.get("bullets", [])
        n = len(bullets)

        all_text = " ".join(bullets).lower()
        if any(kw in all_text for kw in ("步骤", "阶段", "process", "流程", "step", "phase", "pipeline", "workflow")):
            return "process"
        if any(kw in all_text for kw in ("架构", "层级", "hierarchy", "组织", "org chart", "reporting")):
            return "hierarchy"
        if any(kw in all_text for kw in ("里程碑", "timeline", "时间线", "milestone", "roadmap", "路线图")):
            return "timeline"
        if any(kw in all_text for kw in ("swot", "优势", "劣势", "quadrant", "matrix", "对比")):
            return "swot"

        if 2 <= n <= 4 and cards:
            return "cards"
        if n > 8:
            return "bullets"

        return "bullets"

    def needs_content_rerender(self, slide_dna: SlideDNA, plan: PagePlan) -> bool:
        if plan.page_type != "content":
            return False

        if plan.component_type:
            return True

        if plan.cards:
            return True

        if plan.bullets and len(plan.bullets) > 0:
            return True

        body_zones = [z for z in slide_dna.text_zones if z.role == "body"]
        if not body_zones:
            return True

        if plan.layout and plan.layout != "auto" and plan.layout != "bullets":
            return True

        subtitle_zones = [z for z in slide_dna.text_zones if z.role == "subtitle"]
        page_number_zones = [z for z in slide_dna.text_zones if z.role == "page_number"]
        if subtitle_zones or page_number_zones:
            return True

        if plan.bullets and len(plan.bullets) != len(body_zones):
            return True

        if plan.bullets:
            total_text_len = sum(len(b) for b in plan.bullets)
            zone_capacity = len(body_zones) * 30
            if total_text_len > zone_capacity:
                return True

        return False

    def rerender_content_zone(
        self,
        slide,
        slide_dna: SlideDNA,
        plan: PagePlan,
        brand_spec: BrandSpec | None = None,
        card_page_index: int = 0,
        component_lib=None,
        brand_spec_override: BrandSpec | None = None,
    ) -> None:
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer

        all_content_zones = [
            z for z in slide_dna.text_zones
            if z.role in ("body", "subtitle", "page_number")
            and z.shape_name not in {
                tz.shape_name for tz in slide_dna.text_zones if tz.role == "title"
            }
        ]
        body_zones = [z for z in slide_dna.text_zones if z.role == "body"]
        if not all_content_zones:
            all_content_zones = body_zones

        title_zones_list = [z for z in slide_dna.text_zones if z.role == "title"]
        title_y_bottom = max(z.bounds[1] + z.bounds[3] for z in title_zones_list) if title_zones_list else 1.5

        if not all_content_zones:
            content_y_min = title_y_bottom + 0.3
            content_y_max = 7.2
            margin_left = 0.9
            content_bounds = (margin_left, content_y_min, 13.333 - margin_left - 0.9, content_y_max - content_y_min)
        else:
            content_y_min = min(z.bounds[1] for z in all_content_zones)
            if title_y_bottom < content_y_min:
                content_y_min = title_y_bottom
            content_y_max = max(z.bounds[1] + z.bounds[3] for z in all_content_zones)

            if plan.page_type == "content":
                content_y_min = title_y_bottom
                content_y_max = 7.5

        self._clear_text_only_in_zone(slide, slide_dna, content_y_min, content_y_max)

        if plan.page_type == "content" and plan.subtitle and plan.bullets and len(plan.bullets) <= 3 and not plan.cards:
            pass
        elif plan.page_type == "content" and plan.layout in ("cards", "bullets", "numbered_row", "two_col", "hierarchy", "swot", "timeline", "auto"):
            self._selective_clear_decorations(slide, slide_dna, plan, keep_levels={0, 1})

        template_card_groups = self._detect_card_groups(slide_dna)

        body_color = self._extract_body_color(slide_dna)
        accent_color = self._extract_accent_color(slide_dna, body_color)

        if plan.component_type:
            self._selective_clear_decorations(slide, slide_dna, plan, keep_levels={0})
            content_bounds = self._compute_content_bounds(slide_dna, plan)
            effective_brand = brand_spec_override or brand_spec or slide_dna.brand_spec
            if effective_brand is None:
                effective_brand = BrandSpec(source="design_dna")
            self._apply_dna_fonts_to_brand(effective_brand, dna=None, slide_dna=slide_dna)
            comp_precision = PrecisionRenderer(brand_spec=effective_brand)
            self._render_professional_content(
                slide, comp_precision, plan, content_bounds,
                body_color, accent_color, card_page_index, component_lib,
                brand_spec=effective_brand,
            )
            return

        use_template_cards = False

        if use_template_cards:
            self._patch_card_content(slide, slide_dna, plan)
            return

        if template_card_groups and plan.cards and len(template_card_groups) == len(plan.cards):
            pass
        elif template_card_groups and plan.bullets:
            items = [{"title": "", "body": b} for b in plan.bullets]
            if len(items) > len(template_card_groups):
                pass
            elif all(len(b) <= 15 for b in plan.bullets) and len(items) <= len(template_card_groups):
                adapted = PagePlan(
                    page_type=plan.page_type,
                    title=plan.title,
                    subtitle=plan.subtitle,
                    layout="cards",
                    cards=[{"title": b, "body": ""} for b in plan.bullets],
                    bullets=plan.bullets,
                )
                self._patch_card_content(slide, slide_dna, adapted)
                for i in range(len(items), len(template_card_groups)):
                    group = template_card_groups[i]
                    slide_elem = slide._element
                    patched: set = set()
                    for key in ("subtitle", "body"):
                        zone = group.get(key)
                        if zone:
                            self._clear_zone_text(slide_elem, zone, patched)
                return

        all_content_zones = [
            z for z in slide_dna.text_zones
            if z.role in ("body", "subtitle", "page_number")
            and z.shape_name not in {
                tz.shape_name for tz in slide_dna.text_zones if tz.role == "title"
            }
        ]
        body_zones = [z for z in slide_dna.text_zones if z.role == "body"]
        if not all_content_zones:
            all_content_zones = body_zones

        if not all_content_zones:
            title_zones_list = [z for z in slide_dna.text_zones if z.role == "title"]
            title_y_bottom = max(z.bounds[1] + z.bounds[3] for z in title_zones_list) if title_zones_list else 1.5
            content_bounds = (0.9, title_y_bottom + 0.3, 13.333 - 1.8, 7.2 - title_y_bottom - 0.3)
        else:
            content_y_min = min(z.bounds[1] for z in all_content_zones)
            title_zones_list = [z for z in slide_dna.text_zones if z.role == "title"]
            title_y_bottom = max(z.bounds[1] + z.bounds[3] for z in title_zones_list) if title_zones_list else 0
            if title_y_bottom < content_y_min:
                content_y_min = title_y_bottom
            content_y_max = max(z.bounds[1] + z.bounds[3] for z in all_content_zones)

            min_x = min(z.bounds[0] for z in all_content_zones)
            max_x = max(z.bounds[0] + z.bounds[2] for z in all_content_zones)
            content_bounds = (min_x, content_y_min, max_x - min_x, content_y_max - content_y_min)

        effective_brand = brand_spec_override or brand_spec or slide_dna.brand_spec
        if effective_brand is None:
            effective_brand = BrandSpec(source="design_dna")
        self._apply_dna_fonts_to_brand(effective_brand, dna=None, slide_dna=slide_dna)

        precision = PrecisionRenderer(brand_spec=effective_brand)

        body_color = self._extract_body_color(slide_dna)
        accent_color = self._extract_accent_color(slide_dna, body_color)

        self._render_professional_content(
            slide, precision, plan, content_bounds,
            body_color, accent_color, card_page_index, component_lib,
            brand_spec=effective_brand,
        )

    def _render_professional_content(
        self,
        slide,
        precision,
        plan: PagePlan,
        content_bounds: tuple[float, float, float, float],
        body_color: str | None,
        accent_color: str | None,
        card_page_index: int = 0,
        component_lib=None,
        brand_spec: BrandSpec | None = None,
    ) -> None:
        x, y, w, h = content_bounds
        layout = plan.layout if plan.layout != "auto" else self._auto_layout_from_plan(plan)

        dna_font = self._resolve_dna_font(brand_spec)
        dna_accent = accent_color or self._resolve_dna_accent(brand_spec)
        dna_body = body_color or self._resolve_dna_body(brand_spec)

        tried_component = self._try_component_render(slide, plan, layout, content_bounds, dna_body, dna_accent, card_page_index, component_lib, brand_spec=brand_spec)

        if tried_component:
            p_ns = _NS["p"]
            sp_tree_check = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
            grps_check = sp_tree_check.findall(f"{{{p_ns}}}grpSp")
            if len(grps_check) > 0:
                return
            tried_component = False

        if layout == "process" and plan.bullets:
            self._render_chevron_group(slide, plan.bullets, x, y, w, h, dna_body, dna_accent, dna_font)
        elif layout == "cards" and plan.cards:
            variant = card_page_index % 3
            if variant == 0:
                self._render_numbered_card_group(slide, plan.cards, x, y, w, h, dna_body, dna_accent, dna_font)
            elif variant == 1:
                self._render_icon_card_group(slide, plan.cards, x, y, w, h, dna_body, dna_accent, dna_font)
            else:
                self._render_stat_card_group(slide, plan.cards, x, y, w, h, dna_body, dna_accent, dna_font)
        elif layout == "hierarchy" and plan.bullets:
            self._render_hierarchy_group(slide, plan.bullets, x, y, w, h, dna_body, dna_accent, dna_font)
        elif layout == "swot" and plan.bullets:
            self._render_swot_group(slide, plan.bullets, x, y, w, h, dna_body, dna_accent, dna_font)
        elif layout == "timeline" and plan.bullets:
            self._render_timeline_group(slide, plan.bullets, x, y, w, h, dna_body, dna_accent, dna_font)
        elif plan.bullets:
            n = len(plan.bullets)
            if n >= 5:
                self._render_two_col_group(slide, plan.bullets, x, y, w, h, dna_body, dna_accent, dna_font)
            else:
                self._render_numbered_row_group(slide, plan.bullets, x, y, w, h, dna_body, dna_accent, dna_font)

    def _try_component_render(
        self,
        slide,
        plan: PagePlan,
        layout: str,
        bounds: tuple[float, float, float, float],
        body_color: str | None,
        accent_color: str | None,
        card_page_index: int,
        component_lib=None,
        brand_spec: BrandSpec | None = None,
    ) -> bool:
        try:
            from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
            from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
            from ppt_pro_max.enterprise.brand_spec import BrandSpec

            lib = component_lib
            if lib is None:
                db_path = find_db_path()
                if db_path is None:
                    return False
                lib = ComponentLibrary(db_path)
            try:
                category_map = {
                    "chevron": "process",
                    "bullets": "infographic",
                    "cards": "infographic",
                    "process": "process",
                    "hierarchy": "hierarchy",
                    "two_col": "infographic",
                    "numbered_row": "process",
                    "swot": "swot",
                    "timeline": "timeline",
                    "chart": "chart",
                }
                category = plan.component_category or category_map.get(layout, "infographic")

                bullets = plan.bullets or []
                cards = plan.cards or []
                texts = bullets if bullets else [c.get("title", "") + " " + c.get("body", c.get("text", "")) for c in cards]
                if not texts:
                    texts = [plan.title or ""]

                node_count = len(texts)

                element = {
                    "type": "group",
                    "category": category,
                    "texts": texts,
                    "nodes": node_count,
                    "node_count": node_count,
                    "bounds": bounds,
                    "component_fit": "stretch",
                }

                effective_brand = brand_spec
                if effective_brand is None:
                    if accent_color or body_color:
                        effective_brand = BrandSpec(colors={
                            "primary": accent_color or "#6096E6",
                            "accent": accent_color or "#6096E6",
                            "muted": self._lighten_color(accent_color or "#6096E6", 0.85),
                            "foreground": body_color or "#333333",
                            "on-primary": "#FFFFFF",
                            "background": "#FFFFFF",
                        })

                renderer = ComponentRenderer()
                result = renderer.render_group(slide, element, brand_spec=effective_brand, component_lib=lib)
                return result
            finally:
                lib.close()
        except Exception:
            return False

    def _build_grp_xml(self, children_xml: list[str], x_in: float, y_in: float, w_in: float, h_in: float, ch_x: float = None, ch_y: float = None, ch_w: float = None, ch_h: float = None) -> str:
        p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

        sp_tree = self._get_sp_tree_from_current_slide()
        next_id = 2
        if sp_tree is not None:
            for cNv in sp_tree.iter(f"{{{p_ns}}}cNvPr"):
                try:
                    next_id = max(next_id, int(cNv.get("id", "1")) + 1)
                except (ValueError, TypeError):
                    pass

        x_emu = int(x_in * 914400)
        y_emu = int(y_in * 914400)
        w_emu = int(w_in * 914400)
        h_emu = int(h_in * 914400)
        ch_x_emu = int((ch_x if ch_x is not None else x_in) * 914400)
        ch_y_emu = int((ch_y if ch_y is not None else y_in) * 914400)
        ch_w_emu = int((ch_w if ch_w is not None else w_in) * 914400)
        ch_h_emu = int((ch_h if ch_h is not None else h_in) * 914400)

        grp_id = next_id
        next_id += 1

        inner = "".join(children_xml).replace("{GRP_ID_OFFSET}", str(next_id))

        return (
            f'<p:grpSp xmlns:p="{p_ns}" xmlns:a="{a_ns}" '
            f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            f'<p:nvGrpSpPr><p:cNvPr id="{grp_id}" name="Group {grp_id}"/>'
            f'<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f'<p:grpSpPr><a:xfrm>'
            f'<a:off x="{x_emu}" y="{y_emu}"/>'
            f'<a:ext cx="{w_emu}" cy="{h_emu}"/>'
            f'<a:chOff x="{ch_x_emu}" y="{ch_y_emu}"/>'
            f'<a:chExt cx="{ch_w_emu}" cy="{ch_h_emu}"/>'
            f'</a:xfrm></p:grpSpPr>'
            f'{inner}'
            f'</p:grpSp>'
        )

    def _get_sp_tree_from_current_slide(self):
        return None

    def _sp_xml(self, shape_id: int, name: str, x_in: float, y_in: float, w_in: float, h_in: float,
                 fill_hex: str | None = None, no_fill: bool = False,
                 prst: str = "rect", text: str = "", font_size: int = 12,
                 font_color: str | None = None, bold: bool = False,
                 font_name: str = "Microsoft YaHei", align: str = "left",
                 is_textbox: bool = False) -> str:
        p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        x_emu = int(x_in * 914400)
        y_emu = int(y_in * 914400)
        w_emu = int(w_in * 914400)
        h_emu = int(h_in * 914400)
        sz = font_size * 100

        fill_str = ""
        if fill_hex:
            fill_str = f'<a:solidFill><a:srgbClr val="{fill_hex.lstrip("#")}"/></a:solidFill>'
        elif no_fill:
            fill_str = "<a:noFill/>"

        ln_str = ""
        if fill_hex and prst in ("chevron", "roundRect", "ellipse"):
            ln_str = f'<a:ln><a:solidFill><a:srgbClr val="{fill_hex.lstrip("#")}"/></a:solidFill></a:ln>'

        color_str = ""
        if font_color:
            color_str = f'<a:solidFill><a:srgbClr val="{font_color.lstrip("#")}"/></a:solidFill>'

        bold_attr = ' b="1"' if bold else ''
        align_val = {"left": "l", "center": "ctr", "right": "r"}.get(align, "l")
        txbox_attr = ' txBox="1"' if is_textbox else ''

        escaped_text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")

        bodyPr_wrap = "square" if is_textbox else "square"
        bodyPr_anchor = "t" if is_textbox else "ctr"
        bodyPr_inset = ' lIns="45720" tIns="91440" rIns="45720" bIns="91440"' if not is_textbox else ''

        return (
            f'<p:sp xmlns:p="{p_ns}" xmlns:a="{a_ns}">'
            f'<p:nvSpPr><p:cNvPr id="{shape_id}" name="{name}"/>'
            f'<p:cNvSpPr{txbox_attr}/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{x_emu}" y="{y_emu}"/>'
            f'<a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>'
            f'<a:prstGeom prst="{prst}"><a:avLst/></a:prstGeom>'
            f'{fill_str}{ln_str}</p:spPr>'
            f'<p:txBody><a:bodyPr wrap="{bodyPr_wrap}" anchor="{bodyPr_anchor}"{bodyPr_inset}/>'
            f'<a:lstStyle/>'
            f'<a:p><a:pPr algn="{align_val}"/>'
            f'<a:r><a:rPr lang="zh-CN" sz="{sz}"{bold_attr} dirty="0"/>'
            f'{color_str}<a:latin typeface="{font_name}"/><a:ea typeface="{font_name}"/>'
            f'<a:t>{escaped_text}</a:t></a:r></a:p></p:txBody></p:sp>'
        )

    def _inject_grp_xml(self, slide, grp_xml_str: str) -> None:
        from lxml import etree
        p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        sp_tree = slide._element.find(f".//{{{p_ns}}}spTree")
        if sp_tree is None:
            return
        try:
            grp_elem = etree.fromstring(grp_xml_str.encode("utf-8"))
            next_id = 2
            for cNv in sp_tree.iter(f"{{{p_ns}}}cNvPr"):
                try:
                    next_id = max(next_id, int(cNv.get("id", "1")) + 1)
                except (ValueError, TypeError):
                    pass
            for cNv in grp_elem.iter(f"{{{p_ns}}}cNvPr"):
                cNv.set("id", str(next_id))
                next_id += 1
            sp_tree.append(grp_elem)
        except Exception:
            pass

    def _render_chevron_group(self, slide, bullets: list[str], x: float, y: float, w: float, h: float,
                               body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        n = len(bullets)
        if n == 0:
            return
        gap = 0.15
        step_w = (w - gap * (n - 1)) / n if n > 1 else w
        step_h = h * 0.35
        desc_y = y + step_h + 0.2
        desc_h = h - step_h - 0.2

        children = []
        for i, bullet in enumerate(bullets):
            sx = x + i * (step_w + gap)
            num_label = f"STEP {i + 1}"
            children.append(self._sp_xml(0, f"Chevron {i}", sx, y, step_w, step_h,
                                          fill_hex=accent_color, prst="chevron",
                                          text=num_label, font_size=14, font_color="#FFFFFF",
                                          bold=True, align="center", is_textbox=False, font_name=dna_font))
            children.append(self._sp_xml(0, f"Desc {i}", sx, desc_y, step_w, desc_h,
                                          text=bullet, font_size=12, font_color=body_color,
                                          is_textbox=True, font_name=dna_font))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _render_numbered_card_group(self, slide, cards: list[dict], x: float, y: float, w: float, h: float,
                                     body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        n = len(cards)
        if n == 0:
            return
        gap = 0.25
        card_w = (w - gap * (n - 1)) / n if n > 1 else w
        card_h = h
        bg = self._lighten_color(accent_color or "#84AF7D", 0.88)

        children = []
        for i, card in enumerate(cards):
            cx = x + i * (card_w + gap)
            title = card.get("title", "")
            body = card.get("body", card.get("text", ""))

            children.append(self._sp_xml(0, f"CardBg {i}", cx, y, card_w, card_h,
                                          fill_hex=bg, prst="roundRect"))
            children.append(self._sp_xml(0, f"CardHeader {i}", cx, y, card_w, 0.12,
                                          fill_hex=accent_color, prst="roundRect"))
            num_str = f"{i + 1:02d}"
            children.append(self._sp_xml(0, f"CardNum {i}", cx + 0.15, y + 0.2, 0.5, 0.5,
                                          text=num_str, font_size=24, font_color=accent_color, bold=True, font_name=dna_font))
            if title:
                children.append(self._sp_xml(0, f"CardTitle {i}", cx + 0.15, y + 0.75, card_w - 0.3, 0.4,
                                              text=title, font_size=16, font_color=accent_color, bold=True, font_name=dna_font))
            if body:
                children.append(self._sp_xml(0, f"CardBody {i}", cx + 0.15, y + 1.2, card_w - 0.3, card_h - 1.4,
                                              text=body, font_size=12, font_color=body_color, font_name=dna_font))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _render_icon_card_group(self, slide, cards: list[dict], x: float, y: float, w: float, h: float,
                                  body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        n = len(cards)
        if n == 0:
            return
        gap = 0.3
        card_w = (w - gap * (n - 1)) / n if n > 1 else w
        card_h = h

        children = []
        for i, card in enumerate(cards):
            cx = x + i * (card_w + gap)
            title = card.get("title", "")
            body = card.get("body", card.get("text", ""))

            children.append(self._sp_xml(0, f"CardBg {i}", cx, y, card_w, card_h,
                                          fill_hex="#FFFFFF", prst="roundRect"))
            children.append(self._sp_xml(0, f"CardTop {i}", cx, y, card_w, 1.2,
                                          fill_hex=accent_color, prst="roundRect"))
            children.append(self._sp_xml(0, f"CardNum {i}", cx + 0.2, y + 0.1, card_w - 0.4, 0.5,
                                          text=f"0{i + 1}", font_size=28, font_color="#FFFFFF",
                                          bold=True, align="center", font_name=dna_font))
            if title:
                children.append(self._sp_xml(0, f"CardTitle {i}", cx + 0.2, y + 0.65, card_w - 0.4, 0.4,
                                              text=title, font_size=16, font_color="#FFFFFF", bold=True, align="left", font_name=dna_font))
            if body:
                children.append(self._sp_xml(0, f"CardBody {i}", cx + 0.2, y + 1.4, card_w - 0.4, card_h - 1.6,
                                              text=body, font_size=12, font_color=body_color, font_name=dna_font))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _render_stat_card_group(self, slide, cards: list[dict], x: float, y: float, w: float, h: float,
                                   body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        n = len(cards)
        if n == 0:
            return
        gap = 0.3
        card_w = (w - gap * (n - 1)) / n if n > 1 else w
        card_h = h
        stats = ["98%", "40%", "5x"]

        children = []
        for i, card in enumerate(cards):
            cx = x + i * (card_w + gap)
            title = card.get("title", "")
            body = card.get("body", card.get("text", ""))

            children.append(self._sp_xml(0, f"CardBg {i}", cx, y, card_w, card_h,
                                          fill_hex="#FFFFFF", prst="roundRect"))
            children.append(self._sp_xml(0, f"CardTop {i}", cx, y, card_w, 0.08,
                                          fill_hex=accent_color, prst="roundRect"))
            stat = stats[i] if i < len(stats) else f"{(i + 1) * 33}%"
            children.append(self._sp_xml(0, f"CardStat {i}", cx + 0.2, y + 0.25, card_w - 0.4, 0.7,
                                          text=stat, font_size=36, font_color=accent_color,
                                          bold=True, align="center", font_name=dna_font))
            children.append(self._sp_xml(0, f"CardDivider {i}", cx + 0.15, y + 1.1, card_w - 0.3, 0.04,
                                          fill_hex=self._lighten_color(accent_color or "#84AF7D", 0.7)))
            if title:
                children.append(self._sp_xml(0, f"CardTitle {i}", cx + 0.2, y + 1.3, card_w - 0.4, 0.4,
                                              text=title, font_size=16, font_color=accent_color,
                                              bold=True, align="center", font_name=dna_font))
            if body:
                children.append(self._sp_xml(0, f"CardBody {i}", cx + 0.2, y + 1.8, card_w - 0.4, card_h - 2.0,
                                              text=body, font_size=12, font_color=body_color, align="center", font_name=dna_font))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _render_two_col_group(self, slide, bullets: list[str], x: float, y: float, w: float, h: float,
                               body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        col_w = (w - 0.3) / 2
        mid = (len(bullets) + 1) // 2
        bg = self._lighten_color(accent_color or "#84AF7D", 0.88)

        children = []
        for col, start, end in [(0, 0, mid), (1, mid, len(bullets))]:
            cx = x + col * (col_w + 0.3)
            for j, bullet in enumerate(bullets[start:end]):
                by = y + j * 1.0
                idx = start + j + 1
                children.append(self._sp_xml(0, f"RowBg {idx}", cx, by, col_w, 0.85,
                                              fill_hex=bg, prst="roundRect"))
                children.append(self._sp_xml(0, f"RowNum {idx}", cx + 0.1, by + 0.15, 0.4, 0.4,
                                              fill_hex=accent_color, prst="ellipse",
                                              text=str(idx), font_size=12, font_color="#FFFFFF",
                                              bold=True, align="center", font_name=dna_font))
                children.append(self._sp_xml(0, f"RowText {idx}", cx + 0.6, by + 0.15, col_w - 0.75, 0.55,
                                              text=bullet, font_size=12, font_color=body_color, font_name=dna_font))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _render_numbered_row_group(self, slide, bullets: list[str], x: float, y: float, w: float, h: float,
                                    body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        n = len(bullets)
        if n == 0:
            return
        row_h = min(h / n, 1.3)
        bg = self._lighten_color(accent_color or "#84AF7D", 0.88)

        children = []
        for i, bullet in enumerate(bullets):
            by = y + i * row_h
            children.append(self._sp_xml(0, f"RowBg {i}", x, by, w, row_h - 0.08,
                                          fill_hex=bg, prst="roundRect"))
            children.append(self._sp_xml(0, f"RowBar {i}", x, by, 0.08, row_h - 0.08,
                                          fill_hex=accent_color))
            children.append(self._sp_xml(0, f"RowNum {i}", x + 0.2, by + 0.12, 0.45, 0.45,
                                          fill_hex=accent_color, prst="ellipse",
                                          text=str(i + 1), font_size=14, font_color="#FFFFFF",
                                          bold=True, align="center", font_name=dna_font))
            children.append(self._sp_xml(0, f"RowText {i}", x + 0.8, by + 0.1, w - 1.0, row_h - 0.2,
                                          text=bullet, font_size=13, font_color=body_color, font_name=dna_font))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _render_hierarchy_group(self, slide, bullets: list[str], x: float, y: float, w: float, h: float,
                                 body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        n = len(bullets)
        if n == 0:
            return
        top_h = 0.8
        child_h = min((h - top_h - 0.3) / max(n - 1, 1), 1.0)
        child_w = (w - 0.2 * (n - 2)) / max(n - 1, 1) if n > 1 else w
        bg = self._lighten_color(accent_color or "#84AF7D", 0.88)

        children = []
        children.append(self._sp_xml(0, "HRoot", x + (w - 2.5) / 2, y, 2.5, top_h,
                                      fill_hex=accent_color, prst="roundRect",
                                      text=bullets[0], font_size=14, font_color="#FFFFFF",
                                      bold=True, align="center", font_name=dna_font))
        for i in range(1, n):
            cx = x + (i - 1) * (child_w + 0.2)
            cy = y + top_h + 0.3
            children.append(self._sp_xml(0, f"HChild {i}", cx, cy, child_w, child_h,
                                          fill_hex=bg, prst="roundRect",
                                          text=bullets[i], font_size=12, font_color=body_color,
                                          font_name=dna_font))
            connector_x = cx + child_w / 2
            root_cx = x + w / 2
            children.append(self._sp_xml(0, f"HLine {i}", min(connector_x, root_cx) - 0.01, y + top_h,
                                          abs(connector_x - root_cx) + 0.02, 0.3,
                                          fill_hex=accent_color, prst="rect"))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _render_swot_group(self, slide, bullets: list[str], x: float, y: float, w: float, h: float,
                            body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        n = len(bullets)
        if n == 0:
            return
        labels = ["S 优势", "W 劣势", "O 机会", "T 威胁"]
        swot_colors = [accent_color, self._lighten_color(accent_color or "#84AF7D", 0.6),
                       self._lighten_color(accent_color or "#84AF7D", 0.75),
                       self._lighten_color(accent_color or "#84AF7D", 0.45)]
        gap = 0.15
        cell_w = (w - gap) / 2
        cell_h = (h - gap) / 2

        children = []
        for i in range(min(n, 4)):
            col = i % 2
            row = i // 2
            cx = x + col * (cell_w + gap)
            cy = y + row * (cell_h + gap)
            label = labels[i] if i < len(labels) else f"Item {i + 1}"
            color = swot_colors[i] if i < len(swot_colors) else accent_color

            children.append(self._sp_xml(0, f"SWOTBg {i}", cx, cy, cell_w, cell_h,
                                          fill_hex=color, prst="roundRect"))
            children.append(self._sp_xml(0, f"SWOTLabel {i}", cx + 0.15, cy + 0.1, cell_w - 0.3, 0.4,
                                          text=label, font_size=16, font_color="#FFFFFF",
                                          bold=True, font_name=dna_font))
            children.append(self._sp_xml(0, f"SWOTText {i}", cx + 0.15, cy + 0.6, cell_w - 0.3, cell_h - 0.75,
                                          text=bullets[i], font_size=11, font_color="#FFFFFF",
                                          font_name=dna_font))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _render_timeline_group(self, slide, bullets: list[str], x: float, y: float, w: float, h: float,
                                body_color: str | None, accent_color: str | None, dna_font: str = "Microsoft YaHei") -> None:
        n = len(bullets)
        if n == 0:
            return
        line_y = y + 0.6
        node_r = 0.25
        step = w / max(n, 1)
        bg = self._lighten_color(accent_color or "#84AF7D", 0.88)

        children = []
        children.append(self._sp_xml(0, "TLLine", x, line_y - 0.03, w, 0.06,
                                      fill_hex=accent_color, prst="rect"))
        for i, bullet in enumerate(bullets):
            nx = x + i * step + step / 2 - node_r
            children.append(self._sp_xml(0, f"TLNode {i}", nx, line_y - node_r, node_r * 2, node_r * 2,
                                          fill_hex=accent_color, prst="ellipse",
                                          text=str(i + 1), font_size=11, font_color="#FFFFFF",
                                          bold=True, align="center", font_name=dna_font))
            text_y = line_y + node_r + 0.15 if i % 2 == 0 else y
            text_h = h - (text_y - y) - 0.1
            if text_h < 0.5:
                text_y = line_y + node_r + 0.15
                text_h = h - (text_y - y) - 0.1
            children.append(self._sp_xml(0, f"TLText {i}", x + i * step, text_y, step, text_h,
                                          text=bullet, font_size=11, font_color=body_color,
                                          align="center", font_name=dna_font))

        grp_xml = self._build_grp_xml(children, x, y, w, h)
        self._inject_grp_xml(slide, grp_xml)

    def _clear_content_shapes_by_name(
        self,
        slide,
        content_shape_names: set[str],
        slide_dna: SlideDNA,
        content_y_min: float = 0,
        content_y_max: float = 999,
    ) -> None:
        p_ns = _NS["p"]
        a_ns = _NS["a"]

        title_shape_names = set()
        for z in slide_dna.text_zones:
            if z.role == "title":
                title_shape_names.add(z.shape_name)

        preserve_names = title_shape_names

        slide_elem = slide._element
        sp_tree = slide_elem.find(qn("p:cSld")).find(qn("p:spTree"))

        shapes_to_remove = []
        for sp in list(sp_tree.findall(f"{{{p_ns}}}sp")):
            nvSpPr = sp.find(f"{{{p_ns}}}nvSpPr")
            shape_name = ""
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(f"{{{p_ns}}}cNvPr")
                if cNvPr is not None:
                    shape_name = cNvPr.get("name", "")
                    if shape_name in preserve_names:
                        continue
                    if shape_name in content_shape_names:
                        shapes_to_remove.append(sp)
                        continue

            off_elem = self._get_shape_offset(sp)
            if off_elem is not None and shape_name not in preserve_names:
                try:
                    sp_y = int(off_elem.get("y", 0)) / 914400
                    sp_h = 0
                    spPr = sp.find(f"{{{p_ns}}}spPr")
                    if spPr is not None:
                        xfrm = spPr.find(f"{{{a_ns}}}xfrm")
                        if xfrm is not None:
                            ext = xfrm.find(f"{{{a_ns}}}ext")
                            if ext is not None:
                                sp_h = int(ext.get("cy", 0)) / 914400
                    if sp_y < content_y_max and (sp_y + sp_h) > content_y_min:
                        has_text = False
                        text_content = ""
                        for t in sp.iter(f"{{{a_ns}}}t"):
                            if t.text and t.text.strip():
                                has_text = True
                                text_content += t.text.strip()
                        if not has_text:
                            shapes_to_remove.append(sp)
                        elif self._is_placeholder_text(text_content):
                            shapes_to_remove.append(sp)
                except (ValueError, TypeError):
                    pass

        for sp in shapes_to_remove:
            sp_tree.remove(sp)

    def _clear_text_only_in_zone(
        self,
        slide,
        slide_dna: SlideDNA,
        content_y_min: float = 0,
        content_y_max: float = 999,
    ) -> None:
        p_ns = _NS["p"]
        a_ns = _NS["a"]

        title_shape_names = set()
        for z in slide_dna.text_zones:
            if z.role == "title":
                title_shape_names.add(z.shape_name)

        slide_elem = slide._element
        sp_tree = slide_elem.find(qn("p:cSld")).find(qn("p:spTree"))

        for pic in list(sp_tree.findall(f"{{{p_ns}}}pic")):
            off_elem = None
            spPr = pic.find(f"{{{p_ns}}}spPr")
            if spPr is not None:
                xfrm = spPr.find(qn("a:xfrm"))
                if xfrm is not None:
                    off_elem = xfrm.find(qn("a:off"))
            if off_elem is None:
                nvSpPr = pic.find(f"{{{p_ns}}}nvSpPr")
                if nvSpPr is not None:
                    cNvPr = nvSpPr.find(f"{{{p_ns}}}cNvPr")
                    if cNvPr is not None and cNvPr.get("name", "") in title_shape_names:
                        continue
            if off_elem is not None:
                try:
                    sp_y = int(off_elem.get("y", 0)) / 914400
                except (ValueError, TypeError):
                    continue
                if content_y_min <= sp_y <= content_y_max:
                    sp_tree.remove(pic)

        for grp in list(sp_tree.findall(f"{{{p_ns}}}grpSp")):
            grpSpPr = grp.find(f"{{{p_ns}}}grpSpPr")
            if grpSpPr is None:
                continue
            xfrm = grpSpPr.find(qn("a:xfrm"))
            if xfrm is None:
                continue
            off = xfrm.find(qn("a:off"))
            if off is None:
                continue
            try:
                sp_y = int(off.get("y", 0)) / 914400
            except (ValueError, TypeError):
                continue
            if content_y_min <= sp_y <= content_y_max:
                sp_tree.remove(grp)

        for sp in list(sp_tree.findall(f"{{{p_ns}}}sp")):
            nvSpPr = sp.find(f"{{{p_ns}}}nvSpPr")
            shape_name = ""
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(f"{{{p_ns}}}cNvPr")
                if cNvPr is not None:
                    shape_name = cNvPr.get("name", "")
                    if shape_name in title_shape_names:
                        continue

            off_elem = self._get_shape_offset(sp)
            if off_elem is None:
                continue

            try:
                sp_y = int(off_elem.get("y", 0)) / 914400
            except (ValueError, TypeError):
                continue

            if sp_y < content_y_min or sp_y > content_y_max:
                continue

            has_fill = False
            spPr = sp.find(f"{{{p_ns}}}spPr")
            if spPr is not None:
                for fill_tag in ("a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"):
                    if spPr.find(qn(fill_tag)) is not None:
                        has_fill = True
                        break

            if has_fill:
                sp_bounds = self._get_shape_bounds_from_xml(sp)
                if sp_bounds is not None:
                    sw, sh = sp_bounds[2], sp_bounds[3]
                    is_banner = sw >= 10 and sh <= 0.8
                    is_small_deco = sw <= 1.0 and sh <= 1.0
                    if not is_banner and not is_small_deco:
                        sp_tree.remove(sp)
                        continue
                for t in list(sp.iter(f"{{{a_ns}}}t")):
                    if t.text and self._is_placeholder_text(t.text):
                        t.text = ""
                continue

            has_text = False
            text_content = ""
            for t in sp.iter(f"{{{a_ns}}}t"):
                if t.text and t.text.strip():
                    has_text = True
                    text_content += t.text.strip()
            if has_text:
                if self._is_placeholder_text(text_content):
                    sp_tree.remove(sp)
                else:
                    for t in list(sp.iter(f"{{{a_ns}}}t")):
                        t.text = ""
            else:
                pass

    def _clear_all_decorations(self, slide, slide_dna: SlideDNA) -> None:
        p_ns = _NS["p"]
        a_ns = _NS["a"]

        title_shape_names = set()
        for z in slide_dna.text_zones:
            if z.role == "title":
                title_shape_names.add(z.shape_name)

        slide_elem = slide._element
        sp_tree = slide_elem.find(qn("p:cSld")).find(qn("p:spTree"))

        for pic in list(sp_tree.findall(f"{{{p_ns}}}pic")):
            sp_tree.remove(pic)

        for grp in list(sp_tree.findall(f"{{{p_ns}}}grpSp")):
            sp_tree.remove(grp)

        for sp in list(sp_tree.findall(f"{{{p_ns}}}sp")):
            nvSpPr = sp.find(f"{{{p_ns}}}nvSpPr")
            shape_name = ""
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(f"{{{p_ns}}}cNvPr")
                if cNvPr is not None:
                    shape_name = cNvPr.get("name", "")
                    if shape_name in title_shape_names:
                        continue

            spPr = sp.find(f"{{{p_ns}}}spPr")
            if spPr is None:
                continue
            has_fill = False
            for fill_tag in ("a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"):
                if spPr.find(qn(fill_tag)) is not None:
                    has_fill = True
                    break
            if has_fill:
                sp_tree.remove(sp)
                continue

            has_text = False
            for t in sp.iter(f"{{{a_ns}}}t"):
                if t.text and t.text.strip():
                    has_text = True
                    break
            if has_text:
                text_content = ""
                for t in sp.iter(f"{{{a_ns}}}t"):
                    if t.text and t.text.strip():
                        text_content += t.text.strip() + " "
                text_content = text_content.strip()
                if self._is_placeholder_text(text_content):
                    for t in list(sp.iter(f"{{{a_ns}}}t")):
                        t.text = ""
            else:
                sp_tree.remove(sp)

    def _classify_decoration_by_bounds(
        self, x: float, y: float, w: float, h: float,
        has_fill: bool = False, is_picture: bool = False,
    ) -> tuple[int, str]:
        if is_picture:
            if w > 3.0 and h > 3.0:
                return (3, "large_image")
            if w <= 1.0 and h <= 1.0:
                return (1, "icon")
            return (2, "scene_image")
        if has_fill:
            if w <= 1.5 and h >= 5.0:
                return (0, "sidebar")
            if y >= 6.8:
                return (0, "footer")
            if w >= 10.0 and h <= 0.8:
                return (0, "banner")
            if 0.5 <= y <= 2.0 and h <= 0.3 and w <= 8.0:
                return (0, "title_decoration")
            if w <= 2.0 and h <= 2.0 and (x <= 2.0 or x >= 11.0 or y <= 1.0 or y >= 5.5):
                return (1, "corner")
            if w >= 5.0 and h <= 0.1:
                return (1, "divider")
            if w <= 1.0 and h <= 1.0:
                return (1, "small_deco")
            return (3, "large_fill")
        return (3, "unknown")

    def _selective_clear_decorations(
        self, slide, slide_dna: SlideDNA, plan: PagePlan,
        keep_levels: set[int] | None = None,
    ) -> None:
        if keep_levels is None:
            keep_levels = {0}

        p_ns = _NS["p"]
        a_ns = _NS["a"]
        title_shape_names = {z.shape_name for z in slide_dna.text_zones if z.role == "title"}
        slide_elem = slide._element
        sp_tree = slide_elem.find(qn("p:cSld")).find(qn("p:spTree"))

        for pic in list(sp_tree.findall(f"{{{p_ns}}}pic")):
            bounds = self._get_shape_bounds_from_xml(pic)
            if bounds is None:
                sp_tree.remove(pic)
                continue
            level, _ = self._classify_decoration_by_bounds(
                bounds[0], bounds[1], bounds[2], bounds[3],
                has_fill=False, is_picture=True,
            )
            if level not in keep_levels:
                sp_tree.remove(pic)

        for grp in list(sp_tree.findall(f"{{{p_ns}}}grpSp")):
            sp_tree.remove(grp)

        for sp in list(sp_tree.findall(f"{{{p_ns}}}sp")):
            nvSpPr = sp.find(f"{{{p_ns}}}nvSpPr")
            shape_name = ""
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(f"{{{p_ns}}}cNvPr")
                if cNvPr is not None:
                    shape_name = cNvPr.get("name", "")
                    if shape_name in title_shape_names:
                        continue

            bounds = self._get_shape_bounds_from_xml(sp)
            spPr = sp.find(f"{{{p_ns}}}spPr")
            has_fill = False
            if spPr is not None:
                for fill_tag in ("a:solidFill", "a:gradFill", "a:pattFill", "a:blipFill"):
                    if spPr.find(qn(fill_tag)) is not None:
                        has_fill = True
                        break

            if bounds is not None:
                level, _ = self._classify_decoration_by_bounds(
                    bounds[0], bounds[1], bounds[2], bounds[3],
                    has_fill=has_fill,
                )
            else:
                level = 3

            has_text = False
            text_content = ""
            for t in sp.iter(f"{{{a_ns}}}t"):
                if t.text and t.text.strip():
                    has_text = True
                    text_content += t.text.strip() + " "
            text_content = text_content.strip()

            if level not in keep_levels:
                sp_tree.remove(sp)
            elif has_text and self._is_placeholder_text(text_content):
                for t in list(sp.iter(f"{{{a_ns}}}t")):
                    t.text = ""

    def _detect_card_groups(self, slide_dna: SlideDNA) -> list[dict]:
        page_number_zones = sorted(
            [z for z in slide_dna.text_zones if z.role == "page_number"],
            key=lambda z: z.bounds[0],
        )
        if not page_number_zones:
            return []

        groups = []
        for pn in page_number_zones:
            px, py = pn.bounds[0], pn.bounds[1]
            card_zones = {"page_number": pn}

            for z in slide_dna.text_zones:
                if z is pn:
                    continue
                zx, zy = z.bounds[0], z.bounds[1]
                if z.role == "subtitle" and zy > py and abs(zx - px) < 3.0 and zy - py < 2.0:
                    if "subtitle" not in card_zones or zy < card_zones["subtitle"].bounds[1]:
                        card_zones["subtitle"] = z
                elif z.role == "body" and zy > py and abs(zx - px) < 3.0 and zy - py < 2.0:
                    if "body" not in card_zones or zy < card_zones["body"].bounds[1]:
                        card_zones["body"] = z

            if "subtitle" in card_zones or "body" in card_zones:
                groups.append(card_zones)

        return groups if len(groups) >= 2 else []

    def _patch_card_content(self, slide, slide_dna: SlideDNA, plan: PagePlan) -> None:
        card_groups = self._detect_card_groups(slide_dna)
        if not card_groups:
            return

        items = plan.cards if plan.cards else [{"title": b, "body": ""} for b in (plan.bullets or [])]
        n = min(len(card_groups), len(items))

        slide_elem = slide._element

        for i in range(n):
            group = card_groups[i]
            item = items[i]
            patched: set = set()

            pn_zone = group.get("page_number")
            if pn_zone:
                self._patch_zone_text(slide_elem, pn_zone, f"{i + 1:02d}", patched)

            sub_zone = group.get("subtitle")
            if sub_zone:
                title = item.get("title", "") if isinstance(item, dict) else str(item)
                if title and len(title) <= 8:
                    self._patch_zone_text(slide_elem, sub_zone, title, patched)
                else:
                    self._clear_zone_text(slide_elem, sub_zone, patched)

            body_zone = group.get("body")
            if body_zone:
                body = item.get("body", item.get("text", "")) if isinstance(item, dict) else ""
                self._patch_zone_text(slide_elem, body_zone, body, patched)
                self._autosize_zone_text(slide_elem, body_zone, body)

        for i in range(n, len(card_groups)):
            group = card_groups[i]
            patched: set = set()
            for key in ("page_number", "subtitle", "body"):
                zone = group.get(key)
                if zone:
                    if key == "page_number":
                        self._patch_zone_text(slide_elem, zone, f"{i + 1:02d}", patched)
                    else:
                        self._clear_zone_text(slide_elem, zone, patched)

    def _clear_decorative_images(
        self,
        slide,
        slide_dna: SlideDNA,
        content_y_min: float = 0,
        content_y_max: float = 999,
    ) -> None:
        p_ns = _NS["p"]
        a_ns = _NS["a"]
        slide_elem = slide._element
        sp_tree = slide_elem.find(qn("p:cSld")).find(qn("p:spTree"))

        title_shape_names = set()
        for z in slide_dna.text_zones:
            if z.role == "title":
                title_shape_names.add(z.shape_name)

        pics_to_remove = []
        for pic in list(sp_tree.findall(f"{{{p_ns}}}pic")):
            off_elem = None
            spPr = pic.find(f"{{{p_ns}}}spPr")
            if spPr is not None:
                xfrm = spPr.find(f"{{{a_ns}}}xfrm")
                if xfrm is not None:
                    off_elem = xfrm.find(f"{{{a_ns}}}off")

            if off_elem is not None:
                try:
                    sp_y = int(off_elem.get("y", 0)) / 914400
                    sp_h = 0
                    ext = spPr.find(f"{{{a_ns}}}xfrm").find(f"{{{a_ns}}}ext") if spPr.find(f"{{{a_ns}}}xfrm") is not None else None
                    if ext is not None:
                        sp_h = int(ext.get("cy", 0)) / 914400
                    if sp_y < content_y_max and (sp_y + sp_h) > content_y_min:
                        pics_to_remove.append(pic)
                except (ValueError, TypeError, AttributeError):
                    pass

        for pic in pics_to_remove:
            sp_tree.remove(pic)

    def _apply_dna_fonts_to_brand(
        self,
        brand_spec: BrandSpec,
        dna: Any = None,
        slide_dna: SlideDNA | None = None,
    ) -> None:
        heading_font = None
        body_font = None
        cjk_heading = None
        cjk_body = None

        if slide_dna:
            for z in slide_dna.text_zones:
                if z.role == "title" and z.font_name:
                    heading_font = z.font_name
                    if not cjk_heading and any('\u4e00' <= c <= '\u9fff' for c in z.font_name):
                        cjk_heading = z.font_name
                elif z.role in ("body", "subtitle") and z.font_name:
                    body_font = z.font_name
                    if not cjk_body and any('\u4e00' <= c <= '\u9fff' for c in z.font_name):
                        cjk_body = z.font_name

        if brand_spec.fonts is None:
            brand_spec.fonts = {}

        if heading_font:
            brand_spec.fonts["heading"] = heading_font
        if body_font:
            brand_spec.fonts["body"] = body_font
        if cjk_heading:
            brand_spec.fonts["cjk_heading"] = cjk_heading
        if cjk_body:
            brand_spec.fonts["cjk_body"] = cjk_body

        if brand_spec.colors is None:
            brand_spec.colors = {}

        if slide_dna:
            for z in slide_dna.text_zones:
                if z.role == "subtitle" and z.color_hex:
                    brand_spec.colors.setdefault("foreground", z.color_hex)
                    brand_spec.colors.setdefault("body", z.color_hex)
                    break
            for z in slide_dna.text_zones:
                if z.role in ("title", "page_number") and z.color_hex:
                    brand_spec.colors.setdefault("accent", z.color_hex)
                    brand_spec.colors.setdefault("heading", z.color_hex)
                    break

    def _get_shape_offset(self, sp):
        a_ns = _NS["a"]
        p_ns = _NS["p"]
        spPr = sp.find(f"{{{p_ns}}}spPr")
        if spPr is not None:
            xfrm = spPr.find(f"{{{a_ns}}}xfrm")
            if xfrm is not None:
                return xfrm.find(f"{{{a_ns}}}off")
        return None

    def _render_cards_in_zone(
        self,
        slide,
        precision,
        cards: list[dict],
        x: float, y: float, w: float, h: float,
        title_color: str | None = None,
        body_color: str | None = None,
    ) -> None:
        n = len(cards)
        if n == 0:
            return

        gap = 0.2
        card_w = (w - gap * (n - 1)) / n if n > 1 else w
        card_h = h

        for i, card in enumerate(cards):
            cx = x + i * (card_w + gap)
            title = card.get("title", "")
            body = card.get("body", card.get("text", ""))

            if title:
                precision.add_text(
                    slide, title, cx, y, card_w, 0.5,
                    size=14, bold=True, color_hex=title_color, color_role="accent",
                )
            if body:
                precision.add_multiline(
                    slide, [body], cx, y + 0.6, card_w, card_h - 0.6,
                    size=12, color_hex=body_color, color_role="foreground",
                )

    def _render_styled_cards(
        self,
        slide,
        precision,
        cards: list[dict],
        x: float, y: float, w: float, h: float,
        title_color: str | None = None,
        body_color: str | None = None,
        accent_color: str | None = None,
    ) -> None:
        n = len(cards)
        if n == 0:
            return

        gap = 0.3
        card_w = (w - gap * (n - 1)) / n if n > 1 else w
        card_h = h

        for i, card in enumerate(cards):
            cx = x + i * (card_w + gap)
            title = card.get("title", "")
            body = card.get("body", card.get("text", ""))

            bg_color = self._lighten_color(accent_color or "#84AF7D", 0.85)
            precision.add_rounded_rect(
                slide, cx, y, card_w, card_h,
                fill_hex=bg_color,
            )

            precision.add_gradient_line(
                slide, cx + 0.15, y + 0.15, card_w * 0.3, 0.04,
                accent_color or "#84AF7D",
            )

            if title:
                precision.add_text(
                    slide, title, cx + 0.15, y + 0.3, card_w - 0.3, 0.5,
                    size=16, bold=True, color_hex=accent_color, color_role="accent",
                )
            if body:
                precision.add_multiline(
                    slide, [body], cx + 0.15, y + 0.9, card_w - 0.3, card_h - 1.1,
                    size=12, color_hex=body_color, color_role="foreground",
                )

    def _render_numbered_bullets(
        self,
        slide,
        precision,
        bullets: list[str],
        x: float, y: float, w: float, h: float,
        font_size: int = 14,
        body_color: str | None = None,
        accent_color: str | None = None,
    ) -> None:
        n = len(bullets)
        if n == 0:
            return

        line_h = min(h / n, 0.8)
        num_size = max(font_size - 2, 10)

        for i, bullet in enumerate(bullets):
            by = y + i * line_h

            num_str = f"{i + 1:02d}"
            precision.add_text(
                slide, num_str, x, by, 0.6, line_h,
                size=num_size + 4, bold=True, color_hex=accent_color, color_role="accent",
            )

            precision.add_text(
                slide, bullet, x + 0.7, by + 0.05, w - 0.7, line_h - 0.1,
                size=font_size, color_hex=body_color, color_role="foreground",
            )

    def _render_numbered_list(
        self,
        slide,
        precision,
        bullets: list[str],
        x: float, y: float, w: float, h: float,
        font_size: int = 14,
        body_color: str | None = None,
        accent_color: str | None = None,
    ) -> None:
        n = len(bullets)
        if n == 0:
            return

        line_h = min(h / n, 1.0)
        circle_size = 0.4

        for i, bullet in enumerate(bullets):
            by = y + i * line_h
            cx = x + circle_size / 2

            precision.add_oval(
                slide, cx - circle_size / 2, by + 0.05, circle_size, circle_size,
                fill_hex=accent_color, fill_role="accent",
            )

            precision.add_text(
                slide, str(i + 1), cx - circle_size / 2, by + 0.05, circle_size, circle_size,
                size=11, bold=True, color_hex="#FFFFFF", color_role="on-primary", align="center",
            )

            precision.add_text(
                slide, bullet, x + circle_size + 0.15, by + 0.05, w - circle_size - 0.15, line_h - 0.1,
                size=font_size, color_hex=body_color, color_role="foreground",
            )

    def _render_process_flow(
        self,
        slide,
        precision,
        bullets: list[str],
        x: float, y: float, w: float, h: float,
        font_size: int = 14,
        body_color: str | None = None,
        accent_color: str | None = None,
    ) -> None:
        n = len(bullets)
        if n == 0:
            return

        step_w = (w - 0.3 * (n - 1)) / n if n > 1 else w
        step_h = min(h, 1.5)

        for i, bullet in enumerate(bullets):
            sx = x + i * (step_w + 0.3)

            bg_color = self._lighten_color(accent_color or "#84AF7D", 0.85)
            precision.add_rounded_rect(
                slide, sx, y, step_w, step_h,
                fill_hex=bg_color,
            )

            precision.add_text(
                slide, f"STEP {i + 1}", sx + 0.1, y + 0.1, step_w - 0.2, 0.3,
                size=9, bold=True, color_hex=accent_color, color_role="accent",
            )

            precision.add_text(
                slide, bullet, sx + 0.1, y + 0.45, step_w - 0.2, step_h - 0.55,
                size=font_size - 1, color_hex=body_color, color_role="foreground",
            )

            if i < n - 1:
                arrow_x = sx + step_w + 0.05
                arrow_y = y + step_h / 2 - 0.1
                precision.add_text(
                    slide, "→", arrow_x, arrow_y, 0.2, 0.2,
                    size=16, bold=True, color_hex=accent_color, color_role="accent", align="center",
                )

    def _extract_body_color(self, slide_dna: SlideDNA) -> str | None:
        for z in slide_dna.text_zones:
            if z.role == "body" and z.color_hex:
                return z.color_hex
        for z in slide_dna.text_zones:
            if z.role == "subtitle" and z.color_hex:
                return z.color_hex
        for z in slide_dna.text_zones:
            if z.role in ("title", "page_number") and z.color_hex:
                return z.color_hex
        return None

    def _extract_accent_color(self, slide_dna: SlideDNA, body_color: str | None) -> str | None:
        for z in slide_dna.text_zones:
            if z.role in ("page_number", "badge") and z.color_hex:
                return z.color_hex
        return body_color

    def _extract_body_color_from_dna(self, dna) -> str | None:
        for sd in dna.slides:
            c = self._extract_body_color(sd)
            if c:
                return c
        return None

    def _extract_accent_color_from_dna(self, dna, body_color: str | None) -> str | None:
        for sd in dna.slides:
            c = self._extract_accent_color(sd, body_color)
            if c and c != body_color:
                return c
        return body_color

    def _compute_default_content_bounds(self, plan: PagePlan) -> tuple[float, float, float, float]:
        x = 0.9
        y = 1.6
        if plan.title:
            y = 1.6
        w = 11.5
        h = 5.0
        return (x, y, w, h)

    def _resolve_dna_font(self, brand_spec: BrandSpec | None) -> str:
        if brand_spec and brand_spec.fonts:
            heading = brand_spec.fonts.get("heading", "")
            if heading:
                return heading
        if hasattr(self, "_last_dna") and self._last_dna and self._last_dna.font_scheme:
            return self._last_dna.font_scheme
        return "Microsoft YaHei"

    def _resolve_dna_accent(self, brand_spec: BrandSpec | None) -> str:
        if brand_spec and brand_spec.colors:
            for key in ("primary", "accent"):
                val = brand_spec.colors.get(key, "")
                if val and val.startswith("#"):
                    return val
        if hasattr(self, "_last_dna") and self._last_dna and self._last_dna.color_palette:
            for c in self._last_dna.color_palette:
                if c and c.startswith("#") and c not in ("#FFFFFF", "#000000"):
                    return c
        return "#2D5016"

    def _resolve_dna_body(self, brand_spec: BrandSpec | None) -> str:
        if brand_spec and brand_spec.colors:
            for key in ("foreground", "muted-foreground"):
                val = brand_spec.colors.get(key, "")
                if val and val.startswith("#"):
                    return val
        return "#333333"

    def _compute_content_bounds(self, slide_dna, plan: PagePlan) -> tuple[float, float, float, float]:
        all_content_zones = [
            z for z in slide_dna.text_zones
            if z.role in ("body", "subtitle", "page_number")
            and z.shape_name not in {
                tz.shape_name for tz in slide_dna.text_zones if tz.role == "title"
            }
        ]
        body_zones = [z for z in slide_dna.text_zones if z.role == "body"]
        if not all_content_zones:
            all_content_zones = body_zones
        if not all_content_zones:
            return (0.9, 1.6, 11.5, 5.0)

        content_y_min = min(z.bounds[1] for z in all_content_zones)
        title_zones_list = [z for z in slide_dna.text_zones if z.role == "title"]
        title_y_bottom = max(z.bounds[1] + z.bounds[3] for z in title_zones_list) if title_zones_list else 0
        if plan.page_type == "content":
            content_y_min = title_y_bottom
            content_y_max = 7.5
        else:
            if title_y_bottom < content_y_min:
                content_y_min = title_y_bottom
            content_y_max = max(z.bounds[1] + z.bounds[3] for z in all_content_zones)

        min_x = min(z.bounds[0] for z in all_content_zones)
        max_x = max(z.bounds[0] + z.bounds[2] for z in all_content_zones)
        return (min_x, content_y_min, max_x - min_x, content_y_max - content_y_min)

    def _auto_layout_from_plan(self, plan: PagePlan) -> str:
        if plan.cards:
            return "cards"
        if not plan.bullets:
            return "bullets"

        n = len(plan.bullets)
        all_text = " ".join(plan.bullets).lower()

        if any(kw in all_text for kw in ("步骤", "阶段", "process", "流程", "step", "phase", "pipeline", "workflow")):
            return "process"
        if any(kw in all_text for kw in ("架构", "层级", "hierarchy", "组织", "org chart")):
            return "hierarchy"
        if any(kw in all_text for kw in ("里程碑", "timeline", "时间线", "milestone", "roadmap")):
            return "timeline"
        if any(kw in all_text for kw in ("swot", "优势", "劣势", "对比", "comparison")):
            return "swot"
        if n <= 3:
            return "numbered"
        if n <= 5:
            return "numbered"
        return "bullets"

    @staticmethod
    def _lighten_color(hex_color: str, factor: float = 0.7) -> str:
        try:
            hex_color = hex_color.lstrip("#")
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            r = int(r + (255 - r) * factor)
            g = int(g + (255 - g) * factor)
            b = int(b + (255 - b) * factor)
            return f"#{r:02x}{g:02x}{b:02x}"
        except Exception:
            return "#F0F5FF"

    @staticmethod
    def dna_has_visual_content(dna: DesignDNA) -> bool:
        if not dna.slides:
            return False
        for slide in dna.slides:
            if len(slide.text_zones) > 2:
                return True
            if len(slide.image_refs) > 0:
                return True
            if len(slide.background_colors) > 0:
                return True
            if len(slide.group_xmls) > 0:
                return True
        return False

    def _patch_slide(self, slide, slide_dna: SlideDNA, patch: dict[str, Any]) -> None:
        page_type = slide_dna.page_type

        if page_type == "cover":
            self._patch_text_zones(slide, slide_dna, patch, allowed_roles={"title", "subtitle"})
        elif page_type == "toc":
            self._patch_text_zones(slide, slide_dna, patch, allowed_roles={"title", "subtitle", "body"})
            if patch.get("bullets"):
                self._patch_toc_bullets(slide, slide_dna, patch["bullets"])
        elif page_type == "content":
            self._patch_text_zones(slide, slide_dna, patch, allowed_roles={"title", "subtitle", "body"})
            self._handle_content_dual_title(slide, slide_dna, patch)
        elif page_type == "transition":
            self._patch_text_zones(slide, slide_dna, patch, allowed_roles={"title", "subtitle", "body"})
        elif page_type == "back_cover":
            self._patch_text_zones(slide, slide_dna, patch, allowed_roles={"title", "subtitle", "body"})
            if patch.get("title"):
                title_zones = [z for z in slide_dna.text_zones if z.role == "title"]
                if self._is_composite_title(title_zones) and len(title_zones) > 1:
                    slide_elem = slide._element
                    patched: set = set()
                    self._patch_zone_text(slide_elem, title_zones[0], patch["title"], patched)
                    for zone in title_zones[1:]:
                        self._clear_zone_text(slide_elem, zone, patched)
        else:
            self._patch_text_zones(slide, slide_dna, patch, allowed_roles={"title", "subtitle", "body"})

        self._clear_remaining_placeholders(slide)

    def _clear_remaining_placeholders(self, slide) -> None:
        a_ns = _NS["a"]
        p_ns = _NS["p"]
        slide_elem = slide._element
        sp_tree = slide_elem.find(qn("p:cSld")).find(qn("p:spTree"))
        for sp in list(sp_tree.findall(f"{{{p_ns}}}sp")):
            for t in list(sp.iter(f"{{{a_ns}}}t")):
                if t.text and self._is_placeholder_text(t.text):
                    t.text = ""
                    r = t.getparent()
                    if r is not None:
                        rPr = r.find(f"{{{a_ns}}}rPr")
                        if rPr is not None:
                            r.remove(rPr)

    def _patch_text_zones(
        self, slide, slide_dna: SlideDNA, patch: dict[str, Any],
        allowed_roles: set[str] | None = None,
    ) -> None:
        slide_elem = slide._element

        title_zones = [z for z in slide_dna.text_zones if z.role == "title" and (allowed_roles is None or "title" in allowed_roles)]
        subtitle_zones = [z for z in slide_dna.text_zones if z.role == "subtitle" and (allowed_roles is None or "subtitle" in allowed_roles)]

        is_composite_title = self._is_composite_title(title_zones)
        is_composite_subtitle = self._is_composite_title(subtitle_zones)

        patched_t_elems: set = set()

        if is_composite_title and patch.get("title") and title_zones:
            first_zone = title_zones[0]
            self._patch_zone_text(slide_elem, first_zone, patch["title"], patched_t_elems)
            self._autosize_composite_title(slide_elem, first_zone, patch["title"], title_zones)
            for zone in title_zones[1:]:
                self._clear_zone_text(slide_elem, zone, patched_t_elems)

        if is_composite_subtitle and patch.get("subtitle") and subtitle_zones:
            first_zone = subtitle_zones[0]
            self._patch_zone_text(slide_elem, first_zone, patch["subtitle"], patched_t_elems)
            for zone in subtitle_zones[1:]:
                self._clear_zone_text(slide_elem, zone, patched_t_elems)

        for zone in slide_dna.text_zones:
            if allowed_roles and zone.role not in allowed_roles:
                continue

            if is_composite_title and zone.role == "title":
                continue
            if is_composite_subtitle and zone.role == "subtitle":
                continue

            new_text = None
            if zone.role == "title" and patch.get("title"):
                new_text = patch["title"]
            elif zone.role == "subtitle" and patch.get("subtitle"):
                new_text = patch["subtitle"]
            elif zone.role == "body" and patch.get("bullets"):
                bullets = patch["bullets"]
                if isinstance(bullets, list) and zone.zone_id.endswith("_0"):
                    new_text = bullets[0] if bullets else None
                elif isinstance(bullets, list):
                    idx_str = zone.zone_id.rsplit("_", 1)[-1]
                    try:
                        bidx = int(idx_str)
                        new_text = bullets[bidx] if bidx < len(bullets) else None
                    except ValueError:
                        pass
                elif isinstance(bullets, str):
                    new_text = bullets
            elif zone.role == "body" and not patch.get("bullets"):
                if self._is_placeholder_text(zone.text):
                    new_text = ""

            if new_text is None:
                continue

            self._patch_zone_text(slide_elem, zone, new_text, patched_t_elems)

    def _handle_content_dual_title(self, slide, slide_dna: SlideDNA, patch: dict[str, Any]) -> None:
        title_zones = [z for z in slide_dna.text_zones if z.role == "title"]
        if len(title_zones) < 2:
            return

        sorted_by_size = sorted(title_zones, key=lambda z: z.font_size_pt or 0)

        small_zones = [z for z in sorted_by_size if (z.font_size_pt or 0) < 50]
        large_zones = [z for z in sorted_by_size if (z.font_size_pt or 0) >= 50]

        if not small_zones or not large_zones:
            return

        slide_elem = slide._element
        patched: set = set()

        if patch.get("title"):
            self._patch_zone_text(slide_elem, small_zones[0], patch["title"], patched)
            for zone in small_zones[1:]:
                self._clear_zone_text(slide_elem, zone, patched)
            for zone in large_zones:
                self._clear_zone_text(slide_elem, zone, patched)

    @staticmethod
    def _is_placeholder_text(text: str) -> bool:
        if not text:
            return False
        lower = text.lower()
        return any(kw in lower for kw in (
            "单击此处", "请输入", "click", "placeholder", "your text",
            "此处添加", "输入你的正文",
            "您的正文已经简明扼要", "字字珠玑", "提炼思想的精髓",
            "此处可编辑",
        ))

    def _patch_toc_bullets(
        self, slide, slide_dna: SlideDNA, bullets: list[str],
    ) -> None:
        slide_elem = slide._element

        subtitle_zones = [z for z in slide_dna.text_zones if z.role == "subtitle"]
        patched: set = set()

        for i, bullet in enumerate(bullets):
            if i < len(subtitle_zones):
                self._patch_zone_text(slide_elem, subtitle_zones[i], bullet, patched)

        for i in range(len(bullets), len(subtitle_zones)):
            self._clear_zone_text(slide_elem, subtitle_zones[i], patched)

    def _is_composite_title(self, zones: list[TextZone]) -> bool:
        if len(zones) < 2:
            return False
        same_size = len(set(z.font_size_pt for z in zones if z.font_size_pt)) <= 1
        close_horizontally = True
        sorted_zones = sorted(zones, key=lambda z: z.bounds[0])
        for i in range(1, len(sorted_zones)):
            prev_right = sorted_zones[i - 1].bounds[0] + sorted_zones[i - 1].bounds[2]
            curr_left = sorted_zones[i].bounds[0]
            gap = curr_left - prev_right
            if gap > 1.5:
                close_horizontally = False
                break
        short_texts = all(len(z.text) <= 2 for z in zones)
        western_composite = (
            same_size
            and close_horizontally
            and len(zones) <= 5
            and all(len(z.text) <= 15 for z in zones)
        )
        return (same_size or close_horizontally) and (short_texts or western_composite)

    def _autosize_composite_title(
        self, slide_elem, first_zone: TextZone, new_text: str, all_title_zones: list[TextZone],
    ) -> None:
        a_ns = _NS["a"]
        original_total_width = sum(z.bounds[2] for z in all_title_zones)
        original_char_count = sum(len(z.text) for z in all_title_zones)
        if original_char_count == 0 or original_total_width == 0:
            return

        avg_char_width = original_total_width / original_char_count
        needed_width = len(new_text) * avg_char_width
        if needed_width <= original_total_width * 1.1:
            return

        scale = original_total_width / needed_width
        target_shape = self._find_shape_by_name(slide_elem, first_zone.shape_name)
        if target_shape is None:
            return

        for rPr in target_shape.iter(f"{{{a_ns}}}rPr"):
            sz = rPr.get("sz")
            if sz:
                try:
                    old_sz = int(sz)
                    new_sz = max(int(old_sz * scale), 2400)
                    rPr.set("sz", str(new_sz))
                except ValueError:
                    pass

    def _autosize_zone_text(self, slide_elem, zone: TextZone, new_text: str) -> None:
        a_ns = _NS["a"]
        if not new_text or not zone.font_size_pt or zone.bounds[2] <= 0:
            return

        chars_per_line = zone.bounds[2] / 0.22
        needed_lines = max(1, len(new_text) / chars_per_line) if chars_per_line > 0 else 1
        available_lines = zone.bounds[3] / (zone.font_size_pt / 72) if zone.font_size_pt > 0 else 1

        if needed_lines <= available_lines * 1.2:
            return

        target_shape = self._find_shape_by_name(slide_elem, zone.shape_name)
        if target_shape is None:
            return

        scale = min(1.0, available_lines / needed_lines)
        for rPr in target_shape.iter(f"{{{a_ns}}}rPr"):
            sz = rPr.get("sz")
            if sz:
                try:
                    old_sz = int(sz)
                    new_sz = max(int(old_sz * scale), 1000)
                    rPr.set("sz", str(new_sz))
                except ValueError:
                    pass

    def _patch_zone_text(self, slide_elem, zone: TextZone, new_text: str, patched: set) -> None:
        a_ns = _NS["a"]

        target_shape = self._find_shape_by_name(slide_elem, zone.shape_name)
        if target_shape is not None:
            t_elem = self._locate_t_elem(target_shape, zone)
            if t_elem is not None and id(t_elem) not in patched:
                t_elem.text = new_text
                for child in list(t_elem):
                    t_elem.remove(child)
                patched.add(id(t_elem))
                return

        t_elems = list(slide_elem.iter(f"{{{a_ns}}}t"))
        for t_elem in t_elems:
            elem_id = id(t_elem)
            if elem_id in patched:
                continue
            if t_elem.text and t_elem.text.strip() == zone.text.strip():
                t_elem.text = new_text
                for child in list(t_elem):
                    t_elem.remove(child)
                patched.add(elem_id)
                return

    def _clear_zone_text(self, slide_elem, zone: TextZone, patched: set) -> None:
        a_ns = _NS["a"]

        target_shape = self._find_shape_by_name(slide_elem, zone.shape_name)
        if target_shape is not None:
            t_elem = self._locate_t_elem(target_shape, zone)
            if t_elem is not None and id(t_elem) not in patched:
                t_elem.text = ""
                for child in list(t_elem):
                    t_elem.remove(child)
                patched.add(id(t_elem))
                return

        t_elems = list(slide_elem.iter(f"{{{a_ns}}}t"))
        for t_elem in t_elems:
            elem_id = id(t_elem)
            if elem_id in patched:
                continue
            if t_elem.text and t_elem.text.strip() == zone.text.strip():
                t_elem.text = ""
                for child in list(t_elem):
                    t_elem.remove(child)
                patched.add(elem_id)
                return

    def _locate_t_elem(self, shape_elem, zone: TextZone):
        a_ns = _NS["a"]
        txBody = shape_elem.find(f"{{{a_ns}}}txBody")
        if txBody is None:
            for spPr in shape_elem.iter(f"{{{_NS['p']}}}spPr"):
                parent = spPr.getparent()
                if parent is not None:
                    txBody = parent.find(f"{{{a_ns}}}txBody")
                    if txBody is not None:
                        break
            if txBody is None:
                return None

        paras = list(txBody.findall(f"{{{a_ns}}}p"))
        if zone.para_index < len(paras):
            para = paras[zone.para_index]
            runs = list(para.findall(f"{{{a_ns}}}r"))
            if zone.run_index < len(runs):
                return runs[zone.run_index].find(f"{{{a_ns}}}t")

        return None

    def _find_shape_by_name(self, slide_elem, shape_name: str):
        p_ns = _NS["p"]
        for sp in slide_elem.iter(f"{{{p_ns}}}sp"):
            nvSpPr = sp.find(f"{{{p_ns}}}nvSpPr")
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(f"{{{p_ns}}}cNvPr")
                if cNvPr is not None and cNvPr.get("name") == shape_name:
                    return sp
        for grp in slide_elem.iter(f"{{{p_ns}}}grpSp"):
            nvGrpSpPr = grp.find(f"{{{p_ns}}}nvGrpSpPr")
            if nvGrpSpPr is not None:
                cNvPr = nvGrpSpPr.find(f"{{{p_ns}}}cNvPr")
                if cNvPr is not None and cNvPr.get("name") == shape_name:
                    return grp
        return None

    def _inject_new_slide(self, prs, template_dna: SlideDNA, patch: dict[str, Any], dna: DesignDNA) -> None:
        template_slide_idx = None
        for i, sd in enumerate(dna.slides):
            if sd.slide_xml == template_dna.slide_xml:
                template_slide_idx = i
                break

        if template_slide_idx is not None and template_slide_idx < len(prs.slides):
            source_slide = prs.slides[template_slide_idx]
            slide_layout = source_slide.slide_layout
        else:
            slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(slide_layout)
        sp_tree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))

        for child in list(sp_tree):
            tag = child.tag
            if tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
                continue
            sp_tree.remove(child)

        template_root = etree.fromstring(template_dna.slide_xml)
        template_sp_tree = template_root.find(qn("p:cSld")).find(qn("p:spTree"))

        for child in list(template_sp_tree):
            tag = child.tag
            if tag == qn("p:nvGrpSpPr") or tag == qn("p:grpSpPr"):
                continue
            if tag in (qn("p:sp"), qn("p:grpSp"), qn("p:pic"), qn("p:cxnSp")):
                new_child = copy.deepcopy(child)
                sp_tree.append(new_child)

        self._copy_image_rels(slide, template_dna, dna)

        self._patch_slide(slide, template_dna, patch)

    def _copy_image_rels(self, new_slide, template_dna: SlideDNA, dna: DesignDNA) -> None:
        a_ns = _NS["a"]
        r_ns = _NS["r"]

        blip_rids = set()
        slide_elem = new_slide._element
        for blip in slide_elem.iter(f"{{{a_ns}}}blip"):
            embed = blip.get(f"{{{r_ns}}}embed")
            if embed:
                blip_rids.add(embed)

        if not blip_rids:
            return

        source_prs = Presentation(dna.source_path)
        template_slide_idx = template_dna.slide_index
        if template_slide_idx >= len(source_prs.slides):
            return

        source_slide = source_prs.slides[template_slide_idx]
        source_part = source_slide.part

        new_slide_part = new_slide.part

        for rId in blip_rids:
            try:
                rel = source_part.rels.get(rId)
                if rel is None:
                    continue
                if "image" not in rel.reltype:
                    continue

                image_part = rel.target_part

                existing_rid = None
                for existing_rid_key, existing_rel in new_slide_part.rels.items():
                    try:
                        if (existing_rel.reltype == rel.reltype
                                and existing_rel.target_part is image_part):
                            existing_rid = existing_rid_key
                            break
                    except Exception:
                        pass

                if not existing_rid:
                    pass

                if existing_rid:
                    if existing_rid != rId:
                        for blip in slide_elem.iter(f"{{{a_ns}}}blip"):
                            if blip.get(f"{{{r_ns}}}embed") == rId:
                                blip.set(f"{{{r_ns}}}embed", existing_rid)
                else:
                    new_rId = new_slide_part.relate_to(
                        image_part, rel.reltype
                    )
                    if new_rId != rId:
                        for blip in slide_elem.iter(f"{{{a_ns}}}blip"):
                            if blip.get(f"{{{r_ns}}}embed") == rId:
                                blip.set(f"{{{r_ns}}}embed", new_rId)

            except Exception:
                pass

    # ── ZIP-level extraction ──

    def _extract_slides_from_zip(self, z: zipfile.ZipFile, prs: Presentation, dna: DesignDNA) -> None:
        names = z.namelist()
        total_slides = len(prs.slides)

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            slide_path = f"ppt/slides/slide{slide_num}.xml"
            rels_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"

            slide_xml = z.read(slide_path) if slide_path in names else b""
            rels_xml = z.read(rels_path) if rels_path in names else b""

            rels_map: dict[str, str] = {}
            if rels_xml:
                rels_root = etree.fromstring(rels_xml)
                for rel in rels_root.findall(f"{{{_NS['rel']}}}Relationship"):
                    rId = rel.get("Id", "")
                    target = rel.get("Target", "")
                    rels_map[rId] = target

            text_zones = self._extract_text_zones(slide, idx, total_slides)
            image_refs = self._extract_image_refs(slide, rels_map)
            group_xmls = self._extract_group_xmls(slide)
            smartart_refs = self._extract_smartart_refs(slide, rels_map)
            bg_colors = self._extract_slide_bg_colors(slide)
            layout_name = ""
            layout_index = 0
            try:
                layout = slide.slide_layout
                layout_name = layout.name
                for li, lo in enumerate(prs.slide_layouts):
                    if lo.name == layout_name:
                        layout_index = li
                        break
            except Exception:
                pass

            notes_text = ""
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

            has_bg_image = any(ref.get("is_background", False) for ref in image_refs)

            slide_dna = SlideDNA(
                slide_index=idx,
                page_type="unknown",
                slide_xml=slide_xml,
                rels_xml=rels_xml,
                rels_map=rels_map,
                text_zones=text_zones,
                image_refs=image_refs,
                group_xmls=group_xmls,
                smartart_refs=smartart_refs,
                layout_name=layout_name,
                layout_index=layout_index,
                has_background_image=has_bg_image,
                background_colors=bg_colors,
                notes_text=notes_text,
                brand_spec=None,
            )
            dna.slides.append(slide_dna)

    def _extract_text_zones(self, slide, slide_idx: int, total_slides: int) -> list[TextZone]:
        zones: list[TextZone] = []
        slide_elem = slide._element

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            shape_name = shape.name
            bounds = self._get_shape_bounds(shape)
            is_title = self._is_title_shape(shape, slide_idx)
            is_subtitle = self._is_subtitle_shape(shape)

            for pi, para in enumerate(shape.text_frame.paragraphs):
                para_text = para.text.strip()
                if not para_text:
                    continue

                for ri, run in enumerate(para.runs):
                    run_text = run.text.strip()
                    if not run_text:
                        continue

                    font_name = run.font.name
                    font_size = None
                    if run.font.size:
                        font_size = run.font.size / 12700
                    bold = run.font.bold or False
                    color_hex = None
                    try:
                        if run.font.color and run.font.color.rgb:
                            color_hex = f"#{run.font.color.rgb}"
                    except Exception:
                        pass

                    role = self._infer_text_role(
                        run_text, is_title, is_subtitle, pi, ri,
                        slide_idx, total_slides, font_size, shape_name,
                    )

                    xpath = self._compute_xpath(run._r, slide_elem)

                    zone_id = f"{shape_name}_p{pi}_r{ri}"

                    zones.append(TextZone(
                        zone_id=zone_id,
                        role=role,
                        text=run_text,
                        font_name=font_name,
                        font_size_pt=font_size,
                        bold=bold,
                        color_hex=color_hex,
                        xpath=xpath,
                        run_index=ri,
                        para_index=pi,
                        shape_name=shape_name,
                        bounds=bounds,
                    ))

        return zones

    def _infer_text_role(
        self, text: str, is_title: bool, is_subtitle: bool,
        para_idx: int, run_idx: int, slide_idx: int, total_slides: int,
        font_size: float | None, shape_name: str,
    ) -> str:
        text_lower = text.lower().strip()

        if any(kw in text_lower for kw in {"your logo", "logo", "公司标志"}):
            return "logo_text"

        if any(kw in text_lower for kw in {"汇报", "报告人", "presenter", "汇报人"}):
            return "badge"

        if re.match(r"^\d+\s*/\s*\d+$", text) or re.match(r"^\d+$", text) and len(text) <= 3:
            return "page_number"

        if any(kw in text_lower for kw in {"壹", "贰", "叁", "肆", "伍", "陆", "柒", "捌", "玖", "拾"}):
            return "badge"

        if any(kw in text_lower for kw in {"part", "chapter"}):
            return "badge"

        if is_title:
            return "title"

        if is_subtitle:
            return "subtitle"

        if font_size and font_size >= 28:
            return "title"

        if font_size and 16 <= font_size < 28 and para_idx == 0:
            return "subtitle"

        if any(kw in text_lower for kw in {"单击此处", "请输入", "click", "placeholder"}):
            return "body"

        return "body"

    def _is_title_shape(self, shape, slide_idx: int) -> bool:
        try:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (1, 3):
                    return True
        except (ValueError, AttributeError):
            pass
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size >= Emu(Pt(28)):
                        return True
        if slide_idx == 0 and shape.top < Emu(Inches(2)):
            return True
        return False

    def _is_subtitle_shape(self, shape) -> bool:
        try:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == 4:
                    return True
        except (ValueError, AttributeError):
            pass
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and Pt(16) <= run.font.size < Pt(28):
                        return True
        return False

    def _compute_xpath(self, run_elem, slide_elem) -> str:
        parts: list[str] = []
        current = run_elem
        while current is not None and current != slide_elem:
            tag = current.tag
            if "}" in tag:
                tag = tag.split("}")[1]
            parts.append(tag)
            current = current.getparent()
        parts.reverse()
        return "/".join(parts[-8:])

    def _get_shape_bounds(self, shape) -> tuple[float, float, float, float]:
        try:
            return (
                round(shape.left / 914400, 2),
                round(shape.top / 914400, 2),
                round(shape.width / 914400, 2),
                round(shape.height / 914400, 2),
            )
        except Exception:
            return (0.0, 0.0, 0.0, 0.0)

    def _get_shape_bounds_from_xml(self, sp_elem) -> tuple[float, float, float, float] | None:
        a_ns = _NS["a"]
        xfrm = sp_elem.find(f".//{{{a_ns}}}xfrm")
        if xfrm is None:
            spPr = sp_elem.find(f"{{{_NS['p']}}}spPr")
            if spPr is not None:
                xfrm = spPr.find(f"{{{a_ns}}}xfrm")
        if xfrm is None:
            return None
        off = xfrm.find(f"{{{a_ns}}}off")
        ext = xfrm.find(f"{{{a_ns}}}ext")
        if off is None or ext is None:
            return None
        try:
            x = int(off.get("x", 0)) / 914400
            y = int(off.get("y", 0)) / 914400
            w = int(ext.get("cx", 0)) / 914400
            h = int(ext.get("cy", 0)) / 914400
            return (round(x, 2), round(y, 2), round(w, 2), round(h, 2))
        except (ValueError, TypeError):
            return None

    def _extract_image_refs(self, slide, rels_map: dict[str, str]) -> list[dict[str, Any]]:
        refs: list[dict[str, Any]] = []
        a_ns = _NS["a"]

        for shape in slide.shapes:
            if shape.shape_type == 13 or (hasattr(shape, 'image') and shape.shape_type == 13):
                try:
                    img_blob = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    img_filename = f"slide{slide.slide_index + 1}_img_{shape.shape_id}.{ext}"
                    img_path = os.path.join(self._temp_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(img_blob)

                    bounds = self._get_shape_bounds(shape)
                    is_bg = (bounds[2] >= 12.0 and bounds[3] >= 6.5)

                    refs.append({
                        "shape_id": shape.shape_id,
                        "path": img_path,
                        "content_type": shape.image.content_type,
                        "bounds": bounds,
                        "is_background": is_bg,
                    })
                except Exception:
                    pass

        for blip in slide._element.iter(f"{{{a_ns}}}blip"):
            embed = blip.get(f"{{{_NS['r']}}}embed")
            if embed and embed in rels_map:
                target = rels_map[embed]
                if target not in [r.get("target") for r in refs]:
                    refs.append({"rId": embed, "target": target, "is_background": False})

        return refs

    def _extract_group_xmls(self, slide) -> list[bytes]:
        xmls: list[bytes] = []
        p_ns = _NS["p"]
        for grp in slide._element.iter(f"{{{p_ns}}}grpSp"):
            xmls.append(etree.tostring(grp, xml_declaration=False, encoding="unicode").encode("utf-8"))
        return xmls

    def _extract_smartart_refs(self, slide, rels_map: dict[str, str]) -> list[dict[str, Any]]:
        refs: list[dict[str, Any]] = []
        dgm_ns = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
        r_ns = _NS["r"]

        for shape in slide.shapes:
            try:
                elem = shape._element
                rel_ids_elems = elem.findall(f".//{{{dgm_ns}}}relIds")
                for ri in rel_ids_elems:
                    ids = {}
                    for attr in ("dm", "lo", "qs", "cs"):
                        val = ri.get(f"{{{r_ns}}}{attr}")
                        if val:
                            ids[attr] = val
                    if ids:
                        refs.append({
                            "shape_id": shape.shape_id,
                            "rel_ids": ids,
                            "dm_target": rels_map.get(ids.get("dm", ""), ""),
                        })
            except Exception:
                pass
        return refs

    def _extract_slide_bg_colors(self, slide) -> list[str]:
        colors: list[str] = []
        a_ns = _NS["a"]
        try:
            bg_elem = slide.background._element
            for srgb in bg_elem.iter(f"{{{a_ns}}}srgbClr"):
                val = srgb.get("val", "")
                if val:
                    colors.append(f"#{val}")
        except Exception:
            pass
        return colors

    # ── Theme extraction ──

    def _extract_theme_xml(self, prs: Presentation) -> bytes:
        try:
            theme_part = prs.part.part_related_by(RT.THEME)
            return theme_part.blob
        except Exception:
            return b""

    def _extract_theme_scheme(self, prs: Presentation) -> tuple[dict[str, str], dict[str, str], dict[str, str]]:
        colors: dict[str, str] = {}
        fonts: dict[str, str] = {}
        cjk_fonts: dict[str, str] = {}

        try:
            theme_part = prs.part.part_related_by(RT.THEME)
            theme_element = etree.fromstring(theme_part.blob)
            theme_elements = theme_element.find(qn("a:themeElements"))
            if theme_elements is None:
                return colors, fonts, cjk_fonts

            clr_scheme = theme_elements.find(qn("a:clrScheme"))
            if clr_scheme is not None:
                name_map = {
                    "dk1": "foreground", "lt1": "background",
                    "dk2": "primary", "accent1": "accent",
                    "accent2": "secondary", "accent3": "tertiary",
                    "accent4": "quaternary", "accent5": "quinary",
                    "accent6": "senary",
                    "hlink": "hyperlink", "folHlink": "followed_hyperlink",
                }
                for xml_name, role in name_map.items():
                    elem = clr_scheme.find(qn(f"a:{xml_name}"))
                    if elem is not None:
                        colors[role] = self._extract_color_value(elem)

            font_scheme = theme_elements.find(qn("a:fontScheme"))
            if font_scheme is not None:
                for font_type, key in [("a:majorFont", "heading"), ("a:minorFont", "body")]:
                    font_elem = font_scheme.find(qn(font_type))
                    if font_elem is not None:
                        latin = font_elem.find(qn("a:latin"))
                        if latin is not None:
                            fonts[key] = latin.get("typeface", "")
                        ea = font_elem.find(qn("a:ea"))
                        if ea is not None:
                            typeface = ea.get("typeface", "")
                            if typeface and typeface != "+mj-ea" and typeface != "+mn-ea":
                                cjk_fonts[key] = typeface
        except Exception:
            pass

        return colors, fonts, cjk_fonts

    def _extract_color_value(self, clr_element) -> str:
        srgb = clr_element.find(qn("a:srgbClr"))
        if srgb is not None:
            val = srgb.get("val")
            if val:
                return f"#{val}"
            return "#000000"
        sys_clr = clr_element.find(qn("a:sysClr"))
        if sys_clr is not None:
            val = sys_clr.get("val", "")
            last_clr = sys_clr.get("lastClr", "")
            if last_clr:
                return f"#{last_clr}"
            return _SYS_CLR_FALLBACK.get(val, "#000000")
        return "#000000"

    # ── Actual usage statistics ──

    def _extract_actual_usage(self, prs: Presentation, dna: DesignDNA) -> None:
        a_ns = _NS["a"]
        color_counter: Counter = Counter()
        font_counter: Counter = Counter()
        size_counter: Counter = Counter()
        bg_colors: Counter = Counter()

        for slide in prs.slides:
            slide_elem = slide._element
            for srgb in slide_elem.iter(f"{{{a_ns}}}srgbClr"):
                val = srgb.get("val", "")
                if val and len(val) == 6:
                    color_counter[f"#{val}"] += 1

            bg_elem = slide_elem.find(f".//{{{a_ns}}}bgClr")
            if bg_elem is None:
                bg_elem = slide_elem.find(f".//{{{a_ns}}}bgRef")
            if bg_elem is not None:
                for srgb in bg_elem.iter(f"{{{a_ns}}}srgbClr"):
                    val = srgb.get("val", "")
                    if val and len(val) == 6:
                        bg_colors[f"#{val}"] += 1

            p_ns = _NS.get("p", "http://schemas.openxmlformats.org/presentationml/2006/main")
            r_ns = _NS.get("r", "http://schemas.openxmlformats.org/officeDocument/2006/relationships")
            cSld = slide_elem.find(f"{{{p_ns}}}cSld")
            if cSld is not None:
                bg = cSld.find(f"{{{p_ns}}}bg")
                if bg is not None:
                    for srgb in bg.iter(f"{{{a_ns}}}srgbClr"):
                        val = srgb.get("val", "")
                        if val and len(val) == 6:
                            bg_colors[f"#{val}"] += 1
                    for solid in bg.iter(f"{{{a_ns}}}solidFill"):
                        for srgb in solid.iter(f"{{{a_ns}}}srgbClr"):
                            val = srgb.get("val", "")
                            if val and len(val) == 6:
                                bg_colors[f"#{val}"] += 1

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                font_counter[run.font.name] += 1
                            if run.font.size:
                                pt = int(run.font.size / 12700)
                                size_counter[pt] += 1
                            try:
                                if run.font.color and run.font.color.rgb:
                                    color_counter[f"#{run.font.color.rgb}"] += 1
                            except Exception:
                                pass

        dna.actual_colors = dict(color_counter.most_common(50))
        dna.actual_fonts = dict(font_counter.most_common(20))
        dna.actual_font_sizes = dict(size_counter.most_common(20))
        if bg_colors:
            dna.actual_bg_colors = dict(bg_colors.most_common(10))

    # ── Page type classification ──

    def _classify_page_types(self, dna: DesignDNA) -> None:
        total = len(dna.slides)
        for slide_dna in dna.slides:
            idx = slide_dna.slide_index
            zones = slide_dna.text_zones
            all_text = " ".join(z.text for z in zones).lower()

            if idx == 0:
                slide_dna.page_type = "cover"
            elif idx == total - 1:
                title_zones = [z for z in zones if z.role == "title"]
                body_zones = [z for z in zones if z.role == "body"]
                has_large_title = any(z.font_size_pt and z.font_size_pt >= 100 for z in title_zones)
                large_title_count = sum(1 for z in title_zones if z.font_size_pt and z.font_size_pt >= 100)
                has_real_body = any(z.text for z in body_zones if not any(kw in z.text.lower() for kw in {"单击", "请输入", "xxx", "x"}))

                if any(kw in all_text for kw in _CTA_KEYWORDS):
                    slide_dna.page_type = "back_cover"
                elif has_large_title and large_title_count >= 3 and not has_real_body:
                    slide_dna.page_type = "back_cover"
                elif has_large_title and large_title_count >= 3 and len(body_zones) > len(title_zones):
                    slide_dna.page_type = "back_cover"
                elif not has_real_body and len(body_zones) <= 2:
                    slide_dna.page_type = "back_cover"
                else:
                    slide_dna.page_type = "content"
            elif any(kw in all_text for kw in _TOC_KEYWORDS):
                slide_dna.page_type = "toc"
            elif any(kw in all_text for kw in _TRANSITION_KEYWORDS):
                slide_dna.page_type = "transition"
            elif any(kw in all_text for kw in _COVER_KEYWORDS) and idx <= 1:
                slide_dna.page_type = "cover"
            else:
                slide_dna.page_type = "content"

    # ── Logo extraction ──

    def _extract_logo(self, z: zipfile.ZipFile, prs: Presentation, dna: DesignDNA) -> None:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    try:
                        name_lower = shape.name.lower()
                        if "logo" in name_lower or "标志" in name_lower:
                            dna.logo_blob = shape.image.blob
                            dna.logo_content_type = shape.image.content_type
                            return
                    except Exception:
                        pass

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    bounds = self._get_shape_bounds(shape)
                    if bounds[0] >= 10.0 and bounds[1] <= 1.0 and bounds[2] <= 2.0:
                        try:
                            dna.logo_blob = shape.image.blob
                            dna.logo_content_type = shape.image.content_type
                            return
                        except Exception:
                            pass

    # ── Image blobs ──

    def _extract_all_images(self, z: zipfile.ZipFile, dna: DesignDNA) -> None:
        names = z.namelist()
        for name in names:
            if name.startswith("ppt/media/"):
                try:
                    dna.image_blobs[name] = z.read(name)
                except Exception:
                    pass

    # ── Decorative groups ──

    def _extract_decorative_groups(self, prs: Presentation, dna: DesignDNA) -> None:
        from ppt_pro_max.enterprise.group_extractor import GroupExtractor
        ge = GroupExtractor()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'shapes'):
                    try:
                        grp_info = ge.extract(shape)
                        if grp_info:
                            dna.decorative_groups.append(grp_info)
                    except Exception:
                        pass

    # ── BrandSpec builder ──

    def _build_brand_spec(self, prs: Presentation, dna: DesignDNA) -> BrandSpec:
        theme_colors, theme_fonts, theme_cjk = self._extract_theme_scheme(prs)

        actual_colors = dict(dna.actual_colors) if dna.actual_colors else {}
        actual_fonts = dict(dna.actual_fonts) if dna.actual_fonts else {}
        actual_sizes = dict(dna.actual_font_sizes) if dna.actual_font_sizes else {}

        colors = self._map_actual_colors_to_roles(actual_colors, dna.actual_bg_colors)
        fonts = self._map_actual_fonts_to_roles(actual_fonts, actual_sizes, theme_fonts, theme_cjk)

        dark_mode = False
        bg = colors.get("background", "#FFFFFF")
        try:
            r, g, b = int(bg[1:3], 16), int(bg[3:5], 16), int(bg[5:7], 16)
            dark_mode = (r * 0.299 + g * 0.587 + b * 0.114) < 128
        except (ValueError, IndexError):
            pass

        return BrandSpec(
            source="design_dna",
            colors=colors or None,
            fonts=fonts or None,
            dark_mode=dark_mode,
            _dna_actual_colors=actual_colors or None,
        )

    def _map_actual_colors_to_roles(self, actual_colors: dict[str, int], actual_bg_colors: dict[str, int] | None = None) -> dict[str, str]:
        if not actual_colors:
            return {}

        sorted_colors = sorted(actual_colors.items(), key=lambda x: x[1], reverse=True)

        def _luminance(hex_color: str) -> float:
            try:
                r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
                return r * 0.299 + g * 0.587 + b * 0.114
            except (ValueError, IndexError):
                return 0.0

        result: dict[str, str] = {}

        bg_candidates = []
        if actual_bg_colors:
            bg_candidates = sorted(actual_bg_colors.items(), key=lambda x: x[1], reverse=True)
        if bg_candidates:
            result["background"] = bg_candidates[0][0]
        else:
            best_bg = None
            best_lum = -1
            for c, _ in sorted_colors:
                lum = _luminance(c)
                if lum > 200:
                    if best_bg is None or lum > best_lum:
                        best_bg = c
                        best_lum = lum
            result["background"] = best_bg or "#FFFFFF"

        dark_fg = None
        for c, _ in sorted_colors:
            lum = _luminance(c)
            if lum < 50:
                dark_fg = c
                break
        if dark_fg is None:
            darkest = min(actual_colors, key=_luminance)
            if _luminance(darkest) < 80:
                dark_fg = darkest
        result["foreground"] = dark_fg or "#000000"

        bg_lum = _luminance(result.get("background", "#FFFFFF"))
        dark_colors = [(c, cnt) for c, cnt in sorted_colors if _luminance(c) < 128]
        mid_colors = [(c, cnt) for c, cnt in sorted_colors if 128 <= _luminance(c) < 200]
        light_colors = [(c, cnt) for c, cnt in sorted_colors if _luminance(c) >= 200 and _luminance(c) < bg_lum - 10]

        if dark_colors:
            result["primary"] = dark_colors[0][0]
        if len(dark_colors) >= 2:
            result["secondary"] = dark_colors[1][0]
        if len(dark_colors) >= 3:
            result["foreground"] = dark_colors[2][0]
        elif "foreground" not in result:
            result["foreground"] = dark_colors[-1][0] if dark_colors else "#000000"

        if mid_colors:
            result["accent"] = mid_colors[0][0]
            if len(mid_colors) >= 2:
                result["tertiary"] = mid_colors[1][0]
        elif dark_colors:
            brightest_dark = max(dark_colors, key=lambda x: _luminance(x[0]))
            result["accent"] = brightest_dark[0]

        if light_colors:
            result["muted"] = light_colors[-1][0]
            if len(light_colors) >= 2:
                result["quaternary"] = light_colors[0][0]
        elif mid_colors:
            palest_mid = max(mid_colors, key=lambda x: _luminance(x[0]))
            result["muted"] = palest_mid[0]
        elif dark_colors:
            palest_dark = max(dark_colors, key=lambda x: _luminance(x[0]))
            result["muted"] = palest_dark[0]

        if "primary" in result and "on-primary" not in result:
            lum = _luminance(result["primary"])
            result["on-primary"] = "#FFFFFF" if lum < 128 else "#000000"

        return result

    def _map_actual_fonts_to_roles(
        self,
        actual_fonts: dict[str, int],
        actual_sizes: dict[int, int],
        theme_fonts: dict[str, str],
        theme_cjk: dict[str, str],
    ) -> dict[str, str]:
        if not actual_fonts:
            merged = dict(theme_fonts)
            if theme_cjk:
                merged["cjk_heading"] = theme_cjk.get("heading", "")
                merged["cjk_body"] = theme_cjk.get("body", "")
            return merged

        sorted_fonts = sorted(actual_fonts.items(), key=lambda x: x[1], reverse=True)

        result: dict[str, str] = {}

        if len(sorted_fonts) >= 1:
            result["body"] = sorted_fonts[0][0]
            result["cjk_body"] = sorted_fonts[0][0]
        if len(sorted_fonts) >= 2:
            result["heading"] = sorted_fonts[1][0]
            result["cjk_heading"] = sorted_fonts[1][0]
        if len(sorted_fonts) >= 3:
            result["display"] = sorted_fonts[2][0]
            result["cjk_display"] = sorted_fonts[2][0]

        if "heading" not in result and "body" in result:
            result["heading"] = result["body"]
            result["cjk_heading"] = result.get("cjk_body", result["body"])

        avg_sizes = {}
        for sz_pt, cnt in actual_sizes.items():
            band = (sz_pt // 10) * 10
            avg_sizes.setdefault(band, 0)
            avg_sizes[band] += cnt

        if avg_sizes:
            max_band = max(avg_sizes, key=avg_sizes.get)
            if max_band >= 30 and len(sorted_fonts) >= 2:
                result["heading"] = sorted_fonts[1][0]
                result["cjk_heading"] = sorted_fonts[1][0]
            elif max_band >= 30 and len(sorted_fonts) == 1:
                pass

        return result
