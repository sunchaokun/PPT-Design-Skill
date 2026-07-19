"""DeliveryGate — quality check and auto-fix for template-clone PPTs.

13 checks across 5 categories:
  1-4  Content completeness (residual placeholders, title duplicates, TOC mismatch, blank pages)
  5-8  Design consistency (color break, font mismatch, background missing, decoration missing)
  9-10 Material completeness (image placeholder, broken image ref)
  11-12 Layout quality (text overflow, font too small)
  13   Page count match

Fatal items → auto-fix → re-check (max 2 rounds).
Warning items → report to user.
"""

from __future__ import annotations

import os
import re
from dataclasses import dataclass, field
from typing import Any

from pptx import Presentation

_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

_PLACEHOLDER_PATTERNS = [
    re.compile(r"单击此处", re.IGNORECASE),
    re.compile(r"请输入", re.IGNORECASE),
    re.compile(r"click\s+to", re.IGNORECASE),
    re.compile(r"your\s+logo", re.IGNORECASE),
    re.compile(r"^xxx$", re.IGNORECASE),
    re.compile(r"汇报[：:]\s*x+", re.IGNORECASE),
    re.compile(r"placeholder", re.IGNORECASE),
    re.compile(r"请添加", re.IGNORECASE),
    re.compile(r"此处添加", re.IGNORECASE),
    re.compile(r"输入你的正文", re.IGNORECASE),
    re.compile(r"enter\s+text\s+here", re.IGNORECASE),
    re.compile(r"insert\s+title", re.IGNORECASE),
    re.compile(r"your\s+title\s+here", re.IGNORECASE),
    re.compile(r"add\s+your\s+content", re.IGNORECASE),
    re.compile(r"type\s+something", re.IGNORECASE),
    re.compile(r"double\s+click\s+to\s+edit", re.IGNORECASE),
]

_PLACEHOLDER_IMAGE_NAMES = {"placeholder", "stock", "sample", "demo", "default"}

_MIN_FONT_PT = 11
_MIN_FONT_HUNDREDTHS = 1100


@dataclass
class CheckItem:
    category: str
    check_id: str
    severity: str
    slide_index: int
    message: str
    detail: str = ""
    auto_fixable: bool = False


@dataclass
class QualityReport:
    total_slides: int
    total_checks: int
    passed: int
    fatals: list[CheckItem] = field(default_factory=list)
    warnings: list[CheckItem] = field(default_factory=list)
    auto_fixed: list[CheckItem] = field(default_factory=list)

    @property
    def is_passable(self) -> bool:
        return len(self.fatals) == 0


class DeliveryGate:

    _ALL_CHECK_IDS = {
        "residual_placeholder",
        "title_duplicate",
        "toc_mismatch",
        "blank_page",
        "color_break",
        "font_mismatch",
        "background_missing",
        "decoration_missing",
        "image_placeholder",
        "broken_image_ref",
        "text_overflow",
        "font_too_small",
        "page_count_mismatch",
        "component_placeholder_residual",
    }

    def check(
        self,
        pptx_path: str,
        dna: Any,
        plans: list[Any],
    ) -> QualityReport:
        if not os.path.isfile(pptx_path):
            return QualityReport(
                total_slides=0, total_checks=0, passed=0,
                fatals=[CheckItem("content", "file_missing", "fatal", -1, "Output file not found")],
            )

        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)

        template_colors = set()
        template_fonts = set()
        if dna:
            for color in (dna.actual_colors or {}):
                template_colors.add(color.upper())
            for font in (dna.actual_fonts or {}):
                template_fonts.add(font)

        all_items: list[CheckItem] = []

        for idx, slide in enumerate(prs.slides):
            slide_elem = slide._element
            all_text = self._extract_all_text(slide_elem)

            all_items.extend(self._check_residual_placeholders(idx, all_text))
            all_items.extend(self._check_title_duplicate(idx, slide_elem, dna, idx))
            all_items.extend(self._check_blank_page(idx, all_text))
            all_items.extend(self._check_color_break(idx, slide_elem, template_colors))
            all_items.extend(self._check_font_mismatch(idx, slide_elem, template_fonts))
            all_items.extend(self._check_background_missing(idx, slide, dna, idx))
            all_items.extend(self._check_decoration_missing(idx, slide_elem, dna))
            all_items.extend(self._check_image_placeholder(idx, slide))
            all_items.extend(self._check_text_overflow(idx, slide))
            all_items.extend(self._check_font_too_small(idx, slide_elem))
            all_items.extend(self._check_broken_image_ref(idx, slide_elem, prs))
            all_items.extend(self._check_component_quality(idx, slide_elem))

        all_items.extend(self._check_page_count(total_slides, len(plans)))
        all_items.extend(self._check_toc_mismatch(prs, plans))

        fatals = [i for i in all_items if i.severity == "fatal"]
        warnings = [i for i in all_items if i.severity == "warning"]
        passed = len(all_items) - len(fatals) - len(warnings)

        return QualityReport(
            total_slides=total_slides,
            total_checks=len(all_items),
            passed=max(0, passed),
            fatals=fatals,
            warnings=warnings,
        )

    def auto_fix(
        self,
        pptx_path: str,
        dna: Any,
        plans: list[Any],
        report: QualityReport,
    ) -> None:
        if not os.path.isfile(pptx_path):
            return

        prs = Presentation(pptx_path)
        fixed_items: list[CheckItem] = []

        for item in report.fatals + [w for w in report.warnings if w.auto_fixable]:
            if not item.auto_fixable:
                continue

            if item.check_id == "residual_placeholder":
                if self._fix_residual_placeholders(prs, item.slide_index):
                    fixed_items.append(CheckItem(
                        item.category, item.check_id, "auto_fixed",
                        item.slide_index, "Cleared residual placeholder",
                        item.detail, True,
                    ))

            elif item.check_id == "title_duplicate":
                if self._fix_title_duplicate(prs, item.slide_index, dna):
                    fixed_items.append(CheckItem(
                        item.category, item.check_id, "auto_fixed",
                        item.slide_index, "Cleared duplicate title",
                        item.detail, True,
                    ))

            elif item.check_id == "background_missing":
                if self._fix_background_missing(prs, item.slide_index, dna):
                    fixed_items.append(CheckItem(
                        item.category, item.check_id, "auto_fixed",
                        item.slide_index, "Restored background from template",
                        item.detail, True,
                    ))

            elif item.check_id == "font_too_small":
                if self._fix_font_too_small(prs, item.slide_index):
                    fixed_items.append(CheckItem(
                        item.category, item.check_id, "auto_fixed",
                        item.slide_index, f"Increased font size to {_MIN_FONT_PT}pt",
                        item.detail, True,
                    ))

        prs.save(pptx_path)
        report.auto_fixed.extend(fixed_items)

    def format_report(self, report: QualityReport) -> str:
        lines = [
            "=" * 60,
            "  PPT Quality Report",
            "=" * 60,
            f"Total slides: {report.total_slides} | Checks: {report.total_checks} | "
            f"Passed: {report.passed} | Fatals: {len(report.fatals)} | Warnings: {len(report.warnings)}",
            "",
        ]

        categories = {}
        for item in report.fatals + report.warnings:
            categories.setdefault(item.category, []).append(item)

        for cat, items in categories.items():
            has_fatal = any(i.severity == "fatal" for i in items)
            has_warning = any(i.severity == "warning" for i in items)
            if has_fatal:
                lines.append(f"  {cat}: FATAL")
            elif has_warning:
                lines.append(f"  {cat}: WARNING")
            for item in items:
                prefix = "FATAL" if item.severity == "fatal" else "WARN"
                lines.append(f"    - Slide {item.slide_index + 1}: [{prefix}] {item.message}")
                if item.detail:
                    lines.append(f"      Detail: {item.detail}")

        if report.auto_fixed:
            lines.append("")
            lines.append(f"Auto-fixed: {len(report.auto_fixed)} items")
            for item in report.auto_fixed:
                lines.append(f"  - Slide {item.slide_index + 1}: {item.message}")

        lines.append("")
        if report.is_passable:
            lines.append("Result: PASS - ready for delivery")
        else:
            lines.append("Result: FAIL - fatal issues remain")

        return "\n".join(lines)

    # ── Check implementations ──

    def _extract_all_text(self, slide_elem) -> list[str]:
        a_ns = _NS["a"]
        texts = []
        for t in slide_elem.iter(f"{{{a_ns}}}t"):
            if t.text:
                texts.append(t.text.strip())
        return texts

    def _check_residual_placeholders(self, slide_idx: int, all_text: list[str]) -> list[CheckItem]:
        items = []
        for text in all_text:
            for pattern in _PLACEHOLDER_PATTERNS:
                if pattern.search(text):
                    items.append(CheckItem(
                        "content", "residual_placeholder", "fatal", slide_idx,
                        "Residual placeholder text found",
                        detail=text, auto_fixable=True,
                    ))
                    break
        return items

    def _check_title_duplicate(self, slide_idx: int, slide_elem, dna: Any, original_idx: int) -> list[CheckItem]:
        items = []
        if not dna:
            return items

        a_ns = _NS["a"]
        title_texts = []
        for t in slide_elem.iter(f"{{{a_ns}}}t"):
            if t.text and t.text.strip():
                parent_run = t.getparent()
                if parent_run is not None:
                    rPr = parent_run.find(f"{{{a_ns}}}rPr")
                    if rPr is not None:
                        sz = rPr.get("sz")
                        if sz and int(sz) >= 2800:
                            title_texts.append(t.text.strip())

        if len(title_texts) >= 2:
            from collections import Counter
            counts = Counter(title_texts)
            for text, count in counts.items():
                if count >= 2:
                    items.append(CheckItem(
                        "content", "title_duplicate", "fatal", slide_idx,
                        f"Duplicate title text ({count}x)",
                        detail=text, auto_fixable=True,
                    ))
        return items

    def _check_blank_page(self, slide_idx: int, all_text: list[str]) -> list[CheckItem]:
        non_empty = [t for t in all_text if t and len(t) > 1]
        if not non_empty:
            return [CheckItem(
                "content", "blank_page", "fatal", slide_idx,
                "Blank page — no text content",
                auto_fixable=False,
            )]
        return []

    def _check_color_break(self, slide_idx: int, slide_elem, template_colors: set) -> list[CheckItem]:
        if not template_colors:
            return []

        a_ns = _NS["a"]
        items = []
        for srgb in slide_elem.iter(f"{{{a_ns}}}srgbClr"):
            val = srgb.get("val", "").upper()
            if val and val not in template_colors and len(val) == 6:
                items.append(CheckItem(
                    "design", "color_break", "warning", slide_idx,
                    "Non-template color found",
                    detail=f"#{val}", auto_fixable=True,
                ))
        return items[:3]

    def _check_font_mismatch(self, slide_idx: int, slide_elem, template_fonts: set) -> list[CheckItem]:
        if not template_fonts:
            return []

        a_ns = _NS["a"]
        items = []
        for latin in slide_elem.iter(f"{{{a_ns}}}latin"):
            typeface = latin.get("typeface", "")
            if typeface and typeface not in template_fonts and not typeface.startswith("+"):
                items.append(CheckItem(
                    "design", "font_mismatch", "warning", slide_idx,
                    "Non-template font found",
                    detail=typeface, auto_fixable=True,
                ))
        return items[:3]

    def _check_background_missing(self, slide_idx: int, slide, dna: Any, original_idx: int) -> list[CheckItem]:
        if not dna or not dna.slides:
            return []

        has_template_bg = any(
            s.background_colors for s in dna.slides
        )
        if not has_template_bg:
            return []

        try:
            bg_elem = slide.background._element
            a_ns = _NS["a"]
            has_bg_color = bool(list(bg_elem.iter(f"{{{a_ns}}}srgbClr")))
            if not has_bg_color:
                return [CheckItem(
                    "design", "background_missing", "fatal", slide_idx,
                    "Background color missing on cloned slide",
                    auto_fixable=True,
                )]
        except Exception:
            pass
        return []

    def _check_font_too_small(self, slide_idx: int, slide_elem) -> list[CheckItem]:
        a_ns = _NS["a"]
        items = []
        for tag_name in (f"{{{a_ns}}}rPr", f"{{{a_ns}}}defRPr"):
            for rPr in slide_elem.iter(tag_name):
                sz = rPr.get("sz")
                if sz:
                    try:
                        pt = int(sz) / 100
                        if pt < _MIN_FONT_PT:
                            items.append(CheckItem(
                                "layout", "font_too_small", "warning", slide_idx,
                                f"Font size too small ({pt:.0f}pt)",
                                auto_fixable=True,
                            ))
                    except ValueError:
                        pass
        return items[:3]

    def _check_broken_image_ref(self, slide_idx: int, slide_elem, prs) -> list[CheckItem]:
        items = []
        a_ns = _NS["a"]
        r_ns = _NS["r"]

        for blip in slide_elem.iter(f"{{{a_ns}}}blip"):
            embed = blip.get(f"{{{r_ns}}}embed")
            if embed:
                try:
                    slide_part = prs.slides[slide_idx].part
                    rel = slide_part.rels.get(embed)
                    if rel is None:
                        items.append(CheckItem(
                            "material", "broken_image_ref", "fatal", slide_idx,
                            "Broken image reference",
                            detail=embed, auto_fixable=False,
                        ))
                except Exception:
                    pass
        return items

    def _check_text_overflow(self, slide_idx: int, slide) -> list[CheckItem]:
        items = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if not tf.word_wrap:
                continue
            text = tf.text
            if not text or not text.strip():
                continue
            try:
                width_inches = shape.width / 914400
                height_inches = shape.height / 914400
                total_est_lines = 0
                for p in tf.paragraphs:
                    if not p.text or not p.text.strip():
                        continue
                    font_size_pt = 12
                    if p.font.size:
                        font_size_pt = p.font.size / 12700
                    else:
                        for r in p.runs:
                            if r.font.size:
                                font_size_pt = r.font.size / 12700
                                break
                    chars_per_line = max(1, int(width_inches * 72 / font_size_pt * 0.65))
                    total_est_lines += max(1, len(p.text) / chars_per_line)
                if total_est_lines == 0:
                    continue
                avg_font_pt = 12
                for p in tf.paragraphs:
                    if p.font.size:
                        avg_font_pt = p.font.size / 12700
                        break
                    for r in p.runs:
                        if r.font.size:
                            avg_font_pt = r.font.size / 12700
                            break
                    if avg_font_pt != 12:
                        break
                line_height_in = avg_font_pt * 1.4 / 72
                max_lines = height_inches / line_height_in if line_height_in > 0 else 999
                if total_est_lines > max_lines * 2.0:
                    items.append(CheckItem(
                        "layout", "text_overflow", "warning", slide_idx,
                        "Text may overflow textbox",
                        detail=f"est_lines={total_est_lines:.0f} vs max={max_lines:.0f}",
                        auto_fixable=False,
                    ))
            except Exception:
                pass
        return items[:3]

    def _check_decoration_missing(self, slide_idx: int, slide_elem, dna: Any) -> list[CheckItem]:
        if not dna or not dna.slides:
            return []

        template_deco_count = 0
        for s in dna.slides:
            template_deco_count += len([z for z in s.text_zones if z.role == "decoration"])
        if template_deco_count == 0:
            return []

        a_ns = _NS["a"]
        p_ns = _NS["p"]
        output_deco_count = 0
        sp_tree = slide_elem.find(f"{{{p_ns}}}cSld")
        if sp_tree is not None:
            sp_tree = sp_tree.find(f"{{{p_ns}}}spTree")
        if sp_tree is not None:
            for sp in sp_tree:
                tag = sp.tag.split("}")[-1] if "}" in sp.tag else sp.tag
                if tag in ("sp", "pic"):
                    has_text = False
                    for t in sp.iter(f"{{{a_ns}}}t"):
                        if t.text and t.text.strip():
                            has_text = True
                            break
                    if not has_text:
                        output_deco_count += 1

        if output_deco_count == 0:
            return [CheckItem(
                "design", "decoration_missing", "warning", slide_idx,
                "Template has decorations but output slide has none",
                auto_fixable=False,
            )]
        return []

    def _check_image_placeholder(self, slide_idx: int, slide) -> list[CheckItem]:
        items = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                name = shape.name.lower()
                if any(p in name for p in _PLACEHOLDER_IMAGE_NAMES):
                    items.append(CheckItem(
                        "material", "image_placeholder", "warning", slide_idx,
                        "Placeholder/stock image detected",
                        detail=shape.name, auto_fixable=False,
                    ))
        return items

    def _check_component_quality(self, slide_idx: int, slide_elem) -> list[CheckItem]:
        items = []
        a_ns = _NS["a"]
        p_ns = _NS["p"]
        sp_tree = slide_elem.find(f"{{{p_ns}}}cSld")
        if sp_tree is None:
            return items
        sp_tree = sp_tree.find(f"{{{p_ns}}}spTree")
        if sp_tree is None:
            return items

        for grp in sp_tree.findall(f"{{{p_ns}}}grpSp"):
            for t in grp.iter(f"{{{a_ns}}}t"):
                if t.text and t.text.strip():
                    text = t.text.strip()
                    for pattern in _PLACEHOLDER_PATTERNS:
                        if pattern.search(text):
                            items.append(CheckItem(
                                "design", "component_placeholder_residual", "warning", slide_idx,
                                "GroupShape contains residual placeholder text",
                                detail=text, auto_fixable=True,
                            ))
                            break
        return items[:3]

    def _check_page_count(self, total_slides: int, planned_count: int) -> list[CheckItem]:
        if total_slides != planned_count and planned_count > 0:
            return [CheckItem(
                "count", "page_count_mismatch", "fatal", -1,
                f"Page count mismatch: {total_slides} actual vs {planned_count} planned",
                auto_fixable=False,
            )]
        return []

    def _check_toc_mismatch(self, prs, plans: list[Any]) -> list[CheckItem]:
        items = []
        toc_slide = None
        toc_idx = -1
        for idx, slide in enumerate(prs.slides):
            a_ns = _NS["a"]
            for t in slide._element.iter(f"{{{a_ns}}}t"):
                if t.text and any(kw in t.text.lower() for kw in ("目录", "contents", "agenda")):
                    toc_slide = slide
                    toc_idx = idx
                    break
            if toc_slide is not None:
                break

        if toc_slide is None or toc_idx == -1:
            return items

        toc_text = []
        a_ns = _NS["a"]
        for t in toc_slide._element.iter(f"{{{a_ns}}}t"):
            if t.text and t.text.strip():
                toc_text.append(t.text.strip())

        content_sections = sum(1 for p in plans if getattr(p, 'page_type', '') == "transition")

        if content_sections > 0 and len(toc_text) > 1 and content_sections != len(toc_text) - 1:
            items.append(CheckItem(
                "content", "toc_mismatch", "warning", toc_idx,
                f"TOC lists {len(toc_text) - 1} items but content has {content_sections} sections",
                auto_fixable=False,
            ))
        return items

    # ── Auto-fix implementations ──

    def _fix_residual_placeholders(self, prs, slide_idx: int) -> bool:
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            return False

        slide_elem = prs.slides[slide_idx]._element
        a_ns = _NS["a"]
        fixed = False

        for t in list(slide_elem.iter(f"{{{a_ns}}}t")):
            if not t.text:
                continue
            text = t.text.strip()
            for pattern in _PLACEHOLDER_PATTERNS:
                if pattern.search(text):
                    t.text = ""
                    for child in list(t):
                        t.remove(child)
                    fixed = True
                    break
        return fixed

    def _fix_title_duplicate(self, prs, slide_idx: int, dna: Any) -> bool:
        if slide_idx < 0 or slide_idx >= len(prs.slides) or not dna:
            return False

        slide_elem = prs.slides[slide_idx]._element
        a_ns = _NS["a"]

        title_runs = []
        for r in slide_elem.iter(f"{{{a_ns}}}r"):
            rPr = r.find(f"{{{a_ns}}}rPr")
            if rPr is not None:
                sz = rPr.get("sz")
                if sz and int(sz) >= 2800:
                    t = r.find(f"{{{a_ns}}}t")
                    if t is not None and t.text:
                        title_runs.append(t)

        if len(title_runs) < 2:
            return False

        from collections import Counter
        text_counts = Counter(t.text.strip() for t in title_runs)
        fixed = False
        for text, count in text_counts.items():
            if count >= 2:
                seen = False
                for t in title_runs:
                    if t.text and t.text.strip() == text:
                        if seen:
                            t.text = ""
                            for child in list(t):
                                t.remove(child)
                            fixed = True
                        else:
                            seen = True
        return fixed

    def _fix_background_missing(self, prs, slide_idx: int, dna: Any) -> bool:
        if slide_idx < 0 or slide_idx >= len(prs.slides) or not dna:
            return False

        bg_color = None
        for s in dna.slides:
            if s.background_colors:
                bg_color = s.background_colors[0]
                break
        if not bg_color:
            return False

        try:
            from pptx.dml.color import RGBColor
            slide = prs.slides[slide_idx]
            bg = slide.background
            fill = bg.fill
            fill.solid()
            first_color = bg_color.lstrip("#")
            fill.fore_color.rgb = RGBColor.from_string(first_color)
            return True
        except Exception:
            return False

    def _fix_font_too_small(self, prs, slide_idx: int) -> bool:
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            return False

        slide_elem = prs.slides[slide_idx]._element
        a_ns = _NS["a"]
        fixed = False

        for tag_name in (f"{{{a_ns}}}rPr", f"{{{a_ns}}}defRPr"):
            for rPr in slide_elem.iter(tag_name):
                sz = rPr.get("sz")
                if sz:
                    try:
                        pt = int(sz) / 100
                        if pt < _MIN_FONT_PT:
                            rPr.set("sz", str(_MIN_FONT_HUNDREDTHS))
                            fixed = True
                    except ValueError:
                        pass
        return fixed
