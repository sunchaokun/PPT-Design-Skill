"""EnterpriseRenderer — template-driven PPT rendering."""

from __future__ import annotations

import os
from typing import Any

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.enterprise.slide_utils import remove_slide


class EnterpriseRenderer:

    def __init__(self, template_path: str | None = None):
        self._template_path = template_path

    def create_presentation(self, keep_slides: bool = False) -> Presentation:
        if self._template_path and os.path.exists(self._template_path):
            try:
                prs = Presentation(self._template_path)
                if not keep_slides:
                    while len(prs.slides) > 0:
                        remove_slide(prs, 0)
                return prs
            except Exception:
                return Presentation()
        return Presentation()

    def add_slide(
        self,
        prs: Presentation,
        layout_name: str | None = None,
    ):
        if not prs.slide_layouts:
            from pptx import Presentation as _Prs
            prs = _Prs()
        if layout_name:
            for layout in prs.slide_layouts:
                if layout.name == layout_name:
                    return prs.slides.add_slide(layout)
        for layout in prs.slide_layouts:
            if "blank" in layout.name.lower():
                continue
            if layout_name is None and "title" in layout.name.lower():
                return prs.slides.add_slide(layout)
        for layout in prs.slide_layouts:
            if "blank" not in layout.name.lower():
                return prs.slides.add_slide(layout)
        return prs.slides.add_slide(prs.slide_layouts[0])

    def insert_logo(
        self,
        slide,
        logo_path: str,
        logo_spec: dict[str, Any],
        current_goal: str | None = None,
        prs: Presentation | None = None,
    ) -> None:
        skip_slides = logo_spec.get("skip_slides", [])
        if current_goal and current_goal in skip_slides:
            return

        if not os.path.isfile(logo_path):
            return

        position = logo_spec.get("position", "top_right")
        width_inches = logo_spec.get("width_inches", 1.0)

        if prs is not None:
            slide_width = prs.slide_width
            slide_height = prs.slide_height
        else:
            slide_width = Inches(13.333)
            slide_height = Inches(7.5)
        width_emu = Inches(width_inches)

        try:
            with PILImage.open(logo_path) as img:
                img_w, img_h = img.size
                aspect = img_h / img_w if img_w > 0 else 0.5
        except Exception:
            aspect = 0.5
        height_emu = int(width_emu * aspect)

        if position == "top_right":
            left = slide_width - width_emu - Inches(0.3)
            top = Inches(0.3)
        elif position == "top_left":
            left = Inches(0.3)
            top = Inches(0.3)
        elif position == "bottom_right":
            left = slide_width - width_emu - Inches(0.3)
            top = slide_height - height_emu - Inches(0.3)
        else:
            left = slide_width - width_emu - Inches(0.3)
            top = Inches(0.3)

        slide.shapes.add_picture(logo_path, left, top, width=width_emu, height=height_emu)

    def save(self, prs: Presentation, output_path: str) -> None:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        prs.save(output_path)
