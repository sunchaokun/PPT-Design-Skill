"""Debug overlap test case."""
from ppt_pro_max.renderer.svg_compiler import SVGCompiler
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 200">'
svg += '<text x="50" y="50" font-size="14" fill="#000">First text</text>'
svg += '<text x="60" y="55" font-size="14" fill="#000">Overlapping text</text>'
svg += '</svg>'

result = SVGCompiler().compile(svg, slide, (1, 1, 6, 4))
print('warnings:', result.warnings)