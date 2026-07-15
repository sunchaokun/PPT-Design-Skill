"""EnterpriseRenderer — template-driven or free-form PPT rendering.

When a template.pptx is provided, strictly follow the template's layouts,
dimensions, and styles — the user chose that template for a reason.

When no template is provided, use LayoutRegistry's 12 master layouts on a
13.333"×7.5" widescreen canvas for maximum creative freedom.
"""

from __future__ import annotations

import os
from typing import Any

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.enterprise.slide_utils import remove_slide
from ppt_pro_max.renderer.layout_registry import SLIDE_WIDTH, SLIDE_HEIGHT


class EnterpriseRenderer:

    def __init__(self, template_path: str | None = None):
        self._template_path = template_path
        self._has_template = bool(template_path and os.path.exists(template_path))

    @property
    def has_template(self) -> bool:
        return self._has_template

    def create_presentation(self, keep_slides: bool = False) -> Presentation:
        if self._has_template:
            try:
                prs = Presentation(self._template_path)
                if not keep_slides:
                    while len(prs.slides) > 0:
                        remove_slide(prs, 0)
                return prs
            except Exception:
                pass
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        return prs

    def add_slide(
        self,
        prs: Presentation,
        layout_name: str | None = None,
    ):
        if self._has_template:
            if layout_name:
                for layout in prs.slide_layouts:
                    if layout.name == layout_name:
                        return prs.slides.add_slide(layout)
            for layout in prs.slide_layouts:
                if "blank" in layout.name.lower():
                    continue
                if layout_name is None and "title" in layout.name.lower():
                    return prs.slides.add_slide(layout)
            for layout in prs.slide_layouts:
                if "blank" not in layout.name.lower():
                    return prs.slides.add_slide(layout)
            return prs.slides.add_slide(prs.slide_layouts[0])
        blank_layout = None
        for layout in prs.slide_layouts:
            if "blank" in layout.name.lower():
                blank_layout = layout
                break
        if blank_layout is None and prs.slide_layouts:
            blank_layout = prs.slide_layouts[-1]
        if blank_layout is not None:
            return prs.slides.add_slide(blank_layout)
        return prs.slides.add_slide(prs.slide_layouts[0])

    def insert_logo(
        self,
        slide,
        logo_path: str,
        logo_spec: dict[str, Any],
        current_goal: str | None = None,
        prs: Presentation | None = None,
    ) -> None:
        skip_slides = logo_spec.get("skip_slides", [])
        if current_goal and current_goal in skip_slides:
            return

        if not os.path.isfile(logo_path):
            return

        position = logo_spec.get("position", "top_right")
        width_inches = logo_spec.get("width_inches", 1.0)

        if prs is not None:
            slide_width = prs.slide_width
            slide_height = prs.slide_height
        else:
            slide_width = Inches(13.333)
            slide_height = Inches(7.5)
        width_emu = Inches(width_inches)

        try:
            with PILImage.open(logo_path) as img:
                img_w, img_h = img.size
                aspect = img_h / img_w if img_w > 0 else 0.5
        except Exception:
            aspect = 0.5
        height_emu = int(width_emu * aspect)

        if position == "top_right":
            left = slide_width - width_emu - Inches(0.3)
            top = Inches(0.3)
        elif position == "top_left":
            left = Inches(0.3)
            top = Inches(0.3)
        elif position == "bottom_right":
            left = slide_width - width_emu - Inches(0.3)
            top = slide_height - height_emu - Inches(0.3)
        else:
            left = slide_width - width_emu - Inches(0.3)
            top = Inches(0.3)

        slide.shapes.add_picture(logo_path, left, top, width=width_emu, height=height_emu)

    def save(self, prs: Presentation, output_path: str) -> None:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        prs.save(output_path)

    def add_page_numbers(self, prs: Presentation, footer_config: dict[str, Any], brand_spec=None) -> None:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        total = len(prs.slides)
        show_page_number = footer_config.get("show_page_number", False)
        page_number_format = footer_config.get("page_number_format", "{n}")
        page_number_position = footer_config.get("page_number_position", "bottom_right")
        show_footer_text = footer_config.get("show_footer_text", False)
        footer_text = footer_config.get("footer_text", "")
        footer_position = footer_config.get("footer_position", "bottom_center")
        font_size_pt = footer_config.get("font_size_pt", 10)
        skip_pages = footer_config.get("skip_pages", [])

        if not show_page_number and not show_footer_text:
            return

        muted_color = RGBColor(0x99, 0x99, 0x99)
        color_role = footer_config.get("color_role", "muted-foreground")
        if brand_spec and brand_spec.colors:
            role_color = brand_spec.colors.get(color_role) or brand_spec.colors.get("muted-foreground")
            if role_color:
                try:
                    muted_color = RGBColor.from_string(role_color.lstrip("#"))
                except Exception:
                    pass

        position_map = {
            "bottom_left": Inches(0.9),
            "bottom_center": Inches(5.833),
            "bottom_right": Inches(11.433),
        }

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            is_cover = idx == 0
            should_skip = is_cover or slide_num in skip_pages

            if show_page_number and not should_skip:
                page_text = page_number_format.replace("{n}", str(slide_num)).replace("{total}", str(total))
                left = position_map.get(page_number_position, Inches(11.433))
                top = Inches(7.0)

                txBox = slide.shapes.add_textbox(left, top, Inches(2.0), Inches(0.3))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = page_text
                p.font.size = Pt(font_size_pt)
                p.font.color.rgb = muted_color
                p.alignment = PP_ALIGN.RIGHT

            if show_footer_text and footer_text and not should_skip:
                left = position_map.get(footer_position, Inches(5.833))
                top = Inches(7.0)

                txBox = slide.shapes.add_textbox(left, top, Inches(4.0), Inches(0.3))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = footer_text
                p.font.size = Pt(font_size_pt)
                p.font.color.rgb = muted_color
                p.alignment = PP_ALIGN.CENTER

    def add_watermark(self, prs: Presentation, watermark_config: dict[str, Any], brand_spec=None) -> None:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from lxml import etree

        text = watermark_config.get("text", "CONFIDENTIAL")
        opacity = watermark_config.get("opacity", 0.15)
        rotation = watermark_config.get("rotation", -45)
        font_size_pt = watermark_config.get("font_size_pt", 72)
        skip_pages = watermark_config.get("skip_pages", [1])

        muted_hex = "999999"
        color_role = watermark_config.get("color_role", "muted-foreground")
        if brand_spec and brand_spec.colors:
            role_color = brand_spec.colors.get(color_role) or brand_spec.colors.get("muted-foreground")
            if role_color:
                muted_hex = role_color.lstrip("#")

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            if slide_num in skip_pages:
                continue

            txBox = slide.shapes.add_textbox(
                Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5),
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size_pt)
            p.font.color.rgb = RGBColor.from_string(muted_hex)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            sp = txBox._element
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

            solidFill_el = sp.find(".//" + ns + "solidFill")
            if solidFill_el is not None:
                srgbClr_el = solidFill_el.find(ns + "srgbClr")
                if srgbClr_el is not None:
                    alpha_val = str(int(opacity * 100000))
                    alpha_elem = etree.SubElement(srgbClr_el, ns + "alpha")
                    alpha_elem.set("val", alpha_val)

            xfrm_el = sp.find(".//" + ns + "xfrm")
            if xfrm_el is not None:
                rot_val = str(int(rotation * 60000))
                xfrm_el.set("rot", rot_val)
