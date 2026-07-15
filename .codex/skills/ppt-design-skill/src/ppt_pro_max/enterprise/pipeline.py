"""EnterprisePipeline — main orchestration flow."""

from __future__ import annotations

import os
from typing import Any

from ppt_pro_max.enterprise.scanner import ProjectScanner
from ppt_pro_max.enterprise.template_analyzer import TemplateAnalyzer
from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.version_manager import next_version, read_meta, write_meta
from ppt_pro_max.enterprise.enterprise_decider import EnterpriseDesignDecider
from ppt_pro_max.enterprise.review_gate import ReviewGate
from ppt_pro_max.enterprise.content_parser import load_enterprise_content
from ppt_pro_max.enterprise.image_matcher import match_images, assign_images_by_size, auto_generate_image_prompts
from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
from ppt_pro_max.enterprise.density_profile import get_density_profile, apply_density_to_bullets


class EnterprisePipeline:

    def run_beautify(
        self,
        input_pptx: str,
        output_path: str | None = None,
        style: str | None = None,
        brand_spec: BrandSpec | None = None,
        density: int | None = None,
        motion: int | None = None,
        beautify_mode: str | None = None,
        component_library: str | None = None,
        **kwargs,
    ) -> dict[str, Any]:
        from ppt_pro_max.renderer.theme_composer import ThemeComposer

        if not os.path.isfile(input_pptx):
            return {"mode": "beautify", "error": "Input file not found", "num_slides": 0}

        if brand_spec is None:
            if style:
                composer = ThemeComposer()
                composed = composer.compose(style=style)
                brand_spec = BrandSpec(
                    source="theme_composer",
                    colors=composed.get("colors"),
                    fonts=composed.get("typography"),
                    dark_mode=composed.get("dark_mode", False),
                )
            else:
                brand_spec = BrandSpec(source="default")

        if output_path is None:
            base_dir = os.path.dirname(input_pptx)
            output_dir = os.path.join(base_dir, "output")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, "beautified.pptx")

        try:
            from pptx import Presentation as _Prs
            _tmp_prs = _Prs(input_pptx)
            num_slides = len(_tmp_prs.slides)
            del _tmp_prs
        except Exception:
            num_slides = 0

        effective_mode = beautify_mode or "full"

        if effective_mode == "light":
            return self._beautify_light(
                input_pptx, output_path, brand_spec, motion, num_slides
            )

        return self._beautify_full(
            input_pptx, output_path, brand_spec, motion, num_slides, component_library
        )

    def _beautify_light(
        self, input_pptx: str, output_path: str, brand_spec: BrandSpec,
        motion: int | None, num_slides: int,
    ) -> dict[str, Any]:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from lxml import etree
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        import copy
        import shutil

        shutil.copy2(input_pptx, output_path)
        precision = PrecisionRenderer(brand_spec=brand_spec)

        try:
            prs = Presentation(output_path)
        except Exception:
            return {"mode": "beautify", "beautify_mode": "light", "error": "Cannot open input PPT", "num_slides": 0}

        brand_colors = brand_spec.colors or {}
        heading_font = (brand_spec.fonts or {}).get("heading")
        body_font = (brand_spec.fonts or {}).get("body")

        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

        color_map = {}
        for role, hex_val in brand_colors.items():
            color_map[role] = hex_val.lstrip("#")

        scheme_to_brand = {
            "accent1": color_map.get("accent", color_map.get("primary", "")),
            "accent2": color_map.get("secondary", ""),
            "accent3": color_map.get("accent", ""),
            "accent4": color_map.get("accent", ""),
            "accent5": color_map.get("accent", ""),
            "accent6": color_map.get("accent", ""),
            "lt1": color_map.get("foreground", color_map.get("on-primary", "FFFFFF")),
            "dk1": color_map.get("foreground", "000000"),
            "lt2": color_map.get("muted-foreground", ""),
            "dk2": color_map.get("foreground", ""),
            "tx1": color_map.get("foreground", ""),
            "tx2": color_map.get("muted-foreground", ""),
            "bg1": color_map.get("background", ""),
            "bg2": color_map.get("muted", ""),
            "hlink": color_map.get("accent", ""),
            "folHlink": color_map.get("accent", ""),
        }

        for slide in prs.slides:
            slide_elem = slide._element

            for scheme in list(slide_elem.iter(f"{{{a_ns}}}schemeClr")):
                val = scheme.get("val", "")
                if val in scheme_to_brand and scheme_to_brand[val]:
                    parent = scheme.getparent()
                    if parent is not None:
                        srgb = etree.SubElement(parent, f"{{{a_ns}}}srgbClr")
                        srgb.set("val", scheme_to_brand[val])
                        for child in scheme:
                            srgb.append(copy.deepcopy(child))
                        parent.remove(scheme)

            if heading_font or body_font:
                for latin in slide_elem.iter(f"{{{a_ns}}}latin"):
                    typeface = latin.get("typeface", "")
                    if not typeface:
                        continue
                    if typeface == "+mj-lt" and heading_font:
                        latin.set("typeface", heading_font)
                    elif typeface == "+mn-lt" and body_font:
                        latin.set("typeface", body_font)

            try:
                bg = slide.background
                fill = bg.fill
                bg_color = color_map.get("background")
                if bg_color:
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(bg_color)
            except Exception:
                pass

        precision.save(prs, output_path)

        if motion and isinstance(motion, int) and motion > 0:
            for slide in prs.slides:
                self._apply_animations(slide, {"goal": "content"}, motion)

        return {
            "mode": "beautify",
            "beautify_mode": "light",
            "input_pptx": input_pptx,
            "output_path": output_path,
            "num_slides": num_slides,
            "render_mode": "inplace_restyle",
        }

    def _beautify_full(
        self, input_pptx: str, output_path: str, brand_spec: BrandSpec,
        motion: int | None, num_slides: int, component_library: str | None = None,
    ) -> dict[str, Any]:
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor
        from ppt_pro_max.enterprise.content_parser import infer_component_category
        from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path

        try:
            extractor = SlideExtractor()
            extracted_pages = extractor.extract(input_pptx)
        except Exception:
            extracted_pages = []

        if not extracted_pages:
            return self._beautify_light(input_pptx, output_path, brand_spec, motion, num_slides)

        for page in extracted_pages:
            bullets = page.get("bullets") or []
            if bullets and not page.get("component_type"):
                comp_type, comp_cat = infer_component_category(bullets)
                if comp_type:
                    page["component_type"] = comp_type
                    page["component_category"] = comp_cat

        component_lib = None
        db_path = find_db_path(component_library)
        if db_path:
            try:
                component_lib = ComponentLibrary(db_path)
            except Exception:
                component_lib = None

        precision = PrecisionRenderer(brand_spec=brand_spec)
        prs = precision.create_presentation()

        total = len(extracted_pages)
        for i, page in enumerate(extracted_pages):
            precision.render_slide(prs, page, component_lib=component_lib,
                                   page_index=i, total_pages=total)

        if component_lib:
            try:
                component_lib.close()
            except Exception:
                pass

        precision.save(prs, output_path)

        if motion and isinstance(motion, int) and motion > 0:
            for slide in prs.slides:
                self._apply_animations(slide, {"goal": "content"}, motion)

        return {
            "mode": "beautify",
            "beautify_mode": "full",
            "input_pptx": input_pptx,
            "output_path": output_path,
            "num_slides": len(prs.slides),
            "render_mode": "precision_rebuild",
        }

    def run(
        self,
        query: str,
        project_dir: str,
        dry_run: bool = False,
        business_mode: str | None = None,
        density: int | None = None,
        motion: int | None = None,
        review: bool = False,
        review_file: str | None = None,
        output_version: int | None = None,
        from_version: int | None = None,
        pages: str | None = None,
        history: bool = False,
        output: str | None = None,
        content_file: str | None = None,
        component_library: str | None = None,
        **kwargs,
    ) -> dict[str, Any]:
        if history:
            return self._handle_history(project_dir)

        scanner = ProjectScanner()
        asset = scanner.scan(project_dir)

        brand_spec = self._build_brand_spec(asset)

        if dry_run:
            return {
                "dry_run": True,
                "pipeline": "enterprise",
                "project": project_dir,
                "brand_spec": {
                    "source": brand_spec.source,
                    "colors": brand_spec.colors,
                    "fonts": brand_spec.fonts,
                },
                "assets": {
                    "template": asset.template_path is not None,
                    "logo": asset.logo_path is not None,
                    "brand_json": asset.brand_raw is not None,
                    "content_json": asset.content_raw is not None,
                    "image_pool_count": len(asset.image_pool),
                },
            }

        decider = EnterpriseDesignDecider(
            brand_spec=brand_spec,
            business_mode=business_mode,
        )
        effective_density = density if density is not None else decider.suggest_density()

        page_contents = self._load_content(asset, project_dir, content_file)

        story_plan = self._build_story_plan(query, page_contents, effective_density, business_mode)

        page_designs = self._build_page_designs(
            story_plan, decider, asset, brand_spec
        )

        if review:
            gate = ReviewGate()
            proposal = gate.generate_proposal(
                project_dir=project_dir,
                brand_spec={"source": brand_spec.source, "colors": brand_spec.colors, "fonts": brand_spec.fonts},
                story_plan=story_plan,
                page_designs=page_designs,
                assets={
                    "template": asset.template_path is not None,
                    "logo": asset.logo_path is not None,
                    "brand_json": asset.brand_raw is not None,
                    "content_json": asset.content_raw is not None,
                    "image_pool_count": len(asset.image_pool),
                },
            )
            proposal_path = review_file or os.path.join(project_dir, "output", "proposal.json")
            gate.write_proposal(proposal, proposal_path)
            return {
                "pipeline": "enterprise",
                "review": True,
                "proposal_path": proposal_path,
                "proposal": proposal,
            }

        output_dir = os.path.join(project_dir, "output")
        from_meta = None
        if output_version is not None:
            vnum = output_version
        elif from_version is not None:
            vnum = from_version
            from_meta = read_meta(os.path.join(output_dir, f"v{from_version}"))
        else:
            vnum = next_version(output_dir)

        if from_meta and from_meta.get("slides"):
            page_contents = self._load_content_from_meta(from_meta)

        version_dir = os.path.join(output_dir, f"v{vnum}")
        pptx_path = output or os.path.join(version_dir, "presentation.pptx")

        if pages:
            source_pptx = self._find_latest_pptx(output_dir, vnum)
            if not source_pptx:
                source_pptx = asset.template_path
            if not source_pptx:
                return {
                    "pipeline": "enterprise",
                    "error": "--pages requires a template.pptx or existing version in the project directory",
                }
            result = self._handle_page_revision(
                source_pptx, pages, pptx_path, vnum, version_dir
            )
            return result

        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path as _find_db

        component_lib = None
        db_path = _find_db(component_library, project_dir)
        if db_path:
            try:
                component_lib = ComponentLibrary(db_path)
            except Exception:
                component_lib = None

        precision = PrecisionRenderer(brand_spec=brand_spec, template_path=asset.template_path)
        prs = precision.create_presentation()

        page_designs = match_images(asset.image_pool, page_designs)
        page_designs = assign_images_by_size(asset.image_pool, page_designs)
        page_designs = auto_generate_image_prompts(page_designs)
        density_profile = get_density_profile(effective_density)
        for design in page_designs:
            bullets = design.get("bullets")
            if bullets:
                design["bullets"] = apply_density_to_bullets(bullets, density_profile)

        image_fetcher = self._build_image_fetcher(kwargs)

        for design in page_designs:
            image_val = design.get("image")
            if image_val and not os.path.isfile(image_val):
                image_val = None
                design["image"] = None
            if not image_val and image_fetcher is not None:
                prompt = design.get("image_prompt")
                keywords = prompt or design.get("title", query)
                goal = design.get("goal", "")
                try:
                    fetched = image_fetcher.fetch(keywords, goal=goal)
                    if fetched:
                        design["image"] = fetched
                except Exception:
                    pass

        total = len(page_designs)
        for i, design in enumerate(page_designs):
            layout_variant = design.get("layout_variant") if isinstance(design, dict) else None
            precision.render_slide(prs, design, component_lib=component_lib,
                                   layout_variant=layout_variant,
                                   page_index=i, total_pages=total)

        if asset.logo_path:
            for idx, slide in enumerate(prs.slides):
                design_goal = page_designs[idx].get("goal") if idx < len(page_designs) else None
                precision.apply_logo(slide, asset.logo_path, prs, current_goal=design_goal)

        precision.apply_footer(prs)
        precision.apply_watermark(prs)

        effective_motion = motion
        if effective_motion is None:
            effective_motion = kwargs.get("motion")
        if effective_motion and isinstance(effective_motion, int) and effective_motion > 0:
            for slide in prs.slides:
                self._apply_animations(slide, {"goal": "content"}, effective_motion)

        precision.save(prs, pptx_path)

        if component_lib:
            try:
                component_lib.close()
            except Exception:
                pass

        meta = {
            "version": vnum,
            "query": query,
            "business_mode": business_mode,
            "density": effective_density,
            "num_slides": len(page_designs),
            "brand_source": brand_spec.source,
            "render_mode": "precision",
            "slides": [
                {"goal": d.get("goal"), "title": d.get("title", "")}
                for d in page_designs
            ],
        }
        write_meta(version_dir, meta)

        return {
            "pipeline": "enterprise",
            "project": project_dir,
            "output_path": pptx_path,
            "version": vnum,
            "num_slides": len(page_designs),
            "render_mode": "precision",
        }

    def _build_brand_spec(self, asset) -> BrandSpec:
        brand_spec = BrandSpec()
        if asset.template_path:
            analyzer = TemplateAnalyzer()
            template_spec = analyzer.analyze(asset.template_path)
            if asset.brand_raw:
                brand_spec = BrandSpec.merge(template_spec, asset.brand_raw)
            else:
                brand_spec = template_spec
        elif asset.brand_raw:
            brand_spec = BrandSpec.from_brand_json(asset.brand_raw)
        return brand_spec

    def _handle_history(self, project_dir: str) -> dict[str, Any]:
        output_dir = os.path.join(project_dir, "output")
        versions: list[dict[str, Any]] = []

        if os.path.isdir(output_dir):
            for entry in sorted(os.listdir(output_dir)):
                if entry.startswith("v") and entry[1:].isdigit():
                    vnum = int(entry[1:])
                    meta = read_meta(os.path.join(output_dir, entry))
                    versions.append({
                        "version": vnum,
                        "meta": meta,
                    })

        return {"history": True, "project": project_dir, "versions": versions}

    def _load_content(self, asset, project_dir: str, content_file: str | None = None) -> list[dict[str, Any]]:
        if content_file and os.path.isfile(content_file):
            import json
            try:
                with open(content_file, encoding="utf-8") as f:
                    content_raw = json.load(f)
                return load_enterprise_content(content_raw, project_dir)
            except (json.JSONDecodeError, OSError):
                pass
        if asset.content_raw:
            return load_enterprise_content(asset.content_raw, project_dir)
        readme_path = getattr(asset, "readme_path", None)
        if readme_path:
            from ppt_pro_max.enterprise.content_parser import parse_readme
            return parse_readme(readme_path, project_dir)
        return []

    def _load_content_from_meta(self, meta: dict[str, Any]) -> list[dict[str, Any]]:
        return [
            {"goal": s.get("goal", "content"), "title": s.get("title", "")}
            for s in meta.get("slides", [])
        ]

    def _build_story_plan(
        self,
        query: str,
        page_contents: list[dict[str, Any]],
        density: int,
        business_mode: str | None = None,
    ) -> dict[str, Any]:
        if page_contents:
            return {
                "strategy": "content_json",
                "pages": page_contents,
            }

        from ppt_pro_max.planner.story_planner import StoryPlanner
        planner = StoryPlanner()
        strategy_name = self._map_business_mode_to_strategy(business_mode)
        target_pages = max(5, min(20, density + 2))
        plan = planner.plan(query, strategy_override=strategy_name, slide_count_override=target_pages)
        return {
            "strategy": plan.strategy,
            "pages": [
                {"goal": p.goal, "title": p.goal.replace("_", " ").title(), "notes": f"本页目标: {p.goal}。请结合标题展开讲解。"}
                for p in plan.pages
            ],
        }

    def _map_business_mode_to_strategy(self, business_mode: str | None) -> str | None:
        mapping = {
            "pitch": None,
            "education": "Education Course",
            "training": "Training Workshop",
            "report": "Business Report",
        }
        return mapping.get(business_mode)

    def _build_page_designs(
        self,
        story_plan: dict[str, Any],
        decider: EnterpriseDesignDecider,
        asset,
        brand_spec: BrandSpec,
    ) -> list[dict[str, Any]]:
        pages = story_plan.get("pages", [])
        designs: list[dict[str, Any]] = []

        layout_names: dict[int, str] = {}
        if asset.template_path:
            try:
                from pptx import Presentation
                tprs = Presentation(asset.template_path)
                for idx, layout in enumerate(tprs.slide_layouts):
                    layout_names[idx] = layout.name
            except Exception:
                pass

        for page in pages:
            goal = page.get("goal", "content")
            layout_idx = decider.resolve_layout_index(goal)
            design = {
                "goal": goal,
                "title": page.get("title", ""),
                "subtitle": page.get("subtitle"),
                "bullets": page.get("bullets"),
                "image": page.get("image"),
                "cards": page.get("cards"),
                "diagram_type": page.get("diagram_type"),
                "diagram_data": page.get("diagram_data"),
                "code": page.get("code"),
                "exercise": page.get("exercise"),
                "notes": page.get("notes"),
                "links": page.get("links"),
                "chart": page.get("chart"),
                "image_grid": page.get("image_grid"),
                "icons": page.get("icons"),
                "component_type": page.get("component_type"),
                "component_category": page.get("component_category"),
                "component_variant": page.get("component_variant"),
                "template_layout_index": layout_idx,
            }
            if layout_idx is not None and layout_idx in layout_names:
                design["template_layout_name"] = layout_names[layout_idx]
            designs.append(design)
        return designs

    def _apply_animations(self, slide, design: dict[str, Any], motion: int) -> None:
        from ppt_pro_max.renderer.animation import (
            add_slide_transition, add_entrance_animation, TRANSITION_TYPES,
        )
        motion = max(1, min(10, motion))

        transition_list = list(TRANSITION_TYPES.keys())
        transition_type = transition_list[motion % len(transition_list)]
        speed = "fast" if motion >= 7 else ("medium" if motion >= 4 else "slow")
        add_slide_transition(slide, transition_type=transition_type, speed=speed)

        if motion >= 3:
            animatable_ids = []
            for shape in slide.shapes:
                if shape.shape_id and shape.has_text_frame:
                    animatable_ids.append(shape.shape_id)
            if animatable_ids:
                effect = "fade_in" if motion <= 5 else "fly_in"
                interval = max(100, 600 - motion * 50)
                for i, sid in enumerate(animatable_ids[:5]):
                    add_entrance_animation(
                        slide, sid, effect=effect,
                        delay_ms=interval * i if i > 0 else 0,
                        duration_ms=400,
                        click_triggered=(i == 0),
                    )

    def _handle_page_revision(
        self,
        template_path: str,
        pages_str: str,
        pptx_path: str,
        vnum: int,
        version_dir: str,
    ) -> dict[str, Any]:
        engine = PageRevisionEngine(template_path)
        result = engine.revise([pages_str], output_path=pptx_path)

        meta = {
            "version": vnum,
            "mode": "page_revision",
            "pages_op": pages_str,
            "num_slides": result["num_slides"],
        }
        write_meta(version_dir, meta)

        return {
            "pipeline": "enterprise",
            "mode": "page_revision",
            "output_path": result["output_path"],
            "version": vnum,
            "num_slides": result["num_slides"],
            "ops_applied": result.get("ops_applied", []),
        }

    def _build_image_fetcher(self, kwargs: dict[str, Any]):
        image_mode = kwargs.get("image_mode", "placeholder")
        if image_mode == "placeholder":
            return None

        try:
            from ppt_pro_max.renderer.image_fetcher import ImageFetcher

            image_config = dict(kwargs.get("image_config") or {})
            for key in ("llm_provider", "llm_api_key", "llm_base_url", "llm_model"):
                if key not in image_config and kwargs.get(key):
                    image_config[key] = kwargs[key]
            return ImageFetcher(
                mode=image_mode,
                unsplash_access_key=image_config.get("unsplash_access_key"),
                pexels_api_key=image_config.get("pexels_api_key"),
                llm_provider=image_config.get("llm_provider"),
                llm_api_key=image_config.get("llm_api_key"),
                llm_base_url=image_config.get("llm_base_url"),
                llm_model=image_config.get("llm_model"),
            )
        except Exception:
            return None

    def _find_latest_pptx(self, output_dir: str, current_vnum: int) -> str | None:
        if not os.path.isdir(output_dir):
            return None
        for v in range(current_vnum - 1, 0, -1):
            vdir = os.path.join(output_dir, f"v{v}")
            if not os.path.isdir(vdir):
                continue
            for name in ("presentation.pptx", f"v{v}.pptx"):
                candidate = os.path.join(vdir, name)
                if os.path.isfile(candidate):
                    return candidate
        for name in sorted(os.listdir(output_dir)):
            if name.endswith(".pptx") and not name.startswith("~"):
                candidate = os.path.join(output_dir, name)
                if os.path.isfile(candidate):
                    return candidate
        return None
