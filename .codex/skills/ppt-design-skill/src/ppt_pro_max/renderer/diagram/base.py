"""BaseDiagram — abstract base class for all diagram types."""

from __future__ import annotations

from abc import ABC, abstractmethod
from typing import Any

from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
from ppt_pro_max.renderer.diagram.layout_engine import Region


class BaseDiagram(ABC):

    def __init__(self, data: dict[str, Any], style: DiagramStyle, region: Region):
        self.data = data
        self.style = style
        self.region = region
        self._nodes: list[dict[str, Any]] = []
        self._connectors: list[dict[str, Any]] = []

    @abstractmethod
    def compute_layout(self) -> None:
        ...

    def render(self, slide) -> None:
        self.compute_layout()
        self._draw_nodes(slide)
        self._draw_connectors(slide)
        if self.data.get("group_nodes"):
            self._group_rendered_nodes(slide)

    def _group_rendered_nodes(self, slide) -> None:
        try:
            from ppt_pro_max.renderer.shape_utils import group_shapes
            n_nodes = len(self._nodes)
            n_connectors = len(self._connectors)
            if n_nodes < 2:
                return
            total = n_nodes + n_connectors
            existing = len(slide.shapes) - total
            indices = list(range(existing, existing + n_nodes))
            group_shapes(slide, shape_indices=indices)
        except Exception:
            pass

    def _draw_nodes(self, slide) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from ppt_pro_max.renderer.visual_effects import apply_gradient, apply_shadow

        for node in self._nodes:
            left = Inches(node["x"])
            top = Inches(node["y"])
            width = Inches(node["width"])
            height = Inches(node["height"])

            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
            if node.get("shape") == "rectangle":
                shape_type = MSO_SHAPE.RECTANGLE
            elif node.get("shape") == "oval":
                shape_type = MSO_SHAPE.OVAL
            elif node.get("shape") == "hexagon":
                shape_type = MSO_SHAPE.HEXAGON
            elif node.get("shape") == "donut":
                shape_type = MSO_SHAPE.DONUT
            elif node.get("shape") == "diamond":
                shape_type = MSO_SHAPE.DIAMOND
            elif node.get("shape") == "triangle":
                shape_type = MSO_SHAPE.ISOSCELES_TRIANGLE
            elif node.get("shape") == "chevron":
                shape_type = MSO_SHAPE.CHEVRON

            shape = slide.shapes.add_shape(shape_type, left, top, width, height)

            fill_color = self.style.resolve_color(node.get("fill_role", self.style.node_fill))
            use_gradient = node.get("gradient", self.style.node_gradient)
            use_shadow = node.get("shadow", self.style.node_shadow)

            if use_gradient:
                dark = self._darken_color(fill_color, 30)
                grad_type = "path" if shape_type in (MSO_SHAPE.OVAL, MSO_SHAPE.DONUT) else "linear"
                apply_gradient(shape, fill_color.lstrip("#"), dark.lstrip("#"), gradient_type=grad_type)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))

            if use_shadow:
                apply_shadow(shape, blur_pt=4, distance_pt=2, alpha_pct=25)

            border_color = self.style.resolve_color(self.style.node_border)
            shape.line.color.rgb = RGBColor.from_string(border_color.lstrip("#"))
            shape.line.width = Pt(self.style.node_border_width_pt)

            label = node.get("label", "")
            if label:
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = label
                p.font.size = Pt(node.get("font_size_pt", self.style.node_font_size_pt))
                p.alignment = self._resolve_alignment()

                font_color = self.style.resolve_color(node.get("font_color_role", self.style.node_font_color))
                p.font.color.rgb = RGBColor.from_string(font_color.lstrip("#"))

                if node.get("font_weight", self.style.node_font_weight) == "bold":
                    p.font.bold = True

    @staticmethod
    def _darken_color(hex_color: str, amount: int = 30) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = max(0, r - amount)
        g = max(0, g - amount)
        b = max(0, b - amount)
        return f"#{r:02X}{g:02X}{b:02X}"

    def _draw_connectors(self, slide) -> None:
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        for conn in self._connectors:
            x1 = Inches(conn["x1"])
            y1 = Inches(conn["y1"])
            x2 = Inches(conn["x2"])
            y2 = Inches(conn["y2"])

            waypoints = self._route_connector(conn)

            if len(waypoints) <= 2:
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2,
                )
            else:
                for i in range(len(waypoints) - 1):
                    wx1, wy1 = waypoints[i]
                    wx2, wy2 = waypoints[i + 1]
                    connector = slide.shapes.add_connector(
                        MSO_CONNECTOR_TYPE.STRAIGHT,
                        Inches(wx1), Inches(wy1), Inches(wx2), Inches(wy2),
                    )

                    conn_color = self.style.resolve_color(self.style.connector_color)
                    connector.line.color.rgb = RGBColor.from_string(conn_color.lstrip("#"))
                    connector.line.width = Pt(self.style.connector_width_pt)
                continue

            conn_color = self.style.resolve_color(self.style.connector_color)
            connector.line.color.rgb = RGBColor.from_string(conn_color.lstrip("#"))
            connector.line.width = Pt(self.style.connector_width_pt)

    def _resolve_alignment(self):
        from pptx.enum.text import PP_ALIGN
        if self.style.node_text_alignment == "center":
            return PP_ALIGN.CENTER
        elif self.style.node_text_alignment == "right":
            return PP_ALIGN.RIGHT
        return PP_ALIGN.LEFT

    def _route_connector(self, conn: dict[str, Any]) -> list[tuple[float, float]]:
        from ppt_pro_max.renderer.diagram.connector_router import route_connector
        return route_connector(
            conn["x1"], conn["y1"], conn["x2"], conn["y2"],
            self._nodes,
        )
