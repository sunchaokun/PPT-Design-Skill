"""Shape Factory — high-level shape creation with gradient fills and effects.

Creates PPT-7 quality shapes (OVAL, DONUT, HEXAGON, ARC, etc.) with
gradient fills, shadows, and other visual effects pre-applied.

Usage:
    factory = ShapeFactory(brand_colors={"primary": "#1D78FA", ...})
    factory.add_glow_oval(slide, 5.0, 3.0, 1.2, 1.2, label="策略")
    factory.add_ring_node(slide, 5.0, 3.0, 2.0, label="核心")
    factory.add_hexagon_card(slide, 2.0, 2.0, 1.5, label="技术")
"""

from __future__ import annotations

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

from ppt_pro_max.renderer.visual_effects import (
    apply_gradient,
    apply_shadow,
)


class ShapeFactory:

    def __init__(self, brand_colors: dict[str, str] | None = None):
        self._colors = brand_colors or {}
        self._gradient_presets = self._build_gradient_presets()

    def _c(self, role: str, fallback: str = "#000000") -> str:
        return self._colors.get(role, fallback).lstrip("#")

    def _build_gradient_presets(self) -> dict[str, tuple[str, str]]:
        primary = self._c("primary", "#1D78FA")
        accent = self._c("accent", "#FF5500")
        secondary = self._c("secondary", "#64748B")
        return {
            "primary": (primary, self._darken(primary, 30)),
            "accent": (accent, self._darken(accent, 30)),
            "secondary": (secondary, self._darken(secondary, 20)),
            "muted": (self._c("muted", "#F1F5F9"), self._c("border", "#E2E8F0")),
        }

    @staticmethod
    def _darken(hex_color: str, amount: int = 30) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = max(0, r - amount)
        g = max(0, g - amount)
        b = max(0, b - amount)
        return f"{r:02X}{g:02X}{b:02X}"

    @staticmethod
    def _lighten(hex_color: str, amount: int = 30) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = min(255, r + amount)
        g = min(255, g + amount)
        b = min(255, b + amount)
        return f"{r:02X}{g:02X}{b:02X}"

    def add_glow_oval(
        self, slide, x: float, y: float, w: float, h: float,
        label: str = "", font_size: int = 14, color_role: str = "primary",
        gradient_type: str = "linear", angle: int = 5400000,
        shadow: bool = True, font_color: str = "#FFFFFF",
    ) -> object:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        c1, c2 = self._gradient_presets.get(color_role, (self._c(color_role), self._darken(self._c(color_role))))
        apply_gradient(shape, c1, c2, gradient_type=gradient_type, angle=angle)
        if shadow:
            apply_shadow(shape, blur_pt=4, distance_pt=2, alpha_pct=25)
        shape.line.fill.background()
        if label:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor.from_string(font_color.lstrip("#"))
            run.font.bold = True
            run.font.name = "微软雅黑"
        return shape

    def add_ring_node(
        self, slide, x: float, y: float, size: float,
        label: str = "", font_size: int = 18, color_role: str = "primary",
        ring_width_ratio: float = 0.25, shadow: bool = True,
    ) -> object:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.DONUT, Inches(x), Inches(y), Inches(size), Inches(size),
        )
        c1, c2 = self._gradient_presets.get(color_role, (self._c(color_role), self._darken(self._c(color_role))))
        apply_gradient(shape, c1, c2, gradient_type="path")
        if shadow:
            apply_shadow(shape, blur_pt=6, distance_pt=3, alpha_pct=20)
        shape.line.fill.background()
        if label:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor.from_string("FFFFFF")
            run.font.bold = True
            run.font.name = "微软雅黑"
        return shape

    def add_hexagon_card(
        self, slide, x: float, y: float, size: float,
        label: str = "", subtitle: str = "", font_size: int = 13,
        color_role: str = "primary", shadow: bool = True,
    ) -> object:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(size), Inches(size * 0.87),
        )
        c1, c2 = self._gradient_presets.get(color_role, (self._c(color_role), self._darken(self._c(color_role))))
        apply_gradient(shape, c1, c2, gradient_type="linear", angle=5400000)
        if shadow:
            apply_shadow(shape, blur_pt=4, distance_pt=2, alpha_pct=25)
        shape.line.fill.background()
        if label:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor.from_string("FFFFFF")
            run.font.bold = True
            run.font.name = "微软雅黑"
            if subtitle:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                run2 = p2.add_run()
                run2.text = subtitle
                run2.font.size = Pt(max(11, font_size - 3))
                run2.font.color.rgb = RGBColor.from_string("FFFFFF")
                run2.font.name = "微软雅黑"
        return shape

    def add_arc_connector(
        self, slide, x: float, y: float, w: float, h: float,
        color_role: str = "accent", line_width_pt: float = 2.0,
    ) -> object:
        shape = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(x), Inches(y), Inches(w), Inches(h))
        spPr = shape._element.find(qn("p:spPr"))
        for tag in ["a:solidFill", "a:noFill", "a:gradFill"]:
            el = spPr.find(qn(tag))
            if el is not None:
                spPr.remove(el)
        etree.SubElement(spPr, qn("a:noFill"))
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(spPr, qn("a:ln"))
        ln.set("w", str(int(line_width_pt * 12700)))
        for tag in ["a:solidFill", "a:noFill"]:
            el = ln.find(qn(tag))
            if el is not None:
                ln.remove(el)
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(solidFill, qn("a:srgbClr")).set("val", self._c(color_role, "#FF5500"))
        return shape

    def add_rounded_card(
        self, slide, x: float, y: float, w: float, h: float,
        title: str = "", body: str = "", color_role: str = "primary",
        shadow: bool = True, gradient: bool = True,
    ) -> object:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
        )
        if gradient:
            c1 = self._c("muted", "#F1F5F9")
            c2 = self._lighten(c1, 10)
            apply_gradient(shape, c1, c2, gradient_type="linear", angle=2700000)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(self._c("muted", "#F1F5F9"))
        if shadow:
            apply_shadow(shape, blur_pt=4, distance_pt=2, alpha_pct=15)
        shape.line.color.rgb = RGBColor.from_string(self._c("border", "#E2E8F0"))
        shape.line.width = Pt(0.5)
        if title:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = title
            run.font.size = Pt(15)
            run.font.color.rgb = RGBColor.from_string(self._c(color_role, "#1D78FA"))
            run.font.bold = True
            run.font.name = "微软雅黑"
            if body:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.LEFT
                run2 = p2.add_run()
                run2.text = body
                run2.font.size = Pt(11)
                run2.font.color.rgb = RGBColor.from_string(self._c("muted-foreground", "#64748B"))
                run2.font.name = "微软雅黑"
        return shape

    def add_triangle_pointer(
        self, slide, x: float, y: float, w: float, h: float,
        color_role: str = "accent", gradient: bool = True,
    ) -> object:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
        )
        if gradient:
            c1 = self._c(color_role, "#FF5500")
            c2 = self._darken(c1, 30)
            apply_gradient(shape, c1, c2, gradient_type="linear", angle=5400000)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(self._c(color_role, "#FF5500"))
        shape.line.fill.background()
        return shape

    def add_chevron_step(
        self, slide, x: float, y: float, w: float, h: float,
        label: str = "", font_size: int = 12, color_role: str = "primary",
        step_index: int = 0,
    ) -> object:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h),
        )
        c1 = self._c(color_role, "#1D78FA")
        shade = min(step_index * 15, 60)
        c2 = self._darken(c1, shade)
        apply_gradient(shape, c1, c2, gradient_type="linear", angle=0)
        apply_shadow(shape, blur_pt=3, distance_pt=1, alpha_pct=15)
        shape.line.fill.background()
        if label:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor.from_string("FFFFFF")
            run.font.bold = True
            run.font.name = "微软雅黑"
        return shape

    def add_diamond_node(
        self, slide, x: float, y: float, size: float,
        label: str = "", font_size: int = 12, color_role: str = "accent",
    ) -> object:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(size), Inches(size),
        )
        c1, c2 = self._gradient_presets.get(color_role, (self._c(color_role), self._darken(self._c(color_role))))
        apply_gradient(shape, c1, c2, gradient_type="path")
        apply_shadow(shape, blur_pt=3, distance_pt=2, alpha_pct=20)
        shape.line.fill.background()
        if label:
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor.from_string("FFFFFF")
            run.font.bold = True
            run.font.name = "微软雅黑"
        return shape

    def add_block_arc(
        self, slide, x: float, y: float, w: float, h: float,
        color_role: str = "primary", line_width_pt: float = 3.0,
    ) -> object:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.BLOCK_ARC, Inches(x), Inches(y), Inches(w), Inches(h),
        )
        c1 = self._c(color_role, "#1D78FA")
        c2 = self._darken(c1, 20)
        apply_gradient(shape, c1, c2, gradient_type="path")
        shape.line.fill.background()
        return shape

    def add_radial_diagram(
        self, slide, center_x: float, center_y: float,
        center_label: str, center_size: float = 2.5,
        nodes: list[dict[str, str]] | None = None,
        node_size: float = 1.0, radius: float = 3.0,
        center_color_role: str = "primary",
    ) -> list[object]:
        shapes = []
        center = self.add_glow_oval(
            slide, center_x - center_size / 2, center_y - center_size / 2,
            center_size, center_size, label=center_label,
            font_size=20, color_role=center_color_role, gradient_type="path",
        )
        shapes.append(center)

        if not nodes:
            return shapes

        n = len(nodes)
        for i, node in enumerate(nodes):
            angle = 2 * math.pi * i / n - math.pi / 2
            nx = center_x + radius * math.cos(angle) - node_size / 2
            ny = center_y + radius * math.sin(angle) - node_size / 2
            oval = self.add_glow_oval(
                slide, nx, ny, node_size, node_size,
                label=node.get("label", ""),
                font_size=node.get("font_size", 12),
                color_role=node.get("color_role", "accent"),
            )
            shapes.append(oval)
        return shapes

    def add_ring_diagram(
        self, slide, center_x: float, center_y: float,
        center_label: str, ring_size: float = 4.0,
        nodes: list[dict[str, str]] | None = None,
        node_size: float = 0.8, color_role: str = "primary",
    ) -> list[object]:
        shapes = []
        ring = self.add_ring_node(
            slide, center_x - ring_size / 2, center_y - ring_size / 2,
            ring_size, label=center_label, font_size=16, color_role=color_role,
        )
        shapes.append(ring)

        if not nodes:
            return shapes

        n = len(nodes)
        radius = ring_size / 2 + node_size * 0.8
        for i, node in enumerate(nodes):
            angle = 2 * math.pi * i / n - math.pi / 2
            nx = center_x + radius * math.cos(angle) - node_size / 2
            ny = center_y + radius * math.sin(angle) - node_size / 2
            oval = self.add_glow_oval(
                slide, nx, ny, node_size, node_size,
                label=node.get("label", ""),
                font_size=node.get("font_size", 11),
                color_role=node.get("color_role", "accent"),
            )
            shapes.append(oval)
        return shapes
