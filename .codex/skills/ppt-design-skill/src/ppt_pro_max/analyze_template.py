"""analyze_template — extract visual DNA from a PPTX template.

Usage: python -m ppt_pro_max.analyze_template template.pptx
Output: structured text for LLM consumption
"""
from __future__ import annotations

import sys
import io
from collections import Counter

from pptx import Presentation

_NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

_COVER_KEYWORDS = {"汇报", "报告", "方案", "规划", "战略", "年度", "总结", "展望", "report", "plan", "strategy", "annual"}
_TOC_KEYWORDS = {"目录", "contents", "agenda", "议程", "概览", "overview", "壹", "贰", "叁", "肆", "伍"}
_CTA_KEYWORDS = {"谢谢", "感谢", "thank", "thanks", "联系", "contact", "观看", "聆听", "谢观"}
_TRANSITION_KEYWORDS = {"part", "章节", "section"}


def _classify_page_type(slide, slide_index, total_slides):
    all_text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            all_text.append(shape.text_frame.text.strip())
    full_text = ' '.join(all_text).lower()

    if slide_index == 0:
        return 'cover'
    if slide_index == total_slides - 1:
        for kw in _CTA_KEYWORDS:
            if kw in full_text:
                return 'back_cover'
    for kw in _TOC_KEYWORDS:
        if kw in full_text:
            return 'toc'
    if len(all_text) <= 3:
        for kw in _TRANSITION_KEYWORDS:
            if kw in full_text:
                return 'transition'
    for kw in _CTA_KEYWORDS:
        if kw in full_text:
            return 'back_cover'
    return 'content'


def analyze(pptx_path: str) -> str:
    buf = io.StringIO()
    _analyze_to(pptx_path, buf)
    return buf.getvalue()


def _analyze_to(pptx_path: str, out: io.StringIO):
    prs = Presentation(pptx_path)
    W_in = prs.slide_width / 914400
    H_in = prs.slide_height / 914400
    total = len(prs.slides)

    def w(s=''):
        out.write(s + '\n')

    w('=== TEMPLATE ANALYSIS ===')
    w(f'Size: {W_in:.3f} x {H_in:.3f} inches')
    w(f'Slides: {total}')
    w(f'Layouts: {len(prs.slide_layouts)}')
    for i, layout in enumerate(prs.slide_layouts):
        w(f'  Layout {i}: {layout.name}')
    w()

    color_counter = Counter()

    for si, slide in enumerate(prs.slides):
        page_type = _classify_page_type(slide, si, total)
        w(f'--- Slide {si} [{page_type}] ({len(slide.shapes)} shapes) ---')

        for shape in slide.shapes:
            pos = f"({shape.left/914400:.2f}, {shape.top/914400:.2f})"
            size = f"({shape.width/914400:.2f}x{shape.height/914400:.2f})"

            stype = str(shape.shape_type)
            tag = shape._element.tag.split('}')[-1]
            if tag == 'cxnSp':
                stype = 'connector'
            elif tag == 'grpSp':
                stype = 'group'
            elif shape.shape_type == 13:
                stype = 'image'

            info = f"  {stype} {pos} {size}"

            if shape.has_text_frame:
                txt = shape.text_frame.text[:80].replace('\n', ' | ')
                info += f' text="{txt}"'
                if shape.text_frame.paragraphs:
                    p0 = shape.text_frame.paragraphs[0]
                    if p0.font.size:
                        info += f' size={p0.font.size.pt}pt'
                    try:
                        if p0.font.color and p0.font.color.type is not None:
                            info += f' color=#{p0.font.color.rgb}'
                    except Exception:
                        pass
                    if p0.font.name:
                        info += f' font={p0.font.name}'

            sp = shape._element
            prst = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prst is not None:
                info += f' geom={prst.get("prst", "?")}'
                avLst = prst.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if avLst is not None:
                    adjs = [gd.get('fmla', '') for gd in avLst.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}gd')]
                    if adjs:
                        info += f' adj={adjs}'

            for clr in sp.findall('.//a:solidFill/a:srgbClr', _NS):
                alpha = clr.find('a:alpha', _NS)
                if alpha is not None:
                    alpha_val = int(alpha.get('val', '100000'))
                    info += f' alpha={alpha_val/1000:.0f}%'

            for grad in sp.findall('.//a:gradFill', _NS):
                stops = grad.findall('.//a:gs', _NS)
                if stops:
                    grad_info = []
                    for gs in stops:
                        pos_val = gs.get('pos', '?')
                        srgb = gs.find('a:srgbClr', _NS)
                        if srgb is not None:
                            grad_info.append(f"pos={pos_val} #{srgb.get('val', '')}")
                    if grad_info:
                        info += f' gradient=[{" | ".join(grad_info)}]'

            w(info)
        w()

    for slide in prs.slides:
        for shape in slide.shapes:
            sp = shape._element
            for clr in sp.findall('.//a:solidFill/a:srgbClr', _NS):
                val = clr.get('val', '')
                if val:
                    color_counter[val] += 1

    w('=== COLOR FREQUENCY ===')
    for hex_val, count in color_counter.most_common(20):
        w(f'  #{hex_val}: {count} occurrences')

    w()
    w('=== THEME COLORS ===')
    theme = prs.slide_masters[0].element.findall('.//a:clrScheme/a:*', _NS)
    for t in theme:
        tag = t.tag.split('}')[1] if '}' in t.tag else t.tag
        srgb = t.find('a:srgbClr', _NS)
        if srgb is not None:
            w(f'  {tag}: #{srgb.get("val")}')
        else:
            sysclr = t.find('a:sysClr', _NS)
            if sysclr is not None:
                w(f'  {tag}: sys={sysclr.get("val")} lastClr={sysclr.get("lastClr")}')

    font_counter = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.font.name:
                        font_counter[p.font.name] += 1
                    for r in p.runs:
                        if r.font.name:
                            font_counter[r.font.name] += 1

    w()
    w('=== FONT FREQUENCY ===')
    for font_name, count in font_counter.most_common(10):
        w(f'  {font_name}: {count} occurrences')


if __name__ == '__main__':
    path = sys.argv[1] if len(sys.argv) > 1 else 'template.pptx'
    result = analyze(path)
    sys.stdout.buffer.write(result.encode('utf-8'))
