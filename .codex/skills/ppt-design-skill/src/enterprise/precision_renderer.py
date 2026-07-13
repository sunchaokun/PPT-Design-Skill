"""PrecisionRenderer — brand-compliant + pixel-perfect rendering.

Combines Pipeline's brand constraint system (BrandSpec, template layouts,
logo/footer/watermark) with Build Script's rendering precision (run-level
fonts, Pillow cropping, per-element coordinates).

This is the third rendering path:
  - EnterpriseRenderer: template-driven, paragraph-level fonts → low quality
  - PPTRenderer (freestyle fallback): layout-registry driven → medium quality
  - PrecisionRenderer: brand-aware + run-level precision → delivery quality
"""

from __future__ import annotations

import hashlib
import os
import tempfile
from typing import Any

from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.renderer.layout_registry import SLIDE_WIDTH, SLIDE_HEIGHT


class PrecisionRenderer:

    def __init__(self, brand_spec: BrandSpec | None = None, template_path: str | None = None):
        self._brand = brand_spec
        self._template_path = template_path
        self._has_template = bool(template_path and os.path.exists(template_path))
        self._crop_cache_dir = os.path.join(tempfile.gettempdir(), "ppt-precision-crops")
        os.makedirs(self._crop_cache_dir, exist_ok=True)

    @property
    def brand(self) -> BrandSpec:
        return self._brand or BrandSpec()

    def _c(self, role: str, fallback: str = "#000000") -> str:
        if self._brand and self._brand.colors:
            return self._brand.colors.get(role, fallback)
        return fallback

    def _font_h(self) -> str:
        if self._brand and self._brand.fonts:
            return self._brand.fonts.get("heading", "Inter")
        return "Inter"

    def _font_b(self) -> str:
        if self._brand and self._brand.fonts:
            return self._brand.fonts.get("body", "Inter")
        return "Inter"

    def _is_dark(self) -> bool:
        if self._brand and self._brand.dark_mode:
            return True
        bg = self._c("background", "#FFFFFF").lstrip("#").upper()
        return bg in ("060B18", "0A1E3D", "120C1E", "111827", "09090B", "1C1010", "2D2D2D")

    @staticmethod
    def _rgb(h: str) -> RGBColor:
        return RGBColor.from_string(h.lstrip("#"))

    def create_presentation(self) -> Presentation:
        if self._has_template:
            try:
                from ppt_pro_max.enterprise.slide_utils import remove_slide
                prs = Presentation(self._template_path)
                while len(prs.slides) > 0:
                    remove_slide(prs, 0)
                return prs
            except Exception:
                pass
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH)
        prs.slide_height = Inches(SLIDE_HEIGHT)
        return prs

    def add_slide(self, prs: Presentation, layout_name: str | None = None):
        if self._has_template:
            if layout_name:
                for layout in prs.slide_layouts:
                    if layout.name == layout_name:
                        return prs.slides.add_slide(layout)
            for layout in prs.slide_layouts:
                if "blank" in layout.name.lower():
                    return prs.slides.add_slide(layout)
            return prs.slides.add_slide(prs.slide_layouts[0])
        blank = None
        for layout in prs.slide_layouts:
            if "blank" in layout.name.lower():
                blank = layout
                break
        return prs.slides.add_slide(blank or prs.slide_layouts[-1])

    # ── Text helpers (run-level fonts) ──

    def add_text(self, slide, text: str, x: float, y: float, w: float, h: float,
                 font: str | None = None, size: int = 20, color_role: str = "foreground",
                 bold: bool = False, align: str = "left") -> object:
        font = font or self._font_h()
        color = self._c(color_role)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = self._rgb(color)
        run.font.bold = bold
        return tb

    def add_multiline(self, slide, lines: list[str], x: float, y: float, w: float, h: float,
                      font: str | None = None, size: int = 14, color_role: str = "foreground",
                      bold: bool = False, align: str = "left", spacing: int = 6) -> object:
        font = font or self._font_b()
        color = self._c(color_role)
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
            run = p.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = self._rgb(color)
            run.font.bold = bold
            p.space_after = Pt(spacing)
        return tb

    # ── Image helpers (Pillow pre-crop) ──

    def add_image(self, slide, path: str, x: float, y: float, w: float, h: float) -> None:
        if not os.path.isfile(path):
            return
        with PILImage.open(path) as img:
            iw, ih = img.size
            box_ratio = w / h
            img_ratio = iw / ih
            if img_ratio > box_ratio:
                cw, ch = int(ih * box_ratio), ih
                cl, ct = (iw - cw) // 2, 0
            else:
                cw, ch = iw, int(iw / box_ratio)
                cl, ct = 0, (ih - ch) // 2
            cropped = img.crop((cl, ct, cl + cw, ct + ch))
            cp = os.path.join(self._crop_cache_dir,
                              hashlib.md5(f"{path}:{w}x{h}".encode()).hexdigest() + ".png")
            if not os.path.isfile(cp):
                cropped.save(cp, "PNG")
        slide.shapes.add_picture(cp, Inches(x), Inches(y), Inches(w), Inches(h))

    # ── Shape helpers ──

    def add_rect(self, slide, x: float, y: float, w: float, h: float,
                 fill_role: str | None = None, fill_hex: str | None = None,
                 border_role: str | None = None, border_hex: str | None = None) -> object:
        fill = fill_hex or self._c(fill_role or "muted")
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self._rgb(fill)
        border = border_hex or (self._c(border_role) if border_role else None)
        if border:
            sh.line.color.rgb = self._rgb(border)
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    def add_rounded_rect(self, slide, x: float, y: float, w: float, h: float,
                         fill_role: str | None = None, fill_hex: str | None = None,
                         border_role: str | None = None, border_hex: str | None = None) -> object:
        fill = fill_hex or self._c(fill_role or "muted")
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self._rgb(fill)
        border = border_hex or (self._c(border_role) if border_role else None)
        if border:
            sh.line.color.rgb = self._rgb(border)
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        return sh

    # ── Overlay helpers ──

    def add_dark_overlay(self, slide, opacity: float = 0.65) -> None:
        bg_hex = self._c("background", "#060B18" if self._is_dark() else "#000000")
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT))
        ov.fill.solid()
        ov.fill.fore_color.rgb = self._rgb(bg_hex)
        ov.line.fill.background()
        el = ov._element.find(qn("p:spPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if el is not None:
            a = etree.SubElement(el, qn("a:alpha"))
            a.set("val", str(int(opacity * 100000)))

    # ── Brand visual design ──

    def render_slide(self, prs: Presentation, page: dict[str, Any]) -> object:
        goal = page.get("goal", "content")
        title = page.get("title", "")
        subtitle = page.get("subtitle")
        bullets = page.get("bullets") or []
        image_path = page.get("image")
        cards = page.get("cards") or []
        diagram_type = page.get("diagram_type")
        diagram_data = page.get("diagram_data")
        code = page.get("code")
        exercise = page.get("exercise")
        chart = page.get("chart")
        notes = page.get("notes")
        _image_grid = page.get("image_grid")
        _icons = page.get("icons")

        explicit_layout = page.get("layout")
        if explicit_layout:
            goal = self._remap_layout_to_goal(explicit_layout, goal)

        is_hero = goal in ("hook", "cta")

        layout_name = page.get("template_layout_name")
        slide = self.add_slide(prs, layout_name=layout_name)

        if is_hero:
            self.apply_hero_overlay(slide, prs, image_path=image_path if image_path and os.path.isfile(image_path) else None)
            if title:
                self.add_text(slide, title, 1.2, 2.0, 8, 1.5,
                              size=52, color_role="on-primary" if not (image_path and os.path.isfile(image_path)) else "foreground", bold=True)
            if subtitle:
                self.add_text(slide, subtitle, 1.2, 3.6, 8, 0.6,
                              font=self._font_b(), size=22, color_role="accent")
            if bullets:
                bullet_text = "  ·  ".join(bullets[:5])
                self.add_text(slide, bullet_text, 1.2, 4.4, 8, 0.5,
                              font=self._font_b(), size=14, color_role="muted-foreground")
        else:
            self.apply_brand_background(slide, prs, goal=goal)

            if title:
                self.add_text(slide, title, 0.9, 0.5, 11, 0.6,
                              size=28, color_role="foreground", bold=True)
                self.add_rect(slide, 0.9, 1.2, 2, 0.04, fill_role="accent")

            if subtitle:
                self.add_text(slide, subtitle, 0.9, 1.5, 10, 0.5,
                              font=self._font_b(), size=14, color_role="muted-foreground")

            if cards:
                n = len(cards)
                card_w = min(3.6, (11.5 - 0.4 * (n - 1)) / n)
                for i, card in enumerate(cards):
                    xx = 0.9 + i * (card_w + 0.4)
                    card_title = card.get("title", "")
                    card_body = card.get("text", card.get("body", ""))
                    self.add_card(slide, xx, 1.6, card_w, 4.5, card_title, card_body)

            elif diagram_type and diagram_data:
                self._render_diagram_on_slide(slide, diagram_type, diagram_data)

            elif code:
                self._render_code_on_slide(slide, code)

            elif exercise:
                self._render_exercise_on_slide(slide, exercise)

            elif bullets:
                bullet_lines = [f"•  {b}" for b in bullets[:8]]
                self.add_multiline(slide, bullet_lines, 0.9, 1.6, 7, 4.5,
                                   size=12, color_role="foreground", spacing=6)

            if image_path and os.path.isfile(image_path) and not is_hero:
                self.add_image(slide, image_path, 8.3, 1.2, 4.2, 5.3)

            if chart:
                self._render_chart_on_slide(slide, chart)

        if notes:
            self._render_notes_on_slide(slide, notes)

        return slide

    def _remap_layout_to_goal(self, layout_name: str, fallback_goal: str) -> str:
        layout_to_goal = {
            "title-slide": "hook",
            "cta-closing": "cta",
            "content-with-title": "content",
            "two-column": "content",
            "three-column-cards": "features",
            "four-metrics": "traction",
            "big-number": "agitation",
            "quote": "testimonials",
            "chart-focus": "data",
            "image-plus-text": "content",
            "sidebar-left": "overview",
            "exercise-layout": "exercise",
            "code-block": "code",
            "funnel": "funnel",
            "grid-2x2-cards": "features",
            "dense-bullets": "content",
            "swot-matrix": "content",
            "cycle-diagram": "content",
            "timeline-horizontal": "content",
            "table-layout": "content",
            "section-header": "content",
            "blank": "content",
        }
        return layout_to_goal.get(layout_name, fallback_goal)

    def _render_diagram_on_slide(self, slide, diagram_type: str, diagram_data: dict) -> None:
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        style = DiagramStyle.from_brand_spec(self.brand) if self._brand else DiagramStyle()
        region = Region(left=0.9, top=1.5, width=7.0, height=5.0)
        engine = DiagramEngine()
        try:
            engine.render(slide, diagram_type, diagram_data, style, region)
        except Exception:
            pass

    def _render_code_on_slide(self, slide, code_data) -> None:
        code_text = code_data if isinstance(code_data, str) else code_data.get("source", code_data.get("code", ""))
        language = code_data.get("language", "") if isinstance(code_data, dict) else ""
        code_bg = self._c("muted", "#0D152A" if self._is_dark() else "#F8FAFC")
        self.add_rounded_rect(slide, 0.9, 1.5, 11.533, 5.0, fill_hex=code_bg)
        header = f"  {language}" if language else ""
        lines = code_text.split("\n")
        all_lines = []
        if header:
            all_lines.append(header)
        all_lines.extend(f"  {line}" for line in lines[:30])
        self.add_multiline(slide, all_lines, 1.2, 1.7, 11, 4.5,
                           font="Consolas", size=11, color_role="foreground", spacing=4)

    def _render_exercise_on_slide(self, slide, exercise_data) -> None:
        instructions = exercise_data.get("instructions", "") if isinstance(exercise_data, dict) else str(exercise_data)
        duration = exercise_data.get("duration", "") if isinstance(exercise_data, dict) else ""
        steps = exercise_data.get("steps", []) if isinstance(exercise_data, dict) else []
        badge_text = f"Exercise {duration}" if duration else "Exercise"
        self.add_rounded_rect(slide, 0.9, 1.2, 1.8 if duration else 1.2, 0.4, fill_role="primary")
        self.add_text(slide, badge_text, 0.95, 1.22, 1.7, 0.36, size=11, color_role="on-primary", bold=True)
        if instructions:
            self.add_text(slide, instructions, 0.9, 1.8, 11.533, 1.2,
                          font=self._font_b(), size=13, color_role="muted-foreground")
        if steps:
            step_lines = [f"{i+1}. {s}" for i, s in enumerate(steps)]
            self.add_multiline(slide, step_lines, 0.9, 3.2, 11.533, 3.5,
                               size=13, color_role="foreground", spacing=6)

    def _render_chart_on_slide(self, slide, chart_config: dict) -> None:
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        brand_colors = None
        if self._brand and self._brand.colors:
            brand_colors = self._brand.colors
        position = {"x": 1.5, "y": 1.5, "width": 10.333, "height": 4.5}
        builder = ChartBuilder()
        try:
            builder.build(slide, chart_config, position=position, brand_colors=brand_colors)
        except Exception:
            pass

    def _render_notes_on_slide(self, slide, notes_text: str) -> None:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    def apply_brand_background(self, slide, prs: Presentation, goal: str = "content") -> None:
        bg_hex = self._c("background", "#FFFFFF")
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self._rgb(bg_hex)
        except Exception:
            pass

        accent_hex = self._c("accent", self._c("primary", "#2563EB"))
        self.add_rect(slide, 0, 0, 0.06, SLIDE_HEIGHT, fill_hex=accent_hex)
        muted_hex = self._c("muted", "#F1F5F9")
        self.add_rect(slide, 0, SLIDE_HEIGHT - 0.25, SLIDE_WIDTH, 0.25, fill_hex=muted_hex)

    def apply_hero_overlay(self, slide, prs: Presentation, image_path: str | None = None) -> None:
        if image_path and os.path.isfile(image_path):
            self.add_image(slide, image_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
            self.add_dark_overlay(slide, 0.72)
        else:
            primary_hex = self._c("primary", "#2563EB")
            self.add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_hex=primary_hex)

    # ── Card component ──

    def add_card(self, slide, x: float, y: float, w: float, h: float,
                 title: str, body: str, accent_role: str = "primary") -> None:
        accent_hex = self._c(accent_role)
        card_bg = self._c("muted", "#0D152A" if self._is_dark() else "#F8FAFC")
        card_bd = self._c("border", "#1A2A4A" if self._is_dark() else "#E2E8F0")
        self.add_rounded_rect(slide, x, y, w, h, fill_hex=card_bg, border_hex=card_bd)
        self.add_rect(slide, x + 0.15, y + 0.15, 0.4, 0.05, fill_hex=accent_hex)
        self.add_text(slide, title, x + 0.15, y + 0.35, w - 0.3, 0.4,
                      self._font_h(), 15, accent_role, True)
        self.add_multiline(slide, body.split("\n"), x + 0.15, y + 0.8, w - 0.3, h - 1.0,
                           self._font_b(), 11, "muted-foreground")

    # ── Brand compliance: logo, footer, watermark ──

    def apply_logo(self, slide, logo_path: str, prs: Presentation,
                   current_goal: str | None = None) -> None:
        if not self._brand or not self._brand.logo:
            return
        logo_spec = self._brand.logo
        skip_slides = logo_spec.get("skip_slides", [])
        if current_goal and current_goal in skip_slides:
            return
        if not os.path.isfile(logo_path):
            return

        position = logo_spec.get("position", "top_right")
        width_inches = logo_spec.get("width_inches", 1.0)
        width_emu = Inches(width_inches)

        try:
            with PILImage.open(logo_path) as img:
                aspect = img.size[1] / img.size[0] if img.size[0] > 0 else 0.5
        except Exception:
            aspect = 0.5
        height_emu = int(width_emu * aspect)

        slide_w, slide_h = prs.slide_width, prs.slide_height
        if position == "top_right":
            left, top = slide_w - width_emu - Inches(0.3), Inches(0.3)
        elif position == "top_left":
            left, top = Inches(0.3), Inches(0.3)
        elif position == "bottom_right":
            left, top = slide_w - width_emu - Inches(0.3), slide_h - height_emu - Inches(0.3)
        else:
            left, top = slide_w - width_emu - Inches(0.3), Inches(0.3)

        slide.shapes.add_picture(logo_path, left, top, width=width_emu, height=height_emu)

    def apply_footer(self, prs: Presentation) -> None:
        if not self._brand or not self._brand.footer:
            return
        config = self._brand.footer
        total = len(prs.slides)
        show_page_number = config.get("show_page_number", False)
        page_number_format = config.get("page_number_format", "{n}")
        page_number_position = config.get("page_number_position", "bottom_right")
        show_footer_text = config.get("show_footer_text", False)
        footer_text = config.get("footer_text", "")
        footer_position = config.get("footer_position", "bottom_center")
        font_size_pt = config.get("font_size_pt", 10)
        skip_pages = config.get("skip_pages", [])

        if not show_page_number and not show_footer_text:
            return

        muted_color = self._rgb(self._c("muted-foreground", "#999999"))
        position_map = {
            "bottom_left": Inches(0.9),
            "bottom_center": Inches(5.833),
            "bottom_right": Inches(11.433),
        }

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            is_cover = idx == 0
            should_skip = is_cover or slide_num in skip_pages

            if show_page_number and not should_skip:
                page_text = page_number_format.replace("{n}", str(slide_num)).replace("{total}", str(total))
                left = position_map.get(page_number_position, Inches(11.433))
                tb = slide.shapes.add_textbox(left, Inches(7.0), Inches(2.0), Inches(0.3))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT
                run = p.add_run()
                run.text = page_text
                run.font.size = Pt(font_size_pt)
                run.font.color.rgb = muted_color
                run.font.name = self._font_b()

            if show_footer_text and footer_text and not should_skip:
                left = position_map.get(footer_position, Inches(5.833))
                tb = slide.shapes.add_textbox(left, Inches(7.0), Inches(4.0), Inches(0.3))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = footer_text
                run.font.size = Pt(font_size_pt)
                run.font.color.rgb = muted_color
                run.font.name = self._font_b()

    def apply_watermark(self, prs: Presentation) -> None:
        if not self._brand or not self._brand.watermark:
            return
        config = self._brand.watermark
        text = config.get("text", "CONFIDENTIAL")
        opacity = config.get("opacity", 0.15)
        rotation = config.get("rotation", -45)
        font_size_pt = config.get("font_size_pt", 72)
        skip_pages = config.get("skip_pages", [1])

        muted_hex = self._c("muted-foreground", "#999999")

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            if slide_num in skip_pages:
                continue

            tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = self._rgb(muted_hex)
            run.font.bold = True
            run.font.name = self._font_h()

            sp = tb._element
            srgbClr_el = sp.find(".//" + qn("a:srgbClr"))
            if srgbClr_el is not None:
                alpha_elem = etree.SubElement(srgbClr_el, qn("a:alpha"))
                alpha_elem.set("val", str(int(opacity * 100000)))

            xfrm_el = sp.find(".//" + qn("a:xfrm"))
            if xfrm_el is not None:
                xfrm_el.set("rot", str(int(rotation * 60000)))

    # ── Save ──

    def save(self, prs: Presentation, output_path: str) -> None:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        prs.save(output_path)
