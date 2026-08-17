"""Quick font scale test."""
from pptx import Presentation
from pptx.util import Inches
from ppt_pro_max.build_helpers import add_slide, svg_chart

svg = '<svg viewBox="0 0 1280 720" xmlns="http://www.w3.org/2000/svg">'
svg += '<rect width="1280" height="720" fill="#0F172A"/>'
svg += '<text x="64" y="120" fill="#FFFFFF" font-size="42" font-weight="700" font-family="Microsoft YaHei">战略层级与关键举措</text>'
svg += '<text x="64" y="180" fill="#94A3B8" font-size="22" font-family="Microsoft YaHei">从战略意图到执行闭环，五层递进</text>'
svg += '<text x="354" y="329" fill="#FFFFFF" font-size="22" text-anchor="middle" font-family="Microsoft YaHei">战略目标</text>'
svg += '<text x="64" y="400" fill="#64748B" font-size="11" font-family="Cascadia Mono">STRATEGY MAP</text>'
svg += '</svg>'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = add_slide(prs)
r = svg_chart(s, svg, 0, 0, 13.333, 7.5)
print(f'shapes={r.shape_count}, warnings={len(r.warnings)}')
prs.save('output/test_font_scale.pptx')
print('Saved')
