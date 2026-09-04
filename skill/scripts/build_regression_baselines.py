"""Build independent neutral regression baselines from case content."""
from __future__ import annotations

import argparse
from io import BytesIO
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
CASES = {
    "technical-infrastructure": ROOT / "examples/new_examplex/ai_infrastructure_economics/output/ai_infrastructure_economics.pptx",
    "scientific-evidence": ROOT / "examples/new_examplex/car_t_single_cell_paper/output/car_t_single_cell_atlas_blue_editorial.pptx",
    "brand-architecture": ROOT / "examples/new_examplex/louvre_abudhabi/output/louvre_abudhabi_complete.pptx",
}


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.text = text or "[no text extracted]"
    for paragraph in frame.paragraphs:
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(25, 25, 25)


def build(regression_id: str, source: Path, output: Path) -> dict:
    source_prs = Presentation(str(source))
    baseline = Presentation()
    baseline.slide_width = source_prs.slide_width
    baseline.slide_height = source_prs.slide_height
    blank = baseline.slide_layouts[6]
    for index, source_slide in enumerate(source_prs.slides, start=1):
        slide = baseline.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 248, 246)
        texts = [shape.text.strip() for shape in source_slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        add_text(slide, f"BASELINE / {regression_id.upper()} / {index:02d}", 0.55, 0.35, 12.0, 0.25, 8, True)
        add_text(slide, texts[0] if texts else f"Slide {index}", 0.65, 0.85, 5.9, 0.75, 22, True)
        body = "\n".join(texts[1:])
        add_text(slide, body, 0.72, 1.8, 5.55, 4.7, 11)
        pictures = [shape for shape in source_slide.shapes if getattr(shape, "shape_type", None) == 13]
        if pictures:
            image = pictures[0].image
            slide.shapes.add_picture(BytesIO(image.blob), Inches(6.65), Inches(1.8), width=Inches(5.8), height=Inches(4.55))
        else:
            shape = slide.shapes.add_shape(1, Inches(6.65), Inches(1.8), Inches(5.8), Inches(4.55))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(225, 225, 220)
            shape.line.color.rgb = RGBColor(170, 170, 165)
            add_text(slide, "NEUTRAL VISUAL FIELD", 7.1, 3.85, 4.9, 0.4, 12, True)
        add_text(slide, "Independent neutral baseline · same case content and page count", 0.65, 7.08, 12.0, 0.2, 7)
    output.parent.mkdir(parents=True, exist_ok=True)
    baseline.save(str(output))
    controls = {
        "schema_version": 1,
        "regression_id": regression_id,
        "source_case_pptx": str(source.relative_to(ROOT)).replace("\\", "/"),
        "baseline_pptx": str(output.relative_to(ROOT)).replace("\\", "/"),
        "slide_count": len(source_prs.slides),
        "content_source": "case PPTX text and picture blobs; no copied PPTX package parts",
        "visual_mode": "neutral_baseline",
        "output_size": "16:9",
        "font": "Arial",
    }
    output.with_name("controls.json").write_text(json.dumps(controls, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return controls


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--regression-id", choices=sorted(CASES))
    args = parser.parse_args()
    targets = {args.regression_id: CASES[args.regression_id]} if args.regression_id else CASES
    results = []
    for regression_id, source in targets.items():
        output = ROOT / "examples/regression_baselines" / regression_id / "baseline.pptx"
        results.append(build(regression_id, source, output))
    print(json.dumps(results, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
