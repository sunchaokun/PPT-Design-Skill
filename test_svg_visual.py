"""Comprehensive SVG Compiler visual test — covers all new features.

Slides:
1. Rounded rect (rx/ry) — various radii
2. tspan dx/dy offsets — subscript/superscript, staggered text
3. Scaling modes — contain vs cover vs stretch
4. Pentagon radar chart (existing regression)
5. Mixed: rounded rect + gradient + tspan dy
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ppt_pro_max.build_helpers import svg_chart, rect, text, page_header

C = {
    "primary": "#4472C4",
    "secondary": "#5B9BD5",
    "accent": "#ED7D31",
    "background": "#FFFFFF",
    "surface": "#F2F2F2",
    "text": "#333333",
    "text_dark": "#000000",
    "muted": "#999999",
    "muted_foreground": "#666666",
    "border": "#D0D0D0",
    "success": "#70AD47",
    "warning": "#FFC000",
    "danger": "#FF0000",
}

OUT_DIR = os.path.join(os.path.dirname(__file__), "output")
os.makedirs(OUT_DIR, exist_ok=True)


def make_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def save_and_preview(prs, name):
    path = os.path.join(OUT_DIR, name)
    prs.save(path)
    print(f"Saved: {path}")
    return path


# ── Slide 1: Rounded rects ──────────────────────────────────────
def slide_rounded_rects(prs):
    slide = make_slide(prs)
    page_header(slide, "Rounded Rect (rx/ry)", "SVG Compiler Feature Test", C=C)

    # Small rx/ry
    svg1 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 120">
      <rect x="10" y="10" width="180" height="100" rx="10" ry="10" fill="#4472C4"/>
      <text x="100" y="65" text-anchor="middle" font-size="16" fill="#fff">rx=10 ry=10</text>
    </svg>'''
    svg_chart(slide, svg1, 0.5, 1.5, 3, 1.8, C=C)

    # Large rx/ry (pill shape)
    svg2 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 80">
      <rect x="10" y="10" width="180" height="60" rx="30" ry="30" fill="#ED7D31"/>
      <text x="100" y="45" text-anchor="middle" font-size="14" fill="#fff">rx=30 ry=30 (pill)</text>
    </svg>'''
    svg_chart(slide, svg2, 4, 1.5, 3, 1.2, C=C)

    # rx only (ry inherits)
    svg3 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 120">
      <rect x="10" y="10" width="180" height="100" rx="20" fill="#70AD47"/>
      <text x="100" y="65" text-anchor="middle" font-size="16" fill="#fff">rx=20 (ry inherits)</text>
    </svg>'''
    svg_chart(slide, svg3, 7.5, 1.5, 3, 1.8, C=C)

    # Oversized rx (clamped to w/2)
    svg4 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 120">
      <rect x="10" y="10" width="180" height="100" rx="200" fill="#FFC000"/>
      <text x="100" y="65" text-anchor="middle" font-size="16" fill="#333">rx=200 (clamped)</text>
    </svg>'''
    svg_chart(slide, svg4, 0.5, 3.8, 3, 1.8, C=C)

    # Sharp rect (no rx/ry) — should use fast path
    svg5 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 120">
      <rect x="10" y="10" width="180" height="100" fill="#5B9BD5"/>
      <text x="100" y="65" text-anchor="middle" font-size="16" fill="#fff">No rx/ry (sharp)</text>
    </svg>'''
    svg_chart(slide, svg5, 4, 3.8, 3, 1.8, C=C)

    # Rounded rect with gradient
    svg6 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 120">
      <defs>
        <linearGradient id="rg1" x1="0" y1="0" x2="1" y2="1">
          <stop offset="0" stop-color="#4472C4"/>
          <stop offset="1" stop-color="#2E5BA6"/>
        </linearGradient>
      </defs>
      <rect x="10" y="10" width="180" height="100" rx="15" ry="15" fill="url(#rg1)"/>
      <text x="100" y="65" text-anchor="middle" font-size="16" fill="#fff">Gradient + rounded</text>
    </svg>'''
    svg_chart(slide, svg6, 7.5, 3.8, 3, 1.8, C=C)


# ── Slide 2: tspan dx/dy ────────────────────────────────────────
def slide_tspan_dx_dy(prs):
    slide = make_slide(prs)
    page_header(slide, "tspan dx/dy Offsets", "SVG Compiler Feature Test", C=C)

    # Subscript via dy
    svg1 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 60">
      <text x="10" y="40" font-size="24" fill="#333">
        H<tspan dy="8" font-size="16" fill="#4472C4">2</tspan>O
      </text>
    </svg>'''
    svg_chart(slide, svg1, 0.5, 1.5, 4, 0.8, C=C)

    # Superscript via dy
    svg2 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 60">
      <text x="10" y="40" font-size="24" fill="#333">
        E=mc<tspan dy="-10" font-size="16" fill="#ED7D31">2</tspan>
      </text>
    </svg>'''
    svg_chart(slide, svg2, 5, 1.5, 4, 0.8, C=C)

    # dx offset — indented text
    svg3 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 60">
      <text x="10" y="40" font-size="18" fill="#333">
        Label:<tspan dx="20" fill="#4472C4">Value with dx=20</tspan>
      </text>
    </svg>'''
    svg_chart(slide, svg3, 0.5, 2.8, 6, 0.8, C=C)

    # Multi-line with dy (staggered)
    svg4 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 200">
      <text x="10" y="30" font-size="16" fill="#333">
        Line 1<tspan dy="25" x="10">Line 2 (dy=25)</tspan><tspan dy="25" x="10">Line 3 (dy=25)</tspan>
      </text>
    </svg>'''
    svg_chart(slide, svg4, 0.5, 4, 4, 2.5, C=C)

    # Mixed dx + dy
    svg5 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 80">
      <text x="10" y="50" font-size="20" fill="#333">
        A<tspan dx="5" dy="-5" font-size="14" fill="#70AD47">up</tspan><tspan dx="5" dy="10" font-size="14" fill="#ED7D31">down</tspan>
      </text>
    </svg>'''
    svg_chart(slide, svg5, 5, 4, 4, 1, C=C)


# ── Slide 3: Scaling modes ──────────────────────────────────────
def slide_scaling(prs):
    slide = make_slide(prs)
    page_header(slide, "Scaling Modes", "contain vs cover vs stretch", C=C)

    # Wide SVG (2:1 aspect) in square region
    wide_svg = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 200">
      <rect x="0" y="0" width="400" height="200" fill="#4472C4" fill-opacity="0.3"/>
      <rect x="5" y="5" width="390" height="190" rx="10" fill="none" stroke="#4472C4" stroke-width="2"/>
      <text x="200" y="110" text-anchor="middle" font-size="24" fill="#4472C4">400×200</text>
    </svg>'''

    # Tall SVG (1:2 aspect) in square region
    tall_svg = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 400">
      <rect x="0" y="0" width="200" height="400" fill="#ED7D31" fill-opacity="0.3"/>
      <rect x="5" y="5" width="190" height="390" rx="10" fill="none" stroke="#ED7D31" stroke-width="2"/>
      <text x="100" y="210" text-anchor="middle" font-size="24" fill="#ED7D31">200×400</text>
    </svg>'''

    # Contain (default) — wide
    text(slide, 0.3, 1.3, 3, 0.4, "contain (wide→square)", font_size=11, C=C)
    svg_chart(slide, wide_svg, 0.3, 1.7, 3, 3, C=C, scaling="contain")

    # Cover — wide
    text(slide, 3.6, 1.3, 3, 0.4, "cover (wide→square)", font_size=11, C=C)
    svg_chart(slide, wide_svg, 3.6, 1.7, 3, 3, C=C, scaling="cover")

    # Stretch — wide
    text(slide, 6.9, 1.3, 3, 0.4, "stretch (wide→square)", font_size=11, C=C)
    svg_chart(slide, wide_svg, 6.9, 1.7, 3, 3, C=C, scaling="stretch")

    # Contain — tall
    text(slide, 0.3, 5.2, 3, 0.4, "contain (tall→square)", font_size=11, C=C)
    svg_chart(slide, tall_svg, 0.3, 5.5, 3, 3, C=C, scaling="contain")

    # Cover — tall
    text(slide, 3.6, 5.2, 3, 0.4, "cover (tall→square)", font_size=11, C=C)
    svg_chart(slide, tall_svg, 3.6, 5.5, 3, 3, C=C, scaling="cover")

    # Stretch — tall
    text(slide, 6.9, 5.2, 3, 0.4, "stretch (tall→square)", font_size=11, C=C)
    svg_chart(slide, tall_svg, 6.9, 5.5, 3, 3, C=C, scaling="stretch")


# ── Slide 4: Pentagon radar (regression) ─────────────────────────
def slide_pentagon(prs):
    slide = make_slide(prs)
    page_header(slide, "Pentagon Radar", "Regression Test", C=C)

    import math
    R = 150
    N = 5
    cx, cy = 250, 250
    labels = ["创新力", "执行力", "影响力", "协作力", "学习力"]
    values = [0.8, 0.6, 0.9, 0.7, 0.85]
    LABEL_DIST = 48

    pts_outer = []
    pts_inner = []
    label_positions = []
    for i in range(N):
        angle = -math.pi / 2 + 2 * math.pi * i / N
        ox = cx + R * math.cos(angle)
        oy = cy + R * math.sin(angle)
        pts_outer.append(f"{ox:.1f},{oy:.1f}")
        ix = cx + R * values[i] * math.cos(angle)
        iy = cy + R * values[i] * math.sin(angle)
        pts_inner.append(f"{ix:.1f},{iy:.1f}")
        lx = cx + (R + LABEL_DIST) * math.cos(angle)
        ly = cy + (R + LABEL_DIST) * math.sin(angle)
        label_positions.append((lx, ly, labels[i]))

    grid_levels = [0.2, 0.4, 0.6, 0.8, 1.0]
    grid_polys = ""
    grid_colors = ["#E8E8E8", "#D8D8D8", "#C8C8C8", "#B8B8B8", "#A8A8A8"]
    for gi, gl in enumerate(grid_levels):
        gpts = []
        for i in range(N):
            angle = -math.pi / 2 + 2 * math.pi * i / N
            gx = cx + R * gl * math.cos(angle)
            gy = cy + R * gl * math.sin(angle)
            gpts.append(f"{gx:.1f},{gy:.1f}")
        grid_polys += f'<polygon points="{" ".join(gpts)}" fill="none" stroke="{grid_colors[gi]}" stroke-width="1"/>'

    spokes = ""
    for i in range(N):
        angle = -math.pi / 2 + 2 * math.pi * i / N
        ex = cx + R * math.cos(angle)
        ey = cy + R * math.sin(angle)
        spokes += f'<line x1="{cx}" y1="{cy}" x2="{ex}" y2="{ey}" stroke="#C0C0C0" stroke-width="1"/>'

    labels_svg = ""
    for lx, ly, lbl in label_positions:
        anchor = "middle"
        if lx < cx - 10:
            anchor = "end"
        elif lx > cx + 10:
            anchor = "start"
        labels_svg += f'<text x="{lx}" y="{ly}" text-anchor="{anchor}" font-size="13" fill="#333">{lbl}</text>'

    svg = f'''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 500 500">
      {grid_polys}
      {spokes}
      <polygon points="{" ".join(pts_outer)}" fill="none" stroke="#4472C4" stroke-width="1.5" stroke-opacity="0.4"/>
      <polygon points="{" ".join(pts_inner)}" fill="#4472C4" fill-opacity="0.25" stroke="#4472C4" stroke-width="2"/>
      {labels_svg}
    </svg>'''

    svg_chart(slide, svg, 1.5, 1.5, 5.5, 5.5, C=C)


# ── Slide 5: Mixed features ─────────────────────────────────────
def slide_mixed(prs):
    slide = make_slide(prs)
    page_header(slide, "Mixed Features", "Rounded rect + gradient + tspan dy", C=C)

    # Card with rounded rect + gradient + subscript
    svg = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 250">
      <defs>
        <linearGradient id="card_bg" x1="0" y1="0" x2="0" y2="1">
          <stop offset="0" stop-color="#4472C4"/>
          <stop offset="1" stop-color="#2E5BA6"/>
        </linearGradient>
      </defs>
      <rect x="20" y="20" width="360" height="210" rx="20" ry="20" fill="url(#card_bg)"/>
      <text x="200" y="80" text-anchor="middle" font-size="28" fill="#fff">Revenue</text>
      <text x="200" y="140" text-anchor="middle" font-size="42" fill="#fff" font-weight="bold">$4.2M</text>
      <text x="200" y="190" text-anchor="middle" font-size="16" fill="#BDD7EE">
        YoY +23<tspan dy="-6" font-size="11">%</tspan><tspan dy="6"> growth</tspan>
      </text>
    </svg>'''
    svg_chart(slide, svg, 0.5, 1.5, 5, 3.1, C=C)

    # Dashboard gauge with rounded rect frame
    svg2 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 200">
      <rect x="10" y="10" width="280" height="180" rx="15" ry="15" fill="#F5F5F5" stroke="#D0D0D0" stroke-width="1"/>
      <text x="150" y="40" text-anchor="middle" font-size="14" fill="#666">Performance Score</text>
      <text x="150" y="120" text-anchor="middle" font-size="48" fill="#4472C4" font-weight="bold">87</text>
      <text x="150" y="155" text-anchor="middle" font-size="14" fill="#70AD47">
        ▲ 12<tspan dy="-5" font-size="10">pts</tspan>
      </text>
    </svg>'''
    svg_chart(slide, svg2, 6, 1.5, 3.5, 2.3, C=C)

    # Rounded rect with stroke-dasharray
    svg3 = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 100">
      <rect x="10" y="10" width="280" height="80" rx="12" ry="12" fill="none" stroke="#ED7D31" stroke-width="2" stroke-dasharray="8,4"/>
      <text x="150" y="55" text-anchor="middle" font-size="16" fill="#ED7D31">Dashed rounded border</text>
    </svg>'''
    svg_chart(slide, svg3, 6, 4.2, 3.5, 1.2, C=C)


# ── Main ─────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_rounded_rects(prs)
    slide_tspan_dx_dy(prs)
    slide_scaling(prs)
    slide_pentagon(prs)
    slide_mixed(prs)

    path = save_and_preview(prs, "svg_compiler_visual_test.pptx")

    from ppt_pro_max.render_preview import render_preview
    result = render_preview(path, engine="powerpoint")
    print(f"Preview: {result}")
    print("Done!")
