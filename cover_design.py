from ppt_pro_max.build_helpers import *
from ppt_pro_max.renderer.text_effects import apply_text_shadow
from lxml import etree
from pptx.oxml.ns import qn

C = {
    'primary': '#1B2A4A',
    'accent': '#C9A96E',
    'on-primary': '#FFFFFF',
    'secondary': '#2C3E6B',
    'background': '#0D1B2A',
    'foreground': '#FFFFFF',
    'muted': '#1A2D45',
    'muted-foreground': '#8A9BB0',
    'border': '#2A3F5F',
    'success': '#2E7D5B',
    'divider': '#C9A96E',
    'text_dark': '#1B2A4A',
    'text_body': '#3D4F5F',
    'text_muted': '#5A6A7A',
    'white': '#FFFFFF',
    'card_bg': '#FFFFFF',
    'font_heading': '微软雅黑',
    'font_body': '微软雅黑',
    'font_cjk': '微软雅黑',
}

prs = Presentation()
set_widescreen(prs)
set_theme_colors(prs, C)

s = add_slide(prs)


def set_alpha(shp, pct):
    sp = shp._element.find(qn('p:spPr'))
    sf = sp.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', str(int(pct * 1000)))


# ══════════════════════════════════════════════════════════════════
# LAYER 1: Darkest background
# ══════════════════════════════════════════════════════════════════

rect(s, 0, 0, 13.333, 7.5, '#0D1B2A', C=C)

# ══════════════════════════════════════════════════════════════════
# LAYER 2: Mid-tone block — large quarter-circle upper-right
# ══════════════════════════════════════════════════════════════════
# Quarter circle: center at (13.333, 0), radius 6.5"
# → occupies upper-right, curves down into slide
# Implemented as oval positioned so only the bottom-left arc shows
interblock = oval(s, 6.833, -6.5, 13.0, 13.0, '#1A334D', C=C)
#                          ^x   ^y   ^w   ^h — center at (13.333, 0)

# Subtle inner highlight — smaller concentric arc, slightly lighter
inner = oval(s, 8.333, -4.5, 10.0, 10.0, '#22415C', C=C)
set_alpha(inner, 50)

# ══════════════════════════════════════════════════════════════════
# LAYER 3:薄金线 — 从色块内部延伸出来，引导视线回标题
# ═════════════════════════════��════════════════════════════════════
# 金线起点在色块内（右上区域），终点延伸到标题附近
# 上方金线
rect(s, 8.5, 0.45, 4.2, 0.025, '#C9A96E', C=C)

# ══════════════════════════════════════════════════════════════════
# LAYER 4: 标题 — 大字号，左对齐，压在色块弧线边缘
# ══════════════════════════════════════════════════════════════════

# 第一行：在深色背景上
t1 = text(s, 0.8, 2.1, 4.5, 1.3, '企业',
          font_size=66, bold=True, color='#FFFFFF', C=C,
          font_name='微软雅黑')
for p in t1.text_frame.paragraphs:
    for run in p.runs:
        apply_text_shadow(run, blur=8, dist=1, direction=90,
                          color="#000000", alpha=8)

# 第二行 — "数字化" 压在色块弧线边缘，半深半浅
t2 = text(s, 0.8, 3.2, 5.0, 1.3, '数字化',
          font_size=66, bold=True, color='#FFFFFF', C=C,
          font_name='微软雅黑')
for p in t2.text_frame.paragraphs:
    for run in p.runs:
        apply_text_shadow(run, blur=8, dist=1, direction=90,
                          color="#000000", alpha=8)

# 第三行 — "转型" 在色块内部
t3 = text(s, 0.8, 4.3, 5.0, 1.3, '转型',
          font_size=66, bold=True, color='#FFFFFF', C=C,
          font_name='微软雅黑')
for p in t3.text_frame.paragraphs:
    for run in p.runs:
        apply_text_shadow(run, blur=8, dist=1, direction=90,
                          color="#000000", alpha=8)

# ══════════════════════════════════════════════════════════════════
# LAYER 5: 副标题
# ══════════════════════════════════════════════════════════════════

text(s, 0.8, 5.7, 7.0, 0.4, '从战略到落地的全链路智能升级方案',
     font_size=15, color='#8A9BB0', C=C,
     font_name='微软雅黑')

# ══════════════════════════════════════════════════════════════════
# LAYER 6: 右下角金色小字 — 部门/日期
# ══════════════════════════════════════════════════════════════════

text(s, 9.5, 6.8, 3.2, 0.3, '战略发展部',
     font_size=11, bold=True, color='#C9A96E', C=C,
     font_name='微软雅黑', align='right')

text(s, 9.5, 7.1, 3.2, 0.3, '2026年8月',
     font_size=10, color='#5A6A7A', C=C,
     font_name='微软雅黑', align='right')

# ══════════════════════════════════════════════════════════════════
# LAYER 7: 左下角极小英文标识页眉
# ══════════════════════════════════════════════════════════════════

label = text(s, 0.8, 0.5, 5.0, 0.25,
            'AI-DRIVEN ENTERPRISE TRANSFORMATION',
            font_size=9, color='#5A6A7A', C=C,
            font_name='Arial')
for p in label.text_frame.paragraphs:
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', '350')

# ══════════════════════════════════════════════════════════════════
# LAYER 8: Transition
# ══════════════════════════════════════════════════════════════════

slide_transition(s, 'fade')

VERSION = 'v10'
out = f'output/cover_{VERSION}.pptx'
clean_save(prs, out)
print(f'DONE: {out}')
