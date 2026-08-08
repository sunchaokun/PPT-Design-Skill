"""Tests for Phase 1 API exposure in PrecisionRenderer and build_helpers.

Covers:
  1. add_text_with_gradient() — gradient fill on text characters
  2. add_vertical_text() — vertical (top-to-bottom) text column
  3. add_seal_stamp() — traditional Chinese seal stamp element
  4. gradient_text() — build_helpers wrapper
  5. vertical_text() — build_helpers wrapper
  6. seal_stamp() — build_helpers wrapper
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from ppt_pro_max.build_helpers import gradient_text, seal_stamp, vertical_text
from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer


def _make_brand() -> BrandSpec:
    return BrandSpec(
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "accent": "#6366F1",
            "background": "#FFFFFF",
            "foreground": "#0F172A",
            "muted": "#F1F5F9",
            "muted-foreground": "#64748B",
            "border": "#E2E8F0",
        },
        fonts={"heading": "Inter", "body": "Inter"},
    )


def _make_pr() -> PrecisionRenderer:
    return PrecisionRenderer(brand_spec=_make_brand())


def _make_slide():
    pr = _make_pr()
    prs = pr.create_presentation()
    slide = pr.add_slide(prs)
    return pr, prs, slide


def _make_build_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    return prs, slide


# ── PrecisionRenderer: add_text_with_gradient ──


class TestAddTextWithGradient:
    def test_creates_textbox_with_gradient_fill(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_text_with_gradient(slide, 1.0, 1.0, 8.0, 1.5, "Hello",
                                          gradient_preset="gold-shine", font_size=44)
        assert shape is not None
        assert shape.has_text_frame
        assert shape.text_frame.text == "Hello"
        run = shape.text_frame.paragraphs[0].runs[0]
        rPr = run._r.find(qn("a:rPr"))
        gradFill = rPr.find(qn("a:gradFill"))
        assert gradFill is not None, "Text run should have gradient fill"
        gsLst = gradFill.find(qn("a:gsLst"))
        assert gsLst is not None
        stops = gsLst.findall(qn("a:gs"))
        assert len(stops) == 3

    def test_custom_stops_override_preset(self):
        pr, prs, slide = _make_slide()
        custom_stops = [("FF0000", 0), ("0000FF", 100000)]
        shape = pr.add_text_with_gradient(slide, 1.0, 1.0, 8.0, 1.5, "Custom",
                                          gradient_stops=custom_stops, font_size=44)
        run = shape.text_frame.paragraphs[0].runs[0]
        rPr = run._r.find(qn("a:rPr"))
        gradFill = rPr.find(qn("a:gradFill"))
        stops = gradFill.find(qn("a:gsLst")).findall(qn("a:gs"))
        assert len(stops) == 2

    def test_font_size_and_bold(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_text_with_gradient(slide, 1.0, 1.0, 8.0, 1.5, "Bold",
                                          gradient_preset="blue-deep", font_size=36, bold=True)
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(36)
        assert run.font.bold is True

    def test_invalid_preset_raises(self):
        pr, prs, slide = _make_slide()
        with pytest.raises(KeyError, match="Unknown text gradient preset"):
            pr.add_text_with_gradient(slide, 1.0, 1.0, 8.0, 1.5, "X",
                                     gradient_preset="nonexistent")

    def test_default_preset_is_gold_shine(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_text_with_gradient(slide, 1.0, 1.0, 8.0, 1.5, "Default")
        run = shape.text_frame.paragraphs[0].runs[0]
        rPr = run._r.find(qn("a:rPr"))
        gradFill = rPr.find(qn("a:gradFill"))
        first_stop = gradFill.find(qn("a:gsLst")).findall(qn("a:gs"))[0]
        srgb = first_stop.find(qn("a:srgbClr"))
        assert srgb.get("val") == "F5AF19"

    def test_custom_font(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_text_with_gradient(slide, 1.0, 1.0, 8.0, 1.5, "Font",
                                          gradient_preset="gold-shine", font_size=44,
                                          font="Georgia")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Georgia"

    def test_cjk_font_pairing(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_text_with_gradient(slide, 1.0, 1.0, 8.0, 1.5, "测试",
                                          gradient_preset="gold-shine", font_size=44)
        run = shape.text_frame.paragraphs[0].runs[0]
        rPr = run._r.find(qn("a:rPr"))
        ea = rPr.find(qn("a:ea"))
        assert ea is not None
        assert ea.get("typeface") is not None


# ── PrecisionRenderer: add_vertical_text ──


class TestAddVerticalText:
    def test_creates_vertical_textbox(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "诗词",
                                     font_size=24)
        assert shape is not None
        assert shape.has_text_frame
        bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        assert bodyPr is not None
        assert bodyPr.get("vert") == "eaVert"

    def test_custom_direction(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "Text",
                                     direction="270", font_size=24)
        bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        assert bodyPr.get("vert") == "vert270"

    def test_font_name_and_size(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "测试",
                                     font_name="STKaiti", font_size=28)
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(28)
        assert run.font.name == "STKaiti"

    def test_color_role(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "Color",
                                     color_role="foreground", font_size=24)
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb is not None

    def test_color_hex_overrides_role(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "Hex",
                                     color_hex="#FF0000", font_size=24)
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor.from_string("FF0000")

    def test_cjk_font_pairing(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "测试",
                                     font_name="STKaiti", font_size=24)
        run = shape.text_frame.paragraphs[0].runs[0]
        rPr = run._r.find(qn("a:rPr"))
        ea = rPr.find(qn("a:ea"))
        assert ea is not None
        assert ea.get("typeface") is not None


# ── PrecisionRenderer: add_seal_stamp ──


class TestAddSealStamp:
    def test_creates_seal_stamp(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "印")
        assert shape is not None
        assert shape.has_text_frame
        assert "印" in shape.text_frame.text

    def test_seal_is_red_square(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "章")
        spPr = shape._element.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb is not None
        assert srgb.get("val") == "C41E3A"

    def test_seal_text_is_white(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "印")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor.from_string("FFFFFF")

    def test_seal_rotation(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "印", rotation=-15)
        xfrm = shape._element.find(".//" + qn("a:xfrm"))
        assert xfrm is not None
        assert int(xfrm.get("rot", "0")) != 0

    def test_seal_no_rotation(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "印", rotation=0)
        xfrm = shape._element.find(".//" + qn("a:xfrm"))
        assert xfrm is not None
        assert xfrm.get("rot") is None

    def test_custom_fill_color(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "印",
                                  fill_hex="#8B0000")
        spPr = shape._element.find(qn("p:spPr"))
        srgb = spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        assert srgb.get("val") == "8B0000"

    def test_seal_border_is_lighter_than_fill(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "印")
        spPr = shape._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        assert ln is not None
        solidFill = ln.find(qn("a:solidFill"))
        assert solidFill is not None
        border_srgb = solidFill.find(qn("a:srgbClr"))
        assert border_srgb is not None
        border_val = border_srgb.get("val")
        fill_srgb = spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        fill_val = fill_srgb.get("val")
        assert border_val != fill_val, "Border should be lighter than fill"

    def test_seal_font_name(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "印",
                                  font_name="STKaiti")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "STKaiti"

    def test_seal_font_size_minimum_11pt(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 0.5, "印")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.size >= Pt(11)

    def test_seal_has_text_outline(self):
        pr, prs, slide = _make_slide()
        shape = pr.add_seal_stamp(slide, 5.0, 5.0, 1.0, "印")
        run = shape.text_frame.paragraphs[0].runs[0]
        rPr = run._r.find(qn("a:rPr"))
        ln = rPr.find(qn("a:ln"))
        assert ln is not None, "Seal text should have outline"


# ── build_helpers: gradient_text ──


class TestGradientTextBuildHelper:
    def test_creates_gradient_text(self):
        prs, slide = _make_build_slide()
        shape = gradient_text(slide, 1.0, 1.0, 8.0, 1.5, "Hello",
                              preset="gold-shine", font_size=44)
        assert shape is not None
        assert shape.has_text_frame
        assert shape.text_frame.text == "Hello"
        run = shape.text_frame.paragraphs[0].runs[0]
        rPr = run._r.find(qn("a:rPr"))
        assert rPr.find(qn("a:gradFill")) is not None

    def test_custom_stops(self):
        prs, slide = _make_build_slide()
        shape = gradient_text(slide, 1.0, 1.0, 8.0, 1.5, "Custom",
                              stops=[("FF0000", 0), ("0000FF", 100000)],
                              font_size=44)
        run = shape.text_frame.paragraphs[0].runs[0]
        gradFill = run._r.find(qn("a:rPr")).find(qn("a:gradFill"))
        assert len(gradFill.find(qn("a:gsLst")).findall(qn("a:gs"))) == 2

    def test_font_and_bold(self):
        prs, slide = _make_build_slide()
        shape = gradient_text(slide, 1.0, 1.0, 8.0, 1.5, "Bold",
                              preset="blue-deep", font_size=36, bold=True,
                              font_name="Arial")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(36)
        assert run.font.bold is True
        assert run.font.name == "Arial"

    def test_invalid_preset_raises(self):
        prs, slide = _make_build_slide()
        with pytest.raises(KeyError, match="Unknown text gradient preset"):
            gradient_text(slide, 1.0, 1.0, 8.0, 1.5, "X",
                         preset="nonexistent", font_size=44)


# ── build_helpers: vertical_text ──


class TestVerticalTextBuildHelper:
    def test_creates_vertical_text(self):
        prs, slide = _make_build_slide()
        shape = vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "诗词",
                              font_size=24)
        assert shape is not None
        bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        assert bodyPr.get("vert") == "eaVert"

    def test_custom_direction(self):
        prs, slide = _make_build_slide()
        shape = vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "Text",
                              direction="270", font_size=24)
        bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        assert bodyPr.get("vert") == "vert270"

    def test_font_name(self):
        prs, slide = _make_build_slide()
        shape = vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "测试",
                              font_name="STKaiti", font_size=28)
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "STKaiti"

    def test_color(self):
        prs, slide = _make_build_slide()
        shape = vertical_text(slide, 1.0, 1.0, 1.5, 5.0, "Color",
                              color="#FF0000", font_size=24)
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor.from_string("FF0000")


# ── build_helpers: seal_stamp ──


class TestSealStampBuildHelper:
    def test_creates_seal_stamp(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "印")
        assert shape is not None
        assert "印" in shape.text_frame.text

    def test_seal_is_red(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "章")
        spPr = shape._element.find(qn("p:spPr"))
        srgb = spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        assert srgb.get("val") == "C41E3A"

    def test_seal_rotation(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "印", rotation=-15)
        xfrm = shape._element.find(".//" + qn("a:xfrm"))
        assert int(xfrm.get("rot", "0")) != 0

    def test_seal_no_rotation(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "印", rotation=0)
        xfrm = shape._element.find(".//" + qn("a:xfrm"))
        assert xfrm.get("rot") is None

    def test_custom_color(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "印", fill_hex="#8B0000")
        spPr = shape._element.find(qn("p:spPr"))
        srgb = spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        assert srgb.get("val") == "8B0000"

    def test_seal_border_matches_fill_for_zhu_style(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "\u5370")
        spPr = shape._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        border_srgb = ln.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        fill_srgb = spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        assert border_srgb.get("val") == fill_srgb.get("val")

    def test_seal_font_name(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "\u5370", font_name="STKaiti")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "STKaiti"

    def test_seal_default_font_is_zhongsong(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "\u5370")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "STZhongsong"

    def test_seal_font_size_minimum_14pt(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 0.5, "\u5370")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.size >= Pt(14)

    def test_seal_zhu_style_white_text(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "\u5370", style="zhu")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor.from_string("FFFFFF")

    def test_seal_bai_style_red_text(self):
        prs, slide = _make_build_slide()
        shape = seal_stamp(slide, 5.0, 5.0, 1.0, "\u5370", style="bai")
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor.from_string("C41E3A")
