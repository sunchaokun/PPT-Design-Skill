"""Tests for Phase 6: Decoration Library.

Covers:
  1. add_brush_divider() — ink brush stroke divider
  2. add_seal_stamp() — traditional Chinese seal stamp (decoration version)
  3. add_scroll_frame() — scroll/panel border frame
  4. add_neon_border() — glowing neon border
  5. add_grid_background() — subtle grid background
  6. add_glass_panel() — semi-transparent panel
  7. add_ink_splash() — ink splash decoration
  8. API exposure in build_helpers
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_pro_max.renderer.decoration_library import (
    add_brush_divider,
    add_seal_stamp,
    add_scroll_frame,
    add_neon_border,
    add_grid_background,
    add_glass_panel,
    add_ink_splash,
)


def _make_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    return prs, slide


# ── add_brush_divider ──


class TestBrushDivider:
    def test_returns_element(self):
        prs, slide = _make_slide()
        elem = add_brush_divider(slide, 1.0, 3.0, 8.0)
        assert elem is not None

    def test_creates_freeform_shape(self):
        prs, slide = _make_slide()
        elem = add_brush_divider(slide, 1.0, 3.0, 8.0)
        assert elem.tag == qn("p:sp")
        spPr = elem.find(qn("p:spPr"))
        assert spPr is not None
        custGeom = spPr.find(qn("a:custGeom"))
        assert custGeom is not None

    def test_has_fill_color(self):
        prs, slide = _make_slide()
        elem = add_brush_divider(slide, 1.0, 3.0, 8.0, color="#2C2C2C")
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "2C2C2C"

    def test_custom_color(self):
        prs, slide = _make_slide()
        elem = add_brush_divider(slide, 1.0, 3.0, 8.0, color="#8B0000")
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "8B0000"

    def test_has_path_with_bezier_curves(self):
        prs, slide = _make_slide()
        elem = add_brush_divider(slide, 1.0, 3.0, 8.0)
        spPr = elem.find(qn("p:spPr"))
        custGeom = spPr.find(qn("a:custGeom"))
        pathLst = custGeom.find(qn("a:pathLst"))
        assert pathLst is not None
        paths = pathLst.findall(qn("a:path"))
        assert len(paths) >= 1
        beziers = paths[0].findall(qn("a:cubicBezTo"))
        assert len(beziers) >= 2

    def test_custom_thickness(self):
        prs, slide = _make_slide()
        elem = add_brush_divider(slide, 1.0, 3.0, 8.0, thickness=0.12)
        assert elem is not None


# ── add_seal_stamp ──


class TestSealStamp:
    def test_returns_element(self):
        prs, slide = _make_slide()
        elem = add_seal_stamp(slide, 10.0, 6.0, 0.8, "印")
        assert elem is not None

    def test_is_rectangle_shape(self):
        prs, slide = _make_slide()
        elem = add_seal_stamp(slide, 10.0, 6.0, 0.8, "印")
        assert elem.tag == qn("p:sp")
        spPr = elem.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        assert prstGeom is not None
        assert prstGeom.get("prst") == "rect"

    def test_has_red_fill(self):
        prs, slide = _make_slide()
        elem = add_seal_stamp(slide, 10.0, 6.0, 0.8, "印")
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "C41E3A"

    def test_has_white_text(self):
        prs, slide = _make_slide()
        elem = add_seal_stamp(slide, 10.0, 6.0, 0.8, "印")
        txBody = elem.find(qn("p:txBody"))
        assert txBody is not None
        runs = txBody.findall(".//" + qn("a:r"))
        assert len(runs) >= 1
        rPr = runs[0].find(qn("a:rPr"))
        solidFill = rPr.find(qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "FFFFFF"

    def test_has_rotation(self):
        prs, slide = _make_slide()
        elem = add_seal_stamp(slide, 10.0, 6.0, 0.8, "印", rotation=-15)
        spPr = elem.find(qn("p:spPr"))
        xfrm = spPr.find(qn("a:xfrm"))
        assert xfrm is not None
        rot_val = xfrm.get("rot")
        assert rot_val is not None

    def test_custom_fill_color(self):
        prs, slide = _make_slide()
        elem = add_seal_stamp(slide, 10.0, 6.0, 0.8, "印", fill_hex="#000080")
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "000080"


# ── add_scroll_frame ──


class TestScrollFrame:
    def test_returns_element(self):
        prs, slide = _make_slide()
        elem = add_scroll_frame(slide, 2.0, 1.0, 6.0, 5.0)
        assert elem is not None

    def test_creates_rectangle(self):
        prs, slide = _make_slide()
        elem = add_scroll_frame(slide, 2.0, 1.0, 6.0, 5.0)
        assert elem.tag == qn("p:sp")
        spPr = elem.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        assert prstGeom is not None

    def test_has_border_line(self):
        prs, slide = _make_slide()
        elem = add_scroll_frame(slide, 2.0, 1.0, 6.0, 5.0)
        spPr = elem.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        assert ln is not None
        solidFill = ln.find(qn("a:solidFill"))
        assert solidFill is not None

    def test_xuan_style_has_warm_border(self):
        prs, slide = _make_slide()
        elem = add_scroll_frame(slide, 2.0, 1.0, 6.0, 5.0, style="xuan")
        spPr = elem.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        solidFill = ln.find(qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") is not None

    def test_silk_style_has_different_border(self):
        prs, slide = _make_slide()
        elem = add_scroll_frame(slide, 2.0, 1.0, 6.0, 5.0, style="silk")
        spPr = elem.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        assert ln is not None


# ── add_neon_border ──


class TestNeonBorder:
    def test_returns_element(self):
        prs, slide = _make_slide()
        elem = add_neon_border(slide, 1.0, 1.0, 4.0, 3.0, color="#8B5CF6")
        assert elem is not None

    def test_is_rectangle(self):
        prs, slide = _make_slide()
        elem = add_neon_border(slide, 1.0, 1.0, 4.0, 3.0, color="#8B5CF6")
        spPr = elem.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        assert prstGeom is not None

    def test_has_glow_effect(self):
        prs, slide = _make_slide()
        elem = add_neon_border(slide, 1.0, 1.0, 4.0, 3.0, color="#8B5CF6")
        spPr = elem.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        assert effectLst is not None
        glow = effectLst.find(qn("a:glow"))
        assert glow is not None

    def test_has_colored_line(self):
        prs, slide = _make_slide()
        elem = add_neon_border(slide, 1.0, 1.0, 4.0, 3.0, color="#22D3EE")
        spPr = elem.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        solidFill = ln.find(qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "22D3EE"

    def test_no_fill(self):
        prs, slide = _make_slide()
        elem = add_neon_border(slide, 1.0, 1.0, 4.0, 3.0, color="#8B5CF6")
        spPr = elem.find(qn("p:spPr"))
        noFill = spPr.find(qn("a:noFill"))
        assert noFill is not None


# ── add_grid_background ──


class TestGridBackground:
    def test_returns_elements(self):
        prs, slide = _make_slide()
        elems = add_grid_background(slide, spacing=1.0, color="#E0E0E0")
        assert isinstance(elems, list)
        assert len(elems) >= 1

    def test_creates_shapes(self):
        prs, slide = _make_slide()
        elems = add_grid_background(slide, spacing=1.0, color="#E0E0E0")
        for e in elems:
            assert e.tag == qn("p:sp")

    def test_lines_have_color(self):
        prs, slide = _make_slide()
        elems = add_grid_background(slide, spacing=1.0, color="#E0E0E0")
        for e in elems:
            spPr = e.find(qn("p:spPr"))
            solidFill = spPr.find(qn("a:solidFill"))
            assert solidFill is not None
            srgb = solidFill.find(qn("a:srgbClr"))
            assert srgb.get("val") == "E0E0E0"

    def test_custom_alpha(self):
        prs, slide = _make_slide()
        elems = add_grid_background(slide, spacing=1.0, color="#E0E0E0", alpha=20)
        assert len(elems) >= 1


# ── add_glass_panel ──


class TestGlassPanel:
    def test_returns_element(self):
        prs, slide = _make_slide()
        elem = add_glass_panel(slide, 1.0, 1.0, 4.0, 3.0)
        assert elem is not None

    def test_is_rectangle(self):
        prs, slide = _make_slide()
        elem = add_glass_panel(slide, 1.0, 1.0, 4.0, 3.0)
        spPr = elem.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        assert prstGeom is not None

    def test_has_semi_transparent_fill(self):
        prs, slide = _make_slide()
        elem = add_glass_panel(slide, 1.0, 1.0, 4.0, 3.0,
                               tint="#FFFFFF", alpha=15)
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "FFFFFF"
        alpha_el = srgb.find(qn("a:alpha"))
        assert alpha_el is not None
        assert alpha_el.get("val") == "15000"

    def test_has_soft_edge(self):
        prs, slide = _make_slide()
        elem = add_glass_panel(slide, 1.0, 1.0, 4.0, 3.0, soft_edge=8)
        spPr = elem.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        assert effectLst is not None
        softEdge = effectLst.find(qn("a:softEdge"))
        assert softEdge is not None

    def test_custom_tint_color(self):
        prs, slide = _make_slide()
        elem = add_glass_panel(slide, 1.0, 1.0, 4.0, 3.0,
                               tint="#1A1A3A", alpha=20)
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "1A1A3A"


# ── add_ink_splash ──


class TestInkSplash:
    def test_returns_element(self):
        prs, slide = _make_slide()
        elem = add_ink_splash(slide, 5.0, 3.0, 1.0)
        assert elem is not None

    def test_creates_freeform_shape(self):
        prs, slide = _make_slide()
        elem = add_ink_splash(slide, 5.0, 3.0, 1.0)
        spPr = elem.find(qn("p:spPr"))
        custGeom = spPr.find(qn("a:custGeom"))
        assert custGeom is not None

    def test_has_fill_color(self):
        prs, slide = _make_slide()
        elem = add_ink_splash(slide, 5.0, 3.0, 1.0, color="#2C2C2C")
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "2C2C2C"

    def test_custom_alpha(self):
        prs, slide = _make_slide()
        elem = add_ink_splash(slide, 5.0, 3.0, 1.0, color="#2C2C2C", alpha=30)
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        alpha_el = srgb.find(qn("a:alpha"))
        assert alpha_el is not None
        assert alpha_el.get("val") == "30000"

    def test_has_organic_path(self):
        prs, slide = _make_slide()
        elem = add_ink_splash(slide, 5.0, 3.0, 1.0)
        spPr = elem.find(qn("p:spPr"))
        custGeom = spPr.find(qn("a:custGeom"))
        pathLst = custGeom.find(qn("a:pathLst"))
        paths = pathLst.findall(qn("a:path"))
        assert len(paths) >= 1
        beziers = paths[0].findall(qn("a:cubicBezTo"))
        assert len(beziers) >= 4


# ── build_helpers API ──


class TestBuildHelpersDecorationAPI:
    def test_brush_divider_import(self):
        from ppt_pro_max.build_helpers import brush_divider
        assert callable(brush_divider)

    def test_neon_border_import(self):
        from ppt_pro_max.build_helpers import neon_border
        assert callable(neon_border)

    def test_glass_panel_import(self):
        from ppt_pro_max.build_helpers import glass_panel
        assert callable(glass_panel)

    def test_grid_bg_import(self):
        from ppt_pro_max.build_helpers import grid_background
        assert callable(grid_background)

    def test_ink_splash_import(self):
        from ppt_pro_max.build_helpers import ink_splash
        assert callable(ink_splash)

    def test_brush_divider_creates_shape(self):
        from ppt_pro_max.build_helpers import brush_divider
        prs, slide = _make_slide()
        elem = brush_divider(slide, 1.0, 3.0, 8.0)
        assert elem is not None

    def test_neon_border_creates_shape(self):
        from ppt_pro_max.build_helpers import neon_border
        prs, slide = _make_slide()
        elem = neon_border(slide, 1.0, 1.0, 4.0, 3.0, color="#8B5CF6")
        assert elem is not None

    def test_glass_panel_creates_shape(self):
        from ppt_pro_max.build_helpers import glass_panel
        prs, slide = _make_slide()
        elem = glass_panel(slide, 1.0, 1.0, 4.0, 3.0)
        assert elem is not None

    def test_ink_splash_creates_shape(self):
        from ppt_pro_max.build_helpers import ink_splash
        prs, slide = _make_slide()
        elem = ink_splash(slide, 5.0, 3.0, 1.0)
        assert elem is not None
