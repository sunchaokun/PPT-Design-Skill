"""Tests for P10-P11: Beautify API integration and end-to-end tests."""

import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


FIXTURES = Path(__file__).parent / "fixtures"


def _create_multi_slide_pptx(path: str, num_slides: int = 5):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    titles = ["Product Launch", "Key Features", "How It Works", "Market Opportunity", "Thank You"]
    bullet_sets = [
        [],
        ["AI Engine", "Live Dashboard", "Integration"],
        ["Step 1: Analyze", "Step 2: Design", "Step 3: Deploy"],
        ["TAM: $50B", "Growth: 25% YoY"],
        [],
    ]

    for i in range(min(num_slides, len(titles))):
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = titles[i]
        run.font.size = Pt(44) if i in (0, num_slides - 1) else Pt(32)
        run.font.bold = True

        if bullet_sets[i]:
            body = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(4))
            tf = body.text_frame
            tf.word_wrap = True
            for j, text in enumerate(bullet_sets[i]):
                p2 = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = text
                run2.font.size = Pt(18)

    prs.save(path)
    return path


class TestBeautifyAPI:
    def test_generate_ppt_beautify_param(self, tmp_path):
        from ppt_pro_max import generate_ppt

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        _create_multi_slide_pptx(input_path, num_slides=3)

        result = generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="professional",
            output=output_path,
        )

        assert result.get("mode") == "beautify"
        assert os.path.isfile(output_path)

    def test_generate_ppt_beautify_preserves_slide_count(self, tmp_path):
        from ppt_pro_max import generate_ppt

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        _create_multi_slide_pptx(input_path, num_slides=5)

        result = generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="warm elegant",
            output=output_path,
        )

        assert result["num_slides"] == 5
        out_prs = Presentation(output_path)
        assert len(out_prs.slides) == 5

    def test_generate_ppt_beautify_content_preserved(self, tmp_path):
        from ppt_pro_max import generate_ppt

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        _create_multi_slide_pptx(input_path, num_slides=3)

        generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="professional",
            output=output_path,
        )

        out_prs = Presentation(output_path)
        all_text = []
        for slide in out_prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text)
        combined = " ".join(all_text)
        assert "Product Launch" in combined
        assert "Key Features" in combined

    def test_generate_ppt_beautify_with_density(self, tmp_path):
        from ppt_pro_max import generate_ppt

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        _create_multi_slide_pptx(input_path, num_slides=3)

        result = generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="professional",
            output=output_path,
            density=8,
        )

        assert os.path.isfile(output_path)

    def test_generate_ppt_beautify_with_motion(self, tmp_path):
        from ppt_pro_max import generate_ppt

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        _create_multi_slide_pptx(input_path, num_slides=3)

        result = generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="professional",
            output=output_path,
            motion=5,
        )

        assert os.path.isfile(output_path)


class TestBeautifyE2E:
    def test_beautify_multi_slide_all_goals(self, tmp_path):
        from ppt_pro_max import generate_ppt

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        _create_multi_slide_pptx(input_path, num_slides=5)

        result = generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="dark tech",
            output=output_path,
        )

        assert result["mode"] == "beautify"
        assert result["num_slides"] == 5
        assert os.path.isfile(output_path)

        out_prs = Presentation(output_path)
        assert len(out_prs.slides) == 5

        for slide in out_prs.slides:
            assert len(slide.shapes) >= 1

    def test_beautify_output_is_valid_pptx(self, tmp_path):
        from ppt_pro_max import generate_ppt

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        _create_multi_slide_pptx(input_path, num_slides=3)

        generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="professional",
            output=output_path,
        )

        prs = Presentation(output_path)
        assert prs.slide_width > 0
        assert prs.slide_height > 0

    def test_beautify_font_sizes_reasonable(self, tmp_path):
        from ppt_pro_max import generate_ppt
        from pptx.util import Pt

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        _create_multi_slide_pptx(input_path, num_slides=3)

        generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="professional",
            output=output_path,
        )

        prs = Presentation(output_path)
        min_font = 100
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                pt_val = run.font.size / 12700
                                if pt_val > 0 and pt_val < min_font:
                                    min_font = pt_val
        assert min_font >= 8

    def test_beautify_different_styles(self, tmp_path):
        from ppt_pro_max import generate_ppt

        input_path = str(tmp_path / "input.pptx")
        _create_multi_slide_pptx(input_path, num_slides=3)

        for style_name in ["professional", "warm elegant", "dark tech"]:
            output_path = str(tmp_path / f"output_{style_name.replace(' ', '_')}.pptx")
            result = generate_ppt(
                query="beautify test",
                beautify=input_path,
                style=style_name,
                output=output_path,
            )
            assert os.path.isfile(output_path), f"Failed for style: {style_name}"

    def test_beautify_single_slide_pptx(self, tmp_path):
        from ppt_pro_max import generate_ppt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Single Slide"
        run.font.size = Pt(48)

        input_path = str(tmp_path / "single.pptx")
        output_path = str(tmp_path / "output.pptx")
        prs.save(input_path)

        result = generate_ppt(
            query="beautify test",
            beautify=input_path,
            style="professional",
            output=output_path,
        )

        assert result["mode"] == "beautify"
        assert result["num_slides"] >= 1

    def test_beautify_nonexistent_input(self, tmp_path):
        from ppt_pro_max import generate_ppt

        result = generate_ppt(
            query="beautify test",
            beautify=str(tmp_path / "nonexistent.pptx"),
            style="professional",
        )

        assert result.get("mode") == "beautify"
        assert result.get("num_slides") == 0 or result.get("error") is not None
