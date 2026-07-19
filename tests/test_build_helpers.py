import os
import copy
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


@pytest.fixture
def C():
    return {
        'primary': '#2E6504',
        'accent': '#7DA92F',
        'muted': '#84AF7D',
        'light': '#D4E3AC',
        'background': '#FFFFFF',
        'bg_tint': '#F2F8D6',
        'white': '#FFFFFF',
        'text_dark': '#0D4609',
        'text_body': '#2E6504',
        'text_muted': '#466740',
        'divider': '#84AF7D',
        'card_bg': '#F6FAE8',
        'font_heading': 'Microsoft YaHei',
        'font_body': 'Microsoft YaHei',
    }


@pytest.fixture
def prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


@pytest.fixture
def slide(prs):
    from ppt_pro_max.build_helpers import add_slide
    return add_slide(prs)


class TestResolveColor:

    def test_hex_passthrough(self, C):
        from ppt_pro_max.build_helpers import _resolve_color
        assert _resolve_color('#FF0000', C) == '#FF0000'

    def test_role_lookup(self, C):
        from ppt_pro_max.build_helpers import _resolve_color
        assert _resolve_color('primary', C) == '#2E6504'

    def test_missing_role_fallback(self, C):
        from ppt_pro_max.build_helpers import _resolve_color
        assert _resolve_color('nonexistent', C) == '#000000'

    def test_none_fallback(self, C):
        from ppt_pro_max.build_helpers import _resolve_color
        assert _resolve_color(None, C) == '#000000'

    def test_empty_C(self):
        from ppt_pro_max.build_helpers import _resolve_color
        assert _resolve_color('primary', None) == '#000000'
        assert _resolve_color('#FF0000', None) == '#FF0000'


class TestAddSlide:

    def test_adds_slide(self, prs):
        from ppt_pro_max.build_helpers import add_slide
        slide = add_slide(prs)
        assert len(prs.slides) == 1

    def test_specific_layout(self, prs):
        from ppt_pro_max.build_helpers import add_slide
        slide = add_slide(prs, layout_index=6)
        assert len(prs.slides) == 1


class TestRect:

    def test_creates_shape(self, slide, C):
        from ppt_pro_max.build_helpers import rect
        sh = rect(slide, 0, 0, 13.333, 0.08, 'accent', C=C)
        assert sh is not None
        assert len(slide.shapes) == 1

    def test_fill_color_by_role(self, slide, C):
        from ppt_pro_max.build_helpers import rect
        sh = rect(slide, 0, 0, 1, 1, 'primary', C=C)
        assert sh.fill.fore_color.rgb == RGBColor(0x2E, 0x65, 0x04)

    def test_fill_color_by_hex(self, slide, C):
        from ppt_pro_max.build_helpers import rect
        sh = rect(slide, 0, 0, 1, 1, '#FF0000', C=C)
        assert sh.fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_border(self, slide, C):
        from ppt_pro_max.build_helpers import rect
        sh = rect(slide, 0, 0, 1, 1, 'primary', line='accent', C=C)
        assert sh.line.color.rgb == RGBColor(0x7D, 0xA9, 0x2F)


class TestRRect:

    def test_creates_shape(self, slide, C):
        from ppt_pro_max.build_helpers import rrect
        sh = rrect(slide, 1, 1, 3, 2, 'primary', C=C)
        assert sh is not None
        assert len(slide.shapes) == 1


class TestOval:

    def test_creates_shape(self, slide, C):
        from ppt_pro_max.build_helpers import oval
        sh = oval(slide, 5, 5, 1, 1, 'accent', C=C)
        assert sh is not None


class TestText:

    def test_creates_text(self, slide, C):
        from ppt_pro_max.build_helpers import text
        tb = text(slide, 0.65, 0.45, 12, 0.5, 'Hello', font_size=28, color='text_dark', bold=True, C=C)
        assert tb.text_frame.text == 'Hello'

    def test_color_by_role(self, slide, C):
        from ppt_pro_max.build_helpers import text
        tb = text(slide, 0.65, 0.45, 12, 0.5, 'Test', color='primary', C=C)
        p = tb.text_frame.paragraphs[0]
        assert p.font.color.rgb == RGBColor(0x2E, 0x65, 0x04)

    def test_color_by_hex(self, slide, C):
        from ppt_pro_max.build_helpers import text
        tb = text(slide, 0.65, 0.45, 12, 0.5, 'Test', color='#FF0000', C=C)
        p = tb.text_frame.paragraphs[0]
        assert p.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_font_name(self, slide, C):
        from ppt_pro_max.build_helpers import text
        tb = text(slide, 0.65, 0.45, 12, 0.5, 'Test', font_name='Arial', C=C)
        assert tb.text_frame.paragraphs[0].font.name == 'Arial'


class TestMultiline:

    def test_creates_lines(self, slide, C):
        from ppt_pro_max.build_helpers import multiline
        tb = multiline(slide, 0.65, 1.5, 12, 5, ['Line1', 'Line2', 'Line3'], C=C)
        assert len(tb.text_frame.paragraphs) == 3


class TestTopBar:

    def test_creates_bar(self, slide, C):
        from ppt_pro_max.build_helpers import top_bar
        sh = top_bar(slide, 'accent', C=C)
        assert sh is not None
        assert len(slide.shapes) == 1

    def test_custom_width(self, slide, C):
        from ppt_pro_max.build_helpers import top_bar
        sh = top_bar(slide, 'accent', width=10.0, C=C)
        assert sh is not None


class TestPageHeader:

    def test_creates_header(self, slide, C):
        from ppt_pro_max.build_helpers import page_header
        page_header(slide, 'Test Title', 'Test Subtitle', C=C)
        assert len(slide.shapes) >= 2

    def test_no_subtitle(self, slide, C):
        from ppt_pro_max.build_helpers import page_header
        page_header(slide, 'Test Title', C=C)
        assert len(slide.shapes) >= 2


class TestKpiCard:

    def test_creates_card(self, slide, C):
        from ppt_pro_max.build_helpers import kpi_card
        kpi_card(slide, 1.0, 1.5, 3.0, 1.35, '12.8亿', '年度产值', '+8.3%', True, C=C)
        assert len(slide.shapes) >= 3

    def test_no_trend(self, slide, C):
        from ppt_pro_max.build_helpers import kpi_card
        kpi_card(slide, 1.0, 1.5, 3.0, 1.35, '12.8亿', '年度产值', C=C)
        assert len(slide.shapes) >= 3


class TestBarChart:

    def test_creates_bars(self, slide, C):
        from ppt_pro_max.build_helpers import bar_chart
        data = [('A', 0.9, '90'), ('B', 0.5, '50')]
        bar_chart(slide, 1.0, 3.5, data, max_width=5.0, C=C)
        assert len(slide.shapes) > 0


class TestComparisonBars:

    def test_creates_comparison(self, slide, C):
        from ppt_pro_max.build_helpers import comparison_bars
        metrics = [('Revenue', '100', '120', 0.5, 0.6)]
        result = comparison_bars(slide, 1.0, 3.5, metrics, max_width=4.0, C=C)
        assert result > 3.5
        assert len(slide.shapes) > 0


class TestDonutChart:

    def test_creates_donut(self, slide, C):
        from ppt_pro_max.build_helpers import donut_chart
        sectors = [('A', '40%', '#2E6504'), ('B', '60%', '#7DA92F')]
        donut_chart(slide, 3.0, 4.0, 1.3, 0.7, sectors, C=C)
        assert len(slide.shapes) > 0


class TestHighlightCards:

    def test_creates_cards(self, slide, C):
        from ppt_pro_max.build_helpers import highlight_cards
        cards = [('Title1', 'Desc1', '#2E6504'), ('Title2', 'Desc2', '#7DA92F')]
        highlight_cards(slide, 0.65, 5.0, cards, total_width=12.0, C=C)
        assert len(slide.shapes) > 0


class TestCopyDecorations:

    def test_copies_non_text_shapes(self, prs, C):
        from ppt_pro_max.build_helpers import add_slide, rect, copy_decorations
        src = add_slide(prs)
        rect(src, 0, 0, 13.333, 0.08, '#7DA92F', C=C)
        rect(src, 1, 1, 3, 2, '#2E6504', C=C)

        dst = add_slide(prs)
        copy_decorations(dst, src, skip_long_text=True, skip_image=True)
        assert len(dst.shapes) == 2

    def test_skips_long_text(self, prs, C):
        from ppt_pro_max.build_helpers import add_slide, text, rect, copy_decorations
        src = add_slide(prs)
        rect(src, 0, 0, 1, 1, '#2E6504', C=C)
        text(src, 1, 1, 5, 1, 'A' * 100, C=C)

        dst = add_slide(prs)
        copy_decorations(dst, src, skip_long_text=True, skip_image=True)
        assert len(dst.shapes) == 1

    def test_keeps_short_text(self, prs, C):
        from ppt_pro_max.build_helpers import add_slide, text, rect, copy_decorations
        src = add_slide(prs)
        rect(src, 0, 0, 1, 1, '#2E6504', C=C)
        text(src, 1, 1, 2, 0.5, 'PART ONE', C=C)

        dst = add_slide(prs)
        copy_decorations(dst, src, skip_long_text=True, skip_image=True)
        assert len(dst.shapes) == 2


class TestCopyLogo:

    def test_copies_first_group_no_hints(self, prs, C):
        from ppt_pro_max.build_helpers import add_slide, copy_logo
        src = add_slide(prs)
        grp = src.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.53), Inches(2.28), Inches(0.34))
        grp.fill.solid()
        grp.fill.fore_color.rgb = RGBColor(0x2E, 0x65, 0x04)
        # shape_type != 6 for rect, so copy_logo with no hints and no groups finds nothing
        # Test with shape_type=6 by creating a group
        from pptx.oxml.ns import qn
        from lxml import etree
        grpSp = etree.SubElement(src.shapes._spTree, qn('p:grpSp'))
        grpSpPr = etree.SubElement(grpSp, qn('p:grpSpPr'))

        dst = add_slide(prs)
        copy_logo(dst, src)
        assert len(dst.shapes) >= 1
