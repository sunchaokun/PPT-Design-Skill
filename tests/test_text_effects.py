"""Tests for text_effects.py — text-level visual effects on a:rPr.

Covers:
  - apply_text_gradient: gradient fill on text runs
  - apply_text_gradient_preset: preset gradient shortcuts
  - apply_text_outline: text outline/stroke
  - apply_text_shadow: text-level shadow
  - apply_text_glow: text-level glow
  - apply_text_alpha: text transparency
  - apply_text_3d: 3D text via a:sp3d
  - set_vertical_text: vertical text direction
  - set_text_rotation: text/shape rotation
  - apply_letter_spacing: spc as attribute (fix from child element)
"""

from __future__ import annotations

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree


def _make_run(text="Test", font_size=28):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    tb = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(5), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    return run


class TestTextGradient:
    def test_two_stop_gradient(self):
        from ppt_pro_max.renderer.text_effects import apply_text_gradient
        run = _make_run("Gradient")
        apply_text_gradient(run, [("#F5AF19", 0), ("#FFA300", 100000)], angle=5400000)
        rPr = run._r.find(qn("a:rPr"))
        grad = rPr.find(qn("a:gradFill"))
        assert grad is not None
        gsLst = grad.find(qn("a:gsLst"))
        assert gsLst is not None
        stops = gsLst.findall(qn("a:gs"))
        assert len(stops) == 2
        assert stops[0].get("pos") == "0"
        assert stops[1].get("pos") == "100000"

    def test_three_stop_gradient(self):
        from ppt_pro_max.renderer.text_effects import apply_text_gradient
        run = _make_run("Triple")
        apply_text_gradient(run, [("#8B5CF6", 0), ("#6366F1", 50000), ("#3B82F6", 100000)])
        rPr = run._r.find(qn("a:rPr"))
        stops = rPr.findall(f".//{qn('a:gs')}")
        assert len(stops) == 3

    def test_gradient_replaces_existing_fill(self):
        from ppt_pro_max.renderer.text_effects import apply_text_gradient
        run = _make_run("Replace")
        rPr = run._r.get_or_add_rPr()
        solid = etree.SubElement(rPr, qn("a:solidFill"))
        etree.SubElement(solid, qn("a:srgbClr")).set("val", "FF0000")
        apply_text_gradient(run, [("#F5AF19", 0), ("#FFA300", 100000)])
        assert rPr.find(qn("a:solidFill")) is None
        assert rPr.find(qn("a:gradFill")) is not None


class TestTextGradientPreset:
    def test_gold_shine_preset(self):
        from ppt_pro_max.renderer.text_effects import apply_text_gradient_preset
        run = _make_run("Gold")
        apply_text_gradient_preset(run, "gold-shine")
        rPr = run._r.find(qn("a:rPr"))
        grad = rPr.find(qn("a:gradFill"))
        assert grad is not None
        stops = grad.findall(f".//{qn('a:gs')}")
        assert len(stops) == 3

    def test_ink_wash_preset(self):
        from ppt_pro_max.renderer.text_effects import apply_text_gradient_preset
        run = _make_run("Ink")
        apply_text_gradient_preset(run, "ink-wash")
        rPr = run._r.find(qn("a:rPr"))
        stops = rPr.findall(f".//{qn('a:gs')}")
        assert len(stops) >= 2

    def test_unknown_preset_raises(self):
        from ppt_pro_max.renderer.text_effects import apply_text_gradient_preset
        run = _make_run("Bad")
        with pytest.raises((KeyError, ValueError)):
            apply_text_gradient_preset(run, "nonexistent-preset")


class TestTextOutline:
    def test_outline_added(self):
        from ppt_pro_max.renderer.text_effects import apply_text_outline
        run = _make_run("Outline")
        apply_text_outline(run, "#6366F1", 1.5)
        rPr = run._r.find(qn("a:rPr"))
        ln = rPr.find(qn("a:ln"))
        assert ln is not None
        assert int(ln.get("w")) == int(1.5 * 12700)

    def test_outline_color(self):
        from ppt_pro_max.renderer.text_effects import apply_text_outline
        run = _make_run("Color")
        apply_text_outline(run, "#FF0000", 2.0)
        rPr = run._r.find(qn("a:rPr"))
        ln = rPr.find(qn("a:ln"))
        srgb = ln.find(f".//{qn('a:srgbClr')}")
        assert srgb is not None
        assert srgb.get("val") == "FF0000"


class TestTextShadow:
    def test_shadow_added(self):
        from ppt_pro_max.renderer.text_effects import apply_text_shadow
        run = _make_run("Shadow")
        apply_text_shadow(run, blur=6.0, dist=3.0, direction=90, color="#000000", alpha=25)
        rPr = run._r.find(qn("a:rPr"))
        effectLst = rPr.find(qn("a:effectLst"))
        assert effectLst is not None
        shdw = effectLst.find(qn("a:outerShdw"))
        assert shdw is not None

    def test_shadow_parameters(self):
        from ppt_pro_max.renderer.text_effects import apply_text_shadow
        run = _make_run("Params")
        apply_text_shadow(run, blur=8.0, dist=4.0, direction=45, color="#333333", alpha=30)
        rPr = run._r.find(qn("a:rPr"))
        shdw = rPr.find(f".//{qn('a:outerShdw')}")
        assert shdw is not None
        assert int(shdw.get("blurRad")) == int(8.0 * 12700)
        assert int(shdw.get("dist")) == int(4.0 * 12700)


class TestTextGlow:
    def test_glow_added(self):
        from ppt_pro_max.renderer.text_effects import apply_text_glow
        run = _make_run("Glow")
        apply_text_glow(run, radius=8.0, color="#8B5CF6", alpha=40)
        rPr = run._r.find(qn("a:rPr"))
        effectLst = rPr.find(qn("a:effectLst"))
        assert effectLst is not None
        glow = effectLst.find(qn("a:glow"))
        assert glow is not None

    def test_glow_color(self):
        from ppt_pro_max.renderer.text_effects import apply_text_glow
        run = _make_run("Color")
        apply_text_glow(run, radius=10.0, color="#22D3EE", alpha=50)
        rPr = run._r.find(qn("a:rPr"))
        glow = rPr.find(f".//{qn('a:glow')}")
        srgb = glow.find(qn("a:srgbClr"))
        assert srgb is not None
        assert srgb.get("val") == "22D3EE"


class TestTextAlpha:
    def test_alpha_set(self):
        from ppt_pro_max.renderer.text_effects import apply_text_alpha
        run = _make_run("Fade")
        apply_text_alpha(run, 50)
        rPr = run._r.find(qn("a:rPr"))
        solidFill = rPr.find(qn("a:solidFill"))
        assert solidFill is not None
        alpha_el = solidFill.find(f".//{qn('a:alpha')}")
        assert alpha_el is not None
        assert alpha_el.get("val") == "50000"

    def test_full_opacity(self):
        from ppt_pro_max.renderer.text_effects import apply_text_alpha
        run = _make_run("Full")
        apply_text_alpha(run, 100)
        rPr = run._r.find(qn("a:rPr"))
        alpha_el = rPr.find(f".//{qn('a:alpha')}")
        assert alpha_el is not None
        assert alpha_el.get("val") == "100000"


class TestText3D:
    def test_3d_added(self):
        from ppt_pro_max.renderer.text_effects import apply_text_3d
        run = _make_run("3D Text")
        apply_text_3d(run, depth_pt=10, bevel=True, material="powder")
        rPr = run._r.find(qn("a:rPr"))
        sp3d = rPr.find(qn("a:sp3d"))
        assert sp3d is not None

    def test_3d_depth(self):
        from ppt_pro_max.renderer.text_effects import apply_text_3d
        run = _make_run("Depth")
        apply_text_3d(run, depth_pt=15, bevel=False, material="metal")
        rPr = run._r.find(qn("a:rPr"))
        sp3d = rPr.find(qn("a:sp3d"))
        assert sp3d is not None
        assert int(sp3d.get("z")) == int(15 * 12700)


class TestVerticalText:
    def test_ea_vertical(self):
        from ppt_pro_max.renderer.text_effects import set_vertical_text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(5))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "古诗句"
        set_vertical_text(tb.text_frame, "ea")
        bodyPr = tb.text_frame._txBody.find(qn("a:bodyPr"))
        assert bodyPr is not None
        assert bodyPr.get("vert") == "eaVert"

    def test_vert270(self):
        from ppt_pro_max.renderer.text_effects import set_vertical_text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(5))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Rotated"
        set_vertical_text(tb.text_frame, "270")
        bodyPr = tb.text_frame._txBody.find(qn("a:bodyPr"))
        assert bodyPr.get("vert") == "vert270"


class TestTextRotation:
    def test_rotation_set(self):
        from ppt_pro_max.renderer.text_effects import set_text_rotation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "印"
        set_text_rotation(tb, -15)
        xfrm = tb._element.find(f".//{qn('a:xfrm')}")
        assert xfrm is not None
        assert int(xfrm.get("rot")) == int(-15 * 60000)


class TestLetterSpacingFix:
    def test_spc_is_attribute(self):
        from ppt_pro_max.renderer.text_effects import apply_letter_spacing
        run = _make_run("Spaced", 14)
        apply_letter_spacing(run, 0.02, 14)
        rPr = run._r.find(qn("a:rPr"))
        assert rPr.get("spc") is not None, "spc should be an attribute, not a child element"
        expected = str(int(0.02 * 14 * 100))
        assert rPr.get("spc") == expected

    def test_negative_tracking(self):
        from ppt_pro_max.renderer.text_effects import apply_letter_spacing
        run = _make_run("Tight", 28)
        apply_letter_spacing(run, -0.01, 28)
        rPr = run._r.find(qn("a:rPr"))
        assert rPr.get("spc") is not None

    def test_zero_no_spc(self):
        from ppt_pro_max.renderer.text_effects import apply_letter_spacing
        run = _make_run("Normal", 14)
        apply_letter_spacing(run, 0.0, 14)
        rPr = run._r.find(qn("a:rPr"))
        assert rPr.get("spc") is None

    def test_no_spc_child_element(self):
        from ppt_pro_max.renderer.text_effects import apply_letter_spacing
        run = _make_run("NoChild", 14)
        apply_letter_spacing(run, 0.05, 14)
        rPr = run._r.find(qn("a:rPr"))
        spc_child = rPr.find(qn("a:spc"))
        assert spc_child is None, "a:spc should NOT be a child element"
