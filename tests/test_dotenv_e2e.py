"""End-to-end: build.py → build_helpers.ai_image() → ImageFetcher uses .env key.

Simulates a real build.py running from a user project directory where the
skill .env lives at ~/.agents/skills/ppt-design-skill/.env.
"""
import os
from pathlib import Path
from unittest.mock import patch

import pytest


@pytest.fixture
def clean_env(monkeypatch):
    for k in (
        "ARK_API_KEY", "OPENAI_API_KEY", "GEMINI_API_KEY",
        "DASHSCOPE_API_KEY", "MOONSHOT_API_KEY",
        "PPT_IMAGE_LLM_API_KEY", "PPT_IMAGE_LLM_PROVIDER",
        "PPT_IMAGE_LLM_BASE_URL", "PPT_IMAGE_LLM_MODEL",
        "UNSPLASH_ACCESS_KEY", "PEXELS_API_KEY",
    ):
        monkeypatch.delenv(k, raising=False)


@pytest.fixture
def skill_env(clean_env, tmp_path, monkeypatch):
    """Create a fake skill .env and isolate the loader to only find it.

    `Path.is_file` is patched so the ONLY .env the loader discovers is the
    fake skill .env — this prevents interference from the real project .env
    at the repo root (which pytest would otherwise pick up via cwd).

    IMPORTANT: build_helpers' module-level `_load_dotenv()` runs at import
    time, so we must re-clear os.environ AFTER importing but BEFORE calling
    `_load_dotenv()` manually — otherwise the real project .env key (loaded
    during import) wins via override=False.
    """
    skill_dir = tmp_path / ".agents" / "skills" / "ppt-design-skill"
    skill_dir.mkdir(parents=True)
    skill_env_file = skill_dir / ".env"
    skill_env_file.write_text(
        "ARK_API_KEY=skill-ark-key\n"
        "ARK_BASE_URL=https://skill.base.url\n"
        "ARK_IMAGE_MODEL=skill-seedream-model\n"
        "MOONSHOT_API_KEY=skill-moonshot-key\n",
        encoding="utf-8",
    )
    project_dir = tmp_path / "user_project"
    project_dir.mkdir()
    monkeypatch.chdir(project_dir)

    import ppt_pro_max.build_helpers as bh
    bh._LOADED = False
    for k in ("ARK_API_KEY", "ARK_BASE_URL", "ARK_IMAGE_MODEL", "MOONSHOT_API_KEY"):
        os.environ.pop(k, None)
    def _fake_is_file(path, env_path=skill_env_file):
        return str(path).endswith(".env") and str(path) == str(env_path)
    with patch.object(Path, "home", return_value=tmp_path), patch.object(Path, "is_file", new=_fake_is_file):
        bh._load_dotenv()
    return skill_dir


class TestEndToEndSkillDotenv:
    def test_skill_key_loaded_into_environ(self, skill_env):
        assert os.environ.get("ARK_API_KEY") == "skill-ark-key"
        assert os.environ.get("ARK_BASE_URL") == "https://skill.base.url"
        assert os.environ.get("ARK_IMAGE_MODEL") == "skill-seedream-model"
        assert os.environ.get("MOONSHOT_API_KEY") == "skill-moonshot-key"

    def test_ImageFetcher_resolves_skill_key(self, skill_env):
        from ppt_pro_max.renderer.image_fetcher import ImageFetcher
        with patch("ppt_pro_max.adapters.llm_config_adapter.detect_host_llm_config", return_value={}):
            fetcher = ImageFetcher(mode="generate", llm_provider="seedream", auto_detect=True)
        assert fetcher.llm_provider == "seedream"
        assert fetcher.llm_api_key == "skill-ark-key"
        assert fetcher.llm_base_url == "https://skill.base.url"
        assert fetcher.llm_model == "skill-seedream-model"

    def test_fetch_image_uses_skill_key(self, skill_env):
        from ppt_pro_max import fetch_image
        captured = {}
        orig_init = None
        import ppt_pro_max.renderer.image_fetcher as imgf
        def fake_init(self, *args, **kwargs):
            captured["self"] = self
            orig_init(self, *args, **kwargs)
        orig_init = imgf.ImageFetcher.__init__
        imgf.ImageFetcher.__init__ = fake_init
        try:
            with patch("ppt_pro_max.adapters.llm_config_adapter.detect_host_llm_config", return_value={}), patch("ppt_pro_max.renderer.image_fetcher.ImageFetcher.fetch", return_value=None):
                fetch_image("test image", mode="generate", llm_provider="seedream")
        finally:
            imgf.ImageFetcher.__init__ = orig_init
        assert captured["self"].llm_provider == "seedream"
        assert captured["self"].llm_api_key == "skill-ark-key"

    def test_ai_image_forwards_dotenv_key(self, skill_env):
        """build_helpers.ai_image → fetch_image → ImageFetcher chain picks up skill key."""
        from pptx import Presentation
        from pptx.util import Inches

        from ppt_pro_max.build_helpers import ai_image
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        captured = {}
        def fake_fetch(*args, **kwargs):
            captured["kwargs"] = kwargs
            return {"path": None}
        with patch("ppt_pro_max.build_helpers.fetch_image", side_effect=fake_fetch):
            ai_image(slide, 0.5, 0.5, 3.0, 2.0, "test prompt", mode="generate")
        assert captured["kwargs"]["llm_provider"] is None  # auto-detected inside ImageFetcher
        assert captured["kwargs"]["llm_api_key"] is None  # not passed explicitly
        assert captured["kwargs"]["auto_detect"] is True  # relies on .env / host detection

    def test_explicit_key_overrides_dotenv(self, skill_env):
        from ppt_pro_max.renderer.image_fetcher import ImageFetcher
        with patch("ppt_pro_max.adapters.llm_config_adapter.detect_host_llm_config", return_value={}):
            fetcher = ImageFetcher(mode="generate", llm_provider="seedream",
                                   llm_api_key="explicit-key", auto_detect=True)
        assert fetcher.llm_api_key == "explicit-key"

    def test_cwd_env_takes_precedence_over_skill_env(self, clean_env, tmp_path, monkeypatch):
        skill_dir = tmp_path / ".agents" / "skills" / "ppt-design-skill"
        skill_dir.mkdir(parents=True)
        skill_env_file = skill_dir / ".env"
        skill_env_file.write_text("ARK_API_KEY=skill-key\n", encoding="utf-8")
        project_dir = tmp_path / "project"
        project_dir.mkdir()
        cwd_env_file = project_dir / ".env"
        cwd_env_file.write_text("ARK_API_KEY=cwd-key\n", encoding="utf-8")
        monkeypatch.chdir(project_dir)
        import ppt_pro_max.build_helpers as bh
        bh._LOADED = False
        os.environ.pop("ARK_API_KEY", None)
        with patch.object(Path, "home", return_value=tmp_path):
            bh._load_dotenv()
        assert os.environ.get("ARK_API_KEY") == "cwd-key"

    def test_build_helpers_import_alone_triggers_skill_env(self, clean_env, tmp_path, monkeypatch):
        """from ppt_pro_max.build_helpers import *  →  skill .env loaded."""
        skill_dir = tmp_path / ".agents" / "skills" / "ppt-design-skill"
        skill_dir.mkdir(parents=True)
        skill_env_file = skill_dir / ".env"
        skill_env_file.write_text("ARK_API_KEY=import-skill-key\n", encoding="utf-8")
        project_dir = tmp_path / "import_project"
        project_dir.mkdir()
        monkeypatch.chdir(project_dir)
        import importlib

        import ppt_pro_max.build_helpers as bh
        bh._LOADED = False
        os.environ.pop("ARK_API_KEY", None)
        with patch.object(Path, "home", return_value=tmp_path):
            def _fake_is_file(path, env_path=skill_env_file):
                return str(path).endswith(".env") and str(path) == str(env_path)
            with patch.object(Path, "is_file", new=_fake_is_file):
                importlib.reload(bh)
        assert os.environ.get("ARK_API_KEY") == "import-skill-key"
