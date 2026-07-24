"""Deep design audit: font sizes, line spacing, whitespace, visual rhythm."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree

ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def emu_in(emu):
    return round(emu / 914400, 2) if emu else None

def audit_pptx(path, label):
    prs = Presentation(path)
    sw = emu_in(prs.slide_width)
    sh = emu_in(prs.slide_height)
    
    print(f"\n{'='*80}")
    print(f"  {label} | {sw}\" x {sh}\" | {len(prs.slides)} slides")
    print(f"{'='*80}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n  --- Slide {i} ---")
        
        all_shapes_info = []
        for shape in slide.shapes:
            info = {
                'type': str(shape.shape_type),
                'left': emu_in(shape.left),
                'top': emu_in(shape.top),
                'width': emu_in(shape.width),
                'height': emu_in(shape.height),
                'right': emu_in(shape.left + shape.width) if shape.left and shape.width else None,
                'bottom': emu_in(shape.top + shape.height) if shape.top and shape.height else None,
            }
            
            if shape.shape_type == 6:
                info['children'] = len(shape.shapes)
                info['type'] = 'GROUP'
            
            if shape.has_text_frame:
                info['text'] = shape.text_frame.text[:50].replace('\n', '|')
                info['font_sizes'] = []
                info['line_spacing'] = []
                info['font_names'] = []
                
                for para in shape.text_frame.paragraphs:
                    pPr = para._p.find(f'{{{ns_a}}}pPr')
                    def_sz = None
                    def_font = None
                    if pPr is not None:
                        defRPr = pPr.find(f'{{{ns_a}}}defRPr')
                        if defRPr is not None:
                            def_sz = defRPr.get('sz')
                            latin = defRPr.find(f'{{{ns_a}}}latin')
                            if latin is not None:
                                def_font = latin.get('typeface')
                        spc_bef = pPr.find(f'{{{ns_a}}}spcBef')
                        spc_aft = pPr.find(f'{{{ns_a}}}spcAft')
                    
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.find(f'{{{ns_a}}}rPr')
                        sz = None
                        font = None
                        if rPr is not None:
                            sz = rPr.get('sz')
                            latin = rPr.find(f'{{{ns_a}}}latin')
                            if latin is not None:
                                font = latin.get('typeface')
                        
                        pt = int(sz) / 100 if sz else (int(def_sz) / 100 if def_sz else None)
                        fn = font or def_font
                        if pt:
                            info['font_sizes'].append(pt)
                        if fn:
                            info['font_names'].append(fn)
                
                if info['font_sizes']:
                    info['size_range'] = f"{min(info['font_sizes'])}-{max(info['font_sizes'])}pt"
                else:
                    info['size_range'] = '?'
            
            all_shapes_info.append(info)
        
        # Calculate whitespace
        if all_shapes_info:
            min_top = min(s['top'] or 999 for s in all_shapes_info)
            max_bottom = max(s['bottom'] or 0 for s in all_shapes_info)
            min_left = min(s['left'] or 999 for s in all_shapes_info)
            max_right = max(s['right'] or 0 for s in all_shapes_info)
            
            top_margin = min_top
            bottom_margin = sh - max_bottom if max_bottom else None
            left_margin = min_left
            right_margin = sw - max_right if max_right else None
            
            content_height = (max_bottom or 0) - (min_top or 0)
            whitespace_ratio = 1 - (content_height / sh) if sh and content_height else 0
            
            print(f"  Margins: top={top_margin}\" left={left_margin}\" bottom={bottom_margin or '?'}\" right={right_margin or '?'}\"")
            print(f"  Content height: {content_height:.1f}\" / {sh}\" slide = {whitespace_ratio:.0%} whitespace")
        
        # Print shape details
        for s in all_shapes_info:
            if s.get('text'):
                sizes = s.get('font_sizes', [])
                fonts = s.get('font_names', [])
                print(f"  [{s['type'][:5]}] ({s['left']},{s['top']}) {s['width']}x{s['height']}\" | {s['size_range']} | {s['text'][:40]}")
            elif s['type'] == 'GROUP':
                print(f"  [GROUP] ({s['left']},{s['top']}) {s['width']}x{s['height']}\" | {s['children']} children")
            else:
                print(f"  [{s['type'][:5]}] ({s['left']},{s['top']}) {s['width']}x{s['height']}\"")


def audit_typography(path, label):
    prs = Presentation(path)
    all_sizes = []
    all_fonts = {}
    size_by_role = {'hero': [], 'h1': [], 'h2': [], 'body': [], 'caption': [], 'micro': []}
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                pPr = para._p.find(f'{{{ns_a}}}pPr')
                def_sz = def_font = None
                if pPr is not None:
                    defRPr = pPr.find(f'{{{ns_a}}}defRPr')
                    if defRPr is not None:
                        def_sz = defRPr.get('sz')
                        latin = defRPr.find(f'{{{ns_a}}}latin')
                        if latin is not None:
                            def_font = latin.get('typeface')
                
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rPr = run._r.find(f'{{{ns_a}}}rPr')
                    sz = font = None
                    if rPr is not None:
                        sz = rPr.get('sz')
                        latin = rPr.find(f'{{{ns_a}}}latin')
                        if latin is not None:
                            font = latin.get('typeface')
                    
                    pt = int(sz) / 100 if sz else (int(def_sz) / 100 if def_sz else None)
                    fn = font or def_font
                    
                    if pt:
                        all_sizes.append(pt)
                    if fn:
                        all_fonts[fn] = all_fonts.get(fn, 0) + 1
    
    if not all_sizes:
        print(f"\n  {label}: No font data extracted")
        return
    
    sizes_sorted = sorted(set(all_sizes))
    print(f"\n  {label} TYPOGRAPHY AUDIT:")
    print(f"  Font sizes used: {sizes_sorted}")
    print(f"  Size distribution:")
    for sz in sizes_sorted:
        count = all_sizes.count(sz)
        bar = '#' * count
        print(f"    {sz:5.0f}pt: {bar} ({count}x)")
    
    # Check scale ratios
    if len(sizes_sorted) >= 2:
        print(f"  Scale ratios:")
        for i in range(len(sizes_sorted) - 1):
            ratio = sizes_sorted[i+1] / sizes_sorted[i]
            print(f"    {sizes_sorted[i]:.0f} → {sizes_sorted[i+1]:.0f} = {ratio:.2f}x")
    
    # Check for problems
    print(f"  PROBLEMS:")
    if min(all_sizes) < 9:
        print(f"    ⚠ Font too small: {min(all_sizes)}pt (min should be 9pt)")
    if max(all_sizes) > 72:
        print(f"    ⚠ Font too large: {max(all_sizes)}pt (hero only should exceed 60pt)")
    if len(sizes_sorted) < 4:
        print(f"    ⚠ Too few size levels: {len(sizes_sorted)} (need 5-7)")
    if len(sizes_sorted) > 10:
        print(f"    ⚠ Too many size levels: {len(sizes_sorted)} (max 7-8)")
    
    # Check ratio consistency
    ratios = []
    for i in range(len(sizes_sorted) - 1):
        ratios.append(sizes_sorted[i+1] / sizes_sorted[i])
    if ratios:
        avg_ratio = sum(ratios) / len(ratios)
        max_dev = max(abs(r - avg_ratio) for r in ratios)
        if max_dev > 0.3:
            print(f"    ⚠ Inconsistent scale ratios: avg={avg_ratio:.2f}, max_dev={max_dev:.2f}")
    
    print(f"  Fonts: {dict(sorted(all_fonts.items(), key=lambda x: -x[1]))}")


new_dir = os.path.join(tempfile.gettempdir(), 'ppt_build_v2')
for name, label in [('build_mckinsey.pptx', 'McKinsey v2'), ('build_cyberpunk.pptx', 'Cyberpunk v2'), ('build_creative.pptx', 'Creative v2')]:
    path = os.path.join(new_dir, name)
    if os.path.exists(path):
        audit_pptx(path, label)
        audit_typography(path, label)
