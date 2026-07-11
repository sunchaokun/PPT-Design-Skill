"""Tests for PageRevisionEngine.revise() two-pass rebuild."""

from __future__ import annotations

import os
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def _make_template(path: Path, num_slides: int = 3) -> Path:
    prs = Presentation()
    for i in range(num_slides):
        layout = prs.slide_layouts[-1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i + 1}"
    prs.save(str(path))
    return path


def _make_png(path: Path) -> Path:
    img = Image.new("RGB", (100, 50), color="blue")
    img.save(str(path))
    return path


class TestPageRevisionEngineRevise:

    def test_revise_delete_page(self, tmp_path):
        """Delete page 2 from 3-slide template → 2 slides."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 3)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise(["-2"])
        assert result["num_slides"] == 2
        assert os.path.isfile(result["output_path"])

    def test_revise_swap_pages(self, tmp_path):
        """Swap pages 1 and 3 → order becomes 3,2,1."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 3)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise(["1<>3"])
        assert result["num_slides"] == 3
        prs = Presentation(result["output_path"])
        assert prs.slides[0].shapes.title.text == "Slide 3"
        assert prs.slides[2].shapes.title.text == "Slide 1"

    def test_revise_move_page(self, tmp_path):
        """Move page 1 to after page 3 → order 2,3,1."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 3)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise(["1>3"])
        assert result["num_slides"] == 3
        prs = Presentation(result["output_path"])
        assert prs.slides[0].shapes.title.text == "Slide 2"
        assert prs.slides[2].shapes.title.text == "Slide 1"

    def test_revise_no_ops_returns_original(self, tmp_path):
        """No page ops → output identical to template."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 2)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise([])
        assert result["num_slides"] == 2

    def test_revise_delete_all_but_one(self, tmp_path):
        """Delete pages 2,3 from 3-slide → 1 slide."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 3)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise(["-2", "-3"])
        assert result["num_slides"] == 1
        prs = Presentation(result["output_path"])
        assert prs.slides[0].shapes.title.text == "Slide 1"

    def test_revise_output_path_custom(self, tmp_path):
        """Custom output_path should be used."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 2)
        custom_out = str(tmp_path / "custom_output.pptx")
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise([], output_path=custom_out)
        assert result["output_path"] == custom_out
        assert os.path.isfile(custom_out)

    def test_revise_preserves_layout(self, tmp_path):
        """Two-pass rebuild should preserve slide layouts."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 2)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise(["2<>1"])
        prs = Presentation(result["output_path"])
        for slide in prs.slides:
            assert slide.slide_layout is not None

    def test_revise_invalid_page_raises(self, tmp_path):
        """Out-of-range page should raise ValueError."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 2)
        engine = PageRevisionEngine(str(template_path))
        with pytest.raises(ValueError):
            engine.revise(["-5"])

    def test_revise_swap_self_noop(self, tmp_path):
        """Swap page with itself is a no-op."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 3)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise(["2<>2"])
        assert result["num_slides"] == 3
        prs = Presentation(result["output_path"])
        assert prs.slides[1].shapes.title.text == "Slide 2"

    def test_revise_move_self_noop(self, tmp_path):
        """Move page to itself is a no-op."""
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 3)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise(["3>3"])
        assert result["num_slides"] == 3
        prs = Presentation(result["output_path"])
        assert prs.slides[2].shapes.title.text == "Slide 3"
