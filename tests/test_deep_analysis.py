"""Deep analysis: find text_overflow sources and other issues."""

import os
import json
from pptx import Presentation

from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
from ppt_pro_max.enterprise.delivery_gate import DeliveryGate


def test_deep_analysis(tmp_path):
    td = str(tmp_path)
    slides = [
        {"goal": "hook", "title": "AI Platform Launch", "subtitle": "Transforming Enterprise AI"},
        {"goal": "problem", "title": "Current Challenges", "bullets": ["Data silos across departments", "Manual processes waste time", "No unified AI strategy"]},
        {"goal": "solution", "title": "Our Solution", "bullets": ["Unified data pipeline", "Automated ML workflows", "Real-time analytics dashboard"]},
        {"goal": "features", "title": "Key Features", "cards": [
            {"title": "Data Engine", "text": "Process 10TB/day"},
            {"title": "ML Pipeline", "text": "Auto-retrain models"},
            {"title": "Dashboard", "text": "Real-time insights"},
        ]},
        {"goal": "data", "title": "Performance Metrics"},
        {"goal": "code", "title": "Quick Start", "code": {"language": "python", "lines": ["from ai_platform import Pipeline", "pipe = Pipeline(config='prod')", "result = pipe.run(data)"]}},
        {"goal": "exercise", "title": "Hands-On Lab", "exercise": {"duration": "5 min", "steps": ["Clone the repo", "Run the example", "Check the output"]}},
        {"goal": "cta", "title": "Get Started Today", "subtitle": "Contact us for a free trial"},
    ]
    with open(os.path.join(td, "content.json"), "w", encoding="utf-8") as f:
        json.dump({"slides": slides}, f, ensure_ascii=False, indent=2)

    pipeline = EnterprisePipeline()
    result = pipeline.run(query="AI Platform", project_dir=td)
    pptx_path = result.get("output_path", "")
    assert pptx_path and os.path.isfile(pptx_path)

    prs = Presentation(pptx_path)
    gate = DeliveryGate()
    report = gate.check(pptx_path, None, [])

    print("\n=== Detailed Shape Analysis ===")
    for idx, slide in enumerate(prs.slides):
        print(f"\nSlide {idx}: {len(slide.shapes)} shapes")
        for shape in slide.shapes:
            info = f"  {shape.shape_type}: name='{shape.name}'"
            if shape.has_text_frame:
                tf = shape.text_frame
                text_preview = (tf.text or "")[:50].replace("\n", "|")
                w_in = shape.width / 914400
                h_in = shape.height / 914400
                info += f" text='{text_preview}' wrap={tf.word_wrap} size={w_in:.1f}x{h_in:.1f}\""
                if tf.paragraphs and tf.paragraphs[0].font.size:
                    info += f" font={tf.paragraphs[0].font.size / 12700:.0f}pt"
            if hasattr(shape, "image"):
                try:
                    _ = shape.image
                    info += " [IMAGE]"
                except Exception:
                    pass
            print(info)

    print("\n=== Gate Warnings ===")
    for item in report.warnings:
        print(f"  Slide {item.slide_index}: [{item.check_id}] {item.message} {item.detail}")

    print("\n=== Gate Fatals ===")
    for item in report.fatals:
        print(f"  Slide {item.slide_index}: [{item.check_id}] {item.message} {item.detail}")
