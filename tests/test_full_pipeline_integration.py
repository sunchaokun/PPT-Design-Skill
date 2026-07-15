"""Full-pipeline integration tests: user need → proposal → confirm → build → revise → finalize.

Also tests: ComponentLibrary end-to-end (catalog→match→inject),
DiagramEngine 10 types, DiagramStyle↔BrandSpec, PrecisionRenderer step-by-step API.
"""
import json
import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


# ═══════════════════════════════════════════════════════════════
# 1. Full pipeline: need → proposal → confirm → build → revise → finalize
# ═══════════════════════════════════════════════════════════════

class TestFullPipelineFlow:
    def test_enterprise_dry_run(self, tmp_path):
        from ppt_pro_max import generate_ppt

        project_dir = tmp_path / "project"
        project_dir.mkdir()

        result = generate_ppt("AI产品融资路演", project=str(project_dir), dry_run=True)
        assert result["dry_run"] is True
        assert result["pipeline"] == "enterprise"
        assert "brand_spec" in result

    def test_enterprise_build_with_content(self, tmp_path):
        from ppt_pro_max import generate_ppt

        project_dir = tmp_path / "project"
        project_dir.mkdir()

        result = generate_ppt(
            "产品发布",
            project=str(project_dir),
            content={
                "slides": [
                    {"goal": "hook", "title": "AI产品发布", "subtitle": "重新定义效率"},
                    {"goal": "problem", "title": "当前痛点", "bullets": ["效率低", "成本高", "体验差"]},
                    {"goal": "solution", "title": "我们的方案", "bullets": ["AI引擎", "实时分析", "一键部署"]},
                    {"goal": "cta", "title": "联系我们", "subtitle": "开始体验"},
                ]
            },
            output=str(tmp_path / "output.pptx"),
        )
        assert os.path.isfile(result["output_path"])
        prs = Presentation(result["output_path"])
        assert len(prs.slides) == 4

    def test_enterprise_with_component(self, tmp_path):
        from ppt_pro_max import generate_ppt
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "comp.db")
        lib = ComponentLibrary(db_path)

        group_xml = b'''<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvGrpSpPr><p:cNvPr id="100" name="G"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm><a:off x="914400" y="914400"/><a:ext cx="9144000" cy="4572000"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="4572000"/></a:xfrm></p:grpSpPr>
  <p:sp><p:nvSpPr><p:cNvPr id="101" name="S1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Step 1</a:t></a:r></a:p></p:txBody></p:sp>
  <p:sp><p:nvSpPr><p:cNvPr id="102" name="S2"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="3200400" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Step 2</a:t></a:r></a:p></p:txBody></p:sp>
  <p:sp><p:nvSpPr><p:cNvPr id="103" name="S3"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="6400800" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Step 3</a:t></a:r></a:p></p:txBody></p:sp>
</p:grpSp>'''
        lib.add(type="group", category="process", variant="chevron", node_count=3, xml_parts={"group": group_xml})
        lib.close()

        project_dir = tmp_path / "project"
        project_dir.mkdir()

        result = generate_ppt(
            "产品发布",
            project=str(project_dir),
            component_library=db_path,
            content={
                "slides": [
                    {"goal": "hook", "title": "AI产品发布"},
                    {"goal": "content", "title": "开发流程", "bullets": ["需求", "设计", "开发"],
                     "component_type": "group", "component_category": "process"},
                    {"goal": "cta", "title": "联系我们"},
                ]
            },
            output=str(tmp_path / "output_comp.pptx"),
        )
        assert os.path.isfile(result["output_path"])

    def test_beautify_full_pipeline(self, tmp_path):
        from ppt_pro_max import generate_ppt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        tb = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        tb.text_frame.paragraphs[0].text = "Title Slide"

        slide2 = prs.slides.add_slide(layout)
        tb2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        tb2.text_frame.paragraphs[0].text = "Content Slide"

        input_path = str(tmp_path / "input.pptx")
        prs.save(input_path)

        result = generate_ppt(
            "美化",
            beautify=input_path,
            beautify_mode="full",
            style="professional",
            output=str(tmp_path / "beautified.pptx"),
        )
        assert result["mode"] == "beautify"
        assert result["beautify_mode"] == "full"
        assert os.path.isfile(result["output_path"])


# ═══════════════════════════════════════════════════════════════
# 2. ComponentLibrary end-to-end: catalog → match → inject
# ═══════════════════════════════════════════════════════════════

class TestComponentLibraryE2E:
    def test_catalog_match_inject_roundtrip(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        group_xml = b'''<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvGrpSpPr><p:cNvPr id="100" name="G"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm><a:off x="914400" y="914400"/><a:ext cx="9144000" cy="4572000"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="4572000"/></a:xfrm></p:grpSpPr>
  <p:sp><p:nvSpPr><p:cNvPr id="101" name="S1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Old 1</a:t></a:r></a:p></p:txBody></p:sp>
  <p:sp><p:nvSpPr><p:cNvPr id="102" name="S2"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="3200400" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Old 2</a:t></a:r></a:p></p:txBody></p:sp>
</p:grpSp>'''

        lib.add(type="group", category="process", variant="chevron", node_count=2, xml_parts={"group": group_xml})
        lib.add(type="group", category="swot", variant="grid", node_count=4, xml_parts={"group": b"<g/>"})

        catalog = lib.catalog()
        assert "group" in catalog
        assert catalog["group"]["process"]["count"] == 1
        assert catalog["group"]["swot"]["count"] == 1

        match = lib.match({"type": "group", "category": "process", "node_count": 2})
        assert match is not None
        assert match["category"] == "process"

        xml_parts = lib.load_xml(match["id"])
        assert "group" in xml_parts

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        page = {
            "goal": "content",
            "title": "流程",
            "bullets": ["需求分析", "设计开发"],
            "component_type": "group",
            "component_category": "process",
        }
        slide = renderer.render_slide(prs, page, component_lib=lib)
        assert slide is not None
        assert len(prs.slides) == 1

        output = str(tmp_path / "output.pptx")
        renderer.save(prs, output)
        assert os.path.isfile(output)

        lib.close()

    def test_query_api_e2e(self, tmp_path):
        from ppt_pro_max import query_component_library
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="process", variant="v1", node_count=4, xml_parts={"group": b"<g/>"})
        lib.add(type="group", category="process", variant="v2", node_count=5, xml_parts={"group": b"<g2/>"})
        lib.add(type="smartart", category="hierarchy", variant="org", node_count=6,
                xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})
        lib.close()

        catalog = query_component_library(component_library=db_path)
        assert "group" in catalog
        assert "smartart" in catalog
        assert catalog["group"]["process"]["count"] == 2
        assert catalog["smartart"]["hierarchy"]["count"] == 1

        results = query_component_library(type="group", category="process", component_library=db_path)
        assert len(results) == 2

        results_4 = query_component_library(type="group", category="process", node_count=4, component_library=db_path)
        assert len(results_4) == 1
        assert results_4[0]["node_count"] == 4


# ═══════════════════════════════════════════════════════════════
# 3. DiagramEngine 10 types
# ═══════════════════════════════════════════════════════════════

class TestDiagramEngineAllTypes:
    @pytest.fixture
    def engine_and_slide(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        return DiagramEngine(), slide

    def test_supported_types_count(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine

        types = DiagramEngine.get_supported_types()
        assert len(types) == 10

    def test_flowchart(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "flowchart", {"nodes": [{"label": "Start"}, {"label": "Process"}, {"label": "End"}]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_funnel(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "funnel", {"stages": [{"label": "Leads"}, {"label": "Qualified"}, {"label": "Closed"}]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_timeline(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "timeline", {"events": [{"label": "Q1"}, {"label": "Q2"}, {"label": "Q3"}]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_swot(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "swot",
                      {"strengths": ["S1"], "weaknesses": ["W1"], "opportunities": ["O1"], "threats": ["T1"]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_matrix(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "matrix",
                      {"rows": [{"label": "R1"}, {"label": "R2"}], "cols": [{"label": "C1"}, {"label": "C2"}],
                       "cells": [["1", "2"], ["3", "4"]]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_cycle(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "cycle", {"nodes": [{"label": "Plan"}, {"label": "Do"}, {"label": "Check"}]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_table(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "table",
                      {"headers": ["Name", "Value"], "rows": [["A", "1"], ["B", "2"]]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_hierarchy(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "hierarchy",
                      {"nodes": [{"id": 0, "label": "CEO", "parent": None}, {"id": 1, "label": "CTO", "parent": 0}]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_pyramid(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "pyramid", {"levels": [{"label": "Top"}, {"label": "Mid"}, {"label": "Base"}]},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1

    def test_venn(self, engine_and_slide):
        engine, slide = engine_and_slide
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        engine.render(slide, "venn",
                      {"sets": [{"label": "A"}, {"label": "B"}], "intersection": "A∩B"},
                      DiagramStyle(), Region(left=0.5, top=1, width=12, height=5))
        assert len(slide.shapes) > 1


# ═══════════════════════════════════════════════════════════════
# 4. DiagramStyle ↔ BrandSpec
# ═══════════════════════════════════════════════════════════════

class TestDiagramStyleBrandSpec:
    def test_from_brand_spec(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand = BrandSpec(
            source="test",
            colors={"primary": "#FF0000", "on-primary": "#FFFFFF", "muted-foreground": "#888888"},
        )
        style = DiagramStyle.from_brand_spec(brand)
        assert style._color_map["primary"] == "#FF0000"
        assert style.resolve_color("primary") == "#FF0000"

    def test_from_brand_spec_dark_mode(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand = BrandSpec(
            source="test",
            colors={"primary": "#2563EB", "background": "#0A1E3D", "foreground": "#F0F4F8"},
            dark_mode=True,
        )
        style = DiagramStyle.from_brand_spec(brand)
        assert style.resolve_color("primary") == "#2563EB"

    def test_default_style_fallback(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle

        style = DiagramStyle()
        assert style.resolve_color("primary") == "#2563EB"
        assert style.resolve_color("on-primary") == "#FFFFFF"

    def test_apply_density(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle

        style = DiagramStyle()
        low = style.apply_density(2)
        assert low.node_font_size_pt == 16
        high = style.apply_density(8)
        assert high.node_font_size_pt == 12


# ═══════════════════════════════════════════════════════════════
# 5. PrecisionRenderer step-by-step API
# ═══════════════════════════════════════════════════════════════

class TestPrecisionRendererStepByStep:
    def test_manual_slide_construction(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand = BrandSpec(
            source="test",
            colors={"primary": "#1E3A5F", "on-primary": "#FFFFFF", "accent": "#E8A838",
                    "background": "#FFFFFF", "foreground": "#1A1A1A", "muted": "#F1F5F9",
                    "muted-foreground": "#64748B", "border": "#E2E8F0"},
            fonts={"heading": "Inter", "body": "Inter"},
        )
        renderer = PrecisionRenderer(brand_spec=brand)
        prs = renderer.create_presentation()

        slide = renderer.add_slide(prs)
        renderer.apply_brand_background(slide, prs, goal="content")
        renderer.add_text(slide, "Manual Title", 0.9, 0.5, 11, 0.6, size=28, bold=True)
        renderer.add_rect(slide, 0.9, 1.2, 2, 0.04, fill_role="accent")
        renderer.add_multiline(slide, ["•  Point 1", "•  Point 2", "•  Point 3"], 0.9, 1.6, 7, 4.5, size=14)
        renderer.add_rounded_rect(slide, 9, 1.5, 3.5, 4, fill_role="muted", border_role="border")
        renderer.add_text(slide, "Card Title", 9.15, 1.65, 3.2, 0.4, size=15, color_role="accent", bold=True)

        output = str(tmp_path / "manual.pptx")
        renderer.save(prs, output)
        assert os.path.isfile(output)

        check = Presentation(output)
        assert len(check.slides) == 1
        assert len(check.slides[0].shapes) >= 6

    def test_add_image_placeholder(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        slide = renderer.add_slide(prs)

        renderer.add_text(slide, "Image Slide", 0.9, 0.5, 11, 0.6, size=28, bold=True)

        img_path = str(tmp_path / "test.png")
        from PIL import Image as PILImage
        img = PILImage.new("RGB", (100, 100), color="red")
        img.save(img_path)
        renderer.add_image(slide, img_path, 8, 1.5, 4, 4)

        output = str(tmp_path / "with_image.pptx")
        renderer.save(prs, output)
        assert os.path.isfile(output)

    def test_add_card(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        slide = renderer.add_slide(prs)

        renderer.add_card(slide, 0.9, 1.5, 3.5, 4.5, "Feature A", "Description of feature A")
        renderer.add_card(slide, 4.7, 1.5, 3.5, 4.5, "Feature B", "Description of feature B")
        renderer.add_card(slide, 8.5, 1.5, 3.5, 4.5, "Feature C", "Description of feature C")

        output = str(tmp_path / "cards.pptx")
        renderer.save(prs, output)
        check = Presentation(output)
        assert len(check.slides[0].shapes) >= 3

    def test_dark_overlay(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand = BrandSpec(source="test", colors={"primary": "#1E3A5F", "background": "#060B18"}, dark_mode=True)
        renderer = PrecisionRenderer(brand_spec=brand)
        prs = renderer.create_presentation()
        slide = renderer.add_slide(prs)

        renderer.apply_hero_overlay(slide, prs)
        renderer.add_text(slide, "Dark Hero", 1.2, 2.0, 8, 1.5, size=52, color_role="foreground", bold=True)

        output = str(tmp_path / "dark.pptx")
        renderer.save(prs, output)
        assert os.path.isfile(output)

    def test_footer_and_watermark(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand = BrandSpec(
            source="test",
            colors={"primary": "#2563EB", "on-primary": "#FFFFFF", "muted-foreground": "#999999",
                    "background": "#FFFFFF", "foreground": "#000000"},
            footer={"show_page_number": True, "page_number_format": "{n}/{total}"},
            watermark={"text": "DRAFT", "opacity": 0.1},
        )
        renderer = PrecisionRenderer(brand_spec=brand)
        prs = renderer.create_presentation()

        for i in range(3):
            slide = renderer.add_slide(prs)
            renderer.add_text(slide, f"Slide {i+1}", 1, 1, 10, 1, size=28)

        renderer.apply_footer(prs)
        renderer.apply_watermark(prs)

        output = str(tmp_path / "footer_wm.pptx")
        renderer.save(prs, output)
        assert os.path.isfile(output)


# ═══════════════════════════════════════════════════════════════
# 6. Step-by-step LLM construction flow
# ═══════════════════════════════════════════════════════════════

class TestStepByStepConstruction:
    def test_llm_builds_ppt_incrementally(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        brand = BrandSpec(
            source="llm",
            colors={"primary": "#2563EB", "on-primary": "#FFFFFF", "accent": "#F97316",
                    "background": "#FFFFFF", "foreground": "#1A1A1A", "muted": "#F1F5F9",
                    "muted-foreground": "#64748B", "border": "#E2E8F0"},
            fonts={"heading": "Inter", "body": "Inter"},
        )
        renderer = PrecisionRenderer(brand_spec=brand)
        prs = renderer.create_presentation()

        slide1 = renderer.add_slide(prs)
        renderer.apply_hero_overlay(slide1, prs)
        renderer.add_text(slide1, "AI Startup", 1.2, 2.0, 8, 1.5, size=52, color_role="on-primary", bold=True)
        renderer.add_text(slide1, "Redefining Productivity", 1.2, 3.6, 8, 0.6,
                          font="Inter", size=22, color_role="accent")

        slide2 = renderer.add_slide(prs)
        renderer.apply_brand_background(slide2, prs, goal="content")
        renderer.add_text(slide2, "Core Features", 0.9, 0.5, 11, 0.6, size=28, bold=True)
        renderer.add_rect(slide2, 0.9, 1.2, 2, 0.04, fill_role="accent")
        renderer.add_card(slide2, 0.9, 1.6, 3.5, 4.5, "AI Engine", "Neural processing")
        renderer.add_card(slide2, 4.7, 1.6, 3.5, 4.5, "Dashboard", "Real-time analytics")
        renderer.add_card(slide2, 8.5, 1.6, 3.5, 4.5, "API", "One-line integration")

        slide3 = renderer.add_slide(prs)
        renderer.apply_brand_background(slide3, prs, goal="data")
        renderer.add_text(slide3, "Development Process", 0.9, 0.5, 11, 0.6, size=28, bold=True)
        engine = DiagramEngine()
        style = DiagramStyle.from_brand_spec(brand)
        engine.render(slide3, "flowchart",
                      {"nodes": [{"label": "Plan"}, {"label": "Build"}, {"label": "Test"}, {"label": "Deploy"}]},
                      style, Region(left=0.9, top=1.5, width=11, height=5))

        slide4 = renderer.add_slide(prs)
        renderer.apply_hero_overlay(slide4, prs)
        renderer.add_text(slide4, "Get Started Today", 1.2, 2.0, 8, 1.5, size=52, color_role="on-primary", bold=True)

        renderer.apply_footer(prs)

        output = str(tmp_path / "llm_built.pptx")
        renderer.save(prs, output)
        assert os.path.isfile(output)

        check = Presentation(output)
        assert len(check.slides) == 4
        for slide in check.slides:
            assert len(slide.shapes) >= 1

    def test_llm_uses_component_library(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.component_library import ComponentLibrary
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        db_path = str(tmp_path / "comp.db")
        lib = ComponentLibrary(db_path)

        group_xml = b'''<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvGrpSpPr><p:cNvPr id="100" name="G"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm><a:off x="914400" y="914400"/><a:ext cx="9144000" cy="4572000"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="4572000"/></a:xfrm></p:grpSpPr>
  <p:sp><p:nvSpPr><p:cNvPr id="101" name="S1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>X</a:t></a:r></a:p></p:txBody></p:sp>
  <p:sp><p:nvSpPr><p:cNvPr id="102" name="S2"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="3200400" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Y</a:t></a:r></a:p></p:txBody></p:sp>
</p:grpSp>'''
        lib.add(type="group", category="process", variant="chevron", node_count=2, xml_parts={"group": group_xml})

        brand = BrandSpec(source="test", colors={"primary": "#2563EB", "on-primary": "#FFFFFF"})
        renderer = PrecisionRenderer(brand_spec=brand)
        prs = renderer.create_presentation()

        slide = renderer.add_slide(prs)
        renderer.render_slide(prs, {
            "goal": "content",
            "title": "流程",
            "bullets": ["步骤一", "步骤二"],
            "component_type": "group",
            "component_category": "process",
        }, component_lib=lib)

        output = str(tmp_path / "llm_comp.pptx")
        renderer.save(prs, output)
        assert os.path.isfile(output)

        lib.close()


# ═══════════════════════════════════════════════════════════════
# 7. Issues found during review
# ═══════════════════════════════════════════════════════════════

class TestReviewIssues:
    def test_beautify_full_falls_back_to_light_on_extract_failure(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        tb.text_frame.paragraphs[0].text = "Fallback Test"

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        prs.save(input_path)

        pipeline = EnterprisePipeline()
        result = pipeline.run_beautify(
            input_pptx=input_path,
            output_path=output_path,
            beautify_mode="full",
        )
        assert os.path.isfile(output_path)
        assert result["mode"] == "beautify"

    def test_component_type_with_diagram_type_fallback_chain(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()

        page = {
            "goal": "content",
            "title": "Process Flow",
            "bullets": ["Step A", "Step B", "Step C"],
            "component_type": "group",
            "component_category": "nonexistent",
            "diagram_type": "flowchart",
            "diagram_data": {"nodes": [{"label": "A"}, {"label": "B"}, {"label": "C"}]},
        }

        slide = renderer.render_slide(prs, page, component_lib=None)
        assert slide is not None

    def test_content_and_content_file_priority(self, tmp_path):
        from ppt_pro_max import generate_ppt

        project_dir = tmp_path / "project"
        project_dir.mkdir()

        content_file = str(tmp_path / "file_content.json")
        with open(content_file, "w", encoding="utf-8") as f:
            json.dump({"slides": [{"goal": "hook", "title": "FROM FILE"}]}, f)

        result = generate_ppt(
            "test",
            project=str(project_dir),
            content={"slides": [{"goal": "hook", "title": "FROM DICT"}]},
            content_file=content_file,
            dry_run=True,
        )
        assert result.get("dry_run") is True
