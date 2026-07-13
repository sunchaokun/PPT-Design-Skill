"""Tests for cards rendering and new high-density layouts."""

from __future__ import annotations

import os
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def _make_png(path: Path) -> Path:
    img = Image.new("RGB", (100, 50), color="red")
    img.save(str(path))
    return path


class TestCardsRendering:

    def test_cards_rendered(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        cards = [
            {"title": "Card 1", "body": "Body 1"},
            {"title": "Card 2", "body": "Body 2"},
            {"title": "Card 3", "body": "Body 3"},
        ]
        design = {"goal": "features", "title": "Cards Test", "cards": cards}
        precision.render_slide(prs, design)
        slide = prs.slides[-1]
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Card 1" in t for t in texts)
        assert any("Card 2" in t for t in texts)
        assert any("Card 3" in t for t in texts)

    def test_cards_with_images(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        img = _make_png(tmp_path / "card_img.png")
        cards = [
            {"title": "Card 1", "body": "Body", "image": str(img)},
        ]
        design = {"goal": "features", "title": "Cards With Image", "cards": cards, "image": str(img)}
        precision.render_slide(prs, design)
        slide = prs.slides[-1]
        pics = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pics) >= 1

    def test_empty_cards_no_crash(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        design = {"goal": "features", "title": "Empty Cards", "cards": []}
        precision.render_slide(prs, design)

    def test_card_with_nonexistent_image_no_crash(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        cards = [{"title": "Card", "image": "/nonexistent.png"}]
        design = {"goal": "features", "title": "Bad Image", "cards": cards, "image": "/nonexistent.png"}
        precision.render_slide(prs, design)


class TestNewLayouts:

    def test_grid_2x2_exists(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout("grid-2x2-cards")
        assert layout is not None
        assert "card1" in layout["placeholders"]
        assert "card4" in layout["placeholders"]

    def test_dense_bullets_exists(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout("dense-bullets")
        assert layout is not None
        assert "body" in layout["placeholders"]

    def test_exercise_layout_exists(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout("exercise-layout")
        assert layout is not None
        assert "instructions" in layout["placeholders"]

    def test_code_block_exists(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout("code-block")
        assert layout is not None
        assert "code" in layout["placeholders"]

    def test_timeline_exists(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout("timeline-horizontal")
        assert layout is not None

    def test_swot_matrix_exists(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout("swot-matrix")
        assert layout is not None
        assert "card4" in layout["placeholders"]

    def test_funnel_exists(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout("funnel")
        assert layout is not None

    def test_cycle_exists(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout("cycle-diagram")
        assert layout is not None

    def test_all_new_layouts_have_unique_ids(self):
        from ppt_pro_max.renderer.layout_registry import MASTER_LAYOUTS
        ids = [v["id"] for v in MASTER_LAYOUTS.values()]
        assert len(ids) == len(set(ids))

    def test_new_layouts_within_bounds(self):
        from ppt_pro_max.renderer.layout_registry import MASTER_LAYOUTS, SLIDE_WIDTH, SLIDE_HEIGHT
        new_names = [
            "grid-2x2-cards", "dense-bullets", "two-column-dense",
            "table-layout", "sidebar-left", "exercise-layout",
            "code-block", "timeline-horizontal", "swot-matrix",
            "funnel", "cycle-diagram",
        ]
        for name in new_names:
            layout = MASTER_LAYOUTS[name]
            for ph_name, ph in layout.get("placeholders", {}).items():
                if ph.get("type") in ("image", "chart"):
                    continue
                assert ph["x"] >= 0, f"{name}/{ph_name}: x<0"
                assert ph["y"] >= 0, f"{name}/{ph_name}: y<0"
                assert ph["x"] + ph["width"] <= SLIDE_WIDTH + 0.1, f"{name}/{ph_name}: overflow x"
                assert ph["y"] + ph["height"] <= SLIDE_HEIGHT + 0.1, f"{name}/{ph_name}: overflow y"

    def test_goal_mapping_covers_education(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout_by_goal("education")
        assert layout is not None

    def test_goal_mapping_covers_exercise(self):
        from ppt_pro_max.renderer.layout_registry import LayoutRegistry
        reg = LayoutRegistry()
        layout = reg.get_layout_by_goal("exercise")
        assert layout is not None
