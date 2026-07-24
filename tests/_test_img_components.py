"""Find components with images and test injection."""

from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
from ppt_pro_max.enterprise.brand_spec import BrandSpec
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from lxml import etree
import json
import os

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)

lib = ComponentLibrary(find_db_path())
renderer = ComponentRenderer()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GREEN_BRAND = BrandSpec(colors={
    "primary": "#00AA00", "accent": "#66CC00", "muted": "#E6F5E6",
    "foreground": "#1A3A1A", "on-primary": "#FFFFFF", "background": "#FFFFFF",
})

slide_idx = 0

# Find components with images
for cat in ["hierarchy", "infographic", "process", "swot", "timeline", "chart"]:
    results = lib.search(type="group", category=cat, limit=20)
    for comp in results:
        xml_parts = lib.load_xml(comp["id"])
        img_keys = [k for k in xml_parts if k.startswith("img_")]
        if not img_keys:
            continue

        root = etree.fromstring(xml_parts["group"])
        pic_count = len(list(root.iter(f"{{{_P}}}pic")))
        has_3d = len(list(root.iter(f"{{{_A}}}scene3d"))) > 0
        texts = [t.text.strip()[:15] for t in root.iter(f"{{{_A}}}t") if t.text and t.text.strip()]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"{cat} id={comp['id']} (img={len(img_keys)}, pic={pic_count}, 3d={has_3d})"
        from pptx.util import Pt
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 0xAA, 0)

        meta_bytes = xml_parts.get("_meta")
        orig_aspect = 1.0
        if meta_bytes:
            meta = json.loads(meta_bytes)
            obe = meta.get("orig_bounds_emu")
            if obe and len(obe) == 4 and obe[3] > 0:
                orig_aspect = obe[2] / obe[3]

        content_area = (0.5, 0.9, 12.3, 5.8)
        bounds = renderer.compute_component_bounds(content_area, orig_aspect, "contain")
        renderer._inject_group_to_slide(slide, xml_parts, bounds)

        print(f"  Slide {slide_idx+1}: {cat} id={comp['id']}, img={len(img_keys)}, pic={pic_count}, 3d={has_3d}, texts={texts[:3]}")
        slide_idx += 1
        if slide_idx >= 10:
            break
    if slide_idx >= 10:
        break

if slide_idx > 0:
    out = os.path.join(OUT_DIR, "components_with_images.pptx")
    prs.save(out)
    print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")
else:
    print("No components with images found")

lib.close()
