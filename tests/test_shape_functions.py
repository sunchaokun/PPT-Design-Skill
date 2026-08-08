"""Tests for Phase 1: New shape functions and image cropping functions.

Strict verification of:
- Fill color (hex value, C dict role resolution)
- Line/border (present vs absent, color, width)
- Position and size (left, top, width, height)
- Center-based positioning (cx, cy → left, top)
- Save and reload PPTX (valid .pptx output)
- shape() with string shape_type
- C dict color role resolution
- Image cropping functions
"""

from __future__ import annotations

import io
import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from ppt_pro_max.build_helpers import (
    shape, hexagon, pentagon, octagon, diamond, triangle, right_triangle,
    parallelogram, trapezoid, star5, star6, star8, star10, star12,
    donut, heart, cross, arrow, chevron, cloud, lightning, gear, funnel,
    moon, sun, wave, block_arc, callout, flow_process, flow_decision,
    flow_data, flow_document, flow_connector, no_symbol, plaque, frame,
    cube, folded_corner, tear, math_plus, math_multiply, bevel,
    hex_image, star_image, diamond_image, heart_image, shape_image,
    add_slide,
)


C = {
    'primary': '#1D78FA', 'accent': '#FF5500', 'muted': '#64748B',
    'light': '#E2E8F0', 'white': '#FFFFFF', 'background': '#FFFFFF',
    'card_bg': '#F1F5F9', 'text_dark': '#1A1A1A', 'text_body': '#333333',
    'text_muted': '#666666', 'divider': '#CCCCCC',
}


def _prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _slide(prs=None):
    if prs is None:
        prs = _prs()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _save_reload(prs):
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def _get_fill_hex(sh):
    spPr = sh._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        return None
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        return None
    return srgb.get('val')


def _get_line_fill_hex(sh):
    spPr = sh._element.find(qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        return None
    noFill = ln.find(qn('a:noFill'))
    if noFill is not None:
        return None
    solidFill = ln.find(qn('a:solidFill'))
    if solidFill is None:
        return None
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        return None
    return srgb.get('val')


def _get_line_width(sh):
    spPr = sh._element.find(qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        return None
    return ln.get('w')


def _has_no_line(sh):
    spPr = sh._element.find(qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        return True
    noFill = ln.find(qn('a:noFill'))
    return noFill is not None


# ── shape() universal function ──


class TestShapeUniversal:
    def test_shape_with_enum(self):
        prs = _prs()
        s = _slide(prs)
        sh = shape(s, MSO_SHAPE.HEXAGON, 2, 2, 3, 2.6, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'
        assert _has_no_line(sh)

    def test_shape_with_string(self):
        prs = _prs()
        s = _slide(prs)
        sh = shape(s, 'hexagon', 2, 2, 3, 2.6, '#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'
        assert _has_no_line(sh)

    def test_shape_with_string_uppercase(self):
        prs = _prs()
        s = _slide(prs)
        sh = shape(s, 'HEXAGON', 2, 2, 3, 2.6, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_shape_with_invalid_string_fallback(self):
        prs = _prs()
        s = _slide(prs)
        sh = shape(s, 'NONEXISTENT_SHAPE', 2, 2, 3, 2.6, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_shape_with_line(self):
        prs = _prs()
        s = _slide(prs)
        sh = shape(s, MSO_SHAPE.STAR_5_POINT, 2, 2, 3, 3, '#1D78FA', line='#FF5500')
        assert _get_fill_hex(sh) == '1D78FA'
        assert _get_line_fill_hex(sh) == 'FF5500'
        assert _get_line_width(sh) is not None

    def test_shape_with_C_dict(self):
        prs = _prs()
        s = _slide(prs)
        sh = shape(s, MSO_SHAPE.DONUT, 2, 2, 3, 3, 'primary', line='accent', C=C)
        assert _get_fill_hex(sh) == '1D78FA'
        assert _get_line_fill_hex(sh) == 'FF5500'

    def test_shape_position_size(self):
        prs = _prs()
        s = _slide(prs)
        sh = shape(s, MSO_SHAPE.HEXAGON, 2.0, 3.0, 4.0, 3.48, '#1D78FA')
        assert abs(sh.left - Inches(2.0)) < Emu(100)
        assert abs(sh.top - Inches(3.0)) < Emu(100)
        assert abs(sh.width - Inches(4.0)) < Emu(100)
        assert abs(sh.height - Inches(3.48)) < Emu(100)


# ── Centered shape functions ──


class TestCenteredShapes:
    def test_hexagon_fill(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_hexagon_centered_position(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA')
        assert abs(sh.left - Inches(5 - 2 / 2)) < Emu(100)
        assert abs(sh.top - Inches(3.5 - 2 * 0.87 / 2)) < Emu(100)

    def test_hexagon_with_line(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA', line='#FF5500')
        assert _get_line_fill_hex(sh) == 'FF5500'

    def test_hexagon_no_line(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA')
        assert _has_no_line(sh)

    def test_hexagon_with_C_dict(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, 'primary', line='accent', C=C)
        assert _get_fill_hex(sh) == '1D78FA'
        assert _get_line_fill_hex(sh) == 'FF5500'

    def test_pentagon(self):
        prs = _prs()
        s = _slide(prs)
        sh = pentagon(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'
        assert abs(sh.left - Inches(5 - 1)) < Emu(100)
        assert abs(sh.top - Inches(3.5 - 1)) < Emu(100)

    def test_octagon(self):
        prs = _prs()
        s = _slide(prs)
        sh = octagon(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_diamond(self):
        prs = _prs()
        s = _slide(prs)
        sh = diamond(s, 5, 3.5, 2, '#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'

    def test_diamond_with_line(self):
        prs = _prs()
        s = _slide(prs)
        sh = diamond(s, 5, 3.5, 2, '#1D78FA', line='#333333')
        assert _get_line_fill_hex(sh) == '333333'

    def test_star5(self):
        prs = _prs()
        s = _slide(prs)
        sh = star5(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_star6(self):
        prs = _prs()
        s = _slide(prs)
        sh = star6(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_star8(self):
        prs = _prs()
        s = _slide(prs)
        sh = star8(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_star10(self):
        prs = _prs()
        s = _slide(prs)
        sh = star10(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_star12(self):
        prs = _prs()
        s = _slide(prs)
        sh = star12(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_donut(self):
        prs = _prs()
        s = _slide(prs)
        sh = donut(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_heart(self):
        prs = _prs()
        s = _slide(prs)
        sh = heart(s, 5, 3.5, 2, '#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'

    def test_cross(self):
        prs = _prs()
        s = _slide(prs)
        sh = cross(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_gear_6(self):
        prs = _prs()
        s = _slide(prs)
        sh = gear(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_gear_9(self):
        prs = _prs()
        s = _slide(prs)
        sh = gear(s, 5, 3.5, 2, '#1D78FA', teeth=9)
        assert _get_fill_hex(sh) == '1D78FA'

    def test_moon(self):
        prs = _prs()
        s = _slide(prs)
        sh = moon(s, 5, 3.5, 2, '#FFD700')
        assert _get_fill_hex(sh) == 'FFD700'

    def test_sun(self):
        prs = _prs()
        s = _slide(prs)
        sh = sun(s, 5, 3.5, 2, '#FFD700')
        assert _get_fill_hex(sh) == 'FFD700'

    def test_block_arc(self):
        prs = _prs()
        s = _slide(prs)
        sh = block_arc(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_tear(self):
        prs = _prs()
        s = _slide(prs)
        sh = tear(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_math_plus(self):
        prs = _prs()
        s = _slide(prs)
        sh = math_plus(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_math_multiply(self):
        prs = _prs()
        s = _slide(prs)
        sh = math_multiply(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_flow_connector(self):
        prs = _prs()
        s = _slide(prs)
        sh = flow_connector(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_flow_decision(self):
        prs = _prs()
        s = _slide(prs)
        sh = flow_decision(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_no_symbol(self):
        prs = _prs()
        s = _slide(prs)
        sh = no_symbol(s, 5, 3.5, 2, '#FF0000')
        assert _get_fill_hex(sh) == 'FF0000'


# ── Corner-based shape functions ──


class TestCornerShapes:
    def test_triangle_fill(self):
        prs = _prs()
        s = _slide(prs)
        sh = triangle(s, 2, 2, 3, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_triangle_position(self):
        prs = _prs()
        s = _slide(prs)
        sh = triangle(s, 2.0, 3.0, 4.0, 3.0, '#1D78FA')
        assert abs(sh.left - Inches(2.0)) < Emu(100)
        assert abs(sh.top - Inches(3.0)) < Emu(100)

    def test_triangle_with_line(self):
        prs = _prs()
        s = _slide(prs)
        sh = triangle(s, 2, 2, 3, 2, '#1D78FA', line='#FF5500')
        assert _get_line_fill_hex(sh) == 'FF5500'

    def test_triangle_with_C_dict(self):
        prs = _prs()
        s = _slide(prs)
        sh = triangle(s, 2, 2, 3, 2, 'primary', line='accent', C=C)
        assert _get_fill_hex(sh) == '1D78FA'
        assert _get_line_fill_hex(sh) == 'FF5500'

    def test_right_triangle(self):
        prs = _prs()
        s = _slide(prs)
        sh = right_triangle(s, 2, 2, 3, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_parallelogram(self):
        prs = _prs()
        s = _slide(prs)
        sh = parallelogram(s, 2, 2, 4, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_trapezoid(self):
        prs = _prs()
        s = _slide(prs)
        sh = trapezoid(s, 2, 2, 4, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_arrow(self):
        prs = _prs()
        s = _slide(prs)
        sh = arrow(s, 2, 2, 4, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_arrow_with_line(self):
        prs = _prs()
        s = _slide(prs)
        sh = arrow(s, 2, 2, 4, 2, '#1D78FA', line='#FF5500')
        assert _get_line_fill_hex(sh) == 'FF5500'

    def test_chevron(self):
        prs = _prs()
        s = _slide(prs)
        sh = chevron(s, 2, 2, 4, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_cloud(self):
        prs = _prs()
        s = _slide(prs)
        sh = cloud(s, 2, 2, 4, 3, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_lightning(self):
        prs = _prs()
        s = _slide(prs)
        sh = lightning(s, 2, 2, 2, 3, '#FFD700')
        assert _get_fill_hex(sh) == 'FFD700'

    def test_funnel(self):
        prs = _prs()
        s = _slide(prs)
        sh = funnel(s, 2, 2, 3, 4, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_wave(self):
        prs = _prs()
        s = _slide(prs)
        sh = wave(s, 2, 2, 4, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_plaque(self):
        prs = _prs()
        s = _slide(prs)
        sh = plaque(s, 2, 2, 4, 3, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_frame(self):
        prs = _prs()
        s = _slide(prs)
        sh = frame(s, 2, 2, 4, 3, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_cube(self):
        prs = _prs()
        s = _slide(prs)
        sh = cube(s, 2, 2, 4, 3, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_bevel(self):
        prs = _prs()
        s = _slide(prs)
        sh = bevel(s, 2, 2, 4, 3, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_folded_corner(self):
        prs = _prs()
        s = _slide(prs)
        sh = folded_corner(s, 2, 2, 4, 3, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_flow_process(self):
        prs = _prs()
        s = _slide(prs)
        sh = flow_process(s, 2, 2, 4, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_flow_data(self):
        prs = _prs()
        s = _slide(prs)
        sh = flow_data(s, 2, 2, 4, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_flow_document(self):
        prs = _prs()
        s = _slide(prs)
        sh = flow_document(s, 2, 2, 4, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'


# ── Callout function ──


class TestCallout:
    def test_callout_rect(self):
        prs = _prs()
        s = _slide(prs)
        sh = callout(s, 2, 2, 4, 3, '#1D78FA', style='rect')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_callout_round(self):
        prs = _prs()
        s = _slide(prs)
        sh = callout(s, 2, 2, 4, 3, '#1D78FA', style='round')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_callout_oval(self):
        prs = _prs()
        s = _slide(prs)
        sh = callout(s, 2, 2, 4, 3, '#1D78FA', style='oval')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_callout_cloud(self):
        prs = _prs()
        s = _slide(prs)
        sh = callout(s, 2, 2, 4, 3, '#1D78FA', style='cloud')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_callout_default_is_rect(self):
        prs = _prs()
        s = _slide(prs)
        sh = callout(s, 2, 2, 4, 3, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_callout_invalid_style_fallback(self):
        prs = _prs()
        s = _slide(prs)
        sh = callout(s, 2, 2, 4, 3, '#1D78FA', style='nonexistent')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_callout_with_line(self):
        prs = _prs()
        s = _slide(prs)
        sh = callout(s, 2, 2, 4, 3, '#1D78FA', line='#FF5500')
        assert _get_line_fill_hex(sh) == 'FF5500'

    def test_callout_with_C_dict(self):
        prs = _prs()
        s = _slide(prs)
        sh = callout(s, 2, 2, 4, 3, 'primary', line='accent', C=C)
        assert _get_fill_hex(sh) == '1D78FA'
        assert _get_line_fill_hex(sh) == 'FF5500'


# ── Save/Reload integrity ──


class TestSaveReload:
    def test_shape_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        shape(s, MSO_SHAPE.HEXAGON, 2, 2, 3, 2.6, '#1D78FA')
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1

    def test_multiple_shapes_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        hexagon(s, 3, 3, 2, '#1D78FA')
        star5(s, 7, 3, 2, '#FF5500')
        donut(s, 5, 5, 2, '#00B050')
        arrow(s, 9, 2, 3, 2, '#FFD700')
        prs2 = _save_reload(prs)
        slide2 = prs2.slides[0]
        assert len(list(slide2.shapes)) >= 4

    def test_all_centered_shapes_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        for fn in [hexagon, pentagon, octagon, diamond, star5, star6,
                   star8, star10, star12, donut, heart, cross, moon, sun,
                   block_arc, tear, math_plus, math_multiply]:
            fn(s, 5, 3.5, 1, '#1D78FA')
        prs2 = _save_reload(prs)
        assert len(list(prs2.slides[0].shapes)) >= 18

    def test_all_corner_shapes_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        for fn in [triangle, right_triangle, parallelogram, trapezoid,
                   arrow, chevron, cloud, lightning, funnel, wave,
                   plaque, frame, cube, bevel, folded_corner,
                   flow_process, flow_data, flow_document]:
            fn(s, 2, 2, 3, 2, '#1D78FA')
        prs2 = _save_reload(prs)
        assert len(list(prs2.slides[0].shapes)) >= 18

    def test_save_to_file(self, tmp_path):
        prs = _prs()
        s = _slide(prs)
        hexagon(s, 5, 3.5, 2, '#1D78FA')
        star5(s, 8, 3.5, 2, '#FF5500')
        out = tmp_path / "shapes_test.pptx"
        prs.save(str(out))
        assert out.exists()
        assert out.stat().st_size > 0
        prs2 = Presentation(str(out))
        assert len(prs2.slides) == 1


# ── Color resolution edge cases ──


class TestColorResolution:
    def test_fill_with_hash_hex(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    def test_fill_without_hash_resolves_as_C_role(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '1D78FA')
        assert _get_fill_hex(sh) == '000000'

    def test_fill_with_C_role(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, 'primary', C=C)
        assert _get_fill_hex(sh) == '1D78FA'

    def test_fill_with_missing_C_role(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, 'nonexistent', C=C)
        assert _get_fill_hex(sh) == '000000'

    def test_fill_with_none_C(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, 'primary', C=None)
        assert _get_fill_hex(sh) == '000000'

    def test_line_with_C_role(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA', line='accent', C=C)
        assert _get_line_fill_hex(sh) == 'FF5500'

    def test_line_with_hex(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA', line='#333333')
        assert _get_line_fill_hex(sh) == '333333'

    def test_no_line_by_default(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA')
        assert _has_no_line(sh)

    def test_line_width_is_1pt(self):
        prs = _prs()
        s = _slide(prs)
        sh = hexagon(s, 5, 3.5, 2, '#1D78FA', line='#FF5500')
        w = _get_line_width(sh)
        assert w is not None
        assert abs(int(w) - int(Pt(1))) < Emu(100)


# ── Image cropping functions ──


class TestImageCropping:
    @pytest.fixture
    def sample_image(self, tmp_path):
        from PIL import Image
        img = Image.new('RGB', (200, 200), color='red')
        path = tmp_path / 'test_image.png'
        img.save(str(path))
        return str(path)

    def test_hex_image_creates_shape(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = hex_image(s, 5, 3.5, 2, sample_image)
        assert sh is not None

    def test_star_image_default_5point(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = star_image(s, 5, 3.5, 2, sample_image)
        assert sh is not None

    def test_star_image_6point(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = star_image(s, 5, 3.5, 2, sample_image, points=6)
        assert sh is not None

    def test_star_image_8point(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = star_image(s, 5, 3.5, 2, sample_image, points=8)
        assert sh is not None

    def test_star_image_invalid_points_fallback(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = star_image(s, 5, 3.5, 2, sample_image, points=99)
        assert sh is not None

    def test_diamond_image_creates_shape(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = diamond_image(s, 5, 3.5, 2, sample_image)
        assert sh is not None

    def test_heart_image_creates_shape(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = heart_image(s, 5, 3.5, 2, sample_image)
        assert sh is not None

    def test_shape_image_with_enum(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = shape_image(s, MSO_SHAPE.HEXAGON, 2, 2, 3, 3, sample_image)
        assert sh is not None

    def test_shape_image_with_string(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = shape_image(s, 'HEXAGON', 2, 2, 3, 3, sample_image)
        assert sh is not None

    def test_shape_image_with_border(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        sh = shape_image(s, MSO_SHAPE.HEXAGON, 2, 2, 3, 3, sample_image,
                         border_color='#FF5500')
        assert sh is not None
        spPr = sh._element.find(qn('p:spPr'))
        ln = spPr.find(qn('a:ln'))
        assert ln is not None

    def test_image_save_reload(self, sample_image):
        prs = _prs()
        s = _slide(prs)
        hex_image(s, 5, 3.5, 2, sample_image)
        star_image(s, 8, 3.5, 2, sample_image)
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1

    def test_image_invalid_path(self):
        prs = _prs()
        s = _slide(prs)
        sh = hex_image(s, 5, 3.5, 2, '/nonexistent/path.png')
        assert sh is not None


# ── Line width consistency ──


class TestLineWidthConsistency:
    @pytest.mark.parametrize("fn_name,kwargs", [
        ("hexagon", {"cx": 5, "cy": 3.5, "size": 2}),
        ("pentagon", {"cx": 5, "cy": 3.5, "size": 2}),
        ("diamond", {"cx": 5, "cy": 3.5, "size": 2}),
        ("star5", {"cx": 5, "cy": 3.5, "size": 2}),
        ("donut", {"cx": 5, "cy": 3.5, "size": 2}),
        ("heart", {"cx": 5, "cy": 3.5, "size": 2}),
        ("cross", {"cx": 5, "cy": 3.5, "size": 2}),
        ("triangle", {"left": 2, "top": 2, "width": 3, "height": 2}),
        ("arrow", {"left": 2, "top": 2, "width": 3, "height": 2}),
        ("chevron", {"left": 2, "top": 2, "width": 3, "height": 2}),
    ])
    def test_line_width_1pt(self, fn_name, kwargs):
        import ppt_pro_max.build_helpers as bh
        fn = getattr(bh, fn_name)
        prs = _prs()
        s = _slide(prs)
        sh = fn(s, **kwargs, fill='#1D78FA', line='#FF5500')
        w = _get_line_width(sh)
        assert w is not None
        assert abs(int(w) - int(Pt(1))) < Emu(100)


# ── Fill color consistency across all functions ──


class TestFillColorConsistency:
    CENTERED_FNS = [
        hexagon, pentagon, octagon, diamond, star5, star6, star8,
        star10, star12, donut, heart, cross, moon, sun, block_arc,
        tear, math_plus, math_multiply,
    ]

    CORNER_FNS = [
        triangle, right_triangle, parallelogram, trapezoid,
        arrow, chevron, cloud, lightning, funnel, wave,
        plaque, frame, cube, bevel, folded_corner,
        flow_process, flow_data, flow_document,
    ]

    @pytest.mark.parametrize("fn", CENTERED_FNS)
    def test_centered_fill_color(self, fn):
        prs = _prs()
        s = _slide(prs)
        sh = fn(s, 5, 3.5, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    @pytest.mark.parametrize("fn", CORNER_FNS)
    def test_corner_fill_color(self, fn):
        prs = _prs()
        s = _slide(prs)
        sh = fn(s, 2, 2, 3, 2, '#1D78FA')
        assert _get_fill_hex(sh) == '1D78FA'

    @pytest.mark.parametrize("fn", CENTERED_FNS)
    def test_centered_no_line_default(self, fn):
        prs = _prs()
        s = _slide(prs)
        sh = fn(s, 5, 3.5, 2, '#1D78FA')
        assert _has_no_line(sh)

    @pytest.mark.parametrize("fn", CORNER_FNS)
    def test_corner_no_line_default(self, fn):
        prs = _prs()
        s = _slide(prs)
        sh = fn(s, 2, 2, 3, 2, '#1D78FA')
        assert _has_no_line(sh)
