"""Extract font/color from XML level for Build mode PPTs."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

out = os.path.join(tempfile.gettempdir(), 'ppt_build_test')

for fname, label in [('build_mckinsey.pptx', 'McKinsey'), ('build_cyberpunk.pptx', 'Cyberpunk'), ('build_creative.pptx', 'Creative')]:
    path = os.path.join(out, fname)
    if not os.path.exists(path):
        continue
    prs = Presentation(path)
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    
    all_fonts, all_colors, all_sizes = set(), set(), []
    
    print(f"\n{'='*80}")
    print(f"  {label} | {len(prs.slides)} slides")
    print(f"{'='*80}")
    
    for i, slide in enumerate(prs.slides):
        slide_fonts, slide_colors, slide_sizes = [], [], []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rPr = run._r.find(f'{{{ns_a}}}rPr')
                    font = sz = clr = None
                    if rPr is not None:
                        latin = rPr.find(f'{{{ns_a}}}latin')
                        if latin is not None:
                            font = latin.get('typeface')
                        sz = rPr.get('sz')
                        for child in rPr:
                            if 'srgbClr' in child.tag:
                                clr = child.get('val')
                    if font: all_fonts.add(font); slide_fonts.append(font)
                    if sz: 
                        pt = int(sz) / 100
                        all_sizes.append(pt)
                        slide_sizes.append(pt)
                    if clr: all_colors.add(clr); slide_colors.append(clr)
        
        mn = min(slide_sizes) if slide_sizes else "-"
        mx = max(slide_sizes) if slide_sizes else "-"
        print(f"  Slide {i}: {len(slide.shapes)} shapes | {mn}-{mx}pt | {len(set(slide_colors))} colors | fonts={sorted(set(slide_fonts))}")
    
    if all_sizes:
        print(f"\n  Font sizes: {min(all_sizes)}-{max(all_sizes)}pt ({len(set(all_sizes))} unique levels)")
    print(f"  Fonts: {sorted(all_fonts)}")
    print(f"  Colors: {sorted(all_colors)} ({len(all_colors)} unique)")
