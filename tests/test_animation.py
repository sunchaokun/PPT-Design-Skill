"""Tests for animation module — transitions and entrance effects."""

from __future__ import annotations

from pptx import Presentation
from pptx.oxml.ns import qn


class TestSlideTransition:

    def test_add_fade_transition(self):
        from ppt_pro_max.renderer.animation import add_slide_transition
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_slide_transition(slide, transition_type="fade", speed="medium")
        sld = slide._element
        trans = sld.find(qn("p:transition"))
        assert trans is not None
        assert trans.get("spd") == "med"
        fade = trans.find(qn("p:fade"))
        assert fade is not None

    def test_add_push_transition(self):
        from ppt_pro_max.renderer.animation import add_slide_transition
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_slide_transition(slide, transition_type="push", speed="fast")
        sld = slide._element
        trans = sld.find(qn("p:transition"))
        assert trans is not None
        assert trans.get("spd") == "fast"
        push = trans.find(qn("p:push"))
        assert push is not None

    def test_advance_after_time(self):
        from ppt_pro_max.renderer.animation import add_slide_transition
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_slide_transition(slide, transition_type="fade", advance_after_ms=5000)
        sld = slide._element
        trans = sld.find(qn("p:transition"))
        assert trans.get("advTm") == "5000"

    def test_invalid_type_defaults_to_fade(self):
        from ppt_pro_max.renderer.animation import add_slide_transition
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_slide_transition(slide, transition_type="nonexistent")
        sld = slide._element
        trans = sld.find(qn("p:transition"))
        assert trans is not None
        fade = trans.find(qn("p:fade"))
        assert fade is not None

    def test_overwrite_existing_transition(self):
        from ppt_pro_max.renderer.animation import add_slide_transition
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_slide_transition(slide, transition_type="fade")
        add_slide_transition(slide, transition_type="wipe")
        sld = slide._element
        transitions = sld.findall(qn("p:transition"))
        assert len(transitions) == 1
        wipe = transitions[0].find(qn("p:wipe"))
        assert wipe is not None

    def test_all_transition_types(self):
        from ppt_pro_max.renderer.animation import add_slide_transition, TRANSITION_TYPES
        for t_type in TRANSITION_TYPES:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            add_slide_transition(slide, transition_type=t_type)
            sld = slide._element
            trans = sld.find(qn("p:transition"))
            assert trans is not None, f"Transition {t_type} not added"


class TestEntranceAnimation:

    def test_add_fade_in(self):
        from ppt_pro_max.renderer.animation import add_entrance_animation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_entrance_animation(slide, shape_id=2, effect="fade_in")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None

    def test_add_appear(self):
        from ppt_pro_max.renderer.animation import add_entrance_animation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_entrance_animation(slide, shape_id=2, effect="appear")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None

    def test_add_multiple_animations(self):
        from ppt_pro_max.renderer.animation import add_entrance_animation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_entrance_animation(slide, shape_id=2, effect="fade_in", click_triggered=True)
        add_entrance_animation(slide, shape_id=3, effect="fade_in", click_triggered=False)
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None

    def test_invalid_effect_defaults_to_fade_in(self):
        from ppt_pro_max.renderer.animation import add_entrance_animation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_entrance_animation(slide, shape_id=2, effect="nonexistent")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None

    def test_all_entrance_presets(self):
        from ppt_pro_max.renderer.animation import add_entrance_animation, ENTRANCE_PRESETS
        for effect in ENTRANCE_PRESETS:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            add_entrance_animation(slide, shape_id=2, effect=effect)
            sld = slide._element
            timing = sld.find(qn("p:timing"))
            assert timing is not None, f"Effect {effect} failed"


class TestAnimationSequence:

    def test_add_sequence(self):
        from ppt_pro_max.renderer.animation import add_animation_sequence
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_animation_sequence(slide, shape_ids=[2, 3, 4], effect="fade_in")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None

    def test_sequence_with_fly_in(self):
        from ppt_pro_max.renderer.animation import add_animation_sequence
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        add_animation_sequence(slide, shape_ids=[2, 3], effect="fly_in")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None


class TestPipelineAnimations:

    def test_motion_adds_transition(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        pipeline = EnterprisePipeline()
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        design = {"title": "Test", "bullets": ["a", "b"]}
        precision.render_slide(prs, design)
        slide = prs.slides[-1]
        pipeline._apply_animations(slide, design, motion=5)
        sld = slide._element
        trans = sld.find(qn("p:transition"))
        assert trans is not None

    def test_motion_adds_entrance_animation(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        pipeline = EnterprisePipeline()
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        design = {"title": "Test", "bullets": ["a", "b"]}
        precision.render_slide(prs, design)
        slide = prs.slides[-1]
        pipeline._apply_animations(slide, design, motion=5)
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None

    def test_low_motion_no_entrance(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        pipeline = EnterprisePipeline()
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        design = {"title": "Test"}
        precision.render_slide(prs, design)
        slide = prs.slides[-1]
        pipeline._apply_animations(slide, design, motion=1)
        sld = slide._element
        trans = sld.find(qn("p:transition"))
        assert trans is not None
        timing = sld.find(qn("p:timing"))
        assert timing is None
