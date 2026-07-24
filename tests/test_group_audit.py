"""Audit GroupShape internal font sizes."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
path = os.path.join(tempfile.gettempdir(), 'ppt_build_v2', 'build_mckinsey.pptx')
prs = Presentation(path)
slide = prs.slides[2]

for shape in slide.shapes:
    if shape.shape_type != 6:
        continue
    print(f'GROUP at ({shape.left},{shape.top})')
    for child in shape.shapes:
        if not child.has_text_frame:
            continue
        txt = child.text_frame.text[:30]
        for para in child.text_frame.paragraphs:
            pPr = para._p.find(f'{{{ns_a}}}pPr')
            def_sz = def_font = def_clr = None
            if pPr is not None:
                defRPr = pPr.find(f'{{{ns_a}}}defRPr')
                if defRPr is not None:
                    def_sz = defRPr.get('sz')
                    latin = defRPr.find(f'{{{ns_a}}}latin')
                    if latin is not None:
                        def_font = latin.get('typeface')
                    for c in defRPr:
                        if 'srgbClr' in c.tag:
                            def_clr = c.get('val')
            for run in para.runs:
                if not run.text.strip():
                    continue
                rPr = run._r.find(f'{{{ns_a}}}rPr')
                sz = font = clr = None
                if rPr is not None:
                    sz = rPr.get('sz')
                    latin = rPr.find(f'{{{ns_a}}}latin')
                    if latin is not None:
                        font = latin.get('typeface')
                    for c in rPr:
                        if 'srgbClr' in c.tag:
                            clr = c.get('val')
                pt = int(sz)/100 if sz else (int(def_sz)/100 if def_sz else '?')
                fn = font or def_font or '?'
                cl = clr or def_clr or '?'
                print(f'  "{run.text[:20]}" {pt}pt {fn} #{cl}')
