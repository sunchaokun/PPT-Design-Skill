"""python-pptx capability verification — tests every API category from the reference doc."""

import os
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree


PASS = []
FAIL = []


def test(name, fn):
    try:
        fn()
        PASS.append(name)
        print(f"  PASS  {name}")
    except Exception as e:
        FAIL.append((name, str(e)))
        print(f"  FAIL  {name} — {e}")


def make_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def save(prs, name):
    tmp = Path(tempfile.gettempdir()) / "ppt_pro_max_test" / f"{name}.pptx"
    tmp.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(tmp))
    assert tmp.exists() and tmp.stat().st_size > 1000
    return tmp


# ═══════════════════════════════════════════════════════════════════════════════
# 1. SHAPE TYPES
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 1. SHAPE TYPES ===")

def _test_shape_type():
    shapes_to_test = [
        ("RECTANGLE", MSO_SHAPE.RECTANGLE),
        ("ROUNDED_RECTANGLE", MSO_SHAPE.ROUNDED_RECTANGLE),
        ("OVAL", MSO_SHAPE.OVAL),
        ("DIAMOND", MSO_SHAPE.DIAMOND),
        ("HEXAGON", MSO_SHAPE.HEXAGON),
        ("DONUT", MSO_SHAPE.DONUT),
        ("ISOSCELES_TRIANGLE", MSO_SHAPE.ISOSCELES_TRIANGLE),
        ("CHEVRON", MSO_SHAPE.CHEVRON),
        ("STAR_5_POINT", MSO_SHAPE.STAR_5_POINT),
        ("STAR_8_POINT", MSO_SHAPE.STAR_8_POINT),
        ("LIGHTNING_BOLT", MSO_SHAPE.LIGHTNING_BOLT),
        ("HEART", MSO_SHAPE.HEART),
        ("CLOUD", MSO_SHAPE.CLOUD),
        ("CUBE", MSO_SHAPE.CUBE),
        ("FLOWCHART_DECISION", MSO_SHAPE.FLOWCHART_DECISION),
        ("FLOWCHART_PROCESS", MSO_SHAPE.FLOWCHART_PROCESS),
        ("FLOWCHART_DOCUMENT", MSO_SHAPE.FLOWCHART_DOCUMENT),
        ("RECTANGULAR_CALLOUT", MSO_SHAPE.RECTANGULAR_CALLOUT),
        ("CLOUD_CALLOUT", MSO_SHAPE.CLOUD_CALLOUT),
        ("RIGHT_ARROW", MSO_SHAPE.RIGHT_ARROW),
        ("BENT_ARROW", MSO_SHAPE.BENT_ARROW),
        ("CIRCULAR_ARROW", MSO_SHAPE.CIRCULAR_ARROW),
        ("MATH_PLUS", MSO_SHAPE.MATH_PLUS),
        ("GEAR_6", MSO_SHAPE.GEAR_6),
        ("FUNNEL", MSO_SHAPE.FUNNEL),
        ("BEVEL", MSO_SHAPE.BEVEL),
        ("LEFT_BRACE", MSO_SHAPE.LEFT_brace if hasattr(MSO_SHAPE, 'left_brace') else MSO_SHAPE.LEFT_BRACE),
        ("FOLDED_CORNER", MSO_SHAPE.FOLDED_CORNER),
        ("CROSS", MSO_SHAPE.CROSS),
        ("FRAME", MSO_SHAPE.FRAME),
    ]
    prs, slide = make_prs()
    for i, (name, shape_type) in enumerate(shapes_to_test):
        row = i // 10
        col = i % 10
        s = slide.shapes.add_shape(shape_type, Inches(col * 1.3 + 0.2), Inches(row * 1.5 + 0.5), Inches(1.1), Inches(1.1))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
        s.line.fill.background()
    path = save(prs, "shape_types")
    assert len(slide.shapes) == len(shapes_to_test)

test("170+ Shape Types (30 verified)", _test_shape_type)


# ═══════════════════════════════════════════════════════════════════════════════
# 2. FILL TYPES
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 2. FILL TYPES ===")

def _test_solid_fill():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
    save(prs, "solid_fill")

test("Solid Fill", _test_solid_fill)


def _test_gradient_fill_api():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2))
    s.fill.gradient()
    s.fill.gradient_stops[0].color.rgb = RGBColor(0xFF, 0x00, 0x00)
    s.fill.gradient_stops[0].position = 0.0
    s.fill.gradient_stops[1].color.rgb = RGBColor(0x00, 0x00, 0xFF)
    s.fill.gradient_stops[1].position = 1.0
    save(prs, "gradient_fill_api")

test("Gradient Fill (High-Level API)", _test_gradient_fill_api)


def _test_gradient_fill_xml():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1), Inches(3), Inches(2))
    spPr = s._element.spPr
    # Remove existing fill
    for child in list(spPr):
        if child.tag.endswith("solidFill") or child.tag.endswith("noFill"):
            spPr.remove(child)
    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    gs1 = etree.SubElement(gsLst, qn("a:gs"))
    gs1.set("pos", "0")
    srgb1 = etree.SubElement(gs1, qn("a:srgbClr"))
    srgb1.set("val", "FF6600")
    gs2 = etree.SubElement(gsLst, qn("a:gs"))
    gs2.set("pos", "100000")
    srgb2 = etree.SubElement(gs2, qn("a:srgbClr"))
    srgb2.set("val", "FFCC00")
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", "5400000")
    lin.set("scaled", "1")
    save(prs, "gradient_fill_xml")

test("Gradient Fill (XML)", _test_gradient_fill_xml)


def _test_no_fill():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(3), Inches(2))
    s.fill.background()
    s.line.color.rgb = RGBColor(0x00, 0x00, 0xFF)
    s.line.width = Pt(3)
    save(prs, "no_fill")

test("No Fill (Transparent)", _test_no_fill)


def _test_pattern_fill():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4), Inches(3), Inches(2))
    s.fill.patterned()
    from pptx.enum.dml import MSO_PATTERN_TYPE
    s.fill.pattern = MSO_PATTERN_TYPE.DARK_DOWNWARD_DIAGONAL
    s.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0xFF)
    s.fill.back_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    save(prs, "pattern_fill")

test("Pattern Fill", _test_pattern_fill)


# ═══════════════════════════════════════════════════════════════════════════════
# 3. LINE STYLES
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 3. LINE STYLES ===")

def _test_dash_styles():
    prs, slide = make_prs()
    dashes = [
        MSO_LINE_DASH_STYLE.SOLID,
        MSO_LINE_DASH_STYLE.DASH,
        MSO_LINE_DASH_STYLE.DASH_DOT,
        MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
        MSO_LINE_DASH_STYLE.LONG_DASH,
        MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
        MSO_LINE_DASH_STYLE.ROUND_DOT,
        MSO_LINE_DASH_STYLE.SQUARE_DOT,
    ]
    for i, dash in enumerate(dashes):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5 + i * 1.5), Inches(1), Inches(1.3), Inches(1))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        s.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
        s.line.width = Pt(3)
        s.line.dash_style = dash
    save(prs, "dash_styles")

test("Dash Styles (8 types)", _test_dash_styles)


# ═══════════════════════════════════════════════════════════════════════════════
# 4. FONT PROPERTIES
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 4. FONT PROPERTIES ===")

def _test_font_italic_underline():
    prs, slide = make_prs()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]

    r1 = p.add_run()
    r1.text = "Bold "
    r1.font.bold = True
    r1.font.size = Pt(24)

    r2 = p.add_run()
    r2.text = "Italic "
    r2.font.italic = True
    r2.font.size = Pt(24)

    r3 = p.add_run()
    r3.text = "Underline "
    r3.font.underline = True
    r3.font.size = Pt(24)

    r4 = p.add_run()
    r4.text = "Strikethrough "
    r4.font.strikethrough = True
    r4.font.size = Pt(24)

    r5 = p.add_run()
    r5.text = "Superscript "
    r5.font.superscript = True
    r5.font.size = Pt(24)

    r6 = p.add_run()
    r6.text = "Subscript"
    r6.font.subscript = True
    r6.font.size = Pt(24)

    save(prs, "font_props")

test("Font: bold/italic/underline/strikethrough/superscript/subscript", _test_font_italic_underline)


def _test_text_alignment_and_anchor():
    prs, slide = make_prs()
    alignments = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.JUSTIFY, PP_ALIGN.DISTRIBUTE]
    for i, align in enumerate(alignments):
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5 + i * 1.2), Inches(10), Inches(1))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = f"Alignment: {align}"
        r.font.size = Pt(20)
    save(prs, "text_align_anchor")

test("Text Alignment (5) + Vertical Anchor", _test_text_alignment_and_anchor)


def _test_text_auto_size():
    prs, slide = make_prs()
    from pptx.enum.text import MSO_AUTO_SIZE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = "This long text should auto-shrink to fit the shape without overflow"
    save(prs, "auto_size")

test("Auto-Size (TEXT_TO_FIT_SHAPE)", _test_text_auto_size)


def _test_paragraph_spacing():
    prs, slide = make_prs()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    tf = txBox.text_frame
    for i in range(5):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"Line {i+1}"
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        p.font.size = Pt(18)
    save(prs, "para_spacing")

test("Paragraph spacing (space_before/space_after)", _test_paragraph_spacing)


# ═══════════════════════════════════════════════════════════════════════════════
# 5. TABLES
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 5. TABLES ===")

def _test_table():
    prs, slide = make_prs()
    table_shape = slide.shapes.add_table(rows=4, cols=3, left=Inches(1), top=Inches(1), width=Inches(10), height=Inches(4))
    table = table_shape.table

    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(4)
    table.columns[2].width = Inches(3)

    headers = ["Feature", "Description", "Status"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r.font.bold = True

    data = [
        ["Alpha Generation", "AI-powered alpha creation", "Done"],
        ["Backtesting", "Full IS/OS backtest suite", "Done"],
        ["Risk Control", "Drawdown & correlation limits", "WIP"],
    ]
    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val

    table.first_row = True
    table.horz_banding = True

    save(prs, "table")

test("Table (add_table, cell formatting, banding)", _test_table)


# ═══════════════════════════════════════════════════════════════════════════════
# 6. CHARTS
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 6. CHARTS ===")

def _test_chart_types():
    prs, slide = make_prs()
    chart_types = [
        ("Column", XL_CHART_TYPE.COLUMN_CLUSTERED),
        ("Bar", XL_CHART_TYPE.BAR_CLUSTERED),
        ("Line", XL_CHART_TYPE.LINE),
        ("Pie", XL_CHART_TYPE.PIE),
        ("Doughnut", XL_CHART_TYPE.DOUGHNUT),
        ("Area", XL_CHART_TYPE.AREA),
        ("Radar", XL_CHART_TYPE.RADAR_FILLED),
        ("Line Markers", XL_CHART_TYPE.LINE_MARKERS),
    ]
    for i, (name, ct) in enumerate(chart_types):
        chart_data = CategoryChartData()
        chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
        chart_data.add_series("Revenue", (30, 45, 60, 75))
        row = i // 4
        col = i % 4
        chart_frame = slide.shapes.add_chart(ct, Inches(col * 3.2 + 0.3), Inches(row * 3.5 + 0.5), Inches(3), Inches(3), chart_data)
        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    save(prs, "chart_types")

test("Chart Types (8 verified)", _test_chart_types)


# ═══════════════════════════════════════════════════════════════════════════════
# 7. CONNECTORS
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 7. CONNECTORS ===")

def _test_connectors():
    prs, slide = make_prs()
    s1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2), Inches(1.5), Inches(1.5))
    s1.fill.solid()
    s1.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)

    s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(2), Inches(1.5), Inches(1.5))
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(0xE8, 0xA8, 0x38)

    s3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(2), Inches(1.5), Inches(1.5))
    s3.fill.solid()
    s3.fill.fore_color.rgb = RGBColor(0x2E, 0x7D, 0x32)

    straight = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Inches(2.5), Inches(2.75), Inches(5), Inches(2.75))
    straight.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    straight.line.width = Pt(2)

    elbow = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.ELBOW, Inches(6.5), Inches(2.75), Inches(9), Inches(2.75))
    elbow.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    elbow.line.width = Pt(3)

    curve = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.CURVE, Inches(2.5), Inches(3.5), Inches(9), Inches(3.5))
    curve.line.color.rgb = RGBColor(0x00, 0x00, 0xFF)
    curve.line.width = Pt(2)
    curve.line.dash_style = MSO_LINE_DASH_STYLE.DASH_DOT

    save(prs, "connectors")

test("Connectors (straight/elbow/curve + dash)", _test_connectors)


# ═══════════════════════════════════════════════════════════════════════════════
# 8. HYPERLINKS
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 8. HYPERLINKS ===")

def _test_hyperlinks():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(3), Inches(4), Inches(1.5))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    tf = s.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Visit GitHub"
    r.font.size = Pt(24)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.hyperlink.address = "https://github.com"

    s2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5), Inches(4), Inches(1.5))
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(0x2E, 0x7D, 0x32)
    s2.click_action.hyperlink.address = "https://example.com"

    save(prs, "hyperlinks")

test("Hyperlinks (run + shape click)", _test_hyperlinks)


# ═══════════════════════════════════════════════════════════════════════════════
# 9. EFFECTS (Shadow / Glow / 3D)
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 9. EFFECTS ===")

def _test_shadow():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2), Inches(5), Inches(3))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    spPr = s._element.spPr
    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
    outerShdw.set("blurRad", "76200")
    outerShdw.set("dist", "38100")
    outerShdw.set("dir", "5400000")
    srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"))
    srgbClr.set("val", "000000")
    alpha = etree.SubElement(srgbClr, qn("a:alpha"))
    alpha.set("val", "25000")
    save(prs, "shadow")

test("Shadow (outerShdw via XML)", _test_shadow)


def _test_glow():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(2), Inches(4), Inches(4))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x00, 0x70, 0xC0)
    spPr = s._element.spPr
    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    glow = etree.SubElement(effectLst, qn("a:glow"))
    glow.set("rad", "76200")
    srgbClr = etree.SubElement(glow, qn("a:srgbClr"))
    srgbClr.set("val", "00BFFF")
    alpha = etree.SubElement(srgbClr, qn("a:alpha"))
    alpha.set("val", "40000")
    save(prs, "glow")

test("Glow Effect (via XML)", _test_glow)


def _test_3d_effect():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(2), Inches(6), Inches(3))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x2E, 0x7D, 0x32)
    spPr = s._element.spPr
    sp3d = etree.SubElement(spPr, qn("a:sp3d"))
    sp3d.set("z", "254000")
    bevelT = etree.SubElement(sp3d, qn("a:bevelT"))
    bevelT.set("w", "63500")
    bevelT.set("h", "25400")
    prstMaterial = etree.SubElement(sp3d, qn("a:prstMaterial"))
    prstMaterial.set("val", "powder")
    scene3d = etree.SubElement(spPr, qn("a:scene3d"))
    camera = etree.SubElement(scene3d, qn("a:camera"))
    camera.set("prst", "perspectiveFront")
    lightRig = etree.SubElement(scene3d, qn("a:lightRig"))
    lightRig.set("rig", "threePt")
    lightRig.set("dir", "t")
    save(prs, "3d_effect")

test("3D Effect (bevel + material + lighting)", _test_3d_effect)


# ═══════════════════════════════════════════════════════════════════════════════
# 10. SLIDE BACKGROUND + NOTES + METADATA
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 10. BACKGROUND / NOTES / METADATA ===")

def _test_background():
    prs, slide = make_prs()
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x1E, 0x3D)
    save(prs, "background")

test("Solid Background", _test_background)


def _test_speaker_notes():
    prs, slide = make_prs()
    notes = slide.notes_slide
    notes.notes_text_frame.text = "Key talking points for this slide: ... "
    save(prs, "notes")

test("Speaker Notes", _test_speaker_notes)


def _test_core_properties():
    prs, slide = make_prs()
    prs.core_properties.author = "Test Author"
    prs.core_properties.title = "Capability Verification"
    prs.core_properties.subject = "python-pptx API Test"
    prs.core_properties.keywords = "test, verification"
    save(prs, "core_props")
    assert prs.core_properties.author == "Test Author"
    assert prs.core_properties.title == "Capability Verification"

test("Core Properties (author/title/keywords)", _test_core_properties)


# ═══════════════════════════════════════════════════════════════════════════════
# 11. FREEFORM PATHS
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 11. FREEFORM ===")

def _test_freeform():
    prs, slide = make_prs()
    freeform = slide.shapes.build_freeform(Inches(1), Inches(1))
    freeform.add_line_segments([
        (Inches(5), Inches(1)),
        (Inches(5), Inches(5)),
        (Inches(1), Inches(5)),
        (Inches(1), Inches(1)),
    ])
    shape = freeform.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xA8, 0x38)
    shape.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    shape.line.width = Pt(2)
    save(prs, "freeform")

test("Freeform Path (build_freeform + line_to + close)", _test_freeform)


# ═══════════════════════════════════════════════════════════════════════════════
# 12. GROUP SHAPES
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 12. GROUP SHAPES ===")

def _test_group_shape():
    prs, slide = make_prs()
    group = slide.shapes.add_group_shape()
    c1 = group.shapes.add_shape(MSO_SHAPE.OVAL, Emu(0), Emu(0), Inches(1), Inches(1))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
    c2 = group.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Emu(0), Inches(1), Inches(1))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00)
    c3 = group.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(1.1), Inches(1), Inches(1))
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0xFF)
    assert len(group.shapes) == 3
    save(prs, "group_shape")

test("Group Shape (add_group_shape + children)", _test_group_shape)


# ═══════════════════════════════════════════════════════════════════════════════
# 13. GRADIENT OVERLAY (alpha gradient on shape)
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 13. GRADIENT OVERLAY ===")

def _test_gradient_overlay():
    prs, slide = make_prs()
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    spPr = s._element.spPr
    # Add transparent→dark gradient overlay
    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    gs1 = etree.SubElement(gsLst, qn("a:gs"))
    gs1.set("pos", "0")
    srgb1 = etree.SubElement(gs1, qn("a:srgbClr"))
    srgb1.set("val", "000000")
    a1 = etree.SubElement(srgb1, qn("a:alpha"))
    a1.set("val", "0")
    gs2 = etree.SubElement(gsLst, qn("a:gs"))
    gs2.set("pos", "100000")
    srgb2 = etree.SubElement(gs2, qn("a:srgbClr"))
    srgb2.set("val", "000000")
    a2 = etree.SubElement(srgb2, qn("a:alpha"))
    a2.set("val", "80000")
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", "5400000")
    lin.set("scaled", "1")
    save(prs, "gradient_overlay")

test("Gradient Overlay (alpha fade)", _test_gradient_overlay)


# ═══════════════════════════════════════════════════════════════════════════════
# 14. PICTURE CROPPING
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== 14. PICTURE CROPPING ===")

def _test_picture_cropping():
    prs, slide = make_prs()
    # Create a small test image
    from PIL import Image as PILImage
    img = PILImage.new("RGB", (200, 100), color=(0, 112, 192))
    tmp_img = Path(tempfile.gettempdir()) / "ppt_pro_max_test" / "test_crop.png"
    tmp_img.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(tmp_img))
    pic = slide.shapes.add_picture(str(tmp_img), Inches(1), Inches(1), Inches(6), Inches(3))
    pic.crop_left = 0.1
    pic.crop_right = 0.1
    pic.crop_top = 0.05
    pic.crop_bottom = 0.05
    save(prs, "picture_cropping")

test("Picture Cropping (crop_left/right/top/bottom)", _test_picture_cropping)


# ═══════════════════════════════════════════════════════════════════════════════
# SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
print("\n" + "=" * 60)
print(f"RESULTS: {len(PASS)} passed, {len(FAIL)} failed")
print("=" * 60)

if FAIL:
    print("\nFailed tests:")
    for name, err in FAIL:
        print(f"  - {name}: {err}")
    sys.exit(1)
else:
    print("\nAll capabilities verified! Generated PPTX files saved to:")
    print(f"  {Path(tempfile.gettempdir()) / 'ppt_pro_max_test'}")
