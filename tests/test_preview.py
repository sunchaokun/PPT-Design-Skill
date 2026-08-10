"""Tests for build_helpers.preview() — one-call pptx preview helper.

Ensures build.py can render a layout preview via `from build_helpers import *`
without hand-rolled COM/CLI plumbing. Render backends are mocked.
"""
from __future__ import annotations

from unittest.mock import patch

import pytest


@pytest.fixture
def sample_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "preview test slide"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def test_preview_importable_via_star_import():
    ns = {}
    exec("from ppt_pro_max.build_helpers import *", ns)
    assert "preview" in ns
    assert callable(ns["preview"])


def test_preview_passes_engine_and_outdir(sample_pptx, tmp_path):
    from ppt_pro_max.build_helpers import preview

    captured = {}

    def _fake_render_preview(*args, **kwargs):
        captured["args"] = args
        captured["kwargs"] = kwargs
        return {"pngs": [tmp_path / "slide1.png"], "html": tmp_path / "index.html",
                "engine": "libreoffice", "warnings": []}

    out = tmp_path / "pv"
    with patch("ppt_pro_max.render_preview.render_preview", side_effect=_fake_render_preview):
        result = preview(str(sample_pptx), out_dir=str(out), engine="libreoffice")

    assert result["engine"] == "libreoffice"
    assert captured["args"][0] == str(sample_pptx)
    assert captured["kwargs"]["out_dir"] == str(out)
    assert captured["kwargs"]["engine"] == "libreoffice"


def test_preview_auto_engine_default(sample_pptx):
    from ppt_pro_max.build_helpers import preview

    captured = {}

    def _fake_render_preview(*args, **kwargs):
        captured["kwargs"] = kwargs
        return {"pngs": [], "html": "x", "engine": "powerpoint", "warnings": []}

    with patch("ppt_pro_max.render_preview.render_preview", side_effect=_fake_render_preview):
        preview(str(sample_pptx))
    # engine defaults to None → auto fallback chain
    assert captured["kwargs"].get("engine") is None


def test_preview_open_in_browser(sample_pptx):
    from ppt_pro_max.build_helpers import preview

    with patch("ppt_pro_max.render_preview.render_preview",
               return_value={"pngs": [], "html": r"C:\tmp\index.html",
                             "engine": "libreoffice", "warnings": []}), \
         patch("webbrowser.open") as mock_open:
        preview(str(sample_pptx), open_in_browser=True)
    mock_open.assert_called_once()
