"""Tests for image_processor.py extensions and blip effects.

Covers:
  1. Pillow filters: grayscale, sepia, duotone, ink_wash, blur, vignette, edge_fade
  2. OOXML blip effects: duotone, grayscale, brightness/contrast, saturation, artistic
"""

from __future__ import annotations

import os
import tempfile

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_pro_max.renderer.image_processor import (
    apply_blur,
    apply_duotone,
    apply_edge_fade,
    apply_grayscale,
    apply_ink_wash,
    apply_sepia,
    apply_vignette,
)
from ppt_pro_max.renderer.blip_fill import (
    add_image_in_shape,
    apply_blip_artistic,
    apply_blip_brightness_contrast,
    apply_blip_duotone,
    apply_blip_grayscale,
    apply_blip_saturation,
)


@pytest.fixture
def test_image():
    img = PILImage.new("RGB", (200, 200), (255, 100, 50))
    path = os.path.join(tempfile.gettempdir(), "imgproc_test.png")
    img.save(path, "PNG")
    return path


@pytest.fixture
def test_image_rgba():
    img = PILImage.new("RGBA", (200, 200), (255, 100, 50, 255))
    path = os.path.join(tempfile.gettempdir(), "imgproc_test_rgba.png")
    img.save(path, "PNG")
    return path


def _is_valid_image(path):
    assert os.path.isfile(path)
    with PILImage.open(path) as img:
        img.verify()
    return True


# ── Pillow filters ──


class TestApplyGrayscale:
    def test_produces_grayscale_image(self, test_image):
        out = apply_grayscale(test_image)
        assert _is_valid_image(out)
        with PILImage.open(out) as img:
            assert img.mode in ("L", "RGB")
            if img.mode == "RGB":
                r, g, b = img.getpixel((10, 10))
                assert r == g == b, "Grayscale should have equal R=G=B"

    def test_output_is_different(self, test_image):
        out = apply_grayscale(test_image)
        assert out != test_image


class TestApplySepia:
    def test_produces_sepia_image(self, test_image):
        out = apply_sepia(test_image, intensity=0.5)
        assert _is_valid_image(out)

    def test_intensity_zero_returns_original(self, test_image):
        out = apply_sepia(test_image, intensity=0.0)
        assert out == test_image


class TestApplyDuotone:
    def test_produces_duotone_image(self, test_image):
        out = apply_duotone(test_image, "#0000FF", "#FF0000")
        assert _is_valid_image(out)

    def test_output_is_different(self, test_image):
        out = apply_duotone(test_image, "#000000", "#FFFFFF")
        assert out != test_image


class TestApplyInkWash:
    def test_produces_ink_wash_image(self, test_image):
        out = apply_ink_wash(test_image)
        assert _is_valid_image(out)


class TestApplyBlur:
    def test_produces_blurred_image(self, test_image):
        out = apply_blur(test_image, radius=5)
        assert _is_valid_image(out)

    def test_radius_zero_returns_original(self, test_image):
        out = apply_blur(test_image, radius=0)
        assert out == test_image


class TestApplyVignette:
    def test_produces_vignette_image(self, test_image):
        out = apply_vignette(test_image, intensity=0.5)
        assert _is_valid_image(out)

    def test_intensity_zero_returns_original(self, test_image):
        out = apply_vignette(test_image, intensity=0.0)
        assert out == test_image


class TestApplyEdgeFade:
    def test_produces_edge_fade_image(self, test_image):
        out = apply_edge_fade(test_image, margin_pct=0.1, bg_color="#FFFFFF")
        assert _is_valid_image(out)

    def test_rgba_output(self, test_image):
        out = apply_edge_fade(test_image, margin_pct=0.1)
        assert _is_valid_image(out)
        with PILImage.open(out) as img:
            assert img.mode == "RGBA"

    def test_zero_margin(self, test_image):
        out = apply_edge_fade(test_image, margin_pct=0.0)
        assert _is_valid_image(out)


# ── OOXML blip effects ──


def _make_image_shape(test_image):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    shape = add_image_in_shape(slide, MSO_SHAPE.RECTANGLE, 1, 1, 4, 3, test_image)
    return shape


class TestApplyBlipGrayscale:
    def test_adds_grayscl_element(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_grayscale(shape)
        spPr = shape._element.find(qn("p:spPr"))
        blipFill = spPr.find(qn("a:blipFill"))
        blip = blipFill.find(qn("a:blip"))
        assert blip.find(qn("a:grayscl")) is not None


class TestApplyBlipDuotone:
    def test_adds_duotone_element(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_duotone(shape, "#0000FF", "#FF0000")
        spPr = shape._element.find(qn("p:spPr"))
        blipFill = spPr.find(qn("a:blipFill"))
        blip = blipFill.find(qn("a:blip"))
        duotone = blip.find(qn("a:duotone"))
        assert duotone is not None
        colors = duotone.findall(qn("a:srgbClr"))
        assert len(colors) == 2


class TestApplyBlipBrightnessContrast:
    def test_adds_lum_element(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_brightness_contrast(shape, bright_pct=20, contrast_pct=10)
        spPr = shape._element.find(qn("p:spPr"))
        blipFill = spPr.find(qn("a:blipFill"))
        blip = blipFill.find(qn("a:blip"))
        lum = blip.find(qn("a:lum"))
        assert lum is not None
        assert lum.get("bright") is not None
        assert lum.get("contrast") is not None

    def test_default_values(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_brightness_contrast(shape)
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        lum = blip.find(qn("a:lum"))
        assert lum.get("bright") == "0"
        assert lum.get("contrast") == "0"


class TestApplyBlipSaturation:
    def test_adds_sat_element(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_saturation(shape, saturation_pct=50)
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        sat = blip.find(qn("a:sat"))
        assert sat is not None

    def test_zero_saturation(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_saturation(shape, saturation_pct=0)
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        sat = blip.find(qn("a:sat"))
        assert sat.get("val") == "0"


class TestApplyBlipArtistic:
    def test_watercolor_sponge(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_artistic(shape, "watercolor_sponge")
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        effect = blip.find(qn("a:artisticWatercolorSponge"))
        assert effect is not None

    def test_invalid_effect_raises(self, test_image):
        shape = _make_image_shape(test_image)
        with pytest.raises(KeyError, match="Unknown artistic effect"):
            apply_blip_artistic(shape, "nonexistent_effect")

    def test_pencil_grayscale(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_artistic(shape, "pencil_grayscale")
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        assert blip.find(qn("a:artisticPencilGrayscale")) is not None

    def test_blur_effect(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_artistic(shape, "blur")
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        assert blip.find(qn("a:artisticBlur")) is not None

    def test_custom_params(self, test_image):
        shape = _make_image_shape(test_image)
        apply_blip_artistic(shape, "watercolor_sponge",
                            params={"brushSize": "5"})
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        effect = blip.find(qn("a:artisticWatercolorSponge"))
        assert effect is not None
        assert effect.get("brushSize") == "5"


class TestBlipEffectOnNoBlipFill:
    def test_grayscale_on_plain_shape(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        apply_blip_grayscale(shape)

    def test_duotone_on_plain_shape(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        apply_blip_duotone(shape, "#FF0000", "#0000FF")
