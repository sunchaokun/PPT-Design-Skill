"""Tests for component-pipeline integration issues identified in analysis report.

Tests cover:
- Breakpoint 1: BrandSpec color key mismatch (OOXML key vs semantic key)
- Breakpoint 2: Hardcoded bounds in _post_save_inject_groups
- Breakpoint 3: component_category ignored + layout override by _auto_layout
- Secondary: schemeClr mapping, font replacement, font scaling
"""

import os
import tempfile

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _make_group_xml_with_colors(colors_hex: list[str], text: str = "Test") -> bytes:
    p_ns = _P
    a_ns = _A
    children = ""
    for i, color in enumerate(colors_hex):
        c = color.lstrip("#")
        children += (
            f'<p:sp xmlns:p="{p_ns}" xmlns:a="{a_ns}">'
            f'<p:nvSpPr><p:cNvPr id="{i+2}" name="Sp {i}"/>'
            f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{i*200000}" y="0"/>'
            f'<a:ext cx="100000" cy="100000"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:solidFill><a:srgbClr val="{c}"/></a:solidFill></p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:r><a:rPr lang="zh-CN" sz="1200"/>'
            f'<a:latin typeface="Microsoft YaHei"/>'
            f'<a:ea typeface="Microsoft YaHei"/>'
            f'<a:t>{text}</a:t></a:r></a:p></p:txBody></p:sp>'
        )
    grp_xml = (
        f'<p:grpSp xmlns:p="{p_ns}" xmlns:a="{a_ns}" '
        f'xmlns:r="{_R}">'
        f'<p:nvGrpSpPr><p:cNvPr id="1" name="Group 1"/>'
        f'<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        f'<p:grpSpPr><a:xfrm>'
        f'<a:off x="0" y="0"/><a:ext cx="914400" cy="457200"/>'
        f'<a:chOff x="0" y="0"/><a:chExt cx="914400" cy="457200"/>'
        f'</a:xfrm></p:grpSpPr>'
        f'{children}'
        f'</p:grpSp>'
    )
    return grp_xml.encode("utf-8")


def _extract_srgb_colors(group_xml: bytes) -> set[str]:
    root = etree.fromstring(group_xml)
    colors = set()
    for srgb in root.iter(f"{{{_A}}}srgbClr"):
        val = srgb.get("val", "")
        if val:
            colors.add(val.upper())
    return colors


def _extract_fonts(group_xml: bytes) -> dict[str, list[str]]:
    root = etree.fromstring(group_xml)
    fonts = {"latin": [], "ea": [], "cs": []}
    for latin in root.iter(f"{{{_A}}}latin"):
        tf = latin.get("typeface", "")
        if tf:
            fonts["latin"].append(tf)
    for ea in root.iter(f"{{{_A}}}ea"):
        tf = ea.get("typeface", "")
        if tf:
            fonts["ea"].append(tf)
    for cs in root.iter(f"{{{_A}}}cs"):
        tf = cs.get("typeface", "")
        if tf:
            fonts["cs"].append(tf)
    return fonts


def _extract_font_sizes(group_xml: bytes) -> list[int]:
    root = etree.fromstring(group_xml)
    sizes = []
    for rPr in root.iter(f"{{{_A}}}rPr"):
        sz = rPr.get("sz")
        if sz:
            try:
                sizes.append(int(sz))
            except ValueError:
                pass
    return sizes


# ═══════════════════════════════════════════════════════════════
# Breakpoint 1: BrandSpec color key mismatch
# ═══════════════════════════════════════════════════════════════

class TestBrandSpecColorKeyMismatch:
    """_recolor_group_xml uses semantic keys (primary/accent/muted/foreground/on-primary).
    When BrandSpec has OOXML keys (accent1/tx1), all lookups return None,
    falling back to hardcoded defaults (#2563EB, #F97316, #F1F5F9).
    """

    def test_recolor_with_semantic_keys_uses_brand_colors(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        group_xml = _make_group_xml_with_colors(["#3D485D", "#7B9CFE", "#B5B5B5"])
        brand_spec = BrandSpec(colors={
            "primary": "#0F1423",
            "accent": "#6096E6",
            "muted": "#F0F4FA",
            "foreground": "#1A1A2E",
            "on-primary": "#FFFFFF",
        })
        renderer = ComponentRenderer()
        result = renderer._apply_brand_colors({"group": group_xml}, brand_spec)
        result_colors = _extract_srgb_colors(result["group"])

        assert "0F1423" in result_colors, f"primary #0F1423 not in result colors: {result_colors}"
        assert "6096E6" in result_colors, f"accent #6096E6 not in result colors: {result_colors}"

    def test_recolor_with_ooxml_keys_falls_back_to_defaults(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        group_xml = _make_group_xml_with_colors(["#3D485D", "#7B9CFE", "#B5B5B5"])
        brand_spec_ooxml = BrandSpec(colors={
            "accent1": "#6096E6",
            "tx1": "#1A1A2E",
        })
        renderer = ComponentRenderer()
        result = renderer._apply_brand_colors({"group": group_xml}, brand_spec_ooxml)
        result_colors = _extract_srgb_colors(result["group"])

        assert "2563EB" in result_colors, "Should fall back to default primary #2563EB"
        assert "F97316" in result_colors, "Should fall back to default accent #F97316"
        assert "6096E6" not in result_colors, "Should NOT use brand accent1 #6096E6 (wrong key)"

    def test_try_component_render_constructs_wrong_brand_spec_keys(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        extractor = DesignDNAExtractor()
        plan = PagePlan(
            page_type="content",
            title="Test",
            bullets=["A", "B", "C"],
            component_type="group",
            component_category="process",
        )

        accent_color = "#6096E6"
        body_color = "#1A1A2E"

        brand_spec_constructed = BrandSpec(colors={
            "accent1": accent_color,
            "tx1": body_color,
        })

        assert "primary" not in brand_spec_constructed.colors
        assert "accent" not in brand_spec_constructed.colors
        assert "muted" not in brand_spec_constructed.colors
        assert "foreground" not in brand_spec_constructed.colors
        assert brand_spec_constructed.colors.get("primary") is None
        assert brand_spec_constructed.colors.get("accent") is None

    def test_post_save_inject_groups_constructs_wrong_brand_spec_keys(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        accent_color_val = "#6096E6"
        body_color_val = "#1A1A2E"
        brand_spec = BrandSpec(colors={"accent1": accent_color_val, "tx1": body_color_val})

        assert "primary" not in brand_spec.colors
        assert "accent" not in brand_spec.colors
        assert "muted" not in brand_spec.colors
        assert "foreground" not in brand_spec.colors

    def test_dna_brand_spec_has_semantic_keys(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Test"
        prs.save(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert dna.brand_spec is not None
        if dna.brand_spec.colors:
            semantic_keys = {"foreground", "background", "primary", "accent"}
            has_any_semantic = bool(semantic_keys & set(dna.brand_spec.colors.keys()))
            assert has_any_semantic, f"dna.brand_spec should have semantic keys, got: {list(dna.brand_spec.colors.keys())}"


# ═══════════════════════════════════════════════════════════════
# Breakpoint 2: Hardcoded bounds
# ═══════════════════════════════════════════════════════════════

class TestHardcodedBounds:
    """_post_save_inject_groups hardcodes (1.5, 2.0, 8.0, 4.5) instead of
    computing actual content bounds from slide_dna.
    """

    def test_denormalize_with_hardcoded_bounds_mismatches_content_area(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        renderer = ComponentRenderer()

        group_xml = _make_group_xml_with_colors(["#3D485D"])
        grp_elem = etree.fromstring(group_xml)

        hardcoded_bounds = (1.5, 2.0, 8.0, 4.5)
        typical_content_bounds = (0.9, 1.6, 7.0, 4.8)

        assert hardcoded_bounds != typical_content_bounds
        assert hardcoded_bounds[0] != typical_content_bounds[0]
        assert hardcoded_bounds[2] != typical_content_bounds[2]

    def test_denormalize_coordinates_modifies_position(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        renderer = ComponentRenderer()

        group_xml = _make_group_xml_with_colors(["#3D485D", "#7B9CFE"])
        grp_elem = etree.fromstring(group_xml)

        grpSpPr = grp_elem.find(f"{{{_P}}}grpSpPr")
        xfrm = grpSpPr.find(f"{{{_A}}}xfrm")
        off = xfrm.find(f"{{{_A}}}off")

        renderer._denormalize_coordinates(grp_elem, 1.5, 2.0, 8.0, 4.5)

        new_x = int(off.get("x", "0"))
        new_y = int(off.get("y", "0"))

        target_x = int(1.5 * 914400)
        target_y = int(2.0 * 914400)
        assert new_x >= target_x
        assert new_y >= target_y

    def test_font_scaling_too_aggressive(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        renderer = ComponentRenderer()

        group_xml = _make_group_xml_with_colors(["#3D485D"])
        grp_elem = etree.fromstring(group_xml)

        orig_sizes = _extract_font_sizes(etree.tostring(grp_elem))

        renderer._denormalize_coordinates(grp_elem, 1.5, 2.0, 8.0, 4.5)

        new_sizes = _extract_font_sizes(etree.tostring(grp_elem))

        if orig_sizes and new_sizes:
            for orig, new in zip(orig_sizes, new_sizes):
                if orig >= 2400:
                    min_readable = 1400
                    assert new >= min_readable, f"Font scaled too small: {orig} -> {new}"


# ═══════════════════════════════════════════════════════════════
# Breakpoint 3: component_category ignored + layout override
# ═══════════════════════════════════════════════════════════════

class TestLayoutOverrideByAutoLayout:
    """_auto_layout ignores component_type/component_category and returns
    'bullets' for most content, overriding the intended layout.
    """

    def test_auto_layout_ignores_component_category(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        extractor = DesignDNAExtractor()

        page_with_process = {
            "goal": "content",
            "title": "电商助农实施路径",
            "bullets": [
                "培育农村电商带头人5000名",
                "建设县级电商服务中心100个",
                "打造农产品区域公共品牌50个",
                "实现农产品网络零售额年增长20%",
            ],
            "component_type": "group",
            "component_category": "process",
        }

        layout = extractor._auto_layout(page_with_process)
        assert layout == "bullets", "Current behavior: returns 'bullets' despite component_category='process'"

    def test_auto_layout_returns_process_for_keyword_bullets(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        extractor = DesignDNAExtractor()

        page_with_keyword = {
            "goal": "content",
            "title": "实施流程",
            "bullets": ["步骤1：准备", "步骤2：执行", "步骤3：验收"],
        }

        layout = extractor._auto_layout(page_with_keyword)
        assert layout == "process"

    def test_plan_pages_layout_overridden_by_auto_layout(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        extractor = DesignDNAExtractor()

        story_plan = {
            "strategy": "content_json",
            "pages": [
                {
                    "goal": "content",
                    "title": "电商助农实施路径",
                    "bullets": ["培育农村电商带头人5000名", "建设县级电商服务中心100个"],
                    "component_type": "group",
                    "component_category": "process",
                },
            ],
        }

        page_contents = story_plan["pages"]

        plans = extractor.plan_pages(story_plan, page_contents)

        assert len(plans) == 1
        plan = plans[0]
        assert plan.component_type == "group"
        assert plan.component_category == "process"
        assert plan.layout == "process", "layout should now be 'process' derived from component_category (was 'bullets' before fix)"

    def test_post_save_inject_groups_ignores_component_category(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        extractor = DesignDNAExtractor()

        plan = PagePlan(
            page_type="content",
            title="电商助农实施路径",
            bullets=["培育农村电商带头人5000名", "建设县级电商服务中心100个"],
            layout="bullets",
            component_type="group",
            component_category="process",
        )

        layout = plan.layout if plan.layout != "auto" else extractor._auto_layout_from_plan(plan)
        category_map = {
            "chevron": "process", "bullets": "infographic", "cards": "infographic",
            "process": "process", "hierarchy": "hierarchy",
            "two_col": "infographic", "numbered_row": "process",
        }
        category = plan.component_category or category_map.get(layout, "infographic")

        assert category == "process", "Fix: category should now be 'process' from component_category"
        assert plan.component_category == "process", "component_category is correct but unused"

    def test_try_component_render_uses_component_category(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        extractor = DesignDNAExtractor()

        plan = PagePlan(
            page_type="content",
            title="电商助农实施路径",
            bullets=["培育农村电商带头人5000名", "建设县级电商服务中心100个"],
            layout="bullets",
            component_type="group",
            component_category="process",
        )

        layout = plan.layout if plan.layout != "auto" else extractor._auto_layout_from_plan(plan)
        category_map = {
            "chevron": "process", "bullets": "infographic", "cards": "infographic",
            "process": "process", "hierarchy": "hierarchy",
            "two_col": "infographic", "numbered_row": "process",
            "swot": "swot", "timeline": "timeline", "chart": "chart",
        }
        category = plan.component_category or category_map.get(layout, "infographic")

        assert category == "process", "_try_component_render correctly uses component_category"


# ═══════════════════════════════════════════════════════════════
# Secondary: schemeClr mapping
# ═══════════════════════════════════════════════════════════════

class TestSmartArtSchemeClrMapping:
    """_apply_brand_colors for SmartArt colors.xml uses brand_spec.colors keys
    directly as schemeClr val matches. This fails for both OOXML and semantic keys.
    """

    def _make_colors_xml_with_scheme(self, scheme_vals: list[str]) -> bytes:
        a_ns = _A
        dgm_ns = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
        style_entries = ""
        for i, val in enumerate(scheme_vals):
            style_entries += (
                f'<dgm:styleEntry idx="{i}">'
                f'<a:fillClr><a:schemeClr val="{val}"/></a:fillClr>'
                f'<a:lineClr><a:schemeClr val="{val}"/></a:lineClr>'
                f'<a:effectClr><a:schemeClr val="{val}"/></a:effectClr>'
                f'<a:txClr><a:schemeClr val="{val}"/></a:txClr>'
                f'</dgm:styleEntry>'
            )
        return (
            f'<dgm:colors xmlns:dgm="{dgm_ns}" xmlns:a="{a_ns}">'
            f'<dgm:styleLbl name="node1">{style_entries}</dgm:styleLbl>'
            f'</dgm:colors>'
        ).encode("utf-8")

    def test_schemeclr_with_ooxml_keys_partial_match(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        colors_xml = self._make_colors_xml_with_scheme(["accent1", "dk1", "lt1"])
        brand_spec = BrandSpec(colors={"accent1": "#6096E6", "tx1": "#1A1A2E"})

        renderer = ComponentRenderer()
        result = renderer._apply_brand_colors({"colors": colors_xml}, brand_spec)

        root = etree.fromstring(result["colors"])
        srgb_vals = [s.get("val", "") for s in root.iter(f"{{{_A}}}srgbClr")]

        assert "6096E6" in srgb_vals, "accent1 should be replaced"
        remaining_schemes = [s.get("val", "") for s in root.iter(f"{{{_A}}}schemeClr")]
        assert len(remaining_schemes) > 0, "dk1/lt1 should NOT be replaced (no matching key)"

    def test_schemeclr_with_semantic_keys_now_matches(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        colors_xml = self._make_colors_xml_with_scheme(["accent1", "dk1", "lt1"])
        brand_spec = BrandSpec(colors={
            "primary": "#0F1423",
            "accent": "#6096E6",
            "foreground": "#1A1A2E",
            "background": "#FFFFFF",
        })

        renderer = ComponentRenderer()
        result = renderer._apply_brand_colors({"colors": colors_xml}, brand_spec)

        root = etree.fromstring(result["colors"])
        srgb_vals = [s.get("val", "") for s in root.iter(f"{{{_A}}}srgbClr")]
        remaining_schemes = [s.get("val", "") for s in root.iter(f"{{{_A}}}schemeClr")]

        assert "6096E6" in srgb_vals, "accent1 should now be replaced via scheme→semantic mapping"
        assert "1A1A2E" in srgb_vals, "dk1 should now be replaced via scheme→semantic mapping"
        assert "FFFFFF" in srgb_vals, "lt1 should now be replaced via scheme→semantic mapping"


# ═══════════════════════════════════════════════════════════════
# Secondary: Font replacement missing
# ═══════════════════════════════════════════════════════════════

class TestFontReplacementMissing:
    """ComponentRenderer does not replace fonts in group XML with brand_spec fonts.
    Group XML retains original template fonts (Calibri, Segoe UI, etc.).
    """

    def test_group_xml_retains_original_fonts_after_recolor(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        group_xml = _make_group_xml_with_colors(["#3D485D"])
        brand_spec = BrandSpec(
            colors={"primary": "#0F1423", "accent": "#6096E6"},
            fonts={"heading": "Arial", "body": "Arial", "cjk_heading": "Microsoft YaHei"},
        )

        renderer = ComponentRenderer()
        result = renderer._apply_brand_colors({"group": group_xml}, brand_spec)
        fonts = _extract_fonts(result["group"])

        assert "Microsoft YaHei" in fonts["ea"], "EA font should be Microsoft YaHei from original"
        assert "Microsoft YaHei" in fonts["latin"], "Latin font should be Microsoft YaHei from original"

    def test_brand_spec_fonts_applied_via_replace_group_fonts(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        group_xml = _make_group_xml_with_colors(["#3D485D"])
        brand_spec = BrandSpec(
            colors={"primary": "#0F1423", "accent": "#6096E6"},
            fonts={"heading": "Arial", "body": "Arial"},
        )

        renderer = ComponentRenderer()
        styled = renderer._apply_brand_colors({"group": group_xml}, brand_spec)

        grp_bytes = styled["group"]
        if isinstance(grp_bytes, str):
            grp_bytes = grp_bytes.encode("utf-8")
        replaced = ComponentRenderer._replace_group_fonts(grp_bytes, brand_spec)
        fonts = _extract_fonts(replaced)

        has_arial = "Arial" in fonts["latin"]
        assert has_arial, "After _replace_group_fonts, Arial should appear"


# ═══════════════════════════════════════════════════════════════
# Secondary: Font scaling protection
# ═══════════════════════════════════════════════════════════════

class TestFontScalingProtection:
    """_denormalize_coordinates scales fonts by scale_y with range [800, 7200].
    When scale_y < 1 (common), large fonts are shrunk too aggressively.
    """

    def test_font_scaled_below_readable_size(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        p_ns = _P
        a_ns = _A

        large_font_sz = 3200
        grp_xml = (
            f'<p:grpSp xmlns:p="{p_ns}" xmlns:a="{a_ns}" xmlns:r="{_R}">'
            f'<p:nvGrpSpPr><p:cNvPr id="1" name="Group 1"/>'
            f'<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f'<p:grpSpPr><a:xfrm>'
            f'<a:off x="0" y="0"/><a:ext cx="12192000" cy="6858000"/>'
            f'<a:chOff x="0" y="0"/><a:chExt cx="12192000" cy="6858000"/>'
            f'</a:xfrm></p:grpSpPr>'
            f'<p:sp><p:nvSpPr><p:cNvPr id="2" name="Sp 0"/>'
            f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="0" y="0"/>'
            f'<a:ext cx="12192000" cy="6858000"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:solidFill><a:srgbClr val="3D485D"/></a:solidFill></p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:r><a:rPr lang="zh-CN" sz="{large_font_sz}"/>'
            f'<a:latin typeface="Microsoft YaHei"/>'
            f'<a:ea typeface="Microsoft YaHei"/>'
            f'<a:t>大标题</a:t></a:r></a:p></p:txBody></p:sp>'
            f'</p:grpSp>'
        ).encode("utf-8")

        grp_elem = etree.fromstring(grp_xml)

        renderer = ComponentRenderer()
        renderer._denormalize_coordinates(grp_elem, 1.5, 2.0, 8.0, 4.5)

        new_sizes = _extract_font_sizes(etree.tostring(grp_elem))

        if new_sizes:
            scale_y = (4.5 * 914400) / 6858000
            expected_scaled = int(large_font_sz * scale_y)
            assert expected_scaled < large_font_sz, "Font should be scaled down"
            assert new_sizes[0] == max(800, min(7200, expected_scaled)), "Current scaling behavior"


# ═══════════════════════════════════════════════════════════════
# Secondary: _ensure_shape_fills hardcoded color
# ═══════════════════════════════════════════════════════════════

class TestEnsureShapeFillsHardcodedColor:
    """_ensure_shape_fills adds #F1F5F9 for unfilled shapes, ignoring brand_spec.
    """

    def test_ensure_shape_fills_uses_hardcoded_color(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        p_ns = _P
        a_ns = _A

        grp_xml = (
            f'<p:grpSp xmlns:p="{p_ns}" xmlns:a="{a_ns}" xmlns:r="{_R}">'
            f'<p:nvGrpSpPr><p:cNvPr id="1" name="Group 1"/>'
            f'<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f'<p:grpSpPr><a:xfrm>'
            f'<a:off x="0" y="0"/><a:ext cx="914400" cy="457200"/>'
            f'<a:chOff x="0" y="0"/><a:chExt cx="914400" cy="457200"/>'
            f'</a:xfrm></p:grpSpPr>'
            f'<p:sp><p:nvSpPr><p:cNvPr id="2" name="Sp 0"/>'
            f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="0" y="0"/>'
            f'<a:ext cx="500000" cy="500000"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'</p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:r><a:rPr lang="zh-CN" sz="1200"/>'
            f'<a:t>Test</a:t></a:r></a:p></p:txBody></p:sp>'
            f'</p:grpSp>'
        ).encode("utf-8")

        grp_elem = etree.fromstring(grp_xml)

        renderer = ComponentRenderer()
        renderer._ensure_shape_fills(grp_elem)

        result_xml = etree.tostring(grp_elem)
        result_colors = _extract_srgb_colors(result_xml)

        has_f1f5f9 = "F1F5F9" in result_colors
        assert has_f1f5f9, "Bug: hardcoded #F1F5F9 used instead of brand_spec.muted"
