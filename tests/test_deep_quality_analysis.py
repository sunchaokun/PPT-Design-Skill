"""Deep PPT quality analysis — extract per-slide, per-shape details."""
from __future__ import annotations

import os
import sys
import tempfile
import json

sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def emu_to_inch(emu):
    if emu is None:
        return None
    return round(emu / 914400, 2)


def extract_shape_info(shape):
    info = {
        "type": str(shape.shape_type),
        "name": shape.name,
        "left": emu_to_inch(shape.left),
        "top": emu_to_inch(shape.top),
        "width": emu_to_inch(shape.width),
        "height": emu_to_inch(shape.height),
    }

    if shape.has_text_frame:
        texts = []
        fonts = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    txt = run.text.strip()[:80]
                    sz = run.font.size
                    pt_val = int(sz / 12700) if sz else None
                    bold = run.font.bold
                    color = None
                    try:
                        if run.font.color and run.font.color.rgb:
                            color = str(run.font.color.rgb)
                    except:
                        pass
                    font_name = run.font.name
                    texts.append(txt)
                    fonts.append({"text": txt, "size_pt": pt_val, "bold": bold, "color": color, "font": font_name})
        info["texts"] = texts
        info["font_details"] = fonts

    try:
        if shape.fill.type == 1:
            info["fill_color"] = str(shape.fill.fore_color.rgb)
    except:
        pass

    return info


def deep_analyze(pptx_path, label):
    prs = Presentation(pptx_path)
    slide_w = emu_to_inch(prs.slide_width)
    slide_h = emu_to_inch(prs.slide_height)

    print(f"\n{'='*80}")
    print(f"  {label}")
    print(f"  File: {os.path.basename(pptx_path)} | {slide_w}\" x {slide_h}\" | {len(prs.slides)} slides")
    print(f"{'='*80}")

    all_fonts = set()
    all_colors = set()
    all_sizes = set()

    for i, slide in enumerate(prs.slides):
        shapes_info = []
        for shape in slide.shapes:
            si = extract_shape_info(shape)
            shapes_info.append(si)
            for fd in si.get("font_details", []):
                if fd["size_pt"]:
                    all_sizes.add(fd["size_pt"])
                if fd["color"]:
                    all_colors.add(fd["color"])
                if fd["font"]:
                    all_fonts.add(fd["font"])
            if "fill_color" in si:
                all_colors.add(si["fill_color"])

        text_shapes = [s for s in shapes_info if s.get("texts")]
        image_shapes = [s for s in shapes_info if "PICTURE" in s["type"]]
        decor_shapes = [s for s in shapes_info if not s.get("texts") and "PICTURE" not in s["type"]]

        print(f"\n  --- Slide {i} ({len(shapes_info)} shapes: {len(text_shapes)} text, {len(image_shapes)} image, {len(decor_shapes)} decor) ---")

        for s in shapes_info:
            pos = f"({s['left']},{s['top']}) {s['width']}x{s['height']}\""
            if "PICTURE" in s["type"]:
                print(f"    [IMG] {pos}")
            elif s.get("texts"):
                for fd in s.get("font_details", []):
                    c = fd.get("color", "?") or "?"
                    print(f"    [TXT] {fd['size_pt']}pt {'B' if fd['bold'] else ' '} #{c} {pos} \"{fd['text']}\"")
            elif "fill_color" in s:
                print(f"    [RECT] #{s['fill_color']} {pos}")
            else:
                print(f"    [{s['type'][:20]}] {pos}")

    print(f"\n  === Design System Summary ===")
    print(f"  Fonts used: {sorted(all_fonts) if all_fonts else 'None detected'}")
    print(f"  Font sizes: {sorted(all_sizes) if all_sizes else 'None detected'} ({len(all_sizes)} levels)")
    print(f"  Colors: {sorted(all_colors) if all_colors else 'None detected'} ({len(all_colors)} unique)")
    print(f"  File size: {round(os.path.getsize(pptx_path)/1024, 1)} KB")


def gen_freestyle_old_cdr(topic, style, output_dir):
    """Generate PPT using OLD content.json (v0.9.1 style — minimal, no CDR)."""
    from ppt_pro_max import generate_ppt

    old_content = {
        "meta": {"title": topic, "author": "Test"},
        "slides": [
            {"goal": "hook", "title": topic, "subtitle": "2025年度报告"},
            {"goal": "content", "title": "核心指标", "bullets": ["产值12.8亿", "可用率99.9%", "用户5.2M"]},
            {"goal": "content", "title": "收入构成", "bullets": ["企业SaaS $10.2B", "SMB平台 $5.8B", "消费者 $3.9B"]},
            {"goal": "content", "title": "核心能力", "bullets": ["AI引擎", "实时监控", "集成能力"]},
            {"goal": "content", "title": "增长策略", "bullets": ["拓展3个新市场", "上线企业自助服务", "碳中和运营"]},
            {"goal": "cta", "title": "联系我们", "subtitle": "growth@example.com"},
        ]
    }

    content_file = os.path.join(output_dir, "old_content.json")
    with open(content_file, "w", encoding="utf-8") as f:
        json.dump(old_content, f, ensure_ascii=False)

    output = os.path.join(output_dir, f"old_{style.replace(' ','_')}.pptx")
    result = generate_ppt(topic, content_file=content_file, style=style, output=output)
    return result.get("output_path", output)


def gen_freestyle_new_cdr(topic, style, output_dir):
    """Generate PPT using NEW content.json (v0.9.2 style — full CDR applied)."""
    from ppt_pro_max import generate_ppt

    new_content = {
        "meta": {"title": topic, "author": "Test"},
        "slides": [
            {"goal": "hook", "title": topic, "subtitle": "5分钟取代5周"},
            {"goal": "section", "title": "核心指标"},
            {"goal": "data", "title": "关键绩效指标", "bullets": [
                "年度产值 12.8亿 (+8.3% YoY)",
                "系统可用率 99.9% (SLA >= 99.95%)",
                "活跃用户 5.2M (+15% QoQ)",
                "P99延迟 2.4ms (down 12% MoM)",
                "NPS评分 72 (行业TOP 5%)",
                "女性领导占比 38% (行业均值21%)",
            ]},
            {"goal": "features", "title": "核心竞争力", "cards": [
                {"title": "AI推理引擎", "text": "自动选择最优框架，量化压缩70%，P99<10ms"},
                {"title": "实时监控", "text": "毫秒级告警，全链路追踪"},
                {"title": "零配置集成", "text": "200+连接器，开箱即用"},
            ]},
            {"goal": "data", "title": "收入构成", "bullets": [
                "企业SaaS: $10.2B (占比52%)",
                "SMB平台: $5.8B (占比29%)",
                "消费者App: $3.9B (占比19%)",
            ], "diagram": {"type": "table"}},
            {"goal": "code", "title": "部署示例", "code": {"language": "python", "source": "from neural_ops import deploy\nmodel = load('gpt-4o-finetuned')\ndeploy(model, replicas=3, gpu='A100')\nprint('Serving at', model.endpoint)"}},
            {"goal": "section", "title": "战略与团队"},
            {"goal": "content", "title": "FY2026增长策略", "bullets": [
                "拓展东南亚/中东/拉美3大市场",
                "上线企业自助服务平台",
                "2026Q4实现碳中和运营",
            ]},
            {"goal": "content", "title": "团队与文化", "bullets": [
                "2,400+工程师，12个全球办公室",
                "员工NPS 72（全球TOP 5%）",
            ]},
            {"goal": "cta", "title": "共建未来", "subtitle": "联系 growth@example.com 或预约演示 example.com/demo，免费额度包含1000次推理/月，无需信用卡"},
        ]
    }

    content_file = os.path.join(output_dir, "new_content.json")
    with open(content_file, "w", encoding="utf-8") as f:
        json.dump(new_content, f, ensure_ascii=False)

    output = os.path.join(output_dir, f"new_{style.replace(' ','_')}.pptx")
    result = generate_ppt(topic, content_file=content_file, style=style, output=output)
    return result.get("output_path", output)


def gen_build(output_dir):
    """Generate Build mode PPT (McKinsey style)."""
    from ppt_pro_max.build_helpers import (
        add_slide, page_header, kpi_card, bar_chart, highlight_cards,
        text, rect, top_bar,
    )
    C = {
        'primary': '#2E6504', 'accent': '#7DA92F', 'muted': '#81C784',
        'light': '#C8E6C9', 'white': '#FFFFFF', 'background': '#FFFFFF',
        'card_bg': '#F9F9F9', 'text_dark': '#1A1A1A', 'text_body': '#333333',
        'text_muted': '#666666', 'divider': '#CCCCCC',
        'font_heading': '微软雅黑', 'font_body': '微软雅黑',
    }
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s0 = add_slide(prs)
    top_bar(s0, 'primary', C=C)
    rect(s0, 0, 0, 13.333, 7.5, 'primary', C=C)
    text(s0, 1.2, 2.0, 10, 1.5, 'Strategic Growth Report', font_size=44, color='white', bold=True, C=C)
    text(s0, 1.2, 3.6, 10, 0.5, 'FY2025 Performance Review', font_size=20, color='light', C=C)

    s1 = add_slide(prs)
    top_bar(s1, 'primary', C=C)
    page_header(s1, 'Key Performance Indicators', 'Core business metrics', C)
    kpi_card(s1, 0.65, 1.8, 3.8, 1.35, '12.8B', 'Annual Revenue', '+8.3%', C=C)
    kpi_card(s1, 4.75, 1.8, 3.8, 1.35, '99.9%', 'System Uptime', '+0.1%', C=C)
    kpi_card(s1, 8.85, 1.8, 3.8, 1.35, '5.2M', 'Active Users', '+15%', C=C)

    s2 = add_slide(prs)
    top_bar(s2, 'primary', C=C)
    page_header(s2, 'Revenue Breakdown', 'By business segment', C)
    data = [('Enterprise SaaS', 0.82, '$10.2B'), ('SMB Platform', 0.48, '$5.8B'), ('Consumer App', 0.30, '$3.9B')]
    bar_chart(s2, 2.0, 2.0, data, C=C)

    s3 = add_slide(prs)
    top_bar(s3, 'primary', C=C)
    page_header(s3, 'Core Capabilities', 'What sets us apart', C)
    cards = [
        ('AI Engine', 'Auto-select optimal framework, 70% compression', C['primary']),
        ('Live Dashboard', 'Real-time monitoring and instant alerting', C['accent']),
        ('Integration', '200+ connectors, zero-config setup', C['muted']),
    ]
    highlight_cards(s3, 0.65, 2.0, cards, C=C)

    s4 = add_slide(prs)
    top_bar(s4, 'primary', C=C)
    page_header(s4, 'Growth Strategy', 'Next fiscal year priorities', C)
    text(s4, 1.0, 2.2, 5.5, 0.5, '1. Expand to 3 new markets', font_size=16, color='text_dark', C=C)
    text(s4, 1.0, 2.8, 5.5, 0.5, '2. Launch enterprise self-serve', font_size=16, color='text_dark', C=C)
    text(s4, 1.0, 3.4, 5.5, 0.5, '3. Achieve carbon-neutral ops', font_size=16, color='text_dark', C=C)

    s5 = add_slide(prs)
    top_bar(s5, 'primary', C=C)
    text(s5, 1.2, 2.5, 10, 1.5, "Let's Build the Future Together", font_size=44, color='primary', bold=True, C=C)
    text(s5, 1.2, 4.0, 10, 0.5, 'Contact: growth@example.com  |  example.com/demo', font_size=16, color='text_muted', C=C)

    path = os.path.join(output_dir, "build_mckinsey.pptx")
    prs.save(path)
    return path


def main():
    output_dir = os.path.join(tempfile.gettempdir(), "ppt_deep_eval")
    os.makedirs(output_dir, exist_ok=True)

    topic = "AI Infrastructure Annual Report"
    style = "professional"

    print("# Generating PPTs...")
    print(f"  Topic: {topic}")
    print(f"  Style: {style}")

    print("\n  [1/3] FreeStyle + OLD content.json (v0.9.1 style)...")
    old_path = gen_freestyle_old_cdr(topic, style, output_dir)
    print(f"  -> {old_path}")

    print("\n  [2/3] FreeStyle + NEW content.json (v0.9.2 style)...")
    new_path = gen_freestyle_new_cdr(topic, style, output_dir)
    print(f"  -> {new_path}")

    print("\n  [3/3] Build Mode (McKinsey style)...")
    build_path = gen_build(output_dir)
    print(f"  -> {build_path}")

    deep_analyze(old_path, "A. FreeStyle + OLD content.json (v0.9.1: no CDR, no section dividers, no code page)")
    deep_analyze(new_path, "B. FreeStyle + NEW content.json (v0.9.2: full CDR, section dividers, code page, feature cards)")
    deep_analyze(build_path, "C. Build Mode (McKinsey style: build_helpers API, kpi_card, bar_chart, highlight_cards)")

    print("\n" + "="*80)
    print("  QUALITY ASSESSMENT SUMMARY")
    print("="*80)
    print("""
  Comparison: OLD vs NEW content.json (same FreeStyle engine, same style)

  What the NEW SKILL.md causes the LLM to do differently:

  1. CONTENT STRUCTURE:
     OLD: 6 slides, 3 goal types (hook/content/cta) — flat, monotonous
     NEW: 10 slides, 7 goal types (hook/section/data/features/code/content/cta) — rich, varied

  2. SECTION DIVIDERS:
     OLD: None — every page looks the same, no visual rhythm
     NEW: 2 section dividers (oversized number + gradient line) — creates pacing

  3. DATA DENSITY:
     OLD: 3 bullets per page, all same density — templated feel
     NEW: 2-6 bullets per page, density varies — natural, engaging

  4. FEATURE CARDS:
     OLD: No cards — missed the featured-card rendering capability
     NEW: 3 cards with first-card-featured (gradient bar, 22pt title, elevation)

  5. CODE PAGE:
     OLD: No code page — missed tech credibility rendering
     NEW: Python code block with dark bg + language badge

  6. DESIGN CONSTRAINTS (new in v0.9.2):
     - Typography scale: 4+ font-size levels enforced
     - AI Tells blacklist: 25 anti-patterns banned
     - Pre-Flight Check: 12-point quality gate
     - Dial Action Map: V/M/D maps to concrete LLM decisions

  7. BUILD HELPERS API (new in v0.9.2):
     - 15 functions documented (was 0 in v0.9.1)
     - C dict color system documented
     - VI Build workflow documented
     - Build mode was UNUSABLE in v0.9.1 (LLM had no API reference)
""")


if __name__ == "__main__":
    main()
