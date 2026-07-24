"""Build mode e2e: 3 styles with Typography/Spacing/GroupShape system."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt

from ppt_pro_max.build_helpers import (
    add_slide, page_header, kpi_card, bar_chart, comparison_bars,
    donut_chart, highlight_cards, text, multiline, rect, rrect, oval,
    top_bar, copy_decorations, copy_logo, code_block, section_divider,
    hero_slide, cta_slide,
    Typography, Spacing, TYPOGRAPHY, SPACING,
)


def build_mckinsey(output_path):
    C = {
        'primary': '#2E6504', 'accent': '#7DA92F', 'muted': '#81C784',
        'light': '#C8E6C9', 'white': '#FFFFFF', 'background': '#FFFFFF',
        'text_dark': '#1A1A1A', 'text_body': '#333333', 'text_muted': '#666666',
        'divider': '#CCCCCC', 'card_bg': '#F9F9F9', 'bg_tint': '#F5F5F5',
        'font_heading': 'Georgia', 'font_body': 'Calibri',
    }
    t = TYPOGRAPHY['mckinsey']
    sp = SPACING['mckinsey']

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = add_slide(prs)
    hero_slide(s, 'Strategic Growth Report', 'FY2025 Performance Review', C=C, typo=t)

    s = add_slide(prs)
    section_divider(s, 1, 'Core Metrics', C=C, typo=t)

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Key Performance Indicators', 'Core business metrics', C=C, typo=t, spacing=sp)
    kpi_card(s, 0.65, 1.8, 3.8, 1.35, '12.8B', 'Annual Revenue', '+8.3%', C=C, typo=t)
    kpi_card(s, 4.75, 1.8, 3.8, 1.35, '99.9%', 'System Uptime', '+0.1%', C=C, typo=t)
    kpi_card(s, 8.85, 1.8, 3.8, 1.35, '5.2M', 'Active Users', '+15%', C=C, typo=t)
    kpi_card(s, 0.65, 3.5, 3.8, 1.35, '2.4ms', 'P99 Latency', '-12%', trend_up=False, C=C, typo=t)
    kpi_card(s, 4.75, 3.5, 3.8, 1.35, '72', 'NPS Score', 'Top 5%', C=C, typo=t)
    kpi_card(s, 8.85, 3.5, 3.8, 1.35, '38%', 'Women in Leadership', 'vs 21% avg', C=C, typo=t)

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Revenue Breakdown', 'By business segment', C=C, typo=t, spacing=sp)
    bar_chart(s, 2.0, 2.0, [
        ('Enterprise SaaS', 0.82, '$10.2B'),
        ('SMB Platform', 0.48, '$5.8B'),
        ('Consumer App', 0.30, '$3.9B'),
    ], max_width=5.0, bar_height=0.3, C=C, typo=t, spacing=sp)
    donut_chart(s, 9.5, 4.0, 1.8, 1.0, [
        ('SaaS', '52%', C['primary']),
        ('SMB', '29%', C['accent']),
        ('Consumer', '19%', C['muted']),
    ], C=C, typo=t)

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Core Capabilities', 'What sets us apart', C=C, typo=t, spacing=sp)
    highlight_cards(s, 0.65, 2.0, [
        ('AI Inference Engine', 'Auto-select optimal framework, 70% compression, P99<10ms', C['primary']),
        ('Real-time Monitoring', 'Millisecond alerting, full-stack tracing, anomaly detection', C['accent']),
        ('Zero-config Integration', '200+ connectors, out-of-box ready, 5-min setup', C['muted']),
    ], total_width=12.0, C=C, typo=t, spacing=sp)

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Before vs After', 'Operational efficiency gains', C=C, typo=t, spacing=sp)
    comparison_bars(s, 2.5, 2.0, [
        ('Deploy Time', '4 hrs', '12 min', 0.8, 0.05),
        ('Incident Response', '45 min', '3 min', 0.9, 0.06),
        ('Cost/Inference', '$0.12', '$0.02', 0.75, 0.12),
    ], max_width=4.0, C=C, typo=t, spacing=sp)

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Deployment Example', 'Production-ready in 3 lines', C=C, typo=t, spacing=sp)
    code_block(s, 0.65, 2.0, 12.0, 3.5, [
        'from neural_ops import deploy',
        'model = load("gpt-4o-finetuned")',
        'deploy(model, replicas=3, gpu="A100")',
        'print("Serving at", model.endpoint)',
    ], language='python', C=C, typo=t)

    s = add_slide(prs)
    cta_slide(s, "Let's Build the Future Together",
              'Contact: growth@example.com  |  example.com/demo', C=C, typo=t)

    prs.save(output_path)


def build_cyberpunk(output_path):
    C = {
        'primary': '#0A0E27', 'accent': '#00F0FF', 'muted': '#7B2FFF',
        'light': '#1A1F3A', 'white': '#E0E0E0', 'background': '#0A0E27',
        'text_dark': '#E0E0E0', 'text_body': '#B0B0C0', 'text_muted': '#6A6A8A',
        'divider': '#2A2F4A', 'card_bg': '#12162B', 'bg_tint': '#1A1F3A',
        'font_heading': 'Orbitron', 'font_body': 'JetBrains Mono',
    }
    t = TYPOGRAPHY['cyberpunk']
    sp = SPACING['cyberpunk']

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = add_slide(prs)
    hero_slide(s, 'NEURAL OPS // v4.2', 'AI Infrastructure at Scale — 2025 Status Report', C=C, typo=t)

    s = add_slide(prs)
    top_bar(s, C['accent'], C=C)
    page_header(s, 'SYSTEM METRICS', 'Real-time dashboard snapshot', C=C, typo=t, spacing=sp)
    kpi_card(s, 0.8, 1.8, 3.8, 1.35, '12.8B', 'INFERENCE/SEC', '+8.3%', C=C, typo=t)
    kpi_card(s, 5.0, 1.8, 3.8, 1.35, '99.9%', 'UPTIME SLA', '+0.1%', C=C, typo=t)
    kpi_card(s, 9.2, 1.8, 3.8, 1.35, '5.2M', 'ACTIVE NODES', '+15%', C=C, typo=t)
    kpi_card(s, 0.8, 3.5, 3.8, 1.35, '2.4ms', 'P99 LATENCY', '-12%', trend_up=False, C=C, typo=t)
    kpi_card(s, 5.0, 3.5, 3.8, 1.35, '72', 'THREAT SCORE', 'LOW', C=C, typo=t)
    kpi_card(s, 9.2, 3.5, 3.8, 1.35, '38%', 'GPU UTIL', 'OPTIMAL', C=C, typo=t)

    s = add_slide(prs)
    top_bar(s, C['accent'], C=C)
    page_header(s, 'TRAFFIC DISTRIBUTION', 'By service endpoint', C=C, typo=t, spacing=sp)
    bar_chart(s, 2.2, 2.0, [
        ('/api/infer', 0.82, '10.2B'),
        ('/api/embed', 0.48, '5.8B'),
        ('/api/train', 0.30, '3.9B'),
    ], max_width=5.0, bar_height=0.3, C=C, typo=t, spacing=sp)
    donut_chart(s, 9.5, 4.0, 1.8, 1.0, [
        ('Inference', '52%', C['accent']),
        ('Embedding', '29%', C['muted']),
        ('Training', '19%', '#FF2D6B'),
    ], C=C, typo=t)

    s = add_slide(prs)
    top_bar(s, C['accent'], C=C)
    page_header(s, 'CORE MODULES', 'Active service mesh', C=C, typo=t, spacing=sp)
    highlight_cards(s, 0.8, 2.0, [
        ('INFERENCE ENGINE', 'Auto-routing, 70% compression, P99<10ms', C['accent']),
        ('MONITORING MESH', 'Millisecond alerting, full-stack tracing', C['muted']),
        ('CONNECTOR HUB', '200+ integrations, zero-config mesh', '#FF2D6B'),
    ], total_width=12.0, C=C, typo=t, spacing=sp)

    s = add_slide(prs)
    top_bar(s, C['accent'], C=C)
    page_header(s, 'DEPLOY SEQUENCE', 'Production pipeline', C=C, typo=t, spacing=sp)
    code_block(s, 0.8, 2.0, 12.0, 3.5, [
        '$ neural-ops deploy --model gpt-4o-ft --gpu A100',
        '> Replicas: 3  |  Region: us-east-1  |  Auto-scale: ON',
        '> Health check: PASS  |  Endpoint: api.neuralops.ai/v4',
        '> Status: LIVE  |  P99: 2.1ms  |  Throughput: 12.8B/s',
    ], language='shell', C=C, typo=t)

    s = add_slide(prs)
    cta_slide(s, 'INITIALIZE NEXT PHASE', 'neural-ops.ai  |  #build-the-future', C=C, typo=t)

    prs.save(output_path)


def build_creative(output_path):
    C = {
        'primary': '#6C3CE1', 'accent': '#FF6B6B', 'muted': '#4ECDC4',
        'light': '#FFE66D', 'white': '#FFFFFF', 'background': '#FAFAFA',
        'text_dark': '#2D3436', 'text_body': '#636E72', 'text_muted': '#B2BEC3',
        'divider': '#DFE6E9', 'card_bg': '#FFFFFF', 'bg_tint': '#F0F0F5',
        'font_heading': 'Fredoka', 'font_body': 'Nunito',
    }
    t = TYPOGRAPHY['creative']
    sp = SPACING['creative']

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = add_slide(prs)
    hero_slide(s, 'Creative Studio 2025', 'Design. Build. Ship. Repeat.', C=C, typo=t)
    oval(s, 9.0, -1.0, 6.0, 6.0, C['accent'])
    oval(s, -2.0, 4.0, 5.0, 5.0, C['muted'])

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Impact Numbers', 'What we achieved together', C=C, typo=t, spacing=sp)
    kpi_card(s, 0.8, 1.8, 3.8, 1.35, '847', 'Projects Delivered', '+23%', C=C, typo=t)
    kpi_card(s, 5.0, 1.8, 3.8, 1.35, '99.2%', 'Client Satisfaction', '+1.8%', C=C, typo=t)
    kpi_card(s, 9.2, 1.8, 3.8, 1.35, '156', 'Team Members', '+34', C=C, typo=t)

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Service Mix', 'Where we create value', C=C, typo=t, spacing=sp)
    bar_chart(s, 2.2, 2.0, [
        ('Brand Design', 0.72, '61%'),
        ('Product UX', 0.55, '47%'),
        ('Motion Design', 0.38, '32%'),
    ], max_width=5.0, bar_height=0.3, C=C, typo=t, spacing=sp)
    donut_chart(s, 9.5, 4.0, 1.8, 1.0, [
        ('Brand', '40%', C['primary']),
        ('UX', '35%', C['accent']),
        ('Motion', '25%', C['muted']),
    ], C=C, typo=t)

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Superpowers', 'Our creative toolkit', C=C, typo=t, spacing=sp)
    highlight_cards(s, 0.8, 2.0, [
        ('Design System', 'Component library with 500+ tokens, auto-sync Figma', C['primary']),
        ('Rapid Prototyping', 'Idea to clickable prototype in 48 hours', C['accent']),
        ('User Research', 'Data-driven design decisions, A/B testing built-in', C['muted']),
    ], total_width=12.0, C=C, typo=t, spacing=sp)

    s = add_slide(prs)
    top_bar(s, C['primary'], C=C)
    page_header(s, 'Before & After', 'Client transformation stories', C=C, typo=t, spacing=sp)
    comparison_bars(s, 2.5, 2.0, [
        ('Time to Market', '6 mo', '8 wk', 0.9, 0.3),
        ('Design Consistency', '40%', '95%', 0.4, 0.95),
        ('User Retention', '32%', '78%', 0.32, 0.78),
    ], max_width=4.0, C=C, typo=t, spacing=sp)

    s = add_slide(prs)
    cta_slide(s, "Let's Create Together",
              'hello@creativestudio.io  |  book a free consultation', C=C, typo=t)
    oval(s, 10.0, 0.5, 4.0, 4.0, C['accent'])
    oval(s, -1.0, 5.0, 3.0, 3.0, C['muted'])

    prs.save(output_path)


def main():
    out = os.path.join(tempfile.gettempdir(), "ppt_build_v2")
    os.makedirs(out, exist_ok=True)

    print("="*80)
    print("  BUILD MODE v2: Typography + Spacing + GroupShape")
    print("="*80)

    p1 = os.path.join(out, "build_mckinsey.pptx")
    p2 = os.path.join(out, "build_cyberpunk.pptx")
    p3 = os.path.join(out, "build_creative.pptx")

    print("\n[1/3] McKinsey...")
    build_mckinsey(p1)
    print("[2/3] Cyberpunk...")
    build_cyberpunk(p2)
    print("[3/3] Creative...")
    build_creative(p3)

    for name, label in [(p1, "McKinsey"), (p2, "Cyberpunk"), (p3, "Creative")]:
        prs = Presentation(name)
        groups = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 6:
                    groups += 1
        print(f"\n  {label}: {len(prs.slides)} slides, {groups} GroupShapes, {round(os.path.getsize(name)/1024,1)} KB")

    print("\nDone! Files in:", out)


if __name__ == "__main__":
    main()
