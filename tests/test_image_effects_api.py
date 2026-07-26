"""Tests for Phase 2 API exposure in PrecisionRenderer and build_helpers.

Covers:
  1. add_circle_image() — circle-cropped image via PrecisionRenderer
  2. add_image_with_soft_edge() — image with soft/faded edges
  3. add_image_with_duotone() — image with duotone color effect
  4. add_image_with_artistic() — image with artistic effect
  5. circle_image() — build_helpers wrapper
  6. soft_edge_image() — build_helpers wrapper
  7. duotone_image() — build_helpers wrapper
"""

from __future__ import annotations

import os
import tempfile

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_pro_max.build_helpers import circle_image, duotone_image, soft_edge_image
from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer


def _make_brand() -> BrandSpec:
    return BrandSpec(
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "accent": "#6366F1",
            "background": "#FFFFFF",
            "foreground": "#0F172A",
            "muted": "#F1F5F9",
            "muted-foreground": "#64748B",
            "border": "#E2E8F0",
        },
        fonts={"heading": "Inter", "body": "Inter"},
    )


def _make_pr() -> PrecisionRenderer:
    return PrecisionRenderer(brand_spec=_make_brand())


def _make_slide():
    pr = _make_pr()
    prs = pr.create_presentation()
    slide = pr.add_slide(prs)
    return pr, prs, slide


@pytest.fixture
def test_image():
    img = PILImage.new("RGB", (200, 200), (255, 100, 50))
    path = os.path.join(tempfile.gettempdir(), "api_test_img.png")
    img.save(path, "PNG")
    return path


def _make_build_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    return prs, slide


# ── PrecisionRenderer: add_circle_image ──


class TestAddCircleImage:
    def test_creates_circle_image(self, test_image):
        pr, prs, slide = _make_slide()
        shape = pr.add_circle_image(slide, 5.0, 3.0, 1.0, test_image)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:blipFill")) is not None

    def test_nonexistent_image_returns_none(self):
        pr, prs, slide = _make_slide()
        result = pr.add_circle_image(slide, 5.0, 3.0, 1.0, "/nonexistent.png")
        assert result is None

    def test_with_brand_border(self, test_image):
        pr, prs, slide = _make_slide()
        shape = pr.add_circle_image(slide, 5.0, 3.0, 1.0, test_image)
        spPr = shape._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        solidFill = ln.find(qn("a:solidFill")) if ln is not None else None
        assert solidFill is not None, "Should have border from brand border color"


# ── PrecisionRenderer: add_image_with_soft_edge ──


class TestAddImageWithSoftEdge:
    def test_creates_image_with_soft_edge(self, test_image):
        pr, prs, slide = _make_slide()
        shape = pr.add_image_with_soft_edge(slide, test_image, 1.0, 1.0, 4.0, 3.0,
                                            radius_pt=10)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        assert effectLst is not None
        softEdge = effectLst.find(qn("a:softEdge"))
        assert softEdge is not None

    def test_nonexistent_image_returns_none(self):
        pr, prs, slide = _make_slide()
        result = pr.add_image_with_soft_edge(slide, "/nonexistent.png",
                                             1.0, 1.0, 4.0, 3.0)
        assert result is None


# ── PrecisionRenderer: add_image_with_duotone ──


class TestAddImageWithDuotone:
    def test_creates_duotone_image(self, test_image):
        pr, prs, slide = _make_slide()
        shape = pr.add_image_with_duotone(slide, test_image, 1.0, 1.0, 4.0, 3.0,
                                          color1="#0000FF", color2="#FF0000")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        duotone = blip.find(qn("a:duotone"))
        assert duotone is not None


# ── PrecisionRenderer: add_image_with_artistic ──


class TestAddImageWithArtistic:
    def test_creates_artistic_image(self, test_image):
        pr, prs, slide = _make_slide()
        shape = pr.add_image_with_artistic(slide, test_image, 1.0, 1.0, 4.0, 3.0,
                                           effect="watercolor_sponge")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        assert blip.find(qn("a:artisticWatercolorSponge")) is not None

    def test_invalid_effect_raises(self, test_image):
        pr, prs, slide = _make_slide()
        with pytest.raises(KeyError, match="Unknown artistic effect"):
            pr.add_image_with_artistic(slide, test_image, 1.0, 1.0, 4.0, 3.0,
                                       effect="nonexistent")


# ── build_helpers: circle_image ──


class TestCircleImageBuildHelper:
    def test_creates_circle_image(self, test_image):
        prs, slide = _make_build_slide()
        shape = circle_image(slide, 5.0, 3.0, 1.0, test_image)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:blipFill")) is not None


# ── build_helpers: soft_edge_image ──


class TestSoftEdgeImageBuildHelper:
    def test_creates_soft_edge_image(self, test_image):
        prs, slide = _make_build_slide()
        shape = soft_edge_image(slide, 1.0, 1.0, 4.0, 3.0, test_image,
                                soft_radius=10)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        assert effectLst is not None
        assert effectLst.find(qn("a:softEdge")) is not None

    def test_nonexistent_image_returns_none(self):
        prs, slide = _make_build_slide()
        result = soft_edge_image(slide, 1.0, 1.0, 4.0, 3.0, "/nonexistent.png")
        assert result is None


# ── build_helpers: duotone_image ──


class TestDuotoneImageBuildHelper:
    def test_creates_duotone_image(self, test_image):
        prs, slide = _make_build_slide()
        shape = duotone_image(slide, 1.0, 1.0, 4.0, 3.0, test_image,
                              color1="#0000FF", color2="#FF0000")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        assert blip.find(qn("a:duotone")) is not None
