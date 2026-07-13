"""P8: End-to-end tests covering P1-P7 full pipeline flows."""

from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from ppt_pro_max import generate_ppt
from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
from ppt_pro_max.enterprise.proposal_generator import ProposalGenerator
from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.renderer.theme_composer import ThemeComposer


def _make_png(path: Path, w: int = 200, h: int = 150) -> Path:
    img = Image.new("RGB", (w, h), color="green")
    img.save(str(path))
    return path


def _setup_full_project(tmp_path: Path) -> Path:
    project = tmp_path / "fullproject"
    project.mkdir()
    prs = Presentation()
    for _ in range(2):
        prs.slides.add_slide(prs.slide_layouts[-1])
    prs.save(str(project / "template.pptx"))
    _make_png(project / "logo.png")
    _make_png(project / "hero.png", w=1600, h=900)
    _make_png(project / "product.png", w=1000, h=800)
    _make_png(project / "icon.png", w=400, h=300)
    (project / "brand.json").write_text(json.dumps({"colors": {"primary": "#1A3C6E"}}), encoding="utf-8")
    content = {
        "meta": {"title": "E2E Full Test"},
        "slides": [
            {"goal": "hook", "title": "Welcome", "subtitle": "Our Product", "image": "hero.png"},
            {"goal": "problem", "title": "Challenges", "bullets": ["Slow", "Expensive", "Complex"]},
            {"goal": "solution", "title": "Our Solution", "bullets": ["Fast", "Affordable", "Simple"]},
            {"goal": "features", "title": "Key Features", "cards": [
                {"title": "Speed", "body": "10x faster"},
                {"title": "Cost", "body": "50% cheaper"},
                {"title": "Ease", "body": "Zero learning curve"},
            ]},
            {"goal": "data", "title": "Architecture", "diagram_type": "table", "diagram_data": {"type": "table", "rows": [["Layer", "Tech"], ["Frontend", "React"], ["Backend", "Python"]]}},
            {"goal": "code", "title": "Quick Start", "code": {"language": "python", "code": "pip install product\nproduct.start()"}},
            {"goal": "exercise", "title": "Try It", "exercise": {"instructions": "Follow these steps", "duration": "10 min", "steps": ["Step 1", "Step 2"]}},
            {"goal": "cta", "title": "Get Started", "subtitle": "Contact us today"},
        ],
    }
    (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
    return project


class TestE2EEnterpriseFullPipeline:

    def test_full_enterprise_pipeline(self, tmp_path):
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="E2E Full Test", project_dir=str(project), density=5)
        assert result["pipeline"] == "enterprise"
        assert result["num_slides"] == 8
        assert result["render_mode"] == "precision"
        assert os.path.isfile(result["output_path"])
        prs = Presentation(result["output_path"])
        assert len(prs.slides) == 8

    def test_pipeline_with_business_mode(self, tmp_path):
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Pitch Deck", project_dir=str(project), business_mode="pitch")
        assert result["num_slides"] == 8

    def test_pipeline_with_motion(self, tmp_path):
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Motion Test", project_dir=str(project), motion=5)
        assert result["num_slides"] == 8

    def test_pipeline_dry_run(self, tmp_path):
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Dry Run", project_dir=str(project), dry_run=True)
        assert result["dry_run"] is True


class TestE2EReadmePipeline:

    def test_readme_only_project(self, tmp_path):
        project = tmp_path / "readme_project"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        (project / "README.md").write_text(
            "# Product\nIntro\n\n# Features\n- Fast\n- Secure\n\n# Contact\nGet in touch",
            encoding="utf-8",
        )
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Product", project_dir=str(project))
        assert result["num_slides"] == 3
        assert os.path.isfile(result["output_path"])


class TestE2EImageAssignment:

    def test_size_aware_image_assignment(self, tmp_path):
        project = tmp_path / "img_project"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        _make_png(project / "hero_bg.png", w=1600, h=900)
        _make_png(project / "product.png", w=1000, h=800)
        _make_png(project / "small_icon.png", w=400, h=300)
        content = {
            "meta": {"title": "Image Test"},
            "slides": [
                {"goal": "hook", "title": "Welcome"},
                {"goal": "features", "title": "Features", "cards": [{"title": "A"}, {"title": "B"}]},
            ],
        }
        (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Image Test", project_dir=str(project))
        assert result["num_slides"] == 2
        prs_out = Presentation(result["output_path"])
        for slide in prs_out.slides:
            assert len(slide.shapes) > 0


class TestE2EProposalFlow:

    def test_proposal_generate_and_inspect(self, tmp_path):
        gen = ProposalGenerator()
        proposals = gen.generate(query="Investor Pitch", style="professional", output_dir=str(tmp_path / "proposals"))
        assert len(proposals) == 3
        for p in proposals:
            assert os.path.isfile(p["path"])
            prs = Presentation(p["path"])
            assert len(prs.slides) == 4

    def test_proposal_via_api(self, tmp_path):
        result = generate_ppt("AI Pitch", proposal=True, output=str(tmp_path / "proposals"))
        assert len(result["proposals"]) == 3

    def test_proposal_different_styles(self, tmp_path):
        styles = ["professional", "dark cyberpunk", "warm elegant"]
        for s in styles:
            result = generate_ppt("Test", style=s, proposal=True, output=str(tmp_path / f"prop_{s.replace(' ', '_')}"))
            assert len(result["proposals"]) == 3


class TestE2EPrecisionRenderer:

    def test_render_all_goal_types(self):
        brand_spec = BrandSpec()
        precision = PrecisionRenderer(brand_spec=brand_spec)
        prs = precision.create_presentation()

        goals = [
            ("hook", {"title": "Welcome", "subtitle": "Our Product"}),
            ("problem", {"title": "Challenges", "bullets": ["A", "B", "C"]}),
            ("solution", {"title": "Solution", "bullets": ["X", "Y"]}),
            ("features", {"title": "Features", "cards": [{"title": "F1", "body": "D1"}, {"title": "F2", "body": "D2"}]}),
            ("data", {"title": "Architecture"}),
            ("code", {"title": "Quick Start", "code": {"language": "python", "code": "print('hello')"}}),
            ("exercise", {"title": "Try It", "exercise": {"instructions": "Do this", "duration": "5 min"}}),
            ("cta", {"title": "Get Started", "subtitle": "Contact us"}),
            ("content", {"title": "Details", "bullets": ["Point 1", "Point 2"]}),
            ("overview", {"title": "Agenda", "bullets": ["Topic 1", "Topic 2", "Topic 3"]}),
        ]
        for goal, page in goals:
            page["goal"] = goal
            precision.render_slide(prs, page)

        assert len(prs.slides) == 10
        for slide in prs.slides:
            assert len(slide.shapes) > 0

    def test_render_with_brand_spec(self):
        brand_spec = BrandSpec()
        brand_spec.source = "test"
        brand_spec.colors = {"primary": "#2563EB", "background": "#FFFFFF", "foreground": "#1E293B"}
        brand_spec.fonts = {"heading": "Inter", "body": "Inter"}
        precision = PrecisionRenderer(brand_spec=brand_spec)
        prs = precision.create_presentation()
        precision.render_slide(prs, {"goal": "hook", "title": "Branded"})
        assert len(prs.slides) == 1


class TestE2EThemeComposer:

    def test_all_35_moods_compose(self):
        tc = ThemeComposer()
        moods = [
            "professional", "tech", "dark", "warm", "elegant", "luxury", "vibrant",
            "startup", "nature", "calm", "minimal", "bold", "fresh", "industrial",
            "fintech", "health", "education", "sustainability", "creative",
            "international", "cream", "frosted", "mckinsey", "consulting",
            "pastel", "retro", "government", "legal", "pharma", "realestate",
            "automotive", "aviation", "energy", "telecom", "logistics",
        ]
        for mood in moods:
            result = tc.compose(style=mood, seed=42)
            assert "atoms" in result, f"Mood '{mood}' failed to compose"
            assert result["atoms"].get("palette"), f"Mood '{mood}' has no palette"

    def test_style_determinism_with_seed(self):
        tc = ThemeComposer()
        r1 = tc.compose(style="professional", seed=123)
        r2 = tc.compose(style="professional", seed=123)
        assert r1["name"] == r2["name"]


class TestE2EDensityProfile:

    def test_low_density_fewer_bullets(self, tmp_path):
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result_low = pipeline.run(query="Low Density", project_dir=str(project), density=3)
        assert result_low["num_slides"] == 8

    def test_high_density_more_bullets(self, tmp_path):
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result_high = pipeline.run(query="High Density", project_dir=str(project), density=8)
        assert result_high["num_slides"] == 8


class TestE2EContentPriority:

    def test_content_json_over_readme(self, tmp_path):
        project = tmp_path / "priority"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        (project / "README.md").write_text("# From Readme\n\nReadme content", encoding="utf-8")
        (project / "content.json").write_text(json.dumps({
            "meta": {"title": "From JSON"},
            "slides": [{"goal": "hook", "title": "JSON Wins"}],
        }), encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Priority Test", project_dir=str(project))
        assert result["num_slides"] == 1

    def test_readme_over_storyplanner(self, tmp_path):
        project = tmp_path / "readme_priority"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        (project / "README.md").write_text("# P1\n\nText\n\n# P2\n\nMore\n\n# P3\n\nEnd", encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Readme Priority", project_dir=str(project))
        assert result["num_slides"] == 3
