"""Tests for P7: generate_ppt() API integration — proposal/confirmed_proposal/materials_dir."""

from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from ppt_pro_max import generate_ppt


def _make_png(path: Path, w: int = 200, h: int = 150) -> Path:
    img = Image.new("RGB", (w, h), color="blue")
    img.save(str(path))
    return path


def _setup_project(tmp_path: Path) -> Path:
    project = tmp_path / "project"
    project.mkdir()
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[-1])
    prs.save(str(project / "template.pptx"))
    return project


class TestProposalMode:

    def test_proposal_returns_three_proposals(self, tmp_path):
        result = generate_ppt("Test Presentation", proposal=True, output=str(tmp_path / "out"))
        assert "proposals" in result
        assert len(result["proposals"]) == 3

    def test_proposal_ids(self, tmp_path):
        result = generate_ppt("Test Presentation", proposal=True, output=str(tmp_path / "out"))
        ids = [p["id"] for p in result["proposals"]]
        assert ids == ["A", "B", "C"]

    def test_proposal_files_exist(self, tmp_path):
        out_dir = tmp_path / "proposals"
        result = generate_ppt("Test Presentation", proposal=True, output=str(out_dir))
        for p in result["proposals"]:
            assert os.path.isfile(p["path"])

    def test_proposal_each_has_4_slides(self, tmp_path):
        result = generate_ppt("Test Presentation", proposal=True, output=str(tmp_path / "out"))
        for p in result["proposals"]:
            prs = Presentation(p["path"])
            assert len(prs.slides) == 4

    def test_proposal_with_style(self, tmp_path):
        result = generate_ppt("AI Pitch", style="dark cyberpunk", proposal=True, output=str(tmp_path / "out"))
        assert len(result["proposals"]) == 3

    def test_proposal_with_project(self, tmp_path):
        project = _setup_project(tmp_path)
        result = generate_ppt("Enterprise Test", project=str(project), proposal=True)
        assert len(result["proposals"]) == 3


class TestProposalWithProjectDir:

    def test_proposal_with_project_sets_output_dir(self, tmp_path):
        project = _setup_project(tmp_path)
        result = generate_ppt("Enterprise Test", project=str(project), proposal=True)
        for p in result["proposals"]:
            assert project.name in p["path"] or "output" in p["path"]


class TestMaterialsDir:

    def test_materials_dir_with_project(self, tmp_path):
        project = _setup_project(tmp_path)
        materials = tmp_path / "materials"
        materials.mkdir()
        _make_png(materials / "hero.png", w=1600, h=900)
        content = {"meta": {"title": "Materials Test"}, "slides": [{"goal": "hook", "title": "Welcome"}]}
        (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
        result = generate_ppt("Materials Test", project=str(project), materials_dir=str(materials))
        assert result["num_slides"] == 1


class TestFreestyleNotBroken:

    def test_freestyle_still_works(self, tmp_path):
        result = generate_ppt("Simple Test", slides=3, output=str(tmp_path / "freestyle.pptx"))
        assert os.path.isfile(result["output_path"])

    def test_enterprise_still_works(self, tmp_path):
        project = _setup_project(tmp_path)
        content = {"meta": {"title": "Enterprise"}, "slides": [{"goal": "hook", "title": "Welcome"}]}
        (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
        result = generate_ppt("Enterprise", project=str(project))
        assert result["num_slides"] == 1
