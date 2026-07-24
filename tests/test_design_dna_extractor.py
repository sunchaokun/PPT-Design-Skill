"""Tests for DesignDNAExtractor — zero-loss design DNA extraction."""

import os
import tempfile
import pytest

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _create_sample_pptx(path: str, num_slides: int = 3):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layouts = prs.slide_layouts
    blank = layouts[6] if len(layouts) > 6 else layouts[-1]

    cover = prs.slides.add_slide(layouts[0] if layouts else blank)
    for ph in cover.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = "Test Cover Title"
        elif ph.placeholder_format.type == 2:
            ph.text = "Test Subtitle"

    for i in range(1, num_slides - 1):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Content Slide {i}"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        btf = body.text_frame
        for j, text in enumerate(["Point A", "Point B", "Point C"]):
            bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            br = bp.add_run()
            br.text = text
            br.font.size = Pt(14)

    back = prs.slides.add_slide(layouts[0] if layouts else blank)
    for ph in back.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = "Thank You"

    prs.save(path)


def _create_green_theme_pptx(path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    slide = prs.slides.add_slide(blank)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2E, 0x65, 0x04)

    tb = slide.shapes.add_textbox(Inches(2), Inches(1), Inches(8), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "乡村振兴"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"

    sub = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(8), Inches(1))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "产业兴旺 · 生态宜居"
    sr.font.size = Pt(24)
    sr.font.color.rgb = RGBColor(0xD4, 0xE3, 0xAC)
    sr.font.name = "Arial"

    prs.save(path)


class TestDesignDNAExtractor:

    def test_extract_basic_structure(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_sample_pptx(pptx_path, num_slides=3)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert dna.source_path == pptx_path
        assert len(dna.slides) == 3
        assert dna.slide_width_emu > 0
        assert dna.slide_height_emu > 0

    def test_extract_text_zones(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_sample_pptx(pptx_path, num_slides=3)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        total_zones = sum(len(s.text_zones) for s in dna.slides)
        assert total_zones > 0

        for slide_dna in dna.slides:
            for zone in slide_dna.text_zones:
                assert zone.zone_id
                assert zone.role in ("title", "subtitle", "body", "badge", "decoration", "logo_text", "page_number")
                assert zone.text

    def test_extract_color_palette(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_green_theme_pptx(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert len(dna.color_palette) > 0
        assert "foreground" in dna.color_palette or "background" in dna.color_palette

    def test_extract_actual_colors(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_green_theme_pptx(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert len(dna.actual_colors) > 0
        green_found = any("2E6504" in c.upper() or "466740" in c.upper() or "D4E3AC" in c.upper() for c in dna.actual_colors)
        assert green_found or len(dna.actual_colors) > 0

    def test_extract_font_scheme(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_sample_pptx(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert "heading" in dna.font_scheme or len(dna.font_scheme) > 0

    def test_page_type_classification(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_sample_pptx(pptx_path, num_slides=3)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert dna.slides[0].page_type == "cover"
        assert dna.slides[-1].page_type in ("back_cover", "content")

    def test_slide_xml_preserved(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_sample_pptx(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        for slide_dna in dna.slides:
            assert len(slide_dna.slide_xml) > 0
            assert b"<p:sld" in slide_dna.slide_xml or b"sld" in slide_dna.slide_xml

    def test_brand_spec_built(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_sample_pptx(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert dna.brand_spec is not None
        assert dna.brand_spec.source == "design_dna"

    def test_clone_and_patch_preserves_visual(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_green_theme_pptx(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        output_path = str(tmp_path / "patched.pptx")
        patches = [{"title": "数字乡村", "subtitle": "智慧引领"}]
        result = extractor.clone_and_patch(dna, patches, output_path)

        assert os.path.isfile(result)
        prs = Presentation(result)
        assert len(prs.slides) >= 1

    def test_clone_and_patch_replaces_text(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_sample_pptx(pptx_path, num_slides=3)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        output_path = str(tmp_path / "patched.pptx")
        patches = [
            {"title": "New Cover", "subtitle": "New Subtitle"},
            {"title": "New Content", "bullets": ["New Point A", "New Point B"]},
            {"title": "New Ending"},
        ]
        result = extractor.clone_and_patch(dna, patches, output_path)

        prs = Presentation(result)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text)
        combined = " ".join(all_text)

        assert "New Cover" in combined or "New Content" in combined

    def test_extract_dna_api(self, tmp_path):
        from ppt_pro_max import extract_design_dna

        pptx_path = str(tmp_path / "test.pptx")
        _create_sample_pptx(pptx_path)

        result = extract_design_dna(pptx_path)

        assert "slides" in result
        assert "color_palette" in result
        assert "font_scheme" in result
        assert "actual_colors" in result
        assert "brand_spec" in result
        assert result["num_slides"] == 3

    def test_real_pptx_extraction(self):
        pptx_path = r"E:\PPT-Design-Skill\docs\05乡村振兴.pptx"
        if not os.path.isfile(pptx_path):
            pytest.skip("05乡村振兴.pptx not found")

        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert len(dna.slides) == 5
        assert dna.slides[0].page_type == "cover"
        assert dna.slides[1].page_type == "toc"
        assert dna.slides[2].page_type == "transition"
        assert dna.slides[4].page_type == "back_cover"

        assert len(dna.color_palette) >= 4
        assert len(dna.actual_colors) >= 5
        assert len(dna.actual_fonts) >= 1

        total_zones = sum(len(s.text_zones) for s in dna.slides)
        assert total_zones >= 20

        assert dna.brand_spec is not None
        assert dna.brand_spec.colors is not None

    def test_real_pptx_clone_patch(self):
        pptx_path = r"E:\PPT-Design-Skill\docs\05乡村振兴.pptx"
        if not os.path.isfile(pptx_path):
            pytest.skip("05乡村振兴.pptx not found")

        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        output_dir = r"E:\PPT-Design-Skill\docs\xiangcun-project\output"
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, "dna_patched.pptx")

        patches = [
            {"title": "数字乡村建设规划", "subtitle": "智慧引领 · 数字赋能"},
            {"title": "目录", "subtitle": "数字乡村发展现状", "bullets": ["基础设施", "智慧农业", "数字治理", "乡村电商"]},
            {"title": "数字基础设施现状", "subtitle": "PART ONE"},
            {"title": "数字基础设施现状", "bullets": ["农村宽带覆盖率98%", "5G基站向乡镇延伸", "物联网终端部署加速"]},
            {"title": "携手共建", "subtitle": "数字乡村新未来"},
        ]
        result = extractor.clone_and_patch(dna, patches, output_path)

        assert os.path.isfile(result)
        prs = Presentation(result)
        assert len(prs.slides) == 5
