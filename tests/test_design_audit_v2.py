"""Precise design audit of raw PPTs: spacing, alignment, font scale, whitespace."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu

ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def in_(emu):
    return round(emu / 914400, 3) if emu else None

def audit(path, label):
    prs = Presentation(path)
    sw = in_(prs.slide_width)
    sh = in_(prs.slide_height)
    print(f"\n{'='*80}")
    print(f"  {label} | {sw}\" x {sh}\"")
    print(f"{'='*80}")

    for i, slide in enumerate(prs.slides):
        shapes_info = []
        all_texts = []
        all_sizes = []
        
        for shape in slide.shapes:
            L = in_(shape.left)
            T = in_(shape.top)
            W = in_(shape.width)
            H = in_(shape.height)
            R = round(L + W, 3) if L and W else None
            B = round(T + H, 3) if T and H else None
            info = {'L': L, 'T': T, 'W': W, 'H': H, 'R': R, 'B': B, 'type': str(shape.shape_type)}
            
            if shape.has_text_frame:
                txt = shape.text_frame.text.replace('\n', '|')[:50]
                for para in shape.text_frame.paragraphs:
                    pPr = para._p.find(f'{{{ns_a}}}pPr')
                    def_sz = def_fn = None
                    if pPr is not None:
                        dR = pPr.find(f'{{{ns_a}}}defRPr')
                        if dR is not None:
                            def_sz = dR.get('sz')
                            lt = dR.find(f'{{{ns_a}}}latin')
                            if lt is not None: def_fn = lt.get('typeface')
                    for run in para.runs:
                        if not run.text.strip(): continue
                        rPr = run._r.find(f'{{{ns_a}}}rPr')
                        sz = fn = None
                        if rPr is not None:
                            sz = rPr.get('sz')
                            lt = rPr.find(f'{{{ns_a}}}latin')
                            if lt is not None: fn = lt.get('typeface')
                        pt = int(sz)/100 if sz else (int(def_sz)/100 if def_sz else None)
                        fn = fn or def_fn or '?'
                        if pt: all_sizes.append(pt)
                        all_texts.append({'txt': run.text.strip()[:30], 'pt': pt, 'fn': fn, 'L': L, 'T': T})
                info['txt'] = txt
            
            shapes_info.append(info)
        
        # Layout analysis
        content_shapes = [s for s in shapes_info if s['T'] is not None and s['T'] > 0.05 and s.get('txt')]
        if not content_shapes:
            content_shapes = [s for s in shapes_info if s['T'] is not None and s.get('txt')]
        
        top_content = min((s['T'] for s in content_shapes), default=None)
        bottom_content = max((s['B'] for s in content_shapes if s['B']), default=None)
        left_content = min((s['L'] for s in content_shapes if s['L'] and s['L'] > 0), default=None)
        right_content = max((s['R'] for s in content_shapes if s['R']), default=None)
        
        top_margin = top_content
        bottom_margin = round(sh - bottom_content, 2) if bottom_content else None
        left_margin = left_content
        right_margin = round(sw - right_content, 2) if right_content else None
        
        print(f"\n  Slide {i}:")
        print(f"    Margins: top={top_margin}\" bottom={bottom_margin}\" left={left_margin}\" right={right_margin}\"")
        
        if bottom_margin and bottom_margin > 2.0:
            print(f"    ⚠ BOTTOM WASTE: {bottom_margin}\" empty at bottom (should be <1.5\")")
        if top_margin and top_margin > 1.0 and i > 0:
            print(f"    ⚠ TOP WASTE: {top_margin}\" empty at top (should be ~0.4-0.6\")")
        
        # Font size analysis
        if all_sizes:
            unique = sorted(set(all_sizes))
            print(f"    Font sizes: {unique}")
            if min(all_sizes) < 8:
                print(f"    ⚠ TOO SMALL: {min(all_sizes)}pt (min 9pt)")
        
        # Text positioning check
        for t in all_texts[:6]:
            print(f"    \"{t['txt']}\" {t['pt']}pt {t['fn']} @({t['L']},{t['T']})")
        if len(all_texts) > 6:
            print(f"    ... +{len(all_texts)-6} more")

for name, label in [('mckinsey.pptx','McKinsey'), ('cyberpunk.pptx','Cyberpunk'), ('creative.pptx','Creative')]:
    p = os.path.join(os.environ['TEMP'], 'ppt_raw_design', name)
    if os.path.exists(p):
        audit(p, label)
