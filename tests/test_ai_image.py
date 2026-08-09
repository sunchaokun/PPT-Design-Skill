"""Tests for build_helpers.ai_image() 鈥?one-call image generate + place.

Ensures build.py can fetch an image and place it cover-fit without hand-rolled
image-API scripts. Network/COM is mocked 鈥?fetch path is patched.
"""
from __future__ import annotations

from unittest.mock import patch

import pytest
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def slide(tmp_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])


@pytest.fixture
def fake_image(tmp_path):
    """A tiny real PNG on disk."""
    from PIL import Image

    img = Image.new("RGB", (200, 100), (10, 20, 30))
    p = tmp_path / "fake.png"
    img.save(str(p))
    return str(p)


def test_ai_image_generates_and_places(slide, fake_image):
    from ppt_pro_max.build_helpers import ai_image

    def _fake_fetch(*args, **kwargs):
        return {"path": fake_image, "mode": "generate", "provider": "seedream"}

    with patch("ppt_pro_max.build_helpers.fetch_image", side_effect=_fake_fetch):
        shape = ai_image(slide, 1, 1, 4, 3, "protein structure 3D",
                         mode="generate", llm_provider="seedream")

    assert shape is not None
    assert shape.shape_type == 13  # PICTURE
    assert len(slide.shapes) == 1


def test_ai_image_passes_keywords_and_box(slide, fake_image):
    from ppt_pro_max.build_helpers import ai_image

    captured = {}

    def _fake_fetch(*args, **kwargs):
        captured["keywords"] = args[0]
        captured.update(kwargs)
        return {"path": fake_image}

    with patch("ppt_pro_max.build_helpers.fetch_image", side_effect=_fake_fetch):
        ai_image(slide, 2, 3, 5, 4, "lab bench equipment",
                 mode="search", emotion="neutral", goal="content",
                 unsplash_access_key="k1", pexels_api_key="k2")

    assert captured["keywords"] == "lab bench equipment"
    assert captured["mode"] == "search"
    assert captured["emotion"] == "neutral"
    assert captured["goal"] == "content"
    assert captured["unsplash_access_key"] == "k1"
    assert captured["pexels_api_key"] == "k2"
    assert captured["width"] == int(5 * 96)
    assert captured["height"] == int(4 * 96)


def test_ai_image_fallback_placeholder_on_failure(slide):
    from ppt_pro_max.build_helpers import ai_image

    with patch("ppt_pro_max.build_helpers.fetch_image", return_value={"path": None}):
        shape = ai_image(slide, 1, 1, 4, 3, "nothing found")
    # placeholder box + keyword text drawn, picture absent
    assert shape is None
    assert len(slide.shapes) == 2


def test_ai_image_no_fallback_when_disabled(slide):
    from ppt_pro_max.build_helpers import ai_image

    with patch("ppt_pro_max.build_helpers.fetch_image", return_value={"path": None}):
        shape = ai_image(slide, 1, 1, 4, 3, "nothing found",
                         fallback_placeholder=False)
    assert shape is None
    assert len(slide.shapes) == 0


def test_ai_image_ignores_missing_file(slide, tmp_path):
    from ppt_pro_max.build_helpers import ai_image

    with patch("ppt_pro_max.build_helpers.fetch_image", return_value={"path": str(tmp_path / "gone.png")}):
        shape = ai_image(slide, 1, 1, 4, 3, "broken path")
    assert shape is None
    assert len(slide.shapes) == 2  # placeholder only


def test_fetch_image_reachable_from_star_import():
    """from ppt_pro_max.build_helpers import *  must expose ai_image + fetch_image."""
    import ppt_pro_max.build_helpers as bh

    ns = {}
    exec("from ppt_pro_max.build_helpers import *", ns)
    assert "ai_image" in ns
    assert "fetch_image" in ns
    assert callable(ns["ai_image"])
    assert bh.ai_image is ns["ai_image"]

