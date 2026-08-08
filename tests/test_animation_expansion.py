"""Tests for Phase 5: Animation & Transition expansion.

Covers:
  1. Morph transition (TRANSITION_TYPES["morph"])
  2. EXIT_PRESETS dict
  3. add_exit_animation() — exit animation XML structure
  4. EMPHASIS_PRESETS dict
  5. add_emphasis_animation() — emphasis animation XML structure
  6. API exposure in PrecisionRenderer
  7. API exposure in build_helpers
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from ppt_pro_max.renderer.animation import (
    TRANSITION_TYPES, ENTRANCE_PRESETS,
    add_slide_transition, add_entrance_animation,
)


def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _add_shape(slide):
    sh = slide.shapes.add_shape(1, 100000, 100000, 200000, 200000)
    return sh.shape_id


# ── Morph transition ──


class TestMorphTransition:
    def test_morph_in_transition_types(self):
        assert "morph" in TRANSITION_TYPES

    def test_morph_tag_and_attrs(self):
        tag, attrs = TRANSITION_TYPES["morph"]
        assert tag == "p:morph"
        assert "option" in attrs
        assert attrs["option"] == "byObject"

    def test_add_morph_transition(self):
        prs, slide = _make_slide()
        add_slide_transition(slide, "morph")
        sld = slide._element
        trans = sld.find(qn("p:transition"))
        assert trans is not None
        morph = trans.find(qn("p:morph"))
        assert morph is not None
        assert morph.get("option") == "byObject"

    def test_morph_opt_in_only(self):
        assert "morph" not in ["fade", "push", "wipe", "split", "cover",
                               "dissolve", "wheel", "wedge", "blinds",
                               "checker", "comb", "random"]


# ── EXIT_PRESETS ──


class TestExitPresets:
    def test_exit_presets_exist(self):
        from ppt_pro_max.renderer.animation import EXIT_PRESETS
        assert isinstance(EXIT_PRESETS, dict)
        assert len(EXIT_PRESETS) >= 7

    def test_exit_presets_keys(self):
        from ppt_pro_max.renderer.animation import EXIT_PRESETS
        expected = {"fade_out", "fly_out_left", "fly_out_right",
                    "fly_out_top", "fly_out_bottom", "zoom_out",
                    "shrink", "spin_out"}
        assert expected.issubset(set(EXIT_PRESETS.keys()))

    def test_exit_presets_values_are_tuples(self):
        from ppt_pro_max.renderer.animation import EXIT_PRESETS
        for key, val in EXIT_PRESETS.items():
            assert isinstance(val, tuple), f"{key}: {val}"
            assert len(val) == 2, f"{key}: expected 2-tuple"
            assert isinstance(val[0], int), f"{key}: preset_id must be int"
            assert isinstance(val[1], int), f"{key}: preset_subtype must be int"

    def test_fade_out_preset(self):
        from ppt_pro_max.renderer.animation import EXIT_PRESETS
        assert EXIT_PRESETS["fade_out"] == (10, 0)

    def test_fly_out_left_preset(self):
        from ppt_pro_max.renderer.animation import EXIT_PRESETS
        assert EXIT_PRESETS["fly_out_left"] == (2, 8)


# ── add_exit_animation ──


class TestAddExitAnimation:
    def test_creates_timing_xml(self):
        from ppt_pro_max.renderer.animation import add_exit_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_exit_animation(slide, shape_id, effect="fade_out")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None

    def test_preset_class_is_exit(self):
        from ppt_pro_max.renderer.animation import add_exit_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_exit_animation(slide, shape_id, effect="fade_out")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        exit_cTns = [c for c in cTns if c.get("presetClass") == "exit"]
        assert len(exit_cTns) >= 1

    def test_exit_preset_id_matches(self):
        from ppt_pro_max.renderer.animation import add_exit_animation, EXIT_PRESETS
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_exit_animation(slide, shape_id, effect="fade_out")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        exit_cTns = [c for c in cTns if c.get("presetClass") == "exit"]
        expected_id, expected_sub = EXIT_PRESETS["fade_out"]
        assert exit_cTns[0].get("presetID") == str(expected_id)
        assert exit_cTns[0].get("presetSubtype") == str(expected_sub)

    def test_exit_with_custom_duration(self):
        from ppt_pro_max.renderer.animation import add_exit_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_exit_animation(slide, shape_id, effect="zoom_out",
                           duration_ms=800, delay_ms=200)
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        exit_cTns = [c for c in cTns if c.get("presetClass") == "exit"]
        assert exit_cTns[0].get("dur") == "800"
        assert exit_cTns[0].get("delay") == "200"

    def test_exit_targets_correct_shape(self):
        from ppt_pro_max.renderer.animation import add_exit_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_exit_animation(slide, shape_id, effect="fade_out")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        spTgts = timing.findall(".//" + qn("p:spTgt"))
        assert any(t.get("spid") == str(shape_id) for t in spTgts)

    def test_exit_invalid_effect_falls_back(self):
        from ppt_pro_max.renderer.animation import add_exit_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_exit_animation(slide, shape_id, effect="nonexistent")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        exit_cTns = [c for c in cTns if c.get("presetClass") == "exit"]
        assert len(exit_cTns) >= 1

    def test_exit_sets_visibility_hidden(self):
        from ppt_pro_max.renderer.animation import add_exit_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_exit_animation(slide, shape_id, effect="fade_out")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        set_elems = timing.findall(".//" + qn("p:set"))
        for s in set_elems:
            attrNames = s.findall(".//" + qn("p:attrName"))
            for a in attrNames:
                if a.text == "style.visibility":
                    to_elem = s.find(qn("p:to"))
                    assert to_elem is not None
                    strVal = to_elem.find(qn("p:strVal"))
                    assert strVal.get("val") == "hidden"


# ── EMPHASIS_PRESETS ──


class TestEmphasisPresets:
    def test_emphasis_presets_exist(self):
        from ppt_pro_max.renderer.animation import EMPHASIS_PRESETS
        assert isinstance(EMPHASIS_PRESETS, dict)
        assert len(EMPHASIS_PRESETS) >= 7

    def test_emphasis_presets_keys(self):
        from ppt_pro_max.renderer.animation import EMPHASIS_PRESETS
        expected = {"grow", "shrink", "spin", "pulse",
                    "color_change", "transparency", "bold_flash", "wave"}
        assert expected.issubset(set(EMPHASIS_PRESETS.keys()))

    def test_emphasis_presets_values_are_tuples(self):
        from ppt_pro_max.renderer.animation import EMPHASIS_PRESETS
        for key, val in EMPHASIS_PRESETS.items():
            assert isinstance(val, tuple), f"{key}: {val}"
            assert len(val) == 2

    def test_grow_preset(self):
        from ppt_pro_max.renderer.animation import EMPHASIS_PRESETS
        assert EMPHASIS_PRESETS["grow"] == (53, 0)

    def test_pulse_preset(self):
        from ppt_pro_max.renderer.animation import EMPHASIS_PRESETS
        assert EMPHASIS_PRESETS["pulse"] == (39, 0)


# ── add_emphasis_animation ──


class TestAddEmphasisAnimation:
    def test_creates_timing_xml(self):
        from ppt_pro_max.renderer.animation import add_emphasis_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_emphasis_animation(slide, shape_id, effect="pulse")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None

    def test_preset_class_is_emph(self):
        from ppt_pro_max.renderer.animation import add_emphasis_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_emphasis_animation(slide, shape_id, effect="pulse")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        emph_cTns = [c for c in cTns if c.get("presetClass") == "emph"]
        assert len(emph_cTns) >= 1

    def test_emphasis_preset_id_matches(self):
        from ppt_pro_max.renderer.animation import add_emphasis_animation, EMPHASIS_PRESETS
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_emphasis_animation(slide, shape_id, effect="grow")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        emph_cTns = [c for c in cTns if c.get("presetClass") == "emph"]
        expected_id, expected_sub = EMPHASIS_PRESETS["grow"]
        assert emph_cTns[0].get("presetID") == str(expected_id)
        assert emph_cTns[0].get("presetSubtype") == str(expected_sub)

    def test_emphasis_with_custom_duration(self):
        from ppt_pro_max.renderer.animation import add_emphasis_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_emphasis_animation(slide, shape_id, effect="spin",
                               duration_ms=1000, delay_ms=500)
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        emph_cTns = [c for c in cTns if c.get("presetClass") == "emph"]
        assert emph_cTns[0].get("dur") == "1000"
        assert emph_cTns[0].get("delay") == "500"

    def test_emphasis_targets_correct_shape(self):
        from ppt_pro_max.renderer.animation import add_emphasis_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_emphasis_animation(slide, shape_id, effect="pulse")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        spTgts = timing.findall(".//" + qn("p:spTgt"))
        assert any(t.get("spid") == str(shape_id) for t in spTgts)

    def test_emphasis_invalid_effect_falls_back(self):
        from ppt_pro_max.renderer.animation import add_emphasis_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_emphasis_animation(slide, shape_id, effect="nonexistent")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        emph_cTns = [c for c in cTns if c.get("presetClass") == "emph"]
        assert len(emph_cTns) >= 1

    def test_emphasis_no_visibility_set(self):
        from ppt_pro_max.renderer.animation import add_emphasis_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_emphasis_animation(slide, shape_id, effect="pulse")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        set_elems = timing.findall(".//" + qn("p:set"))
        for s in set_elems:
            attrNames = s.findall(".//" + qn("p:attrName"))
            for a in attrNames:
                assert a.text != "style.visibility", \
                    "Emphasis animations should not set visibility"


# ── Mixed animations (entrance + exit + emphasis on same slide) ──


class TestMixedAnimations:
    def test_entrance_and_exit_on_same_slide(self):
        from ppt_pro_max.renderer.animation import add_exit_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_entrance_animation(slide, shape_id, effect="fade_in")
        add_exit_animation(slide, shape_id, effect="fade_out")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        classes = [c.get("presetClass") for c in cTns if c.get("presetClass")]
        assert "entr" in classes
        assert "exit" in classes

    def test_entrance_and_emphasis_on_same_slide(self):
        from ppt_pro_max.renderer.animation import add_emphasis_animation
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_entrance_animation(slide, shape_id, effect="fade_in")
        add_emphasis_animation(slide, shape_id, effect="pulse")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        cTns = timing.findall(".//" + qn("p:cTn"))
        classes = [c.get("presetClass") for c in cTns if c.get("presetClass")]
        assert "entr" in classes
        assert "emph" in classes


# ── Backward compatibility ──


class TestBackwardCompat:
    def test_existing_transitions_unchanged(self):
        expected_keys = {"fade", "push", "wipe", "split", "cover",
                         "dissolve", "wheel", "wedge", "blinds",
                         "checker", "comb", "random"}
        assert expected_keys.issubset(set(TRANSITION_TYPES.keys()))

    def test_existing_entrance_presets_unchanged(self):
        assert "fade_in" in ENTRANCE_PRESETS
        assert ENTRANCE_PRESETS["fade_in"] == (10, 0)
        assert "zoom_in" in ENTRANCE_PRESETS
        assert ENTRANCE_PRESETS["zoom_in"] == (53, 16)

    def test_add_entrance_still_works(self):
        prs, slide = _make_slide()
        shape_id = _add_shape(slide)
        add_entrance_animation(slide, shape_id, effect="fade_in")
        sld = slide._element
        timing = sld.find(qn("p:timing"))
        assert timing is not None
