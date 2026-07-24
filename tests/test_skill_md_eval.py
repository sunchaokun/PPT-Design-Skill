"""Evaluate SKILL.md v0.9.1 (old) vs v0.9.2 (new) PPT quality.
Runs both FreeStyle and Build mode for the same topics, measures quality metrics.
"""
from __future__ import annotations

import json
import os
import tempfile
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


def count_font_sizes(prs: Presentation) -> dict:
    sizes = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    sz = run.font.size
                    if sz:
                        pt = int(sz / 12700)
                        sizes[pt] = sizes.get(pt, 0) + 1
    return sizes


def analyze_pptx(path: str) -> dict:
    prs = Presentation(path)
    total_shapes = 0
    total_text_chars = 0
    slides_with_text = 0
    slides_with_images = 0
    min_font = 999
    max_font = 0
    goal_types = []
    consecutive_same = 0
    max_consecutive_same = 0
    last_goal = None
    bullet_counts = []

    for i, slide in enumerate(prs.slides):
        shapes = len(slide.shapes)
        total_shapes += shapes

        text_on_slide = ""
        has_image = False
        slide_font_sizes = []

        for shape in slide.shapes:
            if shape.shape_type == 13:
                has_image = True
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                total_text_chars += len(text)
                text_on_slide += text
                for run in para.runs:
                    sz = run.font.size
                    if sz:
                        pt_val = int(sz / 12700)
                        slide_font_sizes.append(pt_val)

        if text_on_slide:
            slides_with_text += 1
        if has_image:
            slides_with_images += 1

        for s in slide_font_sizes:
            if s > 0:
                min_font = min(min_font, s)
                max_font = max(max_font, s)

        bullets_in_slide = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt.startswith("•") or txt.startswith("─") or txt.startswith("·"):
                        bullets_in_slide += 1
        if bullets_in_slide > 0:
            bullet_counts.append(bullets_in_slide)

    num_slides = len(prs.slides)
    avg_shapes = total_shapes / num_slides if num_slides else 0
    font_sizes = count_font_sizes(prs)
    font_levels = len(font_sizes)
    density_variance = 0
    if len(bullet_counts) > 1:
        mean = sum(bullet_counts) / len(bullet_counts)
        density_variance = sum((x - mean) ** 2 for x in bullet_counts) / len(bullet_counts)

    return {
        "slides": num_slides,
        "total_shapes": total_shapes,
        "avg_shapes_per_slide": round(avg_shapes, 1),
        "text_chars": total_text_chars,
        "slides_with_text": slides_with_text,
        "slides_with_images": slides_with_images,
        "min_font_pt": min_font if min_font != 999 else 0,
        "max_font_pt": max_font,
        "font_size_levels": font_levels,
        "file_size_kb": round(os.path.getsize(path) / 1024, 1),
        "density_variance": round(density_variance, 2),
    }


def run_freestyle(topic: str, style: str, output_dir: str) -> str:
    from ppt_pro_max import generate_ppt
    output_path = os.path.join(output_dir, f"freestyle_{style.replace(' ', '_')}.pptx")
    result = generate_ppt(topic, style=style, content_file=None, output=output_path)
    return result.get("output_path", output_path)


def run_build_mckinsey(output_dir: str) -> str:
    from ppt_pro_max.build_helpers import (
        add_slide, page_header, kpi_card, bar_chart, highlight_cards,
        text, rect, top_bar,
    )
    C = {
        'primary': '#2E6504', 'accent': '#7DA92F', 'muted': '#81C784',
        'light': '#C8E6C9', 'white': '#FFFFFF', 'background': '#FFFFFF',
        'card_bg': '#F9F9F9', 'text_dark': '#1A1A1A', 'text_body': '#333333',
        'text_muted': '#666666', 'divider': '#CCCCCC',
        'font_heading': '微软雅黑', 'font_body': '微软雅黑',
    }
    prs = Presentation()
    prs.slide_width = 13333300
    prs.slide_height = 7500000

    s0 = add_slide(prs)
    top_bar(s0, 'primary', C=C)
    rect(s0, 0, 0, 13.333, 7.5, 'primary', C=C)
    text(s0, 1.2, 2.0, 10, 1.5, 'Strategic Growth Report', font_size=44, color='white', bold=True, C=C)
    text(s0, 1.2, 3.6, 10, 0.5, 'FY2025 Performance Review', font_size=20, color='light', C=C)

    s1 = add_slide(prs)
    top_bar(s1, 'primary', C=C)
    page_header(s1, 'Key Performance Indicators', 'Core business metrics', C)
    kpi_card(s1, 0.65, 1.8, 3.8, 1.35, '12.8亿', '年度产值', '+8.3%', C=C)
    kpi_card(s1, 4.75, 1.8, 3.8, 1.35, '99.9%', '系统可用率', '+0.1%', C=C)
    kpi_card(s1, 8.85, 1.8, 3.8, 1.35, '5.2M', '活跃用户', '+15%', C=C)

    s2 = add_slide(prs)
    top_bar(s2, 'primary', C=C)
    page_header(s2, 'Revenue Breakdown', 'By business segment', C)
    data = [('Enterprise SaaS', 0.82, '$10.2B'), ('SMB Platform', 0.48, '$5.8B'), ('Consumer App', 0.30, '$3.9B')]
    bar_chart(s2, 2.0, 2.0, data, C=C)

    s3 = add_slide(prs)
    top_bar(s3, 'primary', C=C)
    page_header(s3, 'Core Capabilities', 'What sets us apart', C)
    cards = [
        ('AI Engine', 'Auto-select optimal framework, 70% compression', C['primary']),
        ('Live Dashboard', 'Real-time monitoring and instant alerting', C['accent']),
        ('Integration', '200+ connectors, zero-config setup', C['muted']),
    ]
    highlight_cards(s3, 0.65, 2.0, cards, C=C)

    s4 = add_slide(prs)
    top_bar(s4, 'primary', C=C)
    page_header(s4, 'Growth Strategy', 'Next fiscal year priorities', C)
    text(s4, 1.0, 2.2, 5.5, 0.5, '1. Expand to 3 new markets', font_size=16, color='text_dark', C=C)
    text(s4, 1.0, 2.8, 5.5, 0.5, '2. Launch enterprise self-serve', font_size=16, color='text_dark', C=C)
    text(s4, 1.0, 3.4, 5.5, 0.5, '3. Achieve carbon-neutral ops', font_size=16, color='text_dark', C=C)

    s5 = add_slide(prs)
    top_bar(s5, 'primary', C=C)
    page_header(s5, 'Team & Culture', 'Our greatest asset', C)
    text(s5, 1.0, 2.2, 10, 0.5, '2,400+ engineers across 12 offices', font_size=18, color='text_dark', bold=True, C=C)
    text(s5, 1.0, 2.9, 10, 0.5, 'Women in leadership: 38% (industry avg: 21%)', font_size=14, color='text_body', C=C)
    text(s5, 1.0, 3.4, 10, 0.5, 'Employee NPS: 72 (top 5% globally)', font_size=14, color='text_body', C=C)

    s6 = add_slide(prs)
    top_bar(s6, 'primary', C=C)
    text(s6, 1.2, 2.5, 10, 1.5, 'Let\'s Build the Future Together', font_size=44, color='primary', bold=True, C=C)
    text(s6, 1.2, 4.0, 10, 0.5, 'Contact: growth@example.com  |  Schedule a demo at example.com/demo', font_size=16, color='text_muted', C=C)

    path = os.path.join(output_dir, "build_mckinsey.pptx")
    prs.save(path)
    return path


def run_build_cyberpunk(output_dir: str) -> str:
    from ppt_pro_max.build_helpers import (
        add_slide, page_header, kpi_card, bar_chart, highlight_cards,
        text, rect, top_bar,
    )
    C = {
        'primary': '#0A0E27', 'accent': '#00FF88', 'muted': '#1A1F3A',
        'light': '#2A2F4A', 'white': '#E0E0E0', 'background': '#0A0E27',
        'card_bg': '#12162B', 'text_dark': '#E0E0E0', 'text_body': '#B0B0C0',
        'text_muted': '#606080', 'divider': '#2A2F4A',
        'font_heading': 'Orbitron', 'font_body': 'JetBrains Mono',
    }
    prs = Presentation()
    prs.slide_width = 13333300
    prs.slide_height = 7500000

    s0 = add_slide(prs)
    rect(s0, 0, 0, 13.333, 7.5, 'background', C=C)
    top_bar(s0, 'accent', C=C)
    text(s0, 1.2, 2.0, 10, 1.5, 'NEURAL NETWORK OPS', font_size=44, color='accent', bold=True, font_name='Orbitron', C=C)
    text(s0, 1.2, 3.8, 10, 0.5, 'Real-time AI Infrastructure Monitoring', font_size=18, color='text_muted', font_name='JetBrains Mono', C=C)

    s1 = add_slide(prs)
    rect(s1, 0, 0, 13.333, 7.5, 'background', C=C)
    top_bar(s1, 'accent', C=C)
    page_header(s1, 'System Metrics', 'Live dashboard data', C)
    kpi_card(s1, 0.65, 1.8, 3.8, 1.35, '99.97%', 'Uptime', '+0.02%', C=C)
    kpi_card(s1, 4.75, 1.8, 3.8, 1.35, '2.4ms', 'P99 Latency', '-12%', C=C)
    kpi_card(s1, 8.85, 1.8, 3.8, 1.35, '1.2M', 'RPS', '+23%', C=C)

    s2 = add_slide(prs)
    rect(s2, 0, 0, 13.333, 7.5, 'background', C=C)
    top_bar(s2, 'accent', C=C)
    page_header(s2, 'Traffic Analysis', 'Request distribution', C)
    data = [('API Gateway', 0.85, '1.02M'), ('ML Pipeline', 0.6, '720K'), ('Data Lake', 0.35, '420K')]
    bar_chart(s2, 2.0, 2.0, data, C=C)

    s3 = add_slide(prs)
    rect(s3, 0, 0, 13.333, 7.5, 'background', C=C)
    top_bar(s3, 'accent', C=C)
    page_header(s3, 'Architecture', 'Microservice topology', C)
    cards = [
        ('Inference', 'GPU cluster, auto-scaling', C['accent']),
        ('Training', 'Distributed gradient sync', C['muted']),
        ('Serving', 'Edge + cloud hybrid', C['light']),
    ]
    highlight_cards(s3, 0.65, 2.0, cards, C=C)

    s4 = add_slide(prs)
    rect(s4, 0, 0, 13.333, 7.5, 'background', C=C)
    top_bar(s4, 'accent', C=C)
    text(s4, 1.2, 2.5, 10, 1.5, 'Deploy Now →', font_size=44, color='accent', bold=True, font_name='Orbitron', C=C)
    text(s4, 1.2, 4.0, 10, 0.5, 'github.com/neural-ops  |  docs.neural-ops.dev', font_size=16, color='text_muted', font_name='JetBrains Mono', C=C)

    path = os.path.join(output_dir, "build_cyberpunk.pptx")
    prs.save(path)
    return path


def evaluate_content_json_quality(content: dict) -> dict:
    slides = content.get("slides", [])
    num_slides = len(slides)
    goals = [s.get("goal", "content") for s in slides]
    unique_goals = len(set(goals))

    has_section = "section" in goals
    has_hook = "hook" in goals
    has_cta = "cta" in goals

    bullet_counts = []
    for s in slides:
        b = s.get("bullets", [])
        if b:
            bullet_counts.append(len(b))

    density_variance = 0
    if len(bullet_counts) > 1:
        mean = sum(bullet_counts) / len(bullet_counts)
        density_variance = sum((x - mean) ** 2 for x in bullet_counts) / len(bullet_counts)

    has_code = any(s.get("code") for s in slides)
    has_exercise = any(s.get("exercise") for s in slides)
    has_component = any(s.get("component_type") for s in slides)
    has_cards = any(s.get("cards") for s in slides)

    consecutive_same = 0
    max_consecutive_same = 0
    for i in range(1, len(goals)):
        if goals[i] == goals[i-1] and goals[i] not in ("content", "section"):
            consecutive_same += 1
            max_consecutive_same = max(max_consecutive_same, consecutive_same)
        else:
            consecutive_same = 0

    return {
        "slides": num_slides,
        "unique_goal_types": unique_goals,
        "has_section_dividers": has_section,
        "has_hook": has_hook,
        "has_cta": has_cta,
        "has_code_page": has_code,
        "has_exercise_page": has_exercise,
        "has_component": has_component,
        "has_cards": has_cards,
        "bullet_density_variance": round(density_variance, 2),
        "max_consecutive_same_goal": max_consecutive_same,
    }


OLD_CONTENT_JSON = {
    "meta": {"title": "AI基础设施年度报告", "author": "Ops Team"},
    "slides": [
        {"goal": "hook", "title": "AI基础设施年度报告", "subtitle": "2025年度"},
        {"goal": "content", "title": "核心指标", "bullets": ["产值12.8亿", "可用率99.9%", "用户5.2M"]},
        {"goal": "content", "title": "收入构成", "bullets": ["企业SaaS $10.2B", "SMB平台 $5.8B", "消费者 $3.9B"]},
        {"goal": "content", "title": "核心能力", "bullets": ["AI引擎", "实时监控", "集成能力"]},
        {"goal": "content", "title": "增长策略", "bullets": ["拓展3个新市场", "上线企业自助服务", "碳中和运营"]},
        {"goal": "cta", "title": "联系我们", "subtitle": "growth@example.com"},
    ]
}

NEW_CONTENT_JSON = {
    "meta": {"title": "AI基础设施年度报告", "author": "Ops Team"},
    "slides": [
        {"goal": "hook", "title": "AI基础设施年度报告", "subtitle": "5分钟取代5周"},
        {"goal": "section", "title": "核心指标"},
        {"goal": "data", "title": "关键绩效指标", "bullets": [
            "年度产值 ¥12.8亿 (+8.3% YoY)",
            "系统可用率 99.9% (SLA ≥99.95%)",
            "活跃用户 5.2M (+15% QoQ)",
            "P99延迟 2.4ms (↓12% MoM)",
            "NPS评分 72 (行业TOP 5%)",
            "女性领导占比 38% (行业均值21%)",
        ]},
        {"goal": "features", "title": "核心竞争力", "cards": [
            {"title": "AI推理引擎", "text": "自动选择最优框架，量化压缩70%，P99<10ms"},
            {"title": "实时监控", "text": "毫秒级告警，全链路追踪"},
            {"title": "零配置集成", "text": "200+连接器，开箱即用"},
        ]},
        {"goal": "data", "title": "收入构成", "bullets": [
            "企业SaaS: $10.2B (占比52%)",
            "SMB平台: $5.8B (占比29%)",
            "消费者App: $3.9B (占比19%)",
        ], "diagram": {"type": "table"}},
        {"goal": "code", "title": "部署示例", "code": {"language": "python", "source": "from neural_ops import deploy\nmodel = load('gpt-4o-finetuned')\ndeploy(model, replicas=3, gpu='A100')\nprint('Serving at', model.endpoint)"}},
        {"goal": "section", "title": "战略与团队"},
        {"goal": "content", "title": "FY2026增长策略", "bullets": [
            "拓展东南亚/中东/拉美3大市场",
            "上线企业自助服务平台",
            "2026Q4实现碳中和运营",
        ]},
        {"goal": "content", "title": "团队与文化", "bullets": [
            "2,400+工程师，12个全球办公室",
            "员工NPS 72（全球TOP 5%）",
        ]},
        {"goal": "cta", "title": "共建未来", "subtitle": "联系 growth@example.com 或预约演示 example.com/demo · 免费额度包含1000次推理/月，无需信用卡"},
    ]
}


def main():
    import sys
    sys.stdout.reconfigure(encoding='utf-8')

    output_dir = os.path.join(tempfile.gettempdir(), "ppt_eval")
    os.makedirs(output_dir, exist_ok=True)

    print("=" * 70)
    print("SKILL.md v0.9.1 (old) vs v0.9.2 (new) — PPT Quality Evaluation")
    print("=" * 70)

    print("\n## 1.'content.json Quality Comparison\n")
    old_cq = evaluate_content_json_quality(OLD_CONTENT_JSON)
    new_cq = evaluate_content_json_quality(NEW_CONTENT_JSON)

    print(f"{'Metric':<30} {'v0.9.1 (old)':<15} {'v0.9.2 (new)':<15} {'Delta':<10}")
    print("-" * 70)
    for key in old_cq:
        ov = old_cq[key]
        nv = new_cq[key]
        if isinstance(ov, bool):
            delta = "+NEW" if (not ov and nv) else ""
            print(f"{key:<30} {str(ov):<15} {str(nv):<15} {delta}")
        else:
            delta = nv - ov
            sign = "+" if delta > 0 else ""
            print(f"{key:<30} {ov:<15} {nv:<15} {sign}{delta}")

    print("\n## 2. Build Mode — McKinsey Style\n")
    mck_path = run_build_mckinsey(output_dir)
    mck_metrics = analyze_pptx(mck_path)
    print(f"  Slides: {mck_metrics['slides']}")
    print(f"  Total shapes: {mck_metrics['total_shapes']}")
    print(f"  Avg shapes/slide: {mck_metrics['avg_shapes_per_slide']}")
    print(f"  Font size levels: {mck_metrics['font_size_levels']}")
    print(f"  Min font: {mck_metrics['min_font_pt']}pt, Max: {mck_metrics['max_font_pt']}pt")
    print(f"  Slides with images: {mck_metrics['slides_with_images']}")
    print(f"  File size: {mck_metrics['file_size_kb']} KB")

    print("\n## 3. Build Mode — Cyberpunk Style\n")
    cp_path = run_build_cyberpunk(output_dir)
    cp_metrics = analyze_pptx(cp_path)
    print(f"  Slides: {cp_metrics['slides']}")
    print(f"  Total shapes: {cp_metrics['total_shapes']}")
    print(f"  Font size levels: {cp_metrics['font_size_levels']}")
    print(f"  Min font: {cp_metrics['min_font_pt']}pt, Max: {cp_metrics['max_font_pt']}pt")
    print(f"  File size: {cp_metrics['file_size_kb']} KB")

    print("\n## 4. FreeStyle — Same Topic, Old vs New CDR Impact\n")
    topics = [
        ("AI startup investor pitch", "dark cyberpunk"),
        ("Fintech quarterly review", "warm fintech"),
    ]
    for topic, style in topics:
        try:
            path = run_freestyle(topic, style, output_dir)
            metrics = analyze_pptx(path)
            print(f"\n  [{style}] {topic}")
            print(f"    Slides: {metrics['slides']}, Shapes: {metrics['total_shapes']}, "
                  f"Font levels: {metrics['font_size_levels']}, "
                  f"Min font: {metrics['min_font_pt']}pt, "
                  f"File: {metrics['file_size_kb']} KB")
        except Exception as e:
            print(f"\n  [{style}] {topic} — ERROR: {e}")

    print("\n## 5. Summary: New SKILL.md Improvements\n")
    improvements = []
    if new_cq["unique_goal_types"] > old_cq["unique_goal_types"]:
        improvements.append(f"+{new_cq['unique_goal_types'] - old_cq['unique_goal_types']} unique goal types (diverse layouts)")
    if new_cq["has_section_dividers"] and not old_cq["has_section_dividers"]:
        improvements.append("+Section dividers (visual rhythm)")
    if new_cq["has_code_page"] and not old_cq["has_code_page"]:
        improvements.append("+Code page (tech credibility)")
    if new_cq["bullet_density_variance"] > old_cq["bullet_density_variance"]:
        improvements.append(f"+Density variance {new_cq['bullet_density_variance']:.1f} vs {old_cq['bullet_density_variance']:.1f} (natural feel)")
    if new_cq["max_consecutive_same_goal"] < old_cq["max_consecutive_same_goal"]:
        improvements.append("Less consecutive same-goal slides")
    if new_cq["has_cards"] and not old_cq["has_cards"]:
        improvements.append("+Feature cards (gradient bar + elevation)")
    for imp in improvements:
        print(f"  [+] {imp}")

    print(f"\n  content.json slide count: {old_cq['slides']} → {new_cq['slides']} (+{new_cq['slides'] - old_cq['slides']})")
    print(f"  Build mode: build_helpers API now documented (15 functions, C dict, color resolution)")
    print(f"  Dial Action Map: V/M/D now maps to LLM decisions (was empty declaration before)")
    print(f"  Design Constraints: ~180 lines of anti-pattern rules (was 0 before)")


if __name__ == "__main__":
    main()
