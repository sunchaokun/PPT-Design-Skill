"""Tests for Phase 4: 3D shapes, pattern fill, frosted glass."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ── 3D Shape tests ──


class TestShape3D:
    def test_defaults(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        s = Shape3D()
        assert s.depth_pt == 10.0
        assert s.bevel_top_w == 4.0
        assert s.bevel_top_h == 2.0
        assert s.bevel_bottom_w == 4.0
        assert s.bevel_bottom_h == 2.0
        assert s.material == "powder"
        assert s.extrusion_color == "#000000"
        assert s.extrusion_alpha == 40
        assert s.contour_color is None
        assert s.contour_width_pt == 0.5
        assert s.light_rig == "threePt"
        assert s.light_dir == "t"
        assert s.camera == "perspectiveFront"

    def test_custom_values(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        s = Shape3D(depth_pt=20, material="metal", extrusion_color="#FF0000",
                    light_rig="balanced", light_dir="br", camera="orthoFront")
        assert s.depth_pt == 20
        assert s.material == "metal"
        assert s.extrusion_color == "#FF0000"
        assert s.light_rig == "balanced"
        assert s.light_dir == "br"
        assert s.camera == "orthoFront"

    def test_apply_creates_sp3d(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        s3d = Shape3D(depth_pt=15, material="metal")
        s3d.apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d is not None
        assert sp3d.get("z") == str(int(15 * 12700))

    def test_apply_creates_bevel(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        s3d = Shape3D(bevel_top_w=6, bevel_top_h=3)
        s3d.apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        bevelT = sp3d.find(qn("a:bevelT"))
        assert bevelT is not None
        assert bevelT.get("w") == str(int(6 * 12700))
        assert bevelT.get("h") == str(int(3 * 12700))

    def test_apply_creates_bottom_bevel(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        s3d = Shape3D(bevel_bottom_w=5, bevel_bottom_h=2.5)
        s3d.apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        bevelB = sp3d.find(qn("a:bevelB"))
        assert bevelB is not None
        assert bevelB.get("w") == str(int(5 * 12700))
        assert bevelB.get("h") == str(int(2.5 * 12700))

    def test_apply_creates_material(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        s3d = Shape3D(material="plastic")
        s3d.apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        prstMat = sp3d.find(qn("a:prstMaterial"))
        assert prstMat is not None
        assert prstMat.get("val") == "plastic"

    def test_apply_creates_extrusion_color(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        s3d = Shape3D(extrusion_color="#3366FF", extrusion_alpha=50)
        s3d.apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        extrusionClr = sp3d.find(qn("a:extrusionClr"))
        assert extrusionClr is not None
        srgb = extrusionClr.find(qn("a:srgbClr"))
        assert srgb is not None
        assert srgb.get("val") == "3366FF"
        alpha = srgb.find(qn("a:alpha"))
        assert alpha is not None
        assert alpha.get("val") == str(50 * 1000)

    def test_apply_creates_contour(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        s3d = Shape3D(contour_color="#888888", contour_width_pt=1.0)
        s3d.apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        contourClr = sp3d.find(qn("a:contourClr"))
        assert contourClr is not None
        srgb = contourClr.find(qn("a:srgbClr"))
        assert srgb is not None
        assert srgb.get("val") == "888888"
        assert sp3d.get("contourW") == str(int(1.0 * 12700))

    def test_apply_no_contour_when_none(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        s3d = Shape3D(contour_color=None)
        s3d.apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        contourClr = sp3d.find(qn("a:contourClr"))
        assert contourClr is None

    def test_apply_creates_scene3d(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        s3d = Shape3D(light_rig="soft", light_dir="l", camera="isometricOffAxis1Top")
        s3d.apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        scene3d = spPr.find(qn("a:scene3d"))
        assert scene3d is not None
        camera = scene3d.find(qn("a:camera"))
        assert camera is not None
        assert camera.get("prst") == "isometricOffAxis1Top"
        lightRig = scene3d.find(qn("a:lightRig"))
        assert lightRig is not None
        assert lightRig.get("rig") == "soft"
        assert lightRig.get("dir") == "l"

    def test_apply_replaces_existing_3d(self):
        from ppt_pro_max.renderer.visual_effects import Shape3D
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        Shape3D(depth_pt=5).apply(sh)
        Shape3D(depth_pt=20).apply(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d_els = spPr.findall(qn("a:sp3d"))
        scene3d_els = spPr.findall(qn("a:scene3d"))
        assert len(sp3d_els) == 1
        assert len(scene3d_els) == 1
        assert sp3d_els[0].get("z") == str(int(20 * 12700))


class TestApply3D:
    def test_apply_3d_convenience(self):
        from ppt_pro_max.renderer.visual_effects import apply_3d
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_3d(sh, depth_pt=12, material="metal", extrusion_color="#FF0000")
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d is not None
        assert sp3d.get("z") == str(int(12 * 12700))
        prstMat = sp3d.find(qn("a:prstMaterial"))
        assert prstMat.get("val") == "metal"

    def test_apply_3d_default_args(self):
        from ppt_pro_max.renderer.visual_effects import apply_3d
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_3d(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d is not None
        assert sp3d.get("z") == str(int(10 * 12700))


class TestApplyBevel:
    def test_apply_bevel(self):
        from ppt_pro_max.renderer.visual_effects import apply_bevel
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_bevel(sh, top_w=6, top_h=3, material="plastic")
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d is not None
        assert sp3d.get("z") == "0"
        bevelT = sp3d.find(qn("a:bevelT"))
        assert bevelT is not None
        assert bevelT.get("w") == str(int(6 * 12700))
        assert bevelT.get("h") == str(int(3 * 12700))
        prstMat = sp3d.find(qn("a:prstMaterial"))
        assert prstMat.get("val") == "plastic"

    def test_apply_bevel_default(self):
        from ppt_pro_max.renderer.visual_effects import apply_bevel
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_bevel(sh)
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d.get("z") == "0"
        bevelT = sp3d.find(qn("a:bevelT"))
        assert bevelT.get("w") == str(int(4 * 12700))
        assert bevelT.get("h") == str(int(2 * 12700))


# ── Pattern Fill tests ──


class TestPatternTypes:
    def test_pattern_types_dict_exists(self):
        from ppt_pro_max.renderer.visual_effects import PATTERN_TYPES
        assert isinstance(PATTERN_TYPES, dict)
        assert len(PATTERN_TYPES) >= 15

    def test_common_patterns(self):
        from ppt_pro_max.renderer.visual_effects import PATTERN_TYPES
        assert "cross" in PATTERN_TYPES
        assert "dark_downward_diagonal" in PATTERN_TYPES
        assert "small_checker" in PATTERN_TYPES
        assert "weave" in PATTERN_TYPES
        assert "zigzag" in PATTERN_TYPES
        assert "dotted_grid" in PATTERN_TYPES

    def test_ooxml_abbreviations(self):
        from ppt_pro_max.renderer.visual_effects import PATTERN_TYPES
        assert PATTERN_TYPES["dark_downward_diagonal"] == "dkDnDiag"
        assert PATTERN_TYPES["small_checker"] == "smCheck"
        assert PATTERN_TYPES["light_horizontal"] == "ltHorz"
        assert PATTERN_TYPES["dotted_grid"] == "dotGrid"
        assert PATTERN_TYPES["dotted_diamond"] == "dotDmnd"


class TestApplyPatternFill:
    def test_creates_pattFill(self):
        from ppt_pro_max.renderer.visual_effects import apply_pattern_fill
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_pattern_fill(sh, "cross", fg_color="#FF0000", bg_color="#FFFFFF")
        spPr = sh._element.find(qn("p:spPr"))
        pattFill = spPr.find(qn("a:pattFill"))
        assert pattFill is not None
        assert pattFill.get("prst") == "cross"

    def test_fg_color(self):
        from ppt_pro_max.renderer.visual_effects import apply_pattern_fill
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_pattern_fill(sh, "cross", fg_color="#FF0000", bg_color="#FFFFFF")
        spPr = sh._element.find(qn("p:spPr"))
        pattFill = spPr.find(qn("a:pattFill"))
        fgClr = pattFill.find(qn("a:fgClr"))
        assert fgClr is not None
        srgb = fgClr.find(qn("a:srgbClr"))
        assert srgb.get("val") == "FF0000"

    def test_bg_color(self):
        from ppt_pro_max.renderer.visual_effects import apply_pattern_fill
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_pattern_fill(sh, "cross", fg_color="#FF0000", bg_color="#FFFFFF")
        spPr = sh._element.find(qn("p:spPr"))
        pattFill = spPr.find(qn("a:pattFill"))
        bgClr = pattFill.find(qn("a:bgClr"))
        assert bgClr is not None
        srgb = bgClr.find(qn("a:srgbClr"))
        assert srgb.get("val") == "FFFFFF"

    def test_replaces_existing_fill(self):
        from ppt_pro_max.renderer.visual_effects import apply_pattern_fill
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor_from_str("0000FF")
        apply_pattern_fill(sh, "cross", fg_color="#FF0000", bg_color="#FFFFFF")
        spPr = sh._element.find(qn("p:spPr"))
        solidFills = spPr.findall(qn("a:solidFill"))
        assert len(solidFills) == 0
        pattFill = spPr.find(qn("a:pattFill"))
        assert pattFill is not None

    def test_uses_ooxml_abbreviation(self):
        from ppt_pro_max.renderer.visual_effects import apply_pattern_fill
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_pattern_fill(sh, "dark_downward_diagonal", fg_color="#333333", bg_color="#CCCCCC")
        spPr = sh._element.find(qn("p:spPr"))
        pattFill = spPr.find(qn("a:pattFill"))
        assert pattFill.get("prst") == "dkDnDiag"

    def test_invalid_pattern_raises(self):
        from ppt_pro_max.renderer.visual_effects import apply_pattern_fill
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        with pytest.raises(KeyError, match="Unknown pattern"):
            apply_pattern_fill(sh, "nonexistent_pattern", fg_color="#FF0000", bg_color="#FFFFFF")

    def test_fg_alpha(self):
        from ppt_pro_max.renderer.visual_effects import apply_pattern_fill
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_pattern_fill(sh, "cross", fg_color="#FF0000", bg_color="#FFFFFF",
                           fg_alpha=30)
        spPr = sh._element.find(qn("p:spPr"))
        pattFill = spPr.find(qn("a:pattFill"))
        fgClr = pattFill.find(qn("a:fgClr"))
        srgb = fgClr.find(qn("a:srgbClr"))
        alpha = srgb.find(qn("a:alpha"))
        assert alpha is not None
        assert alpha.get("val") == str(30 * 1000)


# ── Frosted Glass tests ──


class TestApplyFrostedGlass:
    def test_creates_solidFill_with_alpha(self):
        from ppt_pro_max.renderer.visual_effects import apply_frosted_glass
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_frosted_glass(sh, tint_color="#FFFFFF", tint_alpha=15, soft_edge=8)
        spPr = sh._element.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb is not None
        assert srgb.get("val") == "FFFFFF"
        alpha = srgb.find(qn("a:alpha"))
        assert alpha is not None
        assert alpha.get("val") == str(15 * 1000)

    def test_creates_soft_edge(self):
        from ppt_pro_max.renderer.visual_effects import apply_frosted_glass
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_frosted_glass(sh, tint_color="#FFFFFF", tint_alpha=15, soft_edge=8)
        spPr = sh._element.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        assert effectLst is not None
        softEdge = effectLst.find(qn("a:softEdge"))
        assert softEdge is not None
        assert softEdge.get("rad") == str(int(8 * 12700))

    def test_custom_tint_color(self):
        from ppt_pro_max.renderer.visual_effects import apply_frosted_glass
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_frosted_glass(sh, tint_color="#3366FF", tint_alpha=20, soft_edge=5)
        spPr = sh._element.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        srgb = solidFill.find(qn("a:srgbClr"))
        assert srgb.get("val") == "3366FF"
        alpha = srgb.find(qn("a:alpha"))
        assert alpha.get("val") == str(20 * 1000)

    def test_zero_soft_edge(self):
        from ppt_pro_max.renderer.visual_effects import apply_frosted_glass
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        apply_frosted_glass(sh, tint_color="#FFFFFF", tint_alpha=15, soft_edge=0)
        spPr = sh._element.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        if effectLst is not None:
            softEdge = effectLst.find(qn("a:softEdge"))
            assert softEdge is None

    def test_replaces_existing_fill(self):
        from ppt_pro_max.renderer.visual_effects import apply_frosted_glass
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1), Inches(1), Inches(3), Inches(2))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor_from_str("0000FF")
        apply_frosted_glass(sh, tint_color="#FFFFFF", tint_alpha=15, soft_edge=8)
        spPr = sh._element.find(qn("p:spPr"))
        solidFills = spPr.findall(qn("a:solidFill"))
        assert len(solidFills) == 1
        srgb = solidFills[0].find(qn("a:srgbClr"))
        assert srgb.get("val") == "FFFFFF"


# ── API exposure tests (PrecisionRenderer) ──


class TestPrecisionRenderer3DAPI:
    def _make_renderer(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = BrandSpec(
            colors={
                "primary": "#1D4ED8", "on-primary": "#FFFFFF",
                "secondary": "#64748B", "accent": "#F59E0B",
                "background": "#FFFFFF", "foreground": "#1A1A1A",
                "muted": "#F0F4F8", "muted-foreground": "#6B7B8D",
                "border": "#DEE5EF", "destructive": "#DC2626",
            },
            fonts={"heading": "Georgia", "body": "Calibri"},
        )
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        slide = pr.add_slide(prs)
        return pr, prs, slide

    def test_add_3d_shape(self):
        pr, prs, slide = self._make_renderer()
        sh = pr.add_3d_shape(slide, 1, 1, 3, 2, depth_pt=15, material="metal")
        assert sh is not None
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d is not None

    def test_add_bevel_shape(self):
        pr, prs, slide = self._make_renderer()
        sh = pr.add_bevel_shape(slide, 1, 1, 3, 2, top_w=6, top_h=3)
        assert sh is not None
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d is not None
        bevelT = sp3d.find(qn("a:bevelT"))
        assert bevelT is not None

    def test_add_pattern_fill_shape(self):
        pr, prs, slide = self._make_renderer()
        sh = pr.add_pattern_fill_shape(slide, 1, 1, 3, 2,
                                        pattern_type="cross",
                                        fg_color="#FF0000", bg_color="#FFFFFF")
        assert sh is not None
        spPr = sh._element.find(qn("p:spPr"))
        pattFill = spPr.find(qn("a:pattFill"))
        assert pattFill is not None

    def test_add_frosted_panel(self):
        pr, prs, slide = self._make_renderer()
        sh = pr.add_frosted_panel(slide, 1, 1, 3, 2,
                                   tint_color="#FFFFFF", tint_alpha=15)
        assert sh is not None
        spPr = sh._element.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None


# ── API exposure tests (build_helpers) ──


class TestBuildHelpers3DAPI:
    def test_shape_3d(self):
        from ppt_pro_max.build_helpers import shape_3d
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = shape_3d(slide, 1, 1, 3, 2, depth=15, material="metal")
        assert sh is not None
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d is not None

    def test_bevel_shape(self):
        from ppt_pro_max.build_helpers import bevel_shape
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = bevel_shape(slide, 1, 1, 3, 2, top_w=6, top_h=3)
        assert sh is not None
        spPr = sh._element.find(qn("p:spPr"))
        sp3d = spPr.find(qn("a:sp3d"))
        assert sp3d is not None

    def test_pattern_fill(self):
        from ppt_pro_max.build_helpers import pattern_fill
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = pattern_fill(slide, 1, 1, 3, 2, "cross", "#FF0000", "#FFFFFF")
        assert sh is not None
        spPr = sh._element.find(qn("p:spPr"))
        pattFill = spPr.find(qn("a:pattFill"))
        assert pattFill is not None

    def test_frosted_panel(self):
        from ppt_pro_max.build_helpers import frosted_panel
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = frosted_panel(slide, 1, 1, 3, 2, tint="#FFFFFF", alpha=15)
        assert sh is not None
        spPr = sh._element.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None


def RGBColor_from_str(hex_str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str)
