"""Tests for Speaker Notes feature (v0.4 I-A)."""

from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from pptx import Presentation


class TestSpeakerNotes:

    def test_notes_written_to_slide(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        pipeline._populate_slide(
            slide,
            {"title": "Test", "notes": "This is a speaker note"},
            prs,
        )
        assert slide.has_notes_slide
        assert slide.notes_slide.notes_text_frame.text == "This is a speaker note"

    def test_no_notes_no_notes_slide(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        pipeline._populate_slide(slide, {"title": "No Notes"}, prs)
        assert not slide.has_notes_slide

    def test_multi_paragraph_notes(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        pipeline._populate_slide(
            slide,
            {"title": "Multi", "notes": "Line 1\nLine 2\nLine 3"},
            prs,
        )
        assert slide.has_notes_slide
        tf = slide.notes_slide.notes_text_frame
        assert len(tf.paragraphs) == 3

    def test_auto_notes_in_story_plan(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        story_plan = pipeline._build_story_plan(
            query="AI Startup",
            page_contents=[],
            density=4,
            business_mode=None,
        )
        pages = story_plan.get("pages", [])
        assert len(pages) > 0
        for page in pages:
            assert "notes" in page
            assert page["notes"] is not None
            assert len(page["notes"]) > 0

    def test_content_json_notes_passthrough(self, tmp_path):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content
        content = {
            "meta": {"title": "Test"},
            "slides": [
                {"goal": "hook", "title": "Welcome", "notes": "Opening remark"},
                {"goal": "content", "title": "Details"},
            ],
        }
        result = load_enterprise_content(content, str(tmp_path))
        assert result[0].get("notes") == "Opening remark"
        assert result[1].get("notes") is None

    def test_review_shows_notes_preview(self):
        from ppt_pro_max.enterprise.review_gate import ReviewGate
        gate = ReviewGate()
        proposal = {
            "pages": [
                {"goal": "hook", "title": "Welcome", "notes": "This is a long note that should be truncated"},
                {"goal": "content", "title": "No notes here"},
            ],
        }
        display = gate.format_cli(proposal)
        assert "This is a long note" in display
