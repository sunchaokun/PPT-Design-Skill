"""Generate PPT with specific best hierarchy component."""

import json
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
from lxml import etree

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

GREEN_BRAND = BrandSpec(colors={
    "primary": "#00AA00", "accent": "#66CC00", "muted": "#E6F5E6",
    "foreground": "#1A3A1A", "on-primary": "#FFFFFF", "background": "#FFFFFF",
})

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)

lib = ComponentLibrary(find_db_path())
renderer = ComponentRenderer()
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Test multiple hierarchy candidates
for cid in [4458, 3557, 4037, 4081]:
    xml_parts = lib.load_xml(cid)
    if not xml_parts or "group" not in xml_parts:
        continue

    root = etree.fromstring(xml_parts["group"])
    texts = [t.text.strip()[:15] for t in root.iter(f"{{{_A}}}t") if t.text and t.text.strip()]
    custGeom = len(list(root.iter(f"{{{_A}}}custGeom")))
    sp = len(list(root.iter(f"{{{_P}}}sp")))

    filled = renderer._fill_group_data(xml_parts, ["CEO", "CTO", "CFO", "VP Eng", "VP Sales"])
    styled = renderer._apply_brand_colors(filled, GREEN_BRAND)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"hierarchy id={cid} (sp={sp}, custGeom={custGeom})"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0xAA, 0)

    meta_bytes = styled.get("_meta")
    orig_aspect = 1.0
    if meta_bytes:
        meta = json.loads(meta_bytes)
        obe = meta.get("orig_bounds_emu")
        if obe and len(obe) == 4 and obe[3] > 0:
            orig_aspect = obe[2] / obe[3]

    bounds = renderer.compute_component_bounds((0.5, 0.9, 12.3, 5.8), orig_aspect, "contain")
    renderer._inject_group_to_slide(slide, styled, bounds)
    print(f"  id={cid}: sp={sp}, custGeom={custGeom}, texts={len(texts)}, aspect={orig_aspect:.2f}")

out = os.path.join(OUT_DIR, "hierarchy_candidates.pptx")
prs.save(out)
print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")
lib.close()
