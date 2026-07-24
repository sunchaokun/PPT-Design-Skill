"""Quick shape-level analysis of Build mode PPTs."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree

def analyze_deep(path, label):
    prs = Presentation(path)
    print(f"\n{'='*80}")
    print(f"  {label} | {len(prs.slides)} slides | {round(os.path.getsize(path)/1024,1)} KB")
    print(f"{'='*80}")

    all_fonts, all_colors, all_sizes = set(), set(), []
    for i, slide in enumerate(prs.slides):
        shapes = len(slide.shapes)
        slide_info = []
        for shape in slide.shapes:
            info = ""
            if shape.has_text_frame:
                txt = shape.text_frame.text[:60].replace('\n', ' ')
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            sz = run.font.size
                            if sz:
                                pt = int(sz / 12700)
                                all_sizes.append(pt)
                            try:
                                if run.font.color and run.font.color.rgb:
                                    all_colors.add(str(run.font.color.rgb))
                            except: pass
                            if run.font.name: all_fonts.add(run.font.name)
                info = f'"{txt}"'
            else:
                stype = shape.shape_type
                try:
                    fill_c = str(shape.fill.fore_color.rgb) if shape.fill.type == 1 else "?"
                except:
                    fill_c = "?"
                info = f"[{stype}] fill=#{fill_c}"
            slide_info.append(f"    {info}")
        
        print(f"  Slide {i} ({shapes} shapes):")
        for s in slide_info[:8]:
            print(s)
        if len(slide_info) > 8:
            print(f"    ... +{len(slide_info)-8} more")

    if all_sizes:
        print(f"\n  Font sizes: {min(all_sizes)}-{max(all_sizes)}pt ({len(set(all_sizes))} levels)")
    print(f"  Fonts: {sorted(all_fonts)}")
    print(f"  Colors: {sorted(all_colors)} ({len(all_colors)} unique)")


out = os.path.join(tempfile.gettempdir(), "ppt_build_test")
for name, label in [("build_mckinsey.pptx", "McKinsey"), ("build_cyberpunk.pptx", "Cyberpunk"), ("build_creative.pptx", "Creative")]:
    p = os.path.join(out, name)
    if os.path.exists(p):
        analyze_deep(p, label)
