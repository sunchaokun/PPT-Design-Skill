import io
import os
import sys
import tempfile
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


@pytest.fixture
def simple_template(tmp_path):
    pptx_path = str(tmp_path / "template.pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s1 = prs.slides.add_slide(blank)

    tf = s0.shapes.title.text_frame
    tf.paragraphs[0].text = "封面标题"
    run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
    run.text = "封面标题"
    run.font.size = Pt(44)
    run.font.color.rgb = RGBColor(0x2E, 0x65, 0x04)
    run.font.name = "微软雅黑"

    sh = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x7D, 0xA9, 0x2F)
    sh.line.fill.background()

    tb = s1.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "正文标题"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0x0D, 0x46, 0x09)
    p.font.bold = True

    prs.save(pptx_path)
    return pptx_path


@pytest.fixture
def template_with_various_shapes(tmp_path):
    pptx_path = str(tmp_path / "complex.pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    tf = s0.shapes.title.text_frame
    tf.paragraphs[0].text = "汇报报告"

    s1 = prs.slides.add_slide(blank)

    sh = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x4C, 0xAF, 0x50)
    sh.line.fill.background()

    sh2 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(1), Inches(1), Inches(1))
    sh2.fill.solid()
    sh2.fill.fore_color.rgb = RGBColor(0xFF, 0x57, 0x22)
    sh2.line.fill.background()

    tb = s1.shapes.add_textbox(Inches(1), Inches(4), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "这是一段正文内容，用于测试"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(1.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "谢谢观看"
    p2.font.size = Pt(36)

    prs.save(pptx_path)
    return pptx_path


class TestAnalyzeTemplate:

    def test_basic_info(self, simple_template):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(simple_template)
        assert "Size: 13.333 x 7.500 inches" in result
        assert "Slides: 2" in result

    def test_slide_type_cover(self, simple_template):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(simple_template)
        assert "[cover]" in result

    def test_slide_type_back_cover(self, template_with_various_shapes):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(template_with_various_shapes)
        assert "[back_cover]" in result

    def test_color_frequency(self, simple_template):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(simple_template)
        assert "COLOR FREQUENCY" in result
        assert "#2E6504" in result
        assert "#7DA92F" in result

    def test_font_frequency(self, simple_template):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(simple_template)
        assert "FONT FREQUENCY" in result

    def test_shape_info_per_slide(self, simple_template):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(simple_template)
        assert "Slide 0" in result
        assert "Slide 1" in result

    def test_text_content_shown(self, simple_template):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(simple_template)
        assert "封面标题" in result or "正文标题" in result

    def test_geom_type_shown(self, template_with_various_shapes):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(template_with_various_shapes)
        assert "geom=roundRect" in result or "geom=ellipse" in result

    def test_theme_colors(self, simple_template):
        from ppt_pro_max.analyze_template import analyze
        result = analyze(simple_template)
        assert "THEME COLORS" in result

    def test_real_template(self):
        path = r"E:\PPT-Design-Skill\docs\分析脚本\template.pptx"
        if not os.path.isfile(path):
            pytest.skip("Real template not available")
        from ppt_pro_max.analyze_template import analyze
        result = analyze(path)
        assert "Size:" in result
        assert "Slides:" in result
        assert "COLOR FREQUENCY" in result
        assert "FONT FREQUENCY" in result
        assert "#2E6504" in result or "#7DA92F" in result
