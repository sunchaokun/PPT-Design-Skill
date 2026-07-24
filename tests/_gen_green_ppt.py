"""Generate a PPT with real library components + green brand colors for visual inspection."""

import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

GREEN_BRAND = BrandSpec(colors={
    "primary": "#00AA00",
    "accent": "#66CC00",
    "muted": "#E6F5E6",
    "foreground": "#1A3A1A",
    "on-primary": "#FFFFFF",
    "background": "#FFFFFF",
})

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)

lib = ComponentLibrary(find_db_path())
renderer = ComponentRenderer()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

categories = [
    ("hierarchy", ["CEO", "CTO", "CFO", "VP Eng", "VP Sales"], "组织架构"),
    ("process", ["Plan", "Design", "Build", "Launch"], "流程图"),
    ("swot", ["Strengths", "Weaknesses", "Opportunities", "Threats"], "SWOT分析"),
    ("infographic", ["Revenue", "Growth", "Profit", "Market Share"], "信息图"),
    ("timeline", ["Q1", "Q2", "Q3", "Q4"], "时间线"),
]

for cat, texts, title_cn in categories:
    results = lib.search(type="group", category=cat, limit=1)
    if not results:
        print(f"  {cat}: SKIPPED (no components)")
        continue
    comp = results[0]
    xml_parts = lib.load_xml(comp["id"])

    filled = renderer._fill_group_data(xml_parts, texts)
    styled = renderer._apply_brand_colors(filled, GREEN_BRAND)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{title_cn} ({cat}) — Green Brand"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0xAA, 0x00)

    meta_bytes = styled.get("_meta")
    orig_aspect = 1.0
    if meta_bytes:
        meta = json.loads(meta_bytes)
        obe = meta.get("orig_bounds_emu")
        if obe and len(obe) == 4 and obe[3] > 0:
            orig_aspect = obe[2] / obe[3]

    content_area = (0.5, 1.3, 12.3, 5.5)
    bounds = renderer.compute_component_bounds(content_area, orig_aspect, "contain")
    renderer._inject_group_to_slide(slide, styled, bounds)

    print(f"  {cat}: id={comp['id']}, aspect={orig_aspect:.2f}, bounds={bounds}")

out = os.path.join(OUT_DIR, "green_brand_components.pptx")
prs.save(out)
print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")
print(f"Slides: {len(prs.slides)}")

lib.close()
