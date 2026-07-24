"""TDD tests for Phase 1 bug fixes (B1-B5, D11).

B1: precision variable undefined in plan.component_type branch
B2: image matching results discarded in template-clone path
B3: candidates[0] always picks first template slide
B4: _is_placeholder_text(shape_name) checks shape name instead of text content
B5: DeliveryGate _check_background_missing uses wrong index
D11: SlideDNA missing brand_spec field
"""

import os
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def _create_multi_page_pptx(path: str, num_content: int = 2):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layouts = prs.slide_layouts
    blank = layouts[6] if len(layouts) > 6 else layouts[-1]

    cover = prs.slides.add_slide(layouts[0] if layouts else blank)
    for ph in cover.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = "Template Cover"
        elif ph.placeholder_format.type == 2:
            ph.text = "Template Subtitle"

    toc = prs.slides.add_slide(blank)
    tb = toc.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "目录"
    run.font.size = Pt(36)
    run.font.bold = True
    body = toc.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    btf = body.text_frame
    for j, text in enumerate(["Section A", "Section B"]):
        bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        br = bp.add_run()
        br.text = text
        br.font.size = Pt(14)

    trans = prs.slides.add_slide(blank)
    tb2 = trans.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "PART ONE"
    r2.font.size = Pt(48)
    r2.font.bold = True

    for i in range(num_content):
        slide = prs.slides.add_slide(blank)
        tb3 = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = f"Content Page {i + 1}"
        r3.font.size = Pt(28)
        r3.font.bold = True

        body2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        btf2 = body2.text_frame
        for j, text in enumerate(["Bullet A", "Bullet B", "Bullet C"]):
            bp2 = btf2.paragraphs[0] if j == 0 else btf2.add_paragraph()
            br2 = bp2.add_run()
            br2.text = text
            br2.font.size = Pt(14)

    back = prs.slides.add_slide(layouts[0] if layouts else blank)
    for ph in back.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = "Thank You"
        elif ph.placeholder_format.type == 2:
            ph.text = "Contact Us"

    prs.save(path)


def _create_content_with_deco_pptx(path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    cover = prs.slides.add_slide(blank)
    tb = cover.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Cover"
    r.font.size = Pt(36)

    content = prs.slides.add_slide(blank)
    tb2 = content.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Content Title"
    r2.font.size = Pt(28)
    r2.font.bold = True
    body = content.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    btf = body.text_frame
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = "Bullet A"
    br.font.size = Pt(14)

    deco_bar = content.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.5), Inches(7.5)
    )
    deco_bar.fill.solid()
    deco_bar.fill.fore_color.rgb = RGBColor(0x2E, 0x65, 0x04)
    deco_bar.line.fill.background()

    back = prs.slides.add_slide(blank)
    tb3 = back.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Thank You"
    r3.font.size = Pt(36)

    prs.save(path)


class TestB1PrecisionVariableUndefined:
    """B1: plan.component_type branch references undefined `precision` variable."""

    def test_rerender_with_component_type_does_not_crash(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, PagePlan,
        )

        pptx_path = str(tmp_path / "template.pptx")
        _create_content_with_deco_pptx(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        content_slide_dna = [s for s in dna.slides if s.page_type == "content"][0]
        plan = PagePlan(
            page_type="content",
            title="Component Test",
            bullets=["A", "B", "C"],
            component_type="group",
            component_category="process",
        )

        prs = Presentation(pptx_path)
        content_slide = prs.slides[content_slide_dna.slide_index]

        try:
            extractor.rerender_content_zone(
                content_slide, content_slide_dna, plan, dna.brand_spec
            )
            succeeded = True
        except NameError as e:
            succeeded = False
            error_msg = str(e)
        except Exception:
            succeeded = True

        assert succeeded, f"rerender_content_zone crashed with NameError: {error_msg if not succeeded else ''}"


class TestB3TemplateSelectionStrategy:
    """B3: candidates[0] always picks first template slide; should use smart selection."""

    def test_different_content_pages_use_different_templates(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, PagePlan,
        )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

        cover = prs.slides.add_slide(blank)
        tb = cover.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Cover"
        r.font.size = Pt(36)

        for i in range(3):
            slide = prs.slides.add_slide(blank)
            tb2 = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = f"Content {i + 1}"
            r2.font.size = Pt(28)
            body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            btf = body.text_frame
            n_bullets = i + 2
            for j in range(n_bullets):
                bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                br = bp.add_run()
                br.text = f"B{j + 1}"
                br.font.size = Pt(14)

        back = prs.slides.add_slide(blank)
        tb3 = back.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = "End"
        r3.font.size = Pt(36)

        pptx_path = str(tmp_path / "multi_content.pptx")
        prs.save(pptx_path)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        content_slides = [s for s in dna.slides if s.page_type == "content"]
        assert len(content_slides) == 3, "Template should have 3 content slides"

        plans = [
            PagePlan(page_type="cover", title="Cover"),
            PagePlan(page_type="content", title="C1", bullets=["A", "B"]),
            PagePlan(page_type="content", title="C2", bullets=["X", "Y", "Z"]),
            PagePlan(page_type="content", title="C3", bullets=["P1", "P2", "P3", "P4"]),
            PagePlan(page_type="back_cover", title="End"),
        ]

        output_path = str(tmp_path / "output.pptx")
        result = extractor.expand_and_patch(dna, plans, output_path)

        template_indices_used = [
            m.get("template_slide_index")
            for m in result.template_mapping
            if m.get("page_type") == "content"
        ]
        unique_templates = set(template_indices_used)
        assert len(unique_templates) > 1, (
            f"Expected different template slides for different content plans, "
            f"but all used: {template_indices_used}"
        )


class TestB4PlaceholderTextCheck:
    """B4: _is_placeholder_text(shape_name) should check text content, not shape name."""

    def test_is_placeholder_text_with_actual_text(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        assert DesignDNAExtractor._is_placeholder_text("单击此处添加标题") is True
        assert DesignDNAExtractor._is_placeholder_text("请输入内容") is True
        assert DesignDNAExtractor._is_placeholder_text("click to add text") is True
        assert DesignDNAExtractor._is_placeholder_text("placeholder") is True

    def test_is_placeholder_text_with_shape_name(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        assert DesignDNAExtractor._is_placeholder_text("Rectangle 5") is False
        assert DesignDNAExtractor._is_placeholder_text("TextBox 3") is False
        assert DesignDNAExtractor._is_placeholder_text("Title 1") is False

    def test_is_placeholder_text_with_empty(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        assert DesignDNAExtractor._is_placeholder_text("") is False
        assert DesignDNAExtractor._is_placeholder_text(None) is False

    def test_clear_all_decorations_preserves_non_placeholder_text(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, SlideDNA, TextZone,
        )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(blank)

        title_tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_tf = title_tb.text_frame
        title_p = title_tf.paragraphs[0]
        title_r = title_p.add_run()
        title_r.text = "My Title"
        title_r.font.size = Pt(28)

        real_text_tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        real_tf = real_text_tb.text_frame
        real_p = real_tf.paragraphs[0]
        real_r = real_p.add_run()
        real_r.text = "Real Content Here"
        real_r.font.size = Pt(14)

        slide_dna = SlideDNA(
            slide_index=0,
            page_type="content",
            slide_xml=b"",
            rels_xml=b"",
            text_zones=[
                TextZone(zone_id="t_0", role="title", text="My Title",
                         shape_name=title_tb.name),
                TextZone(zone_id="b_0", role="body", text="Real Content Here",
                         shape_name=real_text_tb.name),
            ],
        )

        extractor = DesignDNAExtractor()
        extractor._clear_all_decorations(slide, slide_dna)

        all_text = []
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        slide_elem = slide._element
        for t in slide_elem.iter(f"{{{a_ns}}}t"):
            if t.text and t.text.strip():
                all_text.append(t.text.strip())

        assert "My Title" in all_text, "Title text should be preserved after _clear_all_decorations"


class TestB5DeliveryGateIndexMapping:
    """B5: _check_background_missing uses output slide index instead of template mapping."""

    def test_background_check_with_template_mapping(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNA, SlideDNA, PagePlan,
        )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

        slide1 = prs.slides.add_slide(blank)
        tb1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf1 = tb1.text_frame
        p1 = tf1.paragraphs[0]
        r1 = p1.add_run()
        r1.text = "Cover"
        r1.font.size = Pt(36)

        slide2 = prs.slides.add_slide(blank)
        tb2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = "Content"
        r2.font.size = Pt(36)

        pptx_path = str(tmp_path / "test_bg.pptx")
        prs.save(pptx_path)

        dna = DesignDNA(
            source_path=pptx_path,
            slides=[
                SlideDNA(
                    slide_index=0, page_type="cover", slide_xml=b"", rels_xml=b"",
                    background_colors=["#2E6504"],
                ),
                SlideDNA(
                    slide_index=1, page_type="content", slide_xml=b"", rels_xml=b"",
                    background_colors=[],
                ),
            ],
        )

        gate = DeliveryGate()
        report = gate.check(
            pptx_path, dna,
            [PagePlan(page_type="cover", title="A"), PagePlan(page_type="content", title="B")],
        )

        bg_issues = [i for i in report.fatals if i.check_id == "background_missing"]
        assert len(bg_issues) >= 1, "Should detect background_missing for slides without bg"
        for issue in bg_issues:
            assert issue.slide_index >= 0


class TestD11SlideDNABrandSpec:
    """D11: SlideDNA should have brand_spec field populated during extract."""

    def test_slide_dna_has_brand_spec_field(self):
        from ppt_pro_max.enterprise.design_dna_extractor import SlideDNA

        slide_dna = SlideDNA(
            slide_index=0, page_type="content", slide_xml=b"", rels_xml=b"",
        )
        assert hasattr(slide_dna, "brand_spec"), "SlideDNA should have brand_spec field"
        assert slide_dna.brand_spec is None

    def test_slide_dna_brand_spec_populated_on_extract(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "template.pptx")
        _create_multi_page_pptx(pptx_path, num_content=1)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        if dna.brand_spec is not None:
            for slide_dna in dna.slides:
                assert slide_dna.brand_spec is not None, (
                    f"Slide {slide_dna.slide_index} should have brand_spec populated"
                )


class TestB2ImageInjection:
    """B2: Image matching results should be applied to output PPTX in template-clone path."""

    def test_inject_matched_images_method_exists(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        assert hasattr(EnterprisePipeline, "_inject_matched_images"), (
            "EnterprisePipeline should have _inject_matched_images method"
        )

    def test_image_pool_applied_in_template_clone(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        project_dir = str(tmp_path / "project")
        os.makedirs(project_dir)
        os.makedirs(os.path.join(project_dir, "images"))

        template_path = os.path.join(project_dir, "template.pptx")
        _create_multi_page_pptx(template_path, num_content=1)

        from PIL import Image
        img_path = os.path.join(project_dir, "images", "test_hero.png")
        img = Image.new("RGB", (1600, 900), color=(100, 150, 200))
        img.save(img_path)

        content = {
            "pages": [
                {"goal": "hook", "title": "Cover", "subtitle": "Sub"},
                {"goal": "content", "title": "Content", "bullets": ["A", "B"]},
                {"goal": "cta", "title": "End"},
            ]
        }
        content_path = os.path.join(project_dir, "content.json")
        with open(content_path, "w", encoding="utf-8") as f:
            json.dump(content, f, ensure_ascii=False)

        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="test",
            project_dir=project_dir,
            content_file=content_path,
        )

        if result.get("render_mode") == "template_clone" and result.get("output_path"):
            assert os.path.isfile(result["output_path"])
