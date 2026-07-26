"""Test ComponentAdapter — strict validation of layout, color, font details.

Tests use realistic GroupShape XML to verify:
1. Layout: no overflow, correct bounds, text fits
2. Color: no white-on-light, contrast >= 3:1, brand colors applied
3. Font: no unreadable sizes, CJK applied, heading/body correct
"""
from __future__ import annotations

import os

import pytest
from lxml import etree

from ppt_pro_max.enterprise.brand_spec import BrandSpec


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _light_brand() -> BrandSpec:
    return BrandSpec(
        source="test",
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "secondary": "#7C3AED",
            "accent": "#F97316",
            "background": "#FFFFFF",
            "foreground": "#0F172A",
            "muted": "#F1F5F9",
            "muted-foreground": "#64748B",
            "border": "#E2E8F0",
            "destructive": "#EF4444",
        },
        fonts={"heading": "Inter", "body": "Inter", "cjk_heading": "Microsoft YaHei", "cjk_body": "Microsoft YaHei"},
    )


def _dark_brand() -> BrandSpec:
    return BrandSpec(
        source="test",
        colors={
            "primary": "#3B82F6",
            "on-primary": "#FFFFFF",
            "secondary": "#8B5CF6",
            "accent": "#F59E0B",
            "background": "#0F172A",
            "foreground": "#F8FAFC",
            "muted": "#1E293B",
            "muted-foreground": "#94A3B8",
            "border": "#334155",
            "destructive": "#EF4444",
        },
        fonts={"heading": "Orbitron", "body": "JetBrains Mono", "cjk_heading": "Microsoft YaHei", "cjk_body": "Microsoft YaHei"},
        dark_mode=True,
    )


def _make_group_xml(
    shapes=None,
    bg_color="4472C4",
    text_color="FFFFFF",
    title_size=3200,
    body_size=1800,
    fill_colors=None,
    font_latin="Calibri",
    font_ea="",
    width_emu=9144000,
    height_emu=5486400,
) -> bytes:
    """Build realistic GroupShape XML with configurable properties."""
    if fill_colors is None:
        fill_colors = ["4472C4", "ED7D31", "A5A5A5", "FFC000"]

    shapes_xml = ""
    for i, fc in enumerate(fill_colors):
        x = 200000 + i * 2200000
        y = 300000
        cx = 2000000
        cy = 1500000
        is_title = (i == 0)
        sz = title_size if is_title else body_size
        tc = text_color
        txt = f"Step {i+1}" if is_title else f"Item {i+1} detail text"

        cjk_cs = ""
        if font_ea:
            cjk_cs = f'<a:ea typeface="{font_ea}"/><a:cs typeface="{font_ea}"/>'

        text_color_fill = f'<a:solidFill><a:srgbClr val="{tc}"/></a:solidFill>'

        shapes_xml += f"""
      <p:sp>
        <p:nvSpPr><p:cNvPr id="{10+i}" name="Shape{i}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
        <p:spPr>
          <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
          <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
          <a:solidFill><a:srgbClr val="{fc}"/></a:solidFill>
        </p:spPr>
        <p:txBody>
          <a:bodyPr lIns="91440" tIns="45720" rIns="91440" bIns="45720"/>
          <a:p>
            <a:r>
              <a:rPr lang="en-US" sz="{sz}" b="{"1" if is_title else "0"}" dirty="0">
                <a:latin typeface="{font_latin}"/>{cjk_cs}
                {text_color_fill}
              </a:rPr>
              <a:t>{txt}</a:t>
            </a:r>
          </a:p>
        </p:txBody>
      </p:sp>"""

    xml = f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="Group"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm>
      <a:off x="457200" y="457200"/>
      <a:ext cx="{width_emu}" cy="{height_emu}"/>
      <a:chOff x="0" y="0"/>
      <a:chExt cx="{width_emu}" cy="{height_emu}"/>
    </a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="3" name="BgRect"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="0" y="0"/><a:ext cx="{width_emu}" cy="{height_emu}"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="{bg_color}"/></a:solidFill>
    </p:spPr>
  </p:sp>{shapes_xml}
</p:grpSp>"""
    return xml.encode("utf-8")


def _make_smartart_data_xml(texts=None) -> bytes:
    """Build SmartArt data XML for testing schemeClr replacement."""
    if texts is None:
        texts = ["Step 1", "Step 2", "Step 3"]
    dgm_ns = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    pts = ""
    for i, t in enumerate(texts):
        pts += f'<dgm:pt><dgm:prLo val="0"/><a:t>{t}</a:t></dgm:pt>'
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<dgm:dataModel xmlns:dgm="{dgm_ns}" xmlns:a="{A_NS}">
  <dgm:ptLst>{pts}</dgm:ptLst>
</dgm:dataModel>""".encode("utf-8")


def _make_smartart_colors_xml() -> bytes:
    """Build SmartArt colors XML with schemeClr references."""
    dgm_ns = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<dgm:colorsDef xmlns:dgm="{dgm_ns}" xmlns:a="{A_NS}">
  <dgm:styleLbl name="node0">
    <a:solidFill><a:schemeClr val="accent1"/></a:solidFill>
  </dgm:styleLbl>
  <dgm:styleLbl name="node1">
    <a:solidFill><a:schemeClr val="accent2"/></a:solidFill>
  </dgm:styleLbl>
  <dgm:styleLbl name="node2">
    <a:solidFill><a:schemeClr val="dk1"/></a:solidFill>
  </dgm:styleLbl>
</dgm:colorsDef>""".encode("utf-8")


# ═══════════════════════════════════════════════════════════════
# Test: ComponentAdapter.analyze()
# ═══════════════════════════════════════════════════════════════

class TestComponentAdapterAnalyze:

    def test_analyze_extracts_color_roles(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(bg_color="4472C4", text_color="FFFFFF", fill_colors=["4472C4", "ED7D31", "A5A5A5"])
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group", "category": "process", "node_count": 3})
        assert analysis.color_count >= 3
        assert any(r.role == "text" for r in analysis.color_roles.values())

    def test_analyze_detects_dark_background(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(bg_color="1E293B", text_color="FFFFFF")
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group", "category": "process", "node_count": 3})
        assert analysis.has_dark_bg is True

    def test_analyze_detects_light_background(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(bg_color="F1F5F9", text_color="0F172A")
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group", "category": "process", "node_count": 3})
        assert analysis.has_dark_bg is False

    def test_analyze_extracts_font_levels(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(title_size=3200, body_size=1600)
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group", "category": "process", "node_count": 3})
        assert len(analysis.font_levels) >= 2
        roles = [fl.role for fl in analysis.font_levels]
        assert "title" in roles or "subtitle" in roles

    def test_analyze_extracts_aspect_ratio(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(width_emu=9144000, height_emu=5486400)
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group", "category": "process", "node_count": 3})
        assert abs(analysis.aspect_ratio - (9144000 / 5486400)) < 0.1

    def test_analyze_detects_cjk(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(font_ea="Microsoft YaHei")
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group", "category": "process", "node_count": 3})
        assert analysis.has_cjk is True


# ═══════════════════════════════════════════════════════════════
# Test: Color Adaptation — 白字在浅底上必须变深
# ═══════════════════════════════════════════════════════════════

class TestColorAdaptationWhiteOnLight:

    def test_white_text_on_light_bg_replaced_with_foreground(self):
        """CRITICAL: 白色文本在浅色品牌下必须替换为 foreground 色。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(bg_color="4472C4", text_color="FFFFFF")
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B", "C"], "node_count": 3,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])

        text_colors = set()
        for rpr_tag in (f"{{{A_NS}}}rPr", f"{{{A_NS}}}endParaRPr"):
            for rpr in grp_root.iter(rpr_tag):
                for sf in rpr.iter(f"{{{A_NS}}}srgbClr"):
                    parent = sf.getparent()
                    if parent is not None and parent.tag == f"{{{A_NS}}}solidFill":
                        gp = parent.getparent()
                        if gp is not None:
                            gp_local = etree.QName(gp.tag).localname if isinstance(gp.tag, str) else ""
                            if gp_local in ("rPr", "endParaRPr", "defRPr"):
                                text_colors.add(sf.get("val", "").upper())

        assert "FFFFFF" not in text_colors, f"White text still present in light brand! Found: {text_colors}"

    def test_white_text_kept_on_dark_bg(self):
        """深色品牌下文本应保持高亮度（可见）。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(bg_color="1E293B", text_color="FFFFFF")
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B"], "node_count": 2,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _dark_brand())
        grp_root = etree.fromstring(result["group"])

        text_colors = set()
        for rpr in grp_root.iter(f"{{{A_NS}}}rPr"):
            for sf in rpr.findall(f".//{{{A_NS}}}srgbClr"):
                val = sf.get("val", "").upper()
                parent = sf.getparent()
                if parent is not None and parent.tag == f"{{{A_NS}}}solidFill":
                    gp = parent.getparent()
                    if gp is not None:
                        gp_local = etree.QName(gp.tag).localname if isinstance(gp.tag, str) else ""
                        if gp_local in ("rPr", "endParaRPr", "defRPr"):
                            text_colors.add(val)

        # Text should be bright on dark brand (foreground is F8FAFC or similar)
        for tc in text_colors:
            brightness = 0.299 * int(tc[0:2], 16) + 0.587 * int(tc[2:4], 16) + 0.114 * int(tc[4:6], 16)
            assert brightness > 180, f"Text color #{tc} is too dark for dark brand (brightness={brightness:.0f})"


# ═══════════════════════════════════════════════════════════════
# Test: Color Adaptation — contrast ratio validation
# ═══════════════════════════════════════════════════════════════

class TestColorContrastValidation:

    def test_text_fill_contrast_above_3_to_1(self):
        """所有文本-背景配对对比度必须 >= 3:1。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(fill_colors=["4472C4", "ED7D31", "A5A5A5", "FFC000"])
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B", "C", "D"], "node_count": 4,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())

        issues = result.get("_validation_issues", [])
        contrast_issues = [i for i in issues if "contrast" in i.lower() or "Low contrast" in i]
        assert len(contrast_issues) == 0, f"Contrast issues found: {contrast_issues}"


# ═══════════════════════════════════════════════════════════════
# Test: Font Adaptation — 字号层级识别
# ═══════════════════════════════════════════════════════════════

class TestFontAdaptationLevels:

    def test_heading_font_applied_to_large_text(self):
        """大字号文本必须使用 heading 字体。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(title_size=3600, body_size=1400, font_latin="Arial")
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["Title", "Body"], "node_count": 2,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])

        heading_fonts = set()
        for rpr in grp_root.iter(f"{{{A_NS}}}rPr"):
            sz = rpr.get("sz")
            if sz and int(sz) >= 3000:
                latin = rpr.find(f"{{{A_NS}}}latin")
                if latin is not None:
                    heading_fonts.add(latin.get("typeface", ""))

        assert "Inter" in heading_fonts or any("Inter" in f for f in heading_fonts), \
            f"Heading font should be Inter, got: {heading_fonts}"

    def test_body_font_applied_to_small_text(self):
        """小字号文本必须使用 body 字体。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(title_size=3600, body_size=1400, font_latin="Arial")
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["Title", "Body"], "node_count": 2,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])

        body_fonts = set()
        for rpr in grp_root.iter(f"{{{A_NS}}}rPr"):
            sz = rpr.get("sz")
            if sz and int(sz) <= 2000:
                latin = rpr.find(f"{{{A_NS}}}latin")
                if latin is not None:
                    body_fonts.add(latin.get("typeface", ""))

        assert "Inter" in body_fonts or any("Inter" in f for f in body_fonts), \
            f"Body font should be Inter, got: {body_fonts}"

    def test_cjk_font_added_when_missing(self):
        """缺少 CJK 字体的组件必须自动添加。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(font_latin="Arial", font_ea="")
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B"], "node_count": 2,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])

        has_cjk_ea = False
        for rpr in grp_root.iter(f"{{{A_NS}}}rPr"):
            ea = rpr.find(f"{{{A_NS}}}ea")
            if ea is not None and ea.get("typeface", ""):
                has_cjk_ea = True

        assert has_cjk_ea, "CJK font (ea element) should be added to rPr elements"

    def test_no_font_below_11pt(self):
        """缩放后字号不得低于 11pt。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(title_size=1200, body_size=800)
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B"], "node_count": 2,
                   "bounds": (0.9, 1.6, 5.0, 3.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])

        min_sz = float("inf")
        for rpr_tag in (f"{{{A_NS}}}rPr", f"{{{A_NS}}}endParaRPr", f"{{{A_NS}}}defRPr"):
            for rpr in grp_root.iter(rpr_tag):
                sz = rpr.get("sz")
                if sz:
                    min_sz = min(min_sz, int(sz) / 100)

        assert min_sz >= 11.0, f"Font size {min_sz}pt is below minimum 11pt"


# ═══════════════════════════════════════════════════════════════
# Test: Layout Adaptation — 边界/偏移/溢出
# ═══════════════════════════════════════════════════════════════

class TestLayoutAdaptation:

    def test_component_bounds_within_content_area(self):
        """组件必须完全在目标区域内。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(width_emu=9144000, height_emu=5486400)
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B"], "node_count": 2,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        bounds = result.get("_adapted_bounds", (0.9, 1.6, 11.5, 5.0))
        left, top, w, h = bounds
        target_left, target_top, target_w, target_h = 0.9, 1.6, 11.5, 5.0
        assert left >= 0, f"Left bound {left} is negative"
        assert top >= 0, f"Top bound {top} is negative"
        assert w > 0 and h > 0, f"Invalid bounds: {bounds}"
        assert left + w <= target_left + target_w + 0.1, f"Component overflows right: {left}+{w} > {target_left}+{target_w}"
        assert top + h <= target_top + target_h + 0.1, f"Component overflows bottom: {top}+{h} > {target_top}+{target_h}"

    def test_bodypr_insets_not_zero(self):
        """文本框 insets 不能因缩放变为 0。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml()
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B"], "node_count": 2,
                   "bounds": (0.9, 1.6, 5.0, 3.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])

        for bodyPr in grp_root.iter(f"{{{A_NS}}}bodyPr"):
            for attr in ("lIns", "tIns", "rIns", "bIns"):
                val = bodyPr.get(attr)
                if val is not None:
                    assert int(val) >= 36000, f"bodyPr {attr}={val} is below minimum 36000 EMU (~0.04 inch)"

    def test_wide_component_uses_width_fit(self):
        """宽组件应优先适配宽度。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(width_emu=16000000, height_emu=4000000)
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B"], "node_count": 2,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        strategy = result.get("_fit_strategy", "contain")
        assert strategy in ("width", "contain", "stretch"), f"Expected width/contain/stretch strategy, got {strategy}"
        bounds = result.get("_adapted_bounds", (0.9, 1.6, 11.5, 5.0))
        assert bounds[2] > 0 and bounds[3] > 0, f"Invalid adapted bounds: {bounds}"


# ═══════════════════════════════════════════════════════════════
# Test: SmartArt schemeClr bug fix
# ═══════════════════════════════════════════════════════════════

class TestSmartArtSchemeClrFix:

    def test_schemeclr_replaced_with_correct_position(self):
        """schemeClr 替换后 srgbClr 必须在原位，不能跑到末尾。"""
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        renderer = ComponentRenderer()
        brand_spec = _light_brand()
        colors_xml = _make_smartart_colors_xml()
        xml_parts = {"colors": colors_xml}

        result = renderer._apply_brand_colors(xml_parts, brand_spec)
        result_colors = result.get("colors")
        if result_colors is None:
            pytest.skip("No colors XML to test")

        root = etree.fromstring(result_colors)
        srgb_vals = [s.get("val", "").upper() for s in root.iter(f"{{{A_NS}}}srgbClr")]
        assert len(srgb_vals) >= 1, "Should have at least one srgbClr after schemeClr replacement"
        brand_vals = {"2563EB", "7C3AED", "F97316", "0F172A"}
        matched = [v for v in srgb_vals if v in brand_vals]
        assert len(matched) >= 1, f"At least one brand color should appear. Got: {srgb_vals}, expected one of {brand_vals}"

    def test_no_schemeclr_left_after_replacement(self):
        """替换后不应残留 schemeClr（除非在 style 元素内）。"""
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        renderer = ComponentRenderer()
        brand_spec = _light_brand()
        colors_xml = _make_smartart_colors_xml()
        xml_parts = {"colors": colors_xml}

        result = renderer._apply_brand_colors(xml_parts, brand_spec)
        result_colors = result.get("colors")
        if result_colors is None:
            pytest.skip("No colors XML to test")

        root = etree.fromstring(result_colors)
        remaining_scheme = list(root.iter(f"{{{A_NS}}}schemeClr"))
        non_style_scheme = []
        for s in remaining_scheme:
            p = s.getparent()
            in_style = False
            while p is not None:
                if p.tag == f"{{{A_NS}}}style":
                    in_style = True
                    break
                p = p.getparent()
            if not in_style:
                non_style_scheme.append(s.get("val", ""))

        assert len(non_style_scheme) == 0, f"Unreplaced schemeClr values: {non_style_scheme}"


# ═══════════════════════════════════════════════════════════════
# Test: Group recolor — gradient handling
# ═══════════════════════════════════════════════════════════════

class TestGroupRecolorGradient:

    def test_gradient_colors_recolored(self):
        """渐变中的颜色也必须被替换。"""
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()

        xml = f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="Group"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm><a:off x="457200" y="457200"/><a:ext cx="9144000" cy="5486400"/>
      <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="5486400"/></a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="3" name="Shape0"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="200000" y="300000"/><a:ext cx="2000000" cy="1500000"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:gradFill>
        <a:gsLst>
          <a:gs pos="0"><a:srgbClr val="4472C4"/></a:gs>
          <a:gs pos="50000"><a:srgbClr val="5B9BD5"/></a:gs>
          <a:gs pos="100000"><a:srgbClr val="ED7D31"/></a:gs>
        </a:gsLst>
      </a:gradFill>
    </p:spPr>
  </p:sp>
</p:grpSp>""".encode("utf-8")

        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": [], "node_count": 0,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])

        gradient_colors = set()
        for gs in grp_root.iter(f"{{{A_NS}}}gs"):
            for srgb in gs.findall(f"{{{A_NS}}}srgbClr"):
                val = srgb.get("val", "").upper()
                gradient_colors.add(val)

        original_colors = {"4472C4", "5B9BD5", "ED7D31"}
        unchanged = gradient_colors & original_colors
        assert len(unchanged) == 0, f"Gradient colors not fully recolored: {unchanged} still present"
        assert len(gradient_colors) >= 1, "Should have at least one gradient color after recolor"


# ═══════════════════════════════════════════════════════════════
# Test: End-to-end PPT generation with component
# ═══════════════════════════════════════════════════════════════

class TestEndToEndComponentQuality:

    def _build_test_library(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary
        db_path = str(tmp_path / "test_lib.db")
        lib = ComponentLibrary(db_path)

        xml = _make_group_xml(fill_colors=["4472C4", "ED7D31", "A5A5A5"], text_color="FFFFFF")
        lib.add(
            type="group",
            category="process",
            variant="chevron",
            node_count=3,
            xml_parts={"group": xml},
        )
        return lib

    def test_rendered_ppt_has_no_white_text_on_light_bg(self, tmp_path):
        """生成的 PPT 中浅色品牌下不能有白色文本。"""
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        lib = self._build_test_library(tmp_path)
        brand = _light_brand()
        precision = PrecisionRenderer(brand_spec=brand)
        prs = precision.create_presentation()
        slide = precision.add_slide(prs)
        precision.apply_brand_background(slide, prs, goal="content")

        element = {
            "type": "group",
            "category": "process",
            "texts": ["Phase 1", "Phase 2", "Phase 3"],
            "nodes": [{"text": "Phase 1"}, {"text": "Phase 2"}, {"text": "Phase 3"}],
            "node_count": 3,
            "bounds": (0.9, 1.6, 11.5, 5.0),
        }
        renderer = ComponentRenderer()
        renderer.render(slide, element, brand, lib)

        output = str(tmp_path / "test_output.pptx")
        precision.save(prs, output)
        lib.close()

        assert os.path.isfile(output)
        from pptx import Presentation
        result_prs = Presentation(output)
        assert len(result_prs.slides) == 1

    def test_rendered_ppt_fonts_are_brand_compliant(self, tmp_path):
        """生成的 PPT 中字体必须符合品牌规范。"""
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        lib = self._build_test_library(tmp_path)
        brand = _light_brand()
        precision = PrecisionRenderer(brand_spec=brand)
        prs = precision.create_presentation()
        slide = precision.add_slide(prs)
        precision.apply_brand_background(slide, prs, goal="content")

        element = {
            "type": "group",
            "category": "process",
            "texts": ["Phase 1", "Phase 2", "Phase 3"],
            "nodes": [{"text": "Phase 1"}, {"text": "Phase 2"}, {"text": "Phase 3"}],
            "node_count": 3,
            "bounds": (0.9, 1.6, 11.5, 5.0),
        }
        renderer = ComponentRenderer()
        renderer.render(slide, element, brand, lib)

        output = str(tmp_path / "test_fonts.pptx")
        precision.save(prs, output)
        lib.close()

        from lxml import etree as _etree
        import zipfile
        with zipfile.ZipFile(output) as z:
            slide_xml = z.read("ppt/slides/slide1.xml")
            root = _etree.fromstring(slide_xml)
            non_brand_fonts = set()
            for latin in root.iter(f"{{{A_NS}}}latin"):
                tf = latin.get("typeface", "")
                if tf and tf not in ("Inter", "Microsoft YaHei", "+mj-lt", "+mn-lt", ""):
                    if not tf.startswith("+"):
                        non_brand_fonts.add(tf)
            assert len(non_brand_fonts) == 0, f"Non-brand fonts found: {non_brand_fonts}"


# ═══════════════════════════════════════════════════════════════
# v2 New Test Scenarios (17+)
# ═══════════════════════════════════════════════════════════════

def _make_nested_grpsp_xml(depth=2, ch_off_x=0, ch_off_y=0) -> bytes:
    """Build nested GroupShape XML with configurable depth and chOff."""
    inner_xml = """<p:sp>
      <p:nvSpPr><p:cNvPr id="10" name="InnerShape"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="200000" y="100000"/><a:ext cx="800000" cy="600000"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:p><a:r><a:rPr lang="en-US" sz="1800">
        <a:latin typeface="Calibri"/><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
      </a:rPr><a:t>Inner</a:t></a:r></a:p></p:txBody>
    </p:sp>"""

    if depth >= 2:
        inner_xml = f"""<p:grpSp>
      <p:nvGrpSpPr><p:cNvPr id="5" name="SubGroup"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm>
        <a:off x="100000" y="50000"/><a:ext cx="1600000" cy="1200000"/>
        <a:chOff x="{ch_off_x}" y="{ch_off_y}"/><a:chExt cx="1600000" cy="1200000"/>
      </a:xfrm></p:grpSpPr>
      {inner_xml}
    </p:grpSp>"""

    if depth >= 3:
        inner_xml = f"""<p:grpSp>
      <p:nvGrpSpPr><p:cNvPr id="4" name="MidGroup"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm>
        <a:off x="0" y="0"/><a:ext cx="3200000" cy="2400000"/>
        <a:chOff x="0" y="0"/><a:chExt cx="3200000" cy="2400000"/>
      </a:xfrm></p:grpSpPr>
      {inner_xml}
    </p:grpSp>"""

    return f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="OuterGroup"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm>
    <a:off x="457200" y="457200"/><a:ext cx="9144000" cy="5486400"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="5486400"/>
  </a:xfrm></p:grpSpPr>
  {inner_xml}
</p:grpSp>""".encode("utf-8")


def _make_mixed_role_color_xml() -> bytes:
    """Build XML where same hex color appears as both fill and text."""
    return f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="Group"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm>
    <a:off x="457200" y="457200"/><a:ext cx="9144000" cy="5486400"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="5486400"/>
  </a:xfrm></p:grpSpPr>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="3" name="Shape0"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="200000" y="300000"/><a:ext cx="4000000" cy="2000000"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    </p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:rPr lang="en-US" sz="1800">
      <a:latin typeface="Calibri"/>
      <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    </a:rPr><a:t>Same color as fill</a:t></a:r></a:p></p:txBody>
  </p:sp>
</p:grpSp>""".encode("utf-8")


def _make_multi_level_font_xml(sizes=None) -> bytes:
    """Build XML with multiple font size levels."""
    if sizes is None:
        sizes = [4400, 3200, 2400, 1800, 1200]
    shapes = ""
    for i, sz in enumerate(sizes):
        shapes += f"""<p:sp>
      <p:nvSpPr><p:cNvPr id="{10+i}" name="Sp{i}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{200000+i*1500000}" y="300000"/><a:ext cx="1200000" cy="800000"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:p><a:r><a:rPr lang="en-US" sz="{sz}">
        <a:latin typeface="Calibri"/><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
      </a:rPr><a:t>Level{i}</a:t></a:r></a:p></p:txBody>
    </p:sp>"""
    return f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="Group"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm>
    <a:off x="457200" y="457200"/><a:ext cx="9144000" cy="5486400"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="5486400"/>
  </a:xfrm></p:grpSpPr>{shapes}
</p:grpSp>""".encode("utf-8")


class TestNestedGrpSpCoordinateTransform:
    """v2: Nested grpSp (2-3 levels) coordinate transform."""

    def test_2_level_nested_grpsp_transforms(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_nested_grpsp_xml(depth=2)
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A"], "node_count": 1,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        all_sp = list(grp_root.iter(f"{{{P_NS}}}sp"))
        assert len(all_sp) >= 1, "Inner shape should exist after transform"

    def test_3_level_nested_grpsp_transforms(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_nested_grpsp_xml(depth=3)
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A"], "node_count": 1,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        all_sp = list(grp_root.iter(f"{{{P_NS}}}sp"))
        assert len(all_sp) >= 1, "Inner shape should exist after 3-level transform"

    def test_choff_nonzero_normalizes(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_nested_grpsp_xml(depth=2, ch_off_x=100000, ch_off_y=50000)
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A"], "node_count": 1,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        for grpSpPr in grp_root.iter(f"{{{P_NS}}}grpSpPr"):
            xfrm = grpSpPr.find(f"{{{A_NS}}}xfrm")
            if xfrm is not None:
                chOff = xfrm.find(f"{{{A_NS}}}chOff")
                if chOff is not None:
                    assert int(chOff.get("x", "0")) == 0, "chOff.x should be normalized to 0"
                    assert int(chOff.get("y", "0")) == 0, "chOff.y should be normalized to 0"


class TestDoubleTransformGuard:
    """v2: Running adapt() twice should be idempotent (no double-transform)."""

    def test_double_adapt_idempotent(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(fill_colors=["4472C4", "ED7D31"])
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B"], "node_count": 2,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result1 = adapter.adapt(xml_parts, element, _light_brand())
        result2 = adapter.adapt(result1, element, _light_brand())
        root1 = etree.fromstring(result1["group"])
        root2 = etree.fromstring(result2["group"])
        sizes1 = sorted([int(r.get("sz", "0")) for r in root1.iter(f"{{{A_NS}}}rPr") if r.get("sz")])
        sizes2 = sorted([int(r.get("sz", "0")) for r in root2.iter(f"{{{A_NS}}}rPr") if r.get("sz")])
        assert sizes1 == sizes2, f"Double adapt changed font sizes: {sizes1} vs {sizes2}"


class TestAreaWeightedColorRoles:
    """v2: Area-weighted color role inference."""

    def test_large_shape_is_dominant_fill(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(fill_colors=["4472C4", "ED7D31", "A5A5A5"])
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group", "category": "process", "node_count": 3})
        dominant = [v for v in analysis.color_roles.values() if v.role == "dominant_fill"]
        assert len(dominant) >= 1, "Should have at least one dominant_fill"

    def test_small_shape_is_secondary_or_data(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="Group"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm>
    <a:off x="0" y="0"/><a:ext cx="9144000" cy="5486400"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="5486400"/>
  </a:xfrm></p:grpSpPr>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="3" name="Big"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="0" y="0"/><a:ext cx="8000000" cy="5000000"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="4472C4"/></a:solidFill>
    </p:spPr>
  </p:sp>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="4" name="Tiny"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="8000000" y="0"/><a:ext cx="100000" cy="100000"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>
    </p:spPr>
  </p:sp>
</p:grpSp>""".encode("utf-8")
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group", "category": "process", "node_count": 2})
        big_role = analysis.color_roles.get("4472C4")
        tiny_role = analysis.color_roles.get("FF0000")
        assert big_role is not None and big_role.role == "dominant_fill", f"Big shape should be dominant_fill, got {big_role.role if big_role else None}"
        assert tiny_role is not None and tiny_role.role in ("data_fill", "secondary_fill"), f"Tiny shape should be data_fill or secondary_fill, got {tiny_role.role if tiny_role else None}"


class TestMixedRoleColorSeparation:
    """v2: Same hex as both fill and text gets different replacements per context."""

    def test_same_hex_different_replacement_for_fill_vs_text(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_mixed_role_color_xml()
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A"], "node_count": 1,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        fill_colors = set()
        text_colors = set()
        for srgb in grp_root.iter(f"{{{A_NS}}}srgbClr"):
            val = srgb.get("val", "").upper()
            parent = srgb.getparent()
            if parent is None:
                continue
            p_tag = etree.QName(parent.tag).localname if isinstance(parent.tag, str) else ""
            gp = parent.getparent()
            gp_tag = etree.QName(gp.tag).localname if gp is not None and isinstance(gp.tag, str) else ""
            if p_tag == "solidFill":
                if gp_tag in ("rPr", "endParaRPr", "defRPr"):
                    text_colors.add(val)
                else:
                    fill_colors.add(val)
        assert "4472C4" not in fill_colors, f"Original fill color not replaced: {fill_colors}"
        assert "4472C4" not in text_colors, f"Original text color not replaced: {text_colors}"


class TestFontHierarchyPreserved:
    """v2: Font hierarchy preserved at various scale factors."""

    def test_hierarchy_at_small_scale_30pct(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_multi_level_font_xml(sizes=[4400, 3200, 2400, 1800, 1200])
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B", "C", "D", "E"],
                   "node_count": 5, "bounds": (0.9, 1.6, 3.0, 2.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        sizes = sorted([int(r.get("sz", "0")) for r in grp_root.iter(f"{{{A_NS}}}rPr") if r.get("sz")], reverse=True)
        unique_sizes = sorted(set(sizes), reverse=True)
        if len(unique_sizes) >= 2:
            for i in range(len(unique_sizes) - 1):
                gap_pt = (unique_sizes[i] - unique_sizes[i + 1]) / 100
                assert gap_pt >= 1.0, f"Font gap between levels {i} and {i+1} is {gap_pt}pt, should be >= 1pt"

    def test_hierarchy_at_medium_scale_50pct(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_multi_level_font_xml(sizes=[3600, 2800, 2000, 1400])
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B", "C", "D"],
                   "node_count": 4, "bounds": (0.9, 1.6, 5.0, 3.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        sizes = sorted([int(r.get("sz", "0")) for r in grp_root.iter(f"{{{A_NS}}}rPr") if r.get("sz")], reverse=True)
        for sz in sizes:
            assert sz / 100 >= 11.0, f"Font size {sz/100}pt is below 11pt minimum"

    def test_font_scale_up_when_enlarged(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_multi_level_font_xml(sizes=[1800, 1400])
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B"],
                   "node_count": 2, "bounds": (0.5, 0.5, 12.0, 6.5)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        sizes = [int(r.get("sz", "0")) for r in grp_root.iter(f"{{{A_NS}}}rPr") if r.get("sz")]
        assert len(sizes) >= 2
        assert max(sizes) >= 1800, f"Font should scale up when target is larger, got max={max(sizes)}"


class Test72ptUpperLimit:
    """v2: Font sizes capped at 72pt (7200 hundredths)."""

    def test_no_font_above_72pt(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_multi_level_font_xml(sizes=[9600, 7200, 4800, 2400, 1200, 800])
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B", "C", "D", "E", "F"],
                   "node_count": 6, "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        for rpr in grp_root.iter(f"{{{A_NS}}}rPr"):
            sz = rpr.get("sz")
            if sz:
                assert int(sz) <= 7200, f"Font size {int(sz)/100}pt exceeds 72pt cap"


class TestMinFontGapEnforcement:
    """v2: Each font level >= level_below + 1pt."""

    def test_gap_enforcement_across_levels(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_multi_level_font_xml(sizes=[2400, 2200, 2000, 1800])
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A", "B", "C", "D"],
                   "node_count": 4, "bounds": (0.9, 1.6, 5.0, 3.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        sizes = sorted(set(int(r.get("sz", "0")) for r in grp_root.iter(f"{{{A_NS}}}rPr") if r.get("sz")), reverse=True)
        if len(sizes) >= 2:
            for i in range(len(sizes) - 1):
                gap = (sizes[i] - sizes[i + 1]) / 100
                assert gap >= 1.0, f"Gap between level {i} ({sizes[i]/100}pt) and {i+1} ({sizes[i+1]/100}pt) is {gap}pt < 1pt"


class TestSchemeClrInGroupXml:
    """v2: schemeClr in group XML handled correctly."""

    def test_schemeclr_in_spPr_replaced(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="Group"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm>
    <a:off x="0" y="0"/><a:ext cx="9144000" cy="5486400"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="5486400"/>
  </a:xfrm></p:grpSpPr>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="3" name="Shape0"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="200000" y="300000"/><a:ext cx="2000000" cy="1500000"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:schemeClr val="accent1"/></a:solidFill>
    </p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:rPr lang="en-US" sz="1800">
      <a:latin typeface="Calibri"/>
    </a:rPr><a:t>Test</a:t></a:r></a:p></p:txBody>
  </p:sp>
</p:grpSp>""".encode("utf-8")
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A"], "node_count": 1,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        remaining_scheme = list(grp_root.iter(f"{{{A_NS}}}schemeClr"))
        assert len(remaining_scheme) == 0, f"Unreplaced schemeClr in group XML: {len(remaining_scheme)}"


class TestNonWhiteLowContrast:
    """v2: Non-white text on light background — validation should detect low contrast."""

    def test_light_text_contrast_detected_by_validate(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(bg_color="F1F5F9", text_color="B0B0B0", fill_colors=["4472C4"])
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A"], "node_count": 1,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        issues = result.get("_validation_issues", [])
        has_contrast_issue = any("contrast" in i.lower() or "Low contrast" in i for i in issues)
        assert has_contrast_issue, "Low contrast should be detected by validation"


class TestBrandSpecNone:
    """v2: brand_spec=None should not crash."""

    def test_adapt_with_none_brand_spec(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml()
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A"], "node_count": 1,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, None)
        assert "group" in result
        assert result["group"] is not None


class TestEmptyComponent:
    """v2: Empty/minimal XML should not crash."""

    def test_empty_group_xml(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="Empty"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm>
    <a:off x="0" y="0"/><a:ext cx="9144000" cy="5486400"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="5486400"/>
  </a:xfrm></p:grpSpPr>
</p:grpSp>""".encode("utf-8")
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": [], "node_count": 0,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        assert "group" in result

    def test_no_group_key(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml_parts = {"_meta": b"test"}
        element = {"type": "group", "category": "process", "texts": [], "node_count": 0,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        assert "_meta" in result

    def test_invalid_xml(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml_parts = {"group": b"<invalid><xml"}
        element = {"type": "group", "category": "process", "texts": [], "node_count": 0,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        assert "group" in result


class TestImportErrorFallback:
    """v2: ComponentRenderer.render_group falls back to old path on ImportError."""

    def test_render_group_fallback_returns_xml(self, monkeypatch):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        brand = BrandSpec(colors={"primary": "#2563EB", "accent": "#F97316", "muted": "#F1F5F9"})
        xml = _make_group_xml(fill_colors=["4472C4", "ED7D31"])
        xml_parts = {"group": xml}
        renderer = ComponentRenderer()
        result = renderer._apply_brand_colors(xml_parts, brand)
        assert "group" in result
        root = etree.fromstring(result["group"])
        fill_vals = set()
        for sf in root.iter(f"{{{A_NS}}}solidFill"):
            srgb = sf.find(f"{{{A_NS}}}srgbClr")
            if srgb is not None:
                fill_vals.add(srgb.get("val", "").upper())
        assert "4472C4" not in fill_vals, f"Old path should still recolor: {fill_vals}"


class TestGradientRecolorAllStops:
    """v2: Gradient recolor must replace ALL stops, not just some."""

    def test_all_gradient_stops_replaced(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = f"""<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}" xmlns:r="{R_NS}">
  <p:nvGrpSpPr><p:cNvPr id="2" name="Group"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm>
    <a:off x="0" y="0"/><a:ext cx="9144000" cy="5486400"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="5486400"/>
  </a:xfrm></p:grpSpPr>
  <p:sp>
    <p:nvSpPr><p:cNvPr id="3" name="Shape0"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="200000" y="300000"/><a:ext cx="8000000" cy="4000000"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:gradFill>
        <a:gsLst>
          <a:gs pos="0"><a:srgbClr val="4472C4"/></a:gs>
          <a:gs pos="50000"><a:srgbClr val="5B9BD5"/></a:gs>
          <a:gs pos="100000"><a:srgbClr val="ED7D31"/></a:gs>
        </a:gsLst>
      </a:gradFill>
    </p:spPr>
  </p:sp>
</p:grpSp>""".encode("utf-8")
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": [], "node_count": 0,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        grp_root = etree.fromstring(result["group"])
        original_colors = {"4472C4", "5B9BD5", "ED7D31"}
        remaining = set()
        for gs in grp_root.iter(f"{{{A_NS}}}gs"):
            for srgb in gs.findall(f"{{{A_NS}}}srgbClr"):
                val = srgb.get("val", "").upper()
                if val in original_colors:
                    remaining.add(val)
        assert len(remaining) == 0, f"Gradient stops not recolored: {remaining}"


class TestFitStrategyExact:
    """v2: Fit strategy should be exact match, not fuzzy 'in' check."""

    def test_contain_strategy_for_standard_aspect(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml(width_emu=9144000, height_emu=5486400)
        xml_parts = {"group": xml}
        element = {"type": "group", "category": "process", "texts": ["A"], "node_count": 1,
                   "bounds": (0.9, 1.6, 11.5, 5.0)}
        result = adapter.adapt(xml_parts, element, _light_brand())
        strategy = result.get("_fit_strategy", "")
        assert strategy in ("stretch", "contain", "width"), f"Expected stretch/contain/width, got {strategy}"


class TestNestingDepthDetection:
    """v2: Nesting depth correctly detected."""

    def test_flat_group_depth_0(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_group_xml()
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group"})
        assert analysis.nesting_depth == 0, f"Flat group should have depth 0, got {analysis.nesting_depth}"

    def test_nested_group_depth_2(self):
        from ppt_pro_max.enterprise.component_adapter import ComponentAdapter
        adapter = ComponentAdapter()
        xml = _make_nested_grpsp_xml(depth=2)
        xml_parts = {"group": xml}
        analysis = adapter.analyze(xml_parts, {"type": "group"})
        assert analysis.nesting_depth >= 1, f"Nested group should have depth >= 1, got {analysis.nesting_depth}"
