"""Audit v3: GroupShape-internal font sizes + whitespace."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def audit_deep(path, label):
    prs = Presentation(path)
    all_sizes = []
    all_fonts = {}
    
    print(f"\n{'='*80}")
    print(f"  {label} | {len(prs.slides)} slides")
    print(f"{'='*80}")
    
    for i, slide in enumerate(prs.slides):
        slide_sizes = []
        slide_fonts = []
        
        shapes_to_check = []
        for shape in slide.shapes:
            shapes_to_check.append(shape)
            if shape.shape_type == 6:
                for child in shape.shapes:
                    shapes_to_check.append(child)
        
        for shape in shapes_to_check:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                pPr = para._p.find(f'{{{ns_a}}}pPr')
                def_sz = def_font = None
                if pPr is not None:
                    defRPr = pPr.find(f'{{{ns_a}}}defRPr')
                    if defRPr is not None:
                        def_sz = defRPr.get('sz')
                        latin = defRPr.find(f'{{{ns_a}}}latin')
                        if latin is not None:
                            def_font = latin.get('typeface')
                
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rPr = run._r.find(f'{{{ns_a}}}rPr')
                    sz = font = None
                    if rPr is not None:
                        sz = rPr.get('sz')
                        latin = rPr.find(f'{{{ns_a}}}latin')
                        if latin is not None:
                            font = latin.get('typeface')
                    
                    pt = int(sz)/100 if sz else (int(def_sz)/100 if def_sz else None)
                    fn = font or def_font
                    
                    if pt:
                        all_sizes.append(pt)
                        slide_sizes.append(pt)
                    if fn:
                        all_fonts[fn] = all_fonts.get(fn, 0) + 1
                        slide_fonts.append(fn)
        
        if slide_sizes:
            print(f"  Slide {i}: sizes={sorted(set(slide_sizes))} fonts={sorted(set(slide_fonts))}")
    
    sizes_sorted = sorted(set(all_sizes))
    print(f"\n  ALL FONT SIZES: {sizes_sorted} ({len(sizes_sorted)} levels)")
    print(f"  FONTS: {dict(sorted(all_fonts.items(), key=lambda x: -x[1]))}")
    
    if len(sizes_sorted) >= 2:
        print(f"  SCALE RATIOS:")
        for j in range(len(sizes_sorted) - 1):
            r = sizes_sorted[j+1] / sizes_sorted[j]
            print(f"    {sizes_sorted[j]:.0f} -> {sizes_sorted[j+1]:.0f} = {r:.2f}x")
    
    problems = []
    if min(all_sizes) < 8:
        problems.append(f"Font too small: {min(all_sizes)}pt")
    if len(sizes_sorted) < 5:
        problems.append(f"Too few size levels: {len(sizes_sorted)}")
    if len(sizes_sorted) > 10:
        problems.append(f"Too many size levels: {len(sizes_sorted)}")
    
    ratios = [sizes_sorted[j+1]/sizes_sorted[j] for j in range(len(sizes_sorted)-1)]
    if ratios:
        avg_r = sum(ratios)/len(ratios)
        max_dev = max(abs(r - avg_r) for r in ratios)
        if max_dev > 0.4:
            problems.append(f"Inconsistent ratios: avg={avg_r:.2f}, max_dev={max_dev:.2f}")
    
    if problems:
        print(f"  PROBLEMS:")
        for p in problems:
            print(f"    - {p}")
    else:
        print(f"  OK: Typography scale is healthy")


v3_dir = os.path.join(tempfile.gettempdir(), 'ppt_build_v3')
for name, label in [('build_mckinsey.pptx', 'McKinsey v3'), ('build_cyberpunk.pptx', 'Cyberpunk v3'), ('build_creative.pptx', 'Creative v3')]:
    path = os.path.join(v3_dir, name)
    if os.path.exists(path):
        audit_deep(path, label)
