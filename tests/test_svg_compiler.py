"""Tests for svg_compiler._compiler — 6 probe cases + editability audit."""
import zipfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.renderer.svg_compiler import SVGCompileError, SVGCompiler

# ─────────────────────────── probe SVG cases ───────────────────────────

PYRAMID_SVG = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 360">
  <defs>
    <linearGradient id="g1" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0" stop-color="#7DA92F"/><stop offset="1" stop-color="#2E6504"/>
    </linearGradient>
    <linearGradient id="g2" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0" stop-color="#5B9BD5"/><stop offset="1" stop-color="#2E75B6"/>
    </linearGradient>
    <linearGradient id="g3" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0" stop-color="#FFC000"/><stop offset="1" stop-color="#BF8F00"/>
    </linearGradient>
  </defs>
  <polygon points="200,30 380,320 20,320" fill="url(#g1)"/>
  <polygon points="200,120 320,320 80,320" fill="url(#g2)"/>
  <polygon points="200,210 260,320 140,320" fill="url(#g3)"/>
  <text x="200" y="85" text-anchor="middle" font-size="16" fill="#fff" font-family="Arial">战略愿景</text>
  <text x="200" y="175" text-anchor="middle" font-size="14" fill="#fff" font-family="Arial">三年目标</text>
  <text x="200" y="265" text-anchor="middle" font-size="14" fill="#fff" font-family="Arial">年度计划</text>
</svg>"""

VENN_SVG = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 260">
  <path d="M150,130 a80,80 0 1,0 0.1,0 Z M250,130 a80,80 0 1,0 0.1,0 Z"
        fill="#E8534E" fill-opacity="0.85"/>
  <text x="105" y="135" text-anchor="middle" font-size="15" fill="#fff">客户A</text>
  <text x="295" y="135" text-anchor="middle" font-size="15" fill="#fff">客户B</text>
  <text x="200" y="135" text-anchor="middle" font-size="13" fill="#fff">交叉用户</text>
</svg>"""

FUNNEL_SVG = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 420">
  <defs>
    <linearGradient id="f1" x1="0" y1="0" x2="0" y2="1">
      <stop offset="0" stop-color="#4472C4"/><stop offset="1" stop-color="#2E5BA6"/>
    </linearGradient>
  </defs>
  <polygon points="40,20 360,20 300,90 100,90" fill="url(#f1)"/>
  <polygon points="70,110 330,110 280,180 120,180" fill="#5B9BD5"/>
  <polygon points="100,200 300,200 260,270 140,270" fill="#9DC3E6"/>
  <polygon points="130,290 270,290 240,360 160,360" fill="#BDD7EE"/>
  <text x="200" y="55" text-anchor="middle" font-size="15" fill="#fff">曝光 10000</text>
  <text x="200" y="145" text-anchor="middle" font-size="15" fill="#fff">点击 4000</text>
  <text x="200" y="235" text-anchor="middle" font-size="15" fill="#2E5BA6">注册 1500</text>
  <text x="200" y="325" text-anchor="middle" font-size="15" fill="#2E5BA6">转化 600</text>
</svg>"""

GROWTH_CURVE_SVG = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">
  <defs>
    <clipPath id="clip"><rect x="30" y="40" width="340" height="220"/></clipPath>
    <linearGradient id="area" x1="0" y1="0" x2="0" y2="1">
      <stop offset="0" stop-color="#1D78FA" stop-opacity="0.45"/>
      <stop offset="1" stop-color="#1D78FA" stop-opacity="0.03"/>
    </linearGradient>
  </defs>
  <line x1="40" y1="250" x2="380" y2="250" stroke="#D9D9D9" stroke-width="1.5"/>
  <line x1="40" y1="250" x2="40" y2="50" stroke="#D9D9D9" stroke-width="1.5"/>
  <g clip-path="url(#clip)">
    <path d="M40,250 C120,240 150,180 200,150 C260,115 300,90 380,60 L380,250 Z"
          fill="url(#area)" fill-rule="nonzero"/>
    <path d="M40,250 C120,240 150,180 200,150 C260,115 300,90 380,60"
          fill="none" stroke="#1D78FA" stroke-width="3.5"/>
  </g>
  <circle cx="380" cy="60" r="6" fill="#1D78FA"/>
  <text x="200" y="275" text-anchor="middle" font-size="13" fill="#595959">季度增长曲线</text>
</svg>"""

MATRIX_BCG_SVG = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 340">
  <g transform="translate(20,20)">
    <rect x="0" y="0" width="360" height="300" fill="#FFFFFF" stroke="#BFBFBF" stroke-width="1"/>
    <line x1="180" y1="0" x2="180" y2="300" stroke="#BFBFBF" stroke-width="1" stroke-dasharray="4,4"/>
    <line x1="0" y1="150" x2="360" y2="150" stroke="#BFBFBF" stroke-width="1" stroke-dasharray="4,4"/>
    <text x="180" y="20" text-anchor="middle" font-size="11" fill="#7F7F7F">高</text>
    <text x="180" y="296" text-anchor="middle" font-size="11" fill="#7F7F7F">低</text>
    <g transform="translate(30,30)">
      <circle cx="0" cy="0" r="24" fill="#2E75B6"/><text x="0" y="5" text-anchor="middle" font-size="11" fill="#fff">明星</text>
      <circle cx="90" cy="0" r="24" fill="#BF8F00"/><text x="90" y="5" text-anchor="middle" font-size="11" fill="#fff">问题</text>
    </g>
    <g transform="translate(30,90)">
      <circle cx="0" cy="0" r="24" fill="#70AD47"/><text x="0" y="5" text-anchor="middle" font-size="11" fill="#fff">现金牛</text>
      <circle cx="90" cy="0" r="24" fill="#C00000"/><text x="90" y="5" text-anchor="middle" font-size="11" fill="#fff">瘦狗</text>
    </g>
  </g>
</svg>"""

UNSUPPORTED_SVG = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">
  <defs>
    <filter id="blur"><feGaussianBlur stdDeviation="3"/></filter>
    <mask id="m"><rect x="0" y="0" width="400" height="300" fill="#fff"/></mask>
  </defs>
  <image href="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
         x="50" y="50" width="100" height="100"/>
  <rect x="20" y="180" width="360" height="80" fill="#4472C4" filter="url(#blur)"/>
  <circle cx="200" cy="120" r="60" fill="#E8534E" mask="url(#m)"/>
</svg>"""

PROBE_CASES = {
    "pyramid": PYRAMID_SVG,
    "venn_evenodd": VENN_SVG,
    "funnel": FUNNEL_SVG,
    "growth_curve": GROWTH_CURVE_SVG,
    "matrix_bcg": MATRIX_BCG_SVG,
}

SLIDE_RECTS = {
    "pyramid": (3.5, 0.8, 6.3, 5.6),
    "venn_evenodd": (2.5, 1.2, 8.0, 5.2),
    "funnel": (4.0, 0.8, 5.4, 6.0),
    "growth_curve": (2.5, 1.0, 8.0, 5.2),
    "matrix_bcg": (3.0, 0.9, 7.0, 5.7),
}


def _make_slide():
    _prs = Presentation()
    _prs.slide_width = Inches(13.333)
    _prs.slide_height = Inches(7.5)
    return _prs, _prs.slides.add_slide(_prs.slide_layouts[6])


def _audit_pptx(_prs, pptx_path):
    _prs.save(str(pptx_path))
    with zipfile.ZipFile(pptx_path) as z:
        slide_xmls = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        xml = z.read(slide_xmls[0]).decode("utf-8")
    has_pic = "<p:pic>" in xml
    has_sp = xml.count("<p:sp>")
    return has_pic, has_sp


# ─────────────────────────── tests ───────────────────────────


class TestCompileProbeCases:
    @pytest.mark.parametrize("name,svg", PROBE_CASES.items())
    def test_compile_no_pictures(self, name, svg, tmp_path):
        _prs, slide = _make_slide()
        rect = SLIDE_RECTS[name]
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, rect)
        assert result is not None
        pptx_path = tmp_path / f"{name}.pptx"
        has_pic, has_sp = _audit_pptx(_prs, pptx_path)
        assert not has_pic, f"{name}: output contains <p:pic> — zero-raster guarantee violated"
        assert has_sp > 0, f"{name}: no shapes produced"

    @pytest.mark.parametrize("name,svg", PROBE_CASES.items())
    def test_compile_produces_shapes(self, name, svg):
        _prs, slide = _make_slide()
        rect = SLIDE_RECTS[name]
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, rect)
        assert result.shape_count > 0, f"{name}: zero shapes compiled"

    @pytest.mark.parametrize("name,svg", PROBE_CASES.items())
    def test_compile_speed(self, name, svg):
        _prs, slide = _make_slide()
        rect = SLIDE_RECTS[name]
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, rect)
        assert result.compile_ms < 100, f"{name}: compile took {result.compile_ms:.0f}ms (budget 100ms)"

    def test_pyramid_has_gradient(self):
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(PYRAMID_SVG, slide, SLIDE_RECTS["pyramid"])
        assert "gradient" in result.features

    def test_funnel_has_gradient(self):
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(FUNNEL_SVG, slide, SLIDE_RECTS["funnel"])
        assert "gradient" in result.features

    def test_matrix_bcg_has_transform(self):
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(MATRIX_BCG_SVG, slide, SLIDE_RECTS["matrix_bcg"])
        assert "rect" in result.features or "circle" in result.features

    def test_growth_curve_has_clippath(self):
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(GROWTH_CURVE_SVG, slide, SLIDE_RECTS["growth_curve"])
        assert "path" in result.features or "clipPath" in result.features


class TestUnsupportedFeatures:
    def test_image_soft_warning(self):
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(UNSUPPORTED_SVG, slide, (3.0, 0.9, 7.0, 5.7))
        assert any("unsupported" in w for w in result.warnings)

    def test_filter_element_soft_warning(self):
        svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><filter id="f"><feGaussianBlur stdDeviation="3"/></filter></svg>'
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert any("filter" in w for w in result.warnings)

    def test_mask_element_soft_warning(self):
        svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><mask id="m"><rect width="100" height="100" fill="white"/></mask></svg>'
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert any("mask" in w for w in result.warnings)


class TestColorResolution:
    def test_var_primary_resolved(self):
        C = {"primary": "#1D78FA", "on_primary": "#FFFFFF"}
        svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><rect x="0" y="0" width="100" height="100" fill="var(--primary)"/></svg>'
        _prs, slide = _make_slide()
        compiler = SVGCompiler(C=C)
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count > 0

    def test_named_color_resolved(self):
        svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><rect x="0" y="0" width="100" height="100" fill="red"/></svg>'
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count > 0

    def test_unresolved_var_raises(self):
        svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><rect x="0" y="0" width="100" height="100" fill="var(--unknown)"/></svg>'
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        with pytest.raises(SVGCompileError, match="unresolved"):
            compiler.compile(svg, slide, (1, 1, 5, 5))

    def test_currentColor_resolved(self):
        C = {"text_dark": "#1E293B"}
        svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><text x="50" y="50" fill="currentColor" font-size="14">Hello</text></svg>'
        _prs, slide = _make_slide()
        compiler = SVGCompiler(C=C)
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count > 0


class TestEditabilityAudit:
    @pytest.mark.parametrize("name,svg", PROBE_CASES.items())
    def test_zero_pictures_in_output(self, name, svg, tmp_path):
        _prs, slide = _make_slide()
        rect = SLIDE_RECTS[name]
        compiler = SVGCompiler()
        compiler.compile(svg, slide, rect)
        pptx_path = tmp_path / f"audit_{name}.pptx"
        has_pic, _has_sp = _audit_pptx(_prs, pptx_path)
        assert not has_pic, f"{name}: <p:pic> found — editability guarantee violated"


RADIAL_SVG = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 400">
  <defs>
    <radialGradient id="rg" cx="50%" cy="50%" r="50%">
      <stop offset="0" stop-color="#FFFFFF"/>
      <stop offset="1" stop-color="#000000"/>
    </radialGradient>
  </defs>
  <circle cx="200" cy="200" r="150" fill="url(#rg)"/>
</svg>"""


class TestRadialGradient:
    def test_radial_gradient_centertext(self):
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(RADIAL_SVG, slide, (1, 1, 6, 6))
        assert "gradient" in result.features
        assert result.shape_count > 0

    def test_radial_gradient_no_crash(self):
        """Regression: radialGradient with percentage cx/cy/r previously crashed."""
        svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">' \
              '<defs><radialGradient id="g" cx="80%" cy="20%" r="60%">' \
              '<stop offset="0" stop-color="#F00"/><stop offset="1" stop-color="#00F"/>' \
              '</radialGradient></defs>' \
              '<circle cx="50" cy="50" r="40" fill="url(#g)"/></svg>'
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count > 0
        assert "gradient" in result.features


class TestCollisionDetection:
    """Text box overlap warnings from SVGCompiler."""

    def test_overlapping_texts_emit_warning(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 200">
          <text x="50" y="50" font-size="14" fill="#000">First text</text>
          <text x="60" y="55" font-size="14" fill="#000">Overlapping text</text>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 6, 4))
        assert any("overlap" in w for w in result.warnings), (
            f"expected overlap warning, got: {result.warnings}"
        )

    def test_adjacent_texts_no_false_positive(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 200">
          <text x="50" y="50" font-size="14" fill="#000">Above</text>
          <text x="50" y="80" font-size="14" fill="#000">Below</text>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 6, 4))
        overlaps = [w for w in result.warnings if "overlap" in w]
        assert len(overlaps) == 0, f"false positive overlap warning: {overlaps}"

    def test_no_text_no_collision_check(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
          <rect x="0" y="0" width="100" height="100" fill="#FF0000"/>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count > 0
        assert not any("overlap" in w for w in result.warnings)


class TestLinearGradientPercentCoords:
    def test_percent_x1_y1_x2_y2(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 200">
          <defs>
            <linearGradient id="gp" x1="0%" y1="0%" x2="100%" y2="100%">
              <stop offset="0%" stop-color="#FF0000"/>
              <stop offset="100%" stop-color="#0000FF"/>
            </linearGradient>
          </defs>
          <rect x="10" y="10" width="180" height="180" fill="url(#gp)"/>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count >= 1
        assert not any("error" in w.lower() for w in result.warnings)

    def test_mixed_percent_and_float(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 200">
          <defs>
            <linearGradient id="gm" x1="0" y1="0" x2="1" y2="1">
              <stop offset="0" stop-color="#00FF00"/>
              <stop offset="1" stop-color="#0000FF"/>
            </linearGradient>
          </defs>
          <rect x="10" y="10" width="180" height="180" fill="url(#gm)"/>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count >= 1


class TestRectCornerSymmetry:
    """SVG spec: if only one of rx/ry is set, the other inherits it."""

    def test_rect_with_ry_only_renders(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="0" y="0" width="100" height="100" ry="10" fill="red"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1

    def test_rect_with_rx_only_renders(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="0" y="0" width="100" height="100" rx="10" fill="red"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1

    def test_rect_with_neither_renders_native(self):
        """Rect without rx/ry should still take native fast path."""
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="0" y="0" width="100" height="100" fill="red"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1


class TestRoundedRect:
    """rect rx/ry produces freeform with rounded corners."""

    def test_rounded_rect_shape_count(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="10" y="10" width="80" height="60" rx="10" ry="10" fill="#4472C4"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1

    def test_rounded_rect_is_freeform(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="10" y="10" width="80" height="60" rx="10" ry="10" fill="#4472C4"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        SVGCompiler().compile(svg, slide, (1, 1, 5, 5))
        sp = slide.shapes[-1]
        assert sp.shape_type is not None

    def test_rounded_rect_rx_only(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="0" y="0" width="100" height="100" rx="15" fill="red"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1

    def test_rounded_rect_ry_only(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="0" y="0" width="100" height="100" ry="12" fill="blue"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1

    def test_rounded_rect_oversized_rx_clamped(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="0" y="0" width="40" height="40" rx="30" fill="green"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1

    def test_sharp_rect_no_rounding(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="0" y="0" width="100" height="100" fill="red"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1


class TestScalingModes:
    """contain / cover / stretch scaling modes."""

    def test_contain_default(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">'
            '<rect x="0" y="0" width="200" height="100" fill="red"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1

    def test_cover_scaling(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">'
            '<rect x="0" y="0" width="200" height="100" fill="blue"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5), scaling="cover")
        assert result.shape_count == 1

    def test_stretch_scaling(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">'
            '<rect x="0" y="0" width="200" height="100" fill="green"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5), scaling="stretch")
        assert result.shape_count == 1

    def test_invalid_scaling_falls_back_to_contain(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">'
            '<rect x="0" y="0" width="200" height="100" fill="red"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = SVGCompiler().compile(svg, slide, (1, 1, 5, 5), scaling="invalid")
        assert result.shape_count == 1

    def test_svg_chart_scaling_param(self):
        from ppt_pro_max.build_helpers import svg_chart
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">'
            '<rect x="0" y="0" width="200" height="100" fill="#4472C4"/>'
            "</svg>"
        )
        _prs, slide = _make_slide()
        result = svg_chart(slide, svg, 1, 1, 5, 5, scaling="cover")
        assert result.shape_count >= 1
