"""Tests for P9-P14: beautify pipeline, SmartArt/Group/OLE extractors, ComponentLibrary, ComponentRenderer.

TDD: Tests define expected behavior before implementation.
"""
import os
import sqlite3
import tempfile
import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

FIXTURES = Path(__file__).parent / "fixtures"


# ═══════════════════════════════════════════════════════════════
# P9: Pipeline beautify branch
# ═══════════════════════════════════════════════════════════════

class TestPipelineBeautify:
    def test_beautify_generates_output(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Test Title"
        run.font.size = Pt(36)

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        prs.save(input_path)

        pipeline = EnterprisePipeline()
        result = pipeline.run_beautify(
            input_pptx=input_path,
            output_path=output_path,
            style="professional",
        )

        assert os.path.isfile(output_path)
        assert result["mode"] == "beautify"
        assert result["num_slides"] >= 1

    def test_beautify_preserves_content(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "My Important Title"
        run.font.size = Pt(36)

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        prs.save(input_path)

        pipeline = EnterprisePipeline()
        pipeline.run_beautify(input_pptx=input_path, output_path=output_path, style="professional")

        out_prs = Presentation(output_path)
        all_text = []
        for slide in out_prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text)
        combined = " ".join(all_text)
        assert "My Important Title" in combined

    def test_beautify_with_brand_spec(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Branded Title"
        run.font.size = Pt(36)

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        prs.save(input_path)

        brand = BrandSpec(
            source="test",
            colors={"primary": "#FF0000", "accent": "#00FF00", "foreground": "#000000", "background": "#FFFFFF"},
            fonts={"heading": "Arial", "body": "Calibri"},
        )

        pipeline = EnterprisePipeline()
        result = pipeline.run_beautify(
            input_pptx=input_path,
            output_path=output_path,
            brand_spec=brand,
        )

        assert os.path.isfile(output_path)


# ═══════════════════════════════════════════════════════════════
# P12: SmartArtExtractor
# ═══════════════════════════════════════════════════════════════

class TestSmartArtExtractor:
    def test_extract_from_art_pptx(self):
        from ppt_pro_max.enterprise.smartart_extractor import SmartArtExtractor

        art_path = str(FIXTURES / "art_restyled_v3.pptx")
        if not os.path.isfile(art_path):
            pytest.skip("art.pptx fixture not available")

        ext = SmartArtExtractor()
        results = ext.extract_all(art_path)

        assert len(results) >= 1
        first = results[0]
        assert first["type"] == "smartart"
        assert "category" in first
        assert "nodes" in first
        assert "xml_parts" in first
        assert len(first["nodes"]) >= 1

    def test_extract_returns_category(self):
        from ppt_pro_max.enterprise.smartart_extractor import SmartArtExtractor

        art_path = str(FIXTURES / "art_restyled_v3.pptx")
        if not os.path.isfile(art_path):
            pytest.skip("art.pptx fixture not available")

        ext = SmartArtExtractor()
        results = ext.extract_all(art_path)

        categories = {r["category"] for r in results}
        assert len(categories) >= 1
        for cat in categories:
            assert cat in ("process", "cycle", "hierarchy", "pyramid", "matrix", "relationship", "picture")

    def test_extract_xml_parts_has_4_keys(self):
        from ppt_pro_max.enterprise.smartart_extractor import SmartArtExtractor

        art_path = str(FIXTURES / "art_restyled_v3.pptx")
        if not os.path.isfile(art_path):
            pytest.skip("art.pptx fixture not available")

        ext = SmartArtExtractor()
        results = ext.extract_all(art_path)

        for r in results:
            parts = r["xml_parts"]
            assert "data" in parts
            assert "layout" in parts
            assert "colors" in parts
            assert "quickStyle" in parts

    def test_extract_no_smartart_returns_empty(self, tmp_path):
        from ppt_pro_max.enterprise.smartart_extractor import SmartArtExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "No SmartArt here"
        run.font.size = Pt(36)

        pptx_path = str(tmp_path / "no_smartart.pptx")
        prs.save(pptx_path)

        ext = SmartArtExtractor()
        results = ext.extract_all(pptx_path)
        assert len(results) == 0

    def test_parse_data_xml_extracts_texts(self):
        from ppt_pro_max.enterprise.smartart_extractor import SmartArtExtractor

        ext = SmartArtExtractor()
        data_xml = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<dgm:data xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <dgm:ptList>
    <dgm:pt modelId="0" type="doc"><dgm:prSet/><dgm:spPr/><dgm:t><a:bodyPr/><a:p><a:r><a:t>Root</a:t></a:r></a:p></dgm:t></dgm:pt>
    <dgm:pt modelId="1"><dgm:prSet/><dgm:spPr/><dgm:t><a:bodyPr/><a:p><a:r><a:t>Step 1</a:t></a:r></a:p></dgm:t></dgm:pt>
    <dgm:pt modelId="2"><dgm:prSet/><dgm:spPr/><dgm:t><a:bodyPr/><a:p><a:r><a:t>Step 2</a:t></a:r></a:p></dgm:t></dgm:pt>
    <dgm:pt modelId="3" type="parTrans"><dgm:prSet/><dgm:spPr/><dgm:t><a:bodyPr/><a:p><a:r><a:t>Trans</a:t></a:r></a:p></dgm:t></dgm:pt>
  </dgm:ptList>
</dgm:data>'''

        nodes = ext._parse_data_xml(data_xml)
        assert len(nodes) == 2
        assert nodes[0]["text"] == "Step 1"
        assert nodes[1]["text"] == "Step 2"

    def test_parse_layout_xml_extracts_category(self):
        from ppt_pro_max.enterprise.smartart_extractor import SmartArtExtractor

        ext = SmartArtExtractor()
        layout_xml = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<dgm:layoutDef xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram">
  <dgm:alg type="lin"/>
</dgm:layoutDef>'''

        result = ext._parse_layout_xml(layout_xml)
        assert result["category"] == "process"


# ═══════════════════════════════════════════════════════════════
# P12: GroupExtractor
# ═══════════════════════════════════════════════════════════════

class TestGroupExtractor:
    def test_extract_from_group_shape(self, tmp_path):
        from ppt_pro_max.enterprise.group_extractor import GroupExtractor
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        from pptx.enum.shapes import MSO_SHAPE as SHAPE
        rect1 = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(3), Inches(1.5))
        rect1.text_frame.paragraphs[0].text = "Box A"
        rect2 = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(5), Inches(2), Inches(3), Inches(1.5))
        rect2.text_frame.paragraphs[0].text = "Box B"

        base_path = str(tmp_path / "group_base.pptx")
        prs.save(base_path)

        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

        grp_sp_xml = f'''<p:grpSp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvGrpSpPr><p:cNvPr id="100" name="TestGroup"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm><a:off x="914400" y="914400"/><a:ext cx="9144000" cy="4572000"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="4572000"/></a:xfrm></p:grpSpPr>
  <p:sp><p:nvSpPr><p:cNvPr id="101" name="Inner1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Step 1</a:t></a:r></a:p></p:txBody></p:sp>
  <p:sp><p:nvSpPr><p:cNvPr id="102" name="Inner2"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="3200400" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Step 2</a:t></a:r></a:p></p:txBody></p:sp>
</p:grpSp>'''

        with zipfile.ZipFile(base_path, 'r') as zin:
            files = {n: zin.read(n) for n in zin.namelist()}

        slide1_xml = files["ppt/slides/slide1.xml"]
        root = etree.fromstring(slide1_xml)
        sp_tree = root.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree")
        if sp_tree is not None:
            sp_tree.append(etree.fromstring(grp_sp_xml))
        files["ppt/slides/slide1.xml"] = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

        pptx_path = str(tmp_path / "with_group.pptx")
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for n, c in files.items():
                zout.writestr(n, c)

        prs2 = Presentation(pptx_path)
        for shape in prs2.slides[0].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                ext = GroupExtractor()
                result = ext.extract(shape)
                assert result["type"] == "group"
                text_values = [t["text"] if isinstance(t, dict) else t for t in result["texts"]]
                assert "Step 1" in text_values
                assert "Step 2" in text_values
                break
        else:
            pytest.fail("No GroupShape found in test PPTX")

    def test_extract_returns_structure(self, tmp_path):
        from ppt_pro_max.enterprise.group_extractor import GroupExtractor

        ext = GroupExtractor()
        assert hasattr(ext, 'extract')
        assert hasattr(ext, '_infer_category')


# ═══════════════════════════════════════════════════════════════
# P12: OLEExtractor
# ═══════════════════════════════════════════════════════════════

class TestOLEExtractor:
    def test_extract_from_pptx_without_ole(self, tmp_path):
        from ppt_pro_max.enterprise.ole_extractor import OLEExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "No OLE"
        run.font.size = Pt(36)

        pptx_path = str(tmp_path / "no_ole.pptx")
        prs.save(pptx_path)

        ext = OLEExtractor()
        results = ext.extract_all(pptx_path)
        assert len(results) == 0

    def test_ole_extractor_has_extract_method(self):
        from ppt_pro_max.enterprise.ole_extractor import OLEExtractor
        ext = OLEExtractor()
        assert hasattr(ext, 'extract')
        assert hasattr(ext, 'extract_all')


# ═══════════════════════════════════════════════════════════════
# P13: ComponentLibrary
# ═══════════════════════════════════════════════════════════════

class TestComponentLibrary:
    def test_init_creates_db(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)
        assert os.path.isfile(db_path)

    def test_add_and_get(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        xml_parts = {
            "data": b"<data/>",
            "layout": b"<layout/>",
            "colors": b"<colors/>",
            "quickStyle": b"<quickStyle/>",
        }
        cid = lib.add(
            type="smartart",
            category="process",
            variant="chevron",
            node_count=4,
            xml_parts=xml_parts,
        )
        assert cid > 0

        comp = lib.get(cid)
        assert comp is not None
        assert comp["type"] == "smartart"
        assert comp["category"] == "process"
        assert comp["variant"] == "chevron"
        assert comp["node_count"] == 4

    def test_search_by_type_category(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        lib.add(type="smartart", category="process", variant="chevron", node_count=4, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})
        lib.add(type="smartart", category="hierarchy", variant="orgchart", node_count=6, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})
        lib.add(type="group", category="infographic", variant="timeline", node_count=5, xml_parts={"data": b"<g/>"})

        results = lib.search(type="smartart", category="process")
        assert len(results) == 1
        assert results[0]["category"] == "process"

    def test_search_by_node_count(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        lib.add(type="smartart", category="process", variant="chevron_4", node_count=4, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})
        lib.add(type="smartart", category="process", variant="chevron_5", node_count=5, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})

        results = lib.search(type="smartart", category="process", node_count=4)
        assert len(results) == 1
        assert results[0]["node_count"] == 4

    def test_match_finds_best(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        lib.add(type="smartart", category="process", variant="chevron", node_count=4, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})
        lib.add(type="smartart", category="process", variant="arrows", node_count=4, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})
        lib.add(type="smartart", category="hierarchy", variant="orgchart", node_count=4, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})

        result = lib.match({"type": "smartart", "category": "process", "node_count": 4})
        assert result is not None
        assert result["category"] == "process"
        assert result["node_count"] == 4

    def test_match_returns_none_when_no_match(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        result = lib.match({"type": "smartart", "category": "pyramid", "node_count": 10})
        assert result is None

    def test_load_xml(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        xml_parts = {
            "data": b"<data>hello</data>",
            "layout": b"<layout>world</layout>",
            "colors": b"<colors>brand</colors>",
            "quickStyle": b"<quickStyle>style</quickStyle>",
        }
        cid = lib.add(type="smartart", category="process", variant="chevron", node_count=4, xml_parts=xml_parts)

        loaded = lib.load_xml(cid)
        assert loaded is not None
        assert loaded["data"] == b"<data>hello</data>"
        assert loaded["layout"] == b"<layout>world</layout>"

    def test_remove(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        cid = lib.add(type="smartart", category="process", variant="chevron", node_count=4, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})
        assert lib.remove(cid) is True
        assert lib.get(cid) is None

    def test_stats(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        lib.add(type="smartart", category="process", variant="chevron", node_count=4, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})
        lib.add(type="group", category="infographic", variant="timeline", node_count=5, xml_parts={"data": b"<g/>"})

        stats = lib.stats()
        assert stats["smartart"] == 1
        assert stats["group"] == 1
        assert stats["total"] == 2

    def test_checksum_dedup(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        xml_parts = {"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"}
        cid1 = lib.add(type="smartart", category="process", variant="chevron", node_count=4, xml_parts=xml_parts)
        cid2 = lib.add(type="smartart", category="process", variant="chevron", node_count=4, xml_parts=xml_parts)

        assert cid1 == cid2

    def test_bulk_import(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "No SmartArt"
        run.font.size = Pt(36)

        pptx_path = str(tmp_path / "no_smartart.pptx")
        prs.save(pptx_path)

        result = lib.bulk_import([pptx_path])
        assert "added" in result
        assert "skipped" in result
        assert "errors" in result


# ═══════════════════════════════════════════════════════════════
# P14: ComponentRenderer
# ═══════════════════════════════════════════════════════════════

class TestComponentRenderer:
    def test_fill_data_replaces_texts(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        renderer = ComponentRenderer()
        data_xml = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<dgm:data xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <dgm:ptList>
    <dgm:pt modelId="0" type="doc"><dgm:prSet/><dgm:spPr/><dgm:t><a:bodyPr/><a:p><a:r><a:t>Root</a:t></a:r></a:p></dgm:t></dgm:pt>
    <dgm:pt modelId="1"><dgm:prSet/><dgm:spPr/><dgm:t><a:bodyPr/><a:p><a:r><a:t>Old 1</a:t></a:r></a:p></dgm:t></dgm:pt>
    <dgm:pt modelId="2"><dgm:prSet/><dgm:spPr/><dgm:t><a:bodyPr/><a:p><a:r><a:t>Old 2</a:t></a:r></a:p></dgm:t></dgm:pt>
  </dgm:ptList>
</dgm:data>'''

        nodes = [{"id": 0, "text": "New 1"}, {"id": 1, "text": "New 2"}]
        result = renderer._fill_data({"data": data_xml}, nodes)
        assert b"New 1" in result["data"]
        assert b"New 2" in result["data"]
        assert b"Old 1" not in result["data"]

    def test_apply_brand_colors_replaces_schemeclr(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        renderer = ComponentRenderer()
        colors_xml = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<dgm:colorsDef xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <dgm:styleLbl name="node0">
    <dgm:fill><a:solidFill><a:schemeClr val="accent1"/></a:solidFill></dgm:fill>
    <dgm:ln><a:solidFill><a:schemeClr val="lt1"/></a:solidFill></dgm:ln>
  </dgm:styleLbl>
</dgm:colorsDef>'''

        brand = BrandSpec(
            source="test",
            colors={"accent1": "D97706", "lt1": "FFFFFF", "dk1": "000000", "tx1": "FFFFFF"},
        )

        result = renderer._apply_brand_colors({"colors": colors_xml}, brand)
        assert b"D97706" in result["colors"]
        assert b"FFFFFF" in result["colors"]
        assert b"accent1" not in result["colors"]

    def test_fallback_diagram_returns_bool(self):
        from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

        renderer = ComponentRenderer()
        result = renderer._fallback_diagram(None, {"category": "process", "nodes": [{"text": "A"}]}, None)
        assert isinstance(result, bool)
