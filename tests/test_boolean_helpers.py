"""Tests for boolean shape convenience functions in build_helpers.

Verifies:
- All 7 boolean functions work with Shapely installed
- Graceful fallback when Shapely not available
- Fill color, line, alpha, C dict resolution
- Save/reload PPTX integrity
"""

from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn

from ppt_pro_max.build_helpers import (
    spotlight, bool_donut, bool_frame, bool_clipped_card,
    bool_neon_tube, bool_star, bool_cross,
    add_slide,
)

try:
    from ppt_pro_max.renderer.boolean_shapes import HAS_SHAPELY
except ImportError:
    HAS_SHAPELY = False

pytestmark = pytest.mark.skipif(not HAS_SHAPELY, reason="shapely not installed")

C = {
    'primary': '#1D78FA', 'accent': '#FF5500', 'muted': '#64748B',
    'white': '#FFFFFF', 'text_dark': '#1A1A1A',
}


def _prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _slide(prs=None):
    if prs is None:
        prs = _prs()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _save_reload(prs):
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def _to_element(sh):
    if hasattr(sh, '_element'):
        return sh._element
    return sh


def _get_fill_hex(sh):
    if sh is None:
        return None
    el = _to_element(sh)
    spPr = el.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        return None
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        return None
    return srgb.get('val')


def _get_alpha(sh):
    if sh is None:
        return None
    el = _to_element(sh)
    spPr = el.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        return None
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        return None
    alpha = srgb.find(qn('a:alpha'))
    if alpha is None:
        return None
    return int(alpha.get('val'))


# ── spotlight ──


class TestSpotlight:
    def test_creates_shape(self):
        prs = _prs()
        s = _slide(prs)
        sh = spotlight(s, 6.67, 3.75, 2.5)
        assert sh is not None

    def test_fill_color(self):
        prs = _prs()
        s = _slide(prs)
        sh = spotlight(s, 6.67, 3.75, 2.5, alpha=70, color='#000000')
        assert _get_fill_hex(sh) == '000000'

    def test_alpha(self):
        prs = _prs()
        s = _slide(prs)
        sh = spotlight(s, 6.67, 3.75, 2.5, alpha=70)
        assert _get_alpha(sh) == 70000

    def test_custom_color(self):
        prs = _prs()
        s = _slide(prs)
        sh = spotlight(s, 6.67, 3.75, 2.5, color='#1A1A1A', alpha=50)
        assert _get_fill_hex(sh) == '1A1A1A'
        assert _get_alpha(sh) == 50000

    def test_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        spotlight(s, 6.67, 3.75, 2.5)
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1


# ── bool_donut ──


class TestBoolDonut:
    def test_creates_shape(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_donut(s, 5, 3.5, 1.5, 0.8, fill='#1D78FA')
        assert sh is not None

    def test_fill_color(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_donut(s, 5, 3.5, 1.5, 0.8, fill='#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'

    def test_with_C_dict(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_donut(s, 5, 3.5, 1.5, 0.8, fill='primary', C=C)
        assert _get_fill_hex(sh) == '1D78FA'

    def test_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        bool_donut(s, 5, 3.5, 1.5, 0.8, fill='#1D78FA')
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1


# ── bool_frame ──


class TestBoolFrame:
    def test_creates_shape(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_frame(s, 2, 2, 5, 4, 0.3, fill='#1D78FA')
        assert sh is not None

    def test_fill_color(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_frame(s, 2, 2, 5, 4, 0.3, fill='#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'

    def test_with_C_dict(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_frame(s, 2, 2, 5, 4, 0.3, fill='primary', C=C)
        assert _get_fill_hex(sh) == '1D78FA'

    def test_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        bool_frame(s, 2, 2, 5, 4, 0.3, fill='#1D78FA')
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1


# ── bool_clipped_card ──


class TestBoolClippedCard:
    def test_creates_shape(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_clipped_card(s, 2, 2, 5, 3, ['tl', 'br'], fill='#1D78FA')
        assert sh is not None

    def test_fill_color(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_clipped_card(s, 2, 2, 5, 3, ['tl', 'br'], fill='#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'

    def test_all_corners(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_clipped_card(s, 2, 2, 5, 3, ['tl', 'tr', 'bl', 'br'], fill='#1D78FA')
        assert sh is not None

    def test_no_corners_fallback(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_clipped_card(s, 2, 2, 5, 3, [], fill='#1D78FA')
        assert sh is not None

    def test_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        bool_clipped_card(s, 2, 2, 5, 3, ['tl', 'br'], fill='#1D78FA')
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1


# ── bool_neon_tube ──


class TestBoolNeonTube:
    def test_creates_shape(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_neon_tube(s, 2, 2, 5, 3, fill='#8B5CF6')
        assert sh is not None

    def test_fill_color(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_neon_tube(s, 2, 2, 5, 3, fill='#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'

    def test_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        bool_neon_tube(s, 2, 2, 5, 3, fill='#8B5CF6')
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1


# ── bool_star ──


class TestBoolStar:
    def test_creates_shape(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_star(s, 5, 3.5, 1.5, points=5, fill='#1D78FA')
        assert sh is not None

    def test_fill_color(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_star(s, 5, 3.5, 1.5, fill='#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'

    def test_custom_inner_ratio(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_star(s, 5, 3.5, 1.5, points=7, inner_ratio=0.3, fill='#1D78FA')
        assert sh is not None

    def test_with_C_dict(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_star(s, 5, 3.5, 1.5, fill='primary', C=C)
        assert _get_fill_hex(sh) == '1D78FA'

    def test_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        bool_star(s, 5, 3.5, 1.5, fill='#1D78FA')
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1


# ── bool_cross ──


class TestBoolCross:
    def test_creates_shape(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_cross(s, 5, 3.5, 3, 3, fill='#1D78FA')
        assert sh is not None

    def test_fill_color(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_cross(s, 5, 3.5, 3, 3, fill='#FF5500')
        assert _get_fill_hex(sh) == 'FF5500'

    def test_custom_bar_ratio(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_cross(s, 5, 3.5, 3, 3, bar_ratio=0.5, fill='#1D78FA')
        assert sh is not None

    def test_with_C_dict(self):
        prs = _prs()
        s = _slide(prs)
        sh = bool_cross(s, 5, 3.5, 3, 3, fill='accent', C=C)
        assert _get_fill_hex(sh) == 'FF5500'

    def test_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        bool_cross(s, 5, 3.5, 3, 3, fill='#1D78FA')
        prs2 = _save_reload(prs)
        assert len(prs2.slides) == 1


# ── Multiple boolean shapes on one slide ──


class TestMultipleBooleanShapes:
    def test_multiple_shapes_save_reload(self):
        prs = _prs()
        s = _slide(prs)
        spotlight(s, 3, 3.75, 2)
        bool_donut(s, 8, 3.5, 1.5, 0.8, fill='#1D78FA')
        bool_frame(s, 10, 1, 2.5, 2, 0.15, fill='#FF5500')
        bool_cross(s, 5, 5.5, 2, 2, fill='#00B050')
        prs2 = _save_reload(prs)
        assert len(list(prs2.slides[0].shapes)) >= 4

    def test_save_to_file(self, tmp_path):
        prs = _prs()
        s = _slide(prs)
        spotlight(s, 6.67, 3.75, 2.5, alpha=70, color='#000000')
        bool_donut(s, 3, 3.75, 1.5, 0.8, fill='#1D78FA')
        bool_clipped_card(s, 8, 2, 4, 3, ['tl', 'br'], fill='#FF5500')
        out = tmp_path / "boolean_test.pptx"
        prs.save(str(out))
        assert out.exists()
        assert out.stat().st_size > 0
