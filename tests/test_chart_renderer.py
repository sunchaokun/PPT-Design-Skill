"""Tests for Data Charts feature (v0.4 I-D)."""

from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches


class TestChartBuilderBasic:

    def test_bar_chart_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "bar",
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": [{"name": "2025", "values": [120, 150, 180, 210]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_line_chart_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "line",
            "categories": ["Jan", "Feb", "Mar"],
            "series": [{"name": "Sales", "values": [10, 20, 30]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_pie_chart_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "pie",
            "categories": ["A", "B", "C"],
            "series": [{"name": "Share", "values": [40, 35, 25]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_doughnut_chart_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "doughnut",
            "categories": ["A", "B"],
            "series": [{"name": "Data", "values": [60, 40]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_area_chart_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "area",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Growth", "values": [100, 200]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_scatter_chart_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "scatter",
            "series": [{"name": "Points", "values": [10, 25, 45, 80]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_bar_stacked_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "bar_stacked",
            "categories": ["Q1", "Q2"],
            "series": [
                {"name": "A", "values": [10, 20]},
                {"name": "B", "values": [30, 40]},
            ],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_bar_100_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "bar_100",
            "categories": ["Q1", "Q2"],
            "series": [
                {"name": "A", "values": [10, 20]},
                {"name": "B", "values": [30, 40]},
            ],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_bar_3d_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "bar_3d",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Data", "values": [10, 20]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_line_markers_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "line_markers",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Data", "values": [10, 20]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_pie_3d_renders(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "pie_3d",
            "categories": ["A", "B"],
            "series": [{"name": "Data", "values": [60, 40]}],
        }
        result = builder.build(slide, config)
        assert result is not None

    def test_unknown_type_fallback_to_bar(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "unknown_type",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Data", "values": [10, 20]}],
        }
        result = builder.build(slide, config)
        assert result is not None


class TestChartBuilderStyle:

    def test_brand_color_mapping(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "bar",
            "categories": ["Q1", "Q2"],
            "series": [
                {"name": "2025", "values": [10, 20]},
                {"name": "2026", "values": [15, 25]},
            ],
            "style": {"color_scheme": "brand"},
        }
        brand_colors = {"primary": "#1A3C6E", "secondary": "#64748B", "accent": "#F97316"}
        result = builder.build(slide, config, brand_colors=brand_colors)
        assert result is not None
        chart = result.chart
        plot = chart.plots[0]
        series0 = plot.series[0]
        series0.format.fill.solid()
        assert series0.format.fill.fore_color.rgb is not None

    def test_legend_control(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "bar",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Data", "values": [10, 20]}],
            "style": {"show_legend": True, "legend_position": "bottom"},
        }
        result = builder.build(slide, config)
        assert result is not None
        assert result.chart.has_legend is True

    def test_chart_title(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "bar",
            "title": "Quarterly Revenue",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Data", "values": [10, 20]}],
        }
        result = builder.build(slide, config)
        assert result is not None
        assert result.chart.has_title is True

    def test_data_labels(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {
            "type": "bar",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Data", "values": [10, 20]}],
            "style": {"show_labels": True},
        }
        result = builder.build(slide, config)
        assert result is not None
        assert result.chart.plots[0].has_data_labels is True


class TestChartLegacyCompat:

    def test_legacy_api_still_works(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        result = builder.build(
            slide,
            "Bar Chart Vertical",
            {"labels": ["Q1", "Q2"], "values": [10, 20]},
            {"colors": {"primary": "#2563EB"}},
            {"x": 1.5, "y": 1.5, "width": 5, "height": 4},
        )
        assert result is not None


class TestChartInPipeline:

    def test_chart_in_populate_slide(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        design = {
            "goal": "data",
            "title": "Revenue",
            "chart": {
                "type": "bar",
                "categories": ["Q1", "Q2"],
                "series": [{"name": "2025", "values": [100, 200]}],
            },
        }
        precision.render_slide(prs, design)
        slide = prs.slides[-1]
        chart_frames = [s for s in slide.shapes if hasattr(s, "chart")]
        assert len(chart_frames) == 1

    def test_no_chart_no_side_effects(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        design = {"goal": "content", "title": "No Chart"}
        precision.render_slide(prs, design)
        slide = prs.slides[-1]
        chart_frames = [s for s in slide.shapes if hasattr(s, "chart")]
        assert len(chart_frames) == 0

    def test_content_json_chart_passthrough(self, tmp_path):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content
        content = {
            "meta": {"title": "Test"},
            "slides": [
                {
                    "goal": "metrics",
                    "title": "Q3 Revenue",
                    "chart": {
                        "type": "bar",
                        "categories": ["Q1", "Q2"],
                        "series": [{"name": "2025", "values": [100, 200]}],
                    },
                },
            ],
        }
        result = load_enterprise_content(content, str(tmp_path))
        assert result[0].get("chart") is not None
        assert result[0]["chart"]["type"] == "bar"
