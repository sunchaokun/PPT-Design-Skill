"""Strict tests for BuildQA's three-tier bounds checking.

The critical behavior under test: distinguishing legitimate decorative
bleeds (low-alpha, no text) from real content overflows. A decorative
bleed must be 'review'; a content overflow must be 'warning' (minor)
or 'fatal' (severe) — never misclassified.

Also covers the full check() flow and report correctness.
"""

import sys
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from ppt_pro_max.build_qa import BuildQA

_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


@pytest.fixture
def prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


@pytest.fixture
def slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_alpha(shape, alpha_val):
    for el in shape._element.iter():
        if el.tag.endswith("}spPr"):
            sf = el.find(f"{{{_NS}}}solidFill")
            if sf is not None:
                clr = sf.find(f"{{{_NS}}}srgbClr")
                if clr is not None:
                    a = etree.SubElement(clr, f"{{{_NS}}}alpha")
                    a.set("val", str(alpha_val))
                    return True
    return False


def make_shape(slide, shape_type, x, y, w, h, alpha=None, text=None):
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    if alpha is not None:
        add_alpha(sh, alpha)
    if text:
        sh.text_frame.text = text
    return sh


class TestIsDecorativeBleed:
    """Direct unit tests of the decor-vs-content classifier."""

    def test_low_alpha_oval_no_text(self, slide):
        sh = make_shape(slide, MSO_SHAPE.OVAL, -1, -1, 3, 3, alpha=5000)
        assert BuildQA()._is_decorative_bleed(sh, Inches(13.333), Inches(7.5)) is True

    def test_low_alpha_rect_no_text(self, slide):
        sh = make_shape(slide, MSO_SHAPE.RECTANGLE, 10, 6, 4, 2, alpha=15000)
        assert BuildQA()._is_decorative_bleed(sh, Inches(13.333), Inches(7.5)) is True

    def test_opaque_oval_with_text(self, slide):
        sh = make_shape(slide, MSO_SHAPE.OVAL, -1, 4, 3, 3, text="内容文本")
        assert BuildQA()._is_decorative_bleed(sh, Inches(13.333), Inches(7.5)) is False

    def test_opaque_oval_no_text(self, slide):
        sh = make_shape(slide, MSO_SHAPE.OVAL, 10, -1, 3, 3)
        assert BuildQA()._is_decorative_bleed(sh, Inches(13.333), Inches(7.5)) is False

    def test_high_alpha_oval_no_text(self, slide):
        sh = make_shape(slide, MSO_SHAPE.OVAL, 10, 6, 3, 3, alpha=50000)
        assert BuildQA()._is_decorative_bleed(sh, Inches(13.333), Inches(7.5)) is False

    def test_alpha_boundary_20_percent(self, slide):
        # alpha == 20% is threshold-inclusive (<= 0.20)
        sh = make_shape(slide, MSO_SHAPE.OVAL, -1, -1, 3, 3, alpha=20000)
        assert BuildQA()._is_decorative_bleed(sh, Inches(13.333), Inches(7.5)) is True

    def test_alpha_just_above_threshold(self, slide):
        sh = make_shape(slide, MSO_SHAPE.OVAL, -1, -1, 3, 3, alpha=21000)
        assert BuildQA()._is_decorative_bleed(sh, Inches(13.333), Inches(7.5)) is False

    def test_text_whitespace_is_not_content(self, slide):
        sh = make_shape(slide, MSO_SHAPE.OVAL, -1, -1, 3, 3, alpha=5000, text="   ")
        assert BuildQA()._is_decorative_bleed(sh, Inches(13.333), Inches(7.5)) is True


class TestBoundsSeverity:
    """End-to-end: shapes in slides → correct severity."""

    def _write_and_check(self, prs, tmp_path):
        path = tmp_path / "test.pptx"
        prs.save(path)
        return BuildQA().check(str(path))

    def test_decorative_bleed_is_review(self, prs, slide, tmp_path):
        make_shape(slide, MSO_SHAPE.OVAL, -1, -1, 3, 3, alpha=5000)
        report = self._write_and_check(prs, tmp_path)
        oob = [i for i in report.fatals + report.warnings + report.reviews
               if i.check_id == "element_out_of_bounds"]
        assert len(oob) == 1
        assert oob[0].severity == "review"
        assert "Decorative shape" in oob[0].message

    def test_content_overflow_minor_is_warning(self, prs, slide, tmp_path):
        # Opaque rectangle bleeding 4% (within 5% minor threshold)
        make_shape(slide, MSO_SHAPE.RECTANGLE, 12.9, 0.5, 0.6, 0.5)  # bleeds 0.17/13.33+7.5
        report = self._write_and_check(prs, tmp_path)
        oob = [i for i in report.fatals + report.warnings + report.reviews
               if i.check_id == "element_out_of_bounds"]
        assert len(oob) == 1
        assert oob[0].severity == "warning"

    def test_content_overflow_severe_is_fatal(self, prs, slide, tmp_path):
        # Opaque rectangle bleeding >5%
        make_shape(slide, MSO_SHAPE.RECTANGLE, 12.0, 0.5, 2.5, 6.5)  # bleeds 1.17" ≈ 5.6%
        report = self._write_and_check(prs, tmp_path)
        oob = [i for i in report.fatals if i.check_id == "element_out_of_bounds"]
        assert len(oob) >= 1
        assert oob[0].severity == "fatal"

    def test_in_bounds_content_no_issue(self, prs, slide, tmp_path):
        make_shape(slide, MSO_SHAPE.RECTANGLE, 1, 1, 4, 3)
        report = self._write_and_check(prs, tmp_path)
        oob = [i for i in report.fatals + report.warnings + report.reviews
               if i.check_id == "element_out_of_bounds"]
        assert len(oob) == 0

    def test_rect_list_with_text_is_fatal_not_deco(self, prs, slide, tmp_path):
        # User's exact concern: a rectangle list overflowing must NOT be treated as decorative
        # Place opaque rect bleeding 6.5% (severe) WITH text
        make_shape(slide, MSO_SHAPE.RECTANGLE, 12.3, 0.5, 2.5, 6.5, text="列表项")
        report = self._write_and_check(prs, tmp_path)
        oob = [i for i in report.fatals + report.warnings + report.reviews
               if i.check_id == "element_out_of_bounds"]
        assert len(oob) == 1
        assert oob[0].severity == "fatal"  # content overflow, definitely not review

    def test_low_alpha_with_text_is_content(self, prs, slide, tmp_path):
        # Even low-alpha, if it has real text it's content, not decor
        make_shape(slide, MSO_SHAPE.OVAL, -2, -2, 6, 6, alpha=10000, text="有文字的半透明框")
        report = self._write_and_check(prs, tmp_path)
        oob = [i for i in report.fatals + report.warnings + report.reviews
               if i.check_id == "element_out_of_bounds"]
        assert len(oob) == 1
        assert oob[0].severity != "review"


class TestQAReportFlow:
    def test_missing_file(self, tmp_path):
        report = BuildQA().check(str(tmp_path / "nonexistent.pptx"))
        assert report.is_passable is False
        assert any(i.check_id == "file_missing" for i in report.fatals)

    def test_empty_deck_passes_structural(self, prs, tmp_path):
        # Presentation with no slides -> 0 slides, no fatal
        path = tmp_path / "blank.pptx"
        prs.save(path)
        report = BuildQA().check(str(path))
        assert report.total_slides == 0
        assert report.is_passable is True  # no slides = nothing to fail

    def test_report_counts_consistent(self, prs, slide, tmp_path):
        make_shape(slide, MSO_SHAPE.RECTANGLE, 1, 1, 4, 3)
        path = tmp_path / "ok.pptx"
        prs.save(path)
        report = BuildQA().check(str(path))
        total_flagged = len(report.fatals) + len(report.warnings) + len(report.reviews)
        assert report.passed + total_flagged == report.total_checks

    def test_format_report_string(self, prs, tmp_path):
        path = tmp_path / "f.pptx"
        prs.save(path)
        report = BuildQA().check(str(path))
        out = BuildQA().format_report(report)
        assert "QA Report" in out
        assert "Result:" in out


class TestScientificMode:
    """BuildQA mode='scientific' lowers min font to 8pt (journal captions)."""

    def _slide_with_small_font(self, prs, slide, pt):
        # Add a small caption text using Pt()
        from pptx.util import Inches, Pt
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = 'Figure 1 caption (journal standard)'
        run.font.size = Pt(pt)
        return box

    def test_9pt_caption_passes_in_scientific(self, prs, slide, tmp_path):
        self._slide_with_small_font(prs, slide, 9)
        path = tmp_path / "sci.pptx"
        prs.save(path)
        report = BuildQA().check(str(path), mode="scientific")
        small = [i for i in report.warnings if i.check_id == "font_too_small"]
        assert len(small) == 0

    def test_9pt_caption_fails_in_business(self, prs, slide, tmp_path):
        self._slide_with_small_font(prs, slide, 9)
        path = tmp_path / "biz.pptx"
        prs.save(path)
        report = BuildQA().check(str(path), mode="business")
        small = [i for i in report.warnings if i.check_id == "font_too_small"]
        assert len(small) >= 1

    def test_7pt_fails_even_in_scientific(self, prs, slide, tmp_path):
        # Even scientific mode has a floor — 7pt is too small
        self._slide_with_small_font(prs, slide, 7)
        path = tmp_path / "tiny.pptx"
        prs.save(path)
        report = BuildQA().check(str(path), mode="scientific")
        small = [i for i in report.warnings if i.check_id == "font_too_small"]
        assert len(small) >= 1

    def test_8pt_boundary_passes_in_scientific(self, prs, slide, tmp_path):
        self._slide_with_small_font(prs, slide, 8)
        path = tmp_path / "boundary.pptx"
        prs.save(path)
        report = BuildQA().check(str(path), mode="scientific")
        small = [i for i in report.warnings if i.check_id == "font_too_small"]
        assert len(small) == 0

    def test_default_mode_is_business(self, prs, slide, tmp_path):
        self._slide_with_small_font(prs, slide, 9)
        path = tmp_path / "default.pptx"
        prs.save(path)
        report = BuildQA().check(str(path))  # no mode arg = business
        small = [i for i in report.warnings if i.check_id == "font_too_small"]
        assert len(small) >= 1
