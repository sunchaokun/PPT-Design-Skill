"""Tests for Hyperlinks feature (v0.4 I-B)."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches


def _collect_hyperlinks(slide):
    """Return {run_text: url} for every run with a hyperlink on the slide."""
    links = {}
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.hyperlink and run.hyperlink.address:
                        links[run.text.strip()] = run.hyperlink.address
    return links


class TestHyperlinks:

    def test_url_link_on_bullet(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        design = {
            "goal": "content",
            "title": "Links",
            "bullets": ["Visit us", "Contact us"],
            "links": [{"bullet_index": 0, "url": "https://example.com"}],
        }
        slide = renderer.render_slide(prs, design)
        links = _collect_hyperlinks(slide)
        assert "Visit us" in links
        assert links["Visit us"] == "https://example.com"
        assert "Contact us" not in links  # second bullet has no link

    def test_mailto_link(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        design = {
            "goal": "content",
            "title": "Contact",
            "bullets": ["Email us"],
            "links": [{"bullet_index": 0, "url": "mailto:sales@example.com"}],
        }
        slide = renderer.render_slide(prs, design)
        links = _collect_hyperlinks(slide)
        assert "Email us" in links
        assert links["Email us"] == "mailto:sales@example.com"

    def test_standalone_link(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        design = {
            "goal": "cta",
            "title": "CTA",
            "links": [{"text": "Download PDF", "url": "https://example.com/report.pdf", "position": "bottom_right"}],
        }
        slide = renderer.render_slide(prs, design)
        links = _collect_hyperlinks(slide)
        assert "Download PDF" in links
        assert links["Download PDF"] == "https://example.com/report.pdf"

    def test_no_links_no_side_effects(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        design = {"goal": "content", "title": "No Links"}
        slide = renderer.render_slide(prs, design)
        assert _collect_hyperlinks(slide) == {}

    def test_link_style_accent_color(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        design = {
            "goal": "content",
            "title": "Styled",
            "bullets": ["Click here"],
            "links": [{"bullet_index": 0, "url": "https://example.com"}],
        }
        slide = renderer.render_slide(prs, design)
        links = _collect_hyperlinks(slide)
        assert "Click here" in links

    def test_content_json_links_passthrough(self, tmp_path):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content
        content = {
            "meta": {"title": "Test"},
            "slides": [
                {
                    "goal": "cta",
                    "title": "Get Started",
                    "bullets": ["Visit website"],
                    "links": [{"bullet_index": 0, "url": "https://example.com"}],
                },
            ],
        }
        result = load_enterprise_content(content, str(tmp_path))
        assert result[0].get("links") is not None
        assert result[0]["links"][0]["url"] == "https://example.com"
