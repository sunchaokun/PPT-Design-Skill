"""Tests for Enterprise Pipeline — Phase A.

TDD: tests written BEFORE implementation.
"""

from __future__ import annotations

import os
import json
import tempfile
from pathlib import Path

import pytest


# ============================================================
# Phase A-1: generate_ppt() 分流结构
# ============================================================

class TestGeneratePptDispatch:
    """Test that generate_ppt() dispatches to FreeStyle or Enterprise pipeline."""

    def test_no_project_goes_freestyle(self):
        """Without --project, generate_ppt() should use FreeStyle pipeline (existing behavior)."""
        from ppt_pro_max import generate_ppt
        result = generate_ppt("test query", dry_run=True)
        assert result.get("dry_run") is True
        assert "strategy" in result

    def test_with_project_goes_enterprise(self, tmp_path):
        """With --project, generate_ppt() should delegate to EnterprisePipeline."""
        from ppt_pro_max import generate_ppt
        project_dir = tmp_path / "my_project"
        project_dir.mkdir()
        result = generate_ppt("test query", project=str(project_dir), dry_run=True)
        assert result.get("dry_run") is True
        assert result.get("pipeline") == "enterprise"

    def test_freestyle_result_unchanged(self):
        """FreeStyle pipeline result should be identical to current behavior."""
        from ppt_pro_max import generate_ppt
        result = generate_ppt("test query", dry_run=True)
        assert "page_count" in result or "pages" in result


# ============================================================
# Phase A-2: BrandSpec data structure
# ============================================================

class TestBrandSpec:
    """Test BrandSpec dataclass."""

    def test_default_brand_spec(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        spec = BrandSpec()
        assert spec.source == "none"
        assert spec.colors is None
        assert spec.fonts is None
        assert spec.spacing is None
        assert spec.logo is None
        assert spec.layout_mapping is None
        assert spec.template_layouts is None
        assert spec.dark_mode is False

    def test_brand_spec_from_brand_json(self, tmp_path):
        """BrandSpec can be created from brand.json data."""
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        brand_data = {
            "colors": {"primary": "#1A3C6E", "foreground": "#333333"},
            "fonts": {"heading": "Arial", "body": "Calibri"},
            "spacing": {"title_size_pt": 28, "body_size_pt": 14},
            "logo": {"position": "top_right", "on_every_slide": True},
            "layout_mapping": {"hook": 0, "content": 2},
        }
        spec = BrandSpec.from_brand_json(brand_data)
        assert spec.source == "brand_json"
        assert spec.colors["primary"] == "#1A3C6E"
        assert spec.fonts["heading"] == "Arial"
        assert spec.logo["position"] == "top_right"
        assert spec.layout_mapping["hook"] == 0

    def test_brand_spec_merge_template_overrides(self):
        """brand.json fields override template-extracted fields."""
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        template_spec = BrandSpec(
            source="template",
            colors={"primary": "#000000", "foreground": "#111111"},
            fonts={"heading": "Times", "body": "Arial"},
        )
        brand_data = {"colors": {"primary": "#1A3C6E"}}
        merged = BrandSpec.merge(template_spec, brand_data)
        assert merged.source == "merged"
        assert merged.colors["primary"] == "#1A3C6E"  # brand.json overrides
        assert merged.colors["foreground"] == "#111111"  # template preserved
        assert merged.fonts["heading"] == "Times"  # template preserved


# ============================================================
# Phase A-3: ProjectScanner
# ============================================================

class TestProjectScanner:
    """Test ProjectScanner asset detection."""

    def test_scan_empty_project(self, tmp_path):
        """Empty project directory should return ProjectAsset with all None."""
        from ppt_pro_max.enterprise.scanner import ProjectScanner, ProjectAsset
        scanner = ProjectScanner()
        asset = scanner.scan(str(tmp_path))
        assert asset.project_dir == str(tmp_path)
        assert asset.template_path is None
        assert asset.logo_path is None
        assert asset.brand_raw is None
        assert asset.content_raw is None
        assert asset.image_pool == []

    def test_scan_full_project(self, tmp_path):
        """Project with all assets should detect everything."""
        from ppt_pro_max.enterprise.scanner import ProjectScanner
        # Create assets
        (tmp_path / "template.pptx").write_bytes(b"PK\x03\x04" + b"\x00" * 20)  # minimal zip header
        (tmp_path / "logo.png").write_bytes(b"\x89PNG")
        (tmp_path / "brand.json").write_text('{"colors": {"primary": "#1A3C6E"}}', encoding="utf-8")
        (tmp_path / "content.json").write_text('{"slides": []}', encoding="utf-8")
        (tmp_path / "hero.png").write_bytes(b"\x89PNG")
        (tmp_path / "dashboard.jpg").write_bytes(b"\xff\xd8\xff\xe0")

        scanner = ProjectScanner()
        asset = scanner.scan(str(tmp_path))
        assert asset.template_path is not None
        assert asset.logo_path is not None
        assert asset.brand_raw is not None
        assert asset.content_raw is not None
        assert len(asset.image_pool) >= 1  # hero.png at least (logo excluded from pool)

    def test_scan_logo_detection(self, tmp_path):
        """LOGO detection: filename 'logo', 'logo_xxx', 'xxx_logo'."""
        from ppt_pro_max.enterprise.scanner import ProjectScanner
        scanner = ProjectScanner()

        # Exact: logo.png
        (tmp_path / "logo.png").write_bytes(b"\x89PNG")
        asset = scanner.scan(str(tmp_path))
        assert asset.logo_path is not None

    def test_scan_logo_prefix(self, tmp_path):
        """LOGO detection: logo-company.png."""
        from ppt_pro_max.enterprise.scanner import ProjectScanner
        scanner = ProjectScanner()
        (tmp_path / "logo-company.png").write_bytes(b"\x89PNG")
        asset = scanner.scan(str(tmp_path))
        assert asset.logo_path is not None

    def test_scan_excludes_output_dir(self, tmp_path):
        """output/ subdirectory should be excluded from scanning."""
        from ppt_pro_max.enterprise.scanner import ProjectScanner
        scanner = ProjectScanner()
        (tmp_path / "output").mkdir()
        (tmp_path / "output" / "v1").mkdir()
        (tmp_path / "output" / "v1" / "presentation.pptx").write_bytes(b"PK\x03\x04" + b"\x00" * 20)
        asset = scanner.scan(str(tmp_path))
        assert asset.template_path is None  # output/ .pptx should not be detected as template


# ============================================================
# Phase A-4: TemplateAnalyzer
# ============================================================

class TestTemplateAnalyzer:
    """Test TemplateAnalyzer theme/layout extraction."""

    def test_analyze_default_template(self):
        """Analyzing default template should extract colors and layouts."""
        from ppt_pro_max.enterprise.template_analyzer import TemplateAnalyzer
        from pptx import Presentation
        prs = Presentation()
        prs.save(str(Path(tempfile.gettempdir()) / "_test_default.pptx"))
        analyzer = TemplateAnalyzer()
        spec = analyzer.analyze(str(Path(tempfile.gettempdir()) / "_test_default.pptx"))
        assert spec.source == "template"
        assert spec.colors is not None
        assert spec.fonts is not None
        assert spec.template_layouts is not None
        assert len(spec.template_layouts) > 0

    def test_analyze_extracts_layout_names(self):
        """TemplateAnalyzer should extract layout names from slide_layouts."""
        from ppt_pro_max.enterprise.template_analyzer import TemplateAnalyzer
        from pptx import Presentation
        path = str(Path(tempfile.gettempdir()) / "_test_layouts.pptx")
        prs = Presentation()
        prs.save(path)
        analyzer = TemplateAnalyzer()
        spec = analyzer.analyze(path)
        layout_names = [l["name"] for l in spec.template_layouts]
        assert "Title Slide" in layout_names
        assert "Blank" in layout_names

    def test_analyze_extracts_theme_colors(self):
        """TemplateAnalyzer should extract theme colors via XML parsing."""
        from ppt_pro_max.enterprise.template_analyzer import TemplateAnalyzer
        from pptx import Presentation
        path = str(Path(tempfile.gettempdir()) / "_test_colors.pptx")
        prs = Presentation()
        prs.save(path)
        analyzer = TemplateAnalyzer()
        spec = analyzer.analyze(path)
        assert "primary" in spec.colors or "dk1" in spec.colors


# ============================================================
# Phase C+: PageOp and parse_pages
# ============================================================

class TestPageOp:
    """Test PageOp dataclass and parse_pages()."""

    def test_parse_single_page(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, parse_pages
        ops = parse_pages("3")
        assert len(ops) == 1
        assert ops[0].page == 3
        assert ops[0].action == "modify"

    def test_parse_multiple_pages(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, parse_pages
        ops = parse_pages("3,5,8")
        assert len(ops) == 3
        assert [o.page for o in ops] == [3, 5, 8]
        assert all(o.action == "modify" for o in ops)

    def test_parse_range(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        ops = parse_pages("3-5")
        assert len(ops) == 3
        assert [o.page for o in ops] == [3, 4, 5]

    def test_parse_insert(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, parse_pages
        ops = parse_pages("+5")
        assert len(ops) == 1
        assert ops[0].page == 5
        assert ops[0].action == "insert"
        assert ops[0].insert_after is True

    def test_parse_delete(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, parse_pages
        ops = parse_pages("-6")
        assert len(ops) == 1
        assert ops[0].page == 6
        assert ops[0].action == "delete"

    def test_parse_move(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, parse_pages
        ops = parse_pages("3>5")
        assert len(ops) == 1
        assert ops[0].page == 3
        assert ops[0].action == "move"
        assert ops[0].target == 5

    def test_parse_swap(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, parse_pages
        ops = parse_pages("3<>7")
        assert len(ops) == 1
        assert ops[0].page == 3
        assert ops[0].action == "swap"
        assert ops[0].target == 7

    def test_parse_mixed(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        ops = parse_pages("1,3-5,+6,-8,3>5,3<>7")
        actions = [o.action for o in ops]
        assert "modify" in actions
        assert "insert" in actions
        assert "delete" in actions
        assert "move" in actions
        assert "swap" in actions

    def test_parse_validation_out_of_range(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        with pytest.raises(ValueError, match="超出范围"):
            parse_pages("15", num_slides=10)

    def test_parse_validation_valid_range(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        ops = parse_pages("5", num_slides=10)
        assert len(ops) == 1


# ============================================================
# Phase C+: compute_target_sequence
# ============================================================

class TestComputeTargetSequence:
    """Test sequence transformation algorithm."""

    def test_delete_pages(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, compute_target_sequence
        ops = [PageOp(page=3, action="delete"), PageOp(page=7, action="delete")]
        order, new = compute_target_sequence(10, ops)
        assert order == [0, 1, 3, 4, 5, 7, 8, 9]

    def test_move_page(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, compute_target_sequence
        ops = [PageOp(page=3, action="move", target=5)]
        order, new = compute_target_sequence(8, ops)
        assert order == [0, 1, 3, 4, 2, 5, 6, 7]

    def test_swap_pages(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, compute_target_sequence
        ops = [PageOp(page=3, action="swap", target=7)]
        order, new = compute_target_sequence(8, ops)
        assert order == [0, 1, 6, 3, 4, 5, 2, 7]

    def test_delete_and_insert(self):
        from ppt_pro_max.enterprise.page_revision import PageOp, compute_target_sequence
        ops = [
            PageOp(page=3, action="delete"),
            PageOp(page=7, action="delete"),
            PageOp(page=4, action="insert"),
        ]
        order, new = compute_target_sequence(10, ops)
        assert order == [0, 1, 3, 4, 5, 7, 8, 9]
        assert len(new) > 0  # insert position recorded

    def test_no_ops_returns_identity(self):
        from ppt_pro_max.enterprise.page_revision import compute_target_sequence
        order, new = compute_target_sequence(5, [])
        assert order == [0, 1, 2, 3, 4]
        assert new == {}
