import pytest
from pptx import Presentation
from pptx.util import Inches
from ppt_pro_max.renderer.svg_compiler import SVGCompiler

class TestSVGUseElement:
    def test_use_basic(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
          <defs>
            <rect id="my-rect" x="0" y="0" width="10" height="10" fill="#FF0000"/>
          </defs>
          <use href="#my-rect" x="20" y="30"/>
        </svg>"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1
        assert "use" in result.features
        sh = slide.shapes[0]
        # Bounding box should be shifted by (20, 30)
        # Rect in Inches is scaled from 100x100 viewBox to 5x5 inches.
        # So x=20/100*5 + 1 = 2 inches, y=30/100*5 + 1 = 2.5 inches
        # width = 10/100*5 = 0.5 inches, height = 0.5 inches
        assert abs(sh.left - Inches(2)) < 0.01
        assert abs(sh.top - Inches(2.5)) < 0.01

    def test_use_xlink_href(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" viewBox="0 0 100 100">
          <defs>
            <circle id="my-circle" cx="0" cy="0" r="10" fill="#00FF00"/>
          </defs>
          <use xlink:href="#my-circle" x="50" y="50"/>
        </svg>"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 1
        assert "use" in result.features

    def test_use_nested_g(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
          <defs>
            <g id="my-group">
              <rect x="0" y="0" width="10" height="10" fill="#FF0000"/>
              <circle cx="20" cy="20" r="5" fill="#0000FF"/>
            </g>
          </defs>
          <use href="#my-group" x="10" y="10"/>
        </svg>"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count == 2
        assert "use" in result.features
