"""Tests for template-clone pipeline: expand_and_patch, plan_pages, DeliveryGate, Hybrid rerender."""

import os
import json
import tempfile
import pytest

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _create_multi_page_pptx(path: str, num_content: int = 2):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layouts = prs.slide_layouts
    blank = layouts[6] if len(layouts) > 6 else layouts[-1]

    cover = prs.slides.add_slide(layouts[0] if layouts else blank)
    for ph in cover.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = "Template Cover"
        elif ph.placeholder_format.type == 2:
            ph.text = "Template Subtitle"

    toc = prs.slides.add_slide(blank)
    tb = toc.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "目录"
    run.font.size = Pt(36)
    run.font.bold = True
    body = toc.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    btf = body.text_frame
    for j, text in enumerate(["Section A", "Section B"]):
        bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        br = bp.add_run()
        br.text = text
        br.font.size = Pt(14)

    trans = prs.slides.add_slide(blank)
    tb2 = trans.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "PART ONE"
    r2.font.size = Pt(48)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(0x2E, 0x65, 0x04)

    for i in range(num_content):
        slide = prs.slides.add_slide(blank)
        tb3 = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = f"Content Page {i + 1}"
        r3.font.size = Pt(28)
        r3.font.bold = True

        body2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        btf2 = body2.text_frame
        for j, text in enumerate(["Bullet A", "Bullet B", "Bullet C"]):
            bp2 = btf2.paragraphs[0] if j == 0 else btf2.add_paragraph()
            br2 = bp2.add_run()
            br2.text = text
            br2.font.size = Pt(14)

    back = prs.slides.add_slide(layouts[0] if layouts else blank)
    for ph in back.placeholders:
        if ph.placeholder_format.type == 1:
            ph.text = "Thank You"
        elif ph.placeholder_format.type == 2:
            ph.text = "Contact Us"

    prs.save(path)


def _create_placeholder_pptx(path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "单击此处添加标题"
    r.font.size = Pt(36)

    body = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
    btf = body.text_frame
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = "请输入内容"
    br.font.size = Pt(14)

    prs.save(path)


def _create_composite_title_pptx(path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    slide = prs.slides.add_slide(blank)
    for i, char in enumerate("乡村"):
        tb = slide.shapes.add_textbox(Inches(2 + i * 1.5), Inches(1), Inches(1.5), Inches(2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = char
        r.font.size = Pt(72)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x2E, 0x65, 0x04)

    prs.save(path)


class TestPagePlan:

    def test_page_plan_dataclass(self):
        from ppt_pro_max.enterprise.design_dna_extractor import PagePlan
        plan = PagePlan(page_type="content", title="Test", bullets=["a", "b"])
        assert plan.page_type == "content"
        assert plan.title == "Test"
        assert plan.bullets == ["a", "b"]
        assert plan.layout == "auto"

    def test_page_plan_with_layout(self):
        from ppt_pro_max.enterprise.design_dna_extractor import PagePlan
        plan = PagePlan(page_type="content", title="Test", layout="cards")
        assert plan.layout == "cards"


class TestPlanPages:

    def test_plan_pages_from_content_json(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        story_plan = {
            "strategy": "content_json",
            "pages": [
                {"goal": "hook", "title": "Cover Title", "subtitle": "Cover Sub"},
                {"goal": "overview", "title": "目录", "bullets": ["A", "B", "C"]},
                {"goal": "section", "title": "Part One"},
                {"goal": "features", "title": "Features", "bullets": ["F1", "F2", "F3"], "cards": [{"title": "Card1"}]},
                {"goal": "cta", "title": "Thank You"},
            ],
        }

        extractor = DesignDNAExtractor()
        plans = extractor.plan_pages(story_plan)

        assert len(plans) == 5
        assert plans[0].page_type == "cover"
        assert plans[0].title == "Cover Title"
        assert plans[1].page_type == "toc"
        assert plans[2].page_type == "transition"
        assert plans[3].page_type == "content"
        assert plans[3].layout == "cards"
        assert plans[4].page_type == "back_cover"

    def test_plan_pages_auto_layout_bullets(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        story_plan = {
            "strategy": "auto",
            "pages": [
                {"goal": "content", "title": "Test", "bullets": ["A", "B", "C"]},
            ],
        }

        extractor = DesignDNAExtractor()
        plans = extractor.plan_pages(story_plan)

        assert plans[0].layout == "bullets"

    def test_plan_pages_auto_layout_process(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        story_plan = {
            "strategy": "auto",
            "pages": [
                {"goal": "content", "title": "Test", "bullets": ["步骤一：开始", "步骤二：执行", "步骤三：完成"]},
            ],
        }

        extractor = DesignDNAExtractor()
        plans = extractor.plan_pages(story_plan)

        assert plans[0].layout == "process"

    def test_plan_pages_section_counter(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        story_plan = {
            "strategy": "auto",
            "pages": [
                {"goal": "hook", "title": "Cover"},
                {"goal": "section", "title": "Part 1"},
                {"goal": "content", "title": "Content 1"},
                {"goal": "section", "title": "Part 2"},
                {"goal": "content", "title": "Content 2"},
            ],
        }

        extractor = DesignDNAExtractor()
        plans = extractor.plan_pages(story_plan)

        assert plans[0].section_number == 0
        assert plans[1].section_number == 1
        assert plans[2].section_number == 1
        assert plans[3].section_number == 2
        assert plans[4].section_number == 2

    def test_plan_pages_with_page_contents(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        story_plan = {
            "strategy": "auto",
            "pages": [
                {"goal": "content", "title": "Test"},
            ],
        }
        page_contents = [
            {"bullets": ["X", "Y", "Z"], "layout": "cards"},
        ]

        extractor = DesignDNAExtractor()
        plans = extractor.plan_pages(story_plan, page_contents)

        assert plans[0].bullets == ["X", "Y", "Z"]
        assert plans[0].layout == "cards"


class TestExpandAndPatch:

    def test_expand_same_page_count(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        pptx_path = str(tmp_path / "template.pptx")
        _create_multi_page_pptx(pptx_path, num_content=1)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        plans = [
            PagePlan(page_type="cover", title="New Cover", subtitle="New Sub"),
            PagePlan(page_type="toc", title="New TOC", bullets=["A", "B"]),
            PagePlan(page_type="transition", title="New Section"),
            PagePlan(page_type="content", title="New Content", bullets=["X", "Y"]),
            PagePlan(page_type="back_cover", title="New Ending"),
        ]

        output_path = str(tmp_path / "output.pptx")
        result = extractor.expand_and_patch(dna, plans, output_path)

        assert os.path.isfile(result.output_path)
        assert result.total_slides == 5
        assert result.cloned_count == 0

        prs = Presentation(result.output_path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text)
        combined = " ".join(all_text)
        assert "New Cover" in combined or "New Content" in combined

    def test_expand_clones_content_pages(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        pptx_path = str(tmp_path / "template.pptx")
        _create_multi_page_pptx(pptx_path, num_content=1)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        plans = [
            PagePlan(page_type="cover", title="Cover"),
            PagePlan(page_type="content", title="Content 1", bullets=["A"]),
            PagePlan(page_type="content", title="Content 2", bullets=["B"]),
            PagePlan(page_type="content", title="Content 3", bullets=["C"]),
            PagePlan(page_type="back_cover", title="End"),
        ]

        output_path = str(tmp_path / "output.pptx")
        result = extractor.expand_and_patch(dna, plans, output_path)

        assert result.total_slides == 5
        assert result.cloned_count >= 2

        prs = Presentation(result.output_path)
        assert len(prs.slides) == 5

    def test_expand_deletes_unused_pages(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        pptx_path = str(tmp_path / "template.pptx")
        _create_multi_page_pptx(pptx_path, num_content=2)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        original_slides = len(dna.slides)

        plans = [
            PagePlan(page_type="cover", title="Cover"),
            PagePlan(page_type="content", title="Content 1", bullets=["A"]),
            PagePlan(page_type="back_cover", title="End"),
        ]

        output_path = str(tmp_path / "output.pptx")
        result = extractor.expand_and_patch(dna, plans, output_path)

        assert result.total_slides == 3
        assert len(result.deleted_originals) > 0

    def test_expand_template_mapping(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        pptx_path = str(tmp_path / "template.pptx")
        _create_multi_page_pptx(pptx_path, num_content=1)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        plans = [
            PagePlan(page_type="cover", title="Cover"),
            PagePlan(page_type="content", title="Content", bullets=["A"]),
            PagePlan(page_type="content", title="Content 2", bullets=["B"]),
            PagePlan(page_type="back_cover", title="End"),
        ]

        output_path = str(tmp_path / "output.pptx")
        result = extractor.expand_and_patch(dna, plans, output_path)

        assert len(result.template_mapping) == 4
        non_clone = [m for m in result.template_mapping if not m["is_clone"]]
        clone = [m for m in result.template_mapping if m["is_clone"]]
        assert len(non_clone) >= 2
        assert len(clone) >= 1


class TestNeedsContentRerender:

    def test_no_rerender_for_non_content(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, SlideDNA, PagePlan,
        )

        extractor = DesignDNAExtractor()
        plan = PagePlan(page_type="cover", title="Cover")
        slide_dna = SlideDNA(
            slide_index=0, page_type="cover", slide_xml=b"", rels_xml=b"",
        )
        assert not extractor.needs_content_rerender(slide_dna, plan)

    def test_rerender_when_layout_non_auto(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, SlideDNA, TextZone, PagePlan,
        )

        extractor = DesignDNAExtractor()
        plan = PagePlan(page_type="content", title="Test", layout="cards", cards=[{"title": "C1"}])
        slide_dna = SlideDNA(
            slide_index=0, page_type="content", slide_xml=b"", rels_xml=b"",
            text_zones=[
                TextZone(zone_id="b_0", role="body", text="old"),
            ],
        )
        assert extractor.needs_content_rerender(slide_dna, plan)

    def test_rerender_when_bullet_count_mismatch(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, SlideDNA, TextZone, PagePlan,
        )

        extractor = DesignDNAExtractor()
        plan = PagePlan(page_type="content", title="Test", bullets=["A", "B", "C", "D", "E"])
        slide_dna = SlideDNA(
            slide_index=0, page_type="content", slide_xml=b"", rels_xml=b"",
            text_zones=[
                TextZone(zone_id="b_0", role="body", text="old1"),
                TextZone(zone_id="b_1", role="body", text="old2"),
            ],
        )
        assert extractor.needs_content_rerender(slide_dna, plan)

    def test_rerender_when_has_bullets(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, SlideDNA, TextZone, PagePlan,
        )

        extractor = DesignDNAExtractor()
        plan = PagePlan(page_type="content", title="Test", bullets=["A", "B", "C"])
        slide_dna = SlideDNA(
            slide_index=0, page_type="content", slide_xml=b"", rels_xml=b"",
            text_zones=[
                TextZone(zone_id="b_0", role="body", text="old1"),
                TextZone(zone_id="b_1", role="body", text="old2"),
                TextZone(zone_id="b_2", role="body", text="old3"),
            ],
        )
        assert extractor.needs_content_rerender(slide_dna, plan)


class TestDeliveryGate:

    def test_check_residual_placeholder(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate

        pptx_path = str(tmp_path / "test.pptx")
        _create_placeholder_pptx(pptx_path)

        gate = DeliveryGate()
        report = gate.check(pptx_path, None, [])

        fatal_ids = [i.check_id for i in report.fatals]
        assert "residual_placeholder" in fatal_ids

    def test_auto_fix_residual_placeholder(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate

        pptx_path = str(tmp_path / "test.pptx")
        _create_placeholder_pptx(pptx_path)

        gate = DeliveryGate()
        report = gate.check(pptx_path, None, [])
        gate.auto_fix(pptx_path, None, [], report)

        report2 = gate.check(pptx_path, None, [])
        placeholder_fatals = [i for i in report2.fatals if i.check_id == "residual_placeholder"]
        assert len(placeholder_fatals) == 0

    def test_check_blank_page(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        prs.slides.add_slide(blank)

        pptx_path = str(tmp_path / "blank.pptx")
        prs.save(pptx_path)

        gate = DeliveryGate()
        report = gate.check(pptx_path, None, [])

        fatal_ids = [i.check_id for i in report.fatals]
        assert "blank_page" in fatal_ids

    def test_check_page_count_mismatch(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        from ppt_pro_max.enterprise.design_dna_extractor import PagePlan

        pptx_path = str(tmp_path / "test.pptx")
        _create_multi_page_pptx(pptx_path, num_content=1)

        plans = [PagePlan(page_type="content", title="A") for _ in range(3)]

        gate = DeliveryGate()
        report = gate.check(pptx_path, None, plans)

        fatal_ids = [i.check_id for i in report.fatals]
        assert "page_count_mismatch" in fatal_ids

    def test_check_passes_clean_pptx(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        from ppt_pro_max.enterprise.design_dna_extractor import PagePlan

        pptx_path = str(tmp_path / "test.pptx")
        _create_multi_page_pptx(pptx_path, num_content=2)

        plans = [PagePlan(page_type="content", title="A") for _ in range(6)]

        gate = DeliveryGate()
        report = gate.check(pptx_path, None, plans)

        assert report.total_slides == 6

    def test_format_report(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate, QualityReport, CheckItem

        report = QualityReport(
            total_slides=10,
            total_checks=15,
            passed=12,
            fatals=[CheckItem("content", "blank_page", "fatal", 3, "Blank page")],
            warnings=[CheckItem("design", "color_break", "warning", 5, "Color break")],
        )

        gate = DeliveryGate()
        text = gate.format_report(report)

        assert "FATAL" in text or "fatal" in text.lower()
        assert "WARNING" in text or "warning" in text.lower()
        assert "10" in text

    def test_quality_report_is_passable(self):
        from ppt_pro_max.enterprise.delivery_gate import QualityReport, CheckItem

        report_pass = QualityReport(total_slides=5, total_checks=10, passed=10)
        assert report_pass.is_passable

        report_fail = QualityReport(
            total_slides=5, total_checks=10, passed=8,
            fatals=[CheckItem("content", "blank_page", "fatal", 0, "Blank")],
        )
        assert not report_fail.is_passable


class TestDnaHasVisualContent:

    def test_visual_dna(self, tmp_path):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

        pptx_path = str(tmp_path / "test.pptx")
        _create_multi_page_pptx(pptx_path, num_content=1)

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert DesignDNAExtractor.dna_has_visual_content(dna)

    def test_empty_dna(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNA, DesignDNAExtractor

        dna = DesignDNA(source_path="none")
        assert not DesignDNAExtractor.dna_has_visual_content(dna)


class TestPipelineTemplateClone:

    def test_pipeline_template_clone_branch(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        project_dir = str(tmp_path / "project")
        os.makedirs(project_dir)

        template_path = os.path.join(project_dir, "template.pptx")
        _create_multi_page_pptx(template_path, num_content=1)

        content = {
            "pages": [
                {"goal": "hook", "title": "Test Cover", "subtitle": "Test Sub"},
                {"goal": "content", "title": "Test Content", "bullets": ["A", "B"]},
                {"goal": "cta", "title": "Test End"},
            ]
        }
        content_path = os.path.join(project_dir, "content.json")
        with open(content_path, "w", encoding="utf-8") as f:
            json.dump(content, f, ensure_ascii=False)

        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="test query",
            project_dir=project_dir,
            content_file=content_path,
        )

        assert result.get("render_mode") == "template_clone"
        assert result.get("num_slides", 0) >= 1

    def test_pipeline_no_template_goes_precision(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        project_dir = str(tmp_path / "project")
        os.makedirs(project_dir)

        content = {
            "pages": [
                {"goal": "hook", "title": "Test Cover"},
                {"goal": "content", "title": "Test Content", "bullets": ["A"]},
            ]
        }
        content_path = os.path.join(project_dir, "content.json")
        with open(content_path, "w", encoding="utf-8") as f:
            json.dump(content, f, ensure_ascii=False)

        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="test query",
            project_dir=project_dir,
            content_file=content_path,
        )

        assert result.get("render_mode") == "precision"


class TestExpandAndPatchEndToEnd:

    def test_real_pptx_expand(self, tmp_path):
        pptx_path = r"E:\PPT-Design-Skill\docs\05乡村振兴.pptx"
        if not os.path.isfile(pptx_path):
            pytest.skip("05乡村振兴.pptx not found")

        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        assert len(dna.slides) == 5

        plans = [
            PagePlan(page_type="cover", title="数字乡村建设规划", subtitle="智慧引领 · 数字赋能"),
            PagePlan(page_type="toc", title="目录", bullets=["基础设施", "智慧农业", "数字治理"]),
            PagePlan(page_type="transition", title="数字基础设施", subtitle="PART ONE"),
            PagePlan(page_type="content", title="农村宽带现状", bullets=["覆盖率98%", "5G延伸", "物联网部署", "数据中心建设"]),
            PagePlan(page_type="content", title="智慧农业应用", bullets=["精准种植", "智能灌溉", "无人机巡检"]),
            PagePlan(page_type="back_cover", title="携手共建", subtitle="数字乡村新未来"),
        ]

        output_path = str(tmp_path / "expanded.pptx")

        result = extractor.expand_and_patch(dna, plans, output_path)

        assert os.path.isfile(result.output_path)
        assert result.total_slides == 6
        assert result.cloned_count >= 1

        prs = Presentation(result.output_path)
        assert len(prs.slides) == 6

    def test_real_pptx_with_rerender(self, tmp_path):
        pptx_path = r"E:\PPT-Design-Skill\docs\05乡村振兴.pptx"
        if not os.path.isfile(pptx_path):
            pytest.skip("05乡村振兴.pptx not found")

        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        plans = [
            PagePlan(page_type="cover", title="数字乡村", subtitle="新规划"),
            PagePlan(page_type="content", title="多内容页", layout="bullets",
                     bullets=["第一条", "第二条", "第三条", "第四条", "第五条"]),
            PagePlan(page_type="back_cover", title="谢谢"),
        ]

        output_path = str(tmp_path / "rerendered.pptx")

        result = extractor.expand_and_patch(dna, plans, output_path)

        assert os.path.isfile(result.output_path)
        assert result.total_slides == 3
        assert len(result.deleted_originals) > 0

    def test_delivery_gate_on_expanded(self, tmp_path):
        pptx_path = r"E:\PPT-Design-Skill\docs\05乡村振兴.pptx"
        if not os.path.isfile(pptx_path):
            pytest.skip("05乡村振兴.pptx not found")

        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor, PagePlan
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate

        extractor = DesignDNAExtractor()
        dna = extractor.extract(pptx_path)

        plans = [
            PagePlan(page_type="cover", title="数字乡村", subtitle="规划"),
            PagePlan(page_type="content", title="内容", bullets=["A", "B"]),
            PagePlan(page_type="back_cover", title="谢谢"),
        ]

        output_path = str(tmp_path / "gated.pptx")

        result = extractor.expand_and_patch(dna, plans, output_path)

        gate = DeliveryGate()
        report = gate.check(result.output_path, dna, plans)

        assert report.total_slides == 3
