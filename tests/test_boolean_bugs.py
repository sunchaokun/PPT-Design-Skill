"""Tests for boolean_shapes.py BUG fixes.

BUG 1: bool_symdiff(a, b) calls a.symmetric_difference(a) instead of (b)
       → always returns empty geometry (self-symdiff = empty set)

BUG 2: bool_image(geometry, ...) ignores geometry parameter
       → always renders as MSO_SHAPE.OVAL regardless of custom geometry
"""

from __future__ import annotations

import os
import tempfile

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

try:
    from shapely.geometry import Polygon
    HAS_SHAPELY = True
except ImportError:
    HAS_SHAPELY = False

pytestmark = pytest.mark.skipif(not HAS_SHAPELY, reason="shapely not installed")


def _prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _slide(prs=None):
    if prs is None:
        prs = _prs()
    return prs.slides.add_slide(prs.slide_layouts[6])


@pytest.fixture
def test_image():
    img = PILImage.new("RGB", (200, 200), (255, 100, 50))
    path = os.path.join(tempfile.gettempdir(), "bool_img_test.png")
    img.save(path, "PNG")
    return path


# ── BUG 1: bool_symdiff ──


class TestBoolSymdiffBug:
    def test_symdiff_returns_nonempty_for_overlapping(self):
        from ppt_pro_max.renderer.boolean_shapes import bool_symdiff, poly_rect
        a = poly_rect(0, 0, 4, 4)
        b = poly_rect(2, 2, 4, 4)
        result = bool_symdiff(a, b)
        assert result is not None, "bool_symdiff returned None"
        assert not result.is_empty, "BUG: bool_symdiff(a,b) returns empty — likely a.symmetric_difference(a) instead of (b)"
        expected_area = a.area + b.area - 2 * a.intersection(b).area
        assert result.area == pytest.approx(expected_area, abs=0.1), \
            f"Expected area {expected_area}, got {result.area}"

    def test_symdiff_two_non_overlapping(self):
        from ppt_pro_max.renderer.boolean_shapes import bool_symdiff, poly_rect
        a = poly_rect(0, 0, 2, 2)
        b = poly_rect(5, 5, 2, 2)
        result = bool_symdiff(a, b)
        assert result is not None
        assert not result.is_empty, "Non-overlapping symdiff should not be empty"
        assert result.area == pytest.approx(8.0, abs=0.1)

    def test_symdiff_identical_returns_empty(self):
        from ppt_pro_max.renderer.boolean_shapes import bool_symdiff, poly_rect
        a = poly_rect(0, 0, 4, 4)
        b = poly_rect(0, 0, 4, 4)
        result = bool_symdiff(a, b)
        assert result is not None
        assert result.is_empty, "Identical shapes symdiff should be empty"

    def test_symdiff_none_inputs(self):
        from ppt_pro_max.renderer.boolean_shapes import bool_symdiff
        assert bool_symdiff(None, None) is None

    def test_symdiff_vs_manual(self):
        from ppt_pro_max.renderer.boolean_shapes import bool_symdiff, poly_rect
        a = poly_rect(0, 0, 6, 6)
        b = poly_rect(3, 3, 4, 4)
        result = bool_symdiff(a, b)
        manual = a.symmetric_difference(b)
        assert result.area == pytest.approx(manual.area, abs=0.1)


# ── BUG 2: bool_image ignores geometry ──


class TestBoolImageBug:
    def test_bool_image_uses_custGeom_not_oval(self, test_image):
        from ppt_pro_max.renderer.boolean_shapes import (
            bool_image, poly_rect, poly_circle, bool_subtract,
        )
        prs = _prs()
        s = _slide(prs)
        outer = poly_rect(0, 0, 6, 4)
        inner = poly_circle(3, 2, 1.5)
        geom = bool_subtract(outer, inner)
        result = bool_image(geom, s, 1, 1, 6, 4, test_image)
        assert result is not None, "bool_image returned None"
        el = result._element if hasattr(result, '_element') else result
        spPr = el.find(qn('p:spPr'))
        assert spPr is not None
        custGeom = spPr.find(qn('a:custGeom'))
        assert custGeom is not None, \
            "BUG: bool_image renders as preset shape (OVAL) instead of custGeom — geometry parameter is ignored"
        pathLst = custGeom.find(qn('a:pathLst'))
        assert pathLst is not None
        paths = pathLst.findall(qn('a:path'))
        assert len(paths) >= 2, \
            f"Rect-with-hole should have 2 paths (exterior + hole), got {len(paths)}"

    def test_bool_image_has_blipFill(self, test_image):
        from ppt_pro_max.renderer.boolean_shapes import (
            bool_image, poly_star,
        )
        prs = _prs()
        s = _slide(prs)
        geom = poly_star(3, 3, 2, points=5, inner_ratio=0.4)
        result = bool_image(geom, s, 1, 1, 6, 6, test_image)
        assert result is not None
        el = result._element if hasattr(result, '_element') else result
        spPr = el.find(qn('p:spPr'))
        blipFill = spPr.find(qn('a:blipFill'))
        assert blipFill is not None, \
            "BUG: bool_image should have a:blipFill for image-in-custom-shape"

    def test_bool_image_none_geometry(self, test_image):
        from ppt_pro_max.renderer.boolean_shapes import bool_image
        prs = _prs()
        s = _slide(prs)
        result = bool_image(None, s, 1, 1, 4, 4, test_image)
        assert result is None, "None geometry should return None"

    def test_bool_image_star_shape(self, test_image):
        from ppt_pro_max.renderer.boolean_shapes import (
            bool_image, poly_star,
        )
        prs = _prs()
        s = _slide(prs)
        geom = poly_star(3, 3, 2, points=5, inner_ratio=0.4)
        result = bool_image(geom, s, 1, 1, 6, 6, test_image)
        assert result is not None
        el = result._element if hasattr(result, '_element') else result
        spPr = el.find(qn('p:spPr'))
        custGeom = spPr.find(qn('a:custGeom'))
        assert custGeom is not None, "Star geometry should render as custGeom"

    def test_bool_image_with_border(self, test_image):
        from ppt_pro_max.renderer.boolean_shapes import (
            bool_image, poly_rect,
        )
        prs = _prs()
        s = _slide(prs)
        geom = poly_rect(0, 0, 4, 3)
        result = bool_image(geom, s, 1, 1, 4, 3, test_image, border_color='#FF5500')
        assert result is not None
        el = result._element if hasattr(result, '_element') else result
        spPr = el.find(qn('p:spPr'))
        ln = spPr.find(qn('a:ln'))
        assert ln is not None
        solidFill = ln.find(qn('a:solidFill'))
        assert solidFill is not None, "Border color should produce a:ln/a:solidFill"
