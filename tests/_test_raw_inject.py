"""Test: inject raw components WITHOUT any modification to see if extraction is the problem."""

import json
import os

from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
from ppt_pro_max.enterprise.component_renderer import ComponentRenderer

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)

lib = ComponentLibrary(find_db_path())
renderer = ComponentRenderer()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

categories = ["hierarchy", "process", "swot", "infographic", "timeline", "chart"]

for cat in categories:
    results = lib.search(type="group", category=cat, limit=1)
    if not results:
        print(f"  {cat}: SKIPPED")
        continue
    comp = results[0]
    xml_parts = lib.load_xml(comp["id"])

    # RAW injection - NO recolor, NO fill, just inject as-is
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    meta_bytes = xml_parts.get("_meta")
    orig_aspect = 1.0
    if meta_bytes:
        meta = json.loads(meta_bytes)
        obe = meta.get("orig_bounds_emu")
        if obe and len(obe) == 4 and obe[3] > 0:
            orig_aspect = obe[2] / obe[3]

    content_area = (0.5, 1.0, 12.3, 5.8)
    bounds = renderer.compute_component_bounds(content_area, orig_aspect, "contain")
    renderer._inject_group_to_slide(slide, xml_parts, bounds)
    print(f"  {cat}: id={comp['id']}, aspect={orig_aspect:.2f}, source={comp.get('source_file','?')}")

out = os.path.join(OUT_DIR, "raw_components_no_modify.pptx")
prs.save(out)
print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")
print(f"Slides: {len(prs.slides)}")

lib.close()
