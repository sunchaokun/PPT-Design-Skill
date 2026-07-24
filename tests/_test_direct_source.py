"""Extract grpSp XML DIRECTLY from source PPT (no scheme resolution, no modification) and inject."""

import os
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)

def_dir = r"E:\BaiduNetdiskDownload\DEF"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_count = 0

for root_dir, dirs, files in os.walk(def_dir):
    for fname in files:
        if not fname.endswith(".pptx"):
            continue
        if slide_count >= 8:
            break
        fpath = os.path.join(root_dir, fname)
        try:
            with zipfile.ZipFile(fpath) as z:
                for sf in [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml") and "rels" not in n]:
                    if slide_count >= 8:
                        break
                    slide_root = etree.fromstring(z.read(sf))
                    sp_tree = slide_root.find(f".//{{{_P}}}spTree")
                    if sp_tree is None:
                        continue
                    for grp in sp_tree.iter(f"{{{_P}}}grpSp"):
                        t_count = len([t for t in grp.iter(f"{{{_A}}}t") if t.text and t.text.strip()])
                        if t_count < 3:
                            continue

                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        target_sp_tree = slide._element.find(f".//{{{_P}}}spTree")

                        # Deep copy the grpSp element
                        grp_copy = etree.fromstring(etree.tostring(grp))

                        # Remap cNvPr ids only
                        next_id = 2
                        for sp in target_sp_tree:
                            for cNv in sp.iter(f"{{{_P}}}cNvPr"):
                                try:
                                    next_id = max(next_id, int(cNv.get("id", "1")) + 1)
                                except (ValueError, TypeError):
                                    pass
                        for cNv in grp_copy.iter(f"{{{_P}}}cNvPr"):
                            cNv.set("id", str(next_id))
                            cNv.set("name", f"Comp {next_id}")
                            next_id += 1

                        # Remove blip r:embed refs (can't resolve without rels)
                        for blip in grp_copy.iter(f"{{{_A}}}blip"):
                            embed = blip.get(f"{{{_R}}}embed")
                            if embed:
                                blip.attrib.pop(f"{{{_R}}}embed", None)

                        target_sp_tree.append(grp_copy)

                        texts = [t.text.strip()[:15] for t in grp_copy.iter(f"{{{_A}}}t") if t.text and t.text.strip()]
                        has_3d = len(list(grp_copy.iter(f"{{{_A}}}scene3d"))) > 0
                        has_scheme = len(list(grp_copy.iter(f"{{{_A}}}schemeClr"))) > 0
                        print(f"  Slide {slide_count+1}: {fname[:30]}, texts={t_count}, 3d={has_3d}, scheme={has_scheme}, sample={texts[:3]}")
                        slide_count += 1
                        if slide_count >= 8:
                            break
        except Exception as e:
            pass
    if slide_count >= 8:
        break

out = os.path.join(OUT_DIR, "direct_from_source.pptx")
prs.save(out)
print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")
print(f"Slides: {len(prs.slides)}")
