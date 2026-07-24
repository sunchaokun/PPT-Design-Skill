"""Build mode integration test — validates build_helpers API matches SKILL.md spec."""
from __future__ import annotations

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def build_env():
    from ppt_pro_max.build_helpers import (
        add_slide, page_header, kpi_card, bar_chart, comparison_bars,
        donut_chart, highlight_cards, text, multiline, rect, rrect, oval,
        top_bar, copy_decorations, copy_logo, _resolve_color,
    )
    C = {
        'primary': '#2E6504', 'accent': '#7DA92F', 'muted': '#81C784',
        'light': '#C8E6C9', 'white': '#FFFFFF', 'background': '#FFFFFF',
        'card_bg': '#F9F9F9', 'text_dark': '#1A1A1A', 'text_body': '#333333',
        'text_muted': '#666666', 'divider': '#CCCCCC',
        'font_heading': '微软雅黑', 'font_body': '微软雅黑',
    }
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return {
        'C': C, 'prs': prs,
        'add_slide': add_slide, 'page_header': page_header,
        'kpi_card': kpi_card, 'bar_chart': bar_chart,
        'comparison_bars': comparison_bars, 'donut_chart': donut_chart,
        'highlight_cards': highlight_cards, 'text': text,
        'multiline': multiline, 'rect': rect, 'rrect': rrect,
        'oval': oval, 'top_bar': top_bar,
        'copy_decorations': copy_decorations, 'copy_logo': copy_logo,
        '_resolve_color': _resolve_color,
    }


class TestColorResolution:
    def test_hex_direct(self, build_env):
        assert build_env['_resolve_color']('#2E6504', build_env['C']) == '#2E6504'

    def test_role_name(self, build_env):
        assert build_env['_resolve_color']('primary', build_env['C']) == '#2E6504'

    def test_missing_role_returns_black(self, build_env):
        assert build_env['_resolve_color']('nonexistent', build_env['C']) == '#000000'

    def test_none_returns_black(self, build_env):
        assert build_env['_resolve_color'](None, build_env['C']) == '#000000'


class TestAddSlide:
    def test_adds_slide(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        assert len(build_env['prs'].slides) == 1
        assert slide is not None

    def test_multiple_slides(self, build_env):
        for _ in range(5):
            build_env['add_slide'](build_env['prs'])
        assert len(build_env['prs'].slides) == 5


class TestPageHeader:
    def test_header_with_title_only(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        build_env['page_header'](slide, 'Revenue Overview', C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any('Revenue Overview' in t for t in texts)

    def test_header_with_subtitle(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        build_env['page_header'](slide, 'Title', 'Subtitle text', build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert 'Title' in texts
        assert 'Subtitle text' in texts

    def test_header_creates_divider(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        shape_count_before = len(slide.shapes)
        build_env['page_header'](slide, 'Title', C=build_env['C'])
        assert len(slide.shapes) > shape_count_before


class TestKpiCard:
    def test_kpi_card_renders(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        build_env['kpi_card'](slide, 1.0, 1.5, 3.0, 1.35, '12.8亿', '年度产值', '+8.3%', C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert '12.8亿' in texts
        assert '年度产值' in texts
        assert '+8.3%' in texts

    def test_kpi_card_without_trend(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        build_env['kpi_card'](slide, 1.0, 1.5, 3.0, 1.35, '99.9%', '可用率', C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert '99.9%' in texts

    def test_multiple_kpi_cards(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        metrics = [('12.8亿', '产值', '+8.3%'), ('99.9%', '可用率', '+0.1%'), ('5.2M', '用户', '+15%')]
        for i, (num, label, trend) in enumerate(metrics):
            build_env['kpi_card'](slide, 0.65 + i * 4.1, 1.5, 3.8, 1.35, num, label, trend, C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert '12.8亿' in texts
        assert '99.9%' in texts
        assert '5.2M' in texts


class TestBarChart:
    def test_bar_chart_renders(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        data = [('Revenue', 0.8, '$12.8B'), ('Growth', 0.6, '+8.3%'), ('Users', 0.45, '5.2M')]
        build_env['bar_chart'](slide, 1.5, 1.8, data, C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert 'Revenue' in texts
        assert '$12.8B' in texts


class TestComparisonBars:
    def test_comparison_bars_renders(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        metrics = [
            ('Revenue', '$10B', '$12.8B', 0.6, 0.8),
            ('Users', '3M', '5.2M', 0.3, 0.52),
        ]
        build_env['comparison_bars'](slide, 2.0, 1.8, metrics, C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert 'Revenue' in texts
        assert '$10B' in texts
        assert '$12.8B' in texts


class TestDonutChart:
    def test_donut_chart_renders(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        sectors = [('Product A', '40%', '#2E6504'), ('Product B', '35%', '#7DA92F'), ('Other', '25%', '#81C784')]
        build_env['donut_chart'](slide, 4.0, 3.5, 1.5, 0.8, sectors, C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert '100%' in texts
        all_text = ' '.join(texts)
        assert 'Product A' in all_text


class TestHighlightCards:
    def test_highlight_cards_renders(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        cards = [
            ('AI Engine', 'Auto-select optimal framework', '#2E6504'),
            ('Live Dashboard', 'Real-time monitoring', '#7DA92F'),
            ('Integration', 'Seamless API connection', '#81C784'),
        ]
        build_env['highlight_cards'](slide, 0.65, 1.8, cards, C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert 'AI Engine' in texts
        assert 'Live Dashboard' in texts
        assert 'Integration' in texts


class TestBasicShapes:
    def test_text_shape(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        build_env['text'](slide, 1.0, 1.0, 5.0, 0.5, 'Hello World', font_size=24, color='primary', bold=True, C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert 'Hello World' in texts

    def test_multiline_shape(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        lines = ['Line 1', 'Line 2', 'Line 3']
        build_env['multiline'](slide, 1.0, 1.0, 5.0, 1.5, lines, font_size=14, color='text_body', C=build_env['C'])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        combined = '\n'.join(texts)
        assert 'Line 1' in combined

    def test_rect_shape(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        shape = build_env['rect'](slide, 1.0, 1.0, 3.0, 2.0, 'primary', C=build_env['C'])
        assert shape is not None
        assert len(slide.shapes) == 1

    def test_rrect_shape(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        shape = build_env['rrect'](slide, 1.0, 1.0, 3.0, 2.0, 'card_bg', line='light', C=build_env['C'])
        assert shape is not None

    def test_oval_shape(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        shape = build_env['oval'](slide, 1.0, 1.0, 2.0, 2.0, 'accent', C=build_env['C'])
        assert shape is not None

    def test_top_bar(self, build_env):
        slide = build_env['add_slide'](build_env['prs'])
        shape = build_env['top_bar'](slide, 'primary', C=build_env['C'])
        assert shape is not None


class TestBuildScriptEndToEnd:
    def test_full_mckinsey_style_build(self, build_env):
        C = build_env['C']
        prs = build_env['prs']
        add_slide = build_env['add_slide']
        page_header = build_env['page_header']
        kpi_card = build_env['kpi_card']
        bar_chart = build_env['bar_chart']
        highlight_cards = build_env['highlight_cards']
        top_bar = build_env['top_bar']
        text = build_env['text']

        s0 = add_slide(prs)
        top_bar(s0, 'primary', C=C)
        text(s0, 1.2, 2.0, 10, 1.5, 'Strategic Growth Report', font_size=44, color='primary', bold=True, C=C)
        text(s0, 1.2, 3.5, 10, 0.5, 'FY2025 Performance Review', font_size=20, color='text_muted', C=C)

        s1 = add_slide(prs)
        top_bar(s1, 'primary', C=C)
        page_header(s1, 'Key Performance Indicators', 'Core business metrics', C)
        kpi_card(s1, 0.65, 1.8, 3.8, 1.35, '12.8亿', '年度产值', '+8.3%', C=C)
        kpi_card(s1, 4.75, 1.8, 3.8, 1.35, '99.9%', '系统可用率', '+0.1%', C=C)
        kpi_card(s1, 8.85, 1.8, 3.8, 1.35, '5.2M', '活跃用户', '+15%', C=C)

        s2 = add_slide(prs)
        top_bar(s2, 'primary', C=C)
        page_header(s2, 'Revenue Breakdown', 'By business segment', C)
        data = [('Enterprise', 0.8, '$10.2B'), ('SMB', 0.45, '$5.8B'), ('Consumer', 0.3, '$3.9B')]
        bar_chart(s2, 2.0, 2.0, data, C=C)

        s3 = add_slide(prs)
        top_bar(s3, 'primary', C=C)
        page_header(s3, 'Core Capabilities', 'What sets us apart', C)
        cards = [
            ('AI Engine', 'Auto-select optimal framework, 70% compression', C['primary']),
            ('Live Dashboard', 'Real-time monitoring and alerting', C['accent']),
            ('Integration', 'Seamless API, 200+ connectors', C['muted']),
        ]
        highlight_cards(s3, 0.65, 2.0, cards, C=C)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            path = f.name
        try:
            prs.save(path)
            assert os.path.getsize(path) > 0
            verify_prs = Presentation(path)
            assert len(verify_prs.slides) == 4
            for slide in verify_prs.slides:
                assert len(slide.shapes) >= 2
        finally:
            os.unlink(path)

    def test_full_cyberpunk_style_build(self, build_env):
        C = {
            'primary': '#0A0E27', 'accent': '#00FF88', 'muted': '#1A1F3A',
            'light': '#2A2F4A', 'white': '#E0E0E0', 'background': '#0A0E27',
            'card_bg': '#12162B', 'text_dark': '#E0E0E0', 'text_body': '#B0B0C0',
            'text_muted': '#606080', 'divider': '#2A2F4A',
            'font_heading': 'Orbitron', 'font_body': 'JetBrains Mono',
        }
        prs = build_env['prs']
        add_slide = build_env['add_slide']
        page_header = build_env['page_header']
        text = build_env['text']
        rect = build_env['rect']
        top_bar = build_env['top_bar']
        kpi_card = build_env['kpi_card']
        bar_chart = build_env['bar_chart']

        s0 = add_slide(prs)
        rect(s0, 0, 0, 13.333, 7.5, 'background', C=C)
        top_bar(s0, 'accent', C=C)
        text(s0, 1.2, 2.0, 10, 1.5, 'NEURAL NETWORK OPS', font_size=44, color='accent', bold=True, font_name='Orbitron', C=C)
        text(s0, 1.2, 3.8, 10, 0.5, 'Real-time AI Infrastructure Monitoring', font_size=18, color='text_muted', font_name='JetBrains Mono', C=C)

        s1 = add_slide(prs)
        rect(s1, 0, 0, 13.333, 7.5, 'background', C=C)
        top_bar(s1, 'accent', C=C)
        page_header(s1, 'System Metrics', 'Live dashboard data', C)
        kpi_card(s1, 0.65, 1.8, 3.8, 1.35, '99.97%', 'Uptime', '+0.02%', C=C)
        kpi_card(s1, 4.75, 1.8, 3.8, 1.35, '2.4ms', 'P99 Latency', '-12%', C=C)
        kpi_card(s1, 8.85, 1.8, 3.8, 1.35, '1.2M', 'RPS', '+23%', C=C)

        s2 = add_slide(prs)
        rect(s2, 0, 0, 13.333, 7.5, 'background', C=C)
        top_bar(s2, 'accent', C=C)
        page_header(s2, 'Traffic Analysis', 'Request distribution', C)
        data = [('API Gateway', 0.85, '1.02M'), ('ML Pipeline', 0.6, '720K'), ('Data Lake', 0.35, '420K')]
        bar_chart(s2, 2.0, 2.0, data, C=C)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            path = f.name
        try:
            prs.save(path)
            assert os.path.getsize(path) > 0
            verify_prs = Presentation(path)
            assert len(verify_prs.slides) == 3
        finally:
            os.unlink(path)

    def test_full_creative_style_build(self, build_env):
        C = {
            'primary': '#6C5CE7', 'accent': '#FD79A8', 'muted': '#A29BFE',
            'light': '#DFE6E9', 'white': '#FFFFFF', 'background': '#FFFFFF',
            'card_bg': '#F8F9FA', 'text_dark': '#2D3436', 'text_body': '#636E72',
            'text_muted': '#B2BEC3', 'divider': '#DFE6E9',
            'font_heading': 'Fredoka', 'font_body': 'Nunito',
        }
        prs = build_env['prs']
        add_slide = build_env['add_slide']
        page_header = build_env['page_header']
        highlight_cards = build_env['highlight_cards']
        oval = build_env['oval']
        text = build_env['text']
        top_bar = build_env['top_bar']

        s0 = add_slide(prs)
        top_bar(s0, 'primary', C=C)
        oval(s0, 9.0, 0.5, 4.0, 4.0, 'accent', C=C)
        text(s0, 1.2, 2.0, 8, 1.5, 'Creative Studio', font_size=48, color='primary', bold=True, font_name='Fredoka', C=C)
        text(s0, 1.2, 3.8, 8, 0.5, 'Design that inspires', font_size=22, color='text_muted', font_name='Nunito', C=C)

        s1 = add_slide(prs)
        top_bar(s1, 'primary', C=C)
        page_header(s1, 'Our Services', 'What we create', C)
        cards = [
            ('Brand Identity', 'Complete visual systems', C['primary']),
            ('Web Design', 'Responsive & beautiful', C['accent']),
            ('Motion Graphics', 'Animated storytelling', C['muted']),
        ]
        highlight_cards(s1, 0.65, 2.0, cards, C=C)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            path = f.name
        try:
            prs.save(path)
            assert os.path.getsize(path) > 0
            verify_prs = Presentation(path)
            assert len(verify_prs.slides) == 2
        finally:
            os.unlink(path)


class TestCopyFunctions:
    def test_copy_decorations_skips_images(self, build_env):
        template_prs = Presentation()
        template_prs.slide_width = Inches(13.333)
        template_prs.slide_height = Inches(7.5)
        template_slide = template_prs.slides.add_slide(template_prs.slide_layouts[6])
        from ppt_pro_max.build_helpers import text as bh_text, rect as bh_rect
        bh_text(template_slide, 0.5, 0.1, 2, 0.3, 'Short', font_size=10, color='text_muted', C=build_env['C'])
        bh_text(template_slide, 0.5, 0.5, 5, 0.3, 'A' * 60, font_size=10, color='text_muted', C=build_env['C'])
        bh_rect(template_slide, 0, 7.0, 13.333, 0.08, 'primary', C=build_env['C'])

        target_slide = build_env['add_slide'](build_env['prs'])
        build_env['copy_decorations'](target_slide, template_slide)
        copied_texts = [s.text_frame.text for s in target_slide.shapes if s.has_text_frame]
        assert 'Short' in copied_texts
        assert 'A' * 60 not in copied_texts

    def test_copy_logo_finds_group_shapes(self, build_env):
        template_prs = Presentation()
        template_prs.slide_width = Inches(13.333)
        template_prs.slide_height = Inches(7.5)
        template_slide = template_prs.slides.add_slide(template_prs.slide_layouts[6])

        group = template_slide.shapes.add_group_shape()
        group.left = Inches(11)
        group.top = Inches(0.3)
        group.width = Inches(1.5)
        group.height = Inches(0.8)

        target_slide = build_env['add_slide'](build_env['prs'])
        build_env['copy_logo'](target_slide, template_slide)
        group_shapes = [s for s in target_slide.shapes if s.shape_type == 6]
        assert len(group_shapes) >= 1


class TestBuildHelpersImport:
    def test_all_functions_exist(self):
        import ppt_pro_max.build_helpers as bh
        expected = [
            'add_slide', 'page_header', 'kpi_card', 'bar_chart',
            'comparison_bars', 'donut_chart', 'highlight_cards',
            'text', 'multiline', 'rect', 'rrect', 'oval',
            'top_bar', 'copy_decorations', 'copy_logo',
            '_resolve_color', '_rgb',
        ]
        for name in expected:
            assert hasattr(bh, name), f"Missing function: {name}"
            assert callable(getattr(bh, name)), f"Not callable: {name}"

    def test_star_import_at_module_level(self):
        from ppt_pro_max import build_helpers
        public_names = [n for n in dir(build_helpers) if not n.startswith('_')]
        assert len(public_names) >= 15
