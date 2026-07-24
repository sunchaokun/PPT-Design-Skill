"""Extract with FULL schemeClr resolution (including styles) and inject."""

import os
import zipfile
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_rel = "http://schemas.openxmlformats.org/package/2006/relationships"

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)
def_dir = r"E:\BaiduNetdiskDownload\DEF"


def parse_theme_from_zip(z, slide_rels, slide_path):
    try:
        for rId, target in slide_rels.items():
            if "slideLayout" in target:
                layout_path = os.path.normpath(os.path.join(os.path.dirname(slide_path), target)).replace("\\", "/")
                if layout_path not in z.namelist():
                    continue
                layout_dir = os.path.dirname(layout_path)
                layout_basename = os.path.basename(layout_path)
                layout_rels_path = os.path.join(layout_dir, "_rels", layout_basename + ".rels").replace("\\", "/")
                if layout_rels_path not in z.namelist():
                    continue
                lr_xml = z.read(layout_rels_path)
                lr_root = etree.fromstring(lr_xml)
                for lr in lr_root.findall(f"{{{_rel}}}Relationship"):
                    lr_target = lr.get("Target", "")
                    if "theme" in lr_target.lower():
                        abs_theme = os.path.normpath(os.path.join(os.path.dirname(layout_path), lr_target)).replace("\\", "/")
                        if abs_theme in z.namelist():
                            return parse_theme_xml(z.read(abs_theme))
                    if "slideMaster" in lr_target:
                        master_path = os.path.normpath(os.path.join(os.path.dirname(layout_path), lr_target)).replace("\\", "/")
                        if master_path not in z.namelist():
                            continue
                        master_dir = os.path.dirname(master_path)
                        master_basename = os.path.basename(master_path)
                        master_rels_path = os.path.join(master_dir, "_rels", master_basename + ".rels").replace("\\", "/")
                        if master_rels_path not in z.namelist():
                            continue
                        mr_xml = z.read(master_rels_path)
                        mr_root = etree.fromstring(mr_xml)
                        for mr in mr_root.findall(f"{{{_rel}}}Relationship"):
                            mr_target = mr.get("Target", "")
                            if "theme" in mr_target.lower():
                                abs_theme = os.path.normpath(os.path.join(os.path.dirname(master_path), mr_target)).replace("\\", "/")
                                if abs_theme in z.namelist():
                                    return parse_theme_xml(z.read(abs_theme))
                break
    except Exception:
        pass
    return {}


def parse_theme_xml(theme_blob):
    try:
        root = etree.fromstring(theme_blob)
        color_map = {}
        theme_elements = root.find(f".//{{{_A}}}themeElements")
        if theme_elements is None:
            return {}
        clr_scheme = theme_elements.find(f"{{{_A}}}clrScheme")
        if clr_scheme is None:
            return {}
        for child in clr_scheme:
            tag = child.tag.split("}")[-1]
            srgb = child.find(f"{{{_A}}}srgbClr")
            sys_clr = child.find(f"{{{_A}}}sysClr")
            if srgb is not None:
                color_map[tag] = srgb.get("val", "000000")
            elif sys_clr is not None:
                color_map[tag] = sys_clr.get("lastClr", "000000")
        if "lt1" in color_map and "bg1" not in color_map:
            color_map["bg1"] = color_map["lt1"]
        if "dk1" in color_map and "tx1" not in color_map:
            color_map["tx1"] = color_map["dk1"]
        if "lt2" in color_map and "bg2" not in color_map:
            color_map["bg2"] = color_map["lt2"]
        if "dk2" in color_map and "tx2" not in color_map:
            color_map["tx2"] = color_map["dk2"]
        return color_map
    except Exception:
        return {}


def resolve_all_scheme_colors(grp_elem, theme_colors):
    if not theme_colors:
        return grp_elem
    elem = copy.deepcopy(grp_elem)
    for scheme in list(elem.iter(f"{{{_A}}}schemeClr")):
        parent = scheme.getparent()
        if parent is None:
            continue
        val = scheme.get("val", "")
        if val in theme_colors:
            idx = list(parent).index(scheme)
            srgb = etree.SubElement(parent, f"{{{_A}}}srgbClr")
            srgb.set("val", theme_colors[val])
            for child in scheme:
                srgb.append(copy.deepcopy(child))
            parent.remove(scheme)
            parent.insert(idx, srgb)
    return elem


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_count = 0

for root_dir, dirs, files in os.walk(def_dir):
    for fname in files:
        if not fname.endswith(".pptx"):
            continue
        if slide_count >= 10:
            break
        fpath = os.path.join(root_dir, fname)
        try:
            with zipfile.ZipFile(fpath) as z:
                for sf in [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml") and "rels" not in n]:
                    if slide_count >= 10:
                        break
                    slide_root = etree.fromstring(z.read(sf))
                    sp_tree = slide_root.find(f".//{{{_P}}}spTree")
                    if sp_tree is None:
                        continue

                    rels_path = sf.replace("slides/", "slides/_rels/") + ".rels"
                    slide_rels = {}
                    if rels_path in z.namelist():
                        rels_root = etree.fromstring(z.read(rels_path))
                        for rel in rels_root.findall(f"{{{_rel}}}Relationship"):
                            slide_rels[rel.get("Id", "")] = rel.get("Target", "")

                    theme_colors = parse_theme_from_zip(z, slide_rels, sf)

                    for grp in sp_tree.iter(f"{{{_P}}}grpSp"):
                        t_count = len([t for t in grp.iter(f"{{{_A}}}t") if t.text and t.text.strip()])
                        if t_count < 3:
                            continue

                        # Resolve ALL schemeClr (including in styles)
                        resolved = resolve_all_scheme_colors(grp, theme_colors)

                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        target_sp_tree = slide._element.find(f".//{{{_P}}}spTree")

                        next_id = 2
                        for sp in target_sp_tree:
                            for cNv in sp.iter(f"{{{_P}}}cNvPr"):
                                try:
                                    next_id = max(next_id, int(cNv.get("id", "1")) + 1)
                                except (ValueError, TypeError):
                                    pass
                        for cNv in resolved.iter(f"{{{_P}}}cNvPr"):
                            cNv.set("id", str(next_id))
                            cNv.set("name", f"Comp {next_id}")
                            next_id += 1

                        for blip in resolved.iter(f"{{{_A}}}blip"):
                            embed = blip.get(f"{{{_R}}}embed")
                            if embed:
                                blip.attrib.pop(f"{{{_R}}}embed", None)

                        target_sp_tree.append(resolved)

                        has_3d = len(list(resolved.iter(f"{{{_A}}}scene3d"))) > 0
                        scheme_remaining = len(list(resolved.iter(f"{{{_A}}}schemeClr")))
                        texts = [t.text.strip()[:15] for t in resolved.iter(f"{{{_A}}}t") if t.text and t.text.strip()]
                        tag = "3D" if has_3d else "2D"
                        print(f"  Slide {slide_count+1} [{tag}]: {fname[:25]}, texts={t_count}, scheme_remaining={scheme_remaining}, theme_keys={len(theme_colors)}, sample={texts[:3]}")
                        slide_count += 1
                        break
        except Exception:
            pass
    if slide_count >= 10:
        break

out = os.path.join(OUT_DIR, "resolved_scheme_colors.pptx")
prs.save(out)
print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")
