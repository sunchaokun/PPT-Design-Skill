"""Tests for Boolean Shape Operations (Union / Intersect / Subtract / Combine / Fragment).

Validates that Shapely boolean geometry can be converted to OOXML custGeom
and written into a valid .pptx file that PowerPoint can open.
"""

from __future__ import annotations

import io
import math
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree

try:
    from shapely.geometry import Polygon, MultiPolygon
    from shapely.ops import unary_union
    HAS_SHAPELY = True
except ImportError:
    HAS_SHAPELY = False

pytestmark = pytest.mark.skipif(not HAS_SHAPELY, reason="shapely not installed")


EMU_PER_INCH = 914400


def _polygon_to_path_cmds(poly: Polygon, scale: float = 1.0) -> list[dict]:
    coords = list(poly.exterior.coords)
    cmds = []
    cmds.append({"cmd": "moveTo", "x": coords[0][0] * scale, "y": coords[0][1] * scale})
    for x, y in coords[1:]:
        cmds.append({"cmd": "lnTo", "x": x * scale, "y": y * scale})
    cmds.append({"cmd": "close"})
    return cmds


def _polygon_to_paths_with_holes(poly: Polygon, scale: float = 1.0) -> list[list[dict]]:
    paths = []
    coords = list(poly.exterior.coords)
    cmds = [{"cmd": "moveTo", "x": coords[0][0] * scale, "y": coords[0][1] * scale}]
    for x, y in coords[1:]:
        cmds.append({"cmd": "lnTo", "x": x * scale, "y": y * scale})
    cmds.append({"cmd": "close"})
    paths.append(cmds)
    for interior in poly.interiors:
        hole_coords = list(interior.coords)
        hole_cmds = [{"cmd": "moveTo", "x": hole_coords[0][0] * scale, "y": hole_coords[0][1] * scale}]
        for x, y in hole_coords[1:]:
            hole_cmds.append({"cmd": "lnTo", "x": x * scale, "y": y * scale})
        hole_cmds.append({"cmd": "close"})
        paths.append(hole_cmds)
    return paths


def _multipolygon_to_paths(mpoly, scale: float = 1.0) -> list[list[dict]]:
    paths = []
    if isinstance(mpoly, Polygon):
        paths.extend(_polygon_to_paths_with_holes(mpoly, scale))
    elif isinstance(mpoly, MultiPolygon):
        for poly in mpoly.geoms:
            paths.extend(_polygon_to_paths_with_holes(poly, scale))
    return paths


def _build_custGeom_shape(slide, paths, x_in, y_in, w_in, h_in,
                          fill_color="#4472C4", line_color=None):
    sp_tree = slide.shapes._spTree
    sp = etree.SubElement(sp_tree, qn("p:sp"))

    nvSpPr = etree.SubElement(sp, qn("p:nvSpPr"))
    cNvPr = etree.SubElement(nvSpPr, qn("p:cNvPr"))
    max_id = 1
    for sh in slide.shapes:
        try:
            if sh.shape_id > max_id:
                max_id = sh.shape_id
        except Exception:
            pass
    cNvPr.set("id", str(max_id + 1))
    cNvPr.set("name", "BooleanShape")
    etree.SubElement(nvSpPr, qn("p:cNvSpPr"))
    etree.SubElement(nvSpPr, qn("p:nvPr"))

    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(int(x_in * EMU_PER_INCH)))
    off.set("y", str(int(y_in * EMU_PER_INCH)))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(int(w_in * EMU_PER_INCH)))
    ext.set("cy", str(int(h_in * EMU_PER_INCH)))

    custGeom = etree.SubElement(spPr, qn("a:custGeom"))
    etree.SubElement(custGeom, qn("a:avLst"))
    etree.SubElement(custGeom, qn("a:gdLst"))
    pathLst = etree.SubElement(custGeom, qn("a:pathLst"))

    path_w = int(w_in * EMU_PER_INCH)
    path_h = int(h_in * EMU_PER_INCH)
    off_x_emu = int(x_in * EMU_PER_INCH)
    off_y_emu = int(y_in * EMU_PER_INCH)

    for path_cmds in paths:
        path_el = etree.SubElement(pathLst, qn("a:path"))
        path_el.set("w", str(path_w))
        path_el.set("h", str(path_h))
        for cmd in path_cmds:
            if cmd["cmd"] == "moveTo":
                moveTo = etree.SubElement(path_el, qn("a:moveTo"))
                pt = etree.SubElement(moveTo, qn("a:pt"))
                pt.set("x", str(int(cmd["x"] * EMU_PER_INCH) - off_x_emu))
                pt.set("y", str(int(cmd["y"] * EMU_PER_INCH) - off_y_emu))
            elif cmd["cmd"] == "lnTo":
                lnTo = etree.SubElement(path_el, qn("a:lnTo"))
                pt = etree.SubElement(lnTo, qn("a:pt"))
                pt.set("x", str(int(cmd["x"] * EMU_PER_INCH) - off_x_emu))
                pt.set("y", str(int(cmd["y"] * EMU_PER_INCH) - off_y_emu))
            elif cmd["cmd"] == "close":
                etree.SubElement(path_el, qn("a:close"))

    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr")).set("val", fill_color.lstrip("#"))

    ln = etree.SubElement(spPr, qn("a:ln"))
    if line_color:
        sf = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr")).set("val", line_color.lstrip("#"))
    else:
        etree.SubElement(ln, qn("a:noFill"))

    return sp


def _boolean_to_slide(slide, result_geom, x_in, y_in, w_in, h_in,
                      fill_color="#4472C4", line_color=None):
    paths = _multipolygon_to_paths(result_geom, scale=1.0)
    return _build_custGeom_shape(slide, paths, x_in, y_in, w_in, h_in,
                                 fill_color=fill_color, line_color=line_color)


# ── Shapely Boolean Computation Tests ──


class TestShapelyBooleanComputation:
    def test_union_two_rectangles(self):
        r1 = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        r2 = Polygon([(2, 2), (6, 2), (6, 6), (2, 6)])
        union = r1.union(r2)
        assert union.is_valid
        assert union.area == pytest.approx(28.0, abs=0.1)

    def test_intersect_two_rectangles(self):
        r1 = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        r2 = Polygon([(2, 2), (6, 2), (6, 6), (2, 6)])
        inter = r1.intersection(r2)
        assert inter.is_valid
        assert inter.area == pytest.approx(4.0, abs=0.1)

    def test_subtract_two_rectangles(self):
        r1 = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        r2 = Polygon([(2, 2), (6, 2), (6, 6), (2, 6)])
        diff = r1.difference(r2)
        assert diff.is_valid
        assert diff.area == pytest.approx(12.0, abs=0.1)

    def test_symmetric_difference(self):
        r1 = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        r2 = Polygon([(2, 2), (6, 2), (6, 6), (2, 6)])
        sdiff = r1.symmetric_difference(r2)
        assert sdiff.is_valid
        assert sdiff.area == pytest.approx(24.0, abs=0.1)

    def test_union_circle_and_rectangle(self):
        import math
        circle = Polygon([(math.cos(a) * 2 + 2, math.sin(a) * 2 + 2)
                          for a in [i * math.pi / 36 for i in range(72)]])
        rect = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        union = circle.union(rect)
        assert union.is_valid
        assert union.area >= rect.area

    def test_subtract_circle_from_rectangle(self):
        import math
        circle = Polygon([(math.cos(a) * 1.5 + 2, math.sin(a) * 1.5 + 2)
                          for a in [i * math.pi / 36 for i in range(72)]])
        rect = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        diff = rect.difference(circle)
        assert diff.is_valid
        assert diff.area < rect.area

    def test_non_overlapping_union(self):
        r1 = Polygon([(0, 0), (2, 0), (2, 2), (0, 2)])
        r2 = Polygon([(5, 5), (7, 5), (7, 7), (5, 7)])
        union = r1.union(r2)
        assert union.is_valid
        assert union.area == pytest.approx(8.0, abs=0.1)
        assert isinstance(union, MultiPolygon)

    def test_subtract_creates_hole(self):
        outer = Polygon([(0, 0), (6, 0), (6, 6), (0, 6)])
        inner = Polygon([(2, 2), (4, 2), (4, 4), (2, 4)])
        diff = outer.difference(inner)
        assert diff.is_valid
        assert diff.area == pytest.approx(32.0, abs=0.1)


# ── OOXML custGeom Generation Tests ──


class TestCustGeomGeneration:
    def test_polygon_to_path_cmds(self):
        poly = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        cmds = _polygon_to_path_cmds(poly)
        assert cmds[0]["cmd"] == "moveTo"
        assert cmds[0]["x"] == 0.0
        assert cmds[0]["y"] == 0.0
        assert len([c for c in cmds if c["cmd"] == "lnTo"]) == 4
        assert cmds[-1]["cmd"] == "close"

    def test_multipolygon_to_paths(self):
        r1 = Polygon([(0, 0), (2, 0), (2, 2), (0, 2)])
        r2 = Polygon([(5, 5), (7, 5), (7, 7), (5, 7)])
        mpoly = MultiPolygon([r1, r2])
        paths = _multipolygon_to_paths(mpoly)
        assert len(paths) == 2

    def test_build_custGeom_shape_creates_xml(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        poly = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        paths = _multipolygon_to_paths(poly)
        elem = _build_custGeom_shape(slide, paths, 1, 1, 5, 5)
        assert elem is not None
        spPr = elem.find(qn("p:spPr"))
        assert spPr is not None
        custGeom = spPr.find(qn("a:custGeom"))
        assert custGeom is not None
        pathLst = custGeom.find(qn("a:pathLst"))
        assert pathLst is not None
        paths_xml = pathLst.findall(qn("a:path"))
        assert len(paths_xml) == 1

    def test_build_custGeom_multipolygon(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        r1 = Polygon([(0, 0), (2, 0), (2, 2), (0, 2)])
        r2 = Polygon([(5, 5), (7, 5), (7, 7), (5, 7)])
        mpoly = MultiPolygon([r1, r2])
        paths = _multipolygon_to_paths(mpoly)
        elem = _build_custGeom_shape(slide, paths, 0, 0, 8, 8)
        pathLst = elem.find(qn("p:spPr")).find(qn("a:custGeom")).find(qn("a:pathLst"))
        assert len(pathLst.findall(qn("a:path"))) == 2


# ── End-to-End: Boolean → custGeom → PPTX ──


class TestBooleanToPPTX:
    def _save_and_reload(self, prs):
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return Presentation(buf)

    def test_union_shape_in_pptx(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        r1 = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        r2 = Polygon([(2, 2), (6, 2), (6, 6), (2, 6)])
        union = r1.union(r2)
        _boolean_to_slide(slide, union, 1, 1, 7, 7, fill_color="#1D78FA")
        prs2 = self._save_and_reload(prs)
        slide2 = prs2.slides[0]
        shapes = list(slide2.shapes)
        freeform_found = False
        for sh in shapes:
            if hasattr(sh, 'name') and 'Freeform' in sh.name:
                freeform_found = True
        assert freeform_found or len(shapes) >= 1

    def test_subtract_shape_in_pptx(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        outer = Polygon([(0, 0), (6, 0), (6, 6), (0, 6)])
        inner = Polygon([(2, 2), (4, 2), (4, 4), (2, 4)])
        diff = outer.difference(inner)
        _boolean_to_slide(slide, diff, 1, 1, 7, 7, fill_color="#FF5500")
        prs2 = self._save_and_reload(prs)
        slide2 = prs2.slides[0]
        assert len(list(slide2.shapes)) >= 1

    def test_intersect_shape_in_pptx(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        r1 = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        r2 = Polygon([(2, 2), (6, 2), (6, 6), (2, 6)])
        inter = r1.intersection(r2)
        _boolean_to_slide(slide, inter, 1, 1, 7, 7, fill_color="#00B050")
        prs2 = self._save_and_reload(prs)
        slide2 = prs2.slides[0]
        assert len(list(slide2.shapes)) >= 1

    def test_symmetric_difference_in_pptx(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        r1 = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        r2 = Polygon([(2, 2), (6, 2), (6, 6), (2, 6)])
        sdiff = r1.symmetric_difference(r2)
        _boolean_to_slide(slide, sdiff, 1, 1, 7, 7, fill_color="#7030A0")
        prs2 = self._save_and_reload(prs)
        slide2 = prs2.slides[0]
        assert len(list(slide2.shapes)) >= 1

    def test_circle_subtract_from_rect_in_pptx(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        import math
        circle = Polygon([(math.cos(a) * 1.5 + 2, math.sin(a) * 1.5 + 2)
                          for a in [i * math.pi / 36 for i in range(72)]])
        rect = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        diff = rect.difference(circle)
        _boolean_to_slide(slide, diff, 1, 1, 5, 5, fill_color="#C00000")
        prs2 = self._save_and_reload(prs)
        slide2 = prs2.slides[0]
        assert len(list(slide2.shapes)) >= 1

    def test_non_overlapping_union_in_pptx(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        r1 = Polygon([(0, 0), (2, 0), (2, 2), (0, 2)])
        r2 = Polygon([(5, 5), (7, 5), (7, 7), (5, 7)])
        union = r1.union(r2)
        _boolean_to_slide(slide, union, 0.5, 0.5, 8, 8, fill_color="#0070C0")
        prs2 = self._save_and_reload(prs)
        slide2 = prs2.slides[0]
        assert len(list(slide2.shapes)) >= 1

    def test_multiple_boolean_ops_on_one_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        r1 = Polygon([(0, 0), (3, 0), (3, 3), (0, 3)])
        r2 = Polygon([(1.5, 1.5), (4.5, 1.5), (4.5, 4.5), (1.5, 4.5)])

        _boolean_to_slide(slide, r1.union(r2), 0.5, 0.5, 6, 6, fill_color="#1D78FA")
        _boolean_to_slide(slide, r1.intersection(r2), 0.5, 3, 6, 6, fill_color="#00B050")
        _boolean_to_slide(slide, r1.difference(r2), 5, 0.5, 6, 6, fill_color="#FF5500")
        _boolean_to_slide(slide, r1.symmetric_difference(r2), 5, 3, 6, 6, fill_color="#7030A0")

        prs2 = self._save_and_reload(prs)
        slide2 = prs2.slides[0]
        assert len(list(slide2.shapes)) >= 4

    def test_save_to_file(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        r1 = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        r2 = Polygon([(2, 2), (6, 2), (6, 6), (2, 6)])
        union = r1.union(r2)
        _boolean_to_slide(slide, union, 1, 1, 7, 7, fill_color="#1D78FA")
        out = tmp_path / "boolean_test.pptx"
        prs.save(str(out))
        assert out.exists()
        assert out.stat().st_size > 0
        prs2 = Presentation(str(out))
        assert len(prs2.slides) == 1


# ── Practical Use Case Tests ──


class TestPracticalBooleanUseCases:
    def test_donut_shape(self):
        outer = Polygon([(0, 0), (6, 0), (6, 6), (0, 6)])
        inner = Polygon([(2, 2), (4, 2), (4, 4), (2, 4)])
        donut = outer.difference(inner)
        assert donut.is_valid
        assert donut.area == pytest.approx(32.0, abs=0.1)
        assert donut.interiors is not None and len(donut.interiors) == 1

    def test_donut_shape_in_pptx(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        outer = Polygon([(0, 0), (6, 0), (6, 6), (0, 6)])
        inner = Polygon([(2, 2), (4, 2), (4, 4), (2, 4)])
        donut = outer.difference(inner)
        _boolean_to_slide(slide, donut, 1, 1, 7, 7, fill_color="#FFC000")
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        assert len(list(prs2.slides[0].shapes)) >= 1

    def test_rounded_hole_in_rect(self):
        import math
        rect = Polygon([(0, 0), (6, 0), (6, 4), (0, 4)])
        circle = Polygon([(math.cos(a) * 1 + 3, math.sin(a) * 1 + 2)
                          for a in [i * math.pi / 36 for i in range(72)]])
        diff = rect.difference(circle)
        assert diff.is_valid
        assert diff.area < rect.area

    def test_cross_shape_union(self):
        h_bar = Polygon([(0, 2), (6, 2), (6, 4), (0, 4)])
        v_bar = Polygon([(2, 0), (4, 0), (4, 6), (2, 6)])
        cross = h_bar.union(v_bar)
        assert cross.is_valid
        assert cross.area == pytest.approx(20.0, abs=0.1)

    def test_l_shape_subtract(self):
        big = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        cut = Polygon([(2, 2), (4, 2), (4, 4), (2, 4)])
        l_shape = big.difference(cut)
        assert l_shape.is_valid
        assert l_shape.area == pytest.approx(12.0, abs=0.1)

    def test_star_shape_intersect(self):
        import math
        circle = Polygon([(math.cos(a) * 3 + 3, math.sin(a) * 3 + 3)
                          for a in [i * math.pi / 36 for i in range(72)]])
        diamond = Polygon([(3, 0), (6, 3), (3, 6), (0, 3)])
        inter = circle.intersection(diamond)
        assert inter.is_valid
        assert inter.area > 0

    def test_frame_shape(self):
        outer = Polygon([(0, 0), (8, 0), (8, 6), (0, 6)])
        inner = Polygon([(1, 1), (7, 1), (7, 5), (1, 5)])
        frame = outer.difference(inner)
        assert frame.is_valid
        assert frame.area == pytest.approx(24.0, abs=0.1)


class TestHoleRendering:
    def test_polygon_with_hole_produces_two_paths(self):
        outer = Polygon([(0, 0), (6, 0), (6, 6), (0, 6)])
        inner = Polygon([(2, 2), (4, 2), (4, 4), (2, 4)])
        donut = outer.difference(inner)
        paths = _multipolygon_to_paths(donut)
        assert len(paths) == 2
        assert paths[0][-1]["cmd"] == "close"
        assert paths[1][-1]["cmd"] == "close"

    def test_polygon_with_hole_in_pptx_has_two_path_elements(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        outer = Polygon([(0, 0), (6, 0), (6, 6), (0, 6)])
        inner = Polygon([(2, 2), (4, 2), (4, 4), (2, 4)])
        donut = outer.difference(inner)
        elem = _boolean_to_slide(slide, donut, 1, 1, 7, 7, fill_color="#FFC000")
        assert elem is not None
        pathLst = elem.find(qn("p:spPr")).find(qn("a:custGeom")).find(qn("a:pathLst"))
        path_elements = pathLst.findall(qn("a:path"))
        assert len(path_elements) == 2

    def test_circle_hole_in_rect_produces_two_paths(self):
        import math
        rect = Polygon([(0, 0), (6, 0), (6, 4), (0, 4)])
        circle = Polygon([(math.cos(a) * 1 + 3, math.sin(a) * 1 + 2)
                          for a in [i * math.pi / 36 for i in range(72)]])
        diff = rect.difference(circle)
        paths = _multipolygon_to_paths(diff)
        assert len(paths) == 2
        assert paths[0][-1]["cmd"] == "close"
        assert paths[1][-1]["cmd"] == "close"

    def test_no_hole_produces_one_path(self):
        poly = Polygon([(0, 0), (4, 0), (4, 4), (0, 4)])
        paths = _multipolygon_to_paths(poly)
        assert len(paths) == 1

    def test_multiple_holes_produce_multiple_paths(self):
        outer = Polygon([(0, 0), (10, 0), (10, 10), (0, 10)])
        h1 = Polygon([(1, 1), (3, 1), (3, 3), (1, 3)])
        h2 = Polygon([(5, 5), (7, 5), (7, 7), (5, 7)])
        diff = outer.difference(h1).difference(h2)
        paths = _multipolygon_to_paths(diff)
        assert len(paths) == 3
