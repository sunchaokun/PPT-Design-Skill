"""3 genuinely different PPTs using raw python-pptx. No component system."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _rgb(hex_str):
    return RGBColor.from_string(hex_str.lstrip('#'))


def _add_text(slide, left, top, width, height, txt, size=12, color='#000000',
              bold=False, font='Calibri', align='left', anchor='t'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor == 'ctr':
        tf._txBody.bodyPr.set('anchor', 'ctr')
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = _rgb(color)
    p.font.bold = bold
    p.font.name = font
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    return txBox


def _add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_rrect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_oval(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()
    return shape


def _add_multiline(slide, left, top, width, height, lines, size=12, color='#000000',
                   bold=False, font='Calibri', align='left', spacing=4):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = _rgb(color)
        p.font.bold = bold
        p.font.name = font
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing)
    return txBox


# =============================================================================
# McKinsey: Clean serif, data-dense, left-aligned grid, muted green palette
# =============================================================================
def build_mckinsey(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    PRI = '#2E6504'
    ACC = '#7DA92F'
    MUT = '#81C784'
    LIG = '#C8E6C9'
    BG  = '#FFFFFF'
    DK  = '#1A1A1A'
    BDY = '#333333'
    MUTED = '#666666'
    DIV = '#CCCCCC'
    CARD = '#F5F5F5'

    # Fonts
    FH = 'Georgia'
    FB = 'Calibri'
    FC = 'Consolas'

    # --- Slide 0: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 7.5, PRI)
    _add_rect(s, 0, 7.2, 13.333, 0.3, ACC)
    _add_text(s, 1.5, 1.8, 10, 1.2, 'Strategic Growth Report', size=48, color='#FFFFFF', bold=True, font=FH)
    _add_text(s, 1.5, 3.2, 10, 0.6, 'FY2025 Performance Review', size=22, color=LIG, font=FB)
    _add_rect(s, 1.5, 4.0, 2.5, 0.04, ACC)
    _add_text(s, 1.5, 4.3, 10, 0.4, 'NeuralForge Inc.  |  Confidential', size=13, color=MUTED, font=FB)

    # --- Slide 1: Section divider ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 7.5, PRI)
    _add_text(s, 1.5, 1.5, 3, 2.5, '01', size=96, color=LIG, bold=True, font=FH)
    _add_rect(s, 1.5, 4.2, 2.0, 0.04, ACC)
    _add_text(s, 1.5, 4.6, 10, 1.2, 'Core Metrics', size=44, color='#FFFFFF', bold=True, font=FH)

    # --- Slide 2: KPI dashboard ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.08, PRI)
    _add_text(s, 0.65, 0.4, 12, 0.5, 'Key Performance Indicators', size=28, color=DK, bold=True, font=FH)
    _add_text(s, 0.65, 0.9, 12, 0.3, 'Core business metrics', size=12, color=MUTED, font=FB)
    _add_rect(s, 0.65, 1.25, 12, 0.004, DIV)

    # 6 KPI cards in 3x2 grid
    kpis = [
        ('12.8B', 'Annual Revenue', '+8.3%', True),
        ('99.9%', 'System Uptime', '+0.1%', True),
        ('5.2M', 'Active Users', '+15%', True),
        ('2.4ms', 'P99 Latency', '-12%', False),
        ('72', 'NPS Score', 'Top 5%', True),
        ('38%', 'Women in Leadership', 'vs 21% avg', True),
    ]
    for idx, (num, lbl, trend, up) in enumerate(kpis):
        col = idx % 3
        row = idx // 3
        x = 0.65 + col * 4.1
        y = 1.6 + row * 2.0

        _add_rrect(s, x, y, 3.8, 1.7, CARD, line=DIV)
        _add_rect(s, x, y, 3.8, 0.06, ACC)
        _add_text(s, x + 0.25, y + 0.2, 3.3, 0.6, num, size=36, color=PRI, bold=True, font=FH)
        _add_text(s, x + 0.25, y + 0.8, 3.3, 0.3, lbl, size=13, color=MUTED, font=FB)
        tc = PRI if up else '#C53030'
        _add_text(s, x + 0.25, y + 1.2, 3.3, 0.3, trend, size=11, color=tc, bold=True, font=FB)

    # --- Slide 3: Revenue bar chart ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.08, PRI)
    _add_text(s, 0.65, 0.4, 12, 0.5, 'Revenue Breakdown', size=28, color=DK, bold=True, font=FH)
    _add_text(s, 0.65, 0.9, 12, 0.3, 'By business segment', size=12, color=MUTED, font=FB)
    _add_rect(s, 0.65, 1.25, 12, 0.004, DIV)

    bars = [('Enterprise SaaS', 0.82, '$10.2B'), ('SMB Platform', 0.48, '$5.8B'), ('Consumer App', 0.30, '$3.9B')]
    bar_colors = [PRI, ACC, MUT]
    for i, (lbl, pct, val) in enumerate(bars):
        y = 2.0 + i * 0.7
        _add_rrect(s, 2.5, y, 5.5, 0.4, '#F0F0F0')
        _add_rrect(s, 2.5, y, 5.5 * pct, 0.4, bar_colors[i])
        _add_text(s, 1.0, y + 0.05, 1.4, 0.3, lbl, size=13, color=BDY, font=FB, align='right')
        _add_text(s, 8.2, y + 0.05, 1.0, 0.3, val, size=13, color=DK, bold=True, font=FB)

    # Donut
    _add_oval(s, 9.5, 2.0, 3.6, 3.6, PRI)
    _add_oval(s, 10.3, 2.8, 2.0, 2.0, BG)
    _add_text(s, 10.3, 3.5, 2.0, 0.5, '100%', size=20, color=PRI, bold=True, font=FH, align='center')
    legend = [('SaaS 52%', PRI), ('SMB 29%', ACC), ('Consumer 19%', MUT)]
    for i, (lt, lc) in enumerate(legend):
        ly = 2.2 + i * 0.4
        _add_rrect(s, 9.6, ly, 0.2, 0.2, lc)
        _add_text(s, 9.9, ly - 0.02, 1.5, 0.25, lt, size=11, color=BDY, font=FB)

    # --- Slide 4: Feature cards ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.08, PRI)
    _add_text(s, 0.65, 0.4, 12, 0.5, 'Core Capabilities', size=28, color=DK, bold=True, font=FH)
    _add_text(s, 0.65, 0.9, 12, 0.3, 'What sets us apart', size=12, color=MUTED, font=FB)
    _add_rect(s, 0.65, 1.25, 12, 0.004, DIV)

    cards = [
        ('AI Inference Engine', 'Auto-select optimal framework\n70% compression, P99<10ms', PRI),
        ('Real-time Monitoring', 'Millisecond alerting\nFull-stack tracing, anomaly detection', ACC),
        ('Zero-config Integration', '200+ connectors\nOut-of-box ready, 5-min setup', MUT),
    ]
    for i, (title, desc, accent) in enumerate(cards):
        x = 0.65 + i * 4.1
        _add_rrect(s, x, 1.8, 3.8, 2.5, CARD, line=DIV)
        _add_rect(s, x, 1.8, 3.8, 0.06, accent)
        _add_text(s, x + 0.3, 2.1, 3.2, 0.4, title, size=18, color=DK, bold=True, font=FH)
        _add_text(s, x + 0.3, 2.6, 3.2, 1.5, desc, size=13, color=MUTED, font=FB)

    # --- Slide 5: Code block ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.08, PRI)
    _add_text(s, 0.65, 0.4, 12, 0.5, 'Deployment Example', size=28, color=DK, bold=True, font=FH)
    _add_text(s, 0.65, 0.9, 12, 0.3, 'Production-ready in 3 lines', size=12, color=MUTED, font=FB)
    _add_rect(s, 0.65, 1.25, 12, 0.004, DIV)

    _add_rect(s, 0.65, 1.8, 12, 4.0, '#1E1E1E')
    _add_rrect(s, 0.65, 1.55, 1.2, 0.25, ACC)
    _add_text(s, 0.75, 1.57, 1.0, 0.2, 'Python', size=10, color='#FFFFFF', bold=True, font=FC)
    _add_multiline(s, 1.0, 2.2, 11, 3.0, [
        'from neural_ops import deploy',
        'model = load("gpt-4o-finetuned")',
        'deploy(model, replicas=3, gpu="A100")',
        'print("Serving at", model.endpoint)',
    ], size=14, color='#D4D4D4', font=FC, spacing=6)

    # --- Slide 6: CTA ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 7.5, PRI)
    _add_rect(s, 0, 7.2, 13.333, 0.3, ACC)
    _add_text(s, 1.5, 2.5, 10, 1.2, "Let's Build the Future Together", size=40, color='#FFFFFF', bold=True, font=FH)
    _add_text(s, 1.5, 4.0, 10, 0.5, 'Contact: growth@example.com  |  example.com/demo', size=16, color=LIG, font=FB)

    prs.save(output_path)


# =============================================================================
# Cyberpunk: Dark bg, neon accents, monospace, terminal aesthetic
# =============================================================================
def build_cyberpunk(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    PRI = '#0A0E27'
    ACC = '#00F0FF'
    MUT = '#7B2FFF'
    PINK = '#FF2D6B'
    BG2 = '#12162B'
    TXT = '#E0E0E0'
    DIM = '#6A6A8A'
    FH = 'Orbitron'
    FB = 'JetBrains Mono'
    FC = 'Consolas'

    # --- Slide 0: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 7.5, PRI)
    _add_rect(s, 0, 3.5, 13.333, 0.04, ACC)
    _add_rect(s, 0, 3.6, 13.333, 0.02, MUT)
    _add_text(s, 1.0, 1.2, 11, 1.5, 'NEURAL OPS', size=56, color=ACC, bold=True, font=FH)
    _add_text(s, 1.0, 2.6, 11, 0.5, '// v4.2  |  AI Infrastructure at Scale', size=16, color=DIM, font=FB)
    _add_text(s, 1.0, 4.2, 11, 0.5, '2025 STATUS REPORT', size=20, color=TXT, font=FH)
    _add_rect(s, 1.0, 4.9, 1.5, 0.04, ACC)
    _add_text(s, 1.0, 5.2, 11, 0.4, 'CLASSIFIED  |  INTERNAL USE ONLY', size=10, color=PINK, font=FB)

    # --- Slide 1: KPI grid ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.06, ACC)
    _add_text(s, 0.8, 0.35, 12, 0.5, 'SYSTEM METRICS', size=24, color=ACC, bold=True, font=FH)
    _add_text(s, 0.8, 0.85, 12, 0.3, 'Real-time dashboard snapshot', size=11, color=DIM, font=FB)
    _add_rect(s, 0.8, 1.2, 11.7, 0.004, '#2A2F4A')

    kpis = [
        ('12.8B', 'INFERENCE/SEC', '+8.3%', True),
        ('99.9%', 'UPTIME SLA', '+0.1%', True),
        ('5.2M', 'ACTIVE NODES', '+15%', True),
        ('2.4ms', 'P99 LATENCY', '-12%', False),
        ('72', 'THREAT SCORE', 'LOW', True),
        ('38%', 'GPU UTIL', 'OPTIMAL', True),
    ]
    for idx, (num, lbl, trend, up) in enumerate(kpis):
        col = idx % 3
        row = idx // 3
        x = 0.8 + col * 4.2
        y = 1.6 + row * 2.2

        _add_rrect(s, x, y, 3.8, 1.8, BG2, line='#2A2F4A')
        _add_rect(s, x, y, 3.8, 0.04, ACC if up else PINK)
        _add_text(s, x + 0.3, y + 0.2, 3.2, 0.6, num, size=32, color=ACC, bold=True, font=FH)
        _add_text(s, x + 0.3, y + 0.85, 3.2, 0.3, lbl, size=10, color=DIM, font=FB)
        tc = ACC if up else PINK
        _add_text(s, x + 0.3, y + 1.3, 3.2, 0.3, trend, size=10, color=tc, bold=True, font=FB)

    # --- Slide 2: Traffic bars ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.06, ACC)
    _add_text(s, 0.8, 0.35, 12, 0.5, 'TRAFFIC DISTRIBUTION', size=24, color=ACC, bold=True, font=FH)
    _add_text(s, 0.8, 0.85, 12, 0.3, 'By service endpoint', size=11, color=DIM, font=FB)
    _add_rect(s, 0.8, 1.2, 11.7, 0.004, '#2A2F4A')

    bars = [('/api/infer', 0.82, '10.2B'), ('/api/embed', 0.48, '5.8B'), ('/api/train', 0.30, '3.9B')]
    bar_c = [ACC, MUT, PINK]
    for i, (lbl, pct, val) in enumerate(bars):
        y = 2.0 + i * 0.8
        _add_rrect(s, 2.8, y, 5.5, 0.5, '#1A1F3A')
        _add_rrect(s, 2.8, y, 5.5 * pct, 0.5, bar_c[i])
        _add_text(s, 1.0, y + 0.1, 1.7, 0.3, lbl, size=11, color=DIM, font=FB, align='right')
        _add_text(s, 8.5, y + 0.1, 1.0, 0.3, val, size=11, color=TXT, bold=True, font=FB)

    # Donut
    _add_oval(s, 9.5, 2.0, 3.6, 3.6, ACC)
    _add_oval(s, 10.3, 2.8, 2.0, 2.0, PRI)
    _add_text(s, 10.3, 3.5, 2.0, 0.5, '100%', size=18, color=ACC, bold=True, font=FH, align='center')
    for i, (lt, lc) in enumerate([('Inference 52%', ACC), ('Embedding 29%', MUT), ('Training 19%', PINK)]):
        ly = 2.2 + i * 0.4
        _add_rrect(s, 9.6, ly, 0.2, 0.2, lc)
        _add_text(s, 9.9, ly - 0.02, 1.5, 0.25, lt, size=10, color=TXT, font=FB)

    # --- Slide 3: Module cards ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.06, ACC)
    _add_text(s, 0.8, 0.35, 12, 0.5, 'CORE MODULES', size=24, color=ACC, bold=True, font=FH)
    _add_text(s, 0.8, 0.85, 12, 0.3, 'Active service mesh', size=11, color=DIM, font=FB)
    _add_rect(s, 0.8, 1.2, 11.7, 0.004, '#2A2F4A')

    cards = [
        ('INFERENCE\nENGINE', 'Auto-routing\n70% compression\nP99<10ms', ACC),
        ('MONITORING\nMESH', 'Millisecond alerting\nFull-stack tracing\nAnomaly detection', MUT),
        ('CONNECTOR\nHUB', '200+ integrations\nZero-config mesh\n5-min deploy', PINK),
    ]
    for i, (title, desc, accent) in enumerate(cards):
        x = 0.8 + i * 4.2
        _add_rrect(s, x, 1.8, 3.8, 3.0, BG2, line='#2A2F4A')
        _add_rect(s, x, 1.8, 3.8, 0.04, accent)
        _add_text(s, x + 0.3, 2.1, 3.2, 0.8, title, size=14, color=accent, bold=True, font=FH)
        _add_text(s, x + 0.3, 3.0, 3.2, 1.5, desc, size=11, color=DIM, font=FB)

    # --- Slide 4: Terminal code ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.06, ACC)
    _add_text(s, 0.8, 0.35, 12, 0.5, 'DEPLOY SEQUENCE', size=24, color=ACC, bold=True, font=FH)
    _add_text(s, 0.8, 0.85, 12, 0.3, 'Production pipeline', size=11, color=DIM, font=FB)
    _add_rect(s, 0.8, 1.2, 11.7, 0.004, '#2A2F4A')

    _add_rect(s, 0.8, 1.8, 11.7, 4.5, '#050810')
    _add_rect(s, 0.8, 1.8, 11.7, 0.04, ACC)
    _add_rrect(s, 0.8, 1.55, 1.0, 0.25, ACC)
    _add_text(s, 0.9, 1.57, 0.8, 0.2, 'shell', size=9, color='#000000', bold=True, font=FC)
    _add_multiline(s, 1.2, 2.2, 10.5, 3.5, [
        '$ neural-ops deploy --model gpt-4o-ft --gpu A100',
        '> Replicas: 3  |  Region: us-east-1  |  Auto-scale: ON',
        '> Health check: PASS  |  Endpoint: api.neuralops.ai/v4',
        '> Status: LIVE  |  P99: 2.1ms  |  Throughput: 12.8B/s',
    ], size=13, color=ACC, font=FB, spacing=8)

    # --- Slide 5: CTA ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 7.5, PRI)
    _add_rect(s, 0, 3.5, 13.333, 0.04, ACC)
    _add_rect(s, 0, 3.6, 13.333, 0.02, MUT)
    _add_text(s, 1.0, 2.0, 11, 1.2, 'INITIALIZE', size=52, color=ACC, bold=True, font=FH)
    _add_text(s, 1.0, 3.2, 11, 0.6, 'NEXT PHASE', size=28, color=MUT, font=FH)
    _add_text(s, 1.0, 4.2, 11, 0.5, 'neural-ops.ai  |  #build-the-future', size=14, color=DIM, font=FB)

    prs.save(output_path)


# =============================================================================
# Creative: Warm, playful, rounded, vibrant, generous whitespace
# =============================================================================
def build_creative(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    PRI = '#6C3CE1'
    ACC = '#FF6B6B'
    MUT = '#4ECDC4'
    YEL = '#FFE66D'
    BG  = '#FAFAFA'
    DK  = '#2D3436'
    BDY = '#636E72'
    MUTED = '#B2BEC3'
    CARD = '#FFFFFF'
    FH = 'Fredoka'
    FB = 'Nunito'

    # --- Slide 0: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 7.5, PRI)
    _add_oval(s, 8.5, -1.5, 7, 7, ACC)
    _add_oval(s, -2.5, 4.5, 6, 6, MUT)
    _add_oval(s, 10, 5, 4, 4, YEL)
    _add_text(s, 1.5, 2.0, 8, 1.5, 'Creative Studio', size=52, color='#FFFFFF', bold=True, font=FH)
    _add_text(s, 1.5, 3.5, 8, 0.6, '2025', size=36, color=YEL, bold=True, font=FH)
    _add_text(s, 1.5, 4.3, 8, 0.5, 'Design. Build. Ship. Repeat.', size=18, color='#FFFFFF', font=FB)

    # --- Slide 1: KPI ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.1, PRI)
    _add_text(s, 1.0, 0.5, 11, 0.6, 'Impact Numbers', size=30, color=DK, bold=True, font=FH)
    _add_text(s, 1.0, 1.1, 11, 0.3, 'What we achieved together', size=14, color=MUTED, font=FB)

    kpis = [
        ('847', 'Projects Delivered', '+23%', ACC),
        ('99.2%', 'Client Satisfaction', '+1.8%', MUT),
        ('156', 'Team Members', '+34', YEL),
    ]
    for i, (num, lbl, trend, accent) in enumerate(kpis):
        x = 1.0 + i * 3.9
        _add_rrect(s, x, 2.0, 3.5, 2.5, CARD, line='#E0E0E0')
        _add_oval(s, x + 1.3, 2.2, 0.9, 0.9, accent)
        _add_text(s, x + 0.3, 3.3, 2.9, 0.6, num, size=36, color=DK, bold=True, font=FH, align='center')
        _add_text(s, x + 0.3, 3.9, 2.9, 0.3, lbl, size=14, color=BDY, font=FB, align='center')
        _add_text(s, x + 0.3, 4.3, 2.9, 0.3, trend, size=12, color=accent, bold=True, font=FB, align='center')

    # --- Slide 2: Service mix ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.1, PRI)
    _add_text(s, 1.0, 0.5, 11, 0.6, 'Service Mix', size=30, color=DK, bold=True, font=FH)
    _add_text(s, 1.0, 1.1, 11, 0.3, 'Where we create value', size=14, color=MUTED, font=FB)

    bars = [('Brand Design', 0.72, '61%', PRI), ('Product UX', 0.55, '47%', ACC), ('Motion Design', 0.38, '32%', MUT)]
    for i, (lbl, pct, val, clr) in enumerate(bars):
        y = 2.2 + i * 0.9
        _add_rrect(s, 3.0, y, 5.0, 0.55, '#F0F0F5')
        _add_rrect(s, 3.0, y, 5.0 * pct, 0.55, clr)
        _add_text(s, 1.2, y + 0.1, 1.7, 0.3, lbl, size=14, color=DK, bold=True, font=FB, align='right')
        _add_text(s, 8.2, y + 0.1, 1.0, 0.3, val, size=14, color=DK, bold=True, font=FB)

    # Donut
    _add_oval(s, 9.5, 2.0, 3.6, 3.6, PRI)
    _add_oval(s, 10.3, 2.8, 2.0, 2.0, BG)
    _add_text(s, 10.3, 3.5, 2.0, 0.5, '100%', size=20, color=PRI, bold=True, font=FH, align='center')
    for i, (lt, lc) in enumerate([('Brand 40%', PRI), ('UX 35%', ACC), ('Motion 25%', MUT)]):
        ly = 2.2 + i * 0.4
        _add_rrect(s, 9.6, ly, 0.2, 0.2, lc)
        _add_text(s, 9.9, ly - 0.02, 1.5, 0.25, lt, size=12, color=BDY, font=FB)

    # --- Slide 3: Superpowers ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.1, PRI)
    _add_text(s, 1.0, 0.5, 11, 0.6, 'Superpowers', size=30, color=DK, bold=True, font=FH)
    _add_text(s, 1.0, 1.1, 11, 0.3, 'Our creative toolkit', size=14, color=MUTED, font=FB)

    cards = [
        ('Design System', 'Component library with\n500+ tokens, auto-sync Figma', PRI, '🎨'),
        ('Rapid Prototyping', 'Idea to clickable\nprototype in 48 hours', ACC, '⚡'),
        ('User Research', 'Data-driven design\ndecisions, A/B testing', MUT, '🔬'),
    ]
    for i, (title, desc, accent, emoji) in enumerate(cards):
        x = 1.0 + i * 3.9
        _add_rrect(s, x, 2.0, 3.5, 3.0, CARD, line='#E0E0E0')
        _add_oval(s, x + 1.3, 2.3, 0.9, 0.9, accent)
        _add_text(s, x + 1.3, 2.35, 0.9, 0.8, emoji, size=28, color='#FFFFFF', font=FB, align='center', anchor='ctr')
        _add_text(s, x + 0.3, 3.4, 2.9, 0.4, title, size=18, color=DK, bold=True, font=FH, align='center')
        _add_text(s, x + 0.3, 3.9, 2.9, 1.0, desc, size=13, color=BDY, font=FB, align='center')

    # --- Slide 4: Before & After ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 0.1, PRI)
    _add_text(s, 1.0, 0.5, 11, 0.6, 'Before & After', size=30, color=DK, bold=True, font=FH)
    _add_text(s, 1.0, 1.1, 11, 0.3, 'Client transformation stories', size=14, color=MUTED, font=FB)

    metrics = [
        ('Time to Market', '6 mo', '8 wk', 0.9, 0.3),
        ('Design Consistency', '40%', '95%', 0.4, 0.95),
        ('User Retention', '32%', '78%', 0.32, 0.78),
    ]
    for i, (lbl, v_old, v_new, pct_old, pct_new) in enumerate(metrics):
        y = 2.2 + i * 1.2
        _add_text(s, 1.2, y, 2.0, 0.3, lbl, size=14, color=DK, bold=True, font=FB)
        _add_rrect(s, 3.5, y + 0.3, 5.0, 0.3, '#F0F0F5')
        _add_rrect(s, 3.5, y + 0.3, 5.0 * pct_old, 0.3, MUTED)
        _add_text(s, 8.7, y + 0.3, 1.0, 0.3, v_old, size=12, color=MUTED, font=FB)
        _add_rrect(s, 3.5, y + 0.7, 5.0, 0.3, '#F0F0F5')
        _add_rrect(s, 3.5, y + 0.7, 5.0 * pct_new, 0.3, PRI)
        _add_text(s, 8.7, y + 0.7, 1.0, 0.3, v_new, size=12, color=DK, bold=True, font=FB)

    # --- Slide 5: CTA ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(s, 0, 0, 13.333, 7.5, PRI)
    _add_oval(s, 9.5, 0.5, 5, 5, ACC)
    _add_oval(s, -1.5, 5.5, 4, 4, MUT)
    _add_oval(s, 11, 6, 3, 3, YEL)
    _add_text(s, 1.5, 2.5, 8, 1.2, "Let's Create Together", size=44, color='#FFFFFF', bold=True, font=FH)
    _add_text(s, 1.5, 4.0, 8, 0.5, 'hello@creativestudio.io  |  book a free consultation', size=16, color='#FFFFFF', font=FB)

    prs.save(output_path)


def main():
    out = os.path.join(tempfile.gettempdir(), "ppt_raw_design")
    os.makedirs(out, exist_ok=True)

    print("="*80)
    print("  RAW python-pptx: 3 genuinely different designs")
    print("="*80)

    p1 = os.path.join(out, "mckinsey.pptx")
    p2 = os.path.join(out, "cyberpunk.pptx")
    p3 = os.path.join(out, "creative.pptx")

    print("\n[1/3] McKinsey...")
    build_mckinsey(p1)
    print("[2/3] Cyberpunk...")
    build_cyberpunk(p2)
    print("[3/3] Creative...")
    build_creative(p3)

    for name, label in [(p1, "McKinsey"), (p2, "Cyberpunk"), (p3, "Creative")]:
        sz = round(os.path.getsize(name)/1024, 1)
        prs = Presentation(name)
        print(f"  {label}: {len(prs.slides)} slides, {sz} KB")

    print("\nDone:", out)


if __name__ == "__main__":
    main()
