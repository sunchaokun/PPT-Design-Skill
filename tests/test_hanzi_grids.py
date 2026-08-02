"""Tests for Chinese character writing grid functions."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches


def _count_connectors(slide):
    return sum(1 for s in slide.shapes if int(s.shape_type) == 9)


def _count_textboxes(slide):
    return sum(1 for s in slide.shapes if s.has_text_frame and int(s.shape_type) != 9)


class TestMiziGrid:

    def test_empty_grid(self):
        from ppt_pro_max.build_helpers import mizi_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        mizi_grid(slide, 1.0, 1.5, 2.5)
        assert _count_connectors(slide) == 8
        assert _count_textboxes(slide) == 0

    def test_with_char(self):
        from ppt_pro_max.build_helpers import mizi_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        mizi_grid(slide, 1.0, 1.5, 2.5, char='永')
        assert _count_connectors(slide) == 8
        assert _count_textboxes(slide) == 1

    def test_char_content(self):
        from ppt_pro_max.build_helpers import mizi_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        mizi_grid(slide, 1.0, 1.5, 2.5, char='和')
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert '和' in texts

    def test_custom_colors(self):
        from ppt_pro_max.build_helpers import mizi_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        mizi_grid(slide, 1.0, 1.5, 2.5, char='永',
                  border_color='#FF0000', guide_color='#0000FF',
                  font_color='#333333')
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert '永' in texts


class TestTianGrid:

    def test_empty_grid(self):
        from ppt_pro_max.build_helpers import tian_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        tian_grid(slide, 1.0, 1.5, 2.5)
        assert _count_connectors(slide) == 6
        assert _count_textboxes(slide) == 0

    def test_with_char(self):
        from ppt_pro_max.build_helpers import tian_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        tian_grid(slide, 1.0, 1.5, 2.5, char='永')
        assert _count_connectors(slide) == 6
        assert _count_textboxes(slide) == 1

    def test_fewer_connectors_than_mizi(self):
        from ppt_pro_max.build_helpers import mizi_grid, tian_grid, add_slide
        prs = Presentation()
        s1 = add_slide(prs)
        mizi_grid(s1, 1.0, 1.5, 2.5)
        s2 = add_slide(prs)
        tian_grid(s2, 1.0, 1.5, 2.5)
        assert _count_connectors(s2) == _count_connectors(s1) - 2


class TestPinyinGrid:

    def test_empty_grid(self):
        from ppt_pro_max.build_helpers import pinyin_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        pinyin_grid(slide, 1.0, 2.0, 3.0)
        assert _count_connectors(slide) == 4
        assert _count_textboxes(slide) == 0

    def test_with_pinyin(self):
        from ppt_pro_max.build_helpers import pinyin_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        pinyin_grid(slide, 1.0, 2.0, 3.0, pinyin='yǒng')
        assert _count_connectors(slide) == 4
        assert _count_textboxes(slide) == 1

    def test_pinyin_content(self):
        from ppt_pro_max.build_helpers import pinyin_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        pinyin_grid(slide, 1.0, 2.0, 3.0, pinyin='hé')
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert 'hé' in texts

    def test_custom_baseline(self):
        from ppt_pro_max.build_helpers import pinyin_grid, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        pinyin_grid(slide, 1.0, 2.0, 3.0, pinyin='yǒng', baseline_y=3.0)
        assert _count_connectors(slide) == 4


class TestHanziRow:

    def test_single_char(self):
        from ppt_pro_max.build_helpers import hanzi_row, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        hanzi_row(slide, 1.0, 1.5, 2.0, ['永'])
        assert _count_textboxes(slide) == 1

    def test_multiple_chars_mizi(self):
        from ppt_pro_max.build_helpers import hanzi_row, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        hanzi_row(slide, 1.0, 1.5, 2.0, ['永', None, '和'], grid_type='mizi')
        assert _count_connectors(slide) == 8 * 3
        assert _count_textboxes(slide) == 2

    def test_multiple_chars_tian(self):
        from ppt_pro_max.build_helpers import hanzi_row, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        hanzi_row(slide, 1.0, 1.5, 2.0, ['永', '和'], grid_type='tian')
        assert _count_connectors(slide) == 6 * 2
        assert _count_textboxes(slide) == 2

    def test_none_draws_empty_grid(self):
        from ppt_pro_max.build_helpers import hanzi_row, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        hanzi_row(slide, 1.0, 1.5, 2.0, [None, None])
        assert _count_connectors(slide) == 8 * 2
        assert _count_textboxes(slide) == 0


class TestPinyinHanziBlock:

    def test_single_block(self):
        from ppt_pro_max.build_helpers import pinyin_hanzi_block, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        pinyin_hanzi_block(slide, 1.0, 0.5, 2.0, [('yǒng', '永')])
        assert _count_connectors(slide) == 4 + 8
        assert _count_textboxes(slide) == 2

    def test_multiple_blocks(self):
        from ppt_pro_max.build_helpers import pinyin_hanzi_block, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        pinyin_hanzi_block(slide, 0.5, 0.5, 1.8, [('yǒng', '永'), ('hé', '和')])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert 'yǒng' in texts
        assert '永' in texts
        assert 'hé' in texts
        assert '和' in texts

    def test_empty_block(self):
        from ppt_pro_max.build_helpers import pinyin_hanzi_block, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        pinyin_hanzi_block(slide, 1.0, 0.5, 2.0, [(None, None)])
        assert _count_textboxes(slide) == 0

    def test_pinyin_only_no_char(self):
        from ppt_pro_max.build_helpers import pinyin_hanzi_block, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        pinyin_hanzi_block(slide, 1.0, 0.5, 2.0, [('yǒng', None)])
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert 'yǒng' in texts


class TestGridVisualOutput:

    def test_mizi_produces_valid_pptx(self, tmp_path):
        from ppt_pro_max.build_helpers import mizi_grid, add_slide
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = add_slide(prs)
        mizi_grid(slide, 2.0, 2.0, 3.5, char='永')
        path = str(tmp_path / 'mizi.pptx')
        prs.save(path)
        import os
        assert os.path.getsize(path) > 0

    def test_pinyin_hanzi_produces_valid_pptx(self, tmp_path):
        from ppt_pro_max.build_helpers import pinyin_hanzi_block, add_slide
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = add_slide(prs)
        pinyin_hanzi_block(slide, 0.5, 0.5, 2.0,
                           [('yǒng', '永'), ('hé', '和'), (None, None)])
        path = str(tmp_path / 'pinyin_hanzi.pptx')
        prs.save(path)
        import os
        assert os.path.getsize(path) > 0
