"""Tests for render_preview — pptx → PNG preview tooling.

Backend engines are mocked so tests run on any machine (CI included).
"""
from __future__ import annotations

from unittest.mock import patch

import pytest

from ppt_pro_max.render_preview import (
    _build_html,
    _find_libreoffice,
    _render_libreoffice,
    detect_engine,
    render_preview,
)


# ── fixture: a tiny real pptx via python-pptx ──────────────────────────────
@pytest.fixture
def sample_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for _ in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "preview test slide"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def _make_pngs(d, n):
    paths = []
    for i in range(1, n + 1):
        p = d / f"slide{i}.png"
        p.write_bytes(b"PNG")
        paths.append(p)
    return paths


# ── _find_libreoffice ──────────────────────────────────────────────────────
def test_find_libreoffice_none_when_absent():
    with patch("ppt_pro_max.render_preview.shutil.which", return_value=None), \
         patch("ppt_pro_max.render_preview.os.path.isfile", return_value=False):
        assert _find_libreoffice() is None


def test_find_libreoffice_prefers_soffice_bin():
    """Must prefer soffice.bin (soffice.exe launcher hangs in headless/CI)."""
    with patch("ppt_pro_max.render_preview.shutil.which", return_value=None), \
         patch("ppt_pro_max.render_preview.os.path.isfile",
               side_effect=lambda p: p.endswith("soffice.bin")):
        found = _find_libreoffice()
    assert found is not None
    assert found.endswith("soffice.bin")


# ── detect_engine ──────────────────────────────────────────────────────────
def test_detect_engine_none_when_no_backend():
    with patch("ppt_pro_max.render_preview._powerpoint_available", return_value=False), \
         patch("ppt_pro_max.render_preview._find_libreoffice", return_value=None):
        assert detect_engine() == "none"


def test_detect_engine_powerpoint_when_available():
    with patch("ppt_pro_max.render_preview._powerpoint_available", return_value=True):
        assert detect_engine() == "powerpoint"


def test_detect_engine_libreoffice_when_available():
    with patch("ppt_pro_max.render_preview._powerpoint_available", return_value=False), \
         patch("ppt_pro_max.render_preview._find_libreoffice", return_value="/usr/bin/soffice.bin"):
        assert detect_engine() == "libreoffice"


# ── _render_libreoffice ────────────────────────────────────────────────────
def test_render_libreoffice_uses_soffice_bin_with_norestore(sample_pptx, tmp_path):
    """soffice.bin must be invoked headless+norestore, then PDF→PNG."""
    from pathlib import Path

    import subprocess

    calls = []

    def _fake_run(cmd, **kw):
        calls.append(cmd)
        # create a fake PDF in the outdir so downstream can find it
        pdf = Path(cmd[cmd.index("--outdir") + 1]) / "deck.pdf"
        pdf.write_bytes(b"%PDF-1.7 test")
        return subprocess.CompletedProcess(cmd, 0)

    fake_pngs = tmp_path / "out" / "slide1.png"

    def _fake_convert(pdf_path, out_dir):
        fake_pngs.parent.mkdir(parents=True, exist_ok=True)
        from PIL import Image
        img = Image.new("RGB", (100, 100))
        img.save(fake_pngs, "PNG")
        return [img]

    with patch("ppt_pro_max.render_preview._find_libreoffice",
               return_value=r"C:\LibreOffice\program\soffice.bin"), \
         patch("ppt_pro_max.render_preview.subprocess.run", side_effect=_fake_run), \
         patch("ppt_pro_max.render_preview._pdf_to_pngs", side_effect=_fake_convert):
        out_dir = tmp_path / "out"
        out_dir.mkdir()
        pngs = _render_libreoffice(sample_pptx, out_dir)

    assert pngs
    # soffice.bin invoked with --headless --norestore (skip taskkill cleanup call)
    soffice_call = next(c for c in calls if "--convert-to" in c)
    assert soffice_call[0].endswith("soffice.bin")
    assert "--headless" in soffice_call
    assert "--norestore" in soffice_call


def test_render_libreoffice_raises_when_no_soffice(sample_pptx, tmp_path):
    with patch("ppt_pro_max.render_preview._find_libreoffice", return_value=None):
        with pytest.raises(RuntimeError, match="LibreOffice not found"):
            _render_libreoffice(sample_pptx, tmp_path)


# ── render_preview ─────────────────────────────────────────────────────────
def test_render_preview_powerpoint(sample_pptx, tmp_path):
    out = tmp_path / "preview"
    with patch("ppt_pro_max.render_preview._render_powerpoint",
               side_effect=lambda p, d, w, h: _make_pngs(d, 3)):
        result = render_preview(str(sample_pptx), out_dir=str(out))
    assert result["engine"] == "powerpoint"
    assert len(result["pngs"]) == 3
    assert all(p.exists() for p in result["pngs"])
    assert result["html"].exists()


def test_render_preview_falls_back_to_libreoffice_on_com_failure(sample_pptx, tmp_path):
    """PowerPoint COM fails (e.g. WinError 1312 in sandbox) → auto-fallback to LO."""
    with patch("ppt_pro_max.render_preview._render_powerpoint",
               side_effect=RuntimeError("COM error 0x80070520")), \
         patch("ppt_pro_max.render_preview._render_libreoffice",
               side_effect=lambda p, d, w, h: _make_pngs(d, 3)):
        result = render_preview(str(sample_pptx), out_dir=str(tmp_path / "fb"))
    assert result["engine"] == "libreoffice"
    assert len(result["pngs"]) == 3
    assert len(result["warnings"]) == 1
    assert "powerpoint" in result["warnings"][0]


def test_render_preview_forced_powerpoint_failure_raises(sample_pptx, tmp_path):
    """Explicit engine= that fails must raise (no silent fallback)."""
    with patch("ppt_pro_max.render_preview._render_powerpoint",
               side_effect=RuntimeError("boom")), \
         pytest.raises(RuntimeError):
        render_preview(str(sample_pptx), out_dir=str(tmp_path / "fp"), engine="powerpoint")


def test_render_preview_all_engines_fail(sample_pptx, tmp_path):
    with patch("ppt_pro_max.render_preview._render_powerpoint", side_effect=RuntimeError("a")), \
         patch("ppt_pro_max.render_preview._render_libreoffice", side_effect=RuntimeError("b")), \
         pytest.raises(RuntimeError, match="All rendering engines failed"):
        render_preview(str(sample_pptx), out_dir=str(tmp_path / "af"))


def test_render_preview_unknown_engine(sample_pptx, tmp_path):
    with pytest.raises(RuntimeError):
        render_preview(str(sample_pptx), out_dir=str(tmp_path / "x"), engine="bogus")


def test_render_preview_relative_out_dir_resolved_to_pptx_dir(sample_pptx):
    with patch("ppt_pro_max.render_preview._render_powerpoint",
               side_effect=lambda p, d, w, h: _make_pngs(d, 2)):
        result = render_preview(str(sample_pptx), out_dir="preview_rel")
    expected = sample_pptx.parent / "preview_rel"
    assert result["pngs"][0].parent == expected
    assert result["pngs"][0].exists()


def test_render_preview_relative_pptx_resolved(sample_pptx, monkeypatch):
    monkeypatch.chdir(sample_pptx.parent)
    with patch("ppt_pro_max.render_preview._render_powerpoint",
               side_effect=lambda p, d, w, h: _make_pngs(d, 1)):
        result = render_preview(sample_pptx.name)
    assert result["pngs"][0].exists()


def test_render_preview_missing_file(tmp_path):
    with pytest.raises(FileNotFoundError):
        render_preview(str(tmp_path / "nope.pptx"))


# ── _build_html ────────────────────────────────────────────────────────────
def test_build_html_embeds_every_png(tmp_path):
    d = tmp_path / "preview"
    d.mkdir()
    pngs = [d / "slide1.png", d / "slide2.png"]
    for p in pngs:
        p.write_bytes(b"PNG")
    html = _build_html(pngs, d.parent / "deck.pptx", title="My Deck")
    assert html.exists()
    text = html.read_text(encoding="utf-8")
    assert "slide1.png" in text
    assert "slide2.png" in text
    assert "My Deck" in text
    assert "2 slides" in text
