"""E2E test: generate 20-page rural tourism PPT from xiangcun-project assets, then inspect every slide."""

import os
import json

from pptx import Presentation
from pptx.oxml.ns import qn

from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor

PROJECT_DIR = r"E:\PPT-Design-Skill\docs\xiangcun-project"
OUTPUT_DIR = os.path.join(PROJECT_DIR, "output")


def test_generate_xiangcun_ppt():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    pipeline = EnterprisePipeline()
    result = pipeline.run(
        query="数字乡村建设规划",
        project_dir=PROJECT_DIR,
        style="professional",
    )

    pptx_path = result.get("output_path", "")
    assert pptx_path and os.path.isfile(pptx_path), f"No output: {result}"

    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"\nTotal slides: {total_slides}")
    assert total_slides >= 15, f"Expected >=15 slides, got {total_slides}"

    gate = DeliveryGate()
    gate_report = gate.check(pptx_path, None, [])

    print(f"\nDeliveryGate: passed={gate_report.passed} fatals={len(gate_report.fatals)} warnings={len(gate_report.warnings)} passable={gate_report.is_passable}")

    if gate_report.fatals:
        print("FATALS:")
        for f in gate_report.fatals:
            print(f"  Slide {f.slide_index}: [{f.check_id}] {f.message}")

    if gate_report.warnings:
        print("WARNINGS:")
        for w in gate_report.warnings:
            print(f"  Slide {w.slide_index}: [{w.check_id}] {w.message} {w.detail}")

    with open(os.path.join(PROJECT_DIR, "content.json"), "r", encoding="utf-8") as f:
        content = json.load(f)
    planned_slides = content["slides"]

    print(f"\n{'='*80}")
    print(f"{'SLIDE':<6} {'SHAPES':<8} {'TEXTS':<6} {'IMGS':<5} {'MIN_FT':<8} {'MAX_FT':<8} {'GOAL':<10} TITLE")
    print(f"{'='*80}")

    issues = []

    for idx, slide in enumerate(prs.slides):
        goal = planned_slides[idx]["goal"] if idx < len(planned_slides) else "?"
        title = planned_slides[idx]["title"] if idx < len(planned_slides) else "?"
        shape_count = len(slide.shapes)

        texts = []
        font_sizes = []
        image_count = 0
        has_bg = False
        has_decoration = False

        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text and p.text.strip():
                        texts.append(p.text.strip()[:80])
                    if p.font.size:
                        font_sizes.append(p.font.size / 12700)
                    for r in p.runs:
                        if r.font.size:
                            font_sizes.append(r.font.size / 12700)
            if hasattr(shape, "image"):
                try:
                    _ = shape.image
                    image_count += 1
                except Exception:
                    pass

        bg_elem = slide.background._element
        has_bg = bool(list(bg_elem.iter(f"{{{a_ns}}}srgbClr")))

        sp_tree = slide._element.find(qn("p:cSld"))
        if sp_tree is not None:
            sp_tree = sp_tree.find(qn("p:spTree"))
        if sp_tree is not None:
            for child in sp_tree:
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag in ("sp", "pic"):
                    has_text = any(t.text and t.text.strip() for t in child.iter(f"{{{a_ns}}}t"))
                    if not has_text:
                        has_decoration = True

        min_font = min(font_sizes) if font_sizes else None
        max_font = max(font_sizes) if font_sizes else None

        status = ""
        if shape_count < 2:
            status += " [BLANK?]"
        if not texts:
            status += " [NO_TEXT]"
        if min_font is not None and min_font < 11:
            status += f" [FONT<{min_font:.0f}pt]"
        if not has_bg and goal not in ("section",):
            status += " [NO_BG]"

        if shape_count < 2:
            issues.append((idx, "Blank slide — less than 2 shapes"))
        if not texts:
            issues.append((idx, "No text content on slide"))
        if min_font is not None and min_font < 11:
            issues.append((idx, f"Font too small: {min_font:.0f}pt"))

        min_str = f"{min_font:.0f}" if min_font else "-"
        max_str = f"{max_font:.0f}" if max_font else "-"
        print(f"{idx:<6} {shape_count:<8} {len(texts):<6} {image_count:<5} {min_str:<8} {max_str:<8} {goal:<10} {title[:30]}{status}")

    print(f"\n{'='*80}")
    print(f"PER-SLIDE DETAIL")
    print(f"{'='*80}")

    for idx, slide in enumerate(prs.slides):
        goal = planned_slides[idx]["goal"] if idx < len(planned_slides) else "?"
        title_text = planned_slides[idx]["title"] if idx < len(planned_slides) else "?"
        print(f"\n--- Slide {idx} [{goal}] {title_text} ---")
        for si, shape in enumerate(slide.shapes):
            info = f"  [{si}] {shape.shape_type}: name='{shape.name}'"
            if shape.has_text_frame:
                tf = shape.text_frame
                text_preview = (tf.text or "").replace("\n", "|")[:60]
                w_in = shape.width / 914400
                h_in = shape.height / 914400
                font_info = ""
                for p in tf.paragraphs:
                    if p.font.size:
                        font_info = f" font={p.font.size / 12700:.0f}pt"
                        break
                    for r in p.runs:
                        if r.font.size:
                            font_info = f" font={r.font.size / 12700:.0f}pt"
                            break
                    if font_info:
                        break
                info += f" '{text_preview}' {w_in:.1f}x{h_in:.1f}\"{font_info}"
            if hasattr(shape, "image"):
                try:
                    _ = shape.image
                    info += " [IMG]"
                except Exception:
                    pass
            print(info)

    if issues:
        print(f"\n{'='*80}")
        print(f"ISSUES FOUND: {len(issues)}")
        print(f"{'='*80}")
        for idx, msg in issues:
            print(f"  Slide {idx}: {msg}")
    else:
        print(f"\nAll slides pass quality checks!")

    assert len(issues) == 0, f"Found {len(issues)} quality issues: {issues[:5]}"
