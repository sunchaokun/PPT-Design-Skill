"""E2E test: analyze_template → build_helpers → output PPT"""
import os
import pytest
from pptx import Presentation


@pytest.fixture
def C():
    return {
        'primary': '#2E6504',
        'primary_mid': '#466740',
        'accent': '#7DA92F',
        'muted': '#84AF7D',
        'light': '#D4E3AC',
        'lighter': '#E7F3A6',
        'lightest': '#F2F8D6',
        'background': '#FFFFFF',
        'bg_tint': '#F2F8D6',
        'white': '#FFFFFF',
        'text_dark': '#0D4609',
        'text_body': '#2E6504',
        'text_muted': '#466740',
        'divider': '#84AF7D',
        'card_bg': '#F6FAE8',
        'highlight': '#7DA92F',
        'font_heading': 'Microsoft YaHei',
        'font_body': 'Microsoft YaHei',
    }


class TestE2EBuildFromAnalysis:

    # Removed in V2: analyze_template module removed, project= param removed
    # def test_analyze_then_build(self, tmp_path, C):
    #     ...

    def test_copy_decorations_preserves_template_look(self, tmp_path, C):
        from ppt_pro_max.build_helpers import add_slide, rect, top_bar, page_header, copy_decorations

        template_path = r"E:\PPT-Design-Skill\docs\分析脚本\template.pptx"
        if not os.path.isfile(template_path):
            pytest.skip("Real template not available")

        prs = Presentation(template_path)
        original_count = len(prs.slides)

        content_slide = None
        for slide in prs.slides:
            if len(slide.shapes) >= 5:
                content_slide = slide
                break
        assert content_slide is not None

        s = add_slide(prs)
        rect(s, 0, 0, 13.333, 7.5, 'background', C=C)
        copy_decorations(s, content_slide, skip_long_text=True, skip_image=True)
        page_header(s, '新页面标题', C=C)

        output_path = str(tmp_path / "with_decorations.pptx")
        prs.save(output_path)

        assert os.path.isfile(output_path)
        out_prs = Presentation(output_path)
        assert len(out_prs.slides) == original_count + 1
        new_slide = out_prs.slides[-1]
        assert len(new_slide.shapes) >= 3
