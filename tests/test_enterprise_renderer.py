"""Tests for EnterpriseRenderer and EnterpriseDesignDecider."""

from __future__ import annotations

import json
import tempfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def _make_png(path: Path, w: int = 100, h: int = 50) -> Path:
    img = Image.new("RGB", (w, h), color="red")
    img.save(str(path))
    return path


# ============================================================
# EnterpriseRenderer
# ============================================================

class TestEnterpriseRenderer:

    def test_create_from_template(self, tmp_path):
        """EnterpriseRenderer should create Presentation from template."""
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        template_path = tmp_path / "template.pptx"
        prs = Presentation()
        prs.save(str(template_path))

        renderer = EnterpriseRenderer(template_path=str(template_path))
        prs_out = renderer.create_presentation()
        assert prs_out is not None
        assert len(prs_out.slide_layouts) > 0

    def test_create_from_invalid_template_falls_back(self, tmp_path):
        """Invalid template should fall back to blank Presentation."""
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        bad_path = str(tmp_path / "bad.pptx")
        Path(bad_path).write_bytes(b"not a pptx")

        renderer = EnterpriseRenderer(template_path=bad_path)
        prs_out = renderer.create_presentation()
        assert prs_out is not None

    def test_create_no_template(self):
        """No template should create blank Presentation."""
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        renderer = EnterpriseRenderer()
        prs_out = renderer.create_presentation()
        assert prs_out is not None

    def test_add_slide_with_layout_name(self, tmp_path):
        """EnterpriseRenderer should add slide using layout name matching."""
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        template_path = tmp_path / "template.pptx"
        prs = Presentation()
        prs.save(str(template_path))

        renderer = EnterpriseRenderer(template_path=str(template_path))
        prs_out = renderer.create_presentation()
        slide = renderer.add_slide(prs_out, layout_name="Title Slide")
        assert slide is not None
        assert slide.slide_layout.name == "Title Slide"

    def test_add_slide_fallback_to_blank(self, tmp_path):
        """Unknown layout name should fall back to last layout."""
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        template_path = tmp_path / "template.pptx"
        prs = Presentation()
        prs.save(str(template_path))

        renderer = EnterpriseRenderer(template_path=str(template_path))
        prs_out = renderer.create_presentation()
        slide = renderer.add_slide(prs_out, layout_name="NonExistent Layout")
        assert slide is not None

    def test_insert_logo(self, tmp_path):
        """EnterpriseRenderer should insert LOGO on a slide."""
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        template_path = tmp_path / "template.pptptx"
        prs = Presentation()
        prs.save(str(tmp_path / "template.pptx"))

        logo_path = _make_png(tmp_path / "logo.png")

        renderer = EnterpriseRenderer(template_path=str(tmp_path / "template.pptx"))
        prs_out = renderer.create_presentation()
        slide = renderer.add_slide(prs_out)
        renderer.insert_logo(slide, str(logo_path), {"position": "top_right", "width_inches": 1.0}, prs=prs_out)
        logo_shapes = [s for s in slide.shapes if hasattr(s, 'image')]
        assert len(logo_shapes) == 1

    def test_insert_logo_skip_slide(self, tmp_path):
        """LOGO should be skipped on slides in skip_slides list."""
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        prs.save(str(tmp_path / "template.pptx"))

        logo_path = _make_png(tmp_path / "logo.png")

        renderer = EnterpriseRenderer(template_path=str(tmp_path / "template.pptx"))
        prs_out = renderer.create_presentation()
        slide = renderer.add_slide(prs_out)
        renderer.insert_logo(
            slide, str(logo_path),
            {"position": "top_right", "width_inches": 1.0, "skip_slides": ["hook"]},
            current_goal="hook",
            prs=prs_out,
        )
        logo_shapes = [s for s in slide.shapes if hasattr(s, 'image')]
        assert len(logo_shapes) == 0

    def test_save_presentation(self, tmp_path):
        """EnterpriseRenderer should save presentation to file."""
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        prs.save(str(tmp_path / "template.pptx"))

        renderer = EnterpriseRenderer(template_path=str(tmp_path / "template.pptx"))
        prs_out = renderer.create_presentation()
        out_path = str(tmp_path / "output.pptx")
        renderer.save(prs_out, out_path)
        assert Path(out_path).exists()


# ============================================================
# EnterpriseDesignDecider
# ============================================================

class TestEnterpriseDesignDecider:

    def test_decide_with_brand_spec(self):
        """EnterpriseDesignDecider should use BrandSpec for decisions."""
        from ppt_pro_max.enterprise.enterprise_decider import EnterpriseDesignDecider
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand_spec = BrandSpec(
            source="brand_json",
            colors={"primary": "#1A3C6E", "foreground": "#333333"},
            fonts={"heading": "Arial", "body": "Calibri"},
            layout_mapping={"hook": 0, "content": 2},
        )
        decider = EnterpriseDesignDecider(brand_spec=brand_spec)
        assert decider.brand_spec is brand_spec

    def test_decide_layout_mapping(self):
        """EnterpriseDesignDecider should map goals to layout indices."""
        from ppt_pro_max.enterprise.enterprise_decider import EnterpriseDesignDecider
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand_spec = BrandSpec(
            source="brand_json",
            layout_mapping={"hook": 0, "content": 2, "features": 3},
        )
        decider = EnterpriseDesignDecider(brand_spec=brand_spec)
        assert decider.resolve_layout_index("hook") == 0
        assert decider.resolve_layout_index("content") == 2
        assert decider.resolve_layout_index("features") == 3

    def test_decide_layout_fallback(self):
        """Unmapped goals should fall back to default layout index."""
        from ppt_pro_max.enterprise.enterprise_decider import EnterpriseDesignDecider
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand_spec = BrandSpec(source="none")
        decider = EnterpriseDesignDecider(brand_spec=brand_spec)
        result = decider.resolve_layout_index("unknown_goal")
        assert result is None

    def test_decide_density_from_business_mode(self):
        """EnterpriseDesignDecider should suggest density based on business_mode."""
        from ppt_pro_max.enterprise.enterprise_decider import EnterpriseDesignDecider
        decider = EnterpriseDesignDecider(business_mode="education")
        assert decider.suggest_density() == 7
        decider2 = EnterpriseDesignDecider(business_mode="pitch")
        assert decider2.suggest_density() == 4
