"""Phase 4 tests: DeliveryGate completion — 3 missing checks + component quality + font threshold + placeholder patterns."""

import os

import pytest
from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


def _add_textbox(slide, x, y, w, h, text, font_size=14):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.text = text
    return txBox


def _add_image(slide, tmp_path, name, color, pic_name):
    img_path = os.path.join(str(tmp_path), name)
    img = PILImage.new("RGB", (100, 100), color)
    img.save(img_path)
    img.close()
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1), Inches(2), Inches(2))
    pic.name = pic_name


def _save(prs, tmp_path, name):
    path = os.path.join(str(tmp_path), name)
    prs.save(path)
    return path


def _make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _make_pptx_with_text_overflow(tmp_path):
    prs, slide = _make_slide()
    txBox = _add_textbox(slide, 0.5, 0.5, 1.5, 0.5, "This is a very long text that should definitely overflow the tiny textbox we created", 12)
    txBox.text_frame.word_wrap = True
    return _save(prs, tmp_path, "overflow.pptx")


def _make_pptx_normal(tmp_path):
    prs, slide = _make_slide()
    _add_textbox(slide, 0.5, 0.5, 8, 4, "Normal content", 14)
    return _save(prs, tmp_path, "normal.pptx")


def _make_pptx_with_image(tmp_path, name, color, pic_name):
    prs, slide = _make_slide()
    _add_image(slide, tmp_path, name, color, pic_name)
    return _save(prs, tmp_path, f"{name.rsplit('.', 1)[0]}.pptx")


def _make_pptx_with_font(tmp_path, pt_size, name):
    prs, slide = _make_slide()
    _add_textbox(slide, 0.5, 0.5, 8, 4, "Text", pt_size)
    return _save(prs, tmp_path, name)


def _make_pptx_with_placeholder_text(tmp_path, text):
    prs, slide = _make_slide()
    _add_textbox(slide, 0.5, 0.5, 8, 4, text, 14)
    return _save(prs, tmp_path, "placeholder.pptx")


def _make_pptx_with_group_placeholder_text(tmp_path):
    prs, slide = _make_slide()
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_tree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    grp = etree.SubElement(sp_tree, f"{{{p_ns}}}grpSp")
    etree.SubElement(grp, f"{{{p_ns}}}grpSpPr")
    sp = etree.SubElement(grp, f"{{{p_ns}}}sp")
    etree.SubElement(sp, f"{{{p_ns}}}spPr")
    txBody = etree.SubElement(sp, f"{{{p_ns}}}txBody")
    etree.SubElement(txBody, f"{{{a_ns}}}bodyPr")
    etree.SubElement(txBody, f"{{{a_ns}}}lstStyle")
    p_elem = etree.SubElement(txBody, f"{{{a_ns}}}p")
    r = etree.SubElement(p_elem, f"{{{a_ns}}}r")
    rPr = etree.SubElement(r, f"{{{a_ns}}}rPr")
    rPr.set("lang", "en-US")
    t = etree.SubElement(r, f"{{{a_ns}}}t")
    t.text = "Click to add text"
    return _save(prs, tmp_path, "group_placeholder.pptx")


def _make_dna_with_decoration():
    from ppt_pro_max.enterprise.design_dna_extractor import DesignDNA, SlideDNA, TextZone
    return DesignDNA(
        source_path="test",
        slides=[
            SlideDNA(
                slide_index=0,
                page_type="content",
                slide_xml=b"<sld/>",
                rels_xml=b"<Relationships/>",
                text_zones=[
                    TextZone(zone_id="z1", role="body", bounds=(0.5, 1.5, 8, 4), text="content"),
                    TextZone(zone_id="z2", role="decoration", bounds=(0, 0, 0.5, 10), text=""),
                ],
                background_colors=["#1A1A2E"],
            )
        ],
        color_palette={},
        actual_colors={"#1A1A2E": 10},
        actual_fonts=[],
        brand_spec=None,
    )


def _make_dna_without_decoration():
    from ppt_pro_max.enterprise.design_dna_extractor import DesignDNA, SlideDNA, TextZone
    return DesignDNA(
        source_path="test",
        slides=[
            SlideDNA(
                slide_index=0,
                page_type="content",
                slide_xml=b"<sld/>",
                rels_xml=b"<Relationships/>",
                text_zones=[TextZone(zone_id="z1", role="body", bounds=(0.5, 1.5, 8, 4), text="content")],
                background_colors=[],
            )
        ],
        color_palette={},
        actual_colors={},
        actual_fonts=[],
        brand_spec=None,
    )


def _check(gate, path, dna=None, plans=None):
    report = gate.check(path, dna, plans or [])
    return [i for i in report.warnings + report.fatals]


# ── G1: text_overflow check ─────────────────────────────────────────

class TestTextOverflowCheck:
    def test_overflow_detected(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        items = _check(DeliveryGate(), _make_pptx_with_text_overflow(tmp_path))
        assert any(i.check_id == "text_overflow" for i in items)

    def test_no_overflow_on_normal(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        items = _check(DeliveryGate(), _make_pptx_normal(tmp_path))
        assert not any(i.check_id == "text_overflow" for i in items)

    def test_overflow_is_warning_not_fatal(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        gate = DeliveryGate()
        report = gate.check(_make_pptx_with_text_overflow(tmp_path), None, [])
        assert any(i.check_id == "text_overflow" for i in report.warnings)


# ── G1: decoration_missing check ────────────────────────────────────

class TestDecorationMissingCheck:
    def test_decoration_missing_when_template_has_deco(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        prs, slide = _make_slide()
        _add_textbox(slide, 0.5, 0.5, 8, 4, "Content without decoration", 14)
        path = _save(prs, tmp_path, "no_deco.pptx")
        plans = [type("P", (), {"page_type": "content"})()]
        items = _check(DeliveryGate(), path, _make_dna_with_decoration(), plans)
        assert any(i.check_id == "decoration_missing" for i in items)

    def test_no_warning_when_no_template_decoration(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        prs, slide = _make_slide()
        _add_textbox(slide, 0.5, 0.5, 8, 4, "Content", 14)
        path = _save(prs, tmp_path, "no_deco_template.pptx")
        plans = [type("P", (), {"page_type": "content"})()]
        items = _check(DeliveryGate(), path, _make_dna_without_decoration(), plans)
        assert not any(i.check_id == "decoration_missing" for i in items)


# ── G1: image_placeholder check ─────────────────────────────────────

class TestImagePlaceholderCheck:
    def test_placeholder_image_name_detected(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_image(tmp_path, "placeholder_img.png", "red", "placeholder image")
        items = _check(DeliveryGate(), path)
        assert any(i.check_id == "image_placeholder" for i in items)

    def test_stock_image_name_detected(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_image(tmp_path, "stock_img.png", "blue", "stock photo")
        items = _check(DeliveryGate(), path)
        assert any(i.check_id == "image_placeholder" for i in items)

    def test_normal_image_name_not_flagged(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_image(tmp_path, "normal_img.png", "green", "Product Screenshot")
        items = _check(DeliveryGate(), path)
        assert not any(i.check_id == "image_placeholder" for i in items)

    def test_image_placeholder_is_warning(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        gate = DeliveryGate()
        path = _make_pptx_with_image(tmp_path, "placeholder_img2.png", "red", "placeholder image")
        report = gate.check(path, None, [])
        assert any(i.check_id == "image_placeholder" for i in report.warnings)


# ── G2: Font threshold 9pt→11pt ─────────────────────────────────────

class TestFontThreshold11pt:
    def test_8pt_detected_as_too_small(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_font(tmp_path, 8, "small_font.pptx")
        items = _check(DeliveryGate(), path)
        assert any(i.check_id == "font_too_small" for i in items)

    def test_10pt_detected_as_too_small(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_font(tmp_path, 10, "font10pt.pptx")
        items = _check(DeliveryGate(), path)
        assert any(i.check_id == "font_too_small" for i in items)

    def test_auto_fix_bumps_to_11pt(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_font(tmp_path, 8, "small_font_fix.pptx")
        gate = DeliveryGate()
        report = gate.check(path, None, [])
        gate.auto_fix(path, None, [], report)
        prs = Presentation(path)
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        fixed = False
        for tag in (f"{{{a_ns}}}rPr", f"{{{a_ns}}}defRPr"):
            for rPr in prs.slides[0]._element.iter(tag):
                sz = rPr.get("sz")
                if sz:
                    assert int(sz) >= 1100
                    fixed = True
        assert fixed


# ── G3: Component quality check ─────────────────────────────────────

class TestComponentQualityCheck:
    def test_group_placeholder_text_detected(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_group_placeholder_text(tmp_path)
        items = _check(DeliveryGate(), path)
        assert any(i.check_id == "component_placeholder_residual" for i in items)

    def test_component_placeholder_is_warning(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        gate = DeliveryGate()
        path = _make_pptx_with_group_placeholder_text(tmp_path)
        report = gate.check(path, None, [])
        assert any(i.check_id == "component_placeholder_residual" for i in report.warnings)


# ── G4: Expanded placeholder patterns ───────────────────────────────

class TestExpandedPlaceholderPatterns:
    @pytest.mark.parametrize("text", [
        "Enter text here",
        "Insert title",
        "Your Title Here",
        "Add your content",
        "Type something",
        "Double click to edit",
    ])
    def test_english_placeholders_detected(self, tmp_path, text):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_placeholder_text(tmp_path, text)
        items = _check(DeliveryGate(), path)
        assert any(i.check_id == "residual_placeholder" for i in items)

    def test_normal_text_not_flagged(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        path = _make_pptx_with_placeholder_text(tmp_path, "Our mission is to deliver value")
        items = _check(DeliveryGate(), path)
        assert not any(i.check_id == "residual_placeholder" for i in items)


# ── Check count completeness ────────────────────────────────────────

class TestDeliveryGateCompleteness:
    def test_13_check_types_present(self, tmp_path):
        from ppt_pro_max.enterprise.delivery_gate import DeliveryGate
        prs, slide = _make_slide()
        _add_textbox(slide, 0.5, 0.5, 8, 4, "Content", 14)
        _save(prs, tmp_path, "completeness.pptx")

        gate = DeliveryGate()
        expected_ids = {
            "residual_placeholder",
            "title_duplicate",
            "blank_page",
            "color_break",
            "font_mismatch",
            "background_missing",
            "font_too_small",
            "broken_image_ref",
            "page_count_mismatch",
            "toc_mismatch",
            "text_overflow",
            "decoration_missing",
            "image_placeholder",
        }
        assert expected_ids.issubset(gate._ALL_CHECK_IDS)
