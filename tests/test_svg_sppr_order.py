"""Regression tests for OOXML spPr child element order.

Bug: GradientFill.apply() and set_solid_fill_with_alpha() used
etree.SubElement(spPr, ...) which appended fill AFTER <a:ln>, violating
OOXML schema order (xfrm → custGeom → fill → ln). LibreOffice silently
dropped fills that appeared after <a:ln>.

Fix: _insert_fill_before_ln() helper inserts fill elements before any
existing <a:ln>, or appends if no <a:ln> exists.
"""
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_pro_max.renderer.visual_effects import (
    GradientFill,
    GradientStop,
    _insert_fill_before_ln,
    set_solid_fill_with_alpha,
)


def _make_shape_with_ln():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        1, Inches(1), Inches(1), Inches(4), Inches(3)
    )
    spPr = shape._element.find(qn("p:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
        ln.set("w", "12700")
    return shape, spPr


class TestInsertFillBeforeLn:
    def test_insert_before_existing_ln(self):
        _shape, spPr = _make_shape_with_ln()
        fill_el = etree.Element(qn("a:solidFill"))
        srgb = etree.SubElement(fill_el, qn("a:srgbClr"))
        srgb.set("val", "FF0000")
        _insert_fill_before_ln(spPr, fill_el)
        children = list(spPr)
        ln_idx = next(i for i, c in enumerate(children) if c.tag == qn("a:ln"))
        fill_idx = next(i for i, c in enumerate(children) if c.tag == qn("a:solidFill"))
        assert fill_idx < ln_idx, "fill must precede <a:ln> in spPr child order"

    def test_append_when_no_ln(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(4), Inches(3)
        )
        spPr = shape._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            spPr.remove(ln)
        fill_el = etree.Element(qn("a:solidFill"))
        srgb = etree.SubElement(fill_el, qn("a:srgbClr"))
        srgb.set("val", "00FF00")
        _insert_fill_before_ln(spPr, fill_el)
        assert spPr.find(qn("a:solidFill")) is not None

    def test_multiple_fills_all_before_ln(self):
        _shape, spPr = _make_shape_with_ln()
        for color in ["FF0000", "00FF00", "0000FF"]:
            _remove_existing_fill(spPr)
            fill_el = etree.Element(qn("a:solidFill"))
            srgb = etree.SubElement(fill_el, qn("a:srgbClr"))
            srgb.set("val", color)
            _insert_fill_before_ln(spPr, fill_el)
        children = list(spPr)
        ln_idx = next(i for i, c in enumerate(children) if c.tag == qn("a:ln"))
        fill_tags = [c.tag for c in children[:ln_idx]]
        assert qn("a:solidFill") in fill_tags


def _remove_existing_fill(spPr):
    for tag in ["a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"]:
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)


class TestGradientFillElementOrder:
    def test_gradient_fill_before_ln(self):
        shape, spPr = _make_shape_with_ln()
        gf = GradientFill(
            gradient_type="linear",
            stops=[
                GradientStop(color="#1D78FA", position=0, alpha=100),
                GradientStop(color="#0D47A1", position=100000, alpha=100),
            ],
            angle=5400000,
        )
        gf.apply(shape)
        children = list(spPr)
        ln_idx = next(i for i, c in enumerate(children) if c.tag == qn("a:ln"))
        grad_idx = next(
            i for i, c in enumerate(children) if c.tag == qn("a:gradFill")
        )
        assert grad_idx < ln_idx, "gradFill must precede <a:ln>"

    def test_radial_gradient_fill_before_ln(self):
        shape, spPr = _make_shape_with_ln()
        gf = GradientFill(
            gradient_type="path",
            stops=[
                GradientStop(color="#FFFFFF", position=0, alpha=100),
                GradientStop(color="#000000", position=100000, alpha=100),
            ],
            fill_to_rect={"l": "50000", "t": "50000", "r": "50000", "b": "50000"},
        )
        gf.apply(shape)
        children = list(spPr)
        ln_idx = next(i for i, c in enumerate(children) if c.tag == qn("a:ln"))
        grad_idx = next(
            i for i, c in enumerate(children) if c.tag == qn("a:gradFill")
        )
        assert grad_idx < ln_idx, "gradFill (path) must precede <a:ln>"


class TestSolidFillWithAlphaElementOrder:
    def test_solid_fill_alpha_before_ln(self):
        shape, spPr = _make_shape_with_ln()
        set_solid_fill_with_alpha(shape, "#FF0000", 50)
        children = list(spPr)
        ln_idx = next(i for i, c in enumerate(children) if c.tag == qn("a:ln"))
        fill_idx = next(
            i for i, c in enumerate(children) if c.tag == qn("a:solidFill")
        )
        assert fill_idx < ln_idx, "solidFill must precede <a:ln>"

    def test_solid_fill_alpha_has_alpha_element(self):
        shape, spPr = _make_shape_with_ln()
        set_solid_fill_with_alpha(shape, "#00FF00", 75)
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None
        alpha_el = solidFill.find(qn("a:srgbClr")).find(qn("a:alpha"))
        assert alpha_el is not None
        assert alpha_el.get("val") == "75000"


class TestFreeformGradientElementOrder:
    """Regression: gradient polygon + freeform shapes disappeared in LibreOffice
    because gradFill was appended after <a:ln>."""

    def test_gradient_polygon_fill_before_ln(self, tmp_path):
        from ppt_pro_max.renderer.svg_compiler import SVGCompiler

        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">
          <defs>
            <linearGradient id="g1" x1="0" y1="0" x2="1" y2="1">
              <stop offset="0" stop-color="#7DA92F"/>
              <stop offset="1" stop-color="#2E6504"/>
            </linearGradient>
          </defs>
          <polygon points="200,30 380,280 20,280" fill="url(#g1)"/>
        </svg>"""

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (3.5, 0.8, 6.3, 5.6))
        assert result.shape_count > 0

        pptx_path = tmp_path / "gradient_polygon.pptx"
        prs.save(str(pptx_path))

        import zipfile
        with zipfile.ZipFile(pptx_path) as z:
            slide_xmls = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            xml = z.read(slide_xmls[0]).decode("utf-8")

        assert "<p:pic>" not in xml, "zero-raster guarantee"
        assert "gradFill" in xml, "gradient fill must be present"

        root = etree.fromstring(xml.encode("utf-8"))
        for spPr in root.iter(qn("p:spPr")):
            ln = spPr.find(qn("a:ln"))
            grad = spPr.find(qn("a:gradFill"))
            if ln is not None and grad is not None:
                children = list(spPr)
                ln_idx = next(i for i, c in enumerate(children) if c.tag == qn("a:ln"))
                grad_idx = next(
                    i for i, c in enumerate(children) if c.tag == qn("a:gradFill")
                )
                assert grad_idx < ln_idx, "gradFill must precede <a:ln> in spPr"
