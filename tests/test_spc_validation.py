"""Spc Validation Test — verify which OOXML form for character spacing
PowerPoint actually renders.

Generates two PPTX files:
  1. spc_child_element.pptx  — <a:rPr><a:spc val="1400"/></a:rPr>
  2. spc_attribute.pptx      — <a:rPr spc="1400"/>

Each file has 3 text boxes with different spacing values:
  - "Wide Spacing"    (spc=1400 = 14pt expanded)
  - "Normal Spacing"  (no spc attribute)
  - "Tight Spacing"   (spc=-400 = -4pt condensed)

Open both files in PowerPoint and compare:
  - If spc_child_element.pptx shows NO spacing difference → child element form doesn't work
  - If spc_attribute.pptx shows spacing differences → attribute form is correct
  - If both show identical results → both forms work (unlikely per ECMA-376)
"""
import os
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree


SLIDE_W = 13.333
SLIDE_H = 7.5
FONT_SIZE = 28
FONT_SIZE_HUNDREDTHS = FONT_SIZE * 100

SPACING_CONFIGS = [
    ("Wide Spacing (Expanded 14pt)", 1400),
    ("Normal Spacing (No spc)", None),
    ("Tight Spacing (Condensed 4pt)", -400),
]


def _add_textbox(slide, x, y, w, h, text, font_size=FONT_SIZE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = "Calibri"
    return run


def _generate_child_element_version(out_dir):
    """Generate PPTX using <a:rPr><a:spc val="..."/></a:rPr> (child element form)."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[-1])

    title_run = _add_textbox(slide, 0.5, 0.3, 12, 0.8, "Child Element Form: <a:spc val=\"...\"/>", 32)
    title_run.font.bold = True

    for i, (label, spc_val) in enumerate(SPACING_CONFIGS):
        y = 1.8 + i * 1.8
        desc_run = _add_textbox(slide, 0.5, y, 12, 0.5, f"[spc={spc_val}]", 18)
        from pptx.dml.color import RGBColor
        desc_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        text_run = _add_textbox(slide, 0.5, y + 0.5, 12, 0.8, label, FONT_SIZE)

        if spc_val is not None:
            rPr = text_run._r.get_or_add_rPr()
            spc_el = etree.SubElement(rPr, qn("a:spc"))
            spc_el.set("val", str(spc_val))

    out_path = os.path.join(out_dir, "spc_child_element.pptx")
    prs.save(out_path)
    return out_path


def _generate_attribute_version(out_dir):
    """Generate PPTX using <a:rPr spc="..."/> (attribute form)."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[-1])

    title_run = _add_textbox(slide, 0.5, 0.3, 12, 0.8, "Attribute Form: <a:rPr spc=\"...\"/>", 32)
    title_run.font.bold = True

    for i, (label, spc_val) in enumerate(SPACING_CONFIGS):
        y = 1.8 + i * 1.8
        desc_run = _add_textbox(slide, 0.5, y, 12, 0.5, f"[spc={spc_val}]", 18)
        from pptx.dml.color import RGBColor
        desc_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        text_run = _add_textbox(slide, 0.5, y + 0.5, 12, 0.8, label, FONT_SIZE)

        if spc_val is not None:
            rPr = text_run._r.get_or_add_rPr()
            rPr.set("spc", str(spc_val))

    out_path = os.path.join(out_dir, "spc_attribute.pptx")
    prs.save(out_path)
    return out_path


def test_spc_xml_structure_child():
    """Verify child element form produces correct XML structure."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    run = _add_textbox(slide, 0.5, 0.5, 5, 1, "Test", FONT_SIZE)
    rPr = run._r.get_or_add_rPr()
    spc_el = etree.SubElement(rPr, qn("a:spc"))
    spc_el.set("val", "1400")

    found = rPr.find(qn("a:spc"))
    assert found is not None, "a:spc child element should exist"
    assert found.get("val") == "1400", "a:spc val should be 1400"


def test_spc_xml_structure_attribute():
    """Verify attribute form produces correct XML structure."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[-1])
    run = _add_textbox(slide, 0.5, 0.5, 5, 1, "Test", FONT_SIZE)
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", "1400")

    assert rPr.get("spc") == "1400", "a:rPr spc attribute should be 1400"
    found = rPr.find(qn("a:spc"))
    assert found is None, "a:spc child element should NOT exist in attribute form"


def _get_first_run_rpr(prs):
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.runs:
                    return p.runs[0]._r.find(qn("a:rPr"))
    return None


def test_spc_roundtrip_child():
    """Verify child element form survives save/reload."""
    with tempfile.TemporaryDirectory() as td:
        path = os.path.join(td, "child.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        run = _add_textbox(slide, 0.5, 0.5, 5, 1, "Wide", FONT_SIZE)
        rPr = run._r.get_or_add_rPr()
        spc_el = etree.SubElement(rPr, qn("a:spc"))
        spc_el.set("val", "1400")
        prs.save(path)

        prs2 = Presentation(path)
        rPr2 = _get_first_run_rpr(prs2)
        assert rPr2 is not None
        spc2 = rPr2.find(qn("a:spc"))
        assert spc2 is not None, "a:spc child element should survive roundtrip"
        assert spc2.get("val") == "1400"


def test_spc_roundtrip_attribute():
    """Verify attribute form survives save/reload."""
    with tempfile.TemporaryDirectory() as td:
        path = os.path.join(td, "attr.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        run = _add_textbox(slide, 0.5, 0.5, 5, 1, "Wide", FONT_SIZE)
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", "1400")
        prs.save(path)

        prs2 = Presentation(path)
        rPr2 = _get_first_run_rpr(prs2)
        assert rPr2 is not None
        assert rPr2.get("spc") == "1400", "a:rPr spc attribute should survive roundtrip"


def test_generate_comparison_files():
    """Generate two PPTX files for manual visual comparison in PowerPoint.

    This test ALWAYS passes — it just generates files for human inspection.
    Open the output files in PowerPoint to verify which form renders correctly.
    """
    out_dir = os.path.join(os.path.dirname(__file__), "..", ".cache", "spc_validation")
    os.makedirs(out_dir, exist_ok=True)

    child_path = _generate_child_element_version(out_dir)
    attr_path = _generate_attribute_version(out_dir)

    assert os.path.exists(child_path), f"Child element file not created: {child_path}"
    assert os.path.exists(attr_path), f"Attribute file not created: {attr_path}"

    child_size = os.path.getsize(child_path)
    attr_size = os.path.getsize(attr_path)

    print(f"\n{'='*60}")
    print(f"Spc Validation Files Generated")
    print(f"{'='*60}")
    print(f"  Child element form: {child_path} ({child_size} bytes)")
    print(f"  Attribute form:     {attr_path} ({attr_size} bytes)")
    print(f"{'='*60}")
    print(f"Open BOTH files in PowerPoint and compare:")
    print(f"  1. If child_element shows NO spacing difference → child form doesn't work")
    print(f"  2. If attribute shows spacing differences → attribute form is correct")
    print(f"  3. If both identical → both forms work")
    print(f"{'='*60}")
