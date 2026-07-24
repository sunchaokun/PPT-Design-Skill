"""Truly raw injection - NO coordinate transform, NO 3D strip, NO recolor. Just paste XML as-is."""

import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)

lib = ComponentLibrary(find_db_path())

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

categories = ["hierarchy", "process", "swot", "infographic", "timeline", "chart"]

for cat in categories:
    results = lib.search(type="group", category=cat, limit=1)
    if not results:
        continue
    comp = results[0]
    xml_parts = lib.load_xml(comp["id"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sp_tree = slide._element.find(f".//{{{_P}}}spTree")

    # Parse and inject directly - ZERO modification
    grp_elem = etree.fromstring(xml_parts["group"])

    # Only remap cNvPr ids to avoid conflicts
    next_id = 2
    for sp in sp_tree:
        for cNv in sp.iter(f"{{{_P}}}cNvPr"):
            try:
                next_id = max(next_id, int(cNv.get("id", "1")) + 1)
            except (ValueError, TypeError):
                pass

    for cNv in grp_elem.iter(f"{{{_P}}}cNvPr"):
        cNv.set("id", str(next_id))
        cNv.set("name", f"Comp {next_id}")
        next_id += 1

    # Remove unresolvable image refs (blip r:embed)
    for blip in grp_elem.iter(f"{{{_A}}}blip"):
        embed = blip.get(f"{{{_R}}}embed")
        if embed:
            blip.attrib.pop(f"{{{_R}}}embed", None)

    sp_tree.append(grp_elem)
    print(f"  {cat}: id={comp['id']}, nodes={comp['node_count']}")

out = os.path.join(OUT_DIR, "truly_raw_inject.pptx")
prs.save(out)
print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")

lib.close()
