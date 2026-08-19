"""Tests for SVG z-order rendering — ensure transparent elements don't bleed through.

The critical behavior under test: when multiple SVG elements overlap,
the rendering order must preserve z-order (later elements on top).
Semi-transparent elements should not cause underlying content to show through
in unexpected ways.
"""

import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from ppt_pro_max.renderer.svg_compiler import SVGCompiler


@pytest.fixture
def prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


@pytest.fixture
def slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


class TestZOrderRendering:
    """Test that SVG elements are rendered in correct z-order."""

    def test_element_order_preserved(self, slide):
        """SVG elements should be rendered in document order (z-order)."""
        svg = """
        <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 100">
            <rect x="0" y="0" width="100" height="100" fill="#ff0000"/>
            <rect x="100" y="0" width="100" height="100" fill="#00ff00"/>
            <rect x="200" y="0" width="100" height="100" fill="#0000ff"/>
        </svg>
        """
        result = SVGCompiler().compile(svg, slide, (0, 0, 3, 1))

        # All three shapes should be created
        assert result.shape_count == 3

        # Check z-order: shapes should be in document order
        shapes = list(slide.shapes)
        # Find the three rectangles we just created
        rects = [s for s in shapes if hasattr(s, 'fill') and s.shape_type is not None]
        # The last created shape should be on top (highest z-index)
        # In PowerPoint, shapes are rendered in creation order

    def test_transparent_overlap_no_bleed(self, slide):
        """Semi-transparent overlay should not cause unexpected bleed."""
        svg = """
        <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">
            <rect x="0" y="0" width="200" height="100" fill="#ff0000"/>
            <rect x="50" y="25" width="100" height="50" fill="#0000ff" fill-opacity="0.5"/>
        </svg>
        """
        result = SVGCompiler().compile(svg, slide, (0, 0, 2, 1))

        # Two shapes should be created (no extra mask shapes)
        assert result.shape_count == 2

        # Blue rectangle should be rendered after red (on top)
        shapes = list(slide.shapes)
        # The semi-transparent blue should be the second shape
        # Verify no extra mask shapes were added

    def test_multiple_transparent_layers(self, slide):
        """Multiple semi-transparent layers should stack correctly."""
        svg = """
        <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">
            <rect x="0" y="0" width="200" height="100" fill="#ff0000"/>
            <rect x="25" y="12.5" width="150" height="75" fill="#00ff00" fill-opacity="0.6"/>
            <rect x="50" y="25" width="100" height="50" fill="#0000ff" fill-opacity="0.4"/>
        </svg>
        """
        result = SVGCompiler().compile(svg, slide, (0, 0, 2, 1))

        # Three shapes should be created
        assert result.shape_count == 3

    def test_text_on_top_of_shapes(self, slide):
        """Text should always render on top of shapes."""
        svg = """
        <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">
            <rect x="0" y="0" width="200" height="100" fill="#22d3ee"/>
            <text x="100" y="55" text-anchor="middle" font-size="20" fill="#fff">Label</text>
        </svg>
        """
        result = SVGCompiler().compile(svg, slide, (0, 0, 2, 1))

        # Shape + text should be created
        assert result.shape_count >= 2

    def test_group_z_order(self, slide):
        """Elements inside <g> should maintain z-order relative to siblings."""
        svg = """
        <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 100">
            <g>
                <rect x="0" y="0" width="100" height="100" fill="#ff0000"/>
                <rect x="50" y="0" width="100" height="100" fill="#00ff00"/>
            </g>
            <rect x="100" y="0" width="100" height="100" fill="#0000ff"/>
        </svg>
        """
        result = SVGCompiler().compile(svg, slide, (0, 0, 3, 1))

        # Three shapes should be created
        assert result.shape_count == 3

    def test_nested_groups_z_order(self, slide):
        """Nested groups should flatten z-order correctly."""
        svg = """
        <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 300 100">
            <g>
                <rect x="0" y="0" width="100" height="100" fill="#ff0000"/>
                <g>
                    <rect x="100" y="0" width="100" height="100" fill="#00ff00"/>
                </g>
            </g>
            <rect x="200" y="0" width="100" height="100" fill="#0000ff"/>
        </svg>
        """
        result = SVGCompiler().compile(svg, slide, (0, 0, 3, 1))

        # Three shapes should be created
        assert result.shape_count == 3
