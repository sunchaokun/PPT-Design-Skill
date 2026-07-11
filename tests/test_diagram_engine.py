"""Tests for DiagramEngine, DiagramStyle, Region, TextMeasurer, and diagram types."""

from __future__ import annotations

import math
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches


class TestRegion:

    def test_center_x(self):
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        r = Region(left=1.0, top=2.0, width=4.0, height=2.0)
        assert r.center_x == 3.0

    def test_center_y(self):
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        r = Region(left=0, top=1.0, width=6.0, height=4.0)
        assert r.center_y == 3.0

    def test_right_bottom(self):
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        r = Region(left=1.0, top=2.0, width=3.0, height=4.0)
        assert r.right == 4.0
        assert r.bottom == 6.0

    def test_subregion(self):
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        r = Region(left=0, top=0, width=10.0, height=8.0)
        sub = r.subregion(0.1, 0.1, 0.8, 0.8)
        assert sub.left == 1.0
        assert sub.top == 0.8
        assert abs(sub.width - 8.0) < 0.01
        assert abs(sub.height - 6.4) < 0.01

    def test_inset(self):
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        r = Region(left=1.0, top=2.0, width=5.0, height=4.0)
        ins = r.inset(0.5)
        assert ins.left == 1.5
        assert ins.top == 2.5
        assert abs(ins.width - 4.0) < 0.01
        assert abs(ins.height - 3.0) < 0.01

    def test_frozen(self):
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        r = Region(left=0, top=0, width=1, height=1)
        with pytest.raises(AttributeError):
            r.left = 5

    def test_inset_no_negative(self):
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        r = Region(left=0, top=0, width=0.5, height=0.5)
        ins = r.inset(1.0)
        assert ins.width == 0
        assert ins.height == 0


class TestDiagramStyle:

    def test_default_values(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        s = DiagramStyle()
        assert s.node_font_size_pt == 13
        assert s.connector_width_pt == 1.5

    def test_resolve_color_default(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        s = DiagramStyle()
        assert s.resolve_color("primary") == "#2563EB"
        assert s.resolve_color("on-primary") == "#FFFFFF"

    def test_resolve_color_override(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        s = DiagramStyle(_color_map={"primary": "#FF0000"})
        assert s.resolve_color("primary") == "#FF0000"

    def test_resolve_color_unknown(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        s = DiagramStyle()
        assert s.resolve_color("nonexistent") == "#000000"

    def test_from_theme_dark(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        theme = {"colors": {"background": "#1A1A1A", "primary": "#60A5FA"}}
        s = DiagramStyle.from_theme(theme)
        assert s.node_fill == "muted"
        assert s.node_font_color == "foreground"

    def test_from_theme_education(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        theme = {"colors": {}}
        s = DiagramStyle.from_theme(theme, business_mode="education")
        assert s.node_font_size_pt == 12
        assert s.node_shadow is False

    def test_from_brand_spec(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        spec = BrandSpec(
            source="test",
            colors={"primary": "#1A3C6E"},
            spacing={"body_size_pt": 10, "margins_inches": 0.5},
        )
        s = DiagramStyle.from_brand_spec(spec)
        assert s.node_font_size_pt == 10
        assert abs(s.node_gap_inches - 0.15) < 0.01

    def test_apply_density_low(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        s = DiagramStyle()
        d = s.apply_density(2)
        assert d.node_font_size_pt == 16
        assert d.node_gap_inches == 0.4

    def test_apply_density_mid(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        s = DiagramStyle()
        d = s.apply_density(5)
        assert d.node_font_size_pt == 14
        assert d.node_gap_inches == 0.3

    def test_apply_density_high(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        s = DiagramStyle()
        d = s.apply_density(9)
        assert d.node_font_size_pt == 12
        assert d.node_gap_inches == 0.2

    def test_apply_density_preserves_colors(self):
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        s = DiagramStyle(_color_map={"primary": "#FF0000"})
        d = s.apply_density(5)
        assert d.resolve_color("primary") == "#FF0000"


class TestTextMeasurer:

    def test_empty_text(self):
        from ppt_pro_max.renderer.diagram.text_measurer import estimate_text_size
        w, h = estimate_text_size("", 12, 5.0)
        assert w == 0.0
        assert h == 0.0

    def test_short_text_single_line(self):
        from ppt_pro_max.renderer.diagram.text_measurer import estimate_text_size
        w, h = estimate_text_size("Hello", 12, 10.0)
        assert w > 0
        assert h > 0
        assert w < 10.0

    def test_long_text_wraps(self):
        from ppt_pro_max.renderer.diagram.text_measurer import estimate_text_size as est
        w_narrow, h_narrow = est("A" * 100, 12, 2.0)
        w_wide, h_wide = est("A" * 100, 12, 10.0)
        assert h_narrow > h_wide

    def test_cjk_wider_than_latin(self):
        from ppt_pro_max.renderer.diagram.text_measurer import estimate_text_size
        w_cjk, _ = estimate_text_size("你好世界", 12, 10.0)
        w_lat, _ = estimate_text_size("abcd", 12, 10.0)
        assert w_cjk > w_lat

    def test_estimate_node_size(self):
        from ppt_pro_max.renderer.diagram.text_measurer import estimate_node_size
        w, h = estimate_node_size("Test Label", 14, 3.0)
        assert w > 0
        assert h > 0

    def test_different_font_sizes(self):
        from ppt_pro_max.renderer.diagram.text_measurer import estimate_text_size
        w_small, _ = estimate_text_size("Hello", 10, 10.0)
        w_big, _ = estimate_text_size("Hello", 24, 10.0)
        assert w_big > w_small


class TestFlowchartDiagram:

    def test_horizontal_layout(self):
        from ppt_pro_max.renderer.diagram.flowchart import FlowchartDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"nodes": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}
        style = DiagramStyle()
        region = Region(0, 0, 10, 5)
        d = FlowchartDiagram(data, style, region)
        d.compute_layout()
        assert len(d._nodes) == 3
        assert len(d._connectors) == 2

    def test_vertical_layout_for_many_nodes(self):
        from ppt_pro_max.renderer.diagram.flowchart import FlowchartDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        nodes = [{"label": f"Step {i}"} for i in range(8)]
        data = {"nodes": nodes}
        style = DiagramStyle()
        region = Region(0, 0, 10, 7)
        d = FlowchartDiagram(data, style, region)
        d.compute_layout()
        assert len(d._nodes) == 8

    def test_render_horizontal(self):
        from ppt_pro_max.renderer.diagram.flowchart import FlowchartDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"nodes": [{"label": "A"}, {"label": "B"}], "direction": "horizontal"}
        style = DiagramStyle()
        region = Region(0.5, 0.5, 9.0, 4.0)
        d = FlowchartDiagram(data, style, region)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        d.render(slide)
        shapes = list(slide.shapes)
        assert len(shapes) >= 2

    def test_empty_nodes_no_crash(self):
        from ppt_pro_max.renderer.diagram.flowchart import FlowchartDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"nodes": []}
        style = DiagramStyle()
        region = Region(0, 0, 10, 5)
        d = FlowchartDiagram(data, style, region)
        d.compute_layout()
        assert len(d._nodes) == 0


class TestFunnelDiagram:

    def test_basic_funnel(self):
        from ppt_pro_max.renderer.diagram.funnel import FunnelDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"stages": [
            {"label": "Leads 10K"},
            {"label": "Qualified 3K"},
            {"label": "Proposal 800"},
            {"label": "Closed 200"},
        ]}
        style = DiagramStyle()
        region = Region(0, 0, 10, 5)
        d = FunnelDiagram(data, style, region)
        d.compute_layout()
        assert len(d._nodes) == 4
        first_w = d._nodes[0]["width"]
        last_w = d._nodes[-1]["width"]
        assert first_w > last_w

    def test_funnel_render(self):
        from ppt_pro_max.renderer.diagram.funnel import FunnelDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"stages": [{"label": "A"}, {"label": "B"}]}
        style = DiagramStyle()
        region = Region(0.5, 0.5, 9.0, 4.0)
        d = FunnelDiagram(data, style, region)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        d.render(slide)
        shapes = list(slide.shapes)
        assert len(shapes) >= 2

    def test_funnel_centered(self):
        from ppt_pro_max.renderer.diagram.funnel import FunnelDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"stages": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}
        style = DiagramStyle()
        region = Region(0, 0, 10, 6)
        d = FunnelDiagram(data, style, region)
        d.compute_layout()
        for node in d._nodes:
            node_center = node["x"] + node["width"] / 2
            assert abs(node_center - region.center_x) < 0.01


class TestTimelineDiagram:

    def test_basic_timeline(self):
        from ppt_pro_max.renderer.diagram.timeline import TimelineDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"events": [
            {"date": "2024 Q1", "label": "Start"},
            {"date": "2024 Q3", "label": "MVP"},
            {"date": "2025 Q1", "label": "Launch"},
        ]}
        style = DiagramStyle()
        region = Region(0, 0, 10, 5)
        d = TimelineDiagram(data, style, region)
        d.compute_layout()
        assert len(d._connectors) == 2

    def test_timeline_render(self):
        from ppt_pro_max.renderer.diagram.timeline import TimelineDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"events": [
            {"date": "2024", "label": "Start"},
            {"date": "2025", "label": "Growth"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 0.5, 9.0, 4.0)
        d = TimelineDiagram(data, style, region)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        d.render(slide)
        shapes = list(slide.shapes)
        assert len(shapes) >= 2

    def test_alternating_labels(self):
        from ppt_pro_max.renderer.diagram.timeline import TimelineDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"events": [
            {"date": "Q1", "label": "A"},
            {"date": "Q2", "label": "B"},
            {"date": "Q3", "label": "C"},
            {"date": "Q4", "label": "D"},
        ]}
        style = DiagramStyle()
        region = Region(0, 0, 10, 5)
        d = TimelineDiagram(data, style, region)
        d.compute_layout()
        label_nodes = [n for n in d._nodes if n.get("label")]
        above_y = [n for n in label_nodes if n["y"] < region.center_y]
        below_y = [n for n in label_nodes if n["y"] >= region.center_y]
        assert len(above_y) > 0
        assert len(below_y) > 0


class TestDiagramEngine:

    def test_registry(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        types = DiagramEngine.get_supported_types()
        assert "flowchart" in types
        assert "funnel" in types
        assert "timeline" in types

    def test_render_flowchart(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        engine = DiagramEngine()
        data = {"nodes": [{"label": "A"}, {"label": "B"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 4.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        engine.render(slide, "flowchart", data, style, region)
        assert len(list(slide.shapes)) >= 2

    def test_render_unsupported_as_text(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        engine = DiagramEngine()
        data = {"nodes": [{"label": "Item1"}, {"label": "Item2"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 4.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        engine.render(slide, "venn", data, style, region)
        shapes = list(slide.shapes)
        assert len(shapes) >= 1


class TestDataSplitter:

    def test_no_split_needed(self):
        from ppt_pro_max.renderer.diagram.data_splitter import split_data
        data = {"nodes": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}
        result = split_data(data, max_items_per_page=5)
        assert len(result) == 1

    def test_split_needed(self):
        from ppt_pro_max.renderer.diagram.data_splitter import split_data
        data = {"nodes": [{"label": str(i)} for i in range(7)]}
        result = split_data(data, max_items_per_page=3)
        assert len(result) == 3
        assert len(result[0]["nodes"]) == 3
        assert len(result[1]["nodes"]) == 3
        assert len(result[2]["nodes"]) == 1

    def test_split_stages(self):
        from ppt_pro_max.renderer.diagram.data_splitter import split_data
        data = {"stages": [{"label": str(i)} for i in range(5)]}
        result = split_data(data, max_items_per_page=3)
        assert len(result) == 2

    def test_split_events(self):
        from ppt_pro_max.renderer.diagram.data_splitter import split_data
        data = {"events": [{"date": str(i)} for i in range(6)]}
        result = split_data(data, max_items_per_page=2)
        assert len(result) == 3

    def test_empty_data(self):
        from ppt_pro_max.renderer.diagram.data_splitter import split_data
        data = {}
        result = split_data(data, max_items_per_page=3)
        assert len(result) == 1

    def test_page_index_metadata(self):
        from ppt_pro_max.renderer.diagram.data_splitter import split_data
        data = {"nodes": [{"label": str(i)} for i in range(6)]}
        result = split_data(data, max_items_per_page=3)
        assert result[0]["page_index"] == 1
        assert result[0]["total_pages"] == 2


class TestSwotDiagram:

    def test_render_four_quadrants(self):
        from ppt_pro_max.renderer.diagram.swot import SwotDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {
            "quadrants": [
                {"label": "Strengths", "items": ["Brand", "Tech"]},
                {"label": "Weaknesses", "items": ["Scale"]},
                {"label": "Opportunities", "items": ["Market"]},
                {"label": "Threats", "items": ["Competition"]},
            ]
        }
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = SwotDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) >= 4
        labels = [n["label"] for n in diagram._nodes]
        assert "Strengths" in labels

    def test_render_on_slide(self):
        from ppt_pro_max.renderer.diagram.swot import SwotDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"quadrants": [
            {"label": "S", "items": ["a"]},
            {"label": "W", "items": ["b"]},
            {"label": "O", "items": ["c"]},
            {"label": "T", "items": ["d"]},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        diagram = SwotDiagram(data, style, region)
        diagram.render(slide)
        assert len(slide.shapes) > 0

    def test_default_quadrants_when_empty(self):
        from ppt_pro_max.renderer.diagram.swot import SwotDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = SwotDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) >= 4


class TestMatrixDiagram:

    def test_render_matrix(self):
        from ppt_pro_max.renderer.diagram.matrix import MatrixDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {
            "rows": [{"label": "Row A"}, {"label": "Row B"}],
            "cols": [{"label": "Col 1"}, {"label": "Col 2"}],
            "cells": [["High", "Low"], ["Med", "High"]],
        }
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = MatrixDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) > 4

    def test_render_on_slide(self):
        from ppt_pro_max.renderer.diagram.matrix import MatrixDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {
            "rows": [{"label": "A"}],
            "cols": [{"label": "X"}, {"label": "Y"}],
            "cells": [["v1", "v2"]],
        }
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        diagram = MatrixDiagram(data, style, region)
        diagram.render(slide)
        assert len(slide.shapes) > 0

    def test_empty_data(self):
        from ppt_pro_max.renderer.diagram.matrix import MatrixDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = MatrixDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 0


class TestCycleDiagram:

    def test_render_cycle(self):
        from ppt_pro_max.renderer.diagram.cycle import CycleDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"nodes": [
            {"label": "Plan"}, {"label": "Do"}, {"label": "Check"}, {"label": "Act"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = CycleDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 4
        assert len(diagram._connectors) == 4

    def test_render_on_slide(self):
        from ppt_pro_max.renderer.diagram.cycle import CycleDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"nodes": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        diagram = CycleDiagram(data, style, region)
        diagram.render(slide)
        assert len(slide.shapes) > 0

    def test_too_few_nodes(self):
        from ppt_pro_max.renderer.diagram.cycle import CycleDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"nodes": [{"label": "Only"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = CycleDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 0


class TestTableDiagram:

    def test_render_table(self):
        from ppt_pro_max.renderer.diagram.table import TableDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {
            "headers": ["Name", "Score", "Grade"],
            "rows": [["Alice", "95", "A"], ["Bob", "82", "B"]],
        }
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = TableDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 9

    def test_render_on_slide(self):
        from ppt_pro_max.renderer.diagram.table import TableDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {
            "headers": ["A", "B"],
            "rows": [["1", "2"]],
        }
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        diagram = TableDiagram(data, style, region)
        diagram.render(slide)
        assert len(slide.shapes) > 0

    def test_auto_headers_from_rows(self):
        from ppt_pro_max.renderer.diagram.table import TableDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"rows": [["x", "y"], ["a", "b"]]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = TableDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) > 0


class TestDiagramEngineRegistry:

    def test_all_types_registered(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        supported = DiagramEngine.get_supported_types()
        for t in ("flowchart", "funnel", "timeline", "swot", "matrix", "cycle", "table"):
            assert t in supported, f"{t} not registered"

    def test_render_swot_via_engine(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        engine = DiagramEngine()
        data = {"quadrants": [
            {"label": "S"}, {"label": "W"}, {"label": "O"}, {"label": "T"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        engine.render(slide, "swot", data, style, region)
        assert len(slide.shapes) > 0

    def test_render_cycle_via_engine(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        engine = DiagramEngine()
        data = {"nodes": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        engine.render(slide, "cycle", data, style, region)
        assert len(slide.shapes) > 0


class TestCodeBlockRendering:

    def test_code_string_renders(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        design = {"title": "Demo", "code": "print('hello')"}
        pipeline._populate_slide(slide, design, prs)
        assert len(slide.shapes) > 1

    def test_code_dict_renders(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        design = {"title": "Demo", "code": {"code": "x = 1", "language": "python"}}
        pipeline._populate_slide(slide, design, prs)
        assert len(slide.shapes) > 1


class TestExerciseRendering:

    def test_exercise_dict_renders(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        design = {"title": "Practice", "exercise": {"instructions": "Write a function", "duration": "10 min"}}
        pipeline._populate_slide(slide, design, prs)
        assert len(slide.shapes) > 1

    def test_exercise_string_renders(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        design = {"title": "Practice", "exercise": "Do the task"}
        pipeline._populate_slide(slide, design, prs)
        assert len(slide.shapes) > 1


class TestHierarchyDiagram:

    def test_render_hierarchy(self):
        from ppt_pro_max.renderer.diagram.hierarchy import HierarchyDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"nodes": [
            {"id": "ceo", "label": "CEO", "level": 0},
            {"id": "cto", "label": "CTO", "level": 1, "parent": "ceo"},
            {"id": "cfo", "label": "CFO", "level": 1, "parent": "ceo"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = HierarchyDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 3
        assert len(diagram._connectors) == 2

    def test_render_on_slide(self):
        from ppt_pro_max.renderer.diagram.hierarchy import HierarchyDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"nodes": [
            {"id": "a", "label": "Root", "level": 0},
            {"id": "b", "label": "Child", "level": 1, "parent": "a"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        diagram = HierarchyDiagram(data, style, region)
        diagram.render(slide)
        assert len(slide.shapes) > 0

    def test_empty_data(self):
        from ppt_pro_max.renderer.diagram.hierarchy import HierarchyDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = HierarchyDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 0


class TestPyramidDiagram:

    def test_render_pyramid(self):
        from ppt_pro_max.renderer.diagram.pyramid import PyramidDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"levels": [
            {"label": "Top"},
            {"label": "Middle"},
            {"label": "Base"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = PyramidDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 3
        assert len(diagram._connectors) == 2

    def test_render_on_slide(self):
        from ppt_pro_max.renderer.diagram.pyramid import PyramidDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"levels": [{"label": "A"}, {"label": "B"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        diagram = PyramidDiagram(data, style, region)
        diagram.render(slide)
        assert len(slide.shapes) > 0


class TestVennDiagram:

    def test_render_two_sets(self):
        from ppt_pro_max.renderer.diagram.venn import VennDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"sets": [
            {"label": "A"}, {"label": "B"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = VennDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 2

    def test_render_three_sets(self):
        from ppt_pro_max.renderer.diagram.venn import VennDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"sets": [
            {"label": "A"}, {"label": "B"}, {"label": "C"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = VennDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 3

    def test_render_on_slide(self):
        from ppt_pro_max.renderer.diagram.venn import VennDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"sets": [{"label": "X"}, {"label": "Y"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        diagram = VennDiagram(data, style, region)
        diagram.render(slide)
        assert len(slide.shapes) > 0

    def test_invalid_set_count(self):
        from ppt_pro_max.renderer.diagram.venn import VennDiagram
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        data = {"sets": [{"label": "Only"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        diagram = VennDiagram(data, style, region)
        diagram.compute_layout()
        assert len(diagram._nodes) == 0


class TestDiagramEngineFullRegistry:

    def test_all_phase_h_types_registered(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        supported = DiagramEngine.get_supported_types()
        for t in ("hierarchy", "pyramid", "venn"):
            assert t in supported, f"{t} not registered"

    def test_render_hierarchy_via_engine(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        engine = DiagramEngine()
        data = {"nodes": [
            {"id": "a", "label": "Root", "level": 0},
            {"id": "b", "label": "Child", "level": 1, "parent": "a"},
        ]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        engine.render(slide, "hierarchy", data, style, region)
        assert len(slide.shapes) > 0

    def test_render_pyramid_via_engine(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        engine = DiagramEngine()
        data = {"levels": [{"label": "Top"}, {"label": "Mid"}, {"label": "Base"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        engine.render(slide, "pyramid", data, style, region)
        assert len(slide.shapes) > 0

    def test_render_venn_via_engine(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        engine = DiagramEngine()
        data = {"sets": [{"label": "A"}, {"label": "B"}, {"label": "C"}]}
        style = DiagramStyle()
        region = Region(0.5, 1.0, 9.0, 5.0)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        engine.render(slide, "venn", data, style, region)
        assert len(slide.shapes) > 0
