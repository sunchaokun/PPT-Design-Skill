"""End-to-end smoke tests — full enterprise pipeline with all features."""

from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation


def _make_png(path: Path, w: int = 200, h: int = 150, color: str = "blue") -> Path:
    img = Image.new("RGB", (w, h), color=color)
    img.save(str(path))
    return path


def _setup_full_project(tmp_path: Path) -> Path:
    project = tmp_path / "acme_pitch"
    project.mkdir()

    prs = Presentation()
    for i in range(5):
        layout = prs.slide_layouts[min(i, len(prs.slide_layouts) - 1)]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Template Slide {i + 1}"
    prs.save(str(project / "template.pptx"))

    _make_png(project / "logo.png", 200, 100, "navy")
    _make_png(project / "hero.png", 800, 400, "teal")
    _make_png(project / "dashboard.png", 600, 400, "green")
    _make_png(project / "team_photo.jpg", 400, 300, "orange")
    _make_png(project / "market_chart.png", 500, 350, "purple")

    brand = {
        "colors": {
            "primary": "#1A3C6E",
            "accent": "#C9A227",
            "background": "#FFFFFF",
            "foreground": "#1A1A1A",
        },
        "fonts": {
            "heading": "Georgia",
            "body": "Calibri",
        },
        "logo": {
            "position": "top_right",
            "width_inches": 0.8,
            "skip_slides": ["hook"],
        },
    }
    (project / "brand.json").write_text(json.dumps(brand), encoding="utf-8")

    content = {
        "meta": {
            "title": "ACME Corp — Series A Pitch",
            "subtitle": "Raising $5M to Revolutionize B2B SaaS",
            "author": "Jane Doe",
        },
        "slides": [
            {"goal": "hook", "title": "ACME Corp", "subtitle": "Series A Pitch Deck", "image": "hero.png"},
            {"goal": "problem", "title": "The Problem", "bullets": [
                "Enterprise teams waste 40% of time on manual data entry",
                "Existing solutions cost $50K+/year with poor ROI",
                "No unified platform for cross-department workflows",
            ]},
            {"goal": "solution", "title": "Our Solution", "bullets": [
                "AI-powered workflow automation",
                "Real-time cross-department sync",
                "90% reduction in manual tasks",
            ]},
            {"goal": "product", "title": "Product Demo", "image": "dashboard.png"},
            {"goal": "team", "title": "Our Team", "bullets": [
                "CEO: 15 years enterprise SaaS",
                "CTO: Ex-Google, ML specialist",
                "VP Sales: $100M+ career quota",
            ]},
        ],
    }
    (project / "content.json").write_text(json.dumps(content), encoding="utf-8")

    return project


class TestE2EFullPipeline:

    def test_full_render_all_features(self, tmp_path):
        """Full pipeline: template + brand.json + content.json + logo + image pool."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="ACME Corp Series A",
            project_dir=str(project),
            business_mode="pitch",
            density=5,
        )
        assert result["pipeline"] == "enterprise"
        assert result["num_slides"] == 5
        assert os.path.isfile(result["output_path"])

        prs = Presentation(result["output_path"])
        assert len(prs.slides) == 5

        assert prs.slides[0].shapes.title.text == "ACME Corp"
        assert prs.slides[1].shapes.title.text == "The Problem"
        assert prs.slides[3].shapes.title.text == "Product Demo"

        logo_count = sum(
            1 for slide in prs.slides
            for shape in slide.shapes
            if shape.shape_type == 13
        )
        assert logo_count >= 3

    def test_dry_run_with_all_assets(self, tmp_path):
        """Dry run should report all assets found."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="ACME Corp",
            project_dir=str(project),
            dry_run=True,
        )
        assert result["dry_run"] is True
        assert result["assets"]["template"] is True
        assert result["assets"]["logo"] is True
        assert result["assets"]["brand_json"] is True
        assert result["assets"]["content_json"] is True
        assert result["assets"]["image_pool_count"] >= 3

    def test_review_mode(self, tmp_path):
        """Review mode should produce proposal with all pages."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="ACME Corp",
            project_dir=str(project),
            review=True,
            density=5,
        )
        assert result["review"] is True
        proposal = result["proposal"]
        assert len(proposal["pages"]) == 5

    def test_page_revision_delete(self, tmp_path):
        """Page revision: delete slides 2,4 from 5-slide template."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="Revised",
            project_dir=str(project),
            pages="-2,-4",
        )
        assert result["mode"] == "page_revision"
        assert result["num_slides"] == 3
        prs = Presentation(result["output_path"])
        assert prs.slides[0].shapes.title.text == "Template Slide 1"
        assert prs.slides[1].shapes.title.text == "Template Slide 3"
        assert prs.slides[2].shapes.title.text == "Template Slide 5"

    def test_page_revision_swap(self, tmp_path):
        """Page revision: swap first and last slides."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="Reordered",
            project_dir=str(project),
            pages="1<>5",
        )
        assert result["num_slides"] == 5
        prs = Presentation(result["output_path"])
        assert prs.slides[0].shapes.title.text == "Template Slide 5"
        assert prs.slides[4].shapes.title.text == "Template Slide 1"

    def test_version_meta_written(self, tmp_path):
        """Version meta.json should contain slides[] with goals."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="ACME Corp",
            project_dir=str(project),
            density=5,
        )
        vdir = os.path.join(str(project), "output", f"v{result['version']}")
        meta = json.loads(Path(os.path.join(vdir, "meta.json")).read_text(encoding="utf-8"))
        assert len(meta["slides"]) == 5
        assert meta["slides"][0]["goal"] == "hook"
        assert meta["slides"][3]["goal"] == "product"

    def test_history(self, tmp_path):
        """History should list generated versions."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        pipeline.run(query="V1", project_dir=str(project))
        pipeline.run(query="V2", project_dir=str(project))
        result = pipeline.run(history=True, project_dir=str(project), query="")
        assert len(result["versions"]) == 2

    def test_image_pool_auto_assign(self, tmp_path):
        """Image pool images should auto-assign to slides by goal matching."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _setup_full_project(tmp_path)
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="ACME Corp",
            project_dir=str(project),
            density=5,
        )
        prs = Presentation(result["output_path"])
        slide0_pics = [s for s in prs.slides[0].shapes if s.shape_type == 13]
        assert len(slide0_pics) >= 1

    def test_minimal_project_no_assets(self, tmp_path):
        """Project with no assets should still produce a PPTX."""
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = tmp_path / "empty_project"
        project.mkdir()
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="Simple Presentation",
            project_dir=str(project),
        )
        assert result["pipeline"] == "enterprise"
        assert os.path.isfile(result["output_path"])
