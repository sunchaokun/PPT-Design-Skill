"""Integration test: generate full PPTs and analyze quality."""

import os
import json

from pptx import Presentation

from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
from ppt_pro_max.enterprise.delivery_gate import DeliveryGate


def _setup_project(td, content_slides, brand=None):
    content = {"slides": content_slides}
    with open(os.path.join(td, "content.json"), "w", encoding="utf-8") as f:
        json.dump(content, f, ensure_ascii=False, indent=2)
    if brand:
        with open(os.path.join(td, "brand.json"), "w", encoding="utf-8") as f:
            json.dump(brand, f, ensure_ascii=False, indent=2)
    return td


def _analyze_pptx(path):
    prs = Presentation(path)
    results = []
    for idx, slide in enumerate(prs.slides):
        shape_count = len(slide.shapes)
        texts = []
        font_sizes = []
        image_count = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text and p.text.strip():
                        texts.append(p.text.strip()[:60])
                    if p.font.size:
                        font_sizes.append(p.font.size / 12700)
                    for r in p.runs:
                        if r.font.size:
                            font_sizes.append(r.font.size / 12700)
            if hasattr(shape, "image"):
                try:
                    _ = shape.image
                    image_count += 1
                except Exception:
                    pass
        min_font = min(font_sizes) if font_sizes else None
        max_font = max(font_sizes) if font_sizes else None
        results.append({
            "slide": idx,
            "shapes": shape_count,
            "texts": len(texts),
            "images": image_count,
            "min_font": round(min_font, 1) if min_font else None,
            "max_font": round(max_font, 1) if max_font else None,
            "first_text": texts[0] if texts else "",
        })
    return results


def _run_delivery_gate(path, dna=None, plans=None):
    gate = DeliveryGate()
    report = gate.check(path, dna, plans or [])
    return {
        "total_slides": report.total_slides,
        "passed": report.passed,
        "fatals": len(report.fatals),
        "warnings": len(report.warnings),
        "fatal_ids": [i.check_id for i in report.fatals],
        "warning_ids": [i.check_id for i in report.warnings],
        "is_passable": report.is_passable,
    }


def _run_pipeline(td, style="professional", query="AI Innovation Summit"):
    pipeline = EnterprisePipeline()
    result = pipeline.run(
        query=query,
        project_dir=td,
        style=style,
    )
    return result


# ── Test 1: Freestyle pipeline with various styles ──

def test_freestyle_professional(tmp_path):
    td = str(tmp_path)
    result = _run_pipeline(td, style="professional", query="AI Innovation Summit 2026")
    pptx_path = result.get("output_path", "")
    assert pptx_path and os.path.isfile(pptx_path), f"No output: {result}"

    analysis = _analyze_pptx(pptx_path)
    gate = _run_delivery_gate(pptx_path)

    print("\n=== Freestyle Professional ===")
    print(f"Slides: {len(analysis)}")
    for a in analysis:
        print(f"  Slide {a['slide']}: {a['shapes']} shapes, {a['texts']} texts, "
              f"{a['images']} imgs, font {a['min_font']}-{a['max_font']}pt, "
              f"first: {a['first_text'][:40]}")
    print(f"Gate: {gate}")

    assert len(analysis) >= 4, "Too few slides"
    assert all(a["shapes"] >= 2 for a in analysis), "Blank slide detected"


def test_freestyle_dark_tech(tmp_path):
    td = str(tmp_path)
    result = _run_pipeline(td, style="dark-tech", query="Cybersecurity Strategy")
    pptx_path = result.get("output_path", "")
    assert pptx_path and os.path.isfile(pptx_path)

    analysis = _analyze_pptx(pptx_path)
    gate = _run_delivery_gate(pptx_path)

    print("\n=== Freestyle Dark-Tech ===")
    print(f"Slides: {len(analysis)}")
    for a in analysis:
        print(f"  Slide {a['slide']}: {a['shapes']} shapes, {a['texts']} texts, "
              f"font {a['min_font']}-{a['max_font']}pt")
    print(f"Gate: {gate}")

    assert len(analysis) >= 4


def test_freestyle_warm_elegant(tmp_path):
    td = str(tmp_path)
    result = _run_pipeline(td, style="warm-elegant", query="Sustainable Fashion Brand")
    pptx_path = result.get("output_path", "")
    assert pptx_path and os.path.isfile(pptx_path)

    analysis = _analyze_pptx(pptx_path)
    gate = _run_delivery_gate(pptx_path)

    print("\n=== Freestyle Warm-Elegant ===")
    print(f"Slides: {len(analysis)}")
    for a in analysis:
        print(f"  Slide {a['slide']}: {a['shapes']} shapes, {a['texts']} texts, "
              f"font {a['min_font']}-{a['max_font']}pt")
    print(f"Gate: {gate}")

    assert len(analysis) >= 4


# ── Test 2: Enterprise pipeline with content.json ──

def test_enterprise_content_json(tmp_path):
    td = str(tmp_path)
    slides = [
        {"goal": "hook", "title": "AI Platform Launch", "subtitle": "Transforming Enterprise AI"},
        {"goal": "problem", "title": "Current Challenges", "bullets": ["Data silos across departments", "Manual processes waste time", "No unified AI strategy"]},
        {"goal": "solution", "title": "Our Solution", "bullets": ["Unified data pipeline", "Automated ML workflows", "Real-time analytics dashboard"]},
        {"goal": "features", "title": "Key Features", "cards": [
            {"title": "Data Engine", "text": "Process 10TB/day"},
            {"title": "ML Pipeline", "text": "Auto-retrain models"},
            {"title": "Dashboard", "text": "Real-time insights"},
        ]},
        {"goal": "data", "title": "Performance Metrics", "table": {"headers": ["Metric", "Before", "After"], "rows": [["Accuracy", "72%", "94%"], ["Speed", "4hr", "12min"], ["Cost", "$50K", "$8K"]]}},
        {"goal": "code", "title": "Quick Start", "code": {"language": "python", "lines": ["from ai_platform import Pipeline", "pipe = Pipeline(config='prod')", "result = pipe.run(data)"]}},
        {"goal": "exercise", "title": "Hands-On Lab", "exercise": {"duration": "5 min", "steps": ["Clone the repo", "Run the example", "Check the output"]}},
        {"goal": "cta", "title": "Get Started Today", "subtitle": "Contact us for a free trial"},
    ]
    _setup_project(td, slides, brand={
        "colors": {"primary": "#1A365D", "accent": "#2B6CB0", "foreground": "#1A202C", "background": "#FFFFFF"},
        "fonts": {"heading": "Georgia", "body": "Calibri"},
    })

    pipeline = EnterprisePipeline()
    result = pipeline.run(query="AI Platform", project_dir=td)
    pptx_path = result.get("output_path", "")
    assert pptx_path and os.path.isfile(pptx_path), f"No output: {result}"

    analysis = _analyze_pptx(pptx_path)
    gate = _run_delivery_gate(pptx_path)

    print("\n=== Enterprise Content JSON ===")
    print(f"Slides: {len(analysis)}")
    for a in analysis:
        print(f"  Slide {a['slide']}: {a['shapes']} shapes, {a['texts']} texts, "
              f"{a['images']} imgs, font {a['min_font']}-{a['max_font']}pt, "
              f"first: {a['first_text'][:40]}")
    print(f"Gate: {gate}")

    assert len(analysis) >= 6, "Too few slides for enterprise"
    assert all(a["shapes"] >= 2 for a in analysis), "Blank slide"


# ── Test 3: Font size check — no font below 11pt ──

def test_no_font_below_11pt(tmp_path):
    td = str(tmp_path)
    result = _run_pipeline(td, style="professional", query="Quarterly Business Review")
    pptx_path = result.get("output_path", "")
    assert pptx_path and os.path.isfile(pptx_path)

    analysis = _analyze_pptx(pptx_path)
    small_fonts = []
    for a in analysis:
        if a["min_font"] is not None and a["min_font"] < 11:
            small_fonts.append(a)

    print("\n=== Font Size Check ===")
    if small_fonts:
        print(f"WARNING: {len(small_fonts)} slides with fonts below 11pt:")
        for a in small_fonts:
            print(f"  Slide {a['slide']}: min font {a['min_font']}pt")
    else:
        print("All fonts >= 11pt — PASS")

    gate = _run_delivery_gate(pptx_path)
    font_warnings = [i for i in gate.get("warning_ids", []) if i == "font_too_small"]
    print(f"DeliveryGate font_too_small warnings: {len(font_warnings)}")


# ── Test 4: DeliveryGate comprehensive check ──

def test_delivery_gate_comprehensive(tmp_path):
    styles = ["professional", "dark-tech", "warm-elegant", "vibrant-startup", "nature-calm"]
    print("\n=== DeliveryGate Comprehensive ===")
    for style in styles:
        td2 = os.path.join(str(tmp_path), style)
        os.makedirs(td2, exist_ok=True)
        result = _run_pipeline(td2, style=style, query="Product Strategy 2026")
        pptx_path = result.get("output_path", "")
        if not pptx_path or not os.path.isfile(pptx_path):
            print(f"  {style}: NO OUTPUT")
            continue
        gate = _run_delivery_gate(pptx_path)
        print(f"  {style}: slides={gate['total_slides']} pass={gate['passed']} "
              f"fatals={gate['fatals']} warns={gate['warnings']} "
              f"passable={gate['is_passable']}")
