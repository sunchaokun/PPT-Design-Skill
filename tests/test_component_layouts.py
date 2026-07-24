"""Test 3 layout styles: centered, left+right, left+right+centered.

Creates real PPT files with components from the library.
Run: python tests/test_component_layouts.py
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
from ppt_pro_max.enterprise.brand_spec import BrandSpec

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "_layout_test_output")


def _setup():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    db_path = find_db_path()
    if not db_path:
        print("ERROR: component_library/index.db not found. Run build_library first.")
        sys.exit(1)
    return ComponentLibrary(db_path)


def _make_brand():
    return BrandSpec(
        source="test",
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "accent": "#F97316",
            "foreground": "#1E293B",
            "muted-foreground": "#64748B",
            "background": "#FFFFFF",
            "muted": "#F1F5F9",
            "border": "#E2E8F0",
        },
        fonts={"heading": "Calibri", "body": "Calibri"},
    )


def _add_title(prs, title, subtitle=None):
    from ppt_pro_max.renderer.layout_registry import SLIDE_WIDTH, SLIDE_HEIGHT
    precision = PrecisionRenderer(brand_spec=_make_brand())
    slide = precision.add_slide(prs)
    precision.apply_brand_background(slide, prs)
    precision.add_text(slide, title, 0.9, 0.5, 11, 0.8, size=36, bold=True)
    precision.add_rect(slide, 0.9, 1.3, 2, 0.04, fill_role="accent")
    if subtitle:
        precision.add_text(slide, subtitle, 0.9, 1.5, 10, 0.5,
                           font="Calibri", size=14, color_role="muted-foreground")
    return slide, precision


def _inject_component(slide, precision, lib, category, bounds, fit="contain", texts=None):
    results = lib.search(type="group", category=category, limit=1)
    if not results:
        print(f"  SKIP: no {category} component found")
        return False

    comp = results[0]
    xml_parts = lib.load_xml(comp["id"])
    if not xml_parts or "group" not in xml_parts:
        print(f"  SKIP: {category} component has no XML")
        return False

    from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
    renderer = ComponentRenderer()

    orig_aspect = 1.0
    meta_bytes = xml_parts.get("_meta")
    if meta_bytes:
        try:
            import json
            meta = json.loads(meta_bytes) if isinstance(meta_bytes, bytes) else json.loads(meta_bytes)
            obe = meta.get("orig_bounds_emu")
            if obe and len(obe) == 4 and obe[3] > 0:
                orig_aspect = obe[2] / obe[3]
        except Exception:
            pass

    content_area = tuple(float(v) for v in bounds)
    actual_bounds = renderer.compute_component_bounds(content_area, orig_aspect, fit)

    element = {
        "type": "group",
        "category": category,
        "texts": texts or [],
        "nodes": [{"text": t} for t in (texts or [])],
        "bounds": actual_bounds,
    }

    success = renderer.render(slide, element, _make_brand(), lib)
    if success:
        print(f"  OK: {category} component injected at bounds={actual_bounds}")
    else:
        print(f"  FAIL: {category} component injection failed")
    return success


def _save(prs, name):
    path = os.path.join(OUTPUT_DIR, f"{name}.pptx")
    prs.save(path)
    n_shapes = sum(len(s.shapes) for s in prs.slides)
    print(f"  Saved: {path} ({len(prs.slides)} slides, {n_shapes} shapes)")
    return path


def test_layout_centered(lib):
    """Layout 1: 组件居中占满内容区域"""
    print("\n=== Layout 1: 居中 (Centered) ===")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide, precision = _add_title(prs, "项目流程", "组件居中布局 — contain 模式")

    bounds = precision._compute_content_area({
        "goal": "content",
        "title": "项目流程",
        "subtitle": "组件居中布局",
    })
    print(f"  Auto content area: {bounds}")

    _inject_component(slide, precision, lib, "process", bounds, fit="contain",
                      texts=["需求分析", "方案设计", "开发实现", "测试上线"])

    _save(prs, "layout1_centered")


def test_layout_left_right(lib):
    """Layout 2: 左侧组件 + 右侧文字说明"""
    print("\n=== Layout 2: 左侧+右侧 (Left Component + Right Text) ===")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide, precision = _add_title(prs, "组织架构", "左侧组件 + 右侧说明")

    left_bounds = (0.9, 1.6, 6.0, 5.0)
    _inject_component(slide, precision, lib, "hierarchy", left_bounds, fit="contain",
                      texts=["CEO", "CTO", "CFO", "VP Engineering", "VP Sales"])

    precision.add_text(slide, "团队结构", 7.5, 1.6, 5.0, 0.5, size=20, bold=True)
    precision.add_rect(slide, 7.5, 2.2, 1.5, 0.04, fill_role="accent")
    precision.add_multiline(slide, [
        "•  CEO 统筹全局战略",
        "•  CTO 负责技术方向",
        "•  CFO 管理财务运营",
        "•  VP Engineering 带领研发",
        "•  VP Sales 拓展市场",
    ], 7.5, 2.5, 5.0, 4.0, size=13, color_role="foreground", spacing=8)

    _save(prs, "layout2_left_right")


def test_layout_left_right_centered(lib):
    """Layout 3: 左侧文字 + 右上组件 + 右下组件"""
    print("\n=== Layout 3: 左侧+右上+右下 (Left Text + Right Top + Right Bottom) ===")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide, precision = _add_title(prs, "业务分析", "左侧说明 + 右上流程 + 右下SWOT")

    precision.add_text(slide, "核心业务", 0.9, 1.6, 5.0, 0.5, size=20, bold=True)
    precision.add_rect(slide, 0.9, 2.2, 1.5, 0.04, fill_role="accent")
    precision.add_multiline(slide, [
        "•  战略规划与执行",
        "•  流程优化与监控",
        "•  风险评估与应对",
        "•  资源分配与调度",
        "•  绩效考核与改进",
        "",
        "通过系统化的业务分析，",
        "实现组织目标的高效达成。",
    ], 0.9, 2.5, 5.0, 4.0, size=13, color_role="foreground", spacing=6)

    right_top_bounds = (6.5, 1.6, 6.0, 2.3)
    _inject_component(slide, precision, lib, "process", right_top_bounds, fit="width",
                      texts=["规划", "执行", "监控", "优化"])

    right_bottom_bounds = (6.5, 4.2, 6.0, 2.8)
    _inject_component(slide, precision, lib, "swot", right_bottom_bounds, fit="contain",
                      texts=["技术优势", "市场劣势", "增长机会", "潜在威胁"])

    _save(prs, "layout3_left_right_centered")


def test_layout_quad(lib):
    """Layout 4: 四象限 — 2x2 组件网格"""
    print("\n=== Layout 4: 四象限 (2x2 Grid) ===")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide, precision = _add_title(prs, "全面分析", "四象限组件布局")

    half_w = 5.8
    half_h = 2.4
    gap = 0.4

    tl = (0.9, 1.6, half_w, half_h)
    _inject_component(slide, precision, lib, "process", tl, fit="contain",
                      texts=["步骤1", "步骤2", "步骤3"])

    tr = (0.9 + half_w + gap, 1.6, half_w, half_h)
    _inject_component(slide, precision, lib, "hierarchy", tr, fit="contain",
                      texts=["CEO", "CTO", "CFO"])

    bl = (0.9, 1.6 + half_h + gap, half_w, half_h)
    _inject_component(slide, precision, lib, "infographic", bl, fit="contain",
                      texts=["收入", "利润", "增长", "用户"])

    br = (0.9 + half_w + gap, 1.6 + half_h + gap, half_w, half_h)
    _inject_component(slide, precision, lib, "timeline", br, fit="contain",
                      texts=["Q1", "Q2", "Q3", "Q4"])

    _save(prs, "layout4_quad")


def main():
    lib = _setup()
    print(f"Library: {lib.stats()}")

    test_layout_centered(lib)
    test_layout_left_right(lib)
    test_layout_left_right_centered(lib)
    test_layout_quad(lib)

    lib.close()

    print(f"\n=== ALL DONE ===")
    print(f"Output directory: {OUTPUT_DIR}")
    for f in sorted(os.listdir(OUTPUT_DIR)):
        if f.endswith(".pptx"):
            size = os.path.getsize(os.path.join(OUTPUT_DIR, f))
            print(f"  {f} ({size // 1024} KB)")


if __name__ == "__main__":
    main()
