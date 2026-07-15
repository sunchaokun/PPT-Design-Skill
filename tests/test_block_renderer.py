"""Tests for BlockRenderer — composable blocks layout system.

Covers:
  1. Region preset resolution (string presets, custom dict, auto-compute)
  2. Core block types: cards, bullets, diagram, code, exercise, image
  3. New block types: hexagons, ovals, donuts, metrics, badge, gradient_line, table_chart
  4. Multiple blocks on one slide (composable)
  5. Section divider bug fix (no blank slide)
  6. render_slide() blocks path vs legacy path
  7. Backward compatibility (no blocks key → old elif chain)
  8. content.json blocks passthrough via content_parser
  9. Pipeline integration (blocks from content.json reach render_slide)
"""

from __future__ import annotations

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer


def _make_brand() -> BrandSpec:
    return BrandSpec(
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "accent": "#6366F1",
            "background": "#0A0F1E",
            "foreground": "#F0F4F8",
            "muted": "#1A2238",
            "muted-foreground": "#94A3B8",
            "border": "#1A2A4A",
        },
        fonts={"heading": "Inter", "body": "Inter"},
    )


def _make_light_brand() -> BrandSpec:
    return BrandSpec(
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "accent": "#F97316",
            "background": "#FFFFFF",
            "foreground": "#1E293B",
            "muted": "#F1F5F9",
            "muted-foreground": "#64748B",
            "border": "#E2E8F0",
        },
        fonts={"heading": "Calibri", "body": "Calibri"},
    )


class TestRegionPresets:
    """BlockRenderer resolves region specs correctly."""

    def test_string_preset_full(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        r = br._resolve_region({"region": "full"}, 0, 1, False)
        assert r["x"] == pytest.approx(0.9, abs=0.01)
        assert r["w"] == pytest.approx(11.533, abs=0.01)
        assert r["y"] == pytest.approx(1.6, abs=0.01)
        assert r["h"] == pytest.approx(5.4, abs=0.01)

    def test_string_preset_left_right_gap(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        left = br._resolve_region({"region": "left"}, 0, 2, False)
        right = br._resolve_region({"region": "right"}, 1, 2, False)
        gap = right["x"] - (left["x"] + left["w"])
        assert gap == pytest.approx(0.4, abs=0.01)

    def test_string_preset_top_bottom_gap(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        top = br._resolve_region({"region": "top"}, 0, 2, False)
        bottom = br._resolve_region({"region": "bottom"}, 1, 2, False)
        gap = bottom["y"] - (top["y"] + top["h"])
        assert gap == pytest.approx(0.4, abs=0.01)

    def test_custom_dict_region(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        custom = {"x": 1.2, "y": 2.0, "w": 5.0, "h": 4.5}
        r = br._resolve_region({"region": custom}, 0, 1, False)
        assert r == custom

    def test_auto_region_single_block(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        r = br._auto_region(0, 1, False)
        assert r["x"] == pytest.approx(0.9, abs=0.01)
        assert r["w"] == pytest.approx(11.533, abs=0.01)

    def test_auto_region_two_blocks(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        r0 = br._auto_region(0, 2, False)
        r1 = br._auto_region(1, 2, False)
        assert r0["x"] < r1["x"]
        assert r0["w"] == pytest.approx(r1["w"], abs=0.01)

    def test_auto_region_four_blocks_grid(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        regions = [br._auto_region(i, 4, False) for i in range(4)]
        assert regions[0]["y"] < regions[2]["y"]  # top-left above bottom-left
        assert regions[0]["x"] < regions[1]["x"]  # top-left left of top-right

    def test_all_presets_within_slide_bounds(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        for name, r in BlockRenderer.REGION_PRESETS.items():
            assert r["x"] >= 0, f"Preset {name}: x<0"
            assert r["y"] >= 0, f"Preset {name}: y<0"
            assert r["x"] + r["w"] <= 13.34, f"Preset {name}: right edge exceeds slide"
            assert r["y"] + r["h"] <= 7.51, f"Preset {name}: bottom edge exceeds slide"

    def test_full_region_content_area(self):
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(PrecisionRenderer(brand_spec=_make_brand()))
        r = br._resolve_region({"region": "full"}, 0, 1, False)
        assert r["x"] + r["w"] == pytest.approx(13.333 - 0.9, abs=0.01)
        assert r["y"] + r["h"] == pytest.approx(7.5 - 0.5, abs=0.01)


class TestCoreBlockTypes:
    """Core block types render shapes onto slides."""

    def _render_blocks(self, blocks, brand=None):
        pr = PrecisionRenderer(brand_spec=brand or _make_brand())
        prs = pr.create_presentation()
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(pr)
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        br.render(slide, blocks)
        return prs, slide

    def test_cards_block_creates_shapes(self):
        prs, slide = self._render_blocks([{
            "type": "cards",
            "region": "full",
            "data": {"items": [
                {"title": "Card 1", "text": "Body 1"},
                {"title": "Card 2", "text": "Body 2"},
                {"title": "Card 3", "text": "Body 3"},
            ]}
        }])
        assert len(list(slide.shapes)) >= 6  # 3 cards × (rect + title + body)

    def test_bullets_block_creates_text(self):
        prs, slide = self._render_blocks([{
            "type": "bullets",
            "region": "full",
            "data": {"items": ["Point 1", "Point 2", "Point 3"]}
        }])
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("Point" in t for t in texts)

    def test_code_block_creates_rounded_rect(self):
        prs, slide = self._render_blocks([{
            "type": "code",
            "region": "full",
            "data": {"language": "python", "source": "print('hello')"}
        }])
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("PYTHON" in t.upper() for t in texts)
        assert any("print" in t for t in texts)

    def test_exercise_block_creates_badge(self):
        prs, slide = self._render_blocks([{
            "type": "exercise",
            "region": "full",
            "data": {
                "instructions": "Try it",
                "duration": "5 min",
                "steps": ["Step 1", "Step 2"]
            }
        }])
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("EXERCISE" in t.upper() or "5 MIN" in t.upper() for t in texts)

    def test_diagram_block_creates_shapes(self):
        prs, slide = self._render_blocks([{
            "type": "diagram",
            "region": "full",
            "data": {
                "diagram_type": "flowchart",
                "diagram_data": {
                    "nodes": [{"id": "1", "label": "A"}, {"id": "2", "label": "B"}],
                    "connectors": [{"from": "1", "to": "2"}]
                }
            }
        }])
        assert len(list(slide.shapes)) >= 2

    def test_badge_block_creates_shape(self):
        prs, slide = self._render_blocks([{
            "type": "badge",
            "region": {"x": 0.9, "y": 1.2, "w": 1.5, "h": 0.35},
            "data": {"text": "NEW", "variant": "solid"}
        }])
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("NEW" in t.upper() for t in texts)

    def test_gradient_line_block_creates_shape(self):
        prs, slide = self._render_blocks([{
            "type": "gradient_line",
            "region": {"x": 0.9, "y": 1.15, "w": 3.0, "h": 0.04},
            "data": {"color_role": "accent"}
        }])
        assert len(list(slide.shapes)) >= 1

    def test_metrics_block_creates_shapes(self):
        prs, slide = self._render_blocks([{
            "type": "metrics",
            "region": "top",
            "data": {"items": [
                {"value": "40K+", "label": "Styles"},
                {"value": "28", "label": "Upgrades"},
            ]}
        }])
        assert len(list(slide.shapes)) >= 2
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("40K" in t for t in texts)


class TestNewBlockTypes:
    """New visual block types: hexagons, ovals, donuts."""

    def _render_blocks(self, blocks, brand=None):
        pr = PrecisionRenderer(brand_spec=brand or _make_brand())
        prs = pr.create_presentation()
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(pr)
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        br.render(slide, blocks)
        return prs, slide

    def test_hexagons_block_creates_shapes(self):
        prs, slide = self._render_blocks([{
            "type": "hexagons",
            "region": "left-2-3",
            "data": {"items": [
                {"label": "Plan"}, {"label": "Design"},
                {"label": "Build"}, {"label": "Ship"}
            ]}
        }])
        assert len(list(slide.shapes)) >= 4

    def test_ovals_block_creates_shapes(self):
        prs, slide = self._render_blocks([{
            "type": "ovals",
            "region": "full",
            "data": {"items": [
                {"label": "01", "subtitle": "Plan"},
                {"label": "02", "subtitle": "Design"},
            ]}
        }])
        assert len(list(slide.shapes)) >= 2

    def test_donuts_block_creates_shapes(self):
        prs, slide = self._render_blocks([{
            "type": "donuts",
            "region": "right-1-3",
            "data": {"items": [
                {"label": "40K+", "subtitle": "Styles"},
                {"label": "28", "subtitle": "Upgrades"},
            ]}
        }])
        assert len(list(slide.shapes)) >= 2

    def test_table_chart_block_creates_content(self):
        prs, slide = self._render_blocks([{
            "type": "table_chart",
            "region": "full",
            "data": {
                "headers": ["Engine", "Type"],
                "rows": [["Seedream", "AI"], ["DALL-E", "AI"]]
            }
        }])
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("Engine" in t for t in texts)
        assert any("Seedream" in t for t in texts)


class TestComposableBlocks:
    """Multiple blocks on one slide render without overlap."""

    def test_two_blocks_side_by_side(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(pr)
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        br.render(slide, [
            {"type": "hexagons", "region": "left-2-3", "data": {
                "items": [{"label": "A"}, {"label": "B"}]
            }},
            {"type": "bullets", "region": "right-1-3", "data": {
                "items": ["Point 1", "Point 2"]
            }},
        ])
        assert len(list(slide.shapes)) >= 3

    def test_top_bottom_blocks(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(pr)
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        br.render(slide, [
            {"type": "metrics", "region": "top", "data": {
                "items": [{"value": "40K", "label": "Styles"}]
            }},
            {"type": "cards", "region": "bottom", "data": {
                "items": [{"title": "T1", "text": "Body"}, {"title": "T2", "text": "Body"}]
            }},
        ])
        assert len(list(slide.shapes)) >= 3

    def test_auto_region_two_blocks(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        from ppt_pro_max.enterprise.block_renderer import BlockRenderer
        br = BlockRenderer(pr)
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        br.render(slide, [
            {"type": "bullets", "data": {"items": ["A"]}},
            {"type": "bullets", "data": {"items": ["B"]}},
        ])
        assert len(list(slide.shapes)) >= 2


class TestSectionDividerBugFix:
    """Section divider must not create blank slides."""

    def test_section_creates_exactly_one_slide(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {"goal": "section", "title": "Test Section", "section_number": 1}
        pr.render_slide(prs, page)
        assert len(prs.slides) == 1

    def test_section_slide_has_content(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {"goal": "section", "title": "Test Section", "section_number": 1}
        slide = pr.render_slide(prs, page)
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert len(texts) >= 1
        assert any("Test Section" in t for t in texts)

    def test_section_with_subtitle(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {"goal": "section", "title": "Architecture", "subtitle": "Core Design", "section_number": 2}
        slide = pr.render_slide(prs, page)
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("Architecture" in t for t in texts)

    def test_multiple_sections_no_extra_slides(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        for i in range(3):
            page = {"goal": "section", "title": f"Section {i+1}", "section_number": i+1}
            pr.render_slide(prs, page)
        assert len(prs.slides) == 3


class TestRenderSlideBlocksPath:
    """render_slide() uses blocks when present, legacy when absent."""

    def test_blocks_key_triggers_block_path(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {
            "goal": "content",
            "title": "Test",
            "blocks": [
                {"type": "bullets", "region": "full", "data": {"items": ["A", "B"]}}
            ]
        }
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("A" in t for t in texts)

    def test_no_blocks_key_triggers_legacy_path(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {
            "goal": "content",
            "title": "Test",
            "bullets": ["Legacy bullet 1", "Legacy bullet 2"]
        }
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("Legacy" in t for t in texts)

    def test_hero_with_blocks(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {
            "goal": "hook",
            "title": "Hero Title",
            "subtitle": "Hero Sub",
            "blocks": [
                {"type": "metrics", "region": "bottom", "data": {
                    "items": [{"value": "100+", "label": "Features"}]
                }}
            ]
        }
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1
        texts = []
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("Hero Title" in t for t in texts)

    def test_blocks_with_multiple_types_on_one_slide(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {
            "goal": "content",
            "title": "Mixed Layout",
            "blocks": [
                {"type": "hexagons", "region": "left-2-3", "data": {
                    "items": [{"label": "A"}, {"label": "B"}]
                }},
                {"type": "bullets", "region": "right-1-3", "data": {
                    "items": ["Detail 1", "Detail 2"]
                }}
            ]
        }
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1
        assert len(list(slide.shapes)) >= 4


class TestContentParserBlocksPassthrough:
    """content_parser preserves blocks key from content.json."""

    def test_blocks_preserved_in_parsed_output(self):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content
        content = {
            "slides": [{
                "goal": "content",
                "title": "Test",
                "blocks": [
                    {"type": "bullets", "region": "left", "data": {"items": ["A"]}},
                    {"type": "hexagons", "region": "right", "data": {"items": [{"label": "B"}]}}
                ]
            }]
        }
        result = load_enterprise_content(content, ".")
        assert len(result) == 1
        assert "blocks" in result[0]
        assert len(result[0]["blocks"]) == 2

    def test_old_format_still_works(self):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content
        content = {
            "slides": [{
                "goal": "features",
                "title": "Test",
                "cards": [{"title": "A", "text": "B"}]
            }]
        }
        result = load_enterprise_content(content, ".")
        assert len(result) == 1
        assert result[0].get("blocks") is None
        assert result[0].get("cards") is not None


class TestPipelineBlocksIntegration:
    """Enterprise Pipeline renders blocks from content.json end-to-end."""

    def test_pipeline_generates_pptx_with_blocks(self, tmp_path):
        project_dir = str(tmp_path)
        brand = {
            "colors": {
                "primary": "#00D4FF",
                "accent": "#A855F7",
                "foreground": "#F0F4F8",
                "muted-foreground": "#94A3B8",
                "background": "#050816",
                "muted": "#0F1A2E"
            },
            "fonts": {"heading": "Space Grotesk", "body": "Inter"}
        }
        content = {
            "meta": {"title": "Blocks Test"},
            "slides": [
                {"goal": "hook", "title": "Blocks Test", "subtitle": "Testing blocks"},
                {"goal": "content", "title": "Hexagons + Bullets", "blocks": [
                    {"type": "hexagons", "region": "left-2-3", "data": {
                        "items": [{"label": "Plan"}, {"label": "Design"}, {"label": "Build"}, {"label": "Ship"}]
                    }},
                    {"type": "bullets", "region": "right-1-3", "data": {
                        "items": ["Strategy", "40K+ styles", "AIDA copy"]
                    }}
                ]},
                {"goal": "content", "title": "Metrics + Cards", "blocks": [
                    {"type": "metrics", "region": "top", "data": {
                        "items": [
                            {"value": "40K+", "label": "Styles"},
                            {"value": "28", "label": "Upgrades"}
                        ]
                    }},
                    {"type": "cards", "region": "bottom", "data": {
                        "items": [
                            {"title": "Tier 1", "text": "Visual"},
                            {"title": "Tier 2", "text": "Typography"}
                        ]
                    }}
                ]},
                {"goal": "cta", "title": "Get Started", "subtitle": "Start now"}
            ]
        }
        import json
        with open(os.path.join(project_dir, "brand.json"), "w", encoding="utf-8") as f:
            json.dump(brand, f)
        with open(os.path.join(project_dir, "content.json"), "w", encoding="utf-8") as f:
            json.dump(content, f)

        from ppt_pro_max import generate_ppt
        result = generate_ppt(
            "Blocks Test",
            project=project_dir,
            density=6,
            motion=3,
            output=str(tmp_path / "output.pptx")
        )

        assert not result.get("error"), f"Pipeline error: {result.get('error')}"
        output_path = result.get("output_path", "")
        assert output_path and os.path.isfile(output_path), f"No output file: {output_path}"

        prs = Presentation(output_path)
        assert len(prs.slides) >= 3

        slide2 = prs.slides[1]
        shapes = list(slide2.shapes)
        assert len(shapes) >= 4, f"Slide 2 only has {len(shapes)} shapes (expected blocks)"

    def test_pipeline_old_format_still_works(self, tmp_path):
        project_dir = str(tmp_path)
        brand = {
            "colors": {
                "primary": "#2563EB", "accent": "#F97316",
                "foreground": "#1E293B", "muted-foreground": "#64748B",
                "background": "#FFFFFF", "muted": "#F1F5F9"
            },
            "fonts": {"heading": "Calibri", "body": "Calibri"}
        }
        content = {
            "meta": {"title": "Legacy Test"},
            "slides": [
                {"goal": "hook", "title": "Legacy Test", "subtitle": "Old format"},
                {"goal": "features", "title": "Three Modes", "cards": [
                    {"title": "A", "text": "Mode A"},
                    {"title": "B", "text": "Mode B"},
                    {"title": "C", "text": "Mode C"}
                ]},
                {"goal": "cta", "title": "Done", "subtitle": "End"}
            ]
        }
        import json
        with open(os.path.join(project_dir, "brand.json"), "w", encoding="utf-8") as f:
            json.dump(brand, f)
        with open(os.path.join(project_dir, "content.json"), "w", encoding="utf-8") as f:
            json.dump(content, f)

        from ppt_pro_max import generate_ppt
        result = generate_ppt(
            "Legacy Test",
            project=project_dir,
            output=str(tmp_path / "legacy_output.pptx")
        )

        assert not result.get("error"), f"Pipeline error: {result.get('error')}"
        assert os.path.isfile(result.get("output_path", ""))


class TestBackwardCompatibility:
    """All existing content.json formats must work without modification."""

    def test_cards_only_no_blocks(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {"goal": "features", "title": "T", "cards": [
            {"title": "A", "text": "B"}, {"title": "C", "text": "D"}
        ]}
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1
        assert len(list(slide.shapes)) >= 4

    def test_bullets_only_no_blocks(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {"goal": "content", "title": "T", "bullets": ["A", "B", "C"]}
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1

    def test_code_only_no_blocks(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {"goal": "code", "title": "T", "code": {"language": "python", "source": "x=1"}}
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1

    def test_exercise_only_no_blocks(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {"goal": "content", "title": "T", "exercise": {
            "instructions": "Try", "duration": "5 min", "steps": ["Step 1"]
        }}
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1

    def test_diagram_only_no_blocks(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        page = {
            "goal": "content", "title": "T",
            "diagram_type": "flowchart",
            "diagram_data": {
                "nodes": [{"id": "1", "label": "A"}, {"id": "2", "label": "B"}],
                "connectors": [{"from": "1", "to": "2"}]
            }
        }
        slide = pr.render_slide(prs, page)
        assert len(prs.slides) == 1
