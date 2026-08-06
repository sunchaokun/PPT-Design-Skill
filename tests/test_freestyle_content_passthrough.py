"""Tests for FreeStyle content.json passthrough rendering (agent-driven content flow)."""

from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation

from ppt_pro_max import generate_ppt
from ppt_pro_max import _has_slides, _load_content_json_to_pages
from ppt_pro_max import _build_freestyle_page_dicts


def _write_content(tmp_path: Path, slides: list, name: str = "content.json") -> Path:
    path = tmp_path / name
    path.write_text(json.dumps({"slides": slides}, ensure_ascii=False), encoding="utf-8")
    return path


def _slide_texts(slide) -> list[str]:
    return [sh.text.strip() for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]


def _has_chart(slide) -> bool:
    for sh in slide.shapes:
        if getattr(sh, "chart", None) is not None:
            return True
    return False


def _load_prs(result) -> Presentation:
    return Presentation(result["output_path"])


_SAMPLE_SLIDES = [
    {"goal": "hook", "title": "Hero Title", "subtitle": "Hero Subtitle"},
    {"goal": "problem", "title": "The Problem", "bullets": ["Pain A", "Pain B", "Pain C"]},
    {"goal": "section", "title": "Our Solution", "section_number": "01"},
    {"goal": "features", "title": "Features",
     "cards": [
         {"title": "Feature A", "text": "Description A"},
         {"title": "Feature B", "text": "Description B"},
         {"title": "Feature C", "text": "Description C"},
     ]},
    {"goal": "data", "title": "Metrics",
     "chart": {"type": "bar", "categories": ["A", "B", "C"],
               "series": [{"name": "Acc", "values": [1, 2, 3]}]}},
    {"goal": "cta", "title": "Start Now", "subtitle": "Free tier available."},
]


class TestHasSlides:

    def test_with_slides_list(self, tmp_path):
        cf = _write_content(tmp_path, [{"goal": "hook", "title": "T"}])
        assert _has_slides(str(cf)) is True

    def test_empty_slides_list(self, tmp_path):
        cf = _write_content(tmp_path, [])
        assert _has_slides(str(cf)) is False

    def test_no_slides_key(self, tmp_path):
        cf = tmp_path / "content.json"
        cf.write_text(json.dumps({"meta": {"title": "X"}}), encoding="utf-8")
        assert _has_slides(str(cf)) is False

    def test_missing_file(self, tmp_path):
        assert _has_slides(str(tmp_path / "nope.json")) is False

    def test_bad_json(self, tmp_path):
        cf = tmp_path / "content.json"
        cf.write_text("{not valid json", encoding="utf-8")
        assert _has_slides(str(cf)) is False


class TestLoadContentJsonToPages:

    def test_all_fields_preserved(self, tmp_path):
        cf = _write_content(tmp_path, [
            {"goal": "features", "title": "F", "cards": [{"title": "A", "text": "x"}]},
            {"goal": "content", "title": "D", "diagram": {"type": "flowchart", "data": {}}},
            {"goal": "code", "title": "C", "code": {"language": "python", "source": "print(1)"}},
            {"goal": "exercise", "title": "E", "exercise": {"instructions": "Do it", "steps": ["s1", "s2"]}},
            {"goal": "data", "title": "M", "chart": {"type": "bar"}},
            {"goal": "section", "title": "S", "section_number": "02"},
        ])
        pages = _load_content_json_to_pages(str(cf))
        assert pages is not None
        assert len(pages) == 6
        assert pages[0]["cards"] is not None
        assert pages[1]["diagram_type"] == "flowchart"
        assert pages[1]["diagram_data"] is not None
        assert pages[2]["code"] is not None
        assert pages[3]["exercise"] is not None
        assert pages[4]["chart"] is not None
        assert pages[5]["section_number"] == "02"

    def test_section_number_passthrough(self, tmp_path):
        cf = _write_content(tmp_path, [{"goal": "section", "title": "S", "section_number": "01"}])
        pages = _load_content_json_to_pages(str(cf))
        assert pages[0]["section_number"] == "01"

    def test_section_number_missing_defaults_none(self, tmp_path):
        cf = _write_content(tmp_path, [{"goal": "section", "title": "S"}])
        pages = _load_content_json_to_pages(str(cf))
        assert pages[0]["section_number"] is None

    def test_image_keywords_added(self, tmp_path):
        cf = _write_content(tmp_path, [{"goal": "hook", "title": "H"}])
        pages = _load_content_json_to_pages(str(cf))
        assert pages[0]["image_keywords"]  # non-empty from _GOAL_IMAGE_KEYWORDS

    def test_relative_image_resolved(self, tmp_path):
        img = tmp_path / "hero.png"
        img.write_bytes(b"\x89PNG\r\n\x1a\n")
        cf = _write_content(tmp_path, [{"goal": "hook", "title": "H", "image": "hero.png"}])
        pages = _load_content_json_to_pages(str(cf))
        assert pages[0]["image"] == str(img)

    def test_bad_json_returns_none(self, tmp_path):
        cf = tmp_path / "content.json"
        cf.write_text("{oops", encoding="utf-8")
        assert _load_content_json_to_pages(str(cf)) is None

    def test_missing_file_returns_none(self, tmp_path):
        assert _load_content_json_to_pages(str(tmp_path / "nope.json")) is None

    def test_no_slides_returns_none(self, tmp_path):
        cf = tmp_path / "content.json"
        cf.write_text(json.dumps({"meta": {}}), encoding="utf-8")
        assert _load_content_json_to_pages(str(cf)) is None


class TestPassthroughRendering:

    def test_passthrough_renders_all_slides(self, tmp_path):
        cf = _write_content(tmp_path, _SAMPLE_SLIDES)
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) == 6

    def test_passthrough_query_omitted(self, tmp_path):
        cf = _write_content(tmp_path, _SAMPLE_SLIDES)
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        assert result["page_count"] == 6
        assert result["strategy"] == "generated"

    def test_passthrough_style_sets_theme_atoms(self, tmp_path):
        cf = _write_content(tmp_path, [{"goal": "hook", "title": "H"}])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        assert "theme_atoms" in result

    def test_passthrough_without_style_uses_default_theme(self, tmp_path):
        cf = _write_content(tmp_path, _SAMPLE_SLIDES)
        result = generate_ppt(content_file=str(cf), output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) == 6

    def test_passthrough_content_dict_param(self, tmp_path):
        result = generate_ppt(content={"slides": [{"goal": "hook", "title": "H"}]},
                              style="dark-tech", output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) == 1

    def test_section_number_string_rendered(self, tmp_path):
        cf = _write_content(tmp_path, [{"goal": "section", "title": "Our Solution",
                                        "section_number": "01"}])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        texts = _slide_texts(_load_prs(result).slides[0])
        assert "01" in texts

    def test_section_number_integer_rendered(self, tmp_path):
        cf = _write_content(tmp_path, [{"goal": "section", "title": "Our Solution",
                                        "section_number": 3}])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        texts = _slide_texts(_load_prs(result).slides[0])
        assert "03" in texts

    def test_section_number_missing_renders_index(self, tmp_path):
        cf = _write_content(tmp_path, [{"goal": "section", "title": "Our Solution"}])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        texts = _slide_texts(_load_prs(result).slides[0])
        assert "01" in texts  # page_index + 1

    def test_cards_rendered(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "features", "title": "Features",
            "cards": [{"title": "Alpha", "text": "x"}, {"title": "Beta", "text": "y"},
                      {"title": "Gamma", "text": "z"}],
        }])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        texts = _slide_texts(_load_prs(result).slides[0])
        joined = " ".join(texts)
        assert "Alpha" in joined and "Beta" in joined and "Gamma" in joined

    def test_chart_rendered(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "data", "title": "Metrics",
            "chart": {"type": "bar", "categories": ["A", "B"],
                      "series": [{"name": "S", "values": [1, 2]}]},
        }])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        assert _has_chart(_load_prs(result).slides[0])

    def test_code_rendered(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "code", "title": "Quick Start",
            "code": {"language": "python", "source": "print('hello')"},
        }])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        texts = _slide_texts(_load_prs(result).slides[0])
        assert any("hello" in t for t in texts)

    def test_exercise_rendered(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "exercise", "title": "Try It",
            "exercise": {"instructions": "Complete the steps", "duration": "5 min",
                         "steps": ["Step one", "Step two"]},
        }])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        texts = _slide_texts(_load_prs(result).slides[0])
        joined = " ".join(texts).lower()
        assert "exercise" in joined and "step one" in joined

    def test_diagram_rendered_without_crash(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "content", "title": "Architecture",
            "diagram": {"type": "flowchart",
                        "data": {"nodes": [{"id": "a", "text": "A"}, {"id": "b", "text": "B"}],
                                 "connectors": [["a", "b"]]}},
        }])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) == 1

    def test_passthrough_dry_run(self, tmp_path):
        cf = _write_content(tmp_path, _SAMPLE_SLIDES)
        result = generate_ppt(content_file=str(cf), style="dark-tech", dry_run=True)
        assert result["dry_run"] is True
        assert result["page_count"] == 6


class TestFallbackToLegacy:

    def test_no_slides_key_falls_back(self, tmp_path):
        cf = tmp_path / "content.json"
        cf.write_text(json.dumps({"meta": {"title": "X"}}), encoding="utf-8")
        result = generate_ppt("AI pitch", content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) >= 1

    def test_missing_content_file_falls_back(self, tmp_path):
        result = generate_ppt("AI pitch", content_file=str(tmp_path / "nope.json"),
                              style="dark-tech", output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) >= 1

    def test_malformed_json_falls_back(self, tmp_path):
        cf = tmp_path / "content.json"
        cf.write_text("{broken", encoding="utf-8")
        result = generate_ppt("AI pitch", content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) >= 1

    def test_no_content_file_legacy(self, tmp_path):
        result = generate_ppt("AI pitch", style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) >= 1

    def test_legacy_dry_run(self, tmp_path):
        result = generate_ppt("AI pitch", style="dark-tech", dry_run=True)
        assert result["dry_run"] is True


class TestLayoutVariantPassed:

    def test_layout_variant_and_decoration_passed(self, monkeypatch, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        captured = {}
        original = PrecisionRenderer.render_slide

        def fake(self, prs, page, component_lib=None, layout_variant=None,
                 page_index=0, total_pages=0, proactive_component=True):
            captured["layout_variant"] = layout_variant
            captured["page"] = page
            return original(self, prs, page, component_lib=component_lib,
                            layout_variant=layout_variant,
                            page_index=page_index, total_pages=total_pages,
                            proactive_component=proactive_component)

        monkeypatch.setattr(PrecisionRenderer, "render_slide", fake)
        cf = _write_content(tmp_path, [{"goal": "content", "title": "T", "bullets": ["a"]}])
        generate_ppt(content_file=str(cf), style="dark-tech",
                     output=str(tmp_path / "out.pptx"))
        lv = captured["layout_variant"]
        assert lv is not None
        assert lv.get("content_margin_left") == 0.6  # wide-cards layout
        assert lv.get("decoration_style") == "neon-lines"  # dark-tech decoration

    def test_layout_variant_passed_in_legacy_path(self, monkeypatch, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        captured = {}
        original = PrecisionRenderer.render_slide

        def fake(self, prs, page, component_lib=None, layout_variant=None,
                 page_index=0, total_pages=0, proactive_component=True):
            captured["layout_variant"] = layout_variant
            return original(self, prs, page, component_lib=component_lib,
                            layout_variant=layout_variant,
                            page_index=page_index, total_pages=total_pages,
                            proactive_component=proactive_component)

        monkeypatch.setattr(PrecisionRenderer, "render_slide", fake)
        generate_ppt("AI pitch", style="dark-tech", output=str(tmp_path / "out.pptx"))
        assert captured["layout_variant"] is not None
        assert captured["layout_variant"].get("decoration_style") == "neon-lines"


class TestChartDataKeyFix:

    def test_freestyle_page_dict_uses_chart_key(self):
        from ppt_pro_max.content.content_generator import PageContent

        class _Design:
            goal = "traction"

        content = PageContent(
            position=0, goal="traction",
            title="T", subtitle=None, bullets=None,
            metrics=None, chart_data={"type": "bar", "values": [1, 2]},
            quote=None, image_keywords="growth",
        )
        pages = _build_freestyle_page_dicts([_Design()], [content], None)
        assert "chart" in pages[0]
        assert "chart_data" not in pages[0]
        assert pages[0]["chart"] == content.chart_data


class TestConflictWarnings:

    def test_proposal_conflict_warns(self, tmp_path):
        cf = _write_content(tmp_path, _SAMPLE_SLIDES)
        with pytest.warns(UserWarning, match="proposal"):
            result = generate_ppt(content_file=str(cf), proposal=True,
                                  output=str(tmp_path / "out"))
        assert "proposals" in result

    def test_beautify_conflict_warns(self, tmp_path):
        cf = _write_content(tmp_path, _SAMPLE_SLIDES)
        with pytest.warns(UserWarning, match="beautify"):
            result = generate_ppt(content_file=str(cf), beautify="x.pptx",
                                  output=str(tmp_path / "out"))
        assert isinstance(result, dict)
        assert "num_slides" in result or "output_path" in result or "error" in result


class TestRegression:

    def test_freestyle_component_path_still_works(self, tmp_path):
        result = generate_ppt("AI startup pitch", style="vibrant-startup",
                              output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        assert len(prs.slides) >= 1


class TestProactiveComponentDisabled:

    def test_passthrough_problem_renders_bullets_not_component(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "problem", "title": "The Problem",
            "bullets": ["Pain A", "Pain B", "Pain C"],
        }])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        prs = _load_prs(result)
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        groups = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]
        assert not groups  # no proactive component hijack
        texts = _slide_texts(prs.slides[0])
        assert any("Pain A" in t for t in texts)

    def test_legacy_path_also_disables_proactive_components(self, tmp_path):
        # Both passthrough and legacy one-liner disable proactive component matching;
        # components only render via explicit component_type in content.json.
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        captured = {}
        original = PrecisionRenderer.render_slide

        def fake(self, prs, page, component_lib=None, layout_variant=None,
                 page_index=0, total_pages=0, proactive_component=True):
            captured["proactive_component"] = proactive_component
            return original(self, prs, page, component_lib=component_lib,
                            layout_variant=layout_variant,
                            page_index=page_index, total_pages=total_pages,
                            proactive_component=proactive_component)

        monkeypatch = __import__("pytest").MonkeyPatch()
        monkeypatch.setattr(PrecisionRenderer, "render_slide", fake)
        try:
            generate_ppt("AI pitch", style="dark-tech",
                         output=str(tmp_path / "out.pptx"))
        finally:
            monkeypatch.undo()
        assert captured["proactive_component"] is False


class TestDiagramDataFormats:

    def test_nested_diagram_data_extracted(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "content", "title": "Arch",
            "diagram": {"type": "flowchart",
                        "data": {"nodes": [{"id": "a", "text": "Alpha"},
                                           {"id": "b", "text": "Beta"}],
                                 "connectors": [["a", "b"]]}},
        }])
        pages = _load_content_json_to_pages(str(cf))
        assert pages[0]["diagram_type"] == "flowchart"
        assert pages[0]["diagram_data"]["nodes"]  # inner data, not wrapper
        assert "type" not in pages[0]["diagram_data"]

    def test_flat_diagram_data_preserved(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "content", "title": "Arch",
            "diagram": {"type": "flowchart",
                        "nodes": [{"label": "Start"}]},
        }])
        pages = _load_content_json_to_pages(str(cf))
        assert pages[0]["diagram_data"]["type"] == "flowchart"  # flat format kept

    def test_flowchart_renders_node_text(self, tmp_path):
        cf = _write_content(tmp_path, [{
            "goal": "content", "title": "Architecture",
            "diagram": {"type": "flowchart",
                        "data": {"nodes": [{"id": "a", "text": "Data Sources"},
                                           {"id": "b", "text": "ETL Pipeline"}],
                                 "connectors": [["a", "b"]]}},
        }])
        result = generate_ppt(content_file=str(cf), style="dark-tech",
                              output=str(tmp_path / "out.pptx"))
        texts = _slide_texts(_load_prs(result).slides[0])
        assert any("Data Sources" in t for t in texts)
        assert any("ETL Pipeline" in t for t in texts)
