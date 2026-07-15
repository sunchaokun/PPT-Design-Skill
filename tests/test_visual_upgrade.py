"""Tests for visual_effects, shape_factory, freeform_builder, group_builder."""

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches
from lxml import etree


# ── visual_effects tests ──

class TestGradientFill:
    def test_linear_gradient_on_shape(self):
        from ppt_pro_max.renderer.visual_effects import apply_gradient
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        apply_gradient(shape, "FF5500", "EE9003", gradient_type="linear", angle=5400000)
        spPr = shape._element.find(qn("p:spPr"))
        gradFill = spPr.find(qn("a:gradFill"))
        assert gradFill is not None
        lin = gradFill.find(qn("a:lin"))
        assert lin is not None
        assert lin.get("ang") == "5400000"

    def test_path_gradient_on_shape(self):
        from ppt_pro_max.renderer.visual_effects import apply_gradient
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        apply_gradient(shape, "1D78FA", "0165FF", gradient_type="path")
        spPr = shape._element.find(qn("p:spPr"))
        gradFill = spPr.find(qn("a:gradFill"))
        assert gradFill is not None
        path_el = gradFill.find(qn("a:path"))
        assert path_el is not None
        assert path_el.get("path") == "circle"

    def test_gradient_replaces_solid_fill(self):
        from ppt_pro_max.renderer.visual_effects import apply_gradient
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        shape.fill.solid()
        apply_gradient(shape, "FF5500", "EE9003")
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:solidFill")) is None
        assert spPr.find(qn("a:gradFill")) is not None

    def test_preset_gradient(self):
        from ppt_pro_max.renderer.visual_effects import apply_preset_gradient
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        apply_preset_gradient(shape, "orange-warm")
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None

    def test_3_stop_gradient(self):
        from ppt_pro_max.renderer.visual_effects import GradientFill, GradientStop
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        grad = GradientFill(
            stops=[
                GradientStop(color="FF5500", position=0),
                GradientStop(color="EE9003", position=50000),
                GradientStop(color="CC6600", position=100000),
            ],
            angle=5400000,
        )
        grad.apply(shape)
        spPr = shape._element.find(qn("p:spPr"))
        gsLst = spPr.find(qn("a:gradFill")).find(qn("a:gsLst"))
        stops = gsLst.findall(qn("a:gs"))
        assert len(stops) == 3


class TestShadow:
    def test_shadow_on_shape(self):
        from ppt_pro_max.renderer.visual_effects import apply_shadow
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        apply_shadow(shape, blur_pt=6, distance_pt=3, alpha_pct=25)
        spPr = shape._element.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        assert effectLst is not None
        outerShdw = effectLst.find(qn("a:outerShdw"))
        assert outerShdw is not None
        assert outerShdw.get("blurRad") == str(6 * 12700)

    def test_shadow_alpha(self):
        from ppt_pro_max.renderer.visual_effects import apply_shadow
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        apply_shadow(shape, alpha_pct=30)
        spPr = shape._element.find(qn("p:spPr"))
        alpha_el = spPr.find(".//" + qn("a:alpha"))
        assert alpha_el is not None
        assert alpha_el.get("val") == "30000"


class TestSoftEdge:
    def test_soft_edge_on_shape(self):
        from ppt_pro_max.renderer.visual_effects import apply_soft_edge
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        apply_soft_edge(shape, radius_pt=2.0)
        spPr = shape._element.find(qn("p:spPr"))
        softEdge = spPr.find(".//" + qn("a:softEdge"))
        assert softEdge is not None
        assert softEdge.get("rad") == str(int(2.0 * 12700))


class TestReflection:
    def test_reflection_on_shape(self):
        from ppt_pro_max.renderer.visual_effects import apply_reflection
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        apply_reflection(shape, alpha_pct=20)
        spPr = shape._element.find(qn("p:spPr"))
        reflection = spPr.find(".//" + qn("a:reflection"))
        assert reflection is not None


class TestGlow:
    def test_glow_on_shape(self):
        from ppt_pro_max.renderer.visual_effects import apply_glow
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        apply_glow(shape, radius_pt=8, color="#2563EB", alpha_pct=40)
        spPr = shape._element.find(qn("p:spPr"))
        glow = spPr.find(".//" + qn("a:glow"))
        assert glow is not None


class TestSolidFillWithAlpha:
    def test_alpha_fill(self):
        from ppt_pro_max.renderer.visual_effects import set_solid_fill_with_alpha
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        set_solid_fill_with_alpha(shape, "#000000", 65)
        spPr = shape._element.find(qn("p:spPr"))
        alpha_el = spPr.find(".//" + qn("a:alpha"))
        assert alpha_el is not None
        assert alpha_el.get("val") == "65000"


class TestLineGradient:
    def test_line_gradient(self):
        from ppt_pro_max.renderer.visual_effects import set_line_gradient
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        set_line_gradient(shape, "FF5500", "EE9003", width_pt=2.0)
        spPr = shape._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        assert ln is not None
        gradFill = ln.find(qn("a:gradFill"))
        assert gradFill is not None


# ── shape_factory tests ──

class TestShapeFactory:
    def test_glow_oval(self):
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        factory = ShapeFactory(brand_colors={"primary": "#1D78FA"})
        shape = factory.add_glow_oval(slide, 1, 1, 1.2, 1.2, label="策略")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None
        assert spPr.find(".//" + qn("a:outerShdw")) is not None

    def test_ring_node(self):
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        factory = ShapeFactory()
        shape = factory.add_ring_node(slide, 2, 2, 3.0, label="核心")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None

    def test_hexagon_card(self):
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        factory = ShapeFactory()
        shape = factory.add_hexagon_card(slide, 1, 1, 1.5, label="技术")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None

    def test_radial_diagram(self):
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        factory = ShapeFactory(brand_colors={"primary": "#1D78FA", "accent": "#FF5500"})
        nodes = [
            {"label": "策略"}, {"label": "运营"}, {"label": "技术"},
            {"label": "市场"}, {"label": "产品"}, {"label": "数据"},
        ]
        shapes = factory.add_radial_diagram(slide, 6.65, 3.75, "核心", nodes=nodes)
        assert len(shapes) == 7
        for s in shapes:
            spPr = s._element.find(qn("p:spPr"))
            assert spPr.find(qn("a:gradFill")) is not None

    def test_ring_diagram(self):
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        factory = ShapeFactory()
        nodes = [{"label": "A"}, {"label": "B"}, {"label": "C"}]
        shapes = factory.add_ring_diagram(slide, 6.0, 3.5, "中心", nodes=nodes)
        assert len(shapes) == 4

    def test_chevron_step(self):
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        factory = ShapeFactory()
        shape = factory.add_chevron_step(slide, 1, 2, 2.5, 1.0, label="步骤1", step_index=0)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None

    def test_diamond_node(self):
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        factory = ShapeFactory()
        shape = factory.add_diamond_node(slide, 3, 3, 1.0, label="决策")
        assert shape is not None

    def test_rounded_card_with_gradient(self):
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        factory = ShapeFactory()
        shape = factory.add_rounded_card(slide, 1, 1, 3, 2, title="标题", body="内容")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None


# ── freeform_builder tests ──

class TestFreeformBuilder:
    def test_basic_freeform(self):
        from ppt_pro_max.renderer.freeform_builder import FreeformBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        builder = FreeformBuilder()
        builder.move_to(0, 0)
        builder.line_to(2, 0)
        builder.line_to(2, 1)
        builder.close()
        elem = builder.build(slide, 1, 1, 3, 2, line_color="#FF5500")
        assert elem is not None
        custGeom = elem.find(".//" + qn("a:custGeom"))
        assert custGeom is not None

    def test_arrow_connector(self):
        from ppt_pro_max.renderer.freeform_builder import make_arrow_connector
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        elem = make_arrow_connector(slide, 1, 1, 3, 3, color="#1D78FA")
        assert elem is not None
        tailEnd = elem.find(".//" + qn("a:tailEnd"))
        assert tailEnd is not None
        assert tailEnd.get("type") == "triangle"

    def test_chevron_shape(self):
        from ppt_pro_max.renderer.freeform_builder import make_chevron_shape
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        elem = make_chevron_shape(slide, 1, 1, 3, 1.5, fill_color="#1D78FA")
        assert elem is not None
        spPr = elem.find(qn("p:spPr"))
        solidFill = spPr.find(qn("a:solidFill"))
        assert solidFill is not None


# ── group_builder tests ──

class TestGroupBuilder:
    def test_group_two_shapes(self):
        from ppt_pro_max.renderer.group_builder import group_last_n_shapes
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.enum.shapes import MSO_SHAPE
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(2))
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(1.5), Inches(2), Inches(2))
        result = group_last_n_shapes(slide, 2)
        assert result is not None

    def test_group_too_few_shapes(self):
        from ppt_pro_max.renderer.group_builder import group_last_n_shapes
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        result = group_last_n_shapes(slide, 1)
        assert result is None


# ── PrecisionRenderer integration tests ──

class TestPrecisionRendererUpgrade:
    def test_add_oval_with_gradient(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        slide = renderer.add_slide(prs)
        shape = renderer.add_oval(slide, 1, 1, 2, 2, label="测试")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None
        assert spPr.find(".//" + qn("a:outerShdw")) is not None

    def test_add_donut_with_gradient(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        slide = renderer.add_slide(prs)
        shape = renderer.add_donut(slide, 2, 2, 3.0, label="核心")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None

    def test_add_hexagon_with_gradient(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        slide = renderer.add_slide(prs)
        shape = renderer.add_hexagon(slide, 1, 1, 2.0, label="技术")
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None

    def test_add_rect_with_gradient_and_shadow(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        slide = renderer.add_slide(prs)
        shape = renderer.add_rect(slide, 1, 1, 3, 2, fill_hex="#1D78FA", gradient=True, shadow=True)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:gradFill")) is not None
        assert spPr.find(".//" + qn("a:outerShdw")) is not None

    def test_card_with_gradient_and_shadow(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        slide = renderer.add_slide(prs)
        renderer.add_card(slide, 1, 1, 3, 3, "标题", "内容")
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None and spPr.find(qn("a:gradFill")) is not None:
                break
        else:
            pytest.fail("No gradient found on card shapes")

    def test_shape_factory_property(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.renderer.shape_factory import ShapeFactory
        renderer = PrecisionRenderer()
        factory = renderer.shape_factory
        assert isinstance(factory, ShapeFactory)

    def test_full_ppt_with_visual_effects(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()

        slide = renderer.add_slide(prs)
        renderer.apply_hero_overlay(slide, prs)
        renderer.add_text(slide, "测试标题", 1.2, 2.0, 8, 1.5, size=44, bold=True)

        slide = renderer.add_slide(prs)
        renderer.apply_brand_background(slide, prs)
        renderer.add_text(slide, "功能展示", 0.9, 0.5, 11, 0.6, size=28, bold=True)
        renderer.add_oval(slide, 1, 2, 1.5, 1.5, label="策略", fill_role="primary")
        renderer.add_oval(slide, 3, 2, 1.5, 1.5, label="运营", fill_role="accent")
        renderer.add_donut(slide, 5.5, 2, 2.5, label="核心")
        renderer.add_hexagon(slide, 9, 2, 2.0, label="技术")

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        prs.save(path)

        prs2 = Presentation(path)
        slide2 = prs2.slides[1]
        grad_count = 0
        shadow_count = 0
        for shape in slide2.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                if spPr.find(qn("a:gradFill")) is not None:
                    grad_count += 1
                if spPr.find(".//" + qn("a:outerShdw")) is not None:
                    shadow_count += 1

        assert grad_count >= 3, f"Expected >=3 gradients, got {grad_count}"
        assert shadow_count >= 3, f"Expected >=3 shadows, got {shadow_count}"
        os.unlink(path)


# ── DiagramEngine upgrade tests ──

class TestDiagramEngineUpgrade:
    def test_cycle_with_gradient_and_shadow(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        style = DiagramStyle()
        region = Region(left=0.9, top=1.5, width=7.0, height=5.0)
        engine = DiagramEngine()
        data = {
            "nodes": [
                {"label": "A"}, {"label": "B"}, {"label": "C"}, {"label": "D"},
            ],
        }
        engine.render(slide, "cycle", data, style, region)
        grad_count = 0
        shadow_count = 0
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                if spPr.find(qn("a:gradFill")) is not None:
                    grad_count += 1
                if spPr.find(".//" + qn("a:outerShdw")) is not None:
                    shadow_count += 1
        assert grad_count >= 2, f"Expected >=2 gradients in cycle, got {grad_count}"
        assert shadow_count >= 2, f"Expected >=2 shadows in cycle, got {shadow_count}"

    def test_pyramid_with_gradient(self):
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        style = DiagramStyle()
        region = Region(left=0.9, top=1.5, width=7.0, height=5.0)
        engine = DiagramEngine()
        data = {
            "levels": [
                {"label": "顶层"}, {"label": "中层"}, {"label": "底层"},
            ],
        }
        engine.render(slide, "pyramid", data, style, region)
        grad_count = 0
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None and spPr.find(qn("a:gradFill")) is not None:
                grad_count += 1
        assert grad_count >= 1, f"Expected >=1 gradient in pyramid, got {grad_count}"
