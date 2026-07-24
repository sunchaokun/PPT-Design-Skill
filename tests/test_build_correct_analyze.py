"""Correct Build mode analysis: read defRPr (paragraph-level) + rPr (run-level)."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def extract_text_props(slide):
    results = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            pPr = para._p.find(f'{{{ns_a}}}pPr')
            def_font = def_sz = def_clr = None
            if pPr is not None:
                defRPr = pPr.find(f'{{{ns_a}}}defRPr')
                if defRPr is not None:
                    def_sz = defRPr.get('sz')
                    latin = defRPr.find(f'{{{ns_a}}}latin')
                    if latin is not None:
                        def_font = latin.get('typeface')
                    for child in defRPr:
                        if 'srgbClr' in child.tag:
                            def_clr = child.get('val')
            
            for run in para.runs:
                if not run.text.strip():
                    continue
                rPr = run._r.find(f'{{{ns_a}}}rPr')
                run_font = run_sz = run_clr = None
                if rPr is not None:
                    run_sz = rPr.get('sz')
                    latin = rPr.find(f'{{{ns_a}}}latin')
                    if latin is not None:
                        run_font = latin.get('typeface')
                    for child in rPr:
                        if 'srgbClr' in child.tag:
                            run_clr = child.get('val')
                
                font = run_font or def_font
                sz = run_sz or def_sz
                clr = run_clr or def_clr
                pt = int(sz) / 100 if sz else None
                
                results.append({
                    'text': run.text[:40],
                    'font': font,
                    'size_pt': pt,
                    'color': clr,
                })
    return results


out = os.path.join(tempfile.gettempdir(), 'ppt_build_test')

for fname, label in [('build_mckinsey.pptx', 'McKinsey'), ('build_cyberpunk.pptx', 'Cyberpunk'), ('build_creative.pptx', 'Creative')]:
    path = os.path.join(out, fname)
    if not os.path.exists(path):
        continue
    prs = Presentation(path)
    
    all_fonts, all_colors, all_sizes = set(), set(), []
    
    print(f"\n{'='*80}")
    print(f"  {label} | {len(prs.slides)} slides | {round(os.path.getsize(path)/1024,1)} KB")
    print(f"{'='*80}")
    
    for i, slide in enumerate(prs.slides):
        props = extract_text_props(slide)
        slide_fonts = [p['font'] for p in props if p['font']]
        slide_sizes = [p['size_pt'] for p in props if p['size_pt']]
        slide_colors = [p['color'] for p in props if p['color']]
        
        all_fonts.update(slide_fonts)
        all_colors.update(slide_colors)
        all_sizes.extend(slide_sizes)
        
        mn = min(slide_sizes) if slide_sizes else "-"
        mx = max(slide_sizes) if slide_sizes else "-"
        print(f"  Slide {i}: {len(slide.shapes)} shapes | {mn}-{mx}pt | {len(set(slide_colors))} colors | fonts={sorted(set(slide_fonts))}")
        
        if i <= 1:
            for p in props[:5]:
                print(f"    \"{p['text']}\" font={p['font']} sz={p['size_pt']}pt clr=#{p['color']}")
            if len(props) > 5:
                print(f"    ... +{len(props)-5} more text runs")
    
    if all_sizes:
        print(f"\n  Font sizes: {min(all_sizes)}-{max(all_sizes)}pt ({len(set(all_sizes))} unique levels)")
    print(f"  Fonts: {sorted(all_fonts)}")
    print(f"  Colors: {sorted(all_colors)} ({len(all_colors)} unique)")
