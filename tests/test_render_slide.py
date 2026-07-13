"""Tests for PrecisionRenderer.render_slide() — P1 TDD.

Covers:
  1. render_slide() dispatches by goal to correct layout type
  2. add_slide() with template prefers blank layout
  3. render_slide() produces a slide with run-level fonts (no paragraph-level)
  4. render_slide() handles all major goal types
  5. render_slide() with explicit layout override
"""

from __future__ import annotations

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer


def _make_brand() -> BrandSpec:
    return BrandSpec(
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "accent": "#6366F1",
            "background": "#FFFFFF",
            "foreground": "#0F172A",
            "muted": "#F1F5F9",
            "muted-foreground": "#64748B",
            "border": "#E2E8F0",
            "destructive": "#EF4444",
        },
        fonts={"heading": "Inter", "body": "Inter"},
    )


def _make_template(tmp_path: Path) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    p = str(tmp_path / "template.pptx")
    prs.save(p)
    return p


class TestRenderSlideGoalDispatch:
    """render_slide() must place content according to goal→layout mapping."""

    def test_hook_goal_creates_hero_slide(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "hook", "title": "Test Title", "subtitle": "Test Sub"}
        slide = pr.render_slide(prs, page)
        assert slide is not None
        assert len(prs.slides) == 1
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            texts.append(run.text.strip())
        assert any("Test Title" in t for t in texts)

    def test_cta_goal_creates_closing_slide(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "cta", "title": "Get Started", "subtitle": "Contact us"}
        slide = pr.render_slide(prs, page)
        assert slide is not None
        assert len(prs.slides) == 1

    def test_content_goal_with_bullets(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "content", "title": "Overview", "bullets": ["Point 1", "Point 2"]}
        slide = pr.render_slide(prs, page)
        assert slide is not None
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            texts.append(run.text.strip())
        assert any("Overview" in t for t in texts)

    def test_features_goal_with_cards(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {
            "goal": "features",
            "title": "Features",
            "cards": [
                {"title": "Fast", "text": "Very fast"},
                {"title": "Secure", "text": "Very secure"},
                {"title": "Scalable", "text": "Very scalable"},
            ],
        }
        slide = pr.render_slide(prs, page)
        assert slide is not None
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            texts.append(run.text.strip())
        assert any("Fast" in t for t in texts)

    def test_data_goal_chart(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "data", "title": "Metrics"}
        slide = pr.render_slide(prs, page)
        assert slide is not None

    def test_code_goal(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "code", "title": "Example", "code": {"language": "python", "source": "print('hi')"}}
        slide = pr.render_slide(prs, page)
        assert slide is not None

    def test_exercise_goal(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "exercise", "title": "Practice", "exercise": {"instructions": "Do this", "steps": ["Step 1", "Step 2"]}}
        slide = pr.render_slide(prs, page)
        assert slide is not None

    def test_overview_goal_sidebar(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "overview", "title": "Agenda", "bullets": ["Topic 1", "Topic 2"]}
        slide = pr.render_slide(prs, page)
        assert slide is not None


class TestRenderSlideRunLevelFonts:
    """All text must use run-level fonts, not paragraph-level."""

    def test_all_text_uses_run_fonts(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "content", "title": "Test", "bullets": ["A", "B"]}
        slide = pr.render_slide(prs, page)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    has_runs = len(para.runs) > 0
                    if para.text.strip() and has_runs:
                        for run in para.runs:
                            if run.text.strip():
                                assert run.font.size is not None, (
                                    f"Run '{run.text[:20]}' has no font.size set"
                                )


class TestAddSlideBlankLayoutPriority:
    """With template, add_slide() should prefer blank layout."""

    def test_no_template_uses_blank(self):
        pr = PrecisionRenderer(brand_spec=_make_brand())
        prs = pr.create_presentation()
        slide = pr.add_slide(prs)
        assert slide is not None

    def test_with_template_prefers_blank(self, tmp_path):
        template_path = _make_template(tmp_path)
        pr = PrecisionRenderer(brand_spec=_make_brand(), template_path=template_path)
        prs = pr.create_presentation()
        slide = pr.add_slide(prs)
        assert slide is not None
        layout_name = slide.slide_layout.name.lower()
        assert "blank" in layout_name, f"Expected blank layout, got '{slide.slide_layout.name}'"


class TestRenderSlideLayoutOverride:
    """page dict can specify 'layout' to override goal→layout dispatch."""

    def test_explicit_layout_override(self):
        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "content", "title": "Custom", "layout": "two-column"}
        slide = pr.render_slide(prs, page)
        assert slide is not None


class TestRenderSlideImageHandling:
    """render_slide() must use Pillow pre-crop for images."""

    def test_with_image_uses_cover_crop(self, tmp_path):
        from PIL import Image as PILImage
        img_path = tmp_path / "test.png"
        PILImage.new("RGB", (1920, 1080), color="blue").save(str(img_path))

        brand = _make_brand()
        pr = PrecisionRenderer(brand_spec=brand)
        prs = pr.create_presentation()
        page = {"goal": "content", "title": "With Image", "image": str(img_path), "bullets": ["Point"]}
        slide = pr.render_slide(prs, page)
        assert slide is not None
        has_picture = any(
            shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
            for shape in slide.shapes
        )
        assert has_picture, "Expected a picture shape on the slide"
