"""
Tests to validate component rendering audit findings.

These tests verify the accuracy of each claim in the audit report
by directly testing the component library and renderer behavior.
"""
import sys
import pytest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptx import Presentation
from pptx.util import Emu
from lxml import etree

from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
from ppt_pro_max.enterprise.component_renderer import (
    ComponentRenderer,
    _shape_primary_slots,
)
from ppt_pro_max.enterprise.brand_spec import BrandSpec

_PN = "http://schemas.openxmlformats.org/presentationml/2006/main"
_AN = "http://schemas.openxmlformats.org/drawingml/2006/main"


@pytest.fixture
def comp_lib():
    db_path = find_db_path()
    if db_path is None:
        pytest.skip("No component library DB")
    lib = ComponentLibrary(db_path)
    yield lib
    lib.close()


@pytest.fixture
def renderer():
    return ComponentRenderer()


@pytest.fixture
def brand_spec():
    return BrandSpec(colors={
        "primary": "#2D5016",
        "accent": "#E8A838",
        "muted": "#F0F5E8",
        "foreground": "#1A2E0A",
        "on-primary": "#FFFFFF",
        "background": "#FFFEF7",
    })


def _count_sps_with_text(root):
    sps_with = 0
    total_t = 0
    for sp in root.findall(f"{{{_PN}}}sp"):
        ptxb = sp.find(f"{{{_PN}}}txBody")
        if ptxb is not None:
            found = False
            for t in ptxb.iter(f"{{{_AN}}}t"):
                if t.text and t.text.strip():
                    found = True
                    total_t += 1
            if found:
                sps_with += 1
    return sps_with, total_t


def _is_concentrated(root):
    swt = set()
    for i, sp in enumerate(root.findall(f"{{{_PN}}}sp")):
        ptxb = sp.find(f"{{{_PN}}}txBody")
        if ptxb is not None:
            for t in ptxb.iter(f"{{{_AN}}}t"):
                if t.text and t.text.strip():
                    swt.add(i)
                    break
    return len(swt) <= 1


# ============================================================
# Claim 1: Component library has both A-type (concentrated)
# and B-type (distributed) components
# ============================================================

class TestComponentStructure:
    def test_process3_has_both_types(self, comp_lib):
        results = comp_lib.search(type="group", category="process", node_count=3, limit=50)
        concentrated = 0
        distributed = 0
        for r in results:
            xp = comp_lib.load_xml(r["id"])
            if not xp or "group" not in xp:
                continue
            grp = xp["group"]
            if isinstance(grp, bytes):
                grp = grp.decode("utf-8")
            root = etree.fromstring(grp)
            if _is_concentrated(root):
                concentrated += 1
            else:
                distributed += 1
        assert concentrated > 0, "No concentrated (A-type) process/3 components found"
        assert distributed > 0, "No distributed (B-type) process/3 components found"

    def test_all_categories_have_text(self, comp_lib):
        for cat in ["process", "infographic", "hierarchy", "swot", "timeline"]:
            results = comp_lib.search(type="group", category=cat, limit=10)
            for r in results:
                xp = comp_lib.load_xml(r["id"])
                if not xp or "group" not in xp:
                    continue
                grp = xp["group"]
                if isinstance(grp, bytes):
                    grp = grp.decode("utf-8")
                root = etree.fromstring(grp)
                has_text = False
                for t in root.iter(f"{{{_AN}}}t"):
                    if t.text and t.text.strip():
                        has_text = True
                        break
                assert has_text, f"Component id={r['id']} ({cat}) has no <a:t> text"


# ============================================================
# Claim 2: match() returns A-type (concentrated) for process/3
# ============================================================

class TestMatchBehavior:
    def test_process3_match_returns_distributed(self, comp_lib):
        element = {
            "type": "group",
            "category": "process",
            "texts": ["a", "b", "c"],
            "nodes": 3,
            "node_count": 3,
            "bounds": (0.9, 1.6, 11.5, 5.0),
        }
        match = comp_lib.match(element)
        assert match is not None
        xp = comp_lib.load_xml(match["id"])
        grp = xp["group"]
        if isinstance(grp, bytes):
            grp = grp.decode("utf-8")
        root = etree.fromstring(grp)
        assert not _is_concentrated(root), (
            f"match() returned concentrated component id={match['id']} for process/3, "
            f"expected distributed (B-type)"
        )

    def test_other_categories_match_returns_distributed(self, comp_lib):
        for cat, nc in [("infographic", 3), ("timeline", 4)]:
            element = {
                "type": "group",
                "category": cat,
                "texts": ["a"] * nc,
                "nodes": nc,
                "node_count": nc,
                "bounds": (0.9, 1.6, 11.5, 5.0),
            }
            match = comp_lib.match(element)
            assert match is not None
            xp = comp_lib.load_xml(match["id"])
            grp = xp["group"]
            if isinstance(grp, bytes):
                grp = grp.decode("utf-8")
            root = etree.fromstring(grp)
            assert not _is_concentrated(root), (
                f"match() returned concentrated component for {cat}/{nc}"
            )


# ============================================================
# Claim 3: _fill_group_data on A-type concentrates text in 1 sp
# ============================================================

class TestFillGroupData:
    def test_fill_distributed_spreads_text(self, comp_lib, renderer):
        match = comp_lib.match({
            "type": "group", "category": "process",
            "texts": ["a", "b", "c"], "nodes": 3, "node_count": 3,
        })
        xp = comp_lib.load_xml(match["id"])
        new_texts = ["同质化严重", "基础设施薄弱", "季节性明显"]
        filled = renderer._fill_group_data(dict(xp), new_texts)
        grp = filled["group"]
        if isinstance(grp, bytes):
            grp = grp.decode("utf-8")
        root = etree.fromstring(grp)
        sps_with, total_t = _count_sps_with_text(root)
        assert sps_with >= 2, (
            f"B-type after fill: expected >=2 sps with text, got {sps_with}"
        )
        assert total_t >= 3, f"Expected >=3 text items, got {total_t}"


# ============================================================
# Claim 4: render_group on process/3 produces only 1 visible text child
# ============================================================

class TestRenderGroupVisibility:
    def test_process3_has_multiple_text_children(self, comp_lib, renderer, brand_spec):
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        element = {
            "type": "group", "category": "process",
            "texts": ["同质化严重", "基础设施薄弱", "季节性明显"],
            "nodes": 3, "node_count": 3,
            "bounds": (0.9, 1.6, 11.5, 5.0),
            "component_fit": "stretch",
        }
        result = renderer.render_group(slide, element, brand_spec=brand_spec, component_lib=comp_lib)
        assert result is True

        children_with_text = 0
        for shape in slide.shapes:
            if shape.shape_type == 6:
                for child in shape.shapes:
                    if child.has_text_frame and child.text_frame.text.strip():
                        children_with_text += 1
        assert children_with_text >= 2, (
            f"process/3: expected >=2 children with text (B-type), got {children_with_text}"
        )

    def test_infographic3_has_multiple_text_children(self, comp_lib, renderer, brand_spec):
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        element = {
            "type": "group", "category": "infographic",
            "texts": ["自然山水", "古村落群", "农耕文化"],
            "nodes": 3, "node_count": 3,
            "bounds": (0.9, 1.6, 11.5, 5.0),
            "component_fit": "stretch",
        }
        result = renderer.render_group(slide, element, brand_spec=brand_spec, component_lib=comp_lib)
        assert result is True

        children_with_text = 0
        for shape in slide.shapes:
            if shape.shape_type == 6:
                for child in shape.shapes:
                    if child.has_text_frame and child.text_frame.text.strip():
                        children_with_text += 1
        assert children_with_text >= 2, (
            f"infographic/3: expected >=2 children with text, got {children_with_text}"
        )


# ============================================================
# Claim 5: Fonts are missing after render_group
# ============================================================

class TestFontPreservation:
    def test_process3_fonts_replaced_by_brand(self, comp_lib, renderer, brand_spec):
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        element = {
            "type": "group", "category": "process",
            "texts": ["同质化严重", "基础设施薄弱", "季节性明显"],
            "nodes": 3, "node_count": 3,
            "bounds": (0.9, 1.6, 11.5, 5.0),
            "component_fit": "stretch",
        }
        renderer.render_group(slide, element, brand_spec=brand_spec, component_lib=comp_lib)

        st = slide._element.find(f"{{{_PN}}}cSld").find(f"{{{_PN}}}spTree")
        grps = st.findall(f"{{{_PN}}}grpSp")
        assert len(grps) > 0
        grp = grps[0]

        fonts_found = set()
        for rPr in grp.iter(f"{{{_AN}}}rPr"):
            for tag in (f"{{{_AN}}}latin", f"{{{_AN}}}ea"):
                elem = rPr.find(tag)
                if elem is not None:
                    fn = elem.get("typeface", "")
                    if fn and not fn.startswith("+"):
                        fonts_found.add(fn)
        assert len(fonts_found) > 0, (
            f"process/3: expected fonts after _replace_group_fonts, got none"
        )

    def test_infographic3_fonts_replaced_by_brand(self, comp_lib, renderer, brand_spec):
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        element = {
            "type": "group", "category": "infographic",
            "texts": ["自然山水", "古村落群", "农耕文化"],
            "nodes": 3, "node_count": 3,
            "bounds": (0.9, 1.6, 11.5, 5.0),
            "component_fit": "stretch",
        }
        renderer.render_group(slide, element, brand_spec=brand_spec, component_lib=comp_lib)

        st = slide._element.find(f"{{{_PN}}}cSld").find(f"{{{_PN}}}spTree")
        grps = st.findall(f"{{{_PN}}}grpSp")
        assert len(grps) > 0
        grp = grps[0]

        fonts_found = set()
        for rPr in grp.iter(f"{{{_AN}}}rPr"):
            for tag in (f"{{{_AN}}}latin", f"{{{_AN}}}ea"):
                elem = rPr.find(tag)
                if elem is not None:
                    fn = elem.get("typeface", "")
                    if fn and not fn.startswith("+"):
                        fonts_found.add(fn)
        assert len(fonts_found) > 0, (
            f"infographic/3: expected fonts after _replace_group_fonts, got none"
        )


# ============================================================
# Claim 6: Font sizes may be too small or too large after render
# ============================================================

class TestFontSizeRange:
    def test_all_categories_font_size_within_range(self, comp_lib, renderer, brand_spec):
        test_cases = [
            ("process", 3, ["a", "b", "c"]),
            ("infographic", 3, ["a", "b", "c"]),
            ("hierarchy", 4, ["a", "b", "c", "d"]),
            ("swot", 4, ["a", "b", "c", "d"]),
            ("timeline", 4, ["a", "b", "c", "d"]),
        ]
        for cat, nc, texts in test_cases:
            prs = Presentation()
            prs.slide_width = Emu(12192000)
            prs.slide_height = Emu(6858000)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            element = {
                "type": "group", "category": cat,
                "texts": texts, "nodes": nc, "node_count": nc,
                "bounds": (0.9, 1.6, 11.5, 5.0),
                "component_fit": "stretch",
            }
            renderer.render_group(slide, element, brand_spec=brand_spec, component_lib=comp_lib)

            st = slide._element.find(f"{{{_PN}}}cSld").find(f"{{{_PN}}}spTree")
            grps = st.findall(f"{{{_PN}}}grpSp")
            if not grps:
                continue
            grp = grps[0]

            sizes = []
            for rPr in grp.iter(f"{{{_AN}}}rPr"):
                sz = rPr.get("sz")
                if sz:
                    sizes.append(int(sz) / 100)

            if sizes:
                max_sz = max(sizes)
                assert max_sz <= 72, (
                    f"{cat}/{nc}: max font size {max_sz}pt exceeds 72pt"
                )


# ============================================================
# Claim 7: Fill colors after render - check brand compliance
# ============================================================

class TestColorCompliance:
    def test_fill_colors_are_brand_or_white(self, comp_lib, renderer, brand_spec):
        brand_colors = set(brand_spec.colors.values())
        allowed = brand_colors | {"#FFFFFF", "#ffffff"}

        test_cases = [
            ("infographic", 3, ["a", "b", "c"]),
            ("hierarchy", 4, ["a", "b", "c", "d"]),
        ]
        for cat, nc, texts in test_cases:
            prs = Presentation()
            prs.slide_width = Emu(12192000)
            prs.slide_height = Emu(6858000)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            element = {
                "type": "group", "category": cat,
                "texts": texts, "nodes": nc, "node_count": nc,
                "bounds": (0.9, 1.6, 11.5, 5.0),
                "component_fit": "stretch",
            }
            renderer.render_group(slide, element, brand_spec=brand_spec, component_lib=comp_lib)

            st = slide._element.find(f"{{{_PN}}}cSld").find(f"{{{_PN}}}spTree")
            grps = st.findall(f"{{{_PN}}}grpSp")
            if not grps:
                continue
            grp = grps[0]

            for sp in grp.findall(f".//{{{_PN}}}sp"):
                spPr = sp.find(f"{{{_PN}}}spPr")
                if spPr is not None:
                    sf = spPr.find(f"{{{_AN}}}solidFill")
                    if sf is not None:
                        clr = sf.find(f"{{{_AN}}}srgbClr")
                        if clr is not None:
                            fill = "#" + clr.get("val", "").upper()
                            assert fill in {c.upper() for c in allowed}, (
                                f"{cat}/{nc}: unexpected fill color {fill}, "
                                f"allowed: {allowed}"
                            )
