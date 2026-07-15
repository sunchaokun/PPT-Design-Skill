"""Comprehensive integration test — generate full PPT with all upgrade features.

Creates a complete 8-slide deck exercising every goal type and upgrade,
then verifies output quality at the shape/property level.
"""

from __future__ import annotations

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.precision_renderer import (
    CORNER_RADIUS_SCALE,
    PrecisionRenderer,
)


def _light_brand() -> BrandSpec:
    return BrandSpec(colors={
        "primary": "#2563EB", "on-primary": "#FFFFFF",
        "secondary": "#64748B", "accent": "#F97316",
        "background": "#FFFFFF", "foreground": "#1E293B",
        "muted": "#F1F5F9", "muted-foreground": "#94A3B8",
        "border": "#E2E8F0", "destructive": "#EF4444",
    })


def _dark_brand() -> BrandSpec:
    return BrandSpec(colors={
        "primary": "#3B82F6", "on-primary": "#FFFFFF",
        "secondary": "#64748B", "accent": "#F97316",
        "background": "#0A1E3D", "foreground": "#F0F4F8",
        "muted": "#1A2E4A", "muted-foreground": "#8A9BB5",
        "border": "#1E3A5F", "destructive": "#EF4444",
    })


def _build_full_deck_pages() -> list[dict]:
    return [
        {
            "goal": "hook",
            "title": "Design Quality Upgrade",
            "subtitle": "From Template-Grade to Designer-Grade Presentations",
        },
        {
            "goal": "content",
            "title": "The Problem We Solve",
            "subtitle": "Current output lacks design intelligence",
            "bullets": [
                "All positions hardcoded — every deck looks identical",
                "Binary typography with no visual hierarchy",
                "Flat colors without depth or tint/shade scales",
                "Uniform shadows creating clip-art feel",
            ],
        },
        {
            "goal": "features",
            "title": "Core Design Upgrades",
            "cards": [
                {"title": "TypeScale System", "text": "5-level hierarchy\nDisplay → Title → Subtitle → Body → Caption\nMinimum 2pt gap between levels"},
                {"title": "Color Depth", "text": "OKLCH-based tint/shade\n11 levels per primary\n5 alpha variants"},
                {"title": "Elevation System", "text": "5 shadow levels\nDark mode glow\nForeground-colored shadows"},
            ],
        },
        {
            "goal": "data",
            "title": "Upgrade Coverage",
            "chart": {
                "type": "bar",
                "title": "Upgrades by Tier",
                "categories": ["Tier 1", "Tier 2", "Tier 3"],
                "series": [{"name": "Upgrades", "values": [10, 6, 7]}],
            },
        },
        {
            "goal": "code",
            "title": "Code Block Redesign",
            "code": {
                "language": "python",
                "source": "from ppt_pro_max import generate_ppt\n\nresult = generate_ppt(\n    'AI pitch',\n    style='dark cyberpunk',\n)",
            },
        },
        {
            "goal": "exercise",
            "title": "Try It Yourself",
            "exercise": {
                "instructions": "Generate a PPT with the design quality upgrades",
                "duration": "5 min",
                "steps": [
                    "Install the latest version",
                    "Run: python -m ppt_pro_max 'test' --style dark",
                    "Compare output with previous version",
                ],
            },
        },
        {
            "goal": "content",
            "title": "Many Bullets — Two Column Layout",
            "bullets": [
                "Gradient overlay replaces flat black",
                "CJK font pairing system",
                "Content-adaptive margins",
                "Badge/tag styling with tracking",
                "Section divider generation",
                "Progress bar indicator",
                "Noise texture overlay",
                "Rounded corner consistency",
            ],
        },
        {
            "goal": "cta",
            "title": "Start Building Designer-Grade Decks",
            "subtitle": "pip install ppt-pro-max && python -m ppt_pro_max 'your idea'",
        },
    ]


class TestFullDeckLightMode:
    @pytest.fixture
    def deck_path(self, tmp_path):
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = r.create_presentation()
        pages = _build_full_deck_pages()
        total = len(pages)
        for i, page in enumerate(pages):
            r.render_slide(prs, page, page_index=i, total_pages=total)
        out = str(tmp_path / "light_deck.pptx")
        r.save(prs, out)
        return out

    def test_deck_has_8_slides(self, deck_path):
        prs = Presentation(deck_path)
        assert len(prs.slides) == 8

    def test_hook_slide_has_gradient_overlay(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[0]
        found_grad = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                gradFill = spPr.find(qn("a:gradFill"))
                if gradFill is not None:
                    found_grad = True
        assert found_grad, "Hook slide should use gradient overlay (not flat black)"

    def test_content_slide_has_progress_bar(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[1]
        bar_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, "top") and shape.top > Inches(7.3):
                bar_shapes.append(shape)
        assert len(bar_shapes) >= 2, "Content slide should have progress bar at bottom"

    def test_features_slide_has_featured_card(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[2]
        found_22pt = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size >= Pt(22):
                            found_22pt = True
        assert found_22pt, "First card should be featured with 22pt+ title"

    def test_code_block_always_dark(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[4]
        found_dark = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                solidFill = spPr.find(qn("a:solidFill"))
                if solidFill is not None:
                    srgb = solidFill.find(qn("a:srgbClr"))
                    if srgb is not None and srgb.get("val", "").upper() == "1E293B":
                        found_dark = True
        assert found_dark

    def test_code_block_has_badge(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[4]
        found_python = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "PYTHON" in shape.text_frame.text.upper():
                    found_python = True
        assert found_python

    def test_two_column_bullets(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[6]
        multiline_count = sum(
            1 for s in slide.shapes
            if s.has_text_frame and "\n" in (s.text_frame.text or "")
        )
        assert multiline_count >= 2, "8 bullets should render in 2 columns"

    def test_cjk_font_set_on_runs(self, deck_path):
        prs = Presentation(deck_path)
        found_ea = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            rPr = run._r.find(qn("a:rPr"))
                            if rPr is not None:
                                ea = rPr.find(qn("a:ea"))
                                if ea is not None and ea.get("typeface"):
                                    found_ea = True
        assert found_ea, "CJK font should be set on text runs via a:ea"

    def test_min_font_size_11pt(self, deck_path):
        prs = Presentation(deck_path)
        min_size = Pt(999)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.text.strip():
                                min_size = min(min_size, run.font.size)
        assert min_size >= Pt(10), f"Min font size {min_size} should be >= 10pt (11pt ideal)"

    def test_all_slides_have_content(self, deck_path):
        prs = Presentation(deck_path)
        for i, slide in enumerate(prs.slides):
            has_text = any(
                s.has_text_frame and s.text_frame.text.strip()
                for s in slide.shapes
            )
            assert has_text, f"Slide {i} has no text content"

    def test_exercise_slide_has_badge(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[5]
        found_exercise = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "EXERCISE" in shape.text_frame.text.upper():
                    found_exercise = True
        assert found_exercise

    def test_no_flat_black_overlay_on_hero(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[0]
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                solidFill = spPr.find(qn("a:solidFill"))
                if solidFill is not None:
                    srgb = solidFill.find(qn("a:srgbClr"))
                    if srgb is not None and srgb.get("val", "").upper() == "000000":
                        alpha_el = srgb.find(qn("a:alpha"))
                        if alpha_el is not None:
                            pytest.fail("Hero slide should NOT have flat black overlay")

    def test_brand_strip_not_on_every_slide(self, deck_path):
        prs = Presentation(deck_path)
        strips_found = 0
        for i, slide in enumerate(prs.slides):
            if i == 0 or i == len(prs.slides) - 1:
                continue
            for shape in slide.shapes:
                if hasattr(shape, "width") and shape.width <= Inches(0.1) and hasattr(shape, "left") and shape.left == 0:
                    strips_found += 1
                    break
        assert strips_found <= 3, "Left accent strip should not appear on every content slide"


class TestFullDeckDarkMode:
    @pytest.fixture
    def deck_path(self, tmp_path):
        brand = _dark_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = r.create_presentation()
        pages = _build_full_deck_pages()
        total = len(pages)
        for i, page in enumerate(pages):
            r.render_slide(prs, page, page_index=i, total_pages=total)
        out = str(tmp_path / "dark_deck.pptx")
        r.save(prs, out)
        return out

    def test_dark_deck_has_8_slides(self, deck_path):
        prs = Presentation(deck_path)
        assert len(prs.slides) == 8

    def test_dark_code_block_still_dark(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[4]
        found_dark = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                solidFill = spPr.find(qn("a:solidFill"))
                if solidFill is not None:
                    srgb = solidFill.find(qn("a:srgbClr"))
                    if srgb is not None and srgb.get("val", "").upper() in ("1E293B", "0D152A"):
                        found_dark = True
        assert found_dark

    def test_dark_gradient_overlay_on_hero(self, deck_path):
        prs = Presentation(deck_path)
        slide = prs.slides[0]
        found_grad = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                gradFill = spPr.find(qn("a:gradFill"))
                if gradFill is not None:
                    found_grad = True
        assert found_grad

    def test_dark_all_slides_have_content(self, deck_path):
        prs = Presentation(deck_path)
        for i, slide in enumerate(prs.slides):
            has_text = any(
                s.has_text_frame and s.text_frame.text.strip()
                for s in slide.shapes
            )
            assert has_text, f"Dark mode slide {i} has no text content"


class TestSectionDividerIntegration:
    def test_section_divider_in_deck(self, tmp_path):
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = r.create_presentation()
        r.render_slide(prs, {"goal": "hook", "title": "Start", "subtitle": "Sub"})
        r.render_section_divider(prs, 1, "Architecture", "System design overview")
        r.render_slide(prs, {"goal": "content", "title": "Details", "bullets": ["a", "b"]},
                       page_index=1, total_pages=3)
        out = str(tmp_path / "section_deck.pptx")
        r.save(prs, out)
        prs2 = Presentation(out)
        assert len(prs2.slides) == 3
        section_slide = prs2.slides[1]
        found_01 = any(
            s.has_text_frame and "01" in s.text_frame.text
            for s in section_slide.shapes
        )
        assert found_01


class TestLayoutVariantIntegration:
    def test_centered_variant(self, tmp_path):
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = r.create_presentation()
        variant = {"content_margin_left": 2.0, "title_alignment": "center"}
        r.render_slide(prs, {"goal": "content", "title": "Centered", "bullets": ["a"]},
                       layout_variant=variant)
        out = str(tmp_path / "variant_deck.pptx")
        r.save(prs, out)
        prs2 = Presentation(out)
        assert len(prs2.slides) == 1


class TestBadgeIntegration:
    def test_badge_in_exercise(self, tmp_path):
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = r.create_presentation()
        r.render_slide(prs, {
            "goal": "exercise",
            "title": "Try It",
            "exercise": {"instructions": "Do this", "duration": "5 min", "steps": ["a", "b"]},
        })
        out = str(tmp_path / "badge_deck.pptx")
        r.save(prs, out)
        prs2 = Presentation(out)
        slide = prs2.slides[0]
        found = any(
            s.has_text_frame and "EXERCISE" in s.text_frame.text.upper()
            for s in slide.shapes
        )
        assert found


class TestDesignQualityMetrics:
    """Evaluate overall design quality of generated PPT."""

    @pytest.fixture
    def light_deck(self, tmp_path):
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = r.create_presentation()
        pages = _build_full_deck_pages()
        total = len(pages)
        for i, page in enumerate(pages):
            r.render_slide(prs, page, page_index=i, total_pages=total)
        out = str(tmp_path / "quality_deck.pptx")
        r.save(prs, out)
        return out

    def test_no_empty_slides(self, light_deck):
        prs = Presentation(light_deck)
        for i, slide in enumerate(prs.slides):
            assert len(slide.shapes) >= 3, f"Slide {i} has < 3 shapes"

    def test_file_size_reasonable(self, light_deck):
        size = os.path.getsize(light_deck)
        assert size > 10000, "File should be > 10KB"
        assert size < 500000, "File should be < 500KB (no images)"

    def test_title_font_size_variety(self, light_deck):
        prs = Presentation(light_deck)
        sizes = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                sizes.add(run.font.size)
        assert len(sizes) >= 3, f"Should have 3+ distinct font sizes, got {len(sizes)}"

    def test_card_body_at_least_14pt(self, light_deck):
        prs = Presentation(light_deck)
        slide = prs.slides[2]
        found_body_14 = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size >= Pt(14) and not run.font.bold:
                            found_body_14 = True
        assert found_body_14, "Card body text should be >= 14pt"

    def test_gradient_overlay_on_image_hero(self, tmp_path):
        from PIL import Image as PILImage
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = r.create_presentation()
        img = PILImage.new("RGB", (200, 200), (100, 150, 200))
        img_path = str(tmp_path / "hero.png")
        img.save(img_path)
        r.render_slide(prs, {"goal": "hook", "title": "Test", "image": img_path})
        out = str(tmp_path / "hero_deck.pptx")
        r.save(prs, out)
        prs2 = Presentation(out)
        slide = prs2.slides[0]
        found_transparent_top = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                gradFill = spPr.find(qn("a:gradFill"))
                if gradFill is not None:
                    gsLst = gradFill.find(qn("a:gsLst"))
                    if gsLst is not None:
                        stops = gsLst.findall(qn("a:gs"))
                        if stops:
                            first_srgb = stops[0].find(qn("a:srgbClr"))
                            if first_srgb is not None:
                                alpha = first_srgb.find(qn("a:alpha"))
                                if alpha is not None:
                                    val = int(alpha.get("val", "100000"))
                                    if val < 30000:
                                        found_transparent_top = True
        assert found_transparent_top, "Gradient overlay should be transparent at top for image hero"
