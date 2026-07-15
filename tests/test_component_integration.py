"""Tests for Component Library Integration (P15): catalog, find_db_path, query API,
content.json component fields, PrecisionRenderer component_lib, generate_ppt() integration,
beautify_mode.

TDD: Tests define expected behavior before implementation.
"""
import json
import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


# ═══════════════════════════════════════════════════════════════
# Step 1: ComponentLibrary.catalog() + find_db_path()
# ═══════════════════════════════════════════════════════════════

class TestCatalog:
    def test_catalog_empty_library(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)
        catalog = lib.catalog()
        assert isinstance(catalog, dict)
        assert catalog == {}

    def test_catalog_groups_by_type_and_category(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        lib.add(type="group", category="process", variant="chevron", node_count=4, xml_parts={"group": b"<g/>"})
        lib.add(type="group", category="process", variant="arrows", node_count=5, xml_parts={"group": b"<g2/>"})
        lib.add(type="group", category="swot", variant="grid", node_count=4, xml_parts={"group": b"<g3/>"})
        lib.add(type="smartart", category="hierarchy", variant="orgchart", node_count=6, xml_parts={"data": b"<d/>", "layout": b"<l/>", "colors": b"<c/>", "quickStyle": b"<q/>"})

        catalog = lib.catalog()
        assert "group" in catalog
        assert "smartart" in catalog
        assert "process" in catalog["group"]
        assert "swot" in catalog["group"]
        assert "hierarchy" in catalog["smartart"]

    def test_catalog_includes_count_and_node_range(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        lib.add(type="group", category="process", variant="v1", node_count=3, xml_parts={"group": b"<g1/>"})
        lib.add(type="group", category="process", variant="v2", node_count=7, xml_parts={"group": b"<g2/>"})
        lib.add(type="group", category="process", variant="v3", node_count=5, xml_parts={"group": b"<g3/>"})

        catalog = lib.catalog()
        proc = catalog["group"]["process"]
        assert proc["count"] == 3
        assert proc["min_nodes"] == 3
        assert proc["max_nodes"] == 7

    def test_catalog_caches_result(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        lib.add(type="group", category="process", variant="v1", node_count=4, xml_parts={"group": b"<g/>"})

        c1 = lib.catalog()
        c2 = lib.catalog()
        assert c1 is c2

    def test_catalog_invalidates_on_add(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        lib.add(type="group", category="process", variant="v1", node_count=4, xml_parts={"group": b"<g/>"})
        c1 = lib.catalog()
        assert c1["group"]["process"]["count"] == 1

        lib.add(type="group", category="process", variant="v2", node_count=5, xml_parts={"group": b"<g2/>"})
        c2 = lib.catalog()
        assert c2["group"]["process"]["count"] == 2


class TestFindDbPath:
    def test_find_db_path_explicit(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import find_db_path

        db_path = str(tmp_path / "custom.db")
        with open(db_path, "w") as f:
            f.write("")

        result = find_db_path(component_library=db_path)
        assert result == db_path

    def test_find_db_path_project_dir(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import find_db_path

        cl_dir = tmp_path / "component_library"
        cl_dir.mkdir()
        db_path = cl_dir / "index.db"
        db_path.write_text("")

        result = find_db_path(project_dir=str(tmp_path))
        assert result == str(db_path)

    def test_find_db_path_returns_none_when_not_found(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import find_db_path

        result = find_db_path(project_dir=str(tmp_path))
        assert result is None

    def test_find_db_path_explicit_nonexistent_falls_through(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import find_db_path

        result = find_db_path(component_library="/nonexistent/path.db", project_dir=str(tmp_path))
        assert result is None


# ═══════════════════════════════════════════════════════════════
# Step 2: query_component_library() public API
# ═══════════════════════════════════════════════════════════════

class TestQueryComponentLibrary:
    def test_query_catalog_returns_dict(self, tmp_path, monkeypatch):
        from ppt_pro_max import query_component_library
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="process", variant="v1", node_count=4, xml_parts={"group": b"<g/>"})
        lib.close()

        result = query_component_library(component_library=db_path)
        assert isinstance(result, dict)
        assert "group" in result

    def test_query_search_returns_list(self, tmp_path):
        from ppt_pro_max import query_component_library
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="process", variant="chevron", node_count=4, xml_parts={"group": b"<g/>"})
        lib.close()

        result = query_component_library(type="group", category="process", component_library=db_path)
        assert isinstance(result, list)
        assert len(result) >= 1
        assert result[0]["category"] == "process"

    def test_query_search_with_node_count(self, tmp_path):
        from ppt_pro_max import query_component_library
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="process", variant="v1", node_count=4, xml_parts={"group": b"<g1/>"})
        lib.add(type="group", category="process", variant="v2", node_count=5, xml_parts={"group": b"<g2/>"})
        lib.close()

        result = query_component_library(type="group", category="process", node_count=4, component_library=db_path)
        assert len(result) == 1
        assert result[0]["node_count"] == 4

    def test_query_no_db_returns_empty(self):
        from ppt_pro_max import query_component_library

        result = query_component_library(component_library="/nonexistent/db.db")
        assert result == {} or result == []


# ═══════════════════════════════════════════════════════════════
# Step 3: content.json component fields
# ═══════════════════════════════════════════════════════════════

class TestContentParserComponentFields:
    def test_load_enterprise_content_parses_component_type(self, tmp_path):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content

        content = {
            "slides": [
                {
                    "goal": "content",
                    "title": "开发流程",
                    "bullets": ["需求", "设计", "开发", "测试"],
                    "component_type": "group",
                    "component_category": "process",
                    "component_variant": "chevron",
                }
            ]
        }
        result = load_enterprise_content(content, str(tmp_path))
        assert len(result) == 1
        assert result[0]["component_type"] == "group"
        assert result[0]["component_category"] == "process"
        assert result[0]["component_variant"] == "chevron"

    def test_load_enterprise_content_without_component_fields(self, tmp_path):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content

        content = {
            "slides": [
                {
                    "goal": "content",
                    "title": "普通页面",
                    "bullets": ["A", "B"],
                }
            ]
        }
        result = load_enterprise_content(content, str(tmp_path))
        assert "component_type" not in result[0] or result[0].get("component_type") is None

    def test_build_page_designs_copies_component_fields(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        from ppt_pro_max.enterprise.enterprise_decider import EnterpriseDesignDecider

        pipeline = EnterprisePipeline()
        brand_spec = BrandSpec(source="test")
        decider = EnterpriseDesignDecider(brand_spec=brand_spec)

        class FakeAsset:
            template_path = None
            image_pool = []

        story_plan = {
            "pages": [
                {
                    "goal": "content",
                    "title": "流程",
                    "bullets": ["A", "B", "C"],
                    "component_type": "group",
                    "component_category": "process",
                    "component_variant": "chevron",
                }
            ]
        }

        designs = pipeline._build_page_designs(story_plan, decider, FakeAsset(), brand_spec)
        assert designs[0]["component_type"] == "group"
        assert designs[0]["component_category"] == "process"
        assert designs[0]["component_variant"] == "chevron"


# ═══════════════════════════════════════════════════════════════
# Step 4: PrecisionRenderer component_lib integration
# ═══════════════════════════════════════════════════════════════

class TestPrecisionRendererComponentLib:
    def test_render_slide_accepts_component_lib(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        page = {"goal": "content", "title": "Test", "bullets": ["A", "B"]}

        slide = renderer.render_slide(prs, page, component_lib=None)
        assert slide is not None

    def test_render_slide_with_component_type_uses_renderer(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        group_xml = b'''<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvGrpSpPr><p:cNvPr id="100" name="TestGroup"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm><a:off x="914400" y="914400"/><a:ext cx="9144000" cy="4572000"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="4572000"/></a:xfrm></p:grpSpPr>
  <p:sp><p:nvSpPr><p:cNvPr id="101" name="S1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Step 1</a:t></a:r></a:p></p:txBody></p:sp>
  <p:sp><p:nvSpPr><p:cNvPr id="102" name="S2"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="3200400" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Step 2</a:t></a:r></a:p></p:txBody></p:sp>
</p:grpSp>'''

        lib.add(type="group", category="process", variant="chevron", node_count=2, xml_parts={"group": group_xml})

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        page = {
            "goal": "content",
            "title": "流程",
            "bullets": ["需求分析", "设计"],
            "component_type": "group",
            "component_category": "process",
        }

        slide = renderer.render_slide(prs, page, component_lib=lib)
        assert slide is not None

    def test_render_slide_component_type_priority_over_diagram(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        group_xml = b'''<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvGrpSpPr><p:cNvPr id="100" name="TestGroup"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
  <p:grpSpPr><a:xfrm><a:off x="914400" y="914400"/><a:ext cx="9144000" cy="4572000"/>
    <a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="4572000"/></a:xfrm></p:grpSpPr>
  <p:sp><p:nvSpPr><p:cNvPr id="101" name="S1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="2743200" cy="1371600"/></a:xfrm></p:spPr>
    <p:txBody><a:bodyPr/><a:p><a:r><a:t>Step 1</a:t></a:r></a:p></p:txBody></p:sp>
</p:grpSp>'''

        lib.add(type="group", category="process", variant="chevron", node_count=1, xml_parts={"group": group_xml})

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        page = {
            "goal": "content",
            "title": "流程",
            "bullets": ["Step A"],
            "component_type": "group",
            "component_category": "process",
            "diagram_type": "flowchart",
            "diagram_data": {"items": [{"label": "Step A"}]},
        }

        slide = renderer.render_slide(prs, page, component_lib=lib)
        assert slide is not None

    def test_render_slide_component_type_fallback_to_bullets(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        page = {
            "goal": "content",
            "title": "流程",
            "bullets": ["A", "B", "C"],
            "component_type": "group",
            "component_category": "nonexistent_category",
        }

        slide = renderer.render_slide(prs, page, component_lib=lib)
        assert slide is not None
        all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text)
        combined = " ".join(all_text)
        assert "A" in combined

    def test_render_slide_cards_priority_over_component_type(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        page = {
            "goal": "content",
            "title": "Features",
            "bullets": ["A", "B"],
            "cards": [{"title": "Card1", "text": "Detail1"}],
            "component_type": "group",
            "component_category": "process",
        }

        slide = renderer.render_slide(prs, page, component_lib=lib)
        assert slide is not None
        all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text)
        combined = " ".join(all_text)
        assert "Card1" in combined


# ═══════════════════════════════════════════════════════════════
# Step 5: generate_ppt() component_library parameter
# ═══════════════════════════════════════════════════════════════

class TestGeneratePptComponentLibrary:
    def test_generate_ppt_accepts_component_library_param(self, tmp_path):
        from ppt_pro_max import generate_ppt

        db_path = str(tmp_path / "test.db")
        with open(db_path, "w") as f:
            f.write("")

        result = generate_ppt(
            "test presentation",
            project=str(tmp_path),
            component_library=db_path,
            dry_run=True,
        )
        assert result.get("dry_run") is True

    def test_generate_ppt_content_dict_param(self, tmp_path):
        from ppt_pro_max import generate_ppt

        project_dir = tmp_path / "project"
        project_dir.mkdir()

        result = generate_ppt(
            "test",
            project=str(project_dir),
            content={
                "slides": [
                    {"goal": "hook", "title": "Test Title"},
                    {"goal": "content", "title": "Content", "bullets": ["A", "B"]},
                ]
            },
            dry_run=True,
        )
        assert result.get("dry_run") is True


# ═══════════════════════════════════════════════════════════════
# Step 6: Beautify mode
# ═══════════════════════════════════════════════════════════════

class TestBeautifyMode:
    def test_beautify_light_mode(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Light Mode Test"
        run.font.size = Pt(36)

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        prs.save(input_path)

        pipeline = EnterprisePipeline()
        result = pipeline.run_beautify(
            input_pptx=input_path,
            output_path=output_path,
            style="professional",
            beautify_mode="light",
        )

        assert os.path.isfile(output_path)
        assert result["mode"] == "beautify"
        assert result.get("beautify_mode") == "light"

    def test_beautify_full_mode(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Full Mode Test"
        run.font.size = Pt(36)

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        prs.save(input_path)

        pipeline = EnterprisePipeline()
        result = pipeline.run_beautify(
            input_pptx=input_path,
            output_path=output_path,
            style="professional",
            beautify_mode="full",
        )

        assert os.path.isfile(output_path)
        assert result["mode"] == "beautify"
        assert result.get("beautify_mode") == "full"

    def test_beautify_default_is_full(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Default Mode"
        run.font.size = Pt(36)

        input_path = str(tmp_path / "input.pptx")
        output_path = str(tmp_path / "output.pptx")
        prs.save(input_path)

        pipeline = EnterprisePipeline()
        result = pipeline.run_beautify(
            input_pptx=input_path,
            output_path=output_path,
            style="professional",
        )

        assert result.get("beautify_mode") == "full"

    def test_generate_ppt_beautify_mode_param(self, tmp_path):
        from ppt_pro_max import generate_ppt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "API Test"
        run.font.size = Pt(36)

        input_path = str(tmp_path / "input.pptx")
        prs.save(input_path)

        result = generate_ppt("美化", beautify=input_path, beautify_mode="light", output=str(tmp_path / "out.pptx"))
        assert result.get("beautify_mode") == "light"


# ═══════════════════════════════════════════════════════════════
# Step 6b: Component category inference (beautify full mode)
# ═══════════════════════════════════════════════════════════════

class TestComponentCategoryInference:
    def test_sequential_bullets_infer_process(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        result = infer_component_category(["需求分析", "设计", "开发", "测试"])
        assert result == ("group", "process")

    def test_four_items_with_swot_keywords(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        result = infer_component_category(["优势", "劣势", "机会", "威胁"])
        assert result == ("group", "swot")

    def test_few_bullets_no_inference(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        result = infer_component_category(["Only one"])
        assert result == (None, None)

    def test_too_many_bullets_no_inference(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        result = infer_component_category([f"Item {i}" for i in range(10)])
        assert result == (None, None)

    def test_hierarchy_keywords(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        result = infer_component_category(["CEO", "CTO", "CFO", "VP Engineering", "VP Sales"])
        assert result[1] == "hierarchy"


# ═══════════════════════════════════════════════════════════════
# Review: additional edge case tests
# ═══════════════════════════════════════════════════════════════

class TestCatalogCacheInvalidation:
    def test_catalog_invalidates_on_remove(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        cid = lib.add(type="group", category="process", variant="v1", node_count=4, xml_parts={"group": b"<g/>"})
        assert lib.catalog()["group"]["process"]["count"] == 1

        lib.remove(cid)
        assert lib.catalog() == {}

    def test_catalog_invalidates_on_bulk_import(self, tmp_path):
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test.db")
        lib = ComponentLibrary(db_path)

        assert lib.catalog() == {}

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "No SmartArt"
        run.font.size = Pt(36)

        pptx_path = str(tmp_path / "no_smartart.pptx")
        prs.save(pptx_path)

        lib.bulk_import([pptx_path])
        catalog = lib.catalog()
        assert isinstance(catalog, dict)


class TestBeautifyLightError:
    def test_beautify_light_error_includes_mode(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline

        pipeline = EnterprisePipeline()
        result = pipeline.run_beautify(
            input_pptx="/nonexistent/file.pptx",
            beautify_mode="light",
        )
        assert result["mode"] == "beautify"
        assert "error" in result


class TestContentDictCleanup:
    def test_content_dict_creates_valid_json(self, tmp_path):
        from ppt_pro_max import generate_ppt

        project_dir = tmp_path / "project"
        project_dir.mkdir()

        content_dict = {
            "slides": [
                {"goal": "hook", "title": "Test Title"},
                {"goal": "content", "title": "Content Page", "bullets": ["A", "B"]},
            ]
        }

        result = generate_ppt(
            "test",
            project=str(project_dir),
            content=content_dict,
            dry_run=True,
        )
        assert result.get("dry_run") is True


class TestComponentTypeNone:
    def test_component_type_none_ignored(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        page = {
            "goal": "content",
            "title": "Test",
            "bullets": ["A", "B"],
            "component_type": None,
        }

        slide = renderer.render_slide(prs, page, component_lib=None)
        assert slide is not None
        all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text)
        combined = " ".join(all_text)
        assert "A" in combined

    def test_component_type_without_component_lib_ignored(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer

        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        page = {
            "goal": "content",
            "title": "Test",
            "bullets": ["A", "B"],
            "component_type": "group",
            "component_category": "process",
        }

        slide = renderer.render_slide(prs, page, component_lib=None)
        assert slide is not None
        all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text)
        combined = " ".join(all_text)
        assert "A" in combined
