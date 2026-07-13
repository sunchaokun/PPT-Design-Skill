"""Tests for ProposalGenerator — 2-3 style preview PPTs for user selection."""

from __future__ import annotations

import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.enterprise.proposal_generator import ProposalGenerator


def _make_template(tmp_path: Path) -> str:
    project = tmp_path / "project"
    project.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    p = str(project / "template.pptx")
    prs.save(p)
    return str(project)


class TestProposalGenerator:

    def test_generate_returns_three_proposals(self, tmp_path):
        gen = ProposalGenerator()
        results = gen.generate("AI startup pitch", output_dir=str(tmp_path))
        assert len(results) == 3

    def test_each_proposal_has_pptx_file(self, tmp_path):
        gen = ProposalGenerator()
        results = gen.generate("AI startup pitch", output_dir=str(tmp_path))
        for r in results:
            assert os.path.isfile(r["path"]), f"Missing file: {r['path']}"

    def test_each_proposal_has_4_slides(self, tmp_path):
        gen = ProposalGenerator()
        results = gen.generate("AI startup pitch", output_dir=str(tmp_path))
        for r in results:
            prs = Presentation(r["path"])
            assert len(prs.slides) == 4, f"Proposal {r['id']} has {len(prs.slides)} slides, expected 4"

    def test_proposals_have_different_palettes(self, tmp_path):
        gen = ProposalGenerator()
        results = gen.generate("AI startup pitch", output_dir=str(tmp_path))
        palettes = [r["atoms"]["palette"] for r in results]
        assert len(set(palettes)) >= 2, f"Palettes not differentiated: {palettes}"

    def test_proposal_ids_are_A_B_C(self, tmp_path):
        gen = ProposalGenerator()
        results = gen.generate("AI startup pitch", output_dir=str(tmp_path))
        ids = [r["id"] for r in results]
        assert ids == ["A", "B", "C"]

    def test_proposal_a_matches_style(self, tmp_path):
        gen = ProposalGenerator()
        results = gen.generate("dark cyberpunk pitch", output_dir=str(tmp_path))
        moods_a = results[0]["atoms"]["moods"]
        assert any(m in moods_a for m in ("dark", "tech")), f"Proposal A moods {moods_a} don't match 'dark cyberpunk'"

    def test_custom_output_dir(self, tmp_path):
        custom = tmp_path / "custom_output"
        custom.mkdir()
        gen = ProposalGenerator()
        results = gen.generate("Test query", output_dir=str(custom))
        for r in results:
            assert str(custom) in r["path"]
            assert os.path.isfile(r["path"])

    def test_seed_reproducible(self, tmp_path):
        gen = ProposalGenerator()
        r1 = gen.generate("Test query", seed=42, output_dir=str(tmp_path / "run1"))
        r2 = gen.generate("Test query", seed=42, output_dir=str(tmp_path / "run2"))
        for a, b in zip(r1, r2):
            assert a["atoms"]["palette"] == b["atoms"]["palette"]
            assert a["atoms"]["fonts"] == b["atoms"]["fonts"]
            assert a["atoms"]["decoration"] == b["atoms"]["decoration"]
            assert a["atoms"]["layout"] == b["atoms"]["layout"]

    def test_empty_style_defaults_professional(self, tmp_path):
        gen = ProposalGenerator()
        results = gen.generate("Professional business report", style=None, output_dir=str(tmp_path))
        assert len(results) == 3
        moods_a = results[0]["atoms"]["moods"]
        assert "professional" in moods_a

    def test_with_project_dir(self, tmp_path):
        project_dir = _make_template(tmp_path)
        gen = ProposalGenerator()
        results = gen.generate("Test query", project_dir=project_dir, output_dir=str(tmp_path / "out"))
        assert len(results) == 3
        for r in results:
            assert os.path.isfile(r["path"])

    def test_each_proposal_has_atoms_dict(self, tmp_path):
        gen = ProposalGenerator()
        results = gen.generate("Test query", output_dir=str(tmp_path))
        for r in results:
            atoms = r["atoms"]
            assert "palette" in atoms
            assert "fonts" in atoms
            assert "decoration" in atoms
            assert "layout" in atoms
            assert "moods" in atoms
            assert isinstance(atoms["moods"], list)

    def test_generate_creates_output_dir(self, tmp_path):
        new_dir = str(tmp_path / "nonexistent" / "subdir")
        gen = ProposalGenerator()
        results = gen.generate("Test query", output_dir=new_dir)
        assert os.path.isdir(new_dir)
        for r in results:
            assert os.path.isfile(r["path"])
