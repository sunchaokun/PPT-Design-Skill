"""Generate PPT with rebuilt library components — prioritize 3D+pic components."""

import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.component_library import ComponentLibrary, find_db_path
from ppt_pro_max.enterprise.component_renderer import ComponentRenderer
from lxml import etree

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

GREEN_BRAND = BrandSpec(colors={
    "primary": "#00AA00",
    "accent": "#66CC00",
    "muted": "#E6F5E6",
    "foreground": "#1A3A1A",
    "on-primary": "#FFFFFF",
    "background": "#FFFFFF",
})

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)

lib = ComponentLibrary(find_db_path())
renderer = ComponentRenderer()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Select components: prefer ones with 3D/pic/images
targets = [
    ("hierarchy", ["CEO", "CTO", "CFO", "VP Eng", "VP Sales"]),
    ("process", ["Plan", "Design", "Build", "Launch"]),
    ("swot", ["Strengths", "Weaknesses", "Opportunities", "Threats"]),
    ("infographic", ["Revenue", "Growth", "Profit", "Market Share"]),
    ("timeline", ["Q1", "Q2", "Q3", "Q4"]),
    ("chart", ["Sales", "Revenue", "Profit"]),
]

for cat, texts in targets:
    # Find best component (highest fill_score)
    results = lib.search(type="group", category=cat, limit=50)
    best = None
    best_score = -1
    for comp in results:
        xml_parts = lib.load_xml(comp["id"])
        root = etree.fromstring(xml_parts["group"])
        has_3d = len(list(root.iter(f"{{{_A}}}scene3d"))) > 0
        if has_3d:
            continue
        has_pic = len(list(root.iter(f"{{{_P}}}pic"))) > 0
        img_keys = [k for k in xml_parts if k.startswith("img_")]
        custGeom = len(list(root.iter(f"{{{_A}}}custGeom")))
        sp_count = len(list(root.iter(f"{{{_P}}}sp")))
        cxn_count = len(list(root.iter(f"{{{_P}}}cxnSp")))
        solidFill = sum(1 for _ in root.iter(f"{{{_A}}}solidFill"))
        score = custGeom * 2 + solidFill + cxn_count * 2 + len(img_keys) * 5 + has_pic * 10
        if score > best_score:
            best_score = score
            best = (comp, xml_parts, root)

    if best is None:
        print(f"  {cat}: SKIPPED")
        continue

    comp, xml_parts, root = best
    filled = renderer._fill_group_data(xml_parts, texts)
    styled = renderer._apply_brand_colors(filled, GREEN_BRAND)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{cat} (3d={best_score >= 1}, pic={best_score >= 10}) — Green Brand"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0xAA, 0x00)

    meta_bytes = styled.get("_meta")
    orig_aspect = 1.0
    if meta_bytes:
        meta = json.loads(meta_bytes)
        obe = meta.get("orig_bounds_emu")
        if obe and len(obe) == 4 and obe[3] > 0:
            orig_aspect = obe[2] / obe[3]

    content_area = (0.5, 1.2, 12.3, 5.5)
    bounds = renderer.compute_component_bounds(content_area, orig_aspect, "contain")
    renderer._inject_group_to_slide(slide, styled, bounds)

    has_3d = best_score >= 1
    has_pic = best_score >= 10
    print(f"  {cat}: id={comp['id']}, aspect={orig_aspect:.2f}, 3d={has_3d}, pic={has_pic}, score={best_score}")

out = os.path.join(OUT_DIR, "rebuilt_library_green.pptx")
prs.save(out)
print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")
print(f"Slides: {len(prs.slides)}")

lib.close()
