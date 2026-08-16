"""Tests for svg_compiler._dash — stroke dash, linecap, linejoin.

Validates:
- parse_stroke_style: dasharray, linecap, linejoin, miterlimit, width
- apply_stroke_style: preset dash, custom dash, linecap, linejoin
- _match_dash_preset: pattern matching
"""
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.renderer.svg_compiler._dash import (
    StrokeStyle,
    _match_dash_preset,
    apply_stroke_style,
    parse_stroke_style,
)


def _make_svg_el(svg_str):
    return etree.fromstring(svg_str)


class TestParseStrokeStyle:
    def test_defaults(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg"/>')
        s = parse_stroke_style(el)
        assert s.dash_array is None
        assert s.linecap == "flat"
        assert s.linejoin == "miter"
        assert s.miterlimit == 4.0
        assert s.width == 1.0

    def test_dash_array(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-dasharray="4,4"/>')
        s = parse_stroke_style(el)
        assert s.dash_array == [4.0, 4.0]

    def test_dash_array_spaces(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-dasharray="6 3 2 3"/>')
        s = parse_stroke_style(el)
        assert s.dash_array == [6.0, 3.0, 2.0, 3.0]

    def test_dash_array_none(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-dasharray="none"/>')
        s = parse_stroke_style(el)
        assert s.dash_array is None

    def test_linecap_butt(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-linecap="butt"/>')
        s = parse_stroke_style(el)
        assert s.linecap == "flat"

    def test_linecap_round(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-linecap="round"/>')
        s = parse_stroke_style(el)
        assert s.linecap == "round"

    def test_linecap_square(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-linecap="square"/>')
        s = parse_stroke_style(el)
        assert s.linecap == "square"

    def test_linejoin_miter(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-linejoin="miter"/>')
        s = parse_stroke_style(el)
        assert s.linejoin == "miter"

    def test_linejoin_bevel(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-linejoin="bevel"/>')
        s = parse_stroke_style(el)
        assert s.linejoin == "bevel"

    def test_miterlimit(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-miterlimit="8"/>')
        s = parse_stroke_style(el)
        assert s.miterlimit == 8.0

    def test_stroke_width(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-width="2.5"/>')
        s = parse_stroke_style(el)
        assert s.width == 2.5

    def test_stroke_width_px(self):
        el = _make_svg_el('<rect xmlns="http://www.w3.org/2000/svg" stroke-width="3px"/>')
        s = parse_stroke_style(el)
        assert s.width == 3.0


class TestMatchDashPreset:
    def test_simple_dash(self):
        assert _match_dash_preset([4, 4]) == "dash"

    def test_dot(self):
        assert _match_dash_preset([1, 2]) == "dot"

    def test_dash_dot(self):
        result = _match_dash_preset([4, 2, 1, 2])
        assert result is not None

    def test_no_match(self):
        assert _match_dash_preset([7, 3, 1, 5, 2, 4]) is None

    def test_empty(self):
        assert _match_dash_preset([]) is None


class TestApplyStrokeStyle:
    def _make_shape(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sh = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(3), Inches(2))
        return sh

    def test_preset_dash_applied(self):
        sh = self._make_shape()
        style = StrokeStyle(dash_array=[4, 4])
        apply_stroke_style(sh, style)
        xml = sh._element.xml
        assert "prstDash" in xml

    def test_custom_dash_applied(self):
        sh = self._make_shape()
        style = StrokeStyle(dash_array=[7, 3, 1, 5, 2, 4])
        apply_stroke_style(sh, style)
        xml = sh._element.xml
        assert "custDash" in xml

    def test_round_cap_applied(self):
        sh = self._make_shape()
        style = StrokeStyle(linecap="round")
        apply_stroke_style(sh, style)
        xml = sh._element.xml
        assert 'cap="round"' in xml

    def test_miter_join_with_limit_applied(self):
        sh = self._make_shape()
        style = StrokeStyle(linejoin="miter", miterlimit=6.0)
        apply_stroke_style(sh, style)
        xml = sh._element.xml
        assert "miter" in xml
        assert "lim" in xml

    def test_default_miter_join_not_written(self):
        sh = self._make_shape()
        style = StrokeStyle()
        apply_stroke_style(sh, style)
        xml = sh._element.xml
        assert "miter" not in xml

    def test_bevel_join_applied(self):
        sh = self._make_shape()
        style = StrokeStyle(linejoin="bevel")
        apply_stroke_style(sh, style)
        xml = sh._element.xml
        assert "bevel" in xml

    def test_no_element_attribute_skips(self):
        style = StrokeStyle()
        result = apply_stroke_style(object(), style)
        assert result is None

    def test_flat_cap_not_written(self):
        sh = self._make_shape()
        style = StrokeStyle(linecap="flat")
        apply_stroke_style(sh, style)
        xml = sh._element.xml
        assert 'cap=' not in xml


class TestStrokeStyleIntegration:
    """E2E test to verify stroke styles are wired up and rendered correctly."""

    def test_stroke_dasharray_rect_integration(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect x="10" y="10" width="80" height="80" fill="none" '
            'stroke="#FF0000" stroke-width="2" stroke-dasharray="4,4"/>'
            '</svg>'
        )
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from ppt_pro_max.renderer.svg_compiler import SVGCompiler
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count >= 1

        # Retrieve shape and inspect XML
        rect_shape = slide.shapes[0]
        xml = rect_shape._element.xml
        assert "prstDash" in xml or "custDash" in xml

    def test_stroke_linejoin_bevel_integration(self):
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<polyline points="10,10 50,50 90,10" fill="none" '
            'stroke="#000" stroke-width="3" stroke-linejoin="bevel"/>'
            '</svg>'
        )
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from ppt_pro_max.renderer.svg_compiler import SVGCompiler
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 5, 5))
        assert result.shape_count >= 1

        poly_shape = slide.shapes[0]
        xml = poly_shape._element.xml
        assert "bevel" in xml

