"""Extract 3D components directly from source PPT and inject - ZERO modification."""

import os
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

OUT_DIR = os.path.join(os.path.dirname(__file__), "_test_output")
os.makedirs(OUT_DIR, exist_ok=True)

def_dir = r"E:\BaiduNetdiskDownload\DEF"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_count = 0
found_3d = 0
found_no3d = 0

for root_dir, dirs, files in os.walk(def_dir):
    for fname in files:
        if not fname.endswith(".pptx"):
            continue
        if slide_count >= 10:
            break
        fpath = os.path.join(root_dir, fname)
        try:
            with zipfile.ZipFile(fpath) as z:
                for sf in [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml") and "rels" not in n]:
                    if slide_count >= 10:
                        break
                    slide_root = etree.fromstring(z.read(sf))
                    sp_tree = slide_root.find(f".//{{{_P}}}spTree")
                    if sp_tree is None:
                        continue
                    for grp in sp_tree.iter(f"{{{_P}}}grpSp"):
                        t_count = len([t for t in grp.iter(f"{{{_A}}}t") if t.text and t.text.strip()])
                        if t_count < 3:
                            continue

                        has_3d = len(list(grp.iter(f"{{{_A}}}scene3d"))) > 0

                        # Alternate: get 3D and non-3D components
                        if has_3d and found_3d >= 5:
                            continue
                        if not has_3d and found_no3d >= 5:
                            continue

                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        target_sp_tree = slide._element.find(f".//{{{_P}}}spTree")

                        grp_copy = etree.fromstring(etree.tostring(grp))

                        next_id = 2
                        for sp in target_sp_tree:
                            for cNv in sp.iter(f"{{{_P}}}cNvPr"):
                                try:
                                    next_id = max(next_id, int(cNv.get("id", "1")) + 1)
                                except (ValueError, TypeError):
                                    pass
                        for cNv in grp_copy.iter(f"{{{_P}}}cNvPr"):
                            cNv.set("id", str(next_id))
                            cNv.set("name", f"Comp {next_id}")
                            next_id += 1

                        for blip in grp_copy.iter(f"{{{_A}}}blip"):
                            embed = blip.get(f"{{{_R}}}embed")
                            if embed:
                                blip.attrib.pop(f"{{{_R}}}embed", None)

                        target_sp_tree.append(grp_copy)

                        texts = [t.text.strip()[:15] for t in grp_copy.iter(f"{{{_A}}}t") if t.text and t.text.strip()]
                        tag = "3D" if has_3d else "2D"
                        print(f"  Slide {slide_count+1} [{tag}]: {fname[:25]}, texts={t_count}, sample={texts[:3]}")
                        slide_count += 1
                        if has_3d:
                            found_3d += 1
                        else:
                            found_no3d += 1
                        break  # one grp per slide
        except Exception:
            pass
    if slide_count >= 10:
        break

out = os.path.join(OUT_DIR, "direct_3d_vs_2d.pptx")
prs.save(out)
print(f"\nSaved: {out} ({os.path.getsize(out)/1024:.1f} KB)")
print(f"3D slides: {found_3d}, 2D slides: {found_no3d}")
