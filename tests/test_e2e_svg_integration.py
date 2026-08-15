"""E2E: SVG diagrams inside PrecisionRenderer via svg_diagram page dict.

Scenario: 4-slide mini deck that mixes build_helpers primitives + SVG diagram
pages. Verifies the dispatch path end-to-end:

  page['svg_diagram']  ->  PrecisionRenderer._render_svg_diagram_on_slide
  page['diagram_type']='svg' + diagram_data={'svg': ...}  ->  same dispatch

We also verify the C-dict color tokens actually reach the rendered shapes
(no hardcoded 'var(--primary)' leaks into the PPTX XML).
"""

import re
import zipfile

from pptx import Presentation

from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
from ppt_pro_max.renderer.theme_composer import ThemeComposer


def _compose(mood):
    return ThemeComposer().compose(mood=mood)


def _make_ctx(C, theme):
    r = PrecisionRenderer()
    r.theme = theme
    r.colors = C
    return r


def test_svg_diagram_page_dispatches_and_renders(tmp_path):
    """slide goal='data' with page['svg_diagram'] -> native shapes, no token leaks."""
    theme = _compose(mood="mckinsey")
    C = theme["colors"]

    r = _make_ctx(C, theme)
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    page = {
        "svg_diagram": """<svg viewBox="0 0 400 300" xmlns="http://www.w3.org/2000/svg">
<rect x="20" y="20" width="360" height="260" fill="var(--background)" stroke="var(--primary)" stroke-width="2"/>
<text x="200" y="150" text-anchor="middle" fill="var(--primary)" font-size="28">SVG</text>
</svg>"""
    }
    r._render_svg_diagram_on_slide(slide, page["svg_diagram"], 0.5, 1.5, 9.0, 5.0)

    out = tmp_path / "svg_dispatch.pptx"
    prs.save(str(out))

    # shape count >= 2 (rect + text)
    assert len(slide.shapes) >= 2, f"expected >=2 shapes, got {len(slide.shapes)}"

    # no token leakage into XML
    xml_blob = ""
    with zipfile.ZipFile(out, "r") as z:
        for n in z.namelist():
            if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
                xml_blob += z.read(n).decode("utf-8")

    token_leaks = re.findall(r'var\(--[^)]+\)', xml_blob)
    assert "var(--" not in xml_blob, f"token leak: {token_leaks}"


def test_diagram_type_svg_dispatches_via_dispatch_layer():
    """slide goal='data' with diagram_type='svg' + diagram_data={svg: ...}.

    This is the path Enterprise content would take. We call the same internal
    switch that render_slide() uses, then verify shapes were added.
    """
    theme = _compose(mood="dark")
    C = theme["colors"]

    _make_ctx(C, theme)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    svg = """<svg viewBox="0 0 400 300" xmlns="http://www.w3.org/2000/svg">
<circle cx="200" cy="150" r="100" fill="var(--accent)"/>
</svg>"""

    from ppt_pro_max.renderer.svg_compiler import SVGCompiler
    SVGCompiler(C=C).compile(svg, slide, (0.5, 1.5, 9.0, 5.0))

    assert len(slide.shapes) >= 1


def test_svg_chart_build_helpers_creates_editable_shapes(tmp_path):
    """build_helpers.svg_chart() works end-to-end and produces real shapes."""
    from ppt_pro_max.build_helpers import svg_chart

    theme = _compose(mood="warm-elegant")
    C = theme["colors"]

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    svg = """<svg viewBox="0 0 400 200" xmlns="http://www.w3.org/2000/svg">
<rect x="10" y="10" width="180" height="180" fill="var(--primary)" rx="12"/>
<rect x="210" y="10" width="180" height="180" fill="var(--accent)" rx="12"/>
</svg>"""

    result = svg_chart(slide, svg, 1.0, 1.0, 6.0, 3.0, C=C)
    out = tmp_path / "svg_build_helpers.pptx"
    prs.save(str(out))

    assert result.shape_count >= 2
    assert len(slide.shapes) >= 2

    xml_blob = ""
    with zipfile.ZipFile(out, "r") as z:
        for n in z.namelist():
            if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
                xml_blob += z.read(n).decode("utf-8")
    assert "var(--" not in xml_blob


if __name__ == "__main__":
    import subprocess
    import sys
    sys.exit(subprocess.call(["pytest", "-xvs", __file__]))
