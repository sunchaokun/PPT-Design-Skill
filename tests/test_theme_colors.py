"""Strict tests for set_theme_colors() — writing C palette into PowerPoint theme.

Verifies: mapping correctness, round-trip persistence, edge cases
(empty C, missing keys, unsaved prs, no theme part).
"""

import sys
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from ppt_pro_max.build_helpers import add_slide, set_theme_colors, set_widescreen

_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

FULL_C = {
    'primary': '#1E3A5F', 'secondary': '#37474F', 'tertiary': '#C9A96E',
    'muted': '#5B7BA6', 'light': '#D6E4F0', 'text_dark': '#1A2B3C',
    'background': '#F8FAFC', 'card_bg': '#FFFFFF', 'divider': '#E0E8F0',
}


@pytest.fixture
def prs():
    p = Presentation()
    set_widescreen(p)
    return p


def get_theme(prs):
    for rel in prs.part.rels.values():
        if 'theme' in rel.reltype:
            return rel.target_part
    return None


def get_clr_scheme(prs):
    theme = get_theme(prs)
    assert theme is not None, "no theme part"
    el = etree.fromstring(theme.blob)
    return el.find(f"{{{_NS}}}themeElements/{{{_NS}}}clrScheme")


class TestThemeMapping:
    def test_primary_to_accent1(self, prs):
        set_theme_colors(prs, FULL_C)
        scheme = get_clr_scheme(prs)
        srgb = scheme.find(f"{{{_NS}}}accent1/{{{_NS}}}srgbClr")
        assert srgb.get("val") == "1E3A5F"

    def test_secondary_to_accent2(self, prs):
        set_theme_colors(prs, FULL_C)
        scheme = get_clr_scheme(prs)
        srgb = scheme.find(f"{{{_NS}}}accent2/{{{_NS}}}srgbClr")
        assert srgb.get("val") == "37474F"

    def test_tertiary_to_accent3(self, prs):
        set_theme_colors(prs, FULL_C)
        scheme = get_clr_scheme(prs)
        srgb = scheme.find(f"{{{_NS}}}accent3/{{{_NS}}}srgbClr")
        assert srgb.get("val") == "C9A96E"

    def test_muted_to_accent4(self, prs):
        set_theme_colors(prs, FULL_C)
        scheme = get_clr_scheme(prs)
        srgb = scheme.find(f"{{{_NS}}}accent4/{{{_NS}}}srgbClr")
        assert srgb.get("val") == "5B7BA6"

    def test_light_to_accent5(self, prs):
        set_theme_colors(prs, FULL_C)
        scheme = get_clr_scheme(prs)
        srgb = scheme.find(f"{{{_NS}}}accent5/{{{_NS}}}srgbClr")
        assert srgb.get("val") == "D6E4F0"

    def test_text_dark_to_dk2(self, prs):
        set_theme_colors(prs, FULL_C)
        scheme = get_clr_scheme(prs)
        srgb = scheme.find(f"{{{_NS}}}dk2/{{{_NS}}}srgbClr")
        assert srgb.get("val") == "1A2B3C"

    def test_card_bg_to_lt2(self, prs):
        set_theme_colors(prs, FULL_C)
        scheme = get_clr_scheme(prs)
        srgb = scheme.find(f"{{{_NS}}}lt2/{{{_NS}}}srgbClr")
        assert srgb.get("val") == "FFFFFF"


class TestEdgeCases:
    def test_empty_C_no_crash(self, prs):
        set_theme_colors(prs, {})
        scheme = get_clr_scheme(prs)
        assert scheme is not None  # untouched but intact

    def test_none_C_no_crash(self, prs):
        set_theme_colors(prs, None)
        scheme = get_clr_scheme(prs)
        assert scheme is not None

    def test_partial_C_only_writes_present(self, prs):
        set_theme_colors(prs, {'primary': '#123456'})
        scheme = get_clr_scheme(prs)
        a1 = scheme.find(f"{{{_NS}}}accent1/{{{_NS}}}srgbClr")
        assert a1.get("val") == "123456"
        # accent2 untouched -> still default (C0504D is pptx default accent2)
        a2 = scheme.find(f"{{{_NS}}}accent2/{{{_NS}}}srgbClr")
        assert a2.get("val") == "C0504D"

    def test_hash_optional(self, prs):
        set_theme_colors(prs, {'primary': 'ABCDEF'})  # no leading #
        scheme = get_clr_scheme(prs)
        a1 = scheme.find(f"{{{_NS}}}accent1/{{{_NS}}}srgbClr")
        assert a1.get("val") == "ABCDEF"


class TestPersistence:
    def test_roundtrip_after_save_reload(self, prs, tmp_path):
        set_theme_colors(prs, FULL_C)
        add_slide(prs)
        path = tmp_path / "themed.pptx"
        prs.save(path)

        reloaded = Presentation(str(path))
        scheme = get_clr_scheme(reloaded)
        a1 = scheme.find(f"{{{_NS}}}accent1/{{{_NS}}}srgbClr")
        assert a1.get("val") == "1E3A5F"

    def test_theme_well_formed_xml(self, prs):
        set_theme_colors(prs, FULL_C)
        theme = get_theme(prs)
        etree.fromstring(theme.blob)  # parses = well-formed

    def test_no_orphan_children(self, prs):
        # After set_theme_colors, each mapped clr element has exactly one srgbClr child
        set_theme_colors(prs, FULL_C)
        scheme = get_clr_scheme(prs)
        for tag in ["accent1", "accent2", "accent3", "accent4", "accent5", "dk2", "lt2"]:
            el = scheme.find(f"{{{_NS}}}{tag}")
            children = list(el)
            assert len(children) == 1, f"{tag} has {len(children)} children"
            assert children[0].tag.endswith("}srgbClr")


class TestIntegrationWithBuild:
    def test_full_deck_saves_with_theme(self, tmp_path):
        from ppt_pro_max.build_helpers import (
            TYPOGRAPHY,
            add_slide,
            clean_save,
            hero_slide,
            rect,
        )
        prs = Presentation()
        set_widescreen(prs)
        set_theme_colors(prs, FULL_C)
        t = TYPOGRAPHY['cjk_professional']
        s = add_slide(prs)
        hero_slide(s, '测试主题色', 'set_theme_colors 集成', C=FULL_C, typo=t)
        rect(s, 0, 6.9, 13.333, 0.08, FULL_C['primary'], C=FULL_C)

        path = tmp_path / "deck.pptx"
        clean_save(prs, str(path))
        assert path.exists()

        reloaded = Presentation(str(path))
        scheme = get_clr_scheme(reloaded)
        a1 = scheme.find(f"{{{_NS}}}accent1/{{{_NS}}}srgbClr")
        assert a1.get("val") == "1E3A5F"
