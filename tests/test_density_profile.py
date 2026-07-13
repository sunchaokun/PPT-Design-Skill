"""Tests for DensityProfile."""

from __future__ import annotations

import pytest


class TestDensityProfile:

    def test_get_profile_low(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile
        p = get_density_profile(1)
        assert p.title_size == 36
        assert p.max_bullets == 3

    def test_get_profile_high(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile
        p = get_density_profile(10)
        assert p.title_size == 18
        assert p.max_bullets == 15

    def test_get_profile_clamped_low(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile
        p = get_density_profile(0)
        assert p.title_size == 36

    def test_get_profile_clamped_high(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile
        p = get_density_profile(99)
        assert p.title_size == 18

    def test_apply_density_truncates_bullets(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile, apply_density_to_bullets
        p = get_density_profile(1)
        bullets = ["short", "a" * 100, "medium text"]
        result = apply_density_to_bullets(bullets, p)
        assert len(result) == 3
        assert result[1].endswith("…")

    def test_apply_density_limits_count(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile, apply_density_to_bullets
        p = get_density_profile(1)
        bullets = [f"bullet {i}" for i in range(10)]
        result = apply_density_to_bullets(bullets, p)
        assert len(result) == 3

    def test_apply_density_preserves_short(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile, apply_density_to_bullets
        p = get_density_profile(5)
        result = apply_density_to_bullets(["hello", "world"], p)
        assert result == ["hello", "world"]

    def test_profile_mid_range(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile
        p = get_density_profile(5)
        assert p.title_size == 28
        assert p.bullet_size == 12
        assert p.line_spacing == 1.15

    def test_density_applied_to_design_bullets_then_rendered(self, tmp_path):
        from ppt_pro_max.enterprise.density_profile import get_density_profile, apply_density_to_bullets
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.util import Pt
        profile = get_density_profile(1)
        long_bullets = [f"bullet {i} with lots of text" for i in range(10)]
        truncated = apply_density_to_bullets(long_bullets, profile)
        assert len(truncated) == 3
        design = {"goal": "content", "title": "Test", "bullets": truncated}
        renderer = PrecisionRenderer()
        prs = renderer.create_presentation()
        slide = renderer.render_slide(prs, design)
        found = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text
                    if "bullet 0" in text:
                        found = True
        assert found, "Truncated bullet text should appear in rendered slide"
