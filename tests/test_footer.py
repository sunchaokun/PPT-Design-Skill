"""Tests for Page Numbers, Footer, and Watermark features (v0.4 I-C)."""

from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches


def _make_project(tmp_path: Path, with_footer: bool = False, with_watermark: bool = False) -> Path:
    project = tmp_path / "proj"
    project.mkdir()
    prs = Presentation()
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[-1])
    prs.save(str(project / "template.pptx"))
    brand = {"colors": {"primary": "#1A3C6E"}}
    if with_footer:
        brand["footer"] = {
            "show_page_number": True,
            "page_number_format": "{n}/{total}",
            "page_number_position": "bottom_right",
            "show_footer_text": True,
            "footer_text": "Confidential",
            "footer_position": "bottom_center",
            "font_size_pt": 10,
        }
    if with_watermark:
        brand["watermark"] = {
            "text": "DRAFT",
            "opacity": 0.15,
            "rotation": -45,
            "font_size_pt": 72,
            "skip_pages": [1],
        }
    (project / "brand.json").write_text(json.dumps(brand), encoding="utf-8")
    content = {
        "meta": {"title": "Test"},
        "slides": [
            {"goal": "hook", "title": "Cover"},
            {"goal": "content", "title": "Page 2"},
            {"goal": "content", "title": "Page 3"},
        ],
    }
    (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
    return project


class TestPageNumbers:

    def test_page_numbers_rendered(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        for _ in range(4):
            prs.slides.add_slide(prs.slide_layouts[-1])
        renderer = EnterpriseRenderer()
        renderer.add_page_numbers(prs, {
            "show_page_number": True,
            "page_number_format": "{n}/{total}",
            "page_number_position": "bottom_right",
        })
        slide2 = prs.slides[1]
        texts = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        assert any("2/4" in t for t in texts)

    def test_cover_page_skipped(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[-1])
        renderer = EnterpriseRenderer()
        renderer.add_page_numbers(prs, {
            "show_page_number": True,
            "page_number_format": "{n}",
            "page_number_position": "bottom_right",
        })
        slide0 = prs.slides[0]
        texts = []
        for shape in slide0.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        assert not any("1" == t.strip() for t in texts)

    def test_custom_format(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[-1])
        renderer = EnterpriseRenderer()
        renderer.add_page_numbers(prs, {
            "show_page_number": True,
            "page_number_format": "Page {n} of {total}",
            "page_number_position": "bottom_right",
        })
        slide2 = prs.slides[1]
        texts = []
        for shape in slide2.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        assert any("Page 2 of 3" in t for t in texts)

    def test_footer_text_rendered(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[-1])
        renderer = EnterpriseRenderer()
        renderer.add_page_numbers(prs, {
            "show_page_number": False,
            "show_footer_text": True,
            "footer_text": "Confidential — ACME",
            "footer_position": "bottom_center",
            "font_size_pt": 10,
        })
        slide1 = prs.slides[1]
        texts = []
        for shape in slide1.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        assert any("Confidential" in t for t in texts)

    def test_no_footer_config_no_side_effects(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        for _ in range(2):
            prs.slides.add_slide(prs.slide_layouts[-1])
        renderer = EnterpriseRenderer()
        shape_count_before = len(list(prs.slides[0].shapes))
        renderer.add_page_numbers(prs, {})
        shape_count_after = len(list(prs.slides[0].shapes))
        assert shape_count_after == shape_count_before

    def test_skip_pages(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        for _ in range(4):
            prs.slides.add_slide(prs.slide_layouts[-1])
        renderer = EnterpriseRenderer()
        renderer.add_page_numbers(prs, {
            "show_page_number": True,
            "page_number_format": "{n}",
            "page_number_position": "bottom_right",
            "skip_pages": [2],
        })
        slide1 = prs.slides[1]
        texts = []
        for shape in slide1.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        assert not any("2" == t.strip() for t in texts)


class TestWatermark:

    def test_watermark_rendered(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[-1])
        renderer = EnterpriseRenderer()
        renderer.add_watermark(prs, {
            "text": "CONFIDENTIAL",
            "opacity": 0.15,
            "rotation": -45,
            "font_size_pt": 72,
            "skip_pages": [],
        })
        slide1 = prs.slides[1]
        texts = []
        for shape in slide1.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        assert any("CONFIDENTIAL" in t for t in texts)

    def test_watermark_skip_pages(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[-1])
        renderer = EnterpriseRenderer()
        renderer.add_watermark(prs, {
            "text": "DRAFT",
            "opacity": 0.15,
            "rotation": -45,
            "font_size_pt": 72,
            "skip_pages": [1],
        })
        slide0 = prs.slides[0]
        texts = []
        for shape in slide0.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        assert not any("DRAFT" in t for t in texts)


class TestBrandSpecFooter:

    def test_brand_spec_from_json_with_footer(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        data = {
            "colors": {"primary": "#1A3C6E"},
            "footer": {"show_page_number": True, "page_number_format": "{n}/{total}"},
            "watermark": {"text": "DRAFT"},
        }
        spec = BrandSpec.from_brand_json(data)
        assert spec.footer is not None
        assert spec.footer["show_page_number"] is True
        assert spec.watermark is not None
        assert spec.watermark["text"] == "DRAFT"

    def test_brand_spec_merge_footer(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        template = BrandSpec(source="template", colors={"primary": "#000"})
        merged = BrandSpec.merge(template, {
            "colors": {"primary": "#1A3C6E"},
            "footer": {"show_page_number": True},
        })
        assert merged.footer is not None
        assert merged.footer["show_page_number"] is True

    def test_full_pipeline_with_footer(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = _make_project(tmp_path, with_footer=True)
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Test", project_dir=str(project))
        assert result["num_slides"] == 3
        prs_out = Presentation(result["output_path"])
        slide1 = prs_out.slides[1]
        texts = []
        for shape in slide1.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        has_page_num = any("2/3" in t for t in texts)
        has_footer = any("Confidential" in t for t in texts)
        assert has_page_num or has_footer
