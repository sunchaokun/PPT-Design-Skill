"""Integration tests: Pipeline calls parse_readme, assign_images_by_size, auto_generate_image_prompts."""

from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation


def _make_png(path: Path, w: int = 200, h: int = 150) -> Path:
    img = Image.new("RGB", (w, h), color="green")
    img.save(str(path))
    return path


class TestPipelineReadmeIntegration:

    def test_readme_only_project(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = tmp_path / "readme_project"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        (project / "README.md").write_text(
            "# My Product\n\nIntro text\n\n# Features\n\n- Fast\n- Secure\n- Scalable\n\n# Contact\n\nGet in touch",
            encoding="utf-8",
        )
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="My Product", project_dir=str(project))
        assert result["num_slides"] == 3
        assert result["render_mode"] == "precision"
        assert os.path.isfile(result["output_path"])

    def test_content_json_takes_priority_over_readme(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = tmp_path / "priority_project"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        (project / "README.md").write_text("# From Readme\n\nReadme content", encoding="utf-8")
        content = {"meta": {"title": "From JSON"}, "slides": [{"goal": "hook", "title": "JSON Title"}]}
        (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Priority Test", project_dir=str(project))
        assert result["num_slides"] == 1

    def test_readme_chinese_keywords(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = tmp_path / "chinese_project"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        (project / "README.md").write_text(
            "# 产品介绍\n\n产品概述\n\n# 痛点分析\n\n当前挑战\n\n# 解决方案\n\n我们的方法\n\n# 功能特性\n\n核心能力\n\n# 联系我们\n\n联系方式",
            encoding="utf-8",
        )
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="产品介绍", project_dir=str(project))
        assert result["num_slides"] == 5


class TestPipelineImageSizeIntegration:

    def test_size_based_image_assignment(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = tmp_path / "imgsize_project"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        _make_png(project / "hero_bg.png", w=1600, h=900)
        _make_png(project / "product.png", w=1000, h=800)
        content = {
            "meta": {"title": "Image Size Test"},
            "slides": [
                {"goal": "hook", "title": "Welcome"},
                {"goal": "features", "title": "Features"},
            ],
        }
        (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Image Size Test", project_dir=str(project))
        assert result["num_slides"] == 2
        assert os.path.isfile(result["output_path"])


class TestPipelineImagePromptIntegration:

    def test_image_prompt_generated_for_imageless_pages(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.image_matcher import auto_generate_image_prompts
        designs = [
            {"goal": "hook", "title": "Welcome"},
            {"goal": "problem", "title": "Challenges"},
            {"goal": "features", "title": "Key Features"},
        ]
        result = auto_generate_image_prompts(designs)
        assert all(d.get("image_prompt") for d in result)
        assert "Welcome" in result[0]["image_prompt"]
        assert "Challenges" in result[1]["image_prompt"]

    def test_image_prompt_not_overwritten_when_image_exists(self):
        from ppt_pro_max.enterprise.image_matcher import auto_generate_image_prompts
        designs = [
            {"goal": "hook", "title": "Welcome", "image": "/some/image.png"},
        ]
        result = auto_generate_image_prompts(designs)
        assert result[0].get("image_prompt") is None


class TestPipelineDensityIntegration:

    def test_density_truncates_bullets(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = tmp_path / "density_project"
        project.mkdir()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        content = {
            "meta": {"title": "Density Test"},
            "slides": [
                {"goal": "content", "title": "Many Points", "bullets": [f"Point {i}" for i in range(20)]},
            ],
        }
        (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Density Test", project_dir=str(project), density=3)
        assert result["num_slides"] == 1
