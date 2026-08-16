"""Tests for svg_compiler._paint — gradient + paint resolution.

Validates:
- GradientDef dataclass defaults
- collect_linear_gradient: stops, x1/y1/x2/y2, spreadMethod
- collect_radial_gradient: cx/cy/r, stops
- apply_gradient: linear angle, radial fillToRect
- resolve_paint: solid/grad/none
- unsupported spreadMethod raises SVGCompileError
"""
import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_pro_max.renderer.svg_compiler import SVGCompileError
from ppt_pro_max.renderer.svg_compiler._paint import (
    GradientDef,
    _parse_percent_or_float,
    apply_gradient,
    collect_linear_gradient,
    collect_radial_gradient,
    resolve_paint,
)


def _resolve_color_fn(v, C, fb):
    if v and v.startswith("#"):
        return v
    if v and v.lower() in ("red",):
        return "#FF0000"
    return fb or None


def _make_svg_el(svg_str):
    from lxml import etree
    return etree.fromstring(svg_str)


class TestParsePercentOrFloat:
    def test_percent(self):
        assert _parse_percent_or_float("50%") == 0.5

    def test_float(self):
        assert _parse_percent_or_float("0.75") == 0.75

    def test_integer(self):
        assert _parse_percent_or_float("1") == 1.0

    def test_empty_default(self):
        assert _parse_percent_or_float("", 0.3) == 0.3

    def test_none_default(self):
        assert _parse_percent_or_float(None, 0.7) == 0.7


class TestCollectLinearGradient:
    def test_basic_linear(self):
        svg = (
            '<linearGradient xmlns="http://www.w3.org/2000/svg" id="g1" x1="0" y1="0" x2="1" y2="1">'
            '<stop offset="0" stop-color="#FF0000"/><stop offset="1" stop-color="#0000FF"/>'
            "</linearGradient>"
        )
        el = _make_svg_el(svg)
        grad = collect_linear_gradient(el, {}, _resolve_color_fn)
        assert grad.gradient_type == "linear"
        assert len(grad.stops) == 2
        assert grad.stops[0] == (0.0, "#FF0000", 1.0)
        assert grad.stops[1] == (1.0, "#0000FF", 1.0)
        assert grad.x1 == 0.0
        assert grad.y2 == 1.0

    def test_percent_offsets(self):
        svg = (
            '<linearGradient xmlns="http://www.w3.org/2000/svg" id="g1">'
            '<stop offset="25%" stop-color="#FF0000"/><stop offset="75%" stop-color="#0000FF"/>'
            "</linearGradient>"
        )
        el = _make_svg_el(svg)
        grad = collect_linear_gradient(el, {}, _resolve_color_fn)
        assert grad.stops[0][0] == 0.25
        assert grad.stops[1][0] == 0.75

    def test_stop_opacity(self):
        svg = (
            '<linearGradient xmlns="http://www.w3.org/2000/svg" id="g1">'
            '<stop offset="0" stop-color="#FF0000" stop-opacity="0.5"/>'
            "</linearGradient>"
        )
        el = _make_svg_el(svg)
        grad = collect_linear_gradient(el, {}, _resolve_color_fn)
        assert grad.stops[0][2] == 0.5

    def test_unsupported_spread_raises(self):
        svg = (
            '<linearGradient xmlns="http://www.w3.org/2000/svg" id="g1" spreadMethod="reflect">'
            '<stop offset="0" stop-color="#FF0000"/>'
            "</linearGradient>"
        )
        el = _make_svg_el(svg)
        with pytest.raises(SVGCompileError, match="spreadMethod"):
            collect_linear_gradient(el, {}, _resolve_color_fn)

    def test_gradient_transform(self):
        svg = (
            '<linearGradient xmlns="http://www.w3.org/2000/svg" id="g1" gradientTransform="rotate(45)">'
            '<stop offset="0" stop-color="#FF0000"/>'
            "</linearGradient>"
        )
        el = _make_svg_el(svg)
        grad = collect_linear_gradient(el, {}, _resolve_color_fn)
        assert grad.transform is not None


class TestCollectRadialGradient:
    def test_basic_radial(self):
        svg = (
            '<radialGradient xmlns="http://www.w3.org/2000/svg" id="rg" cx="50%" cy="50%" r="50%">'
            '<stop offset="0" stop-color="#FFFFFF"/><stop offset="1" stop-color="#000000"/>'
            "</radialGradient>"
        )
        el = _make_svg_el(svg)
        grad = collect_radial_gradient(el, {}, _resolve_color_fn)
        assert grad.gradient_type == "radial"
        assert grad.cx == 0.5
        assert grad.cy == 0.5
        assert grad.r == 0.5
        assert len(grad.stops) == 2

    def test_numeric_cx_cy_r(self):
        svg = (
            '<radialGradient xmlns="http://www.w3.org/2000/svg" id="rg" cx="0.3" cy="0.7" r="0.4">'
            '<stop offset="0" stop-color="#FF0000"/>'
            "</radialGradient>"
        )
        el = _make_svg_el(svg)
        grad = collect_radial_gradient(el, {}, _resolve_color_fn)
        assert grad.cx == 0.3
        assert grad.cy == 0.7
        assert grad.r == 0.4


class TestApplyGradient:
    def _make_slide(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs.slides.add_slide(prs.slide_layouts[6])

    def test_linear_gradient_applies(self):
        slide = self._make_slide()
        sh = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        grad = GradientDef(
            stops=[(0.0, "#FF0000", 1.0), (1.0, "#0000FF", 1.0)],
            x1=0, y1=0, x2=1, y2=0,
        )

        def wrap_fn(elem):
            if hasattr(elem, "_element"):
                return elem
            class _P:
                _element = elem
            return _P(elem)

        apply_gradient(sh, grad, wrap_fn)
        xml = sh._element.xml
        assert "gradFill" in xml

    def test_radial_gradient_applies(self):
        slide = self._make_slide()
        sh = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        grad = GradientDef(
            stops=[(0.0, "#FFFFFF", 1.0), (1.0, "#000000", 1.0)],
            gradient_type="radial",
            cx=0.5, cy=0.5, r=0.5,
        )

        def wrap_fn(elem):
            if hasattr(elem, "_element"):
                return elem
            class _P:
                _element = elem
            return _P(elem)

        apply_gradient(sh, grad, wrap_fn)
        xml = sh._element.xml
        assert "gradFill" in xml
        assert "path" in xml

    def test_alpha_in_stops(self):
        slide = self._make_slide()
        sh = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        grad = GradientDef(
            stops=[(0.0, "#FF0000", 0.5), (1.0, "#0000FF", 1.0)],
            x1=0, y1=0, x2=1, y2=0,
        )

        def wrap_fn(elem):
            if hasattr(elem, "_element"):
                return elem
            class _P:
                _element = elem
            return _P(elem)

        apply_gradient(sh, grad, wrap_fn)
        xml = sh._element.xml
        assert "alpha" in xml

    def test_linear_gradient_transform_rotate_45(self):
        """gradientTransform=rotate(45) rotates x1/y1/x2/y2 by 45 degrees."""
        from ppt_pro_max.renderer.svg_compiler._affine import Affine, parse_transform
        import math
        slide = self._make_slide()
        sh = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        # Horizontal gradient x1=0,y1=0,x2=1,y2=0 (angle 0)
        # After rotate(45), new vector is (cos45, sin45) -> angle 45
        tf = parse_transform("rotate(45)")
        grad = GradientDef(
            stops=[(0.0, "#FF0000", 1.0), (1.0, "#0000FF", 1.0)],
            x1=0, y1=0, x2=1, y2=0,
            transform=tf,
        )

        def wrap_fn(elem):
            if hasattr(elem, "_element"):
                return elem
            class _P:
                _element = elem
            return _P(elem)

        apply_gradient(sh, grad, wrap_fn)
        xml = sh._element.xml
        assert "gradFill" in xml
        # The angle after rotate(45) should be 45*60000 = 2700000
        assert "2700000" in xml

    def test_linear_gradient_transform_translate(self):
        """gradientTransform=translate shifts the gradient vector."""
        from ppt_pro_max.renderer.svg_compiler._affine import Affine
        slide = self._make_slide()
        sh = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        # Translate by (0.3, 0.3): the gradient direction (1,0) becomes (1,0) (translation
        # doesn't affect the direction vector). However the start/end points shift.
        # The angle should remain 0.
        grad = GradientDef(
            stops=[(0.0, "#FF0000", 1.0), (1.0, "#0000FF", 1.0)],
            x1=0, y1=0, x2=1, y2=0,
            transform=Affine(1.0, 0.0, 0.0, 1.0, 0.3, 0.3),
        )

        def wrap_fn(elem):
            if hasattr(elem, "_element"):
                return elem
            class _P:
                _element = elem
            return _P(elem)

        apply_gradient(sh, grad, wrap_fn)
        xml = sh._element.xml
        assert "gradFill" in xml
        # Angle should still be 0 (translation doesn't affect direction)
        assert "00000" in xml

    def test_radial_gradient_transform_translate(self):
        """gradientTransform=translate shifts radial gradient cx/cy/r."""
        from ppt_pro_max.renderer.svg_compiler._affine import Affine
        slide = self._make_slide()
        sh = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        # Centered radial: cx=0.5, cy=0.5, r=0.5
        # After translate(0.2, 0.1): cx=0.7, cy=0.6, r=0.5
        grad = GradientDef(
            stops=[(0.0, "#FFFFFF", 1.0), (1.0, "#000000", 1.0)],
            gradient_type="radial",
            cx=0.5, cy=0.5, r=0.5,
            transform=Affine(1.0, 0.0, 0.0, 1.0, 0.2, 0.1),
        )

        def wrap_fn(elem):
            if hasattr(elem, "_element"):
                return elem
            class _P:
                _element = elem
            return _P(elem)

        apply_gradient(sh, grad, wrap_fn)
        xml = sh._element.xml
        assert "gradFill" in xml
        # New fillToRect should reflect translated center:
        # cx=70000, cy=60000, r=50000
        # l = 20000, t = 10000, r = 120000, b = 110000
        assert "l=\"20000\"" in xml
        assert "t=\"10000\"" in xml
        assert "r=\"120000\"" in xml
        assert "b=\"110000\"" in xml


class TestResolvePaint:
    def test_solid_color(self):
        from lxml import etree
        el = etree.fromstring('<rect xmlns="http://www.w3.org/2000/svg" fill="#FF0000"/>')
        features = set()
        kind, val, alpha = resolve_paint(el, "fill", {}, {}, _resolve_color_fn, features)
        assert kind == "solid"
        assert val == "#FF0000"
        assert alpha == 100

    def test_none_fill(self):
        from lxml import etree
        el = etree.fromstring('<rect xmlns="http://www.w3.org/2000/svg" fill="none"/>')
        features = set()
        kind, _, _ = resolve_paint(el, "fill", {}, {}, _resolve_color_fn, features)
        assert kind == "none"

    def test_missing_attr(self):
        from lxml import etree
        el = etree.fromstring('<rect xmlns="http://www.w3.org/2000/svg"/>')
        features = set()
        kind, _, _ = resolve_paint(el, "fill", {}, {}, _resolve_color_fn, features)
        assert kind == "none"

    def test_gradient_ref(self):
        from lxml import etree
        el = etree.fromstring('<rect xmlns="http://www.w3.org/2000/svg" fill="url(#g1)"/>')
        grad = GradientDef(stops=[(0, "#FF0000", 1.0)])
        features = set()
        kind, val, _ = resolve_paint(el, "fill", {"g1": grad}, {}, _resolve_color_fn, features)
        assert kind == "grad"
        assert val is grad
        assert "gradient" in features

    def test_unknown_gradient_ref(self):
        from lxml import etree
        el = etree.fromstring('<rect xmlns="http://www.w3.org/2000/svg" fill="url(#missing)"/>')
        features = set()
        kind, _, _ = resolve_paint(el, "fill", {}, {}, _resolve_color_fn, features)
        assert kind == "none"

    def test_fill_opacity(self):
        from lxml import etree
        el = etree.fromstring('<rect xmlns="http://www.w3.org/2000/svg" fill="#FF0000" fill-opacity="0.5"/>')
        features = set()
        _, _, alpha = resolve_paint(el, "fill", {}, {}, _resolve_color_fn, features)
        assert alpha == 50
