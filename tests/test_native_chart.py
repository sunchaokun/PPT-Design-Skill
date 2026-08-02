"""Tests for native_chart() build_helpers wrapper and expanded ChartBuilder features."""
from __future__ import annotations

from pptx import Presentation


class TestNativeChartBasic:

    def test_bar_chart(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2', 'Q3', 'Q4'],
                              series=[{'name': 'Revenue', 'values': [30, 45, 60, 75]}])
        assert result is not None
        chart = result.chart
        assert len(chart.plots[0].series) == 1

    def test_line_chart(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'line',
                              categories=['Jan', 'Feb', 'Mar'],
                              series=[{'name': 'Sales', 'values': [10, 20, 30]}])
        assert result is not None

    def test_pie_chart(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'pie',
                              categories=['A', 'B', 'C'],
                              series=[{'name': 'Share', 'values': [40, 35, 25]}])
        assert result is not None

    def test_doughnut_chart(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'doughnut',
                              categories=['A', 'B'],
                              series=[{'name': 'Data', 'values': [60, 40]}])
        assert result is not None

    def test_area_chart(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'area',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Growth', 'values': [100, 200]}])
        assert result is not None

    def test_scatter_chart(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'scatter',
                              series=[{'name': 'Points', 'values': [[1, 10], [2, 25], [3, 18]]}])
        assert result is not None

    def test_radar_chart(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'radar',
                              categories=['Speed', 'Reliability', 'Comfort'],
                              series=[{'name': 'Model A', 'values': [80, 90, 70]}])
        assert result is not None

    def test_bar_horizontal(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar_horizontal',
                              categories=['A', 'B'],
                              series=[{'name': 'Data', 'values': [10, 20]}])
        assert result is not None

    def test_bar_stacked(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar_stacked',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'A', 'values': [10, 20]},
                                      {'name': 'B', 'values': [30, 40]}])
        assert result is not None

    def test_bar_3d_fallback(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar_3d',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}])
        assert result is not None

    def test_pie_3d_fallback(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'pie_3d',
                              categories=['A', 'B'],
                              series=[{'name': 'Data', 'values': [60, 40]}])
        assert result is not None

    def test_pie_exploded(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'pie_exploded',
                              categories=['A', 'B', 'C'],
                              series=[{'name': 'Share', 'values': [40, 35, 25]}])
        assert result is not None

    def test_doughnut_exploded(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'doughnut_exploded',
                              categories=['A', 'B'],
                              series=[{'name': 'Data', 'values': [60, 40]}])
        assert result is not None

    def test_line_markers(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'line_markers',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}])
        assert result is not None

    def test_line_stacked(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'line_stacked',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'A', 'values': [10, 20]},
                                      {'name': 'B', 'values': [30, 40]}])
        assert result is not None

    def test_area_stacked(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'area_stacked',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'A', 'values': [10, 20]},
                                      {'name': 'B', 'values': [30, 40]}])
        assert result is not None

    def test_scatter_smooth(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'scatter_smooth',
                              series=[{'name': 'Data', 'values': [[1, 10], [2, 25], [3, 18]]}])
        assert result is not None

    def test_radar_markers(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'radar_markers',
                              categories=['A', 'B', 'C'],
                              series=[{'name': 'Data', 'values': [80, 90, 70]}])
        assert result is not None

    def test_default_params(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar')
        assert result is not None

    def test_multi_series(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2', 'Q3', 'Q4'],
                              series=[{'name': 'Revenue', 'values': [30, 45, 60, 75]},
                                      {'name': 'Cost', 'values': [20, 30, 35, 40]}])
        assert result is not None
        assert len(result.chart.plots[0].series) == 2


class TestNativeChartStyle:

    def test_brand_colors_from_C(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        C = {'primary': '#1A3C6E', 'accent': '#F97316'}
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'color_scheme': 'brand'}, C=C)
        assert result is not None

    def test_custom_color_list(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'color_scheme': ['#FF0000', '#00FF00']})
        assert result is not None

    def test_legend_control(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'show_legend': True, 'legend_position': 'bottom'})
        assert result is not None
        assert result.chart.has_legend is True

    def test_no_legend(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'show_legend': False})
        assert result is not None
        assert result.chart.has_legend is False

    def test_chart_title(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'title': 'Quarterly Revenue'})
        assert result is not None
        assert result.chart.has_title is True

    def test_data_labels(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'show_labels': True})
        assert result is not None
        assert result.chart.plots[0].has_data_labels is True

    def test_pie_sector_colors(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        colors = ['#FF0000', '#00FF00', '#0000FF']
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'pie',
                              categories=['A', 'B', 'C'],
                              series=[{'name': 'Share', 'values': [40, 35, 25]}],
                              style={'color_scheme': colors})
        assert result is not None
        series = result.chart.plots[0].series[0]
        point0 = series.points[0]
        point0.format.fill.solid()
        assert str(point0.format.fill.fore_color.rgb) == 'FF0000'

    def test_value_axis_title(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'value_axis_title': 'Revenue ($M)'})
        assert result is not None
        assert result.chart.value_axis.has_title is True

    def test_gridlines(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'gridlines': 'major_y'})
        assert result is not None
        assert result.chart.value_axis.has_major_gridlines is True

    def test_no_gridlines(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'gridlines': 'none'})
        assert result is not None
        assert result.chart.value_axis.has_major_gridlines is False

    def test_chart_style(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'bar',
                              categories=['Q1', 'Q2'],
                              series=[{'name': 'Data', 'values': [10, 20]}],
                              style={'chart_style': 2})
        assert result is not None

    def test_pie_percentage_labels(self):
        from ppt_pro_max.build_helpers import native_chart, add_slide
        prs = Presentation()
        slide = add_slide(prs)
        result = native_chart(slide, 1.0, 1.5, 7.0, 4.5, 'pie',
                              categories=['A', 'B', 'C'],
                              series=[{'name': 'Share', 'values': [40, 35, 25]}],
                              style={'show_labels': True, 'show_percentage': True})
        assert result is not None
        assert result.chart.plots[0].has_data_labels is True


class TestDonutChartHybrid:

    def test_multi_sector_routes_to_native(self):
        from ppt_pro_max.build_helpers import donut_chart, add_slide
        from pptx.shapes.graphfrm import GraphicFrame
        prs = Presentation()
        slide = add_slide(prs)
        sectors = [('A', '40%', '#2E6504'), ('B', '35%', '#7DA92F'), ('C', '25%', '#81C784')]
        result = donut_chart(slide, 4.0, 3.5, 1.5, 0.8, sectors, C={'primary': '#2E6504'})
        assert isinstance(result, GraphicFrame)

    def test_single_sector_shape_fallback(self):
        from ppt_pro_max.build_helpers import donut_chart, add_slide
        from pptx.shapes.group import GroupShape
        prs = Presentation()
        slide = add_slide(prs)
        sectors = [('Total', '100%', '#2E6504')]
        result = donut_chart(slide, 4.0, 3.5, 1.5, 0.8, sectors, C={'primary': '#2E6504'})
        assert isinstance(result, GroupShape)

    def test_native_false_forces_shape(self):
        from ppt_pro_max.build_helpers import donut_chart, add_slide
        from pptx.shapes.group import GroupShape
        prs = Presentation()
        slide = add_slide(prs)
        sectors = [('A', '40%', '#2E6504'), ('B', '60%', '#7DA92F')]
        result = donut_chart(slide, 4.0, 3.5, 1.5, 0.8, sectors, C={'primary': '#2E6504'}, native=False)
        assert isinstance(result, GroupShape)


class TestChartBuilderExpanded:

    def test_bar_horizontal_type(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "bar_horizontal", "categories": ["A", "B"], "series": [{"name": "Data", "values": [10, 20]}]}
        result = builder.build(slide, config)
        assert result is not None

    def test_line_stacked_type(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "line_stacked", "categories": ["Q1", "Q2"],
                  "series": [{"name": "A", "values": [10, 20]}, {"name": "B", "values": [30, 40]}]}
        result = builder.build(slide, config)
        assert result is not None

    def test_area_stacked_type(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "area_stacked", "categories": ["Q1", "Q2"],
                  "series": [{"name": "A", "values": [10, 20]}, {"name": "B", "values": [30, 40]}]}
        result = builder.build(slide, config)
        assert result is not None

    def test_radar_markers_type(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "radar_markers", "categories": ["A", "B", "C"],
                  "series": [{"name": "Data", "values": [80, 90, 70]}]}
        result = builder.build(slide, config)
        assert result is not None

    def test_pie_exploded_type(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "pie_exploded", "categories": ["A", "B", "C"],
                  "series": [{"name": "Data", "values": [40, 35, 25]}]}
        result = builder.build(slide, config)
        assert result is not None

    def test_doughnut_exploded_type(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "doughnut_exploded", "categories": ["A", "B"],
                  "series": [{"name": "Data", "values": [60, 40]}]}
        result = builder.build(slide, config)
        assert result is not None

    def test_scatter_smooth_type(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "scatter_smooth",
                  "series": [{"name": "Data", "values": [[1, 10], [2, 25], [3, 18]]}]}
        result = builder.build(slide, config)
        assert result is not None

    def test_scatter_lines_type(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "scatter_lines",
                  "series": [{"name": "Data", "values": [[1, 10], [2, 25], [3, 18]]}]}
        result = builder.build(slide, config)
        assert result is not None

    def test_pie_sector_point_colors(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "pie", "categories": ["A", "B", "C"],
                  "series": [{"name": "Share", "values": [40, 35, 25]}],
                  "style": {"color_scheme": ["#FF0000", "#00FF00", "#0000FF"]}}
        result = builder.build(slide, config)
        assert result is not None
        series = result.chart.plots[0].series[0]
        point0 = series.points[0]
        point0.format.fill.solid()
        assert str(point0.format.fill.fore_color.rgb) == 'FF0000'

    def test_data_labels_with_percentage(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "pie", "categories": ["A", "B"],
                  "series": [{"name": "Data", "values": [60, 40]}],
                  "style": {"show_labels": True, "show_percentage": True}}
        result = builder.build(slide, config)
        assert result is not None
        assert result.chart.plots[0].has_data_labels is True

    def test_value_axis_title(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "bar", "categories": ["Q1", "Q2"],
                  "series": [{"name": "Data", "values": [10, 20]}],
                  "style": {"value_axis_title": "Revenue ($M)"}}
        result = builder.build(slide, config)
        assert result is not None
        assert result.chart.value_axis.has_title is True

    def test_category_axis_title(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "bar", "categories": ["Q1", "Q2"],
                  "series": [{"name": "Data", "values": [10, 20]}],
                  "style": {"category_axis_title": "Quarter"}}
        result = builder.build(slide, config)
        assert result is not None
        assert result.chart.category_axis.has_title is True

    def test_gridlines_none(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "bar", "categories": ["Q1", "Q2"],
                  "series": [{"name": "Data", "values": [10, 20]}],
                  "style": {"gridlines": "none"}}
        result = builder.build(slide, config)
        assert result is not None
        assert result.chart.value_axis.has_major_gridlines is False

    def test_chart_style_number(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "bar", "categories": ["Q1", "Q2"],
                  "series": [{"name": "Data", "values": [10, 20]}],
                  "style": {"chart_style": 2}}
        result = builder.build(slide, config)
        assert result is not None

    def test_number_format(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "bar", "categories": ["Q1", "Q2"],
                  "series": [{"name": "Data", "values": [1000, 2000]}],
                  "style": {"show_labels": True, "number_format": "#,##0"}}
        result = builder.build(slide, config)
        assert result is not None

    def test_tick_number_format(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "bar", "categories": ["Q1", "Q2"],
                  "series": [{"name": "Data", "values": [1000, 2000]}],
                  "style": {"tick_number_format": "#,##0"}}
        result = builder.build(slide, config)
        assert result is not None

    def test_unknown_type_fallback(self):
        from ppt_pro_max.renderer.chart_builder import ChartBuilder
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        builder = ChartBuilder()
        config = {"type": "nonexistent", "categories": ["Q1", "Q2"],
                  "series": [{"name": "Data", "values": [10, 20]}]}
        result = builder.build(slide, config)
        assert result is not None
