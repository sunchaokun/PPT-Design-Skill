"""Tests for SlideExtractor — extract content + layout from existing PPT.

TDD: These tests define the expected behavior of SlideExtractor before implementation.
"""
import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


FIXTURES = Path(__file__).parent / "fixtures"
FIXTURES.mkdir(parents=True, exist_ok=True)


def _create_simple_pptx(path: str, num_slides: int = 3):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    for i in range(num_slides):
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"Slide {i+1} Title"
        run.font.size = Pt(36)
        run.font.bold = True

        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
        tf = body.text_frame
        tf.word_wrap = True
        for j, text in enumerate([f"Bullet {j+1}" for j in range(3)]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = text
            run.font.size = Pt(14)

    prs.save(path)
    return path


def _create_pptx_with_image(path: str):
    from PIL import Image as PILImage
    import io

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Slide With Image"
    run.font.size = Pt(36)

    img = PILImage.new("RGB", (200, 150), color="blue")
    img_buf = io.BytesIO()
    img.save(img_buf, format="PNG")
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, Inches(8), Inches(2), Inches(4), Inches(3))

    prs.save(path)
    return path


def _create_pptx_with_table(path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Slide With Table"
    run.font.size = Pt(36)

    rows, cols = 3, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(3))
    table = table_shape.table
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Row 1A"
    table.cell(1, 1).text = "Row 1B"
    table.cell(2, 0).text = "Row 2A"
    table.cell(2, 1).text = "Row 2B"

    prs.save(path)
    return path


def _create_pptx_with_notes(path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Slide With Notes"
    run.font.size = Pt(36)

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker notes for this slide"

    prs.save(path)
    return path


def _create_pptx_with_subtitle(path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Main Title"
    run.font.size = Pt(44)

    sub = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Subtitle text here"
    run2.font.size = Pt(20)

    prs.save(path)
    return path


class TestSlideExtractorBasic:
    def test_extract_simple_pptx(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        pptx_path = str(tmp_path / "simple.pptx")
        _create_simple_pptx(pptx_path, num_slides=3)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)

        assert len(pages) == 3
        for i, page in enumerate(pages):
            assert page["title"] == f"Slide {i+1} Title"
            assert len(page["bullets"]) == 3
            assert page["bullets"][0] == "Bullet 1"

    def test_extract_returns_page_dict_structure(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        pptx_path = str(tmp_path / "simple.pptx")
        _create_simple_pptx(pptx_path, num_slides=1)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        page = pages[0]

        required_keys = {"goal", "title", "subtitle", "bullets", "image", "cards",
                         "diagram_type", "diagram_data", "code", "exercise", "chart",
                         "notes", "links", "layout_hint", "complex_elements"}
        assert required_keys.issubset(set(page.keys()))

    def test_extract_layout_hint(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        pptx_path = str(tmp_path / "simple.pptx")
        _create_simple_pptx(pptx_path, num_slides=1)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        hint = pages[0]["layout_hint"]

        assert "title_pos" in hint
        assert "body_pos" in hint
        assert "has_image" in hint
        assert "bullet_count" in hint
        assert "is_full_bleed" in hint
        assert isinstance(hint["has_image"], bool)
        assert isinstance(hint["bullet_count"], int)

    def test_extract_with_image(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        pptx_path = str(tmp_path / "with_image.pptx")
        _create_pptx_with_image(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)

        assert len(pages) == 1
        page = pages[0]
        assert page["image"] is not None
        assert os.path.isfile(page["image"])
        assert page["layout_hint"]["has_image"] is True

    def test_extract_with_table(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        pptx_path = str(tmp_path / "with_table.pptx")
        _create_pptx_with_table(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)

        assert len(pages) == 1
        page = pages[0]
        assert page["diagram_type"] == "table"
        assert page["diagram_data"] is not None

    def test_extract_with_notes(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        pptx_path = str(tmp_path / "with_notes.pptx")
        _create_pptx_with_notes(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)

        assert len(pages) == 1
        assert pages[0]["notes"] == "Speaker notes for this slide"

    def test_extract_with_subtitle(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        pptx_path = str(tmp_path / "with_subtitle.pptx")
        _create_pptx_with_subtitle(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)

        assert len(pages) == 1
        assert pages[0]["title"] == "Main Title"
        assert pages[0]["subtitle"] == "Subtitle text here"


class TestSlideExtractorGoalInference:
    def test_first_page_inferred_as_hook(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Product Launch"
        run.font.size = Pt(48)

        pptx_path = str(tmp_path / "cover.pptx")
        prs.save(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        assert pages[0]["goal"] == "hook"

    def test_last_page_with_thanks_inferred_as_cta(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        tb1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p1 = tb1.text_frame.paragraphs[0]
        run1 = p1.add_run()
        run1.text = "Content Slide"
        run1.font.size = Pt(36)

        slide2 = prs.slides.add_slide(layout)
        tb2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
        p2 = tb2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = "Thank You"
        run2.font.size = Pt(48)

        pptx_path = str(tmp_path / "cta.pptx")
        prs.save(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        assert pages[-1]["goal"] == "cta"

    def test_problem_keyword_inferred_goal(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Current Challenges"
        run.font.size = Pt(36)

        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(3))
        tf = body.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        run2 = p2.add_run()
        run2.text = "Problem: High costs"
        run2.font.size = Pt(14)

        pptx_path = str(tmp_path / "problem.pptx")
        prs.save(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        assert pages[0]["goal"] == "problem"

    def test_features_keyword_inferred_goal(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Key Features"
        run.font.size = Pt(36)

        pptx_path = str(tmp_path / "features.pptx")
        prs.save(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        assert pages[0]["goal"] == "features"

    def test_default_goal_is_content(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        tb1 = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
        p1 = tb1.text_frame.paragraphs[0]
        run1 = p1.add_run()
        run1.text = "Cover"
        run1.font.size = Pt(48)

        slide2 = prs.slides.add_slide(layout)
        tb2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p2 = tb2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = "Random Content"
        run2.font.size = Pt(36)

        body = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(3))
        tf = body.text_frame
        tf.word_wrap = True
        p3 = tf.paragraphs[0]
        run3 = p3.add_run()
        run3.text = "Some detail here"
        run3.font.size = Pt(14)

        pptx_path = str(tmp_path / "random.pptx")
        prs.save(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        assert pages[1]["goal"] == "content"


class TestSlideExtractorComplexElements:
    def test_group_shape_detected_as_complex_element(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor
        from lxml import etree
        import zipfile

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Group Test"
        run.font.size = Pt(36)

        base_path = str(tmp_path / "group_base.pptx")
        prs.save(base_path)

        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

        grp_sp_xml = f'''<p:grpSp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvGrpSpPr>
    <p:cNvPr id="100" name="TestGroup"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm>
      <a:off x="914400" y="914400"/>
      <a:ext cx="9144000" cy="4572000"/>
      <a:chOff x="0" y="0"/>
      <a:chExt cx="9144000" cy="4572000"/>
    </a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="101" name="Box1"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="0" y="0"/>
        <a:ext cx="2743200" cy="1371600"/>
      </a:xfrm>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:p><a:r><a:t>Step 1</a:t></a:r></a:p>
    </p:txBody>
  </p:sp>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr id="102" name="Box2"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="3200400" y="0"/>
        <a:ext cx="2743200" cy="1371600"/>
      </a:xfrm>
    </p:spPr>
    <p:txBody>
      <a:bodyPr/>
      <a:p><a:r><a:t>Step 2</a:t></a:r></a:p>
    </p:txBody>
  </p:sp>
</p:grpSp>'''

        with zipfile.ZipFile(base_path, 'r') as zin:
            files = {}
            for name in zin.namelist():
                files[name] = zin.read(name)

        slide1_xml = files["ppt/slides/slide1.xml"]
        root = etree.fromstring(slide1_xml)
        ns = {"p": NS_P}
        sp_tree = root.find(".//p:spTree", ns)
        if sp_tree is not None:
            grp_elem = etree.fromstring(grp_sp_xml)
            sp_tree.append(grp_elem)

        files["ppt/slides/slide1.xml"] = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

        pptx_path = str(tmp_path / "with_group.pptx")
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name, content in files.items():
                zout.writestr(name, content)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)

        assert len(pages) == 1
        ce = pages[0]["complex_elements"]
        group_elements = [e for e in ce if e["type"] == "group"]
        assert len(group_elements) >= 1
        assert "Step 1" in group_elements[0]["texts"]
        assert "Step 2" in group_elements[0]["texts"]

    def test_smartart_detected_as_complex_element(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        art_path = str(FIXTURES / "art_restyled_v3.pptx")
        if not os.path.isfile(art_path):
            pytest.skip("art.pptx fixture not available")

        ext = SlideExtractor()
        pages = ext.extract(art_path)

        has_smartart = any(
            e["type"] == "smartart"
            for page in pages
            for e in page.get("complex_elements", [])
        )
        assert has_smartart


class TestSlideExtractorEdgeCases:
    def test_empty_pptx(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        pptx_path = str(tmp_path / "empty.pptx")
        prs.save(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        assert len(pages) == 0

    def test_nonexistent_file(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        ext = SlideExtractor()
        pages = ext.extract(str(tmp_path / "nonexistent.pptx"))
        assert len(pages) == 0

    def test_slide_with_no_text(self, tmp_path):
        from ppt_pro_max.enterprise.slide_extractor import SlideExtractor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(5), Inches(3))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0, 0, 255)

        pptx_path = str(tmp_path / "no_text.pptx")
        prs.save(pptx_path)

        ext = SlideExtractor()
        pages = ext.extract(pptx_path)
        assert len(pages) == 1
        assert pages[0]["title"] == ""
