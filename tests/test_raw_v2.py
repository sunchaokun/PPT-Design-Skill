"""3 PPTs with proper vertical centering, full-height content, no bottom waste."""
from __future__ import annotations
import os, sys, tempfile
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def _rgb(h): return RGBColor.from_string(h.lstrip('#'))

def _txt(s, l, t, w, h, txt, sz=12, c='#000', b=False, fn='Calibri', a='left', an='t'):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    if an == 'ctr': tf._txBody.bodyPr.set('anchor', 'ctr')
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = _rgb(c); p.font.bold = b; p.font.name = fn
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[a]
    return bx

def _rect(s, l, t, w, h, fill, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(fill)
    if line: sh.line.color.rgb = _rgb(line); sh.line.width = Pt(1)
    else: sh.line.fill.background()
    return sh

def _rrect(s, l, t, w, h, fill, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(fill)
    if line: sh.line.color.rgb = _rgb(line); sh.line.width = Pt(1)
    else: sh.line.fill.background()
    return sh

def _oval(s, l, t, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(fill); sh.line.fill.background()
    return sh

def _ml(s, l, t, w, h, lines, sz=12, c='#000', b=False, fn='Calibri', a='left', sp=4):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = _rgb(c)
        p.font.bold = b; p.font.name = fn
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[a]
        p.space_before = Pt(sp); p.space_after = Pt(sp)
    return bx

# Slide dimensions
SW, SH = 13.333, 7.5

# =============================================================================
# McKinsey: Serif, data-dense, left grid, muted green
# =============================================================================
def build_mckinsey(out):
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    P='#2E6504'; A='#7DA92F'; M='#81C784'; L='#C8E6C9'; W='#FFFFFF'
    D='#1A1A1A'; BD='#333333'; MU='#666666'; DV='#CCCCCC'; CD='#F5F5F5'
    FH='Georgia'; FB='Calibri'; FC='Consolas'

    # Cover — vertically centered
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,SH,P)
    _rect(s,0,7.2,SW,0.3,A)
    _txt(s,1.5,2.4,10,1.2,'Strategic Growth Report',sz=48,c=W,b=True,fn=FH)
    _txt(s,1.5,3.7,10,0.5,'FY2025 Performance Review',sz=22,c=L,fn=FB)
    _rect(s,1.5,4.4,2.5,0.04,A)
    _txt(s,1.5,4.7,10,0.4,'NeuralForge Inc.  |  Confidential',sz=13,c=MU,fn=FB)

    # Section divider — vertically centered
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,SH,P)
    _txt(s,1.5,1.8,3,2.5,'01',sz=96,c=L,b=True,fn=FH)
    _rect(s,1.5,4.5,2.0,0.04,A)
    _txt(s,1.5,4.9,10,1.2,'Core Metrics',sz=44,c=W,b=True,fn=FH)

    # KPI — content fills full height (header 0.4-1.3, cards 1.6-6.8)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.08,P)
    _txt(s,0.65,0.4,12,0.5,'Key Performance Indicators',sz=28,c=D,b=True,fn=FH)
    _txt(s,0.65,0.9,12,0.3,'Core business metrics',sz=12,c=MU,fn=FB)
    _rect(s,0.65,1.25,12,0.004,DV)
    kpis=[('12.8B','Annual Revenue','+8.3%',True),('99.9%','System Uptime','+0.1%',True),
          ('5.2M','Active Users','+15%',True),('2.4ms','P99 Latency','-12%',False),
          ('72','NPS Score','Top 5%',True),('38%','Women in Leadership','vs 21% avg',True)]
    for idx,(num,lbl,tr,up) in enumerate(kpis):
        col=idx%3; row=idx//3; x=0.65+col*4.1; y=1.6+row*2.7
        _rrect(s,x,y,3.8,2.3,CD,line=DV)
        _rect(s,x,y,3.8,0.06,A)
        _txt(s,x+0.3,y+0.3,3.2,0.7,num,sz=40,c=P,b=True,fn=FH)
        _txt(s,x+0.3,y+1.1,3.2,0.3,lbl,sz=14,c=MU,fn=FB)
        tc=P if up else '#C53030'
        _txt(s,x+0.3,y+1.6,3.2,0.3,tr,sz=12,c=tc,b=True,fn=FB)

    # Revenue — bars + donut fill full height
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.08,P)
    _txt(s,0.65,0.4,12,0.5,'Revenue Breakdown',sz=28,c=D,b=True,fn=FH)
    _txt(s,0.65,0.9,12,0.3,'By business segment',sz=12,c=MU,fn=FB)
    _rect(s,0.65,1.25,12,0.004,DV)
    bars=[('Enterprise SaaS',0.82,'$10.2B'),('SMB Platform',0.48,'$5.8B'),('Consumer App',0.30,'$3.9B')]
    bc=[P,A,M]
    for i,(lbl,pct,val) in enumerate(bars):
        y=2.0+i*1.2
        _rrect(s,2.5,y,5.5,0.7,'#F0F0F0')
        _rrect(s,2.5,y,5.5*pct,0.7,bc[i])
        _txt(s,1.0,y+0.15,1.4,0.4,lbl,sz=14,c=BD,fn=FB,a='right')
        _txt(s,8.2,y+0.15,1.0,0.4,val,sz=14,c=D,b=True,fn=FB)
    _oval(s,9.0,2.0,4.5,4.5,P)
    _oval(s,10.0,3.0,2.5,2.5,W)
    _txt(s,10.0,3.8,2.5,0.6,'100%',sz=22,c=P,b=True,fn=FH,a='center')
    for i,(lt,lc) in enumerate([('SaaS 52%',P),('SMB 29%',A),('Consumer 19%',M)]):
        ly=2.5+i*0.5
        _rrect(s,9.3,ly,0.25,0.25,lc)
        _txt(s,9.7,ly-0.02,1.8,0.3,lt,sz=12,c=BD,fn=FB)

    # Feature cards — tall cards fill height
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.08,P)
    _txt(s,0.65,0.4,12,0.5,'Core Capabilities',sz=28,c=D,b=True,fn=FH)
    _txt(s,0.65,0.9,12,0.3,'What sets us apart',sz=12,c=MU,fn=FB)
    _rect(s,0.65,1.25,12,0.004,DV)
    cards=[('AI Inference Engine','Auto-select optimal framework\n70% compression, P99<10ms\nServing 12.8B requests/sec',P),
           ('Real-time Monitoring','Millisecond alerting\nFull-stack tracing, anomaly detection\nZero blind spots',A),
           ('Zero-config Integration','200+ connectors\nOut-of-box ready, 5-min setup\nNo vendor lock-in',M)]
    for i,(title,desc,accent) in enumerate(cards):
        x=0.65+i*4.1
        _rrect(s,x,1.6,3.8,5.2,CD,line=DV)
        _rect(s,x,1.6,3.8,0.06,accent)
        _txt(s,x+0.3,1.9,3.2,0.5,title,sz=20,c=D,b=True,fn=FH)
        _txt(s,x+0.3,2.5,3.2,3.5,desc,sz=14,c=MU,fn=FB)

    # Code block — fills height
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.08,P)
    _txt(s,0.65,0.4,12,0.5,'Deployment Example',sz=28,c=D,b=True,fn=FH)
    _txt(s,0.65,0.9,12,0.3,'Production-ready in 3 lines',sz=12,c=MU,fn=FB)
    _rect(s,0.65,1.25,12,0.004,DV)
    _rect(s,0.65,1.6,12,5.2,'#1E1E1E')
    _rrect(s,0.65,1.35,1.2,0.25,A)
    _txt(s,0.75,1.37,1.0,0.2,'Python',sz=10,c=W,b=True,fn=FC)
    _ml(s,1.0,2.2,11,4.0,['from neural_ops import deploy','model = load("gpt-4o-finetuned")',
        'deploy(model, replicas=3, gpu="A100")','print("Serving at", model.endpoint)'],sz=16,c='#D4D4D4',fn=FC,sp=10)

    # CTA — vertically centered
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,SH,P)
    _rect(s,0,7.2,SW,0.3,A)
    _txt(s,1.5,2.8,10,1.2,"Let's Build the Future Together",sz=40,c=W,b=True,fn=FH)
    _txt(s,1.5,4.2,10,0.5,'Contact: growth@example.com  |  example.com/demo',sz=16,c=L,fn=FB)

    prs.save(out)

# =============================================================================
# Cyberpunk: Dark, neon, monospace, terminal
# =============================================================================
def build_cyberpunk(out):
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    P='#0A0E27'; A='#00F0FF'; M='#7B2FFF'; PK='#FF2D6B'; B2='#12162B'
    T='#E0E0E0'; DM='#6A6A8A'; FH='Orbitron'; FB='JetBrains Mono'; FC='Consolas'

    # Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,SH,P)
    _rect(s,0,3.5,SW,0.04,A); _rect(s,0,3.6,SW,0.02,M)
    _txt(s,1.0,1.5,11,1.5,'NEURAL OPS',sz=56,c=A,b=True,fn=FH)
    _txt(s,1.0,3.0,11,0.4,'// v4.2  |  AI Infrastructure at Scale',sz=16,c=DM,fn=FB)
    _txt(s,1.0,4.2,11,0.5,'2025 STATUS REPORT',sz=20,c=T,fn=FH)
    _rect(s,1.0,4.9,1.5,0.04,A)
    _txt(s,1.0,5.3,11,0.4,'CLASSIFIED  |  INTERNAL USE ONLY',sz=10,c=PK,fn=FB)

    # KPI — full height
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.06,A)
    _txt(s,0.8,0.35,12,0.5,'SYSTEM METRICS',sz=24,c=A,b=True,fn=FH)
    _txt(s,0.8,0.85,12,0.3,'Real-time dashboard snapshot',sz=11,c=DM,fn=FB)
    _rect(s,0.8,1.2,11.7,0.004,'#2A2F4A')
    kpis=[('12.8B','INFERENCE/SEC','+8.3%',True),('99.9%','UPTIME SLA','+0.1%',True),
          ('5.2M','ACTIVE NODES','+15%',True),('2.4ms','P99 LATENCY','-12%',False),
          ('72','THREAT SCORE','LOW',True),('38%','GPU UTIL','OPTIMAL',True)]
    for idx,(num,lbl,tr,up) in enumerate(kpis):
        col=idx%3; row=idx//3; x=0.8+col*4.2; y=1.6+row*2.8
        _rrect(s,x,y,3.8,2.4,B2,line='#2A2F4A')
        _rect(s,x,y,3.8,0.04,A if up else PK)
        _txt(s,x+0.3,y+0.3,3.2,0.7,num,sz=36,c=A,b=True,fn=FH)
        _txt(s,x+0.3,y+1.1,3.2,0.3,lbl,sz=11,c=DM,fn=FB)
        tc=A if up else PK
        _txt(s,x+0.3,y+1.7,3.2,0.3,tr,sz=11,c=tc,b=True,fn=FB)

    # Traffic — full height
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.06,A)
    _txt(s,0.8,0.35,12,0.5,'TRAFFIC DISTRIBUTION',sz=24,c=A,b=True,fn=FH)
    _txt(s,0.8,0.85,12,0.3,'By service endpoint',sz=11,c=DM,fn=FB)
    _rect(s,0.8,1.2,11.7,0.004,'#2A2F4A')
    bars=[('/api/infer',0.82,'10.2B'),('/api/embed',0.48,'5.8B'),('/api/train',0.30,'3.9B')]
    bc=[A,M,PK]
    for i,(lbl,pct,val) in enumerate(bars):
        y=2.0+i*1.4
        _rrect(s,2.8,y,5.5,0.8,'#1A1F3A')
        _rrect(s,2.8,y,5.5*pct,0.8,bc[i])
        _txt(s,1.0,y+0.2,1.7,0.4,lbl,sz=12,c=DM,fn=FB,a='right')
        _txt(s,8.5,y+0.2,1.0,0.4,val,sz=12,c=T,b=True,fn=FB)
    _oval(s,9.0,2.0,4.5,4.5,A)
    _oval(s,10.0,3.0,2.5,2.5,P)
    _txt(s,10.0,3.8,2.5,0.6,'100%',sz=20,c=A,b=True,fn=FH,a='center')
    for i,(lt,lc) in enumerate([('Inference 52%',A),('Embedding 29%',M),('Training 19%',PK)]):
        ly=2.5+i*0.5
        _rrect(s,9.3,ly,0.25,0.25,lc)
        _txt(s,9.7,ly-0.02,1.8,0.3,lt,sz=11,c=T,fn=FB)

    # Modules — tall cards
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.06,A)
    _txt(s,0.8,0.35,12,0.5,'CORE MODULES',sz=24,c=A,b=True,fn=FH)
    _txt(s,0.8,0.85,12,0.3,'Active service mesh',sz=11,c=DM,fn=FB)
    _rect(s,0.8,1.2,11.7,0.004,'#2A2F4A')
    cards=[('INFERENCE\nENGINE','Auto-routing\n70% compression\nP99<10ms\n12.8B req/sec',A),
           ('MONITORING\nMESH','Millisecond alerting\nFull-stack tracing\nAnomaly detection\nZero blind spots',M),
           ('CONNECTOR\nHUB','200+ integrations\nZero-config mesh\n5-min deploy\nNo vendor lock-in',PK)]
    for i,(title,desc,accent) in enumerate(cards):
        x=0.8+i*4.2
        _rrect(s,x,1.6,3.8,5.2,B2,line='#2A2F4A')
        _rect(s,x,1.6,3.8,0.04,accent)
        _txt(s,x+0.3,1.9,3.2,0.8,title,sz=16,c=accent,b=True,fn=FH)
        _txt(s,x+0.3,2.9,3.2,3.5,desc,sz=12,c=DM,fn=FB)

    # Terminal — full height
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.06,A)
    _txt(s,0.8,0.35,12,0.5,'DEPLOY SEQUENCE',sz=24,c=A,b=True,fn=FH)
    _txt(s,0.8,0.85,12,0.3,'Production pipeline',sz=11,c=DM,fn=FB)
    _rect(s,0.8,1.2,11.7,0.004,'#2A2F4A')
    _rect(s,0.8,1.6,11.7,5.2,'#050810')
    _rect(s,0.8,1.6,11.7,0.04,A)
    _rrect(s,0.8,1.35,1.0,0.25,A)
    _txt(s,0.9,1.37,0.8,0.2,'shell',sz=9,c='#000000',b=True,fn=FC)
    _ml(s,1.2,2.2,10.5,4.0,['$ neural-ops deploy --model gpt-4o-ft --gpu A100',
        '> Replicas: 3  |  Region: us-east-1  |  Auto-scale: ON',
        '> Health check: PASS  |  Endpoint: api.neuralops.ai/v4',
        '> Status: LIVE  |  P99: 2.1ms  |  Throughput: 12.8B/s'],sz=14,c=A,fn=FB,sp=12)

    # CTA
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,SH,P)
    _rect(s,0,3.5,SW,0.04,A); _rect(s,0,3.6,SW,0.02,M)
    _txt(s,1.0,2.2,11,1.2,'INITIALIZE',sz=52,c=A,b=True,fn=FH)
    _txt(s,1.0,3.4,11,0.6,'NEXT PHASE',sz=28,c=M,fn=FH)
    _txt(s,1.0,4.5,11,0.5,'neural-ops.ai  |  #build-the-future',sz=14,c=DM,fn=FB)

    prs.save(out)

# =============================================================================
# Creative: Warm, rounded, vibrant, generous whitespace
# =============================================================================
def build_creative(out):
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    P='#6C3CE1'; A='#FF6B6B'; M='#4ECDC4'; Y='#FFE66D'; W='#FFFFFF'
    D='#2D3436'; BD='#636E72'; MU='#B2BEC3'; CD='#FFFFFF'; BG='#FAFAFA'
    FH='Fredoka'; FB='Nunito'

    # Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,SH,P)
    _oval(s,8.5,-1.5,7,7,A); _oval(s,-2.5,4.5,6,6,M); _oval(s,10,5,4,4,Y)
    _txt(s,1.5,2.2,8,1.5,'Creative Studio',sz=52,c=W,b=True,fn=FH)
    _txt(s,1.5,3.7,8,0.6,'2025',sz=36,c=Y,b=True,fn=FH)
    _txt(s,1.5,4.5,8,0.5,'Design. Build. Ship. Repeat.',sz=18,c=W,fn=FB)

    # KPI — centered with generous spacing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.1,P)
    _txt(s,1.0,0.5,11,0.6,'Impact Numbers',sz=30,c=D,b=True,fn=FH)
    _txt(s,1.0,1.1,11,0.3,'What we achieved together',sz=14,c=MU,fn=FB)
    kpis=[('847','Projects Delivered','+23%',A),('99.2%','Client Satisfaction','+1.8%',M),('156','Team Members','+34',Y)]
    for i,(num,lbl,tr,accent) in enumerate(kpis):
        x=1.0+i*3.9
        _rrect(s,x,2.0,3.5,4.5,CD,line='#E0E0E0')
        _oval(s,x+1.3,2.3,0.9,0.9,accent)
        _txt(s,x+0.3,3.5,2.9,0.7,num,sz=40,c=D,b=True,fn=FH,a='center')
        _txt(s,x+0.3,4.3,2.9,0.3,lbl,sz=15,c=BD,fn=FB,a='center')
        _txt(s,x+0.3,4.8,2.9,0.3,tr,sz=13,c=accent,b=True,fn=FB,a='center')

    # Service mix — full height
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.1,P)
    _txt(s,1.0,0.5,11,0.6,'Service Mix',sz=30,c=D,b=True,fn=FH)
    _txt(s,1.0,1.1,11,0.3,'Where we create value',sz=14,c=MU,fn=FB)
    bars=[('Brand Design',0.72,'61%',P),('Product UX',0.55,'47%',A),('Motion Design',0.38,'32%',M)]
    for i,(lbl,pct,val,clr) in enumerate(bars):
        y=2.2+i*1.4
        _rrect(s,3.0,y,5.0,0.8,'#F0F0F5')
        _rrect(s,3.0,y,5.0*pct,0.8,clr)
        _txt(s,1.2,y+0.2,1.7,0.4,lbl,sz=15,c=D,b=True,fn=FB,a='right')
        _txt(s,8.2,y+0.2,1.0,0.4,val,sz=15,c=D,b=True,fn=FB)
    _oval(s,9.0,2.0,4.5,4.5,P)
    _oval(s,10.0,3.0,2.5,2.5,BG)
    _txt(s,10.0,3.8,2.5,0.6,'100%',sz=22,c=P,b=True,fn=FH,a='center')
    for i,(lt,lc) in enumerate([('Brand 40%',P),('UX 35%',A),('Motion 25%',M)]):
        ly=2.5+i*0.5
        _rrect(s,9.3,ly,0.25,0.25,lc)
        _txt(s,9.7,ly-0.02,1.8,0.3,lt,sz=13,c=BD,fn=FB)

    # Superpowers — tall cards
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.1,P)
    _txt(s,1.0,0.5,11,0.6,'Superpowers',sz=30,c=D,b=True,fn=FH)
    _txt(s,1.0,1.1,11,0.3,'Our creative toolkit',sz=14,c=MU,fn=FB)
    cards=[('Design System','Component library with\n500+ tokens\nAuto-sync Figma\nDesign-to-code pipeline',P),
           ('Rapid Prototyping','Idea to clickable\nprototype in 48 hours\nUser testing built-in\nIterate at speed',A),
           ('User Research','Data-driven design\ndecisions\nA/B testing built-in\nMeasure everything',M)]
    for i,(title,desc,accent) in enumerate(cards):
        x=1.0+i*3.9
        _rrect(s,x,1.8,3.5,5.0,CD,line='#E0E0E0')
        _oval(s,x+1.3,2.1,0.9,0.9,accent)
        _txt(s,x+0.3,3.3,2.9,0.5,title,sz=20,c=D,b=True,fn=FH,a='center')
        _txt(s,x+0.3,3.9,2.9,2.5,desc,sz=14,c=BD,fn=FB,a='center')

    # Before & After — full height
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,0.1,P)
    _txt(s,1.0,0.5,11,0.6,'Before & After',sz=30,c=D,b=True,fn=FH)
    _txt(s,1.0,1.1,11,0.3,'Client transformation stories',sz=14,c=MU,fn=FB)
    metrics=[('Time to Market','6 mo','8 wk',0.9,0.3),('Design Consistency','40%','95%',0.4,0.95),('User Retention','32%','78%',0.32,0.78)]
    for i,(lbl,vo,vn,po,pn) in enumerate(metrics):
        y=2.2+i*1.6
        _txt(s,1.2,y,2.0,0.4,lbl,sz=15,c=D,b=True,fn=FB)
        _rrect(s,3.5,y+0.4,5.0,0.4,'#F0F0F5')
        _rrect(s,3.5,y+0.4,5.0*po,0.4,MU)
        _txt(s,8.7,y+0.4,1.0,0.4,vo,sz=13,c=MU,fn=FB)
        _rrect(s,3.5,y+0.9,5.0,0.4,'#F0F0F5')
        _rrect(s,3.5,y+0.9,5.0*pn,0.4,P)
        _txt(s,8.7,y+0.9,1.0,0.4,vn,sz=13,c=D,b=True,fn=FB)

    # CTA
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s,0,0,SW,SH,P)
    _oval(s,9.5,0.5,5,5,A); _oval(s,-1.5,5.5,4,4,M); _oval(s,11,6,3,3,Y)
    _txt(s,1.5,2.8,8,1.2,"Let's Create Together",sz=44,c=W,b=True,fn=FH)
    _txt(s,1.5,4.2,8,0.5,'hello@creativestudio.io  |  book a free consultation',sz=16,c=W,fn=FB)

    prs.save(out)


def main():
    d=os.path.join(tempfile.gettempdir(),"ppt_raw_v2"); os.makedirs(d,exist_ok=True)
    print("Generating...")
    build_mckinsey(os.path.join(d,'mckinsey.pptx'))
    build_cyberpunk(os.path.join(d,'cyberpunk.pptx'))
    build_creative(os.path.join(d,'creative.pptx'))
    for n in ['mckinsey.pptx','cyberpunk.pptx','creative.pptx']:
        p=os.path.join(d,n); print(f"  {n}: {round(os.path.getsize(p)/1024,1)} KB")
    print("Done:",d)

if __name__=="__main__": main()
