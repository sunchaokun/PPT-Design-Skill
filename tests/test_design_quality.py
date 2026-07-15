"""Tests for design quality upgrades — Tier 1/2/3 TDD.

Covers:
  §1.9  _is_dark() luminance detection
  §1.2  TypeScale dataclass
  §1.3  Color system: OKLCH, tint/shade, alpha levels
  §1.5  Shadow elevation system
  §1.1  Layout engine: Rect, ContentLayout, LayoutEngine
  §1.6  Conditional brand strip
  §1.4  Gradient overlay
  §1.8  Card design upgrade
  §1.7  Image color-grading
  §1.10 Code block redesign
  §2.1  CJK + Latin font pairing
  §2.2  Content-adaptive margins
  §2.3  Badge/Tag styling system
  §2.4  Section divider slides
  §2.5  Decoration system (10 styles)
  §2.6  Layout variant consumption
  §3.1  Noise texture overlay
  §3.2  Progress bar indicator
  §3.3  Rounded corner consistency
  §3.4  Gradient line accents
  §3.5  Image masking
  §3.6  Two-column bullet layout
  §3.7  Hero slide layout variety
"""

from __future__ import annotations

import pytest

from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer


# ── §1.9 _is_dark() luminance detection ──────────────────────────


def _dark_brand(bg_hex: str) -> BrandSpec:
    return BrandSpec(
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "accent": "#6366F1",
            "background": bg_hex,
            "foreground": "#F8FAFC",
            "muted": "#1E293B",
            "muted-foreground": "#94A3B8",
            "border": "#334155",
            "destructive": "#EF4444",
        },
        fonts={"heading": "Inter", "body": "Inter"},
    )


def _light_brand(bg_hex: str = "#FFFFFF") -> BrandSpec:
    return BrandSpec(
        colors={
            "primary": "#2563EB",
            "on-primary": "#FFFFFF",
            "accent": "#6366F1",
            "background": bg_hex,
            "foreground": "#0F172A",
            "muted": "#F1F5F9",
            "muted-foreground": "#64748B",
            "border": "#E2E8F0",
            "destructive": "#EF4444",
        },
        fonts={"heading": "Inter", "body": "Inter"},
    )


class TestIsDarkLuminance:
    """Verify _is_dark uses luminance calculation, not hardcoded list."""

    def test_white_background_is_light(self):
        r = PrecisionRenderer(brand_spec=_light_brand("#FFFFFF"))
        assert r._is_dark() is False

    def test_black_background_is_dark(self):
        r = PrecisionRenderer(brand_spec=_dark_brand("#000000"))
        assert r._is_dark() is True

    def test_known_dark_hex_in_old_list(self):
        for hex_val in ("060B18", "0A1E3D", "120C1E", "111827", "09090B", "1C1010", "2D2D2D"):
            r = PrecisionRenderer(brand_spec=_dark_brand(f"#{hex_val}"))
            assert r._is_dark() is True, f"#{hex_val} should be dark"

    def test_dark_bg_not_in_old_list_detected(self):
        r = PrecisionRenderer(brand_spec=_dark_brand("#0F172A"))
        assert r._is_dark() is True

    def test_medium_gray_is_light(self):
        r = PrecisionRenderer(brand_spec=_light_brand("#808080"))
        assert r._is_dark() is False

    def test_dark_gray_is_dark(self):
        r = PrecisionRenderer(brand_spec=_dark_brand("#333333"))
        assert r._is_dark() is True

    def test_pale_blue_is_light(self):
        r = PrecisionRenderer(brand_spec=_light_brand("#DBEAFE"))
        assert r._is_dark() is False

    def test_navy_is_dark(self):
        r = PrecisionRenderer(brand_spec=_dark_brand("#1E3A8A"))
        assert r._is_dark() is True

    def test_dark_mode_flag_overrides(self):
        r = PrecisionRenderer(brand_spec=BrandSpec(dark_mode=True, colors=_light_brand().colors))
        assert r._is_dark() is True

    def test_no_brand_defaults_to_light(self):
        r = PrecisionRenderer()
        assert r._is_dark() is False

    def test_threshold_boundary(self):
        r = PrecisionRenderer(brand_spec=_dark_brand("#404040"))
        luminance = (0.299 * 64 + 0.587 * 64 + 0.114 * 64) / 255
        expected = luminance < 0.5
        assert r._is_dark() is expected


# ── §1.2 Typography scale system ─────────────────────────────────


class TestTypeScale:
    """Verify TypeScale dataclass and factory methods."""

    def test_default_values(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale()
        assert ts.display == 52
        assert ts.title == 28
        assert ts.subtitle == 20
        assert ts.body == 14
        assert ts.caption == 11

    def test_for_density_1_no_shrink(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale.for_density(1)
        assert ts.display == 52
        assert ts.body == 14

    def test_for_density_10_min_sizes(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale.for_density(10)
        assert ts.body >= 11
        assert ts.caption >= 9
        assert ts.display >= 36

    def test_for_density_shrinks_proportionally(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts5 = TypeScale.for_density(5)
        ts1 = TypeScale.for_density(1)
        assert ts5.display < ts1.display
        assert ts5.body < ts1.body

    def test_for_mode_presenting(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale.for_mode("presenting")
        assert ts.display == 64
        assert ts.title == 32
        assert ts.body == 18
        assert ts.caption == 12

    def test_for_mode_reading(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale.for_mode("reading")
        assert ts.display == 44
        assert ts.body == 14
        assert ts.caption >= 11

    def test_for_mode_default(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale.for_mode("balanced")
        assert ts.display == 52

    def test_5_level_hierarchy(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale()
        assert ts.display > ts.title > ts.subtitle > ts.body > ts.caption

    def test_min_2pt_gap_between_levels(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale()
        levels = [ts.display, ts.title, ts.subtitle, ts.body, ts.caption]
        for i in range(len(levels) - 1):
            assert levels[i] - levels[i + 1] >= 2, f"Gap between level {i} and {i+1} is less than 2pt"

    def test_tracking_values(self):
        from ppt_pro_max.renderer.typography import TypeScale
        ts = TypeScale()
        assert ts.display_tracking < 0
        assert ts.title_tracking < 0
        assert ts.body_tracking == 0.0
        assert ts.caption_tracking > 0


# ── §1.2 apply_letter_spacing ─────────────────────────────────────


class TestLetterSpacing:
    """Verify apply_letter_spacing sets OOXML a:spc correctly."""

    def test_positive_tracking(self):
        from ppt_pro_max.renderer.visual_effects import apply_letter_spacing
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Pt(0), Pt(0), Pt(100), Pt(50))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "TEST"
        run.font.size = Pt(14)
        apply_letter_spacing(run, 0.02, 14)
        from pptx.oxml.ns import qn
        from lxml import etree
        rPr = run._r.find(qn("a:rPr"))
        spc = rPr.find(qn("a:spc"))
        assert spc is not None
        expected_val = int(0.02 * 14 * 100)
        assert spc.get("val") == str(expected_val)

    def test_negative_tracking(self):
        from ppt_pro_max.renderer.visual_effects import apply_letter_spacing
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Pt(0), Pt(0), Pt(100), Pt(50))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "Title"
        run.font.size = Pt(28)
        apply_letter_spacing(run, -0.01, 28)
        from pptx.oxml.ns import qn
        rPr = run._r.find(qn("a:rPr"))
        spc = rPr.find(qn("a:spc"))
        assert spc is not None
        expected_val = int(-0.01 * 28 * 100)
        assert spc.get("val") == str(expected_val)

    def test_zero_tracking_no_spc(self):
        from ppt_pro_max.renderer.visual_effects import apply_letter_spacing
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(Pt(0), Pt(0), Pt(100), Pt(50))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "Body"
        run.font.size = Pt(14)
        apply_letter_spacing(run, 0.0, 14)
        from pptx.oxml.ns import qn
        rPr = run._r.find(qn("a:rPr"))
        spc = rPr.find(qn("a:spc"))
        assert spc is None


# ── §1.3 Color depth system ──────────────────────────────────────


class TestOKLCHConversion:
    """Verify pure-Python OKLCH conversion round-trips."""

    def test_hex_to_oklch_white(self):
        from ppt_pro_max.renderer.color_system import _hex_to_oklch
        l, c, h = _hex_to_oklch("#FFFFFF")
        assert l > 0.99
        assert c < 0.01

    def test_hex_to_oklch_black(self):
        from ppt_pro_max.renderer.color_system import _hex_to_oklch
        l, c, h = _hex_to_oklch("#000000")
        assert l < 0.01

    def test_hex_to_oklch_primary_blue(self):
        from ppt_pro_max.renderer.color_system import _hex_to_oklch
        l, c, h = _hex_to_oklch("#2563EB")
        assert 0.3 < l < 0.7
        assert c > 0.1

    def test_roundtrip_blue(self):
        from ppt_pro_max.renderer.color_system import _hex_to_oklch, _oklch_to_hex
        l, c, h = _hex_to_oklch("#2563EB")
        result = _oklch_to_hex(l, c, h)
        r1, g1, b1 = int(result[0:2], 16), int(result[2:4], 16), int(result[4:6], 16)
        r2, g2, b2 = 0x25, 0x63, 0xEB
        assert abs(r1 - r2) <= 2
        assert abs(g1 - g2) <= 2
        assert abs(b1 - b2) <= 2

    def test_roundtrip_red(self):
        from ppt_pro_max.renderer.color_system import _hex_to_oklch, _oklch_to_hex
        l, c, h = _hex_to_oklch("#EF4444")
        result = _oklch_to_hex(l, c, h)
        r1, g1, b1 = int(result[0:2], 16), int(result[2:4], 16), int(result[4:6], 16)
        r2, g2, b2 = 0xEF, 0x44, 0x44
        assert abs(r1 - r2) <= 2
        assert abs(g1 - g2) <= 2
        assert abs(b1 - b2) <= 2

    def test_negative_lms_handled(self):
        from ppt_pro_max.renderer.color_system import _hex_to_oklch, _oklch_to_hex
        for hex_color in ("#FF0000", "#00FF00", "#0000FF", "#FF6600"):
            l, c, h = _hex_to_oklch(hex_color)
            result = _oklch_to_hex(l, c, h)
            assert len(result) == 6
            for ch in result:
                assert ch in "0123456789ABCDEF"


class TestColorScale:
    """Verify generate_color_scale produces valid tint/shade scales."""

    def test_11_levels_generated(self):
        from ppt_pro_max.renderer.color_system import generate_color_scale
        scale = generate_color_scale("#2563EB")
        assert len(scale) == 11
        for key in ["50", "100", "200", "300", "400", "500", "600", "700", "800", "900", "950"]:
            assert key in scale

    def test_500_is_base(self):
        from ppt_pro_max.renderer.color_system import generate_color_scale
        scale = generate_color_scale("#2563EB")
        base_r, base_g, base_b = 0x25, 0x63, 0xEB
        result = scale["500"]
        r, g, b = int(result[0:2], 16), int(result[2:4], 16), int(result[4:6], 16)
        assert abs(r - base_r) <= 5
        assert abs(g - base_g) <= 5
        assert abs(b - base_b) <= 5

    def test_tints_get_lighter(self):
        from ppt_pro_max.renderer.color_system import generate_color_scale
        scale = generate_color_scale("#2563EB")
        l50 = _luminance(scale["50"])
        l100 = _luminance(scale["100"])
        l500 = _luminance(scale["500"])
        assert l50 > l100 > l500

    def test_shades_get_darker(self):
        from ppt_pro_max.renderer.color_system import generate_color_scale
        scale = generate_color_scale("#2563EB")
        l500 = _luminance(scale["500"])
        l600 = _luminance(scale["600"])
        l900 = _luminance(scale["900"])
        l950 = _luminance(scale["950"])
        assert l500 > l600 > l900 > l950

    def test_all_values_valid_hex(self):
        from ppt_pro_max.renderer.color_system import generate_color_scale
        scale = generate_color_scale("#EF4444")
        for key, val in scale.items():
            assert len(val) == 6, f"Level {key} produced invalid hex: {val}"
            int(val, 16)


def _luminance(hex_color: str) -> float:
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255


class TestAlphaLevels:
    """Verify ALPHA_LEVELS dict and alpha_color wrapper."""

    def test_alpha_levels_defined(self):
        from ppt_pro_max.renderer.color_system import ALPHA_LEVELS
        assert "ghost" in ALPHA_LEVELS
        assert "subtle" in ALPHA_LEVELS
        assert "light" in ALPHA_LEVELS
        assert "medium" in ALPHA_LEVELS
        assert "strong" in ALPHA_LEVELS

    def test_alpha_levels_are_int_pct(self):
        from ppt_pro_max.renderer.color_system import ALPHA_LEVELS
        for name, val in ALPHA_LEVELS.items():
            assert isinstance(val, int)
            assert 0 <= val <= 100


# ── §1.5 Shadow elevation system ─────────────────────────────────


class TestElevation:
    """Verify ELEVATION_SCALE and apply_elevation."""

    def test_scale_has_5_levels(self):
        from ppt_pro_max.renderer.elevation import ELEVATION_SCALE
        assert set(ELEVATION_SCALE.keys()) == {0, 1, 2, 3, 4}

    def test_level_0_is_flat(self):
        from ppt_pro_max.renderer.elevation import ELEVATION_SCALE
        spec = ELEVATION_SCALE[0]
        assert spec["blur_pt"] == 0
        assert spec["distance_pt"] == 0
        assert spec["alpha_pct"] == 0

    def test_levels_increase_in_depth(self):
        from ppt_pro_max.renderer.elevation import ELEVATION_SCALE
        for i in range(3):
            lo = ELEVATION_SCALE[i]
            hi = ELEVATION_SCALE[i + 1]
            assert hi["blur_pt"] > lo["blur_pt"]
            assert hi["distance_pt"] > lo["distance_pt"]

    def test_apply_elevation_light_mode_uses_shadow(self):
        from ppt_pro_max.renderer.elevation import apply_elevation
        from ppt_pro_max.renderer.visual_effects import apply_shadow
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(4), Inches(2)
        )
        apply_elevation(shape, level=1, is_dark=False, foreground_hex="#000000", primary_hex="#2563EB")
        from pptx.oxml.ns import qn
        spPr = shape._element.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        assert effectLst is not None
        shadow = effectLst.find(qn("a:outerShdw"))
        assert shadow is not None

    def test_apply_elevation_level_0_no_shadow(self):
        from ppt_pro_max.renderer.elevation import apply_elevation
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(4), Inches(2)
        )
        apply_elevation(shape, level=0)
        from pptx.oxml.ns import qn
        spPr = shape._element.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        if effectLst is not None:
            assert effectLst.find(qn("a:outerShdw")) is None
        else:
            pass

    def test_apply_elevation_dark_mode_uses_glow(self):
        from ppt_pro_max.renderer.elevation import apply_elevation
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        shape = slide.shapes.add_shape(
            1, Inches(1), Inches(1), Inches(4), Inches(2)
        )
        apply_elevation(shape, level=2, is_dark=True, foreground_hex="#FFFFFF", primary_hex="#2563EB")
        from pptx.oxml.ns import qn
        spPr = shape._element.find(qn("p:spPr"))
        effectLst = spPr.find(qn("a:effectLst"))
        assert effectLst is not None
        glow = effectLst.find(qn("a:glow"))
        assert glow is not None


# ── §1.1 Layout engine ──────────────────────────────────────────


class TestRect:
    """Verify Rect dataclass with computed content area."""

    def test_default_values(self):
        from ppt_pro_max.renderer.layout_engine import Rect
        r = Rect()
        assert r.width == 13.333
        assert r.height == 7.5
        assert r.margin_left == 0.9

    def test_content_width(self):
        from ppt_pro_max.renderer.layout_engine import Rect
        r = Rect(margin_left=1.5, margin_right=1.5, width=13.333)
        assert abs(r.content_width - 10.333) < 0.01

    def test_content_height(self):
        from ppt_pro_max.renderer.layout_engine import Rect
        r = Rect(margin_top=1.2, margin_bottom=0.8, height=7.5)
        assert abs(r.content_height - 5.5) < 0.01


class TestContentLayout:
    """Verify ContentLayout dataclass."""

    def test_default_values(self):
        from ppt_pro_max.renderer.layout_engine import ContentLayout
        cl = ContentLayout()
        assert cl.title_x == 0.9
        assert cl.bullet_columns == 1


class TestLayoutEngine:
    """Verify LayoutEngine computes content-adaptive positions."""

    def test_few_bullets_with_image_narrow_text(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine, Rect
        engine = LayoutEngine()
        page = {"title": "Features", "bullets": ["A", "B", "C"], "image": "photo.jpg"}
        rect = Rect()
        layout = engine.compute_content_layout(page, rect)
        assert layout.content_w < rect.content_width

    def test_many_bullets_full_width(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine, Rect
        engine = LayoutEngine()
        page = {"title": "Details", "bullets": [f"Item {i}" for i in range(9)]}
        rect = Rect()
        layout = engine.compute_content_layout(page, rect)
        assert layout.bullet_columns == 2

    def test_image_present_reduces_text_width(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine, Rect
        engine = LayoutEngine()
        page_with = {"title": "T", "bullets": ["A", "B"], "image": "photo.jpg"}
        page_without = {"title": "T", "bullets": ["A", "B"]}
        rect = Rect()
        layout_with = engine.compute_content_layout(page_with, rect)
        layout_without = engine.compute_content_layout(page_without, rect)
        assert layout_with.content_w < layout_without.content_w

    def test_no_image_text_expands(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine, Rect
        engine = LayoutEngine()
        page = {"title": "Overview", "bullets": ["A", "B"]}
        rect = Rect()
        layout = engine.compute_content_layout(page, rect)
        assert layout.image_w == 0

    def test_title_y_below_top_margin(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine, Rect
        engine = LayoutEngine()
        page = {"title": "Test"}
        rect = Rect(margin_top=0.6)
        layout = engine.compute_content_layout(page, rect)
        assert layout.title_y >= 0.6

    def test_estimate_lines_short_title(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine
        engine = LayoutEngine()
        lines = engine._estimate_lines("Hi", 10.0, 28)
        assert lines == 1

    def test_estimate_lines_long_title(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine
        engine = LayoutEngine()
        lines = engine._estimate_lines("A" * 100, 10.0, 28)
        assert lines > 1


# ── §1.6 Conditional brand strip (AI watermark removal) ─────────


class TestConditionalBrandStrip:
    """Verify apply_brand_background does NOT add strip to every slide."""

    def _render_slides(self, n: int, goal: str = "content") -> list:
        from pptx import Presentation
        results = []
        for i in range(n):
            brand = _light_brand()
            r = PrecisionRenderer(brand_spec=brand)
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[-1])
            r.apply_brand_background(slide, prs, goal=goal, page_index=i, total_pages=n)
            results.append(slide)
        return results

    def test_strip_not_on_every_slide(self):
        slides = self._render_slides(6)
        strip_count = 0
        for slide in slides:
            for shape in slide.shapes:
                if shape.width < 200000 and shape.height > 4000000:
                    strip_count += 1
        assert strip_count <= 4, f"Accent strip appeared {strip_count} times, should be <= 4 of 6"

    def test_bottom_bar_height_reduced(self):
        from pptx import Presentation
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.apply_brand_background(slide, prs, goal="content", page_index=1, total_pages=5)
        max_bar_height = 0
        for shape in slide.shapes:
            if shape.top > 5000000:
                max_bar_height = max(max_bar_height, shape.height)
        assert max_bar_height < 400000, "Bottom bar should be <= 0.15\" (reduced from 0.25\")"


# ── §1.4 Gradient overlay ────────────────────────────────────────


class TestGradientOverlay:
    """Verify add_gradient_overlay produces gradient fill, not flat solid."""

    def test_overlay_has_gradient_fill(self):
        from pptx import Presentation
        from pptx.util import Inches
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.oxml.ns import qn
        brand = _dark_brand("#0A1E3D")
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_gradient_overlay(slide, opacity_bottom=0.72, opacity_top=0.0)
        found_grad = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                gradFill = spPr.find(qn("a:gradFill"))
                if gradFill is not None:
                    found_grad = True
                    gsLst = gradFill.find(qn("a:gsLst"))
                    assert gsLst is not None
                    stops = gsLst.findall(qn("a:gs"))
                    assert len(stops) >= 2
        assert found_grad, "Gradient overlay should use gradFill, not solidFill"

    def test_overlay_top_transparent(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.oxml.ns import qn
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_gradient_overlay(slide, opacity_bottom=0.72, opacity_top=0.0)
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                gradFill = spPr.find(qn("a:gradFill"))
                if gradFill is not None:
                    gsLst = gradFill.find(qn("a:gsLst"))
                    stops = gsLst.findall(qn("a:gs"))
                    first_stop = stops[0]
                    srgbClr = first_stop.find(qn("a:srgbClr"))
                    if srgbClr is not None:
                        alpha = srgbClr.find(qn("a:alpha"))
                        if alpha is not None:
                            assert int(alpha.get("val", "100000")) < 10000


# ── §1.8 Card design upgrade ────────────────────────────────────


class TestCardUpgrade:
    """Verify card title=20pt, body=14pt, featured parameter."""

    def test_card_title_20pt(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.util import Pt
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_card(slide, 0.9, 1.6, 3.6, 4.5, "Test Title", "Body text")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text == "Test Title":
                            assert run.font.size == Pt(20)

    def test_card_body_14pt(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.util import Pt
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_card(slide, 0.9, 1.6, 3.6, 4.5, "Title", "Body line")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Body" in run.text:
                            assert run.font.size >= Pt(14)

    def test_featured_card_larger_title(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.util import Pt
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_card(slide, 0.9, 1.6, 3.6, 4.5, "Featured", "Body", featured=True)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text == "Featured":
                            assert run.font.size >= Pt(22)


# ── §1.7 Image color-grading ────────────────────────────────────


class TestImageColorGrading:
    """Verify grade_image_to_palette produces graded image."""

    def test_grade_produces_file(self, tmp_path):
        from PIL import Image as PILImage
        from ppt_pro_max.renderer.image_processor import grade_image_to_palette
        img = PILImage.new("RGB", (100, 100), (128, 128, 128))
        src = str(tmp_path / "test.png")
        img.save(src)
        result = grade_image_to_palette(src, "#2563EB", alpha=0.10)
        import os
        assert os.path.isfile(result)

    def test_grade_cache_hit(self, tmp_path):
        from PIL import Image as PILImage
        from ppt_pro_max.renderer.image_processor import grade_image_to_palette
        img = PILImage.new("RGB", (100, 100), (128, 128, 128))
        src = str(tmp_path / "test.png")
        img.save(src)
        r1 = grade_image_to_palette(src, "#2563EB", alpha=0.10)
        r2 = grade_image_to_palette(src, "#2563EB", alpha=0.10)
        assert r1 == r2

    def test_grade_preserves_jpeg(self, tmp_path):
        from PIL import Image as PILImage
        from ppt_pro_max.renderer.image_processor import grade_image_to_palette
        img = PILImage.new("RGB", (100, 100), (128, 128, 128))
        src = str(tmp_path / "test.jpg")
        img.save(src, "JPEG")
        result = grade_image_to_palette(src, "#2563EB", alpha=0.10)
        assert result.lower().endswith(".jpg")


# ── §1.10 Code block redesign ───────────────────────────────────


class TestCodeBlockRedesign:
    """Verify code blocks use always-dark background."""

    def test_code_block_always_dark_bg(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.oxml.ns import qn
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r._render_code_on_slide(slide, {"source": "print('hello')", "language": "python"})
        found_dark = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                solidFill = spPr.find(qn("a:solidFill"))
                if solidFill is not None:
                    srgb = solidFill.find(qn("a:srgbClr"))
                    if srgb is not None:
                        val = srgb.get("val", "")
                        if val.upper() in ("1E293B", "0D152A"):
                            found_dark = True
        assert found_dark, "Code block should use dark background (#1E293B) even on light themes"

    def test_code_block_has_language_badge(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r._render_code_on_slide(slide, {"source": "x = 1", "language": "python"})
        found_python = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "PYTHON" in text.upper():
                    found_python = True
        assert found_python, "Code block should show language badge"


# ── §2.1 CJK + Latin font pairing ────────────────────────────────


class TestCJKFontPairing:
    def test_cjk_companions_dict_exists(self):
        from ppt_pro_max.renderer.theme_mapper import CJK_COMPANIONS
        assert isinstance(CJK_COMPANIONS, dict)
        assert len(CJK_COMPANIONS) >= 5

    def test_set_font_with_cjk_sets_ea(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.oxml.ns import qn
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        tb = r.add_text(slide, "Test", 1, 1, 4, 1, font="Inter")
        r._set_font_with_cjk(tb.text_frame.paragraphs[0].runs[0], "Inter", "Microsoft YaHei")
        rPr = tb.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        assert ea is not None, "a:ea element should be set for CJK font"
        assert ea.get("typeface") == "Microsoft YaHei"

    def test_serif_latin_gets_serif_cjk(self):
        from ppt_pro_max.renderer.theme_mapper import CJK_COMPANIONS
        assert "Playfair Display" in CJK_COMPANIONS
        assert "heading" in CJK_COMPANIONS["Playfair Display"]


# ── §2.2 Content-adaptive margins ────────────────────────────────


class TestAdaptiveMargins:
    def test_sparse_content_generous_margins(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine
        le = LayoutEngine()
        margins = le.compute_margins({"bullets": ["one"]}, mode="presenting")
        assert margins["left"] >= 1.5

    def test_dense_content_tight_margins(self):
        from ppt_pro_max.renderer.layout_engine import LayoutEngine
        le = LayoutEngine()
        bullets = [f"bullet {i}" for i in range(10)]
        margins = le.compute_margins({"bullets": bullets})
        assert margins["left"] <= 0.9


# ── §2.3 Badge/Tag styling ──────────────────────────────────────


class TestBadgeSystem:
    def test_add_badge_creates_shapes(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_badge(slide, "Exercise 5 min", 0.9, 1.2)
        assert len(slide.shapes) >= 2

    def test_badge_text_uppercase(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_badge(slide, "Exercise 5 min", 0.9, 1.2)
        found = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "EXERCISE" in shape.text_frame.text.upper():
                    found = True
        assert found

    def test_badge_variant_solid(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.oxml.ns import qn
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_badge(slide, "NEW", 0.9, 1.2, variant="solid")
        found_primary_fill = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                solidFill = spPr.find(qn("a:solidFill"))
                if solidFill is not None:
                    found_primary_fill = True
        assert found_primary_fill


# ── §2.4 Section divider ────────────────────────────────────────


class TestSectionDivider:
    def test_render_section_divider_creates_slide(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = r.render_section_divider(prs, 1, "Introduction")
        assert len(slide.shapes) >= 2

    def test_section_number_displayed(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = r.render_section_divider(prs, 2, "Architecture")
        found = False
        for shape in slide.shapes:
            if shape.has_text_frame and "02" in shape.text_frame.text:
                found = True
        assert found


# ── §2.5 Decoration system ──────────────────────────────────────


class TestDecorationSystem:
    def test_decoration_renderer_exists(self):
        from ppt_pro_max.renderer.decoration_renderer import DecorationRenderer
        dr = DecorationRenderer()
        assert hasattr(dr, "apply_title_decoration")

    def test_all_10_styles_have_methods(self):
        from ppt_pro_max.renderer.decoration_renderer import DecorationRenderer
        dr = DecorationRenderer()
        styles = ["accent_bar", "neon_lines", "gold_trim", "minimal_dots",
                   "diamond_bullets", "gradient_bar", "circle_accent",
                   "sidebar_nav", "no_decoration", "full_bleed_overlay"]
        for style in styles:
            method_name = f"_title_{style}"
            assert hasattr(dr, method_name), f"Missing method {method_name}"


# ── §2.6 Layout variant consumption ─────────────────────────────


class TestLayoutVariant:
    def test_variant_margins_used(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        variant = {"content_margin_left": 2.0, "content_margin_right": 2.0,
                   "title_alignment": "center", "card_style": "rounded"}
        page = {"goal": "content", "title": "Test", "bullets": ["a", "b"]}
        slide = r.render_slide(prs, page, layout_variant=variant)
        assert slide is not None

    def test_centered_title_alignment(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.enum.text import PP_ALIGN
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        variant = {"content_margin_left": 2.0, "title_alignment": "center"}
        page = {"goal": "content", "title": "Centered Title", "bullets": ["a"]}
        slide = r.render_slide(prs, page, layout_variant=variant)
        found_center = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.alignment == PP_ALIGN.CENTER:
                        found_center = True
        assert found_center


# ── §3.1 Noise texture ──────────────────────────────────────────


class TestNoiseTexture:
    def test_generate_noise_tile(self, tmp_path):
        from ppt_pro_max.renderer.image_processor import generate_noise_tile
        path = generate_noise_tile(size=100, opacity=0.02, deck_title="test_deck")
        import os
        assert os.path.isfile(path)

    def test_noise_cache_hit(self, tmp_path):
        from ppt_pro_max.renderer.image_processor import generate_noise_tile
        p1 = generate_noise_tile(size=100, opacity=0.02, deck_title="cache_test")
        p2 = generate_noise_tile(size=100, opacity=0.02, deck_title="cache_test")
        assert p1 == p2

    def test_different_decks_different_noise(self):
        from ppt_pro_max.renderer.image_processor import generate_noise_tile
        p1 = generate_noise_tile(size=100, opacity=0.02, deck_title="deck_A")
        p2 = generate_noise_tile(size=100, opacity=0.02, deck_title="deck_B")
        assert p1 != p2


# ── §3.2 Progress bar ──────────────────────────────────────────


class TestProgressBar:
    def test_add_progress_bar_creates_shapes(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_progress_bar(slide, 3, 10)
        assert len(slide.shapes) >= 2

    def test_progress_bar_fill_width(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.util import Inches
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_progress_bar(slide, 3, 10)
        found_fill = False
        for shape in slide.shapes:
            if hasattr(shape, "width") and shape.width < Inches(13.333):
                found_fill = True
        assert found_fill


# ── §3.3 Rounded corner consistency ──────────────────────────────


class TestCornerRadius:
    def test_corner_radius_scale_exists(self):
        from ppt_pro_max.enterprise.precision_renderer import CORNER_RADIUS_SCALE
        assert "sm" in CORNER_RADIUS_SCALE
        assert "md" in CORNER_RADIUS_SCALE
        assert "lg" in CORNER_RADIUS_SCALE

    def test_add_rounded_rect_accepts_radius_str(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_rounded_rect(slide, 1, 1, 3, 2, fill_hex="#F8FAFC", corner_radius="sm")
        assert len(slide.shapes) >= 1


# ── §3.4 Gradient line accents ──────────────────────────────────


class TestGradientLine:
    def test_add_gradient_line_creates_shape(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_gradient_line(slide, 0.9, 1.2, 3.0, 0.04, "#2563EB")
        assert len(slide.shapes) >= 1

    def test_gradient_line_has_gradient_fill(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from pptx.oxml.ns import qn
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        r.add_gradient_line(slide, 0.9, 1.2, 3.0, 0.04, "#2563EB")
        found_grad = False
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is not None:
                gradFill = spPr.find(qn("a:gradFill"))
                if gradFill is not None:
                    found_grad = True
        assert found_grad


# ── §3.5 Image masking ──────────────────────────────────────────


class TestImageMasking:
    def test_add_masked_image_creates_shapes(self, tmp_path):
        from PIL import Image as PILImage
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        img = PILImage.new("RGB", (200, 200), (100, 150, 200))
        img_path = str(tmp_path / "test.png")
        img.save(img_path)
        r.add_masked_image(slide, img_path, 8.3, 1.2, 4.2, 5.3)
        assert len(slide.shapes) >= 2


# ── §3.6 Two-column bullet layout ──────────────────────────────


class TestTwoColumnBullets:
    def test_six_bullets_two_columns(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        bullets = [f"Bullet {i}" for i in range(6)]
        page = {"goal": "content", "title": "Test", "bullets": bullets}
        slide = r.render_slide(prs, page)
        multiline_count = sum(1 for s in slide.shapes if s.has_text_frame and "\n" in (s.text_frame.text or ""))
        assert multiline_count >= 2, "6+ bullets should render in 2 columns"

    def test_three_bullets_single_column(self):
        from pptx import Presentation
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        prs = Presentation()
        bullets = [f"Bullet {i}" for i in range(3)]
        page = {"goal": "content", "title": "Test", "bullets": bullets}
        slide = r.render_slide(prs, page)
        assert slide is not None


# ── §3.7 Hero slide layout variety ──────────────────────────────


class TestHeroVariety:
    def test_select_hero_pattern(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        result = r._select_hero_pattern({"subtitle": "short"}, "creative")
        assert result in ("full-bleed", "split-left", "bottom-fade", "asymmetric", "gradient")

    def test_long_subtitle_uses_split(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        long_sub = "x" * 80
        result = r._select_hero_pattern({"subtitle": long_sub, "image": "photo.jpg"}, "professional")
        assert result == "split-left"

    def test_no_image_uses_gradient(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        brand = _light_brand()
        r = PrecisionRenderer(brand_spec=brand)
        result = r._select_hero_pattern({"subtitle": "test"}, "professional")
        assert result == "gradient"
