"""Tests for PrecisionRenderer.render_slide content rendering."""

from __future__ import annotations

import os
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def _make_png(path: Path, w: int = 200, h: int = 150) -> Path:
    img = Image.new("RGB", (w, h), color="green")
    img.save(str(path))
    return path


def _has_text_in_slide(slide, expected: str) -> bool:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if expected in para.text:
                    return True
    return False


def _count_pictures(slide) -> int:
    return sum(1 for s in slide.shapes if s.shape_type == 13)


class TestPopulateSlide:

    def test_title_populated(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        precision.render_slide(prs, {"title": "Hello World"})
        slide = prs.slides[-1]
        assert _has_text_in_slide(slide, "Hello World")

    def test_subtitle_populated(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        precision.render_slide(prs, {"title": "Hi", "subtitle": "Sub text"})
        slide = prs.slides[-1]
        assert _has_text_in_slide(slide, "Hi")
        assert _has_text_in_slide(slide, "Sub text")

    def test_bullets_populated(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        precision.render_slide(prs, {"title": "Points", "bullets": ["First", "Second", "Third"]})
        slide = prs.slides[-1]
        assert _has_text_in_slide(slide, "Points")
        assert _has_text_in_slide(slide, "First")
        assert _has_text_in_slide(slide, "Second")
        assert _has_text_in_slide(slide, "Third")

    def test_bullets_with_existing_markers(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        precision.render_slide(prs, {"title": "T", "bullets": ["• Already marked", "— Dash marked"]})
        slide = prs.slides[-1]
        assert _has_text_in_slide(slide, "• Already marked")
        assert _has_text_in_slide(slide, "— Dash marked")

    def test_image_inserted(self, tmp_path):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        img_path = _make_png(tmp_path / "content.png")
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        precision.render_slide(prs, {"title": "With Image", "image": str(img_path)})
        slide = prs.slides[-1]
        assert _count_pictures(slide) == 1

    def test_image_nonexistent_no_crash(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        precision.render_slide(prs, {"title": "No Image", "image": "/nonexistent/img.png"})
        slide = prs.slides[-1]
        assert _count_pictures(slide) == 0

    def test_empty_design_no_crash(self):
        from ppt_pro_max.enterprise.precision_renderer import PrecisionRenderer
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        precision = PrecisionRenderer(brand_spec=BrandSpec())
        prs = precision.create_presentation()
        precision.render_slide(prs, {})

    def test_full_pipeline_with_bullets_and_image(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        import json
        project = tmp_path / "fullproject"
        project.mkdir()
        prs = Presentation()
        for _ in range(2):
            prs.slides.add_slide(prs.slide_layouts[-1])
        prs.save(str(project / "template.pptx"))
        _make_png(project / "logo.png")
        _make_png(project / "hero.png")
        (project / "brand.json").write_text(json.dumps({"colors": {"primary": "#1A3C6E"}}), encoding="utf-8")
        content = {
            "meta": {"title": "Full Test"},
            "slides": [
                {"goal": "hook", "title": "Welcome", "image": "hero.png"},
                {"goal": "problem", "title": "Pain Points", "bullets": ["Slow", "Expensive", "Complex"]},
            ],
        }
        (project / "content.json").write_text(json.dumps(content), encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(query="Full Test", project_dir=str(project))
        assert result["num_slides"] == 2
        prs_out = Presentation(result["output_path"])
        assert _has_text_in_slide(prs_out.slides[0], "Welcome")
        assert _has_text_in_slide(prs_out.slides[1], "Pain Points")
