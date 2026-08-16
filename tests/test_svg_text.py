"""Tests for svg_compiler._text — SVG text rendering.

Validates:
- Font family resolution (Arial, Georgia, etc.)
- Baseline mapping (dominant-baseline → MSO_ANCHOR)
- Text anchor (start/middle/end → PP_ALIGN)
- Pillow measurement fallback (when font unavailable)
- Simple text (no tspan children)
- <tspan> multi-line: each tspan with x/y → new paragraph
- <tspan> inline: tspan without x/y → same paragraph, different run style
- <tspan> per-span font-size, font-weight, font-style, fill
- Empty tspan content skipped
- _collect_spans and _group_spans_into_lines unit tests
"""
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from ppt_pro_max.renderer.svg_compiler._text import (
    _ANCHOR_MAP,
    _BASELINE_MAP,
    _BASELINE_OFFSET,
    _FONT_MAP,
    _SpanSpec,
    TextMetrics,
    _collect_spans,
    _compute_text_top,
    _group_spans_into_lines,
    _has_cjk,
    _measure_text,
    _resolve_baseline,
    _resolve_baseline_mode,
    _resolve_font_family,
    _parse_font_size,
    _parse_font_weight,
    render_svg_text,
)


class TestFontFamilyResolution:
    def test_arial_lowercase(self):
        assert _resolve_font_family("arial") == "Arial"

    def test_helvetica_alias(self):
        assert _resolve_font_family("Helvetica") == "Arial"

    def test_sans_serif_alias(self):
        assert _resolve_font_family("sans-serif") == "Arial"

    def test_serif_alias(self):
        assert _resolve_font_family("serif") == "Times New Roman"

    def test_monospace_alias(self):
        assert _resolve_font_family("monospace") == "Courier New"

    def test_consolas_direct(self):
        assert _resolve_font_family("Consolas") == "Consolas"

    def test_georgia(self):
        assert _resolve_font_family("georgia") == "Georgia"

    def test_empty_falls_back(self):
        assert _resolve_font_family(None) == "Calibri"
        assert _resolve_font_family("") == "Calibri"

    def test_unknown_returns_first(self):
        assert _resolve_font_family("Helvetica, Roboto") == "Arial"

    def test_font_map_has_at_least_10_fonts(self):
        assert len(_FONT_MAP) >= 10


class TestBaselineResolution:
    def test_auto_maps_to_middle(self):
        assert _resolve_baseline_baseline("auto") == MSO_ANCHOR.MIDDLE

    def test_alphabetic_maps_to_bottom(self):
        assert _resolve_baseline_baseline("alphabetic") == MSO_ANCHOR.BOTTOM

    def test_text_before_edge_maps_to_top(self):
        assert _resolve_baseline_baseline("text-before-edge") == MSO_ANCHOR.TOP

    def test_central_maps_to_middle(self):
        assert _resolve_baseline_baseline("central") == MSO_ANCHOR.MIDDLE

    def test_hanging_maps_to_top(self):
        assert _resolve_baseline_baseline("hanging") == MSO_ANCHOR.TOP

    def test_unknown_falls_back_to_middle(self):
        assert _resolve_baseline_baseline("xyz") == MSO_ANCHOR.MIDDLE

    def test_baseline_map_size(self):
        assert len(_BASELINE_MAP) >= 7


def _resolve_baseline_baseline(value: str):
    class FakeEl:
        def get(self, _, default=None):
            return value
    return _resolve_baseline(FakeEl())


class TestTextAnchor:
    def test_middle_maps_to_center(self):
        assert _ANCHOR_MAP["middle"] == PP_ALIGN.CENTER

    def test_end_maps_to_right(self):
        assert _ANCHOR_MAP["end"] == PP_ALIGN.RIGHT

    def test_start_maps_to_left(self):
        assert _ANCHOR_MAP["start"] == PP_ALIGN.LEFT


class TestRenderSvgTextIntegration:
    """Integration: full pipeline text rendering."""

    def test_simple_text_renders(self):
        from lxml import etree
        from pptx import Presentation
        from pptx.util import Inches

        from ppt_pro_max.renderer.svg_compiler._affine import Affine
        from ppt_pro_max.renderer.svg_compiler._text import render_svg_text

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        svg_text = '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100" font-size="16" fill="#FF0000">Hello</text>'
        el = etree.fromstring(svg_text)

        features = set()

        def resolve_color_fn(v, C, fb):
            return v if v.startswith("#") else None

        render_svg_text(
            el=el,
            tf=Affine(),
            to_inches_fn=lambda x, y: (x / 96, y / 96),
            slide=slide,
            C={},
            resolve_color_fn=resolve_color_fn,
            features=features,
        )

        assert "text" in features
        assert len(slide.shapes) > 0

    def test_empty_content_skipped(self):
        from lxml import etree
        from pptx import Presentation
        from pptx.util import Inches

        from ppt_pro_max.renderer.svg_compiler._affine import Affine
        from ppt_pro_max.renderer.svg_compiler._text import render_svg_text

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        svg_text = '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100">   </text>'
        el = etree.fromstring(svg_text)

        features = set()

        def resolve_color_fn(v, C, fb):
            return None

        render_svg_text(
            el=el,
            tf=Affine(),
            to_inches_fn=lambda x, y: (x / 96, y / 96),
            slide=slide,
            C={},
            resolve_color_fn=resolve_color_fn,
            features=features,
        )

        assert len(slide.shapes) == 0

    def test_grad_fill_falls_back_to_solid(self):
        from lxml import etree
        from pptx import Presentation
        from pptx.util import Inches

        from ppt_pro_max.renderer.svg_compiler._affine import Affine
        from ppt_pro_max.renderer.svg_compiler._text import render_svg_text

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        svg_text = '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100" fill="url(#g1)">Hello</text>'
        el = etree.fromstring(svg_text)

        features = set()

        def resolve_color_fn(v, C, fb):
            return "#000000"

        render_svg_text(
            el=el,
            tf=Affine(),
            to_inches_fn=lambda x, y: (x / 96, y / 96),
            slide=slide,
            C={},
            resolve_color_fn=resolve_color_fn,
            features=features,
        )

        assert len(slide.shapes) > 0


class TestParseFontSize:
    def test_numeric_string(self):
        assert _parse_font_size("16", 14.0) == 16.0

    def test_px_suffix(self):
        assert _parse_font_size("20px", 14.0) == 20.0

    def test_pt_suffix(self):
        assert _parse_font_size("12pt", 14.0) == 12.0

    def test_none_falls_back(self):
        assert _parse_font_size(None, 14.0) == 14.0

    def test_empty_falls_back(self):
        assert _parse_font_size("", 14.0) == 14.0

    def test_garbage_falls_back(self):
        assert _parse_font_size("abc", 14.0) == 14.0


class TestParseFontWeight:
    def test_bold(self):
        assert _parse_font_weight("bold") is True

    def test_700(self):
        assert _parse_font_weight("700") is True

    def test_normal(self):
        assert _parse_font_weight("normal") is False

    def test_200(self):
        assert _parse_font_weight("200") is False

    def test_none(self):
        assert _parse_font_weight(None) is False


class TestCollectSpans:
    def _make_el(self, svg_str):
        from lxml import etree
        return etree.fromstring(svg_str)

    def test_no_tspan_direct_text(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">Hello</text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert len(spans) == 1
        assert spans[0].text == "Hello"
        assert spans[0].is_new_line is False

    def test_single_tspan_with_xy(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            '<tspan x="10" y="40">Line1</tspan></text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert len(spans) == 1
        assert spans[0].text == "Line1"
        assert spans[0].is_new_line is True
        assert spans[0].x == 10.0
        assert spans[0].y == 40.0

    def test_two_tspan_lines(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            '<tspan x="10" y="40">Line1</tspan>'
            '<tspan x="10" y="60">Line2</tspan></text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert len(spans) == 2
        assert spans[0].is_new_line is True
        assert spans[1].is_new_line is True

    def test_inline_tspan_no_xy(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            'Hello <tspan font-weight="bold">World</tspan></text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert len(spans) == 2
        assert spans[0].text == "Hello"
        assert spans[0].is_new_line is False
        assert spans[1].text == "World"
        assert spans[1].bold is True
        assert spans[1].is_new_line is False

    def test_tspan_font_size_override(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            '<tspan font-size="24">Big</tspan></text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert spans[0].font_size == 24.0

    def test_tspan_fill_override(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            '<tspan fill="#FF0000">Red</tspan></text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert spans[0].fill == "#FF0000"

    def test_tspan_italic(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            '<tspan font-style="italic">Em</tspan></text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert spans[0].italic is True

    def test_tspan_dx_dy(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            '<tspan dx="5" dy="3">Shifted</tspan></text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert spans[0].dx == 5.0
        assert spans[0].dy == 3.0

    def test_tail_text_after_tspan(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            '<tspan>Bold</tspan> tail</text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        assert len(spans) == 2
        assert spans[0].text == "Bold"
        assert spans[1].text == "tail"

    def test_non_tspan_child_ignored(self):
        el = self._make_el(
            '<text xmlns="http://www.w3.org/2000/svg" x="10" y="20">'
            '<rect width="10" height="10"/>Text</text>'
        )
        spans = _collect_spans(el, 14.0, "Arial", "#000000", {}, lambda v, C, fb: v)
        texts = [s.text for s in spans if s.text]
        assert "Text" in texts


class TestGroupSpansIntoLines:
    def test_single_line_no_newline(self):
        spans = [_SpanSpec(text="A"), _SpanSpec(text="B")]
        lines = _group_spans_into_lines(spans)
        assert len(lines) == 1
        assert len(lines[0]) == 2

    def test_two_lines(self):
        spans = [
            _SpanSpec(text="A", is_new_line=True),
            _SpanSpec(text="B", is_new_line=False),
            _SpanSpec(text="C", is_new_line=True),
        ]
        lines = _group_spans_into_lines(spans)
        assert len(lines) == 2
        assert lines[0][0].text == "A"
        assert lines[0][1].text == "B"
        assert lines[1][0].text == "C"

    def test_empty_spans(self):
        lines = _group_spans_into_lines([])
        assert lines == [[]]

    def test_all_inline(self):
        spans = [_SpanSpec(text="A"), _SpanSpec(text="B"), _SpanSpec(text="C")]
        lines = _group_spans_into_lines(spans)
        assert len(lines) == 1
        assert len(lines[0]) == 3


class TestTspanIntegration:
    def _render(self, svg_str):
        from lxml import etree
        from pptx import Presentation
        from pptx.util import Inches
        from ppt_pro_max.renderer.svg_compiler._affine import Affine

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        el = etree.fromstring(svg_str)
        features = set()

        def resolve_color_fn(v, C, fb):
            return v if v.startswith("#") else None

        render_svg_text(
            el=el,
            tf=Affine(),
            to_inches_fn=lambda x, y: (x / 96, y / 96),
            slide=slide,
            C={},
            resolve_color_fn=resolve_color_fn,
            features=features,
        )
        return slide

    def test_tspan_two_lines_creates_two_paragraphs(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100">'
            '<tspan x="100" y="120">Line1</tspan>'
            '<tspan x="100" y="140">Line2</tspan></text>'
        )
        assert len(slide.shapes) == 1
        tb = slide.shapes[0]
        assert len(tb.text_frame.paragraphs) == 2

    def test_tspan_inline_same_paragraph(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100">'
            'Hello <tspan font-weight="bold">World</tspan></text>'
        )
        assert len(slide.shapes) == 1
        tb = slide.shapes[0]
        assert len(tb.text_frame.paragraphs) == 1
        assert len(tb.text_frame.paragraphs[0].runs) >= 2

    def test_tspan_bold_run(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100">'
            '<tspan font-weight="bold">Bold</tspan></text>'
        )
        tb = slide.shapes[0]
        bold_found = any(r.font.bold for r in tb.text_frame.paragraphs[0].runs if r.text.strip())
        assert bold_found

    def test_tspan_italic_run(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100">'
            '<tspan font-style="italic">Italic</tspan></text>'
        )
        tb = slide.shapes[0]
        italic_found = any(r.font.italic for r in tb.text_frame.paragraphs[0].runs if r.text.strip())
        assert italic_found

    def test_tspan_different_font_size(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100" font-size="14">'
            '<tspan font-size="24">Big</tspan></text>'
        )
        tb = slide.shapes[0]
        sizes = [r.font.size for r in tb.text_frame.paragraphs[0].runs if r.text.strip()]
        assert any(s and s.pt == 24.0 for s in sizes)

    def test_tspan_different_fill(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100" fill="#000000">'
            '<tspan fill="#FF0000">Red</tspan></text>'
        )
        tb = slide.shapes[0]
        colors = [r.font.color.rgb for r in tb.text_frame.paragraphs[0].runs if r.text.strip()]
        assert any(str(c) == "FF0000" for c in colors)

    def test_tspan_mixed_lines_and_inline(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100">'
            '<tspan x="100" y="120">Line1</tspan>'
            '<tspan x="100" y="140">Line2</tspan>'
            '<tspan font-weight="bold">BoldTail</tspan></text>'
        )
        assert len(slide.shapes) == 1
        tb = slide.shapes[0]
        assert len(tb.text_frame.paragraphs) == 2

    def test_empty_tspan_skipped(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100">'
            '<tspan x="100" y="120">  </tspan></text>'
        )
        assert len(slide.shapes) == 0

    def test_simple_text_no_tspan_still_works(self):
        slide = self._render(
            '<text xmlns="http://www.w3.org/2000/svg" x="100" y="100" font-size="16">Simple</text>'
        )
        assert len(slide.shapes) == 1
        tb = slide.shapes[0]
        assert len(tb.text_frame.paragraphs) == 1


class TestBaselineOffsetMode:
    """Baseline mode resolution — precise y-anchor for vertical positioning."""

    def test_auto_maps_to_middle(self):
        class FakeEl:
            def get(self, _, default=None):
                return "auto"
        assert _resolve_baseline_mode(FakeEl()) == "middle"

    def test_alphabetic_maps_to_baseline(self):
        class FakeEl:
            def get(self, _, default=None):
                return "alphabetic"
        assert _resolve_baseline_mode(FakeEl()) == "baseline"

    def test_hanging_maps_to_top(self):
        class FakeEl:
            def get(self, _, default=None):
                return "hanging"
        assert _resolve_baseline_mode(FakeEl()) == "top"

    def test_central_maps_to_middle(self):
        class FakeEl:
            def get(self, _, default=None):
                return "central"
        assert _resolve_baseline_mode(FakeEl()) == "middle"

    def test_unknown_falls_back_to_middle(self):
        class FakeEl:
            def get(self, _, default=None):
                return "xyz"
        assert _resolve_baseline_mode(FakeEl()) == "middle"

    def test_baseline_offset_map_size(self):
        assert len(_BASELINE_OFFSET) >= 8


class TestComputeTextTop:
    """Precise textbox top computation from SVG y + baseline."""

    METRICS = TextMetrics(
        width_inches=0.5,
        height_inches=0.3,
        ascent_ratio=0.75,
        descent_ratio=0.25,
    )

    def test_middle_centers_on_y(self):
        # top = iy - (h*0.75)/2 = 1.0 - 0.1125
        top = _compute_text_top(1.0, self.METRICS, MSO_ANCHOR.MIDDLE, "middle")
        assert abs(top - (1.0 - 0.3 * 0.75 / 2)) < 1e-6

    def test_alphabetic_baseline_y(self):
        # top = iy - h*0.75 = 1.0 - 0.225
        top = _compute_text_top(1.0, self.METRICS, MSO_ANCHOR.BOTTOM, "baseline")
        assert abs(top - (1.0 - 0.3 * 0.75)) < 1e-6

    def test_top_edge_y(self):
        # top = iy
        top = _compute_text_top(1.0, self.METRICS, MSO_ANCHOR.TOP, "top")
        assert abs(top - 1.0) < 1e-6

    def test_descent_y(self):
        # top = iy - h*0.75 - h*0.25 = iy - h
        top = _compute_text_top(1.0, self.METRICS, MSO_ANCHOR.BOTTOM, "descent")
        assert abs(top - (1.0 - 0.3)) < 1e-6

    def test_metrics_consistency(self):
        # ascent + descent should equal height ratio ~1.0
        assert abs(self.METRICS.ascent_ratio + self.METRICS.descent_ratio - 1.0) < 1e-6


class TestHasCJK:
    def test_cjk_hanzi(self):
        assert _has_cjk("战略") is True

    def test_cjk_fullwidth(self):
        assert _has_cjk("Ａ") is True

    def test_ascii_only(self):
        assert _has_cjk("Hello") is False

    def test_mixed(self):
        assert _has_cjk("Hello世界") is True

    def test_empty(self):
        assert _has_cjk("") is False

    def test_cjk_punctuation(self):
        assert _has_cjk("、") is True


class TestMeasureTextCJK:
    def test_cjk_wider_than_ascii(self):
        m_cjk = _measure_text("战略愿景", 14.0, "Arial", 8.0)
        m_ascii = _measure_text("ABCD", 14.0, "Arial", 8.0)
        assert m_cjk.width_inches > m_ascii.width_inches

    def test_cjk_not_clamped_to_minimum(self):
        m = _measure_text("核心竞争力构建", 14.0, "Arial", 8.0)
        assert m.width_inches > 0.6

    def test_cjk_height_reasonable(self):
        m = _measure_text("战略愿景", 14.0, "Arial", 8.0)
        assert m.height_inches >= 0.3

    def test_ascii_uses_pil_or_fallback(self):
        m = _measure_text("Hello", 14.0, "Arial", 8.0)
        assert m.width_inches >= 0.5

    def test_cjk_more_chars_wider(self):
        m4 = _measure_text("战略", 14.0, "Arial", 8.0)
        m6 = _measure_text("战略愿景", 14.0, "Arial", 8.0)
        assert m6.width_inches > m4.width_inches
