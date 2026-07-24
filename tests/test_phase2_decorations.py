"""TDD tests for Phase 2: decoration classification, BrandColorContext, hardcoded color replacement, goal/layout fixes."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _create_content_with_deco_pptx(path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    cover = prs.slides.add_slide(blank)
    tb = cover.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = tb.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "Cover"
    r.font.size = Pt(36)

    content = prs.slides.add_slide(blank)
    tb2 = content.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf2 = tb2.text_frame
    r2 = tf2.paragraphs[0].add_run()
    r2.text = "Content Title"
    r2.font.size = Pt(28)
    r2.font.bold = True
    body = content.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    btf = body.text_frame
    br = btf.paragraphs[0].add_run()
    br.text = "Bullet A"
    br.font.size = Pt(14)

    deco_bar = content.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.5), Inches(7.5)
    )
    deco_bar.fill.solid()
    deco_bar.fill.fore_color.rgb = RGBColor(0x2E, 0x65, 0x04)
    deco_bar.line.fill.background()

    footer_bar = content.shapes.add_shape(
        1, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(0x2E, 0x65, 0x04)
    footer_bar.line.fill.background()

    back = prs.slides.add_slide(blank)
    tb3 = back.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf3 = tb3.text_frame
    r3 = tf3.paragraphs[0].add_run()
    r3.text = "Thank You"
    r3.font.size = Pt(36)

    prs.save(path)


class TestDecorationClassification:
    """DecorationElement 4-level classification system."""

    def test_classify_sidebar_bar(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        result = extractor._classify_decoration_by_bounds(0.0, 0.0, 0.5, 7.5, has_fill=True)
        assert result[0] == 0
        assert result[1] == "sidebar"

    def test_classify_footer_bar(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        result = extractor._classify_decoration_by_bounds(0.0, 6.8, 13.333, 0.7, has_fill=True)
        assert result[0] == 0
        assert result[1] == "footer"

    def test_classify_corner_decoration(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        result = extractor._classify_decoration_by_bounds(11.5, 0.0, 1.5, 1.5, has_fill=True)
        assert result[0] == 1
        assert result[1] == "corner"

    def test_classify_divider_line(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        result = extractor._classify_decoration_by_bounds(1.0, 3.0, 8.0, 0.05, has_fill=True)
        assert result[0] == 1
        assert result[1] == "divider"

    def test_classify_title_decoration(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        result = extractor._classify_decoration_by_bounds(1.0, 1.5, 6.0, 0.08, has_fill=True)
        assert result[0] == 0
        assert result[1] == "title_decoration"

    def test_classify_large_content_image(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        result = extractor._classify_decoration_by_bounds(7.0, 1.5, 5.0, 5.0, has_fill=False, is_picture=True)
        assert result[0] == 3
        assert result[1] == "large_image"

    def test_classify_small_icon(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        result = extractor._classify_decoration_by_bounds(9.0, 2.0, 0.8, 0.8, has_fill=False, is_picture=True)
        assert result[0] == 1
        assert result[1] == "icon"

    def test_selective_clear_preserves_sidebar(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, SlideDNA, TextZone, PagePlan,
        )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(blank)

        title_tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_r = title_tb.text_frame.paragraphs[0].add_run()
        title_r.text = "Title"
        title_r.font.size = Pt(28)

        sidebar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = RGBColor(0x2E, 0x65, 0x04)
        sidebar.line.fill.background()

        slide_dna = SlideDNA(
            slide_index=0, page_type="content", slide_xml=b"", rels_xml=b"",
            text_zones=[
                TextZone(zone_id="t_0", role="title", text="Title", shape_name=title_tb.name),
            ],
        )
        plan = PagePlan(page_type="content", title="Title", bullets=["A", "B", "C"], layout="cards")

        extractor = DesignDNAExtractor()
        extractor._selective_clear_decorations(slide, slide_dna, plan, keep_levels={0, 1})

        sidebar_still_exists = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.width < Inches(1):
                sidebar_still_exists = True
        assert sidebar_still_exists, "Sidebar decoration should be preserved at keep_levels {0,1}"


class TestBrandColorContext:
    """BrandColorContext: semantic color view from DesignDNA."""

    def test_from_dna_basic(self):
        from ppt_pro_max.enterprise.brand_color_context import BrandColorContext
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNA

        dna = DesignDNA(
            source_path="none",
            color_palette={"accent1": "#3366CC", "dk1": "#333333", "lt1": "#FFFFFF"},
            actual_colors={"#3366CC": 10, "#333333": 5, "#F0F0F0": 3},
        )
        ctx = BrandColorContext.from_dna(dna)
        assert ctx.primary is not None
        assert ctx.foreground is not None
        assert ctx.background is not None

    def test_from_dna_dark_template(self):
        from ppt_pro_max.enterprise.brand_color_context import BrandColorContext
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNA

        dna = DesignDNA(
            source_path="none",
            color_palette={"accent1": "#FF6600", "dk1": "#FFFFFF", "lt1": "#1A1A2E"},
            actual_colors={"#1A1A2E": 20, "#FF6600": 8, "#FFFFFF": 5},
        )
        ctx = BrandColorContext.from_dna(dna)
        assert ctx.primary is not None
        assert ctx.muted is not None

    def test_from_brand_spec(self):
        from ppt_pro_max.enterprise.brand_color_context import BrandColorContext
        from ppt_pro_max.enterprise.brand_spec import BrandSpec

        brand = BrandSpec(source="brand_json", colors={"accent1": "#E8A838", "tx1": "#2E6504"})
        ctx = BrandColorContext.from_brand_spec(brand)
        assert ctx.accent1 == "#E8A838"
        assert ctx.foreground == "#2E6504"


class TestHardcodedColorReplacement:
    """Verify hardcoded colors are replaced with BrandColorContext."""

    def test_post_save_inject_uses_brand_colors(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        import inspect
        source = inspect.getsource(DesignDNAExtractor._post_save_inject_groups)
        assert '"#84AF7D"' not in source or "brand_ctx" in source or "body_color" in source, \
            "_post_save_inject_groups should not hardcode #84AF7D"

    def test_try_component_render_uses_brand_colors(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        import inspect
        source = inspect.getsource(DesignDNAExtractor._try_component_render)
        assert '"#6096E6"' not in source or "brand_ctx" in source or "accent_color" in source, \
            "_try_component_render should not hardcode #6096E6"


class TestD1DecorationPreservation:
    """D1: decoration preservation condition too narrow — extend to more layouts."""

    def test_cards_layout_preserves_shell(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, SlideDNA, TextZone, PagePlan,
        )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(blank)

        title_tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        r = title_tb.text_frame.paragraphs[0].add_run()
        r.text = "Title"
        r.font.size = Pt(28)

        sidebar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = RGBColor(0x2E, 0x65, 0x04)
        sidebar.line.fill.background()

        body_tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        br = body_tb.text_frame.paragraphs[0].add_run()
        br.text = "Old body"
        br.font.size = Pt(14)

        slide_dna = SlideDNA(
            slide_index=0, page_type="content", slide_xml=b"", rels_xml=b"",
            text_zones=[
                TextZone(zone_id="t_0", role="title", text="Title", shape_name=title_tb.name,
                         bounds=(1.0, 0.5, 8.0, 1.0)),
                TextZone(zone_id="b_0", role="body", text="Old body", shape_name=body_tb.name,
                         bounds=(1.0, 2.0, 8.0, 4.0)),
            ],
        )
        plan = PagePlan(page_type="content", title="Title", layout="cards",
                        cards=[{"title": "C1"}, {"title": "C2"}, {"title": "C3"}])

        extractor = DesignDNAExtractor()
        needs_rerender = extractor.needs_content_rerender(slide_dna, plan)
        assert needs_rerender, "cards layout should trigger content rerender"

    def test_process_layout_preserves_shell(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, SlideDNA, TextZone, PagePlan,
        )

        slide_dna = SlideDNA(
            slide_index=0, page_type="content", slide_xml=b"", rels_xml=b"",
            text_zones=[
                TextZone(zone_id="t_0", role="title", text="Title"),
                TextZone(zone_id="b_0", role="body", text="Step 1"),
            ],
        )
        plan = PagePlan(page_type="content", title="Title", layout="process",
                        bullets=["Step 1", "Step 2", "Step 3"])

        extractor = DesignDNAExtractor()
        assert extractor.needs_content_rerender(slide_dna, plan)


class TestD12GoalInference:
    """D12: architecture should not map to 'data'."""

    def test_architecture_not_data(self):
        from ppt_pro_max.enterprise.content_parser import _infer_goal
        goal = _infer_goal("Architecture Overview", 2, 5)
        assert goal != "data", f"'Architecture Overview' should not infer as 'data', got '{goal}'"

    def test_architecture_maps_to_overview_or_content(self):
        from ppt_pro_max.enterprise.content_parser import _infer_goal
        goal = _infer_goal("System Architecture", 2, 5)
        assert goal in ("overview", "content", "features"), \
            f"'System Architecture' should map to overview/content/features, got '{goal}'"


class TestD13AutoLayoutEnglish:
    """D13: _auto_layout should recognize English keywords."""

    def test_steps_keyword(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        page = {"bullets": ["Step 1: Setup", "Step 2: Build", "Step 3: Deploy"]}
        layout = extractor._auto_layout(page)
        assert layout == "process", f"'steps' bullets should infer 'process', got '{layout}'"

    def test_phases_keyword(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        page = {"bullets": ["Phase 1: Discovery", "Phase 2: Design", "Phase 3: Delivery"]}
        layout = extractor._auto_layout(page)
        assert layout == "process", f"'phases' bullets should infer 'process', got '{layout}'"

    def test_org_chart_keyword(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        page = {"bullets": ["org chart: CEO", "CTO", "VP Engineering", "VP Sales"]}
        layout = extractor._auto_layout(page)
        assert layout == "hierarchy", f"'org chart' bullets should infer 'hierarchy', got '{layout}'"

    def test_milestone_keyword(self):
        from ppt_pro_max.enterprise.design_dna_extractor import DesignDNAExtractor
        extractor = DesignDNAExtractor()
        page = {"bullets": ["Milestone 1: MVP", "Milestone 2: Beta", "Milestone 3: Launch"]}
        layout = extractor._auto_layout(page)
        assert layout == "timeline", f"'milestone' bullets should infer 'timeline', got '{layout}'"


class TestD14CompositeTitleWestern:
    """D14: _is_composite_title should support Western multi-textbox titles."""

    def test_western_composite_title(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, TextZone,
        )
        zones = [
            TextZone(zone_id="t_0", role="title", text="Our",
                     bounds=(2.0, 1.0, 1.5, 1.0), font_size_pt=36),
            TextZone(zone_id="t_1", role="title", text="Mission",
                     bounds=(3.6, 1.0, 2.0, 1.0), font_size_pt=36),
        ]
        extractor = DesignDNAExtractor()
        assert extractor._is_composite_title(zones), \
            "Western composite title 'Our / Mission' should be detected"

    def test_chinese_composite_title_still_works(self):
        from ppt_pro_max.enterprise.design_dna_extractor import (
            DesignDNAExtractor, TextZone,
        )
        zones = [
            TextZone(zone_id="t_0", role="title", text="乡",
                     bounds=(2.0, 1.0, 1.5, 2.0), font_size_pt=72),
            TextZone(zone_id="t_1", role="title", text="村",
                     bounds=(3.6, 1.0, 1.5, 2.0), font_size_pt=72),
        ]
        extractor = DesignDNAExtractor()
        assert extractor._is_composite_title(zones), \
            "Chinese composite title '乡/村' should still be detected"
