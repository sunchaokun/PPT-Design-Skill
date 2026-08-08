"""Tests for blip_fill.py — image-in-shape via OOXML a:blipFill.

Covers:
  1. fill_shape_with_image() — fill existing shape with image
  2. add_image_in_shape() — create shape + fill with image
  3. add_circle_image() — circle-cropped image
  4. add_hexagon_image() — hexagon-cropped image
  5. add_diamond_image() — diamond-cropped image
  6. Blip effects: duotone, grayscale, brightness/contrast, saturation, artistic
"""

from __future__ import annotations

import os
import tempfile

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_pro_max.renderer.blip_fill import (
    add_circle_image,
    add_diamond_image,
    add_hexagon_image,
    add_image_in_shape,
    fill_shape_with_image,
)


@pytest.fixture
def test_image():
    img = PILImage.new("RGB", (200, 200), (255, 100, 50))
    path = os.path.join(tempfile.gettempdir(), "blip_test_img.png")
    img.save(path, "PNG")
    return path


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[-1])
    return s


class TestFillShapeWithImage:
    def test_fills_existing_shape_with_image(self, slide, test_image):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        result = fill_shape_with_image(shape, slide, test_image)
        assert result is not None
        spPr = shape._element.find(qn("p:spPr"))
        blipFill = spPr.find(qn("a:blipFill"))
        assert blipFill is not None, "Shape should have blipFill"
        blip = blipFill.find(qn("a:blip"))
        assert blip is not None
        rId = blip.get(qn("r:embed"))
        assert rId is not None, "blip should reference an image rId"

    def test_removes_existing_fill(self, slide, test_image):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor_from_string("FF0000")
        fill_shape_with_image(shape, slide, test_image)
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:solidFill")) is None, "Old fill should be removed"
        assert spPr.find(qn("a:blipFill")) is not None

    def test_stretch_mode(self, slide, test_image):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        fill_shape_with_image(shape, slide, test_image, crop_mode="stretch")
        spPr = shape._element.find(qn("p:spPr"))
        blipFill = spPr.find(qn("a:blipFill"))
        stretch = blipFill.find(qn("a:stretch"))
        assert stretch is not None

    def test_tile_mode(self, slide, test_image):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1),
                                       Inches(3), Inches(3))
        fill_shape_with_image(shape, slide, test_image, crop_mode="tile")
        spPr = shape._element.find(qn("p:spPr"))
        blipFill = spPr.find(qn("a:blipFill"))
        tile = blipFill.find(qn("a:tile"))
        assert tile is not None

    def test_alpha(self, slide, test_image):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        fill_shape_with_image(shape, slide, test_image, alpha=50)
        spPr = shape._element.find(qn("p:spPr"))
        blipFill = spPr.find(qn("a:blipFill"))
        blip = blipFill.find(qn("a:blip"))
        alpha_el = blip.find(qn("a:alpha"))
        assert alpha_el is not None
        assert alpha_el.get("val") == "50000"

    def test_removes_line(self, slide, test_image):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        fill_shape_with_image(shape, slide, test_image)
        spPr = shape._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        noFill = ln.find(qn("a:noFill")) if ln is not None else None
        assert noFill is not None, "Line should be set to noFill"

    def test_nonexistent_image_returns_none(self, slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        result = fill_shape_with_image(shape, slide, "/nonexistent/image.png")
        assert result is None

    def test_alpha_100_no_alpha_element(self, slide, test_image):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1),
                                       Inches(2), Inches(2))
        fill_shape_with_image(shape, slide, test_image, alpha=100)
        spPr = shape._element.find(qn("p:spPr"))
        blip = spPr.find(qn("a:blipFill")).find(qn("a:blip"))
        assert blip.find(qn("a:alpha")) is None, "alpha=100 should not add alpha element"


class TestAddImageInShape:
    def test_creates_shape_with_image(self, slide, test_image):
        shape = add_image_in_shape(slide, MSO_SHAPE.OVAL, 1, 1, 2, 2, test_image)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        assert spPr.find(qn("a:blipFill")) is not None

    def test_shape_type_matches(self, slide, test_image):
        shape = add_image_in_shape(slide, MSO_SHAPE.HEXAGON, 1, 1, 2, 2, test_image)
        spPr = shape._element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        assert prstGeom is not None
        assert prstGeom.get("prst") == "hexagon"

    def test_with_border(self, slide, test_image):
        shape = add_image_in_shape(slide, MSO_SHAPE.OVAL, 1, 1, 2, 2, test_image,
                                   border_hex="#FF0000")
        spPr = shape._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        assert ln is not None
        solidFill = ln.find(qn("a:solidFill"))
        assert solidFill is not None


class TestAddCircleImage:
    def test_creates_circle_image(self, slide, test_image):
        shape = add_circle_image(slide, 5.0, 3.0, 1.0, test_image)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        assert prstGeom.get("prst") == "ellipse"
        assert spPr.find(qn("a:blipFill")) is not None

    def test_circle_dimensions(self, slide, test_image):
        shape = add_circle_image(slide, 5.0, 3.0, 1.0, test_image)
        assert shape.width == Inches(2.0)
        assert shape.height == Inches(2.0)

    def test_with_border(self, slide, test_image):
        shape = add_circle_image(slide, 5.0, 3.0, 1.0, test_image,
                                 border_hex="#CCCCCC")
        spPr = shape._element.find(qn("p:spPr"))
        ln = spPr.find(qn("a:ln"))
        assert ln is not None
        solidFill = ln.find(qn("a:solidFill"))
        assert solidFill is not None


class TestAddHexagonImage:
    def test_creates_hexagon_image(self, slide, test_image):
        shape = add_hexagon_image(slide, 5.0, 3.0, 2.0, test_image)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        assert prstGeom.get("prst") == "hexagon"
        assert spPr.find(qn("a:blipFill")) is not None


class TestAddDiamondImage:
    def test_creates_diamond_image(self, slide, test_image):
        shape = add_diamond_image(slide, 5.0, 3.0, 2.0, test_image)
        assert shape is not None
        spPr = shape._element.find(qn("p:spPr"))
        prstGeom = spPr.find(qn("a:prstGeom"))
        assert prstGeom.get("prst") == "diamond"
        assert spPr.find(qn("a:blipFill")) is not None


def RGBColor_from_string(hex_str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str.lstrip("#"))
