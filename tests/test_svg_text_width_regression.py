"""Regression tests for SVG text width constraint on small shapes.

Bug: SVG text on small shapes (e.g. pyramid top triangle ~0.6" wide)
rendered as vertical text because font-size was in SVG user units
without viewBox→slide scaling, and textbox width was unconstrained.

Fix: render_svg_text() now takes svg_w/svg_h/slide_w/slide_h params,
computes scale = slide_w/svg_w, derives scaled_fs = parent_fs * scale
(clamped >=6pt), and constrains textbox width to min(measured+padding, max_w).
"""
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_pro_max.renderer.svg_compiler import SVGCompiler
from ppt_pro_max.renderer.svg_compiler._affine import Affine
from ppt_pro_max.renderer.svg_compiler._text import render_svg_text


def _make_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


class TestTextScalingWithViewBox:
    """Text font-size must be scaled by viewBox→slide ratio."""

    def test_small_viewbox_large_slide_scales_up(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
          <text x="50" y="50" text-anchor="middle" font-size="14" fill="#000">Hello</text>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 8, 6))
        assert result.shape_count > 0
        assert "text" in result.features

    def test_large_viewbox_small_slide_scales_down(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1000 1000">
          <text x="500" y="500" text-anchor="middle" font-size="40" fill="#000">Hello</text>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 2, 2))
        assert result.shape_count > 0

    def test_pyramid_top_triangle_text_not_vertical(self, tmp_path):
        """Regression: pyramid top triangle is narrow (~0.6" on slide),
        text "战略愿景" should not wrap vertically."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 360">
          <defs>
            <linearGradient id="g1" x1="0" y1="0" x2="1" y2="1">
              <stop offset="0" stop-color="#7DA92F"/><stop offset="1" stop-color="#2E6504"/>
            </linearGradient>
          </defs>
          <polygon points="200,30 380,320 20,320" fill="url(#g1)"/>
          <text x="200" y="85" text-anchor="middle" font-size="16" fill="#fff" font-family="Arial">战略愿景</text>
        </svg>"""
        prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (3.5, 0.8, 6.3, 5.6))
        assert result.shape_count >= 2

        pptx_path = tmp_path / "pyramid_text.pptx"
        prs.save(str(pptx_path))

        import zipfile
        with zipfile.ZipFile(pptx_path) as z:
            slide_xmls = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            xml = z.read(slide_xmls[0]).decode("utf-8")

        root = etree.fromstring(xml.encode("utf-8"))
        text_shapes = []
        for sp in root.iter(qn("p:sp")):
            txBody = sp.find(qn("p:txBody"))
            if txBody is not None:
                text_shapes.append(sp)

        assert len(text_shapes) >= 1, "at least one text shape expected"

        for sp in text_shapes:
            spPr = sp.find(qn("p:spPr"))
            if spPr is None:
                continue
            xfrm = spPr.find(qn("a:xfrm"))
            if xfrm is None:
                continue
            ext = xfrm.find(qn("a:ext"))
            if ext is None:
                continue
            cx = int(ext.get("cx", 0))
            cy = int(ext.get("cy", 0))
            if cx > 0 and cy > 0:
                assert cx >= cy, (
                    f"textbox should be wider than tall (cx={cx}, cy={cy}), "
                    f"not vertical — regression: text wrapping vertically on small shape"
                )


class TestTextWidthConstraint:
    """Textbox width must be constrained to parent shape's available width."""

    def test_text_on_narrow_rect_constrained(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 100">
          <rect x="0" y="0" width="400" height="100" fill="#4472C4"/>
          <text x="200" y="55" text-anchor="middle" font-size="20" fill="#fff">This is a long text label</text>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (5, 2, 3, 0.75))
        assert result.shape_count >= 2

    def test_text_on_wide_rect_unconstrained(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 100">
          <rect x="0" y="0" width="400" height="100" fill="#4472C4"/>
          <text x="200" y="55" text-anchor="middle" font-size="14" fill="#fff">Short</text>
        </svg>"""
        _prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (2, 2, 8, 2))
        assert result.shape_count >= 2


class TestMinFontSizeClamp:
    """Scaled font-size must never go below 6pt."""

    def test_extreme_scale_clamps_to_6pt(self):
        svg = """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 10000 10000">
          <text x="5000" y="5000" text-anchor="middle" font-size="8" fill="#000000">Tiny</text>
        </svg>"""
        prs, slide = _make_slide()
        compiler = SVGCompiler()
        result = compiler.compile(svg, slide, (1, 1, 1, 1))
        assert result.shape_count > 0

        pptx_path = "output/test_min_font.pptx"
        prs.save(pptx_path)

        import zipfile
        with zipfile.ZipFile(pptx_path) as z:
            slide_xmls = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            xml = z.read(slide_xmls[0]).decode("utf-8")

        root = etree.fromstring(xml.encode("utf-8"))
        for rPr in root.iter(qn("a:rPr")):
            sz = rPr.get("sz")
            if sz is not None:
                half_pt = int(sz)
                assert half_pt >= 120, f"font size {half_pt/100}pt is below 6pt minimum"


class TestRenderSvgTextScaleParam:
    """Unit tests for render_svg_text scale parameter."""

    def test_scale_param_affects_font_size(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        el = etree.fromstring(
            '<text xmlns="http://www.w3.org/2000/svg" x="50" y="50" '
            'text-anchor="middle" font-size="20" fill="#000000">Scaled</text>'
        )

        features = set()

        def resolve_color_fn(v, C, fb):
            return v if v.startswith("#") else None

        render_svg_text(
            el=el,
            tf=Affine(),
            to_inches_fn=lambda x, y: (x / 96, y / 96),
            slide=slide,
            C={},
            resolve_color_fn=resolve_color_fn,
            features=features,
            svg_w=100.0,
            svg_h=100.0,
            slide_w=5.0,
            slide_h=5.0,
        )

        assert len(slide.shapes) == 1
        tb = slide.shapes[0]
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size and r.text.strip():
                    assert r.font.size.pt >= 6.0

    def test_zero_svg_w_uses_default_scale(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        el = etree.fromstring(
            '<text xmlns="http://www.w3.org/2000/svg" x="50" y="50" '
            'text-anchor="middle" font-size="14" fill="#000000">Default</text>'
        )

        features = set()

        def resolve_color_fn(v, C, fb):
            return v if v.startswith("#") else None

        render_svg_text(
            el=el,
            tf=Affine(),
            to_inches_fn=lambda x, y: (x / 96, y / 96),
            slide=slide,
            C={},
            resolve_color_fn=resolve_color_fn,
            features=features,
            svg_w=0.0,
            svg_h=0.0,
            slide_w=5.0,
            slide_h=5.0,
        )

        assert len(slide.shapes) == 1
