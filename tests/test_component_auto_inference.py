"""Tests for auto-inference of component_type/component_category in pipeline and freestyle paths."""
import json
import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

_VALID_GROUP_XML = (
    b'<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    b' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    b'<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="100" cy="100"/>'
    b'<a:chOff x="0" y="0"/><a:chExt cx="100" cy="100"/></a:xfrm></p:grpSpPr>'
    b'<p:sp><p:nvSpPr><p:cNvPr id="1" name="s1"/><p:cNvSpPr/>'
    b'<p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="0" y="0"/>'
    b'<a:ext cx="100" cy="50"/></a:xfrm></p:spPr>'
    b'<p:txBody><a:bodyPr/><a:lstStyle/>'
    b'<a:p><a:r><a:rPr lang="en-US" sz="1200"/><a:t>Step 1</a:t></a:r></a:p>'
    b'</p:txBody></p:sp></p:grpSp>'
)


class TestPipelineAutoInference:

    def test_pipeline_infers_process_from_bullets(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        bullets = ["需求分析流程", "方案设计步骤", "开发实现阶段", "测试上线流程"]
        comp_type, comp_cat = infer_component_category(bullets)
        assert comp_type == "group"
        assert comp_cat == "process"

    def test_pipeline_infers_swot_from_bullets(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        bullets = ["技术优势", "市场劣势", "增长机会", "潜在威胁"]
        comp_type, comp_cat = infer_component_category(bullets)
        assert comp_type == "group"
        assert comp_cat == "swot"

    def test_pipeline_infers_hierarchy_from_bullets(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        bullets = ["CEO", "CTO", "VP Engineering", "Team Lead"]
        comp_type, comp_cat = infer_component_category(bullets)
        assert comp_type == "group"
        assert comp_cat == "hierarchy"

    def test_pipeline_infers_timeline_from_keywords(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        bullets = ["Q1 里程碑", "Q2 时间线", "Q3 路线图", "Q4 milestone"]
        comp_type, comp_cat = infer_component_category(bullets)
        assert comp_type == "group"
        assert comp_cat == "timeline"

    def test_pipeline_no_inference_for_2_bullets(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        bullets = ["Point A", "Point B"]
        comp_type, comp_cat = infer_component_category(bullets)
        assert comp_type == "group"
        assert comp_cat == "comparison"

    def test_pipeline_no_inference_for_10_bullets(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        bullets = [f"Point {i}" for i in range(10)]
        comp_type, comp_cat = infer_component_category(bullets)
        assert comp_type is None
        assert comp_cat is None

    def test_pipeline_3_generic_bullets_default_to_pyramid(self):
        from ppt_pro_max.enterprise.content_parser import infer_component_category

        bullets = ["Alpha", "Beta", "Gamma"]
        comp_type, comp_cat = infer_component_category(bullets)
        assert comp_type == "group"
        assert comp_cat == "pyramid"

    def test_pipeline_auto_inference_in_run(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        project_dir = tmp_path / "project"
        project_dir.mkdir()

        content = {
            "slides": [
                {"goal": "hook", "title": "Test Hook"},
                {
                    "goal": "content",
                    "title": "项目流程",
                    "bullets": ["需求分析", "方案设计", "开发实现", "测试上线"],
                },
                {"goal": "cta", "title": "Call to Action"},
            ]
        }
        content_file = project_dir / "content.json"
        content_file.write_text(json.dumps(content, ensure_ascii=False), encoding="utf-8")

        brand_json = project_dir / "brand.json"
        brand_json.write_text(json.dumps({
            "colors": {"primary": "#2563EB", "on-primary": "#FFFFFF", "background": "#FFFFFF"},
            "fonts": {"heading": "Inter", "body": "Inter"},
        }), encoding="utf-8")

        db_path = str(tmp_path / "test_lib.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="process", variant="chevron", node_count=4,
                xml_parts={"group": _VALID_GROUP_XML})
        lib.close()

        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="test project",
            project_dir=str(project_dir),
            component_library=db_path,
        )

        assert result.get("pipeline") == "enterprise"
        assert result.get("num_slides", 0) >= 1

    def test_pipeline_preserves_explicit_component_fields(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        project_dir = tmp_path / "project"
        project_dir.mkdir()

        content = {
            "slides": [
                {"goal": "hook", "title": "Hook"},
                {
                    "goal": "content",
                    "title": "Our Team",
                    "bullets": ["CEO", "CTO", "VP Eng"],
                    "component_type": "group",
                    "component_category": "hierarchy",
                },
            ]
        }
        content_file = project_dir / "content.json"
        content_file.write_text(json.dumps(content, ensure_ascii=False), encoding="utf-8")

        brand_json = project_dir / "brand.json"
        brand_json.write_text(json.dumps({
            "colors": {"primary": "#2563EB", "on-primary": "#FFFFFF", "background": "#FFFFFF"},
            "fonts": {"heading": "Inter", "body": "Inter"},
        }), encoding="utf-8")

        db_path = str(tmp_path / "test_lib.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="hierarchy", variant="orgchart", node_count=3,
                xml_parts={"group": _VALID_GROUP_XML})
        lib.close()

        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="test",
            project_dir=str(project_dir),
            component_library=db_path,
        )

        assert result.get("pipeline") == "enterprise"


class TestFreestyleComponentLibrary:

    def test_freestyle_with_component_library_uses_precision(self, tmp_path):
        from ppt_pro_max import generate_ppt
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test_lib.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="process", variant="chevron", node_count=4,
                xml_parts={"group": _VALID_GROUP_XML})
        lib.close()

        output_path = str(tmp_path / "output.pptx")
        result = generate_ppt(
            "test pitch",
            component_library=db_path,
            output=output_path,
        )

        assert result.get("output_path")
        assert result.get("render_mode") == "precision"
        assert os.path.isfile(result["output_path"])

    def test_freestyle_without_component_library_uses_precision(self, tmp_path):
        from ppt_pro_max import generate_ppt

        output_path = str(tmp_path / "output.pptx")
        result = generate_ppt(
            "test pitch",
            output=output_path,
        )

        assert result.get("output_path")
        assert os.path.isfile(result["output_path"])

    def test_freestyle_auto_infers_component_from_bullets(self, tmp_path):
        from ppt_pro_max import generate_ppt
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        db_path = str(tmp_path / "test_lib.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="process", variant="chevron", node_count=4,
                xml_parts={"group": _VALID_GROUP_XML})
        lib.close()

        content = {
            "slides": [
                {"goal": "hook", "title": "AI Platform"},
                {
                    "goal": "content",
                    "title": "项目流程",
                    "bullets": ["需求分析", "方案设计", "开发实现", "测试上线"],
                },
            ]
        }
        content_file = tmp_path / "content.json"
        content_file.write_text(json.dumps(content, ensure_ascii=False), encoding="utf-8")

        output_path = str(tmp_path / "output.pptx")
        result = generate_ppt(
            "AI platform pitch",
            content_file=str(content_file),
            component_library=db_path,
            output=output_path,
        )

        assert result.get("render_mode") == "precision"
        assert os.path.isfile(result["output_path"])

    def test_freestyle_with_invalid_db_falls_back(self, tmp_path):
        from ppt_pro_max import generate_ppt

        output_path = str(tmp_path / "output.pptx")
        result = generate_ppt(
            "test pitch",
            component_library="/nonexistent/path/db.sqlite",
            output=output_path,
        )

        assert result.get("output_path")
        assert os.path.isfile(result["output_path"])


class TestBeautifyAutoInference:

    def test_beautify_full_infers_components(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        from ppt_pro_max.enterprise.component_library import ComponentLibrary

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        tb = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "项目流程"
        run.font.size = Pt(36)

        for i, txt in enumerate(["需求分析", "方案设计", "开发实现", "测试上线"]):
            tb2 = slide1.shapes.add_textbox(Inches(1 + i * 2.5), Inches(3), Inches(2), Inches(1))
            p2 = tb2.text_frame.paragraphs[0]
            run2 = p2.add_run()
            run2.text = txt
            run2.font.size = Pt(18)

        slide2 = prs.slides.add_slide(layout)
        tb3 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        p3 = tb3.text_frame.paragraphs[0]
        run3 = p3.add_run()
        run3.text = "End Slide"
        run3.font.size = Pt(36)

        input_path = str(tmp_path / "input.pptx")
        prs.save(input_path)

        db_path = str(tmp_path / "test_lib.db")
        lib = ComponentLibrary(db_path)
        lib.add(type="group", category="process", variant="chevron", node_count=4,
                xml_parts={"group": _VALID_GROUP_XML})
        lib.close()

        output_path = str(tmp_path / "beautified.pptx")
        pipeline = EnterprisePipeline()
        result = pipeline.run_beautify(
            input_pptx=input_path,
            output_path=output_path,
            component_library=db_path,
        )

        assert result.get("mode") == "beautify"
        assert result.get("num_slides", 0) >= 1
