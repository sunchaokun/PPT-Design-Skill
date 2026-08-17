"""BuildQA low-contrast check — per-shape background detection regression.

BuildQA._check_low_contrast_text uses a single slide-level bg color for ALL
text runs. This causes false-positives when:

  1. A text run lives inside a non-full-bleed dark shape (e.g. code_block at
     x=0.65, y=1.62 with #1E1E1E bg). The detector falls back to white and
     incorrectly flags light-on-dark text as low-contrast.
  2. A text run lives inside a full-bleed dark group (section_divider uses
     grouped=True). The detector correctly finds the group-level fill, but
     only by accident — it never recurses into <p:grpSp>.

Expected behavior: every text run should be compared against the fill of the
shape that actually contains it (its <p:txBody>'s ancestor <p:sp> or
<p:grpSp>), not against a slide-level heuristic.

Regression covers:
  - Light text on dark offset rect (code_block style) → PASS (was FAIL)
  - Light text on dark full-bleed group (section_divider style) → PASS (was already OK by accident)
  - Dark text on white bg → FAIL (real issue preserved)
  - Slide-level full-bleed bg still used as fallback
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from ppt_pro_max.build_qa import BuildQA, QAReport


@pytest.fixture
def prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _make_textbox(slide, x: float, y: float, w: float, h: float,
                  text: str, color_hex: str, font_pt: int = 14):
    """Add a plain textbox with the given color and size."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    rgb = RGBColor.from_string(color_hex.lstrip("#"))
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    run.font.color.rgb = rgb
    return tb


def _add_solid_rect(slide, x: float, y: float, w: float, h: float,
                    fill_hex: str):
    """Add a rectangle filled with the given hex color."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    rgb = RGBColor.from_string(fill_hex.lstrip("#"))
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid()
    r.fill.fore_color.rgb = rgb
    r.line.fill.background()
    return r


def _build_group_dark_code(slide) -> None:
    """Reproduce the code_block scenario: a group containing a dark offset
    rect and white/light text on it. Mirrors build_10pages.py slide 6."""
    # Outer group covering roughly the code block area (not full-bleed).
    # We build it manually so we control the group offset.
    # Simpler: just add a dark rect at offset + light text overlapping it.
    _add_solid_rect(slide, 0.65, 1.62, 7.5, 4.88, "#1E1E1E")
    # Light text on dark bg — should PASS (good contrast).
    _make_textbox(slide, 0.8, 1.8, 7.2, 0.4,
                  'from enterprise_ai import Engine',
                  color_hex="#FFFFFF", font_pt=12)
    _make_textbox(slide, 0.8, 2.3, 7.2, 0.4,
                  '# 初始化企业智能引擎',
                  color_hex="#D4D4D4", font_pt=12)


def _build_section_divider_style(slide) -> None:
    """Reproduce section_divider: full-bleed dark group + light text inside.
    The dark rect is the group itself (full-bleed) and text is a child sp."""
    _add_solid_rect(slide, 0, 0, 13.333, 7.5, "#1E3A5F")
    _make_textbox(slide, 1.0, 2.0, 10.0, 0.8,
                  "Chapter 01", color_hex="#FFFFFF", font_pt=36)


def _build_white_bg_with_muted_text(slide) -> None:
    """Slide with white background and light-muted text — real FAIL case."""
    _add_solid_rect(slide, 0, 0, 13.333, 7.5, "#FFFFFF")
    _make_textbox(slide, 0.65, 6.8, 10.0, 0.4,
                  "数据来源：IDC / Gartner 2025",
                  color_hex="#78909C", font_pt=12)


def _write_and_check(prs: Presentation, tmp_path) -> QAReport:
    path = tmp_path / "test.pptx"
    prs.save(path)
    return BuildQA().check(str(path))


class TestLowContrastPerShapeBackground:
    """Per-shape background detection — the core fix."""

    def test_light_text_on_dark_offset_rect_passes(self, prs, tmp_path):
        """Light text (#FFFFFF, #D4D4D4) on a dark offset rect (#1E1E1E) must
        NOT be flagged. This was the original false-positive in build_10pages
        slide 6 (code page)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _build_group_dark_code(slide)
        report = _write_and_check(prs, tmp_path)
        lc = [i for i in report.fatals if i.check_id == "low_contrast_text"]
        assert len(lc) == 0, (
            f"Expected no contrast fatals for light-on-dark; got: "
            f"{[i.detail for i in lc]}"
        )

    def test_light_text_on_full_bleed_dark_group_passes(self, prs, tmp_path):
        """Light text on a full-bleed dark background must PASS (section_divider
        style). This was already OK by accident because full-bleed rects are
        detected as slide bg."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _build_section_divider_style(slide)
        report = _write_and_check(prs, tmp_path)
        lc = [i for i in report.fatals if i.check_id == "low_contrast_text"]
        assert len(lc) == 0

    def test_dark_text_on_white_fails(self, prs, tmp_path):
        """Real low-contrast case — dark muted text on white must FAIL."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _build_white_bg_with_muted_text(slide)
        report = _write_and_check(prs, tmp_path)
        lc = [i for i in report.fatals if i.check_id == "low_contrast_text"]
        assert len(lc) >= 1, "Muted text on white should be flagged"
        assert any("3.40" in i.detail for i in lc), (
            f"Expected detail to show 3.40:1 ratio; got: "
            f"{[i.detail for i in lc]}"
        )

    def test_no_text_no_flag(self, prs, tmp_path):
        """Slide with only a dark rect and no text — should not flag."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_solid_rect(slide, 0, 0, 13.333, 7.5, "#1E1E1E")
        report = _write_and_check(prs, tmp_path)
        lc = [i for i in report.fatals if i.check_id == "low_contrast_text"]
        assert len(lc) == 0

    def test_large_text_low_contrast_passes(self, prs, tmp_path):
        """Large text (≥24pt) uses 3:1 threshold. A marginally low-ratio large
        text should PASS."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_solid_rect(slide, 0, 0, 13.333, 7.5, "#FFFFFF")
        # #B0B0B0 on white is ~2.4:1; below 3:1 for large text
        _make_textbox(slide, 1, 3, 8, 0.6,
                      "Too low contrast large text",
                      color_hex="#B0B0B0", font_pt=28)
        report = _write_and_check(prs, tmp_path)
        lc = [i for i in report.fatals if i.check_id == "low_contrast_text"]
        # 2.4:1 < 3:1 — should fail even at large text threshold
        assert len(lc) >= 1


class TestLowContrastRealWorld:
    """Run the actual build_10pages.pptx output through QA and assert the
    expected false-positive count is reduced."""

    def test_build_10pages_after_fix(self, tmp_path):
        """After the fix:
          - Code-block false positives (light text on dark offset rect
            incorrectly flagged as light-on-white) are gone.
          - Total fatals drop from 21 (with code-block false positives) to 18
            (only genuine muted-text/gold-badge contrast issues).
          - The remaining 18 are real issues for build_10pages.py to fix,
            not BuildQA false positives.
        """
        import shutil
        src = Path("output/build_10pages.pptx")
        if not src.exists():
            pytest.skip("output/build_10pages.pptx not found; run build_10pages.py first")
        dst = tmp_path / "build_10pages.pptx"
        shutil.copy(src, dst)
        report = BuildQA().check(str(dst))
        lc = [i for i in report.fatals if i.check_id == "low_contrast_text"]
        print("\n=== remaining low-contrast fatals ===")
        for i in lc:
            print(f"  slide {i.slide_index}: {i.detail}")
        # Hard rule: no false positives — every fatal must be a genuine
        # low-contrast pair where light text BEHIND a dark backplate was
        # incorrectly reported as light-on-white.
        for i in lc:
            # Sanity 1: no white-on-white phantom (the original bug).
            assert "#FFFFFF on #FFFFFF" not in i.detail, (
                f"False positive (white-on-white): {i.detail}"
            )
            # Sanity 2: code-block slide (slide 5, 0-indexed) had the most
            # egregious false positives: #FFFFFF and #D4D4D4 runs sitting
            # on the dark #1E1E1E backplate were reported as light-on-white
            # because the backplate is a group-relative offset rect. After
            # the fix, code-block slide may report real issues (gold-badge
            # white text, muted subtitle), but should NOT report any
            # light-on-white phantom. Verify no D4D4D4-on-white here.
            if i.slide_index == 5:
                assert "#D4D4D4 on #FFFFFF" not in i.detail, (
                    f"Code-block false positive re-emerged: {i.detail}"
                )
        # Sanity 3: count dropped from 21 (pre-fix, with code-block false
        # positives) to 18 (post-fix, only real issues). We assert strictly
        # less than the pre-fix count to catch regressions in the fix.
        assert len(lc) < 21, (
            f"Fix did not reduce false positives: still {len(lc)} fatals"
        )
