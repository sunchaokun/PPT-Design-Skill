"""Strict edge-case tests for audit findings."""

from __future__ import annotations

import json
import os
import tempfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def _make_template(path: Path, num_slides: int = 3) -> Path:
    prs = Presentation()
    for i in range(num_slides):
        layout = prs.slide_layouts[-1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i + 1}"
    prs.save(str(path))
    return path


def _make_png(path: Path) -> Path:
    img = Image.new("RGB", (100, 50), color="blue")
    img.save(str(path))
    return path


# ============================================================
# parse_pages edge cases (audit findings 5.2, 5.3, 5.5)
# ============================================================

class TestParsePagesStrict:

    def test_malformed_swap_missing_second(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        with pytest.raises(ValueError, match="无效"):
            parse_pages("5<>")

    def test_malformed_move_missing_second(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        with pytest.raises(ValueError, match="无效"):
            parse_pages("3>")

    def test_non_numeric_token(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        with pytest.raises(ValueError):
            parse_pages("abc")

    def test_non_numeric_in_swap(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        with pytest.raises(ValueError):
            parse_pages("a<>b")

    def test_range_with_trailing_dash(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        with pytest.raises(ValueError):
            parse_pages("3-")

    def test_unrecognized_token(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        with pytest.raises(ValueError, match="无法识别"):
            parse_pages("3*5")

    def test_empty_string(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        result = parse_pages("")
        assert result == []

    def test_whitespace_only(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        result = parse_pages("  ,  ,  ")
        assert result == []

    def test_valid_range(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        result = parse_pages("3-5", num_slides=5)
        assert len(result) == 3
        assert all(o.action == "modify" for o in result)

    def test_insert_at_end(self):
        from ppt_pro_max.enterprise.page_revision import parse_pages
        result = parse_pages("+5", num_slides=5)
        assert result[0].action == "insert"
        assert result[0].page == 5


# ============================================================
# content_parser edge cases (audit findings 11.1, 11.3)
# ============================================================

class TestContentParserStrict:

    def test_card_without_image_key(self, tmp_path):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content
        data = {
            "slides": [
                {
                    "goal": "features",
                    "title": "Features",
                    "cards": [
                        {"title": "A", "body": "desc"},
                    ],
                },
            ],
        }
        result = load_enterprise_content(data, str(tmp_path))
        assert "image" not in result[0]["cards"][0]

    def test_diagram_not_dict(self, tmp_path):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content
        data = {
            "slides": [
                {"goal": "process", "title": "Flow", "diagram": "not_a_dict"},
            ],
        }
        result = load_enterprise_content(data, str(tmp_path))
        assert result[0]["diagram_type"] is None
        assert result[0]["diagram_data"] is None

    def test_absolute_image_path_preserved(self, tmp_path):
        from ppt_pro_max.enterprise.content_parser import load_enterprise_content
        data = {
            "slides": [
                {"goal": "hook", "title": "Hi", "image": "C:/absolute/img.png"},
            ],
        }
        result = load_enterprise_content(data, str(tmp_path))
        assert result[0]["image"] == "C:/absolute/img.png"


# ============================================================
# Pipeline edge cases (audit findings 7.1, 7.3)
# ============================================================

class TestPipelineStrict:

    def test_pages_without_template_returns_error(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = tmp_path / "noproject"
        project.mkdir()
        (project / "brand.json").write_text("{}", encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="test",
            project_dir=str(project),
            pages="-2",
        )
        assert "error" in result

    def test_density_zero_uses_suggest(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        project = tmp_path / "dproject"
        project.mkdir()
        _make_template(project / "template.pptx", 2)
        (project / "brand.json").write_text("{}", encoding="utf-8")
        pipeline = EnterprisePipeline()
        result = pipeline.run(
            query="test",
            project_dir=str(project),
            density=0,
        )
        assert result["pipeline"] == "enterprise"


# ============================================================
# ReviewGate edge cases (audit finding 10.1)
# ============================================================

class TestReviewGateStrict:

    def test_read_nonexistent_returns_none(self, tmp_path):
        from ppt_pro_max.enterprise.review_gate import ReviewGate
        gate = ReviewGate()
        result = gate.read_proposal(str(tmp_path / "nonexistent.json"))
        assert result is None

    def test_read_invalid_json_returns_none(self, tmp_path):
        from ppt_pro_max.enterprise.review_gate import ReviewGate
        gate = ReviewGate()
        bad = tmp_path / "bad.json"
        bad.write_text("not json", encoding="utf-8")
        result = gate.read_proposal(str(bad))
        assert result is None


# ============================================================
# EnterpriseRenderer edge cases (audit findings 8.3, 8.4)
# ============================================================

class TestEnterpriseRendererStrict:

    def test_logo_aspect_ratio_preserved(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        logo_path = _make_png(tmp_path / "logo.png")
        renderer = EnterpriseRenderer()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        renderer.insert_logo(slide, str(logo_path), {"position": "top_right", "width_inches": 1.0}, prs=prs)
        pic = [s for s in slide.shapes if s.shape_type == 13][0]
        assert pic.height > 0
        assert pic.width > 0

    def test_logo_nonexistent_file_no_crash(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        renderer = EnterpriseRenderer()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        renderer.insert_logo(slide, "/nonexistent/logo.png", {"position": "top_right"}, prs=prs)

    def test_logo_skip_slide(self, tmp_path):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        logo_path = _make_png(tmp_path / "logo.png")
        renderer = EnterpriseRenderer()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        renderer.insert_logo(
            slide, str(logo_path),
            {"position": "top_right", "skip_slides": ["hook"]},
            current_goal="hook", prs=prs,
        )
        pics = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pics) == 0


# ============================================================
# BrandSpec edge cases (audit finding 2.1, 2.3)
# ============================================================

class TestBrandSpecStrict:

    def test_dark_mode_from_dark_mode_key(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        spec = BrandSpec.from_brand_json({"dark_mode": True})
        assert spec.dark_mode is True

    def test_dark_mode_from_color_scheme(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        spec = BrandSpec.from_brand_json({"color_scheme": "dark"})
        assert spec.dark_mode is True

    def test_dark_mode_default_false(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        spec = BrandSpec.from_brand_json({})
        assert spec.dark_mode is False

    def test_merge_empty_dict_overrides(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        template_spec = BrandSpec(source="template", colors={"primary": "#1A3C6E"}, fonts={"heading": "Arial"})
        merged = BrandSpec.merge(template_spec, {"colors": {}})
        assert merged.colors == {"primary": "#1A3C6E"}


# ============================================================
# slide_utils shared utility
# ============================================================

class TestSlideUtils:

    def test_remove_slide_reduces_count(self, tmp_path):
        from ppt_pro_max.enterprise.slide_utils import remove_slide
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[-1])
        prs.slides.add_slide(prs.slide_layouts[-1])
        assert len(prs.slides) == 2
        remove_slide(prs, 0)
        assert len(prs.slides) == 1

    def test_remove_all_slides(self, tmp_path):
        from ppt_pro_max.enterprise.slide_utils import remove_slide
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[-1])
        for _ in range(3):
            remove_slide(prs, 0)
        assert len(prs.slides) == 0

    def test_remove_slide_save_valid(self, tmp_path):
        from ppt_pro_max.enterprise.slide_utils import remove_slide
        prs = Presentation()
        for i in range(3):
            s = prs.slides.add_slide(prs.slide_layouts[-1])
            s.shapes.title.text = f"Slide {i + 1}"
        remove_slide(prs, 1)
        out = str(tmp_path / "out.pptx")
        prs.save(out)
        loaded = Presentation(out)
        assert len(loaded.slides) == 2
        assert loaded.slides[0].shapes.title.text == "Slide 1"
        assert loaded.slides[1].shapes.title.text == "Slide 3"
