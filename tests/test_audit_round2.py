"""Strict edge-case tests for audit round 2 findings."""

from __future__ import annotations

import json
import os
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def _make_png(path: Path) -> Path:
    img = Image.new("RGB", (100, 50), color="red")
    img.save(str(path))
    return path


def _make_template(path: Path, num_slides: int = 3) -> Path:
    prs = Presentation()
    for i in range(num_slides):
        layout = prs.slide_layouts[min(i + 1, len(prs.slide_layouts) - 1)]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i + 1}"
    prs.save(str(path))
    return path


class TestAddSlideFallback:

    def test_add_slide_no_layout_name_gets_blank_layout(self):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        renderer = EnterpriseRenderer()
        prs = Presentation()
        slide = renderer.add_slide(prs, layout_name=None)
        assert slide is not None

    def test_add_slide_unknown_layout_name_fallback(self):
        from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
        renderer = EnterpriseRenderer()
        prs = Presentation()
        slide = renderer.add_slide(prs, layout_name="NonExistent Layout")
        assert slide is not None


class TestImageMatcherWordBoundary:

    def test_no_false_positive_task_vs_ask(self, tmp_path):
        from ppt_pro_max.enterprise.image_matcher import match_images
        task_img = _make_png(tmp_path / "task_force.png")
        hero_img = _make_png(tmp_path / "invest.png")
        pool = [str(task_img), str(hero_img)]
        designs = [{"goal": "ask", "title": "Investment Ask"}]
        result = match_images(pool, designs)
        assert result[0]["image"] == str(hero_img)

    def test_no_false_positive_spain_vs_pain(self, tmp_path):
        from ppt_pro_max.enterprise.image_matcher import match_images
        spain_img = _make_png(tmp_path / "spain_map.png")
        pain_img = _make_png(tmp_path / "pain_point.png")
        pool = [str(spain_img), str(pain_img)]
        designs = [{"goal": "problem", "title": "The Pain Point"}]
        result = match_images(pool, designs)
        assert result[0]["image"] == str(pain_img)

    def test_exact_match_hero(self, tmp_path):
        from ppt_pro_max.enterprise.image_matcher import match_images
        hero_img = _make_png(tmp_path / "hero.png")
        pool = [str(hero_img)]
        designs = [{"goal": "hook", "title": "Welcome"}]
        result = match_images(pool, designs)
        assert result[0]["image"] == str(hero_img)


class TestBrandSpecDarkMode:

    def test_dark_mode_truthy_int(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        spec = BrandSpec.from_brand_json({"dark_mode": 1})
        assert spec.dark_mode is True

    def test_dark_mode_truthy_string(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        spec = BrandSpec.from_brand_json({"dark_mode": "yes"})
        assert spec.dark_mode is True

    def test_color_scheme_dark(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        spec = BrandSpec.from_brand_json({"color_scheme": "dark"})
        assert spec.dark_mode is True

    def test_dark_mode_false_explicit(self):
        from ppt_pro_max.enterprise.brand_spec import BrandSpec
        spec = BrandSpec.from_brand_json({"dark_mode": False})
        assert spec.dark_mode is False


class TestScannerOutputDirExclusion:

    def test_pptx_in_output_dir_excluded(self, tmp_path):
        from ppt_pro_max.enterprise.scanner import ProjectScanner
        project = tmp_path / "proj"
        project.mkdir()
        output = project / "output"
        output.mkdir()
        _make_template(project / "real_template.pptx", 1)
        _make_template(output / "v1_presentation.pptx", 1)
        scanner = ProjectScanner()
        asset = scanner.scan(str(project))
        assert asset.template_path is not None
        assert "real_template" in asset.template_path

    def test_template_name_preferred(self, tmp_path):
        from ppt_pro_max.enterprise.scanner import ProjectScanner
        project = tmp_path / "proj2"
        project.mkdir()
        _make_template(project / "random.pptx", 1)
        _make_template(project / "template.pptx", 2)
        scanner = ProjectScanner()
        asset = scanner.scan(str(project))
        assert "template" in asset.template_path


class TestPipelineImageFetcher:

    def test_build_image_fetcher_placeholder_returns_none(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        result = pipeline._build_image_fetcher({"image_mode": "placeholder"})
        assert result is None

    def test_build_image_fetcher_search_mode_works(self):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        result = pipeline._build_image_fetcher({"image_mode": "search", "image_config": {}})
        assert result is not None


class TestCLIPageRequiresProject:

    def test_cli_pages_without_project_errors(self):
        import subprocess
        result = subprocess.run(
            ["python", "-m", "ppt_pro_max.cli", "test", "--pages", "3,5"],
            capture_output=True, text=True,
            timeout=10,
        )
        assert result.returncode != 0
        assert "--project" in result.stderr or "project" in result.stderr.lower()


class TestDensityProfileEdgeCases:

    def test_truncate_single_char_max(self):
        from ppt_pro_max.enterprise.density_profile import DensityProfile, apply_density_to_bullets
        p = DensityProfile(title_size=18, subtitle_size=11, body_size=9, bullet_size=7,
                           line_spacing=0.9, max_bullets=15, max_bullet_chars=1, image_width_ratio=0.2)
        result = apply_density_to_bullets(["abc"], p)
        assert len(result) == 1

    def test_zero_density_clamped(self):
        from ppt_pro_max.enterprise.density_profile import get_density_profile
        p = get_density_profile(0)
        assert p.title_size == 36


class TestPageRevisionInsert:

    def test_insert_after_page(self, tmp_path):
        from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
        template_path = _make_template(tmp_path / "template.pptx", 3)
        engine = PageRevisionEngine(str(template_path))
        result = engine.revise(["+2"])
        assert result["num_slides"] == 4


class TestDivisonByZeroGuard:

    @pytest.mark.skip(reason="_insert_content_image removed from Pipeline; image insertion now in PrecisionRenderer.add_image()")
    def test_zero_width_image_guard(self, tmp_path):
        from ppt_pro_max.enterprise.pipeline import EnterprisePipeline
        pipeline = EnterprisePipeline()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        img_path = _make_png(tmp_path / "img.png")
        pipeline._insert_content_image(slide, str(img_path), prs)
        pics = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pics) == 1
