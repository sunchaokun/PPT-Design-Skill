"""Inspect text box positions in the infographic PPTX."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation('output/test_infographic_full.pptx')
slide = prs.slides[0]
for i, shp in enumerate(slide.shapes):
    if shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        x = shp.left / 914400.0
        y = shp.top / 914400.0
        w = shp.width / 914400.0
        h = shp.height / 914400.0
        txt = shp.text_frame.text[:30] if shp.has_text_frame else ''
        print(f'  [{i:2d}] ({x:.2f},{y:.2f},{x+w:.2f},{y+h:.2f}) h={h:.2f} "{txt}"')