"""Report per-slide layout geometry to verify unified grid alignment."""
import sys

sys.path.insert(0, "src")
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(sys.argv[1])
print(f"slide size: {prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} in")
for i, slide in enumerate(prs.slides):
    rows = []
    for s in slide.shapes:
        try:
            l = s.left / 914400
            t = s.top / 914400
            w = s.width / 914400
            h = s.height / 914400
        except Exception:
            continue
        txt = ""
        if hasattr(s, "text") and s.text.strip():
            txt = s.text.strip().replace("\n", " / ")[:28]
        kind = "GROUP" if s.shape_type == MSO_SHAPE_TYPE.GROUP else str(s.shape_type)
        rows.append(f"    ({l:5.2f},{t:5.2f} {w:5.2f}x{h:4.2f}) {kind:8s} {txt}")
    print(f"--- slide {i} ---")
    for r in rows[:14]:
        print(r)
