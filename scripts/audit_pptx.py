"""Quality audit for a generated PPTX (min font, shapes per slide, text presence)."""
import sys

sys.path.insert(0, "src")
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx import Presentation


def audit(path):
    print(f"===== AUDIT {path} =====")
    prs = Presentation(path)
    issues = []
    for i, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        texts = []
        min_size = None
        max_size = 0
        for s in shapes:
            if not (hasattr(s, "text") and s.text.strip()):
                continue
            texts.append(s.text.strip())
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        sz = r.font.size.pt
                        min_size = sz if min_size is None else min(min_size, sz)
                        max_size = max(max_size, sz)
        if len(shapes) < 3:
            issues.append(f"slide {i}: only {len(shapes)} shapes")
        if not texts:
            issues.append(f"slide {i}: no text content")
        if min_size is not None and min_size < 11:
            issues.append(f"slide {i}: min font {min_size}pt < 11pt")
        print(f"slide {i}: shapes={len(shapes)} texts={len(texts)} font_range={min_size}-{max_size}pt")
    print("ISSUES:", issues if issues else "none")
    return issues


if __name__ == "__main__":
    audit(sys.argv[1])
