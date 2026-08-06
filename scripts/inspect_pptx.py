"""Inspect a generated PPTX in detail."""
import sys

sys.path.insert(0, "src")
from pptx import Presentation


def dump(path, label, slide_indices=None):
    print(f"===== {label} =====")
    prs = Presentation(path)
    print(f"slides: {len(prs.slides)}")
    for idx, slide in enumerate(prs.slides):
        if slide_indices is not None and idx not in slide_indices:
            continue
        print(f"--- slide {idx} ---")
        for s in slide.shapes:
            if s.shape_type == 6:
                print(f"  [GROUP] {s.name}")
                continue
            if not (hasattr(s, "text") and s.text.strip()):
                continue
            sizes = set()
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        sizes.add(r.font.size.pt)
            txt = s.text.strip().replace("\n", " | ")
            print(f"  [{s.name}] size={sorted(sizes)} :: {txt[:90]}")


if __name__ == "__main__":
    target = sys.argv[1] if len(sys.argv) > 1 else "test_output/freestyle_dark.pptx"
    idxs = None
    if len(sys.argv) > 2:
        idxs = set(int(x) for x in sys.argv[2].split(","))
    dump(target, target, idxs)
