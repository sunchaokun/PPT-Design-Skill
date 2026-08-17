"""Full-page infographic: 5-layer pyramid + 5 action items with connectors."""
from pptx import Presentation
from pptx.util import Inches
from ppt_pro_max.build_helpers import add_slide, svg_chart

svg = '''<svg viewBox="0 0 1280 720" xmlns="http://www.w3.org/2000/svg">
  <!-- Background -->
  <rect width="1280" height="720" fill="#0F172A"/>
  <!-- Subtle grid -->
  <line x1="0" y1="180" x2="1280" y2="180" stroke="#1E293B" stroke-width="1"/>
  <line x1="640" y1="0" x2="640" y2="720" stroke="#1E293B" stroke-width="0.5" stroke-dasharray="4,4"/>

  <!-- Title -->
  <text x="64" y="80" fill="#FFFFFF" font-size="42" font-weight="700" font-family="Microsoft YaHei">战略层级与关键举措</text>
  <text x="64" y="130" fill="#94A3B8" font-size="22" font-family="Microsoft YaHei">从战略意图到执行闭环，五层递进</text>
  <text x="64" y="165" fill="#64748B" font-size="11" font-family="Cascadia Mono">STRATEGY MAP</text>

  <!-- Pyramid: 5 layers, centered at x=440, top-narrow bottom-wide -->
  <!-- Layer 5 (top) - narrowest: top hw=25, bottom hw=50 -->
  <polygon points="415,220 465,220 490,280 390,280" fill="#22D3EE" opacity="0.9"/>
  <text x="440" y="258" fill="#0F172A" font-size="16" font-weight="700" text-anchor="middle" font-family="Microsoft YaHei">愿景使命</text>

  <!-- Layer 4: top hw=50, bottom hw=80 -->
  <polygon points="390,285 490,285 520,345 360,345" fill="#22D3EE" opacity="0.75"/>
  <text x="440" y="323" fill="#0F172A" font-size="16" font-weight="700" text-anchor="middle" font-family="Microsoft YaHei">战略目标</text>

  <!-- Layer 3: top hw=80, bottom hw=115 -->
  <polygon points="360,350 520,350 555,410 325,410" fill="#22D3EE" opacity="0.6"/>
  <text x="440" y="388" fill="#0F172A" font-size="16" font-weight="700" text-anchor="middle" font-family="Microsoft YaHei">关键举措</text>

  <!-- Layer 2: top hw=115, bottom hw=155 -->
  <polygon points="325,415 555,415 595,475 285,475" fill="#22D3EE" opacity="0.45"/>
  <text x="440" y="453" fill="#FFFFFF" font-size="16" font-weight="700" text-anchor="middle" font-family="Microsoft YaHei">执行路径</text>

  <!-- Layer 1 (bottom) - widest: top hw=155, bottom hw=200 -->
  <polygon points="285,480 595,480 640,540 240,540" fill="#22D3EE" opacity="0.3"/>
  <text x="440" y="518" fill="#FFFFFF" font-size="16" font-weight="700" text-anchor="middle" font-family="Microsoft YaHei">运营支撑</text>

  <!-- Connectors: pyramid right edge → action items -->
  <line x1="478" y1="250" x2="700" y2="250" stroke="#A78BFA" stroke-width="1.5" stroke-dasharray="6,3"/>
  <line x1="505" y1="315" x2="700" y2="315" stroke="#A78BFA" stroke-width="1.5" stroke-dasharray="6,3"/>
  <line x1="538" y1="380" x2="700" y2="380" stroke="#A78BFA" stroke-width="1.5" stroke-dasharray="6,3"/>
  <line x1="575" y1="445" x2="700" y2="445" stroke="#A78BFA" stroke-width="1.5" stroke-dasharray="6,3"/>
  <line x1="618" y1="510" x2="700" y2="510" stroke="#A78BFA" stroke-width="1.5" stroke-dasharray="6,3"/>

  <!-- Right-side: 5 action items (badge+title in same <text> via tspan for baseline alignment) -->
  <!-- Item 1 -->
  <rect x="700" y="225" width="540" height="50" rx="6" fill="#1E293B" stroke="#22D3EE" stroke-width="1"/>
  <text x="720" y="254" fill="#22D3EE" font-size="14" font-weight="700" font-family="Cascadia Mono">01 <tspan fill="#FFFFFF" font-size="14" font-weight="600" font-family="Microsoft YaHei">定义北极星指标与3年愿景</tspan></text>
  <text x="750" y="269" fill="#94A3B8" font-size="12" font-family="Microsoft YaHei">明确方向，凝聚团队共识</text>

  <!-- Item 2 -->
  <rect x="700" y="290" width="540" height="50" rx="6" fill="#1E293B" stroke="#22D3EE" stroke-width="1"/>
  <text x="720" y="319" fill="#22D3EE" font-size="14" font-weight="700" font-family="Cascadia Mono">02 <tspan fill="#FFFFFF" font-size="14" font-weight="600" font-family="Microsoft YaHei">设定OKR与关键结果</tspan></text>
  <text x="750" y="334" fill="#94A3B8" font-size="12" font-family="Microsoft YaHei">量化目标，可追踪可衡量</text>

  <!-- Item 3 -->
  <rect x="700" y="355" width="540" height="50" rx="6" fill="#1E293B" stroke="#FBBF24" stroke-width="1"/>
  <text x="720" y="384" fill="#FBBF24" font-size="14" font-weight="700" font-family="Cascadia Mono">03 <tspan fill="#FFFFFF" font-size="14" font-weight="600" font-family="Microsoft YaHei">识别关键举措与优先级</tspan></text>
  <text x="750" y="399" fill="#94A3B8" font-size="12" font-family="Microsoft YaHei">聚焦高杠杆行动，避免资源分散</text>

  <!-- Item 4 -->
  <rect x="700" y="420" width="540" height="50" rx="6" fill="#1E293B" stroke="#22D3EE" stroke-width="1"/>
  <text x="720" y="449" fill="#22D3EE" font-size="14" font-weight="700" font-family="Cascadia Mono">04 <tspan fill="#FFFFFF" font-size="14" font-weight="600" font-family="Microsoft YaHei">制定执行路径与里程碑</tspan></text>
  <text x="750" y="464" fill="#94A3B8" font-size="12" font-family="Microsoft YaHei">分解为季度Sprint，持续交付</text>

  <!-- Item 5 -->
  <rect x="700" y="485" width="540" height="50" rx="6" fill="#1E293B" stroke="#22D3EE" stroke-width="1"/>
  <text x="720" y="514" fill="#22D3EE" font-size="14" font-weight="700" font-family="Cascadia Mono">05 <tspan fill="#FFFFFF" font-size="14" font-weight="600" font-family="Microsoft YaHei">建立运营支撑与反馈闭环</tspan></text>
  <text x="750" y="529" fill="#94A3B8" font-size="12" font-family="Microsoft YaHei">数据驱动，快速迭代优化</text>

  <!-- Bottom accent line -->
  <rect x="64" y="680" width="200" height="3" rx="1.5" fill="#22D3EE"/>
  <text x="64" y="705" fill="#475569" font-size="10" font-family="Cascadia Mono">CONFIDENTIAL · Q3 2026</text>
</svg>'''

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = add_slide(prs)
r = svg_chart(s, svg, 0, 0, 13.333, 7.5)
print(f'shapes={r.shape_count}, warnings={len(r.warnings)}, compile_ms={r.compile_ms:.0f}')
for w in r.warnings:
    print(f'  WARN: {w}')
prs.save('output/test_infographic_full.pptx')
print('Saved output/test_infographic_full.pptx')
