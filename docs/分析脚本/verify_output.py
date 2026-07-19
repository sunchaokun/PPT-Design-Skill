import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches

prs = Presentation(r'E:\简历\output\template_with_data_analysis.pptx')
print(f'Total slides: {len(prs.slides)}')
print(f'Total shapes: {sum(len(s.shapes) for s in prs.slides)}')
print()

for si, slide in enumerate(prs.slides):
    shape_count = len(slide.shapes)
    texts = []
    min_font = 999
    max_font = 0
    has_image = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    texts.append(p.text.strip()[:50])
                if p.font.size:
                    pt = p.font.size.pt
                    min_font = min(min_font, pt)
                    max_font = max(max_font, pt)
                for r in p.runs:
                    if r.font.size:
                        pt = r.font.size.pt
                        min_font = min(min_font, pt)
                        max_font = max(max_font, pt)
    print(f'Slide {si}: {shape_count} shapes, font range: {min_font:.0f}pt-{max_font:.0f}pt')
    for t in texts[:5]:
        print(f'  - {t}')
    if len(texts) > 5:
        print(f'  ... and {len(texts)-5} more text elements')
    print()

import os
size = os.path.getsize(r'E:\简历\output\template_with_data_analysis.pptx')
print(f'File size: {size/1024:.1f} KB')
