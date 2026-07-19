from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os
import copy

C = {
    'bg': '#FFFFFF',
    'bg_tint': '#F2F8D6',
    'primary': '#2E6504',
    'primary_mid': '#466740',
    'accent': '#7DA92F',
    'muted': '#84AF7D',
    'light': '#D4E3AC',
    'lighter': '#E7F3A6',
    'lightest': '#F2F8D6',
    'white': '#FFFFFF',
    'text_dark': '#0D4609',
    'text_body': '#2E6504',
    'text_muted': '#466740',
    'divider': '#84AF7D',
    'card_bg': '#F6FAE8',
    'highlight': '#7DA92F',
}

template_path = r'E:\简历\template.pptx'
prs = Presentation(template_path)
W = prs.slide_width
H = prs.slide_height
L = Inches(0.65)
R = W - Inches(0.65)
CW = R - L


def _rgb(hex_str):
    return RGBColor.from_string(hex_str.lstrip('#'))


def _rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _rrect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _oval(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _text(slide, left, top, width, height, text, font_size=12,
          color='#2E6504', bold=False, align=PP_ALIGN.LEFT,
          font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = _rgb(color)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return txBox


def _multiline(slide, left, top, width, height, lines, font_size=12,
               color='#2E6504', bold=False, align=PP_ALIGN.LEFT,
               font_name='Microsoft YaHei', line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb(color)
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return txBox


def _copy_logo_group(slide):
    source_slide = prs.slides[1]
    for shape in source_slide.shapes:
        if shape.shape_type == 6:
            sp = shape._element
            if sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr[@val="2E6504"]') is not None or sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr[@val="7DA92F"]') is not None:
                el = copy.deepcopy(sp)
                slide.shapes._spTree.append(el)
                return


def _add_logo(slide):
    _text(slide, Inches(0.65), Inches(0.53), Inches(2.28), Inches(0.34),
          'your logo', font_size=10, color=C['text_body'], bold=True)


def _add_top_bar(slide):
    _rect(slide, 0, 0, W, Inches(0.08), C['accent'])


def _page_header(slide, title, subtitle=''):
    _text(slide, L, Inches(0.45), CW, Inches(0.5),
          title, font_size=28, color=C['text_dark'], bold=True)
    if subtitle:
        _text(slide, L, Inches(0.95), CW, Inches(0.25),
              subtitle, font_size=11, color=C['text_muted'])
    _rect(slide, L, Inches(1.25), CW, Pt(3), C['divider'])


def _kpi_card(slide, left, top, width, height, number, label, trend='', trend_up=True):
    _rrect(slide, left, top, width, height, C['white'], line=C['light'])
    _rect(slide, left, top, width, Inches(0.06), C['accent'])

    _text(slide, left + Inches(0.2), top + Inches(0.25), width - Inches(0.4), Inches(0.5),
          number, font_size=28, color=C['primary'], bold=True)
    _text(slide, left + Inches(0.2), top + Inches(0.75), width - Inches(0.4), Inches(0.3),
          label, font_size=11, color=C['text_muted'])
    if trend:
        trend_color = '#2E6504' if trend_up else '#C53030'
        _text(slide, left + Inches(0.2), top + Inches(1.05), width - Inches(0.4), Inches(0.25),
              trend, font_size=10, color=trend_color, bold=True)


def _bar(slide, left, top, width, max_width, height, fill, label='', value=''):
    _rrect(slide, left, top, max_width, height, C['bg_tint'])
    bar_w = int(max_width * width)
    if bar_w > 0:
        _rrect(slide, left, top, bar_w, height, fill)
    if label:
        _text(slide, left - Inches(0.9), top - Inches(0.03), Inches(0.85), height,
              label, font_size=10, color=C['text_body'], align=PP_ALIGN.RIGHT)
    if value:
        _text(slide, left + max_width + Inches(0.08), top - Inches(0.03), Inches(0.6), height,
              value, font_size=10, color=C['text_dark'], bold=True)


# ============================================================
# PAGE 6 — DATA ANALYSIS: Key Metrics Dashboard
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
_rect(s6, 0, 0, W, H, C['bg'])
_add_top_bar(s6)
_add_logo(s6)
_page_header(s6, '乡村振兴核心数据总览', '2024年度关键指标汇总')

kpi_w = (CW - Inches(0.35) * 3) / 4
kpi_data = [
    ('12.8亿', '年度农业总产值', '同比 +8.3%', True),
    ('1,247万', '新增就业人数', '同比 +15.2%', True),
    ('3,580', '示范村建设数', '同比 +22.6%', True),
    ('96.5%', '基础设施覆盖率', '同比 +3.1%', True),
]
kx = L
for num, label, trend, up in kpi_data:
    _kpi_card(s6, kx, Inches(1.55), kpi_w, Inches(1.35), num, label, trend, up)
    kx += kpi_w + Inches(0.35)

_text(s6, L, Inches(3.15), Inches(5), Inches(0.3),
      '各区域产业投入分布', font_size=14, color=C['text_dark'], bold=True)

bar_left = L + Inches(1.0)
bar_max_w = Inches(5.0)
bar_data = [
    ('东部地区', 0.92, '460亿'),
    ('中部地区', 0.75, '375亿'),
    ('西部地区', 0.68, '340亿'),
    ('东北地区', 0.48, '240亿'),
]
by = Inches(3.55)
bar_colors = [C['primary'], C['accent'], C['muted'], C['light']]
for i, (label, pct, val) in enumerate(bar_data):
    _bar(s6, bar_left, by, pct, bar_max_w, Inches(0.3), bar_colors[i], label, val)
    by += Inches(0.5)

chart_left = L + Inches(7.0)
chart_w = Inches(5.0)
_text(s6, chart_left, Inches(3.15), chart_w, Inches(0.3),
      '季度农民收入增长趋势', font_size=14, color=C['text_dark'], bold=True)

_rect(s6, chart_left, Inches(3.55), chart_w, Inches(3.2), C['bg_tint'], line=C['light'])

quarters = ['Q1', 'Q2', 'Q3', 'Q4']
values = [8200, 9500, 11200, 12800]
max_val = 14000
chart_base_y = Inches(6.45)
chart_h = Inches(2.6)
bx_start = chart_left + Inches(0.5)
bar_unit_w = Inches(0.8)
bar_gap = Inches(0.4)

for i, (q, v) in enumerate(zip(quarters, values)):
    bx = bx_start + i * (bar_unit_w + bar_gap)
    bar_h = int(chart_h * v / max_val)
    bar_top = chart_base_y - bar_h

    _rrect(s6, bx, bar_top, bar_unit_w, bar_h, C['accent'] if i < 3 else C['primary'])
    _text(s6, bx, bar_top - Inches(0.25), bar_unit_w, Inches(0.25),
          f'{v//1000}k', font_size=10, color=C['text_dark'], bold=True, align=PP_ALIGN.CENTER)
    _text(s6, bx, chart_base_y + Inches(0.05), bar_unit_w, Inches(0.2),
          q, font_size=10, color=C['text_muted'], align=PP_ALIGN.CENTER)

_rect(s6, chart_left + Inches(0.3), chart_base_y, chart_w - Inches(0.6), Pt(2), C['divider'])

_text(s6, L, Inches(6.85), CW, Inches(0.25),
      '数据来源：国家统计局 · 农业农村部 | 统计周期：2024年1月-12月',
      font_size=9, color=C['text_muted'], align=PP_ALIGN.RIGHT)


# ============================================================
# PAGE 7 — DATA ANALYSIS: Detailed Comparison & Breakdown
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
_rect(s7, 0, 0, W, H, C['bg'])
_add_top_bar(s7)
_add_logo(s7)
_page_header(s7, '乡村振兴成效对比分析', '2023 vs 2024 各维度数据对比')

_text(s7, L, Inches(1.55), Inches(5), Inches(0.3),
      '五大核心指标对比', font_size=14, color=C['text_dark'], bold=True)

metrics = [
    ('农业产值', '11.8亿', '12.8亿', 0.85, 0.92),
    ('就业人数', '1083万', '1247万', 0.65, 0.75),
    ('示范村数', '2920', '3580', 0.55, 0.68),
    ('设施覆盖', '93.4%', '96.5%', 0.78, 0.82),
    ('人均收入', '1.8万', '2.1万', 0.60, 0.70),
]

my = Inches(1.95)
mx = L + Inches(1.2)
mmax = Inches(4.0)
for label, v2023, v2024, pct23, pct24 in metrics:
    _text(s7, L, my - Inches(0.02), Inches(1.1), Inches(0.2),
          label, font_size=11, color=C['text_body'], bold=True, align=PP_ALIGN.RIGHT)

    _rrect(s7, mx, my, mmax, Inches(0.18), C['bg_tint'])
    bar23 = int(mmax * pct23)
    if bar23 > 0:
        _rrect(s7, mx, my, bar23, Inches(0.18), C['muted'])

    _rrect(s7, mx, my + Inches(0.22), mmax, Inches(0.18), C['bg_tint'])
    bar24 = int(mmax * pct24)
    if bar24 > 0:
        _rrect(s7, mx, my + Inches(0.22), bar24, Inches(0.18), C['primary'])

    _text(s7, mx + mmax + Inches(0.1), my - Inches(0.03), Inches(0.8), Inches(0.2),
          v2023, font_size=9, color=C['text_muted'])
    _text(s7, mx + mmax + Inches(0.1), my + Inches(0.19), Inches(0.8), Inches(0.2),
          v2024, font_size=9, color=C['text_dark'], bold=True)

    my += Inches(0.55)

legend_y = Inches(4.85)
_rrect(s7, L + Inches(1.5), legend_y, Inches(0.25), Inches(0.15), C['muted'])
_text(s7, L + Inches(1.8), legend_y - Inches(0.02), Inches(0.6), Inches(0.2),
      '2023', font_size=9, color=C['text_muted'])
_rrect(s7, L + Inches(2.5), legend_y, Inches(0.25), Inches(0.15), C['primary'])
_text(s7, L + Inches(2.8), legend_y - Inches(0.02), Inches(0.6), Inches(0.2),
      '2024', font_size=9, color=C['text_dark'], bold=True)

right_x = L + Inches(7.0)
right_w = Inches(5.0)
_text(s7, right_x, Inches(1.55), right_w, Inches(0.3),
      '产业投资结构分布', font_size=14, color=C['text_dark'], bold=True)

sectors = [
    ('现代农业', '32%', C['primary']),
    ('乡村旅游', '24%', C['accent']),
    ('电商物流', '18%', C['muted']),
    ('绿色能源', '14%', C['light']),
    ('其他产业', '12%', '#B8D4A8'),
]

donut_cx = right_x + Inches(1.8)
donut_cy = Inches(3.5)
donut_r = Inches(1.3)
donut_inner = Inches(0.7)

start_angle = 0
sector_angles = []
for name, pct_str, clr in sectors:
    pct_val = int(pct_str.replace('%', ''))
    angle = pct_val * 3.6
    sector_angles.append((name, pct_str, clr, start_angle, start_angle + angle))
    start_angle += angle

import math
for name, pct_str, clr, sa, ea in sector_angles:
    mid = math.radians((sa + ea) / 2)
    seg_w = int(donut_r * abs(math.cos(math.radians(ea - sa) / 2)) * 2)
    seg_h = int(donut_r * abs(math.cos(math.radians(ea - sa) / 2)) * 2)

    sx = int(donut_cx + donut_r * 0.85 * math.cos(mid) - seg_w / 2)
    sy = int(donut_cy + donut_r * 0.85 * math.sin(mid) - seg_h / 2)

    _oval(s7, int(donut_cx - donut_r), int(donut_cy - donut_r),
          int(donut_r * 2), int(donut_r * 2), clr)

_oval(s7, int(donut_cx - donut_inner), int(donut_cy - donut_inner),
      int(donut_inner * 2), int(donut_inner * 2), C['bg'])

_text(s7, int(donut_cx - Inches(0.5)), int(donut_cy - Inches(0.2)),
      Inches(1.0), Inches(0.4),
      '100%', font_size=20, color=C['primary'], bold=True, align=PP_ALIGN.CENTER)

legend_x = right_x + Inches(3.8)
ly = Inches(2.1)
for name, pct_str, clr, _, _ in sector_angles:
    _rrect(s7, legend_x, ly, Inches(0.2), Inches(0.2), clr)
    _text(s7, legend_x + Inches(0.3), ly - Inches(0.02), Inches(1.5), Inches(0.25),
          f'{name}  {pct_str}', font_size=11, color=C['text_body'])
    ly += Inches(0.35)

_text(s7, L, Inches(5.3), CW, Inches(0.3),
      '增长亮点分析', font_size=14, color=C['text_dark'], bold=True)

highlights = [
    ('乡村旅游增速领跑', '同比增长34.2%，成为乡村振兴新引擎', C['primary']),
    ('电商覆盖突破', '村级电商服务站覆盖率达78.5%', C['accent']),
    ('绿色转型加速', '清洁能源应用比例提升至42%', C['muted']),
]

hx = L
hw = (CW - Inches(0.35) * 2) / 3
for title, desc, clr in highlights:
    _rrect(s7, hx, Inches(5.7), hw, Inches(1.2), C['card_bg'], line=C['light'])
    _rect(s7, hx, Inches(5.7), hw, Inches(0.06), clr)
    _text(s7, hx + Inches(0.2), Inches(5.88), hw - Inches(0.4), Inches(0.3),
          title, font_size=12, color=C['text_dark'], bold=True)
    _text(s7, hx + Inches(0.2), Inches(6.22), hw - Inches(0.4), Inches(0.5),
          desc, font_size=10, color=C['text_muted'])
    hx += hw + Inches(0.35)

_text(s7, L, Inches(7.0), CW, Inches(0.25),
      '数据来源：国家统计局 · 农业农村部 | 统计周期：2023-2024年度',
      font_size=9, color=C['text_muted'], align=PP_ALIGN.RIGHT)


# ============================================================
# SAVE
# ============================================================
output_dir = r'E:\简历\output'
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, 'template_with_data_analysis.pptx')
prs.save(output_path)
total_shapes = sum(len(s.shapes) for s in prs.slides)
print(f'Saved: {output_path}')
print(f'Slides: {len(prs.slides)}')
print(f'Shapes: {total_shapes}')
