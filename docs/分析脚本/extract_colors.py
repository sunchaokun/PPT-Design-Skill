import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree

prs = Presentation(r'E:\简历\template.pptx')

nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

for si, slide in enumerate(prs.slides):
    print(f'=== Slide {si} ===')
    for shape in slide.shapes:
        sp = shape._element
        fills = sp.findall('.//a:solidFill/a:srgbClr', nsmap)
        for f in fills:
            val = f.get('val', '')
            print(f'  Shape fill color: #{val}')
        fills2 = sp.findall('.//a:solidFill/a:schemeClr', nsmap)
        for f in fills2:
            val = f.get('val', '')
            print(f'  Shape scheme color: {val}')
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    rf = r._r.findall('.//a:solidFill/a:srgbClr', nsmap)
                    for f in rf:
                        val = f.get('val', '')
                        print(f'  Text color: #{val} text="{r.text[:30]}"')

# Also check slide master / theme colors
print('\n=== Theme Colors ===')
theme = prs.slide_masters[0].element.findall('.//a:clrScheme/a:*', nsmap)
for t in theme:
    tag = t.tag.split('}')[1] if '}' in t.tag else t.tag
    srgb = t.find('a:srgbClr', nsmap)
    if srgb is not None:
        print(f'  {tag}: #{srgb.get("val")}')
    else:
        sysclr = t.find('a:sysClr', nsmap)
        if sysclr is not None:
            print(f'  {tag}: sys={sysclr.get("val")} lastClr={sysclr.get("lastClr")}')
