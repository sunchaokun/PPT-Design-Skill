import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

prs = Presentation(r'E:\简历\template.pptx')
print(f'Slide width (inches): {prs.slide_width / 914400:.3f}, height (inches): {prs.slide_height / 914400:.3f}')
print(f'Slides: {len(prs.slides)}')
print(f'Slide layouts: {len(prs.slide_layouts)}')
for i, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {i}: {layout.name}')
print()
for si, slide in enumerate(prs.slides):
    print(f'--- Slide {si} ---')
    for shape in slide.shapes:
        pos = f'({shape.left/914400:.2f}, {shape.top/914400:.2f})'
        size = f'({shape.width/914400:.2f}x{shape.height/914400:.2f})'
        info = f'{pos} {size}'
        if shape.has_text_frame:
            txt = shape.text_frame.text[:80].replace(chr(10), ' | ')
            info += f' text="{txt}"'
            if shape.text_frame.paragraphs:
                p0 = shape.text_frame.paragraphs[0]
                if p0.font.size:
                    info += f' size={p0.font.size.pt}pt'
                try:
                    if p0.font.color and p0.font.color.type is not None:
                        info += f' color=#{p0.font.color.rgb}'
                except:
                    pass
                if p0.font.name:
                    info += f' font={p0.font.name}'
        print(f'  {shape.shape_type}: {info}')
