"""probe_svg2pptx.py — SVG → editable-pptx feasibility probe.

Self-contained experiment (scratch/, NEVER imported by src/).  Answers:
  1. Can a subset of SVG compile to NATIVE EDITABLE pptx shapes
     (custGeom / MSO shapes / groups), not pictures?
  2. Where is the boundary (which SVG features degrade or fail)?
  3. Is it fast?  Is rendered geometry faithful (pixel-compare)?
  4. Does it scale to "long-tail" professional business charts?

Pipeline per case:
    SVG ──(Edge headless)──▶ ground-truth PNG
    SVG ──(probe compiler)──▶ editable .pptx ──(LibreOffice)──▶ result PNG
    pixel-compare + editability audit + boundary report + timing

Reuses existing primitives (freeform_builder, boolean_shapes,
text_measurer, build_helpers) via import; creates NO files under src/.
"""

from __future__ import annotations

import math
import os
import re
import subprocess
import sys
import time
import zipfile
from pathlib import Path

import numpy as np
from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ppt_pro_max.renderer.freeform_builder import FreeformBuilder
from ppt_pro_max.renderer.boolean_shapes import bool_shape
from ppt_pro_max.renderer.text_measurer import estimate_text_size

HERE = Path(__file__).resolve().parent
OUT = HERE / "out"
CASES = HERE / "cases"

SVG = "{http://www.w3.org/2000/svg}"
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# ─────────────────────────── test data ───────────────────────────

CASES_SVG = {
    "pyramid": """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 360">
  <defs>
    <linearGradient id="g1" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0" stop-color="#7DA92F"/><stop offset="1" stop-color="#2E6504"/>
    </linearGradient>
    <linearGradient id="g2" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0" stop-color="#5B9BD5"/><stop offset="1" stop-color="#2E75B6"/>
    </linearGradient>
    <linearGradient id="g3" x1="0" y1="0" x2="1" y2="1">
      <stop offset="0" stop-color="#FFC000"/><stop offset="1" stop-color="#BF8F00"/>
    </linearGradient>
  </defs>
  <polygon points="200,30 380,320 20,320" fill="url(#g1)"/>
  <polygon points="200,120 320,320 80,320" fill="url(#g2)"/>
  <polygon points="200,210 260,320 140,320" fill="url(#g3)"/>
  <text x="200" y="85" text-anchor="middle" font-size="16" fill="#fff" font-family="Arial">战略愿景</text>
  <text x="200" y="175" text-anchor="middle" font-size="14" fill="#fff" font-family="Arial">三年目标</text>
  <text x="200" y="265" text-anchor="middle" font-size="14" fill="#fff" font-family="Arial">年度计划</text>
</svg>""",

    "venn_evenodd": """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 260">
  <path d="M150,130 a80,80 0 1,0 0.1,0 Z M250,130 a80,80 0 1,0 0.1,0 Z"
        fill="#E8534E" fill-opacity="0.85"/>
  <text x="105" y="135" text-anchor="middle" font-size="15" fill="#fff">客户A</text>
  <text x="295" y="135" text-anchor="middle" font-size="15" fill="#fff">客户B</text>
  <text x="200" y="135" text-anchor="middle" font-size="13" fill="#fff">交叉用户</text>
</svg>""",

    "funnel": """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 420">
  <defs>
    <linearGradient id="f1" x1="0" y1="0" x2="0" y2="1">
      <stop offset="0" stop-color="#4472C4"/><stop offset="1" stop-color="#2E5BA6"/>
    </linearGradient>
  </defs>
  <polygon points="40,20 360,20 300,90 100,90" fill="url(#f1)"/>
  <polygon points="70,110 330,110 280,180 120,180" fill="#5B9BD5"/>
  <polygon points="100,200 300,200 260,270 140,270" fill="#9DC3E6"/>
  <polygon points="130,290 270,290 240,360 160,360" fill="#BDD7EE"/>
  <text x="200" y="55" text-anchor="middle" font-size="15" fill="#fff">曝光 10000</text>
  <text x="200" y="145" text-anchor="middle" font-size="15" fill="#fff">点击 4000</text>
  <text x="200" y="235" text-anchor="middle" font-size="15" fill="#2E5BA6">注册 1500</text>
  <text x="200" y="325" text-anchor="middle" font-size="15" fill="#2E5BA6">转化 600</text>
</svg>""",

    "growth_curve": """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">
  <defs>
    <clipPath id="clip"><rect x="30" y="40" width="340" height="220"/></clipPath>
    <linearGradient id="area" x1="0" y1="0" x2="0" y2="1">
      <stop offset="0" stop-color="#1D78FA" stop-opacity="0.45"/>
      <stop offset="1" stop-color="#1D78FA" stop-opacity="0.03"/>
    </linearGradient>
  </defs>
  <line x1="40" y1="250" x2="380" y2="250" stroke="#D9D9D9" stroke-width="1.5"/>
  <line x1="40" y1="250" x2="40" y2="50" stroke="#D9D9D9" stroke-width="1.5"/>
  <g clip-path="url(#clip)">
    <path d="M40,250 C120,240 150,180 200,150 C260,115 300,90 380,60 L380,250 Z"
          fill="url(#area)" fill-rule="nonzero"/>
    <path d="M40,250 C120,240 150,180 200,150 C260,115 300,90 380,60"
          fill="none" stroke="#1D78FA" stroke-width="3.5"/>
  </g>
  <circle cx="380" cy="60" r="6" fill="#1D78FA"/>
  <text x="200" y="275" text-anchor="middle" font-size="13" fill="#595959">季度增长曲线</text>
</svg>""",

    "matrix_bcg": """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 340">
  <g transform="translate(20,20)">
    <rect x="0" y="0" width="360" height="300" fill="#FFFFFF" stroke="#BFBFBF" stroke-width="1"/>
    <line x1="180" y1="0" x2="180" y2="300" stroke="#BFBFBF" stroke-width="1" stroke-dasharray="4,4"/>
    <line x1="0" y1="150" x2="360" y2="150" stroke="#BFBFBF" stroke-width="1" stroke-dasharray="4,4"/>
    <text x="180" y="20" text-anchor="middle" font-size="11" fill="#7F7F7F">高</text>
    <text x="180" y="296" text-anchor="middle" font-size="11" fill="#7F7F7F">低</text>
    <g transform="translate(30,30)">
      <circle cx="0" cy="0" r="24" fill="#2E75B6"/><text x="0" y="5" text-anchor="middle" font-size="11" fill="#fff">明星</text>
      <circle cx="90" cy="0" r="24" fill="#BF8F00"/><text x="90" y="5" text-anchor="middle" font-size="11" fill="#fff">问题</text>
    </g>
    <g transform="translate(30,90)">
      <circle cx="0" cy="0" r="24" fill="#70AD47"/><text x="0" y="5" text-anchor="middle" font-size="11" fill="#fff">现金牛</text>
      <circle cx="90" cy="0" r="24" fill="#C00000"/><text x="90" y="5" text-anchor="middle" font-size="11" fill="#fff">瘦狗</text>
    </g>
  </g>
</svg>""",

    # boundary probe: features we deliberately do NOT support (mask/filter/image)
    "unsupported": """<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 400 300">
  <defs>
    <filter id="blur"><feGaussianBlur stdDeviation="3"/></filter>
    <mask id="m"><rect x="0" y="0" width="400" height="300" fill="#fff"/></mask>
  </defs>
  <image href="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
         x="50" y="50" width="100" height="100"/>
  <rect x="20" y="180" width="360" height="80" fill="#4472C4" filter="url(#blur)"/>
  <circle cx="200" cy="120" r="60" fill="#E8534E" mask="url(#m)"/>
</svg>""",
}

# ─────────────────────────── affine ───────────────────────────

class Affine:
    def __init__(self, a=1.0, b=0.0, c=0.0, d=1.0, e=0.0, f=0.0):
        self.a, self.b, self.c, self.d, self.e, self.f = a, b, c, d, e, f

    def compose(self, o):
        return Affine(self.a * o.a + self.b * o.c, self.a * o.b + self.b * o.d,
                      self.c * o.a + self.d * o.c, self.c * o.b + self.d * o.d,
                      self.e * o.a + self.f * o.c + o.e, self.e * o.b + self.f * o.d + o.f)

    def apply(self, x, y):
        return (self.a * x + self.c * y + self.e, self.b * x + self.d * y + self.f)


def parse_transform(s):
    if not s:
        return Affine()
    out = Affine()
    for m in re.finditer(r"([a-z]+)\(([^)]*)\)", s):
        op = m.group(1)
        args = [float(v) for v in m.group(2).replace(",", " ").split()]
        if op == "translate":
            t = Affine(1, 0, 0, 1, args[0], args[1] if len(args) > 1 else 0)
        elif op == "scale":
            t = Affine(args[0], 0, 0, args[1] if len(args) > 1 else args[0], 0, 0)
        elif op == "rotate":
            ang = math.radians(args[0]); cx, cy = args[1], args[2] if len(args) > 2 else 0, (args[2] if len(args) > 2 else 0)
            cy = args[2] if len(args) > 2 else 0
            t = (Affine(1, 0, 0, 1, cx, cy)
                 .compose(Affine(math.cos(ang), math.sin(ang), -math.sin(ang), math.cos(ang), 0, 0))
                 .compose(Affine(1, 0, 0, 1, -cx, -cy)))
        elif op == "matrix":
            t = Affine(*args)
        else:
            t = Affine()
        out = out.compose(t)
    return out


# ─────────────────────────── path parsing ───────────────────────────

_TOKEN = re.compile(r"([MmLlHhVvCcSsQqTtAaZz])|([-+]?(?:\d*\.\d+|\d+\.?)(?:[eE][-+]?\d+)?)")


def arc_to_cubics(x0, y0, rx, ry, rot, large, sweep, x1, y1):
    rx, ry = abs(rx), abs(ry)
    if rx < 1e-9 or ry < 1e-9:
        return [((x0, y0), (x0, y0), (x1, y1), (x1, y1))]
    phi = math.radians(rot); cp, sp = math.cos(phi), math.sin(phi)
    dx, dy = (x0 - x1) / 2, (y0 - y1) / 2
    x1p, y1p = cp * dx + sp * dy, -sp * dx + cp * dy
    lam = x1p * x1p / (rx * rx) + y1p * y1p / (ry * ry)
    if lam > 1:
        r = math.sqrt(lam); rx *= r; ry *= r
    num = rx * rx * ry * ry - rx * rx * y1p * y1p - ry * ry * x1p * x1p
    den = rx * rx * y1p * y1p + ry * ry * x1p * x1p
    coef = (-1.0 if large == sweep else 1.0) * math.sqrt(max(0.0, num / den))
    cxp = coef * (rx * y1p / ry); cyp = coef * (-ry * x1p / rx)
    cx = cp * cxp - sp * cyp + (x0 + x1) / 2; cy = sp * cxp + cp * cyp + (y0 + y1) / 2
    ux, uy = (x1p - cxp) / rx, (y1p - cyp) / ry
    vx, vy = (-x1p - cxp) / rx, (-y1p - cyp) / ry
    th1 = math.atan2(uy, ux)
    dth = math.atan2(vy, vx) - th1
    if sweep == 0 and dth > 0: dth -= 2 * math.pi
    if sweep == 1 and dth < 0: dth += 2 * math.pi
    # large-arc must span > half circle; the "full circle" idiom has nearly
    # coincident endpoints, so force the long way around the circle.
    if large == 1 and abs(dth) < math.pi:
        dth = dth - 2 * math.pi if dth >= 0 else dth + 2 * math.pi
    if large == 0 and abs(dth) > math.pi:
        dth = dth + 2 * math.pi if dth < 0 else dth - 2 * math.pi
    n = max(1, math.ceil(abs(dth) / (math.pi / 2)))
    segs = []
    for i in range(n):
        a1 = th1 + dth * i / n; a2 = th1 + dth * (i + 1) / n
        alpha = math.sin(a2 - a1) * (math.sqrt(4 + 3 * math.tan((a2 - a1) / 2) ** 2) - 1) / 3
        p1 = (cx + rx * math.cos(a1) - alpha * rx * math.sin(a1), cy + ry * math.sin(a1) + alpha * ry * math.cos(a1))
        p2 = (cx + rx * math.cos(a2) + alpha * rx * math.sin(a2), cy + ry * math.sin(a2) - alpha * ry * math.cos(a2))
        p3 = (cx + rx * math.cos(a2), cy + ry * math.sin(a2))
        segs.append(((x0, y0), p1, p2, p3))
        x0, y0 = p3
    return segs


def parse_path(d):
    seq = []
    for m in _TOKEN.finditer(d):
        seq.append(m.group(1) or float(m.group(2)))
    cmds = []
    cur = "M"; last = None; start = None; prev_c = None; prev_cmd = None
    pos = 0
    while pos < len(seq):
        tok = seq[pos]
        if isinstance(tok, str):
            cur = tok; pos += 1
            if cur in "Zz":
                cmds.append(("Z", [])); continue
        args = []
        while pos < len(seq) and isinstance(seq[pos], float):
            args.append(seq[pos]); pos += 1
        if cur in "Zz": continue
        C = cur.upper(); rel = cur.islower()
        if C == "M":
            for k in range(0, len(args), 2):
                x, y = args[k], args[k + 1]
                nx, ny = (x + last[0], y + last[1]) if rel and last else (x, y)
                cmds.append(("M", [nx, ny])); last = (nx, ny); start = (nx, ny)
                prev_cmd = "M"
        elif C == "L":
            for k in range(0, len(args), 2):
                x, y = args[k], args[k + 1]
                nx, ny = (x + last[0], y + last[1]) if rel else (x, y)
                cmds.append(("L", [nx, ny])); last = (nx, ny); prev_cmd = "L"
        elif C == "H":
            for x in args:
                nx = x + last[0] if rel else x
                cmds.append(("L", [nx, last[1]])); last = (nx, last[1]); prev_cmd = "L"
        elif C == "V":
            for y in args:
                ny = y + last[1] if rel else y
                cmds.append(("L", [last[0], ny])); last = (last[0], ny); prev_cmd = "L"
        elif C == "C":
            for k in range(0, len(args), 6):
                x1, y1, x2, y2, x, y = args[k:k + 6]
                if rel:
                    x1 += last[0]; y1 += last[1]; x2 += last[0]; y2 += last[1]; x += last[0]; y += last[1]
                cmds.append(("C", [x1, y1, x2, y2, x, y])); prev_c = (x2, y2); last = (x, y); prev_cmd = "C"
        elif C == "S":
            for k in range(0, len(args), 4):
                x2, y2, x, y = args[k:k + 4]
                if rel:
                    x2 += last[0]; y2 += last[1]; x += last[0]; y += last[1]
                if prev_cmd in ("C", "S"):
                    x1, y1 = 2 * last[0] - prev_c[0], 2 * last[1] - prev_c[1]
                else:
                    x1, y1 = last
                cmds.append(("C", [x1, y1, x2, y2, x, y])); prev_c = (x2, y2); last = (x, y); prev_cmd = "S"
        elif C == "Q":
            for k in range(0, len(args), 4):
                qx, qy, x, y = args[k:k + 4]
                if rel:
                    qx += last[0]; qy += last[1]; x += last[0]; y += last[1]
                x1 = last[0] + 2 / 3 * (qx - last[0]); y1 = last[1] + 2 / 3 * (qy - last[1])
                x2 = x + 2 / 3 * (qx - x); y2 = y + 2 / 3 * (qy - y)
                cmds.append(("C", [x1, y1, x2, y2, x, y])); prev_c = (qx, qy); last = (x, y); prev_cmd = "Q"
        elif C == "T":
            for k in range(0, len(args), 2):
                x, y = args[k], args[k + 1]
                if rel:
                    x += last[0]; y += last[1]
                if prev_cmd in ("Q", "T"):
                    qx, qy = 2 * last[0] - prev_c[0], 2 * last[1] - prev_c[1]
                else:
                    qx, qy = last
                x1 = last[0] + 2 / 3 * (qx - last[0]); y1 = last[1] + 2 / 3 * (qy - last[1])
                x2 = x + 2 / 3 * (qx - x); y2 = y + 2 / 3 * (qy - y)
                cmds.append(("C", [x1, y1, x2, y2, x, y])); prev_c = (qx, qy); last = (x, y); prev_cmd = "T"
        elif C == "A":
            for k in range(0, len(args), 7):
                rx, ry, rot, la, sw, x, y = args[k:k + 7]
                nx, ny = (x + last[0], y + last[1]) if rel else (x, y)
                cmds.append(("A", [rx, ry, rot, int(la), int(sw), nx, ny])); last = (nx, ny); prev_cmd = "A"
    return cmds, start


def to_beziers(cmds):
    """→ list of subpaths; each subpath = list of cubic segs ((p0,c1,c2,p1))."""
    subs, cur, start, last = [], [], None, None
    for cmd, args in cmds:
        if cmd == "M":
            if cur: subs.append(cur)
            cur = []; last = start = tuple(args)
        elif cmd == "L":
            if last: cur.append((last, last, tuple(args), tuple(args))); last = tuple(args)
        elif cmd == "C":
            if last:
                cur.append((last, (args[0], args[1]), (args[2], args[3]), (args[4], args[5]))); last = (args[4], args[5])
        elif cmd == "A":
            if last:
                for seg in arc_to_cubics(last[0], last[1], *args):
                    cur.append(seg)
                last = (args[5], args[6])
        elif cmd == "Z":
            if last and start and last != start:
                cur.append((last, last, start, start))
            last = start
            if cur: subs.append(cur)
            cur = []
    if cur: subs.append(cur)
    return subs


def flatten_bezier(seg, n=12):
    (x0, y0), (x1, y1), (x2, y2), (x3, y3) = seg
    out = []
    for i in range(n + 1):
        t = i / n; mt = 1 - t
        out.append((mt ** 3 * x0 + 3 * mt * mt * t * x1 + 3 * mt * t * t * x2 + t ** 3 * x3,
                    mt ** 3 * y0 + 3 * mt * mt * t * y1 + 3 * mt * t * t * y2 + t ** 3 * y3))
    return out


def normalize_hex(color):
    if not color:
        return None
    h = color.strip().lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return None
    return "#" + h


# ─────────────────────────── compiler ───────────────────────────

class ProbeCompiler:
    def __init__(self, slide, target_rect, vb):
        self.slide = slide
        self.rect = target_rect
        self.vb = vb
        self.grads = {}
        self.clips = {}
        self.warnings = []
        self.features = set()
        self.shape_count = 0

    def to_inches(self, x, y):
        lx, ly, w, h = self.rect
        vx, vy, vw, vh = self.vb
        s = min(w / vw, h / vh)
        ox = lx + (w - vw * s) / 2; oy = ly + (h - vh * s) / 2
        return ox + (x - vx) * s, oy + (y - vy) * s

    def _add_native(self, x, y, w, h, fill, fill_alpha, stroke, sw):
        sh = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        fill = normalize_hex(fill)
        if fill and fill != "none":
            sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fill.lstrip("#"))
            if fill_alpha < 100: self._apply_alpha(sh._element, fill_alpha)
        else:
            sh.fill.background()
        stroke = normalize_hex(stroke)
        if stroke and stroke != "none":
            sh.line.color.rgb = RGBColor.from_string(stroke.lstrip("#"))
            sh.line.width = Pt(max(sw, 0.5))
        else:
            sh.line.fill.background()
        self.shape_count += 1
        return sh

    def _add_freeform(self, local_subs, minx, miny, w, h, fill, fill_alpha, stroke, sw):
        b = FreeformBuilder()
        for sub in local_subs:
            b.move_to(sub[0][0], sub[0][1])
            for p in sub[1:]:
                b.line_to(p[0], p[1])
            b.close()
        elem = b.build(self.slide, minx, miny, w, h, no_fill=(not fill), fill_color=normalize_hex(fill),
                       line_color=normalize_hex(stroke), line_width_pt=max(sw, 0.0))
        if fill_alpha < 100:
            self._apply_alpha(elem, fill_alpha)
        self.shape_count += 1
        return elem

    def _apply_alpha(self, elem, alpha):
        spPr = self._find_spPr(elem)
        if spPr is None: return
        sf = spPr.find(f"{A}solidFill")
        if sf is None: return
        clr = sf.find(f"{A}srgbClr")
        if clr is None: return
        e = etree.SubElement(clr, f"{A}alpha"); e.set("val", str(int(alpha * 1000)))

    @staticmethod
    def _find_spPr(elem):
        P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        spPr = elem.find(f".//{A}spPr")
        if spPr is None:
            spPr = elem.find(f".//{P}spPr")
        return spPr

    def _apply_gradient(self, elem, grad):
        spPr = self._find_spPr(elem)
        if spPr is None: return
        for sf in spPr.findall(f"{A}solidFill"):
            spPr.remove(sf)
        for ln in spPr.findall(f"{A}ln"):
            for sf in ln.findall(f"{A}solidFill"):
                if sf.get("noFill") is None and sf.find(f"{A}noFill") is not None:
                    pass
        gradFill = etree.SubElement(spPr, f"{A}gradFill")
        gsLst = etree.SubElement(gradFill, f"{A}gsLst")
        for pos, col, op in grad["stops"]:
            gs = etree.SubElement(gsLst, f"{A}gs"); gs.set("pos", str(int(pos * 100000)))
            clr = etree.SubElement(gs, f"{A}srgbClr"); clr.set("val", col.lstrip("#"))
            if op < 1.0:
                e = etree.SubElement(clr, f"{A}alpha"); e.set("val", str(int(op * 100000)))
        lin = etree.SubElement(gradFill, f"{A}lin")
        lin.set("ang", str(int(math.degrees(math.atan2(grad["dy"], grad["dx"])) * 60000)))
        lin.set("scaled", "1")

    def _paint(self, el, which):
        v = el.get(which)
        op = el.get(which + "-opacity")
        alpha = int(float(op) * 100) if op else 100
        if v is None:
            return "none", None, alpha
        if v.startswith("url(#"):
            gid = v[v.index("#") + 1:-1]
            if gid in self.grads:
                self.features.add("gradient")
                return "grad", self.grads[gid], alpha
            return "none", None, alpha
        if v == "none":
            return "none", None, alpha
        return "solid", v, alpha

    def collect_defs(self, root):
        for g in root.iter(SVG + "linearGradient"):
            stops = []
            for s in g.iter(SVG + "stop"):
                off = s.get("offset", "0")
                if off.endswith("%"):
                    pos = float(off.rstrip("%")) / 100
                else:
                    pos = float(off)
                col = normalize_hex(s.get("stop-color", "#000000")) or "#000000"
                op = float(s.get("stop-opacity", "1"))
                stops.append((pos, col, op))
            x1, y1 = float(g.get("x1", "0")), float(g.get("y1", "0"))
            x2, y2 = float(g.get("x2", "1")), float(g.get("y2", "1"))
            self.grads[g.get("id")] = {"stops": stops, "dx": x2 - x1, "dy": y2 - y1}
        for c in root.iter(SVG + "clipPath"):
            polys = []
            for child in c:
                poly = self._svg_polygon(child, Affine())
                if poly:
                    polys.extend(poly)
            self.clips[c.get("id")] = polys

    def _svg_polygon(self, el, tf):
        """SVG-space polygon points (no viewBox mapping)."""
        tag = el.tag.split("}")[-1]
        if tag == "rect":
            x, y = float(el.get("x", 0)), float(el.get("y", 0))
            w, h = float(el.get("width")), float(el.get("height"))
            pts = [(x, y), (x + w, y), (x + w, y + h), (x, y + h)]
        elif tag == "circle":
            cx, cy, r = float(el.get("cx", 0)), float(el.get("cy", 0)), float(el.get("r", 0))
            pts = [(cx + r * math.cos(2 * math.pi * i / 64), cy + r * math.sin(2 * math.pi * i / 64)) for i in range(64)]
        elif tag == "ellipse":
            cx, cy = float(el.get("cx", 0)), float(el.get("cy", 0))
            rx, ry = float(el.get("rx", 0)), float(el.get("ry", 0))
            pts = [(cx + rx * math.cos(2 * math.pi * i / 64), cy + ry * math.sin(2 * math.pi * i / 64)) for i in range(64)]
        elif tag in ("polygon", "polyline"):
            flat = [float(v) for v in el.get("points", "").replace(",", " ").split()]
            pts = list(zip(flat[0::2], flat[1::2]))
        elif tag == "line":
            pts = [(float(el.get("x1")), float(el.get("y1"))), (float(el.get("x2")), float(el.get("y2")))]
        elif tag == "path":
            subs = to_beziers(parse_path(el.get("d", ""))[0])
            out = []
            for sub in subs:
                pts = []
                for seg in sub:
                    pts.extend(flatten_bezier(seg, 16))
                out.append([tf.apply(px, py) for px, py in pts])
            return out
        else:
            return []
        return [[tf.apply(px, py) for px, py in pts]]

    def compile(self, svg_text):
        root = etree.fromstring(svg_text.encode("utf-8"))
        self.collect_defs(root)
        self._walk(root, Affine(), [])
        return self.warnings, self.features, self.shape_count

    def _walk(self, el, tf, clip_stack):
        tag = el.tag.split("}")[-1] if el.tag else ""
        tf = tf.compose(parse_transform(el.get("transform")))
        c = el.get("clip-path")
        if c and c.startswith("url(#"):
            self.features.add("clipPath")
            gid = c[c.index("#") + 1:-1]
            polys = self.clips.get(gid, [])
            clip_stack = clip_stack + polys
        if tag == "g":
            for ch in el:
                self._walk(ch, tf, clip_stack)
        elif tag == "svg":
            for ch in el:
                self._walk(ch, tf, clip_stack)
        elif tag == "image":
            self.features.add("image")
            self.warnings.append("image element: UNSUPPORTED (would degrade to picture, refusing)")
        elif tag == "filter" or tag == "mask":
            self.features.add(tag)
            self.warnings.append(f"{tag} element: UNSUPPORTED (decorative effect, refusing)")
        elif tag == "text":
            self._render_text(el, tf)
        elif tag in ("rect", "circle", "ellipse", "polygon", "polyline", "line", "path"):
            self._render_shape(el, tag, tf, clip_stack)

    def _render_shape(self, el, tag, tf, clip_stack):
        self.features.add(tag)
        for bad_attr in ("filter", "mask"):
            if el.get(bad_attr):
                self.warnings.append(
                    f"{tag} has {bad_attr}=<{el.get(bad_attr)}>: UNSUPPORTED (refusing, no silent degrade)")
        fkind, fval, fa = self._paint(el, "fill")
        skind, sval, _ = self._paint(el, "stroke")
        sw = float(el.get("stroke-width", "1"))
        subpaths = self._svg_polygon(el, tf)
        if not subpaths:
            self.warnings.append(f"{tag}: empty geometry"); return
        allp = [p for sub in subpaths for p in sub]
        minx = min(p[0] for p in allp); maxx = max(p[0] for p in allp)
        miny = min(p[1] for p in allp); maxy = max(p[1] for p in allp)
        # map to inches
        corners_in = [self.to_inches(*p) for p in [(minx, miny), (maxx, maxy)]]
        ix0, iy0 = corners_in[0]; ix1, iy1 = corners_in[1]
        iw, ih = ix1 - ix0, iy1 - iy0
        # rebuild local subpaths in inches, relative to (ix0, iy0)
        local = []
        for sub in subpaths:
            in_pts = [self.to_inches(px, py) for px, py in sub]
            local.append([(px - ix0, py - iy0) for px, py in in_pts])

        needs_bool = bool(clip_stack) or el.get("fill-rule") == "evenodd"
        if needs_bool:
            from shapely.geometry import Polygon as ShapelyPoly
            from shapely.validation import make_valid
            poly = None
            for sub in subpaths:
                pts = [self.to_inches(px, py) for px, py in sub]
                if len(pts) >= 3:
                    try:
                        p = ShapelyPoly(pts).buffer(0)
                        p = make_valid(p)
                    except Exception:
                        p = None
                    if p is not None and not p.is_empty:
                        poly = p if poly is None else poly.union(p)
            for cpoly in clip_stack:
                cpts = [self.to_inches(px, py) for px, py in cpoly]
                if len(cpts) >= 3:
                    try:
                        clip_poly = make_valid(ShapelyPoly(cpts).buffer(0))
                        poly = poly.intersection(clip_poly)
                    except Exception:
                        self.warnings.append(f"{tag}: clip failed")
                        poly = None
                    if poly is None or poly.is_empty:
                        self.warnings.append(f"{tag}: clipped to empty")
                        return
            if poly is None or poly.is_empty:
                self.warnings.append(f"{tag}: boolean empty")
                return
            if fkind == "grad":
                geom = poly
                from shapely.geometry import mapping
                m = mapping(geom)
                self._add_boolean_with_grad(geom, ix0, iy0, iw, ih, fval, fa, sval if sval != "none" else None)
            else:
                self._add_boolean(poly, ix0, iy0, iw, ih, fval, fa, sval if sval != "none" else None)
            return

        if tag == "rect" and not el.get("rx") and fkind == "solid" and sval in (None, "none") and not clip_stack:
            self._add_native(ix0, iy0, iw, ih, fval, fa, None, sw)
            return
        if fkind == "solid":
            fill_hex = normalize_hex(fval)
        else:
            fill_hex = None
        elem = self._add_freeform(local, ix0, iy0, iw, ih,
                                  fill_hex if fkind == "solid" else "#FFFFFF",
                                  fa, sval if sval != "none" else None, sw)
        if fkind == "grad":
            self._apply_gradient(elem, fval)

    def _add_boolean(self, poly, x, y, w, h, fill, fill_alpha, stroke):
        elem = bool_shape(poly, self.slide, x, y, w, h, fill=normalize_hex(fill), line=normalize_hex(stroke), alpha=fill_alpha)
        self.shape_count += 1
        return elem

    def _add_boolean_with_grad(self, poly, x, y, w, h, grad, fill_alpha, stroke):
        elem = bool_shape(poly, self.slide, x, y, w, h, fill="#FFFFFF", line=normalize_hex(stroke), alpha=fill_alpha)
        self._apply_gradient(elem, grad)
        return elem

    def _render_text(self, el, tf):
        self.features.add("text")
        x, y = float(el.get("x", 0)), float(el.get("y", 0))
        ix, iy = self.to_inches(*tf.apply(x, y))
        content = "".join(el.itertext()).strip()
        if not content: return
        fs = float(el.get("font-size", "14").replace("px", ""))
        anchor = el.get("text-anchor", "start")
        fkind, fval, _ = self._paint(el, "fill")
        if fkind == "solid":
            fval = normalize_hex(fval) or "#000000"
        else:
            fval = "#000000"
        w_est, h_est = estimate_text_size(content, max(8, int(fs)), 8.0)
        tb = self.slide.shapes.add_textbox(Inches(ix - 4.0), Inches(iy - h_est - fs / 72), Inches(8.0), Inches(h_est * 2))
        tf_el = tb.text_frame
        tf_el.word_wrap = False
        tf_el.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf_el.paragraphs[0]
        p.alignment = { "middle": PP_ALIGN.CENTER, "end": PP_ALIGN.RIGHT }.get(anchor, PP_ALIGN.LEFT)
        run = p.add_run(); run.text = content
        run.font.size = Pt(fs)
        if fval and fval != "none":
            run.font.color.rgb = RGBColor.from_string(fval.lstrip("#"))
        self.shape_count += 1


# ─────────────────────────── rendering ───────────────────────────

EDGE = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
EDGE_PROFILE = str(Path(os.environ.get("TEMP", ".")) / "edge_probe_profile")


def _kill_edge():
    try:
        subprocess.run(["taskkill", "/F", "/IM", "msedge.exe"], capture_output=True)
    except Exception:
        pass


def render_svg_edge(svg_path, png_path, w_px, h_px):
    if png_path.exists():
        png_path.unlink()
    _kill_edge()
    cmd = [EDGE, "--headless", "--disable-gpu", "--no-first-run",
           f"--user-data-dir={EDGE_PROFILE}", f"--screenshot={png_path}",
           f"--window-size={w_px},{h_px}", "--default-background-color=FFFFFFFF",
           f"file:///{svg_path.as_posix()}"]
    subprocess.run(cmd, capture_output=True, timeout=90)
    for _ in range(60):
        if png_path.exists():
            return png_path
        time.sleep(0.5)
    raise FileNotFoundError(f"Edge did not produce {png_path}")


def render_pptx(pptx_path):
    from ppt_pro_max import build_helpers
    ppt_out = OUT / "ppt"
    ppt_out.mkdir(parents=True, exist_ok=True)
    res = build_helpers.preview(str(pptx_path), out_dir=str(ppt_out), engine="libreoffice", title="probe")
    return Path(res["pngs"][0])


def pixel_compare(gt_png, res_png, rect, dpi=110):
    l, t, w, h = rect
    gt = np.array(Image.open(gt_png).convert("RGB"), dtype=float)
    res = np.array(Image.open(res_png).convert("RGB"), dtype=float)
    x0, y0 = int(l * dpi), int(t * dpi)
    x1, y1 = int((l + w) * dpi), int((t + h) * dpi)
    crop = res[y0:y1, x0:x1]
    hh = min(gt.shape[0], crop.shape[0]); ww = min(gt.shape[1], crop.shape[1])
    gt_c, res_c = gt[:hh, :ww], crop[:hh, :ww]
    diff = np.abs(gt_c - res_c).mean()
    gt_ink = gt_c.mean(axis=2) < 240
    res_ink = res_c.mean(axis=2) < 240
    iou = (gt_ink & res_ink).sum() / ((gt_ink | res_ink).sum() + 1e-6)
    return diff, iou


# ─────────────────────────── main ───────────────────────────

SLIDE_RECTS = {
    "pyramid": (3.5, 0.8, 6.3, 5.6),
    "venn_evenodd": (2.5, 1.2, 8.0, 5.2),
    "funnel": (4.0, 0.8, 5.4, 6.0),
    "growth_curve": (2.5, 1.0, 8.0, 5.2),
    "matrix_bcg": (3.0, 0.9, 7.0, 5.7),
    "unsupported": (3.0, 0.9, 7.0, 5.7),
}
DPI = 110


def run_probe():
    OUT.mkdir(parents=True, exist_ok=True)
    CASES.mkdir(parents=True, exist_ok=True)
    report = ["# SVG → editable-pptx probe report", ""]
    total_compile = 0.0

    for name, svg in CASES_SVG.items():
        rect = SLIDE_RECTS[name]
        vb = (0, 0, 400, 300)
        m = re.search(r'viewBox="([^"]+)"', svg)
        if m:
            vals = [float(v) for v in m.group(1).replace(",", " ").split()]
            vb = (vals[0], vals[1], vals[2], vals[3])

        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        t0 = time.perf_counter()
        compiler = ProbeCompiler(slide, rect, vb)
        warnings, features, n_shapes = compiler.compile(svg)
        t_compile = time.perf_counter() - t0
        total_compile += t_compile

        pptx_path = OUT / f"{name}.pptx"
        prs.save(str(pptx_path))

        with zipfile.ZipFile(pptx_path) as z:
            slide_xml = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            xml = z.read(slide_xml[0]).decode("utf-8")
        has_pic = "<p:pic>" in xml
        has_sp = xml.count("<p:sp>")
        has_grp = "<p:grpSp>" in xml
        editable = (not has_pic) and has_sp > 0

        svg_path = CASES / f"{name}.svg"
        svg_path.write_text(svg, encoding="utf-8")
        gt_dir = OUT / "gt"
        gt_dir.mkdir(parents=True, exist_ok=True)
        gt_png = gt_dir / f"{name}_gt.png"
        render_svg_edge(svg_path, gt_png, int(rect[2] * DPI), int(rect[3] * DPI))
        res_png = render_pptx(pptx_path)
        diff, iou = pixel_compare(gt_png, res_png, rect, DPI)

        report.append(f"## {name}")
        report.append(f"- compile: {t_compile*1000:.0f} ms | {n_shapes} shapes")
        report.append(f"- editable: {editable} | pictures={has_pic} p:sp={has_sp} groups={has_grp}")
        report.append(f"- features: {', '.join(sorted(features))}")
        report.append(f"- pixel diff(mean abs)={diff:.1f}  ink-IoU={iou:.2f}")
        if warnings:
            report.append(f"- boundaries: {warnings}")
        report.append("")

    report.insert(1, f"Total compile (all cases): {total_compile*1000:.0f} ms")
    out_md = OUT / "REPORT.md"
    out_md.write_text("\n".join(report), encoding="utf-8")
    print("\n".join(report))
    return out_md


if __name__ == "__main__":
    run_probe()
