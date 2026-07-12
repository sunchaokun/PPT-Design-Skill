"""EnterprisePipeline — main orchestration flow."""

from __future__ import annotations

import os
from typing import Any

from ppt_pro_max.enterprise.scanner import ProjectScanner
from ppt_pro_max.enterprise.template_analyzer import TemplateAnalyzer
from ppt_pro_max.enterprise.brand_spec import BrandSpec
from ppt_pro_max.enterprise.version_manager import next_version, read_meta, write_meta
from ppt_pro_max.enterprise.enterprise_decider import EnterpriseDesignDecider
from ppt_pro_max.enterprise.enterprise_renderer import EnterpriseRenderer
from ppt_pro_max.enterprise.review_gate import ReviewGate
from ppt_pro_max.enterprise.content_parser import load_enterprise_content
from ppt_pro_max.enterprise.image_matcher import match_images
from ppt_pro_max.enterprise.page_revision import PageRevisionEngine
from ppt_pro_max.enterprise.density_profile import get_density_profile, apply_density_to_bullets


class EnterprisePipeline:

    def run(
        self,
        query: str,
        project_dir: str,
        dry_run: bool = False,
        business_mode: str | None = None,
        density: int | None = None,
        motion: int | None = None,
        review: bool = False,
        review_file: str | None = None,
        output_version: int | None = None,
        from_version: int | None = None,
        pages: str | None = None,
        history: bool = False,
        output: str | None = None,
        content_file: str | None = None,
        **kwargs,
    ) -> dict[str, Any]:
        if history:
            return self._handle_history(project_dir)

        scanner = ProjectScanner()
        asset = scanner.scan(project_dir)

        brand_spec = self._build_brand_spec(asset)

        if dry_run:
            return {
                "dry_run": True,
                "pipeline": "enterprise",
                "project": project_dir,
                "brand_spec": {
                    "source": brand_spec.source,
                    "colors": brand_spec.colors,
                    "fonts": brand_spec.fonts,
                },
                "assets": {
                    "template": asset.template_path is not None,
                    "logo": asset.logo_path is not None,
                    "brand_json": asset.brand_raw is not None,
                    "content_json": asset.content_raw is not None,
                    "image_pool_count": len(asset.image_pool),
                },
            }

        decider = EnterpriseDesignDecider(
            brand_spec=brand_spec,
            business_mode=business_mode,
        )
        effective_density = density if density is not None else decider.suggest_density()

        page_contents = self._load_content(asset, project_dir, content_file)

        story_plan = self._build_story_plan(query, page_contents, effective_density, business_mode)

        page_designs = self._build_page_designs(
            story_plan, decider, asset, brand_spec
        )

        if review:
            gate = ReviewGate()
            proposal = gate.generate_proposal(
                project_dir=project_dir,
                brand_spec={"source": brand_spec.source, "colors": brand_spec.colors, "fonts": brand_spec.fonts},
                story_plan=story_plan,
                page_designs=page_designs,
                assets={
                    "template": asset.template_path is not None,
                    "logo": asset.logo_path is not None,
                    "brand_json": asset.brand_raw is not None,
                    "content_json": asset.content_raw is not None,
                    "image_pool_count": len(asset.image_pool),
                },
            )
            proposal_path = review_file or os.path.join(project_dir, "output", "proposal.json")
            gate.write_proposal(proposal, proposal_path)
            return {
                "pipeline": "enterprise",
                "review": True,
                "proposal_path": proposal_path,
                "proposal": proposal,
            }

        output_dir = os.path.join(project_dir, "output")
        from_meta = None
        if output_version is not None:
            vnum = output_version
        elif from_version is not None:
            vnum = from_version
            from_meta = read_meta(os.path.join(output_dir, f"v{from_version}"))
        else:
            vnum = next_version(output_dir)

        if from_meta and from_meta.get("slides"):
            page_contents = self._load_content_from_meta(from_meta)

        version_dir = os.path.join(output_dir, f"v{vnum}")
        pptx_path = output or os.path.join(version_dir, "presentation.pptx")

        if pages:
            source_pptx = self._find_latest_pptx(output_dir, vnum)
            if not source_pptx:
                source_pptx = asset.template_path
            if not source_pptx:
                return {
                    "pipeline": "enterprise",
                    "error": "--pages requires a template.pptx or existing version in the project directory",
                }
            result = self._handle_page_revision(
                source_pptx, pages, pptx_path, vnum, version_dir
            )
            return result

        renderer = EnterpriseRenderer(template_path=asset.template_path)
        prs = renderer.create_presentation()

        page_designs = match_images(asset.image_pool, page_designs)
        density_profile = get_density_profile(effective_density)

        image_fetcher = self._build_image_fetcher(kwargs)

        for design in page_designs:
            image_val = design.get("image")
            if image_val and not os.path.isfile(image_val):
                image_val = None
                design["image"] = None
            if not image_val and image_fetcher is not None:
                keywords = design.get("title", query)
                goal = design.get("goal", "")
                try:
                    fetched = image_fetcher.fetch(keywords, goal=goal)
                    if fetched:
                        design["image"] = fetched
                except Exception:
                    pass

            layout_name = design.get("template_layout_name")
            slide = renderer.add_slide(prs, layout_name=layout_name)
            self._populate_slide(slide, design, prs, density_profile, brand_spec)

            if asset.logo_path:
                logo_spec = brand_spec.logo or {"position": "top_right", "width_inches": 1.0}
                renderer.insert_logo(
                    slide, asset.logo_path, logo_spec,
                    current_goal=design.get("goal"), prs=prs,
                )

            effective_motion = motion
            if effective_motion is None:
                effective_motion = kwargs.get("motion")
            if effective_motion and isinstance(effective_motion, int) and effective_motion > 0:
                self._apply_animations(slide, design, effective_motion)

        if brand_spec.footer:
            renderer.add_page_numbers(prs, brand_spec.footer, brand_spec=brand_spec)

        if brand_spec.watermark:
            renderer.add_watermark(prs, brand_spec.watermark, brand_spec=brand_spec)

        renderer.save(prs, pptx_path)

        meta = {
            "version": vnum,
            "query": query,
            "business_mode": business_mode,
            "density": effective_density,
            "num_slides": len(page_designs),
            "brand_source": brand_spec.source,
            "slides": [
                {"goal": d.get("goal"), "title": d.get("title", "")}
                for d in page_designs
            ],
        }
        write_meta(version_dir, meta)

        return {
            "pipeline": "enterprise",
            "project": project_dir,
            "output_path": pptx_path,
            "version": vnum,
            "num_slides": len(page_designs),
        }

    def _build_brand_spec(self, asset) -> BrandSpec:
        brand_spec = BrandSpec()
        if asset.template_path:
            analyzer = TemplateAnalyzer()
            template_spec = analyzer.analyze(asset.template_path)
            if asset.brand_raw:
                brand_spec = BrandSpec.merge(template_spec, asset.brand_raw)
            else:
                brand_spec = template_spec
        elif asset.brand_raw:
            brand_spec = BrandSpec.from_brand_json(asset.brand_raw)
        return brand_spec

    def _handle_history(self, project_dir: str) -> dict[str, Any]:
        output_dir = os.path.join(project_dir, "output")
        versions: list[dict[str, Any]] = []

        if os.path.isdir(output_dir):
            for entry in sorted(os.listdir(output_dir)):
                if entry.startswith("v") and entry[1:].isdigit():
                    vnum = int(entry[1:])
                    meta = read_meta(os.path.join(output_dir, entry))
                    versions.append({
                        "version": vnum,
                        "meta": meta,
                    })

        return {"history": True, "project": project_dir, "versions": versions}

    def _load_content(self, asset, project_dir: str, content_file: str | None = None) -> list[dict[str, Any]]:
        if content_file and os.path.isfile(content_file):
            import json
            try:
                with open(content_file, encoding="utf-8") as f:
                    content_raw = json.load(f)
                return load_enterprise_content(content_raw, project_dir)
            except (json.JSONDecodeError, OSError):
                pass
        if asset.content_raw:
            return load_enterprise_content(asset.content_raw, project_dir)
        return []

    def _load_content_from_meta(self, meta: dict[str, Any]) -> list[dict[str, Any]]:
        return [
            {"goal": s.get("goal", "content"), "title": s.get("title", "")}
            for s in meta.get("slides", [])
        ]

    def _build_story_plan(
        self,
        query: str,
        page_contents: list[dict[str, Any]],
        density: int,
        business_mode: str | None = None,
    ) -> dict[str, Any]:
        if page_contents:
            return {
                "strategy": "content_json",
                "pages": page_contents,
            }

        from ppt_pro_max.planner.story_planner import StoryPlanner
        planner = StoryPlanner()
        strategy_name = self._map_business_mode_to_strategy(business_mode)
        target_pages = max(5, min(20, density + 2))
        plan = planner.plan(query, strategy_override=strategy_name, slide_count_override=target_pages)
        return {
            "strategy": plan.strategy,
            "pages": [
                {"goal": p.goal, "title": p.goal.replace("_", " ").title(), "notes": f"本页目标: {p.goal}。请结合标题展开讲解。"}
                for p in plan.pages
            ],
        }

    def _map_business_mode_to_strategy(self, business_mode: str | None) -> str | None:
        mapping = {
            "pitch": None,
            "education": "Education Course",
            "training": "Training Workshop",
            "report": "Business Report",
        }
        return mapping.get(business_mode)

    def _build_page_designs(
        self,
        story_plan: dict[str, Any],
        decider: EnterpriseDesignDecider,
        asset,
        brand_spec: BrandSpec,
    ) -> list[dict[str, Any]]:
        pages = story_plan.get("pages", [])
        designs: list[dict[str, Any]] = []

        layout_names: dict[int, str] = {}
        if asset.template_path:
            try:
                from pptx import Presentation
                tprs = Presentation(asset.template_path)
                for idx, layout in enumerate(tprs.slide_layouts):
                    layout_names[idx] = layout.name
            except Exception:
                pass

        for page in pages:
            goal = page.get("goal", "content")
            layout_idx = decider.resolve_layout_index(goal)
            design = {
                "goal": goal,
                "title": page.get("title", ""),
                "subtitle": page.get("subtitle"),
                "bullets": page.get("bullets"),
                "image": page.get("image"),
                "cards": page.get("cards"),
                "diagram_type": page.get("diagram_type"),
                "diagram_data": page.get("diagram_data"),
                "code": page.get("code"),
                "exercise": page.get("exercise"),
                "notes": page.get("notes"),
                "links": page.get("links"),
                "chart": page.get("chart"),
                "image_grid": page.get("image_grid"),
                "icons": page.get("icons"),
                "template_layout_index": layout_idx,
            }
            if layout_idx is not None and layout_idx in layout_names:
                design["template_layout_name"] = layout_names[layout_idx]
            designs.append(design)
        return designs

    def _populate_slide(self, slide, design: dict[str, Any], prs, density_profile=None, brand_spec=None) -> None:
        from pptx.util import Pt

        if density_profile is None:
            from ppt_pro_max.enterprise.density_profile import get_density_profile
            density_profile = get_density_profile(4)

        self._apply_visual_design(slide, design, prs, brand_spec)

        if slide.shapes.title and design.get("title"):
            slide.shapes.title.text = design["title"]
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(density_profile.title_size)
                    self._apply_brand_font_color(run, "foreground", brand_spec)

        if design.get("subtitle"):
            from pptx.enum.shapes import PP_PLACEHOLDER
            for ph in slide.placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.SUBTITLE or ph.placeholder_format.idx == 2:
                    ph.text = design["subtitle"]
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(density_profile.subtitle_size)
                            self._apply_brand_font_color(run, "muted-foreground", brand_spec)
                    break

        if design.get("bullets"):
            bullets = apply_density_to_bullets(design["bullets"], density_profile)
            body_text = "\n".join(f"• {b}" if not b.startswith(("• ", "- ", "— ")) else b for b in bullets)
            from pptx.enum.shapes import PP_PLACEHOLDER
            for ph in slide.placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) or ph.placeholder_format.idx == 1:
                    ph.text = body_text
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(density_profile.bullet_size)
                            self._apply_brand_font_color(run, "foreground", brand_spec)
                    break

        if design.get("image") and os.path.isfile(design["image"]):
            self._insert_content_image(slide, design["image"], prs, density_profile)

        if design.get("chart"):
            self._render_chart(slide, design["chart"], prs, density_profile, brand_spec)

        if design.get("cards"):
            self._render_cards(slide, design["cards"], prs, density_profile)

        if design.get("diagram_type") and design.get("diagram_data"):
            self._render_diagram(slide, design, prs, density_profile)

        if design.get("code"):
            self._render_code_block(slide, design["code"], prs, density_profile)

        if design.get("exercise"):
            self._render_exercise(slide, design["exercise"], prs, density_profile)

        if design.get("notes"):
            self._render_notes(slide, design["notes"])

        if design.get("links"):
            self._render_links(slide, design["links"], design, prs=prs, brand_spec=brand_spec)

    _BASELINE_IMAGE_RATIO = 0.38

    def _apply_visual_design(self, slide, design: dict[str, Any], prs, brand_spec=None) -> None:
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        colors = {}
        if brand_spec and brand_spec.colors:
            colors = brand_spec.colors

        bg_color = colors.get("background", "#FFFFFF")
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip("#"))
        except Exception:
            pass

        goal = design.get("goal", "content")
        is_hero = goal in ("hook", "cta")

        if is_hero:
            primary_hex = colors.get("primary", "#2563EB")
            try:
                bg_rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0), prs.slide_width, prs.slide_height,
                )
                bg_rect.fill.solid()
                bg_rect.fill.fore_color.rgb = RGBColor.from_string(primary_hex.lstrip("#"))
                bg_rect.line.fill.background()

                from pptx.oxml.ns import qn
                from lxml import etree
                sp_pr = bg_rect._element.find(qn("p:spPr"))
                solid_fill = sp_pr.find(qn("a:solidFill"))
                if solid_fill is not None:
                    srgb = solid_fill.find(qn("a:srgbClr"))
                    if srgb is not None:
                        alpha = etree.SubElement(srgb, qn("a:alpha"))
                        alpha.set("val", "90000")
            except Exception:
                pass
        else:
            accent_hex = colors.get("accent", colors.get("primary", "#2563EB"))
            try:
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0), Inches(0.08), prs.slide_height,
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = RGBColor.from_string(accent_hex.lstrip("#"))
                bar.line.fill.background()
            except Exception:
                pass

            muted_hex = colors.get("muted", "#F1F5F9")
            try:
                bottom_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(7.2), prs.slide_width, Inches(0.3),
                )
                bottom_bar.fill.solid()
                bottom_bar.fill.fore_color.rgb = RGBColor.from_string(muted_hex.lstrip("#"))
                bottom_bar.line.fill.background()
            except Exception:
                pass

    def _apply_brand_font_color(self, run, color_role: str, brand_spec=None) -> None:
        from pptx.dml.color import RGBColor

        if brand_spec is None or not brand_spec.colors:
            return
        hex_color = brand_spec.colors.get(color_role)
        if hex_color:
            try:
                run.font.color.rgb = RGBColor.from_string(hex_color.lstrip("#"))
            except Exception:
                pass

    def _insert_content_image(self, slide, image_path: str, prs, density_profile=None) -> None:
        from pptx.util import Inches
        from PIL import Image as PILImage

        if density_profile is None:
            from ppt_pro_max.enterprise.density_profile import get_density_profile
            density_profile = get_density_profile(4)

        slide_width = prs.slide_width
        target_width = Inches(4.0 * density_profile.image_width_ratio / self._BASELINE_IMAGE_RATIO)
        try:
            with PILImage.open(image_path) as img:
                img_w, img_h = img.size
                aspect = img_h / img_w if img_w > 0 else 0.75
        except Exception:
            aspect = 0.75
        target_height = int(target_width * aspect)

        left = slide_width - target_width - Inches(0.5)
        top = Inches(1.5)

        slide.shapes.add_picture(image_path, left, top, width=target_width, height=target_height)

    def _render_chart(self, slide, chart_config: dict[str, Any], prs, density_profile=None, brand_spec=None) -> None:
        from ppt_pro_max.renderer.chart_builder import ChartBuilder

        if density_profile is None:
            from ppt_pro_max.enterprise.density_profile import get_density_profile
            density_profile = get_density_profile(4)

        chart_height = 4.5 * density_profile.image_width_ratio / self._BASELINE_IMAGE_RATIO
        position = {
            "x": 1.5,
            "y": 1.5,
            "width": 10.333,
            "height": chart_height,
        }

        brand_colors = None
        if brand_spec and brand_spec.colors:
            brand_colors = brand_spec.colors

        builder = ChartBuilder()
        try:
            builder.build(slide, chart_config, position=position, brand_colors=brand_colors)
        except Exception:
            pass

    def _render_cards(self, slide, cards: list[dict[str, Any]], prs, density_profile=None) -> None:
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        if not cards:
            return

        if density_profile is None:
            from ppt_pro_max.enterprise.density_profile import get_density_profile
            density_profile = get_density_profile(4)

        slide_width = prs.slide_width
        n = len(cards)

        card_width = min(Inches(3.5), (slide_width - Inches(0.9 * 2 + 0.4 * (n - 1))) // n)
        card_height = Inches(3.0)
        gap = Inches(0.4)
        total_width = card_width * n + gap * (n - 1)
        start_left = (slide_width - total_width) // 2
        top = Inches(2.0)

        for i, card in enumerate(cards):
            left = start_left + i * (card_width + gap)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, card_width, card_height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
            shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            shape.line.width = Pt(1)

            tf = shape.text_frame
            tf.word_wrap = True

            title = card.get("title", "")
            body = card.get("body", "")
            text = f"{title}\n{body}" if body else title

            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(density_profile.bullet_size)

            if card.get("image") and os.path.isfile(card["image"]):
                img_left = left + Inches(0.2)
                img_top = top + card_height - Inches(0.8)
                img_width = Inches(0.6)
                try:
                    slide.shapes.add_picture(card["image"], img_left, img_top, width=img_width)
                except Exception:
                    pass

    def _render_diagram(self, slide, design: dict[str, Any], prs, density_profile=None) -> None:
        from ppt_pro_max.renderer.diagram_engine import DiagramEngine
        from ppt_pro_max.renderer.diagram.diagram_style import DiagramStyle
        from ppt_pro_max.renderer.diagram.layout_engine import Region

        if density_profile is None:
            from ppt_pro_max.enterprise.density_profile import get_density_profile
            density_profile = get_density_profile(4)

        diagram_type = design["diagram_type"]
        diagram_data = design["diagram_data"]
        if isinstance(diagram_data, dict) and "data" in diagram_data:
            diagram_data = diagram_data["data"]
        diagram_data = self._normalize_diagram_data(diagram_type, diagram_data)

        style = DiagramStyle()
        density_val = max(1, min(10, density_profile.title_size // 4))
        style = style.apply_density(density_val)

        region = Region(
            left=0.9,
            top=1.5,
            width=11.533,
            height=5.5,
        )

        engine = DiagramEngine()
        try:
            engine.render(slide, diagram_type, diagram_data, style, region)
        except Exception as exc:
            import logging
            logging.getLogger(__name__).warning("Diagram render failed for %s: %s", diagram_type, exc)

    def _normalize_diagram_data(self, diagram_type: str, data: dict[str, Any]) -> dict[str, Any]:
        if not isinstance(data, dict):
            return data
        if diagram_type == "funnel" and "items" in data and "stages" not in data:
            data = dict(data, stages=[
                {"label": item.get("text", ""), "fill_role": item.get("fill_role")}
                for item in data["items"]
            ])
        elif diagram_type == "timeline" and "items" in data and "events" not in data:
            data = dict(data, events=[
                {"label": item.get("text", ""), "position": item.get("position", "middle")}
                for item in data["items"]
            ])
        elif diagram_type == "swot" and "strengths" in data and "quadrants" not in data:
            data = dict(data, quadrants=[
                {"label": "Strengths", "items": data.get("strengths", [])},
                {"label": "Weaknesses", "items": data.get("weaknesses", [])},
                {"label": "Opportunities", "items": data.get("opportunities", [])},
                {"label": "Threats", "items": data.get("threats", [])},
            ])
        elif diagram_type == "pyramid" and "items" in data and "levels" not in data and "stages" not in data:
            data = dict(data, levels=[
                {"label": item.get("text", "")} for item in data["items"]
            ])
        if "nodes" in data:
            normalized_nodes = []
            for node in data["nodes"]:
                if "text" in node and "label" not in node:
                    node = dict(node, label=node["text"])
                normalized_nodes.append(node)
            data = dict(data, nodes=normalized_nodes)
        return data

    def _render_code_block(self, slide, code_data, prs, density_profile=None) -> None:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        if density_profile is None:
            from ppt_pro_max.enterprise.density_profile import get_density_profile
            density_profile = get_density_profile(4)

        code_text = code_data if isinstance(code_data, str) else code_data.get("source", code_data.get("code", ""))
        language = code_data.get("language", "") if isinstance(code_data, dict) else ""

        left = Inches(0.9)
        top = Inches(1.5)
        width = Inches(11.533)
        height = Inches(5.0)

        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
        bg_shape.line.fill.background()

        tf = bg_shape.text_frame
        tf.word_wrap = True

        header = f"  {language}" if language else ""
        if header:
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            p.font.bold = True

        lines = code_text.split("\n")
        for i, line in enumerate(lines[:30]):
            if i == 0 and not header:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"  {line}"
            p.font.size = Pt(density_profile.bullet_size - 2)
            p.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)
            p.font.name = "Consolas"

    def _render_exercise(self, slide, exercise_data, prs, density_profile=None) -> None:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        if density_profile is None:
            from ppt_pro_max.enterprise.density_profile import get_density_profile
            density_profile = get_density_profile(4)

        instructions = exercise_data.get("instructions", "") if isinstance(exercise_data, dict) else str(exercise_data)
        duration = exercise_data.get("duration", "") if isinstance(exercise_data, dict) else ""
        steps = exercise_data.get("steps", []) if isinstance(exercise_data, dict) else []

        badge_left = Inches(0.9)
        badge_top = Inches(1.2)
        badge_w = Inches(1.8) if duration else Inches(1.2)
        badge_h = Inches(0.4)

        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, badge_left, badge_top, badge_w, badge_h,
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = f"Exercise {duration}" if duration else "Exercise"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if instructions:
            inst_left = Inches(0.9)
            inst_top = Inches(1.8)
            inst_w = Inches(11.533)
            inst_h = Inches(1.2)

            txBox = slide.shapes.add_textbox(inst_left, inst_top, inst_w, inst_h)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = instructions
            p.font.size = Pt(density_profile.bullet_size)
            p.font.italic = True

        if steps:
            steps_left = Inches(0.9)
            steps_top = Inches(3.2)
            steps_w = Inches(11.533)
            steps_h = Inches(3.5)

            txBox = slide.shapes.add_textbox(steps_left, steps_top, steps_w, steps_h)
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, step in enumerate(steps):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"{i + 1}. {step}"
                p.font.size = Pt(density_profile.bullet_size)
                p.space_after = Pt(4)

    def _render_notes(self, slide, notes_text: str) -> None:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    def _render_links(self, slide, links: list[dict[str, Any]], design: dict[str, Any], prs=None, brand_spec=None) -> None:
        from pptx.dml.color import RGBColor

        accent_color = RGBColor(0x25, 0x63, 0xEB)
        if brand_spec and brand_spec.colors:
            accent_hex = brand_spec.colors.get("accent") or brand_spec.colors.get("primary")
            if accent_hex:
                accent_color = RGBColor.from_string(accent_hex.lstrip("#"))

        for link in links:
            bullet_index = link.get("bullet_index")
            url = link.get("url", "")
            text = link.get("text", "")

            if bullet_index is not None:
                self._attach_link_to_bullet(slide, bullet_index, url, accent_color, prs=prs)
            elif text and url:
                self._add_standalone_link(slide, link, accent_color, prs=prs)

    def _attach_link_to_bullet(self, slide, bullet_index: int, url: str, accent_color, prs=None) -> None:
        from pptx.enum.shapes import PP_PLACEHOLDER
        for ph in slide.placeholders:
            ph_type = ph.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) or ph.placeholder_format.idx == 1:
                paras = list(ph.text_frame.paragraphs)
                if bullet_index < len(paras):
                    for run in paras[bullet_index].runs:
                        if url.startswith("slide://"):
                            slide_num = int(url.replace("slide://", ""))
                            if prs is not None and 1 <= slide_num <= len(prs.slides):
                                try:
                                    run.hyperlink.target_slide = prs.slides[slide_num - 1]
                                except Exception:
                                    pass
                        else:
                            run.hyperlink.address = url
                        run.font.color.rgb = accent_color
                        run.font.underline = True
                break

    def _add_standalone_link(self, slide, link: dict[str, Any], accent_color, prs=None) -> None:
        from pptx.util import Inches, Pt

        text = link.get("text", "Link")
        url = link.get("url", "")
        position = link.get("position", "bottom_right")

        position_map = {
            "bottom_right": (Inches(11.433), Inches(6.8)),
            "bottom_left": (Inches(0.9), Inches(6.8)),
            "bottom_center": (Inches(5.5), Inches(6.8)),
        }
        left, top = position_map.get(position, (Inches(11.433), Inches(6.8)))

        txBox = slide.shapes.add_textbox(left, top, Inches(2.0), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(10)
        run.font.color.rgb = accent_color
        run.font.underline = True

        if url.startswith("slide://"):
            slide_num = int(url.replace("slide://", ""))
            if prs is not None and 1 <= slide_num <= len(prs.slides):
                try:
                    run.hyperlink.target_slide = prs.slides[slide_num - 1]
                except Exception:
                    pass
        else:
            run.hyperlink.address = url

    def _apply_animations(self, slide, design: dict[str, Any], motion: int) -> None:
        from ppt_pro_max.renderer.animation import (
            add_slide_transition, add_entrance_animation, TRANSITION_TYPES,
        )
        motion = max(1, min(10, motion))

        transition_list = list(TRANSITION_TYPES.keys())
        transition_type = transition_list[motion % len(transition_list)]
        speed = "fast" if motion >= 7 else ("medium" if motion >= 4 else "slow")
        add_slide_transition(slide, transition_type=transition_type, speed=speed)

        if motion >= 3:
            animatable_ids = []
            for shape in slide.shapes:
                if shape.shape_id and shape.has_text_frame:
                    animatable_ids.append(shape.shape_id)
            if animatable_ids:
                effect = "fade_in" if motion <= 5 else "fly_in"
                interval = max(100, 600 - motion * 50)
                for i, sid in enumerate(animatable_ids[:5]):
                    add_entrance_animation(
                        slide, sid, effect=effect,
                        delay_ms=interval * i if i > 0 else 0,
                        duration_ms=400,
                        click_triggered=(i == 0),
                    )

    def _handle_page_revision(
        self,
        template_path: str,
        pages_str: str,
        pptx_path: str,
        vnum: int,
        version_dir: str,
    ) -> dict[str, Any]:
        engine = PageRevisionEngine(template_path)
        result = engine.revise([pages_str], output_path=pptx_path)

        meta = {
            "version": vnum,
            "mode": "page_revision",
            "pages_op": pages_str,
            "num_slides": result["num_slides"],
        }
        write_meta(version_dir, meta)

        return {
            "pipeline": "enterprise",
            "mode": "page_revision",
            "output_path": result["output_path"],
            "version": vnum,
            "num_slides": result["num_slides"],
            "ops_applied": result.get("ops_applied", []),
        }

    def _build_image_fetcher(self, kwargs: dict[str, Any]):
        image_mode = kwargs.get("image_mode", "placeholder")
        if image_mode == "placeholder":
            return None

        try:
            from ppt_pro_max.renderer.image_fetcher import ImageFetcher

            image_config = dict(kwargs.get("image_config") or {})
            for key in ("llm_provider", "llm_api_key", "llm_base_url", "llm_model"):
                if key not in image_config and kwargs.get(key):
                    image_config[key] = kwargs[key]
            return ImageFetcher(
                mode=image_mode,
                unsplash_access_key=image_config.get("unsplash_access_key"),
                pexels_api_key=image_config.get("pexels_api_key"),
                llm_provider=image_config.get("llm_provider"),
                llm_api_key=image_config.get("llm_api_key"),
                llm_base_url=image_config.get("llm_base_url"),
                llm_model=image_config.get("llm_model"),
            )
        except Exception:
            return None

    def _find_latest_pptx(self, output_dir: str, current_vnum: int) -> str | None:
        if not os.path.isdir(output_dir):
            return None
        for v in range(current_vnum - 1, 0, -1):
            vdir = os.path.join(output_dir, f"v{v}")
            if not os.path.isdir(vdir):
                continue
            for name in ("presentation.pptx", f"v{v}.pptx"):
                candidate = os.path.join(vdir, name)
                if os.path.isfile(candidate):
                    return candidate
        for name in sorted(os.listdir(output_dir)):
            if name.endswith(".pptx") and not name.startswith("~"):
                candidate = os.path.join(output_dir, name)
                if os.path.isfile(candidate):
                    return candidate
        return None
