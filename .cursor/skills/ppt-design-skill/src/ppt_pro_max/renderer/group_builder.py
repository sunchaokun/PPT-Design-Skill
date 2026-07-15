"""Group Builder — create shape groups for visual hierarchy.

Enables PPT-7 style grouped elements like:
  - Glowing button: Group(OVAL gradient background + OVAL highlight)
  - Labeled card: Group(ROUNDED_RECT + TEXT_BOX + accent bar)
  - Connected nodes: Group(OVAL + OVAL + LINE connectors)

Uses python-pptx's add_group_shape() when available, falls back to
direct XML manipulation for complex nesting.
"""

from __future__ import annotations

from typing import Any

from pptx.oxml.ns import qn
from lxml import etree


EMU_PER_INCH = 914400


class GroupBuilder:

    def __init__(self):
        self._shapes: list[dict[str, Any]] = []
        self._child_groups: list[GroupBuilder] = []

    def add_shape_ref(self, shape) -> GroupBuilder:
        self._shapes.append({"type": "ref", "shape": shape})
        return self

    def build_from_shapes(self, slide, shapes: list) -> object | None:
        if len(shapes) < 2:
            return None

        try:
            group = slide.shapes.add_group_shape(shapes)
            return group
        except Exception:
            return self._build_via_xml(slide, shapes)

    def _build_via_xml(self, slide, shapes: list) -> object | None:
        if not shapes:
            return None

        sp_tree = slide.shapes._spTree
        shape_elements = []
        for s in shapes:
            try:
                shape_elements.append(s._element)
            except Exception:
                continue

        if len(shape_elements) < 2:
            return None

        min_x = min_y = float("inf")
        max_x = max_y = float("-inf")

        for elem in shape_elements:
            xfrm = elem.find(".//" + qn("a:xfrm"))
            if xfrm is None:
                continue
            off = xfrm.find(qn("a:off"))
            ext = xfrm.find(qn("a:ext"))
            if off is not None and ext is not None:
                x, y = int(off.get("x", 0)), int(off.get("y", 0))
                cx, cy = int(ext.get("cx", 0)), int(ext.get("cy", 0))
                min_x = min(min_x, x)
                min_y = min(min_y, y)
                max_x = max(max_x, x + cx)
                max_y = max(max_y, y + cy)

        if min_x == float("inf"):
            return None

        grp_cx = max_x - min_x
        grp_cy = max_y - min_y

        grpSp = etree.SubElement(sp_tree, qn("p:grpSp"))

        nvGrpSpPr = etree.SubElement(grpSp, qn("p:nvGrpSpPr"))
        cNvPr = etree.SubElement(nvGrpSpPr, qn("p:cNvPr"))
        cNvPr.set("id", str(self._next_id(slide)))
        cNvPr.set("name", "Group")
        etree.SubElement(nvGrpSpPr, qn("p:cNvSpPr"))
        etree.SubElement(nvGrpSpPr, qn("p:nvPr"))

        grpSpPr = etree.SubElement(grpSp, qn("p:grpSpPr"))
        xfrm = etree.SubElement(grpSpPr, qn("a:xfrm"))
        off = etree.SubElement(xfrm, qn("a:off"))
        off.set("x", str(min_x))
        off.set("y", str(min_y))
        ext = etree.SubElement(xfrm, qn("a:ext"))
        ext.set("cx", str(grp_cx))
        ext.set("cy", str(grp_cy))
        chOff = etree.SubElement(xfrm, qn("a:chOff"))
        chOff.set("x", str(min_x))
        chOff.set("y", str(min_y))
        chExt = etree.SubElement(xfrm, qn("a:chExt"))
        chExt.set("cx", str(grp_cx))
        chExt.set("cy", str(grp_cy))

        for elem in shape_elements:
            try:
                sp_tree.remove(elem)
            except Exception:
                pass
            grpSp.append(elem)

        return grpSp

    @staticmethod
    def _next_id(slide) -> int:
        max_id = 1
        for shape in slide.shapes:
            try:
                sid = shape.shape_id
                if sid > max_id:
                    max_id = sid
            except Exception:
                pass
        return max_id + 1


def group_shapes_on_slide(slide, shape_indices: list[int] | None = None, shape_ids: list[int] | None = None) -> object | None:
    builder = GroupBuilder()

    if shape_indices is not None:
        shapes = []
        for idx in shape_indices:
            try:
                shapes.append(slide.shapes[idx])
            except IndexError:
                continue
        return builder.build_from_shapes(slide, shapes)

    if shape_ids is not None:
        id_map = {s.shape_id: s for s in slide.shapes}
        shapes = [id_map[sid] for sid in shape_ids if sid in id_map]
        return builder.build_from_shapes(slide, shapes)

    return None


def group_last_n_shapes(slide, n: int) -> object | None:
    total = len(slide.shapes)
    if total < n or n < 2:
        return None
    indices = list(range(total - n, total))
    return group_shapes_on_slide(slide, shape_indices=indices)
