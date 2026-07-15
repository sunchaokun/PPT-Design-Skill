"""TemplateAnalyzer — extract brand spec from template .pptx."""

from __future__ import annotations

from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

from ppt_pro_max.enterprise.brand_spec import BrandSpec


_SYS_CLR_FALLBACK = {
    "windowText": "#000000",
    "WindowText": "#000000",
    "window": "#FFFFFF",
    "Window": "#FFFFFF",
}


class TemplateAnalyzer:

    def analyze(self, template_path: str) -> BrandSpec:
        try:
            prs = Presentation(template_path)
        except Exception:
            return BrandSpec(source="template_fallback")

        slide_master = prs.slide_masters[0] if prs.slide_masters else None
        if slide_master is None:
            return BrandSpec(source="template_no_master")

        colors = self._extract_colors(prs, slide_master)
        fonts = self._extract_fonts(prs, slide_master)
        layouts = self._extract_layouts(slide_master)
        dark_mode = self._detect_dark_mode(colors)

        return BrandSpec(
            source="template",
            colors=colors,
            fonts=fonts,
            template_layouts=layouts,
            dark_mode=dark_mode,
        )

    def _extract_colors(self, prs, slide_master) -> dict[str, str]:
        colors = {}
        try:
            theme_part = prs.part.part_related_by(RT.THEME)
            theme_element = etree.fromstring(theme_part.blob)
            theme_elements = theme_element.find(qn("a:themeElements"))
            if theme_elements is None:
                return colors
            clr_scheme = theme_elements.find(qn("a:clrScheme"))
            if clr_scheme is None:
                return colors

            name_map = {
                "dk1": "foreground",
                "lt1": "background",
                "dk2": "primary",
                "accent1": "accent",
            }
            for xml_name, role in name_map.items():
                elem = clr_scheme.find(qn(f"a:{xml_name}"))
                if elem is not None:
                    colors[role] = self._extract_color_value(elem)
        except Exception:
            pass

        if not colors:
            clr_map = slide_master.element.find(qn("p:clrMap"))
            if clr_map is not None:
                colors["foreground"] = clr_map.get("tx1", "dk1")
                colors["background"] = clr_map.get("bg1", "lt1")

        return colors

    def _extract_color_value(self, clr_element) -> str:
        srgb = clr_element.find(qn("a:srgbClr"))
        if srgb is not None:
            val = srgb.get("val")
            if val:
                return f"#{val}"
            return "#000000"
        sys_clr = clr_element.find(qn("a:sysClr"))
        if sys_clr is not None:
            val = sys_clr.get("val")
            if val:
                return _SYS_CLR_FALLBACK.get(val, "#000000")
            return "#000000"
        return "#000000"

    def _extract_fonts(self, prs, slide_master) -> dict[str, str]:
        fonts = {}
        try:
            theme_part = prs.part.part_related_by(RT.THEME)
            theme_element = etree.fromstring(theme_part.blob)
            theme_elements = theme_element.find(qn("a:themeElements"))
            if theme_elements is None:
                return fonts
            font_scheme = theme_elements.find(qn("a:fontScheme"))
            if font_scheme is None:
                return fonts

            major_font = font_scheme.find(qn("a:majorFont"))
            minor_font = font_scheme.find(qn("a:minorFont"))
            if major_font is not None:
                latin = major_font.find(qn("a:latin"))
                if latin is not None:
                    fonts["heading"] = latin.get("typeface", "")
            if minor_font is not None:
                latin = minor_font.find(qn("a:latin"))
                if latin is not None:
                    fonts["body"] = latin.get("typeface", "")
        except Exception:
            pass

        return fonts

    def _extract_layouts(self, slide_master) -> list[dict[str, Any]]:
        layouts = []
        for i, layout in enumerate(slide_master.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                })
            layouts.append({
                "index": i,
                "name": layout.name,
                "placeholders": placeholders,
            })
        return layouts

    def _detect_dark_mode(self, colors: dict[str, str]) -> bool:
        bg = colors.get("background", "#FFFFFF")
        try:
            r, g, b = int(bg[1:3], 16), int(bg[3:5], 16), int(bg[5:7], 16)
            return (r * 0.299 + g * 0.587 + b * 0.114) < 128
        except (ValueError, IndexError):
            return False
